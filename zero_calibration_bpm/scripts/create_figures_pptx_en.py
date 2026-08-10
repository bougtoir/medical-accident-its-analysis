#!/usr/bin/env python3
"""Editable English figure deck (.pptx): one figure per slide (16:9)."""

import os

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
FIGDIR = os.path.join(SCRIPT_DIR, "..", "figures")
OUTDIR = os.path.join(SCRIPT_DIR, "..", "manuscripts")
OUTPATH = os.path.join(OUTDIR, "TIM_Figures_EN.pptx")
os.makedirs(OUTDIR, exist_ok=True)

# Seven inline figures for the IEEE TIM manuscript.
SLIDES = [
    ("figure1_signal_decomposition.png", "Figure 1. Signal decomposition",
     "A DC offset shifts the arterial pressure baseline but leaves pulse "
     "pressure unchanged; a gain error scales the whole waveform and changes "
     "pulse pressure. Zeroing corrects the offset but not the gain."),
    ("figure2_scenarios_concordance.png", "Figure 2. Concordance by scenario",
     "Device-versus-reference concordance for the four static scenarios "
     "(dashed line = identity). Each panel shows the CCC, the CCC scale shift "
     "v, and the Bland\u2013Altman mean bias. S4 has a near-zero mean bias despite "
     "a 10% gain error."),
    ("figure3_detection_panel.png", "Figure 3. Detection pattern",
     "Which reported analysis detects the error in each scenario (green = "
     "detected, grey = missed). In S4 the mean-bias summary misses the gain "
     "error that every proportional-bias\u2013aware analysis detects."),
    ("figure4_ba_masked_gain.png", "Figure 4. Bland\u2013Altman regression",
     "Difference-versus-mean plots. In S2 the regression slope is flat; in S4 "
     "the slope is clearly positive despite a near-zero mean bias, revealing "
     "the masked proportional (gain) error."),
    ("figure5_dynamic_response.png", "Figure 5. Dynamic response",
     "(A) Frequency response of optimal, under-damped and over-damped systems "
     "with arterial harmonics overlaid. (B) Example waveforms. (C) Measured-"
     "versus-true pulse pressure and the mean PP ratio for each system."),
    ("figure6_range_dependence.png", "Figure 6. Range-dependence of the CCC",
     "For one fixed device, the CCC and the bias-correction factor C_b "
     "increase as the sampled pressure range widens, while the scale shift v "
     "is less range-dependent and approaches the true gain ratio at wider ranges."),
    ("figure7_real_validation.png", "Figure 7. Real-waveform validation",
     "(A) A 10-s segment of the SNUADC/ART waveform from the VitalDB Open "
     "Dataset with detected beats. (B) Bland\u2013Altman plot for the gain-masked "
     "scenario (R4), showing a near-zero mean bias with a clear proportional "
     "bias."),
]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SW, SH = 13.333, 7.5
for fname, title, caption in SLIDES:
    slide = prs.slides.add_slide(blank)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(SW - 1.0),
                                  Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = title
    run.font.size = Pt(24); run.font.bold = True

    path = os.path.join(FIGDIR, fname)
    top = 1.05
    avail_h = 5.0
    avail_w = SW - 1.0
    if os.path.exists(path):
        iw, ih = Image.open(path).size
        aspect = iw / ih
        w = avail_w
        h = w / aspect
        if h > avail_h:
            h = avail_h
            w = h * aspect
        left = (SW - w) / 2.0
        slide.shapes.add_picture(path, Inches(left), Inches(top),
                                 width=Inches(w), height=Inches(h))

    cb = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(SW - 1.0),
                                  Inches(1.0))
    ctf = cb.text_frame
    ctf.word_wrap = True
    cp = ctf.paragraphs[0]
    crun = cp.add_run(); crun.text = caption
    crun.font.size = Pt(13); crun.font.italic = True

prs.save(OUTPATH)
print(f"Figure deck saved: {OUTPATH}")
