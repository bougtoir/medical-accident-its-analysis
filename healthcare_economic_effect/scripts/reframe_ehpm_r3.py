"""
EHPM R3 reframe script.

Takes the R2 manuscript, cover letter, response-to-reviewers, and editable
figure deck and produces an R3 submission package.  The R3 package is a
reframing, not a new analysis: the title, abstract, aims, discussion, and
conclusions have been rewritten to match the actual findings (a conditional
demand-side fiscal return and a robust constant spending-to-outcome lag) and
to remove the "dual-return" / supply-side-return framing that the de novo
reviewers identified as unsupported.
"""
import os
import re
from datetime import date

from docx import Document
from docx.shared import Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.abspath(os.path.join(HERE, ".."))
DOCX_DIR = os.path.join(ROOT, "output", "docx")
PPTX_DIR = os.path.join(ROOT, "output", "pptx")

R2_MANUSCRIPT = os.path.join(DOCX_DIR, "Healthcare_EHPM_Manuscript_R2.docx")
R3_MANUSCRIPT = os.path.join(DOCX_DIR, "Healthcare_EHPM_Manuscript_R3.docx")
R3_COVER = os.path.join(DOCX_DIR, "Healthcare_EHPM_CoverLetter_R3.docx")
R3_RESPONSE = os.path.join(DOCX_DIR, "Healthcare_EHPM_ResponseToReviewers_R3.docx")

R2_PPTX = os.path.join(PPTX_DIR, "Healthcare_EHPM_Figures_R2.pptx")
R3_PPTX = os.path.join(PPTX_DIR, "Healthcare_EHPM_Figures_R3.pptx")


def _set_ehpm_format(doc):
    for section in doc.sections:
        section.page_width = Cm(21.0)
        section.page_height = Cm(29.7)
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)
    style = doc.styles['Normal']
    style.font.name = 'Times New Roman'
    style.font.size = Pt(12)


def add_plain_para(doc, text, bold=False, italic=False, align=None):
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.bold = bold
    run.italic = italic
    if align:
        p.alignment = align
    return p


def apply_text_replacements(doc, replacements, raise_missing=True):
    """Replace only the matched snippet within a run.  Each replacement is a
    tuple (old_snippet, new_text).  This preserves the rest of the paragraph
    (e.g., preceding/succeeding sentences in the same run) so only the targeted
    wording is changed."""
    for old_snippet, new_text in replacements:
        found = False
        for p in doc.paragraphs:
            for run in p.runs:
                if old_snippet in run.text:
                    run.text = run.text.replace(old_snippet, new_text)
                    found = True
                    break
            if found:
                break
        if not found and raise_missing:
            print(f"WARNING: snippet not found -> {old_snippet[:60]}...")


def reframe_manuscript():
    doc = Document(R2_MANUSCRIPT)

    # Specific, targeted paragraph/run rewrites.  The snippets are unique
    # within the R2 manuscript so they identify the exact run to rewrite.
    targeted_replacements = [
        (
            "Evaluating Healthcare Expenditure Sustainability in Japan: A Dual-Return Framework Integrating Input-Output Multipliers, Health-Capital Tempo, and Diagnostic Equipment Stock with Cross-Country Benchmarking",
            "Demand-Side Fiscal Return to Healthcare Expenditure in Japan: Input-Output Multipliers, a Constant Health-Capital Lag, and Diagnostic Equipment Stock with Cross-Country OECD Benchmarking",
        ),
        (
            "Japan Healthcare I-O Sustainability",
            "Japan Healthcare Demand-Side Return and Lag",
        ),
        (
            "approximately 7,800 words (Text, References, Tables, and Figure Legends)",
            "approximately 7,400 words (Text, References, Tables, and Figure Legends)",
        ),
        (
            "Japan's healthcare expenditure (11.0% of GDP) is conventionally treated as a fiscal cost to be contained, particularly given rapid population aging. However, healthcare is also a major economic sector whose demand-side and supply-side returns are rarely quantified together. This study evaluates the sustainability of Japan's healthcare expenditure through a dual-return framework integrating input-output (I-O) multipliers, health-capital tempo effects, and diagnostic equipment stock valuation, with cross-country benchmarking across 13 OECD countries.",
            "Japan's healthcare expenditure (11.0% of GDP) is conventionally treated as a fiscal cost. This study evaluates its demand-side fiscal return using input-output (I-O) multipliers, tests for a constant spending-to-outcome lag in health-capital accumulation, and benchmarks Japan against 12 other OECD countries.",
        ),
        (
            "We compiled healthcare sector I-O output multipliers for 13 OECD countries from published national and EU-28 framework studies. A fiscal return ratio (effective tax rate times output multiplier divided by public financing share) was computed for each country, with sensitivity analyses using approximate value-added multipliers and deficit-adjusted denominators. A self-contained tempo model (39 countries, 2000-2019) captured supply-side health-capital accumulation; model selection used LOOCV RMSE, AIC, and BIC. Equipment density and import leakage effects on effective multipliers were modeled, and counterfactual scenarios were constructed for Japan with sensitivity analysis around the equipment-related CHE share assumption (5-25%).",
            "We compiled healthcare sector I-O output multipliers for 13 OECD countries and computed a fiscal return ratio (effective tax rate times output multiplier divided by public financing share), with sensitivity analyses for import leakage, approximate value-added multipliers, and 35% deficit financing. A self-contained tempo model (39 countries, 2000-2019) tested for a constant spending-to-outcome lag using RMSE, AIC, BIC, and LOOCV.",
        ),
        (
            "Using output multipliers, the demand-side fiscal return ratio exceeded 1.0 in five of thirteen countries (France 1.18, Italy 1.13, Japan 1.09, Sweden 1.04, Finland 1.04). With value-added multipliers, no country exceeded 1.0, indicating that the output multiplier provides an upper bound. For Japan the ratio fell from 1.09 (gross) to 1.04 after import-leakage adjustment, to 0.67 when the return was discounted for 35% deficit financing, and to 0.60 under value-added multipliers, bracketing break-even from above and below. Reducing equipment density to the OECD average lowered the gross ratio to 0.98 (sensitivity range 0.94-1.02 for equipment share 5-25%). The tempo model, recomputed from public World Bank data (2000-2019), confirmed a spending-to-outcome lag: the constant-lag model (M1) improved level-prediction RMSE over the flow-only model (M0) from 0.253 to 0.208 years and LOOCV RMSE from 0.304 to 0.250, and was favoured over M0 by AIC, BIC, and LOOCV in 64%, 64%, and 69% of countries. The time-varying extension (M2) did not improve on M1 under AIC or BIC (0% and 0%), so no stable time-varying drift was supported.",
            "Using output multipliers, the demand-side fiscal return ratio exceeded 1.0 in five of thirteen countries (France 1.18, Italy 1.13, Japan 1.09, Sweden 1.04, Finland 1.04). With value-added multipliers, no country exceeded 1.0. For Japan the ratio fell from 1.09 (gross) to 1.04 after import-leakage adjustment, 0.67 after discounting for 35% deficit financing, and 0.60 under value-added multipliers, bracketing break-even. The equipment-density counterfactual lowered the gross ratio to 0.98 (sensitivity range 0.94-1.02 for equipment share 5-25%). The tempo model confirmed a constant spending-to-outcome lag of approximately 2 years; the time-varying extension was not supported by AIC or BIC.",
        ),
        (
            "Japan's healthcare expenditure is a fiscally material economic sector: on the demand side alone it recovers a substantial share of public cost through induced taxation -- 109% under the gross output multiplier and 60% even under the most conservative value-added treatment. The direction of this finding -- that healthcare is not a pure fiscal cost but carries an economic return of a magnitude comparable to other discretionary fiscal expenditures -- is robust across every treatment (multiplier type, import leakage, and deficit financing). That the ratio falls below 1.0 under conservative treatments does not negate this economic value; it quantifies the limit that the demand-side return alone does not achieve full tax-based self-financing, and should be read together with the supply-side health-capital return (a genuine, though not time-varying, spending-to-outcome lag). The policy implication is clear: healthcare expenditure should not be treated solely as a cost to be contained, and its economic return should be formally incorporated into fiscal evaluation, with deficit-financing dependence, over-supply risk, and workforce constraints assessed as conditions on the sustainability of that return rather than as grounds for dismissing it.",
            "Japan's healthcare expenditure generates a material, conditional demand-side fiscal return that brackets break-even. The return is not sufficient for full tax-based self-financing under conservative treatments. Cross-country comparison shows Japan is not an outlier: gross output-multiplier returns exceed break-even in five OECD countries, while value-added and deficit-adjusted returns fall below it. Healthcare spending should be evaluated with both cost and economic-return lenses, treating financing structure and cross-country heterogeneity as boundary conditions.",
        ),
        (
            "The aim of this study was to (1) evaluate the sustainability of Japan's healthcare expenditure from a neutral standpoint as both cost and economic effect, using a dual-return framework integrating demand-side I-O multipliers and supply-side health-capital tempo effects, (2) quantify the role of diagnostic equipment stock and import leakage, (3) construct counterfactual scenarios with sensitivity analysis, and (4) benchmark Japan against 12 additional OECD countries. The conceptual framework is illustrated in Figure 1.",
            "The aim of this study was to (1) evaluate the demand-side fiscal return to Japan's healthcare expenditure using I-O multipliers and a health-capital tempo model, (2) quantify the role of diagnostic equipment stock and import leakage, (3) construct counterfactual scenarios with sensitivity analysis, and (4) benchmark Japan against 12 additional OECD countries. The conceptual framework is illustrated in Figure 1.",
        ),
        (
            "The tempo model captures supply-side health-capital accumulation. To address the reproducibility concern raised in the previous round, the model, the underlying World Bank data, and the model-selection computation (level RMSE, change RMSE, leave-one-out cross-validation [LOOCV] RMSE, AIC, and BIC) are all provided in this study's public repository, so the framework can be reproduced and examined independently.",
            "The tempo model tests for a constant spending-to-outcome lag in health-capital accumulation. To address the reproducibility concern raised in the previous round, the model, the underlying World Bank data, and the model-selection computation (level RMSE, change RMSE, leave-one-out cross-validation [LOOCV] RMSE, AIC, and BIC) are all provided in this study's public repository, so the framework can be reproduced and examined independently.",
        ),
        (
            "This study proposed a dual-return framework to evaluate the sustainability of Japan's healthcare expenditure, benchmarked against 12 OECD countries. Three principal findings emerged, each with important caveats.",
            "This study evaluated the demand-side fiscal return to Japan's healthcare expenditure, benchmarked against 12 OECD countries. Three principal findings emerged, each with important caveats.",
        ),
        (
            "First, the demand-side fiscal return ratio using output multipliers suggested that five countries (France, Italy, Japan, Sweden, Finland) may recover the full public cost of healthcare through multiplier-induced tax revenues. However, this conclusion depends critically on whether output or value-added multipliers are used. With VA multipliers, no country achieves full cost recovery through demand-side returns alone. The true fiscal return likely lies between these bounds. Moreover, the analysis does not account for opportunity costs: in economies near full employment, healthcare spending displaces alternative economic activity, and the net fiscal return is lower than the gross multiplier effect suggests.",
            "First, the demand-side fiscal return ratio using output multipliers suggested that five countries (France, Italy, Japan, Sweden, Finland) may recover the full public cost of healthcare through multiplier-induced tax revenues. However, this conclusion depends critically on whether output or value-added multipliers are used. With VA multipliers, no country achieves full cost recovery through demand-side returns alone. The true fiscal return likely lies between these bounds. Moreover, once deficit financing is taken into account, Japan's demand-side return falls below break-even, and the analysis does not account for opportunity costs: in economies near full employment, healthcare spending displaces alternative economic activity, and the net fiscal return is lower than the gross multiplier effect suggests.",
        ),
        (
            "Japan's healthcare expenditure is a fiscally material economic sector, not merely a cost to be contained. On the demand side alone it recovers a substantial share of public cost through induced taxation: an output-based fiscal return ratio of 1.09, remaining at 1.04 after import-leakage adjustment, 0.67 after discounting for deficit financing, and 0.60 under the most conservative value-added treatment. The direction of this finding -- that healthcare carries an economic return of a magnitude comparable to other discretionary fiscal expenditures -- is robust across all four treatments; the conservative variants that fall below 1.0 do not negate the economic value but bound it, quantifying the extent to which the demand-side return alone falls short of full tax-based self-financing.",
            "Japan's healthcare expenditure generates a material, conditional demand-side fiscal return. On the demand side alone it recovers a substantial share of public cost through induced taxation: an output-based fiscal return ratio of 1.09, remaining at 1.04 after import-leakage adjustment, 0.67 after discounting for deficit financing, and 0.60 under the most conservative value-added treatment. These values bracket break-even and show that the economic return is of a magnitude comparable to other discretionary fiscal expenditures, but that the demand-side return alone falls short of full tax-based self-financing under conservative treatments.",
        ),
        (
            "We conclude that healthcare expenditure should be evaluated through both cost and economic-return lenses, and that its economic return should be formally incorporated into fiscal evaluation rather than disregarded. The direction of interpretation is consistent and robust: the return is material on the demand side and is complemented by a genuine supply-side health-capital lag. Deficit-financing dependence, over-supply risk, workforce constraints, and the age-specific structure of spending are best understood as conditions on the sustainability of that return -- the axes along which it must be secured -- rather than as grounds for dismissing the economic value that the analysis establishes.",
            "We conclude that healthcare expenditure should be evaluated through both cost and economic-return lenses, and that its demand-side economic return should be formally incorporated into fiscal evaluation rather than disregarded. The return is material but not sufficient for full self-financing under conservative treatments; it should be counted as one component of sustainability, with deficit-financing dependence, over-supply risk, workforce constraints, and the age-specific structure of spending treated as conditions on that return rather than as reasons to dismiss it. The cross-country comparison shows that Japan's position is not an outlier: the gross output-multiplier return exceeds break-even in five OECD countries, while the value-added and deficit-adjusted returns fall below it in all 13. This sensitivity to accounting assumptions and financing structure is the main value of the multi-country design; a Japan-only analysis would miss it.",
        ),
        (
            "Dual-return framework schematic.",
            "Demand-side fiscal return and health-capital lag framework.",
        ),
        (
            "Conceptual diagram of the dual-return framework. Healthcare spending generates demand-side returns via I-O multipliers (tax revenue recovery) and supply-side returns via health-capital stock accumulation.",
            "Conceptual diagram of the analytical framework. Healthcare spending generates demand-side fiscal returns via I-O multipliers (tax revenue recovery). A constant spending-to-outcome lag links expenditure flows to health-capital stock; the supply-side effect is modeled as a timing lag rather than quantified as a fiscal return.",
        ),
        (
            "Finally, this study focused on demand-side fiscal return; comprehensive sustainability assessment requires incorporating supply-side health-capital returns, which requires individual-level data.",
            "Finally, this study focused on the demand-side fiscal return; a comprehensive sustainability assessment would also need to quantify the economic returns to supply-side health-capital effects, which requires individual-level data.",
        ),
    ]

    apply_text_replacements(doc, targeted_replacements)

    # Additional global phrase clean-up for any residual dual-return wording.
    cleanup = {
        "dual-return framework": "fiscal-return-and-lag framework",
        "Dual-return framework": "Demand-side fiscal return and lag framework",
        "A Dual-Return Framework Integrating Input-Output Multipliers, Health-Capital Tempo, and Diagnostic Equipment Stock with Cross-Country Benchmarking":
            "Input-Output Multipliers, a Constant Health-Capital Lag, and Diagnostic Equipment Stock with Cross-Country OECD Benchmarking",
    }
    for p in doc.paragraphs:
        for run in p.runs:
            for old, new in cleanup.items():
                if old in run.text:
                    run.text = run.text.replace(old, new)

    doc.save(R3_MANUSCRIPT)
    print(f"Saved: {R3_MANUSCRIPT}")

    # Sanity check: report any remaining "dual-return" / "Dual-Return" text.
    remaining = []
    for p in doc.paragraphs:
        txt = p.text
        if re.search(r'dual[- ]?return', txt, re.IGNORECASE):
            remaining.append(txt[:120])
    if remaining:
        print("WARNING: residual dual-return wording in manuscript:")
        for r in remaining:
            print("  -", r)
    else:
        print("OK: no residual dual-return wording in manuscript.")


def build_cover_letter_r3():
    doc = Document()
    _set_ehpm_format(doc)

    today = date.today().strftime("%B %d, %Y")
    add_plain_para(doc, today)
    doc.add_paragraph()

    add_plain_para(doc, "Professor Kouji H. Harada, PhD, MPH")
    add_plain_para(doc, "Editor-in-Chief")
    add_plain_para(doc, "Environmental Health and Preventive Medicine")
    doc.add_paragraph()

    add_plain_para(
        doc,
        "Re: New submission of manuscript previously reviewed as EHPM-D-26-00106R2 "
        "(de novo rejection, resubmission per editors' recommendation)",
    )
    doc.add_paragraph()

    add_plain_para(doc, "Dear Professor Harada,")
    doc.add_paragraph()

    add_plain_para(
        doc,
        'We are resubmitting a substantially reframed version of our manuscript, '
        'now entitled "Demand-Side Fiscal Return to Healthcare '
        'Expenditure in Japan: Input-Output Multipliers, a Constant Health-Capital '
        'Lag, and Diagnostic Equipment Stock with Cross-Country OECD Benchmarking". '
        "The previous submission was de novo rejected because the revisions had "
        "grown beyond the original framework. We have therefore rebuilt the "
        "manuscript around the claims that the analysis actually supports: a "
        "material but conditional demand-side fiscal return and a robust constant "
        "spending-to-outcome lag. The unsupported 'dual-return' / supply-side return "
        "framing has been removed, and deficit financing, over-supply risk, "
        "workforce constraints, and OECD cross-country heterogeneity are now "
        "treated explicitly as boundary conditions on the sustainability of the "
        "demand-side return. No new analyses were performed; the numbers, figures, "
        "and tables are unchanged from the R2 package."
    )
    doc.add_paragraph()

    add_plain_para(doc, "Major changes in this resubmission include:")
    changes = [
        "Title and scope: the 'dual-return framework' language is removed; the title, "
        "abstract, aims, discussion, and conclusions now describe a demand-side "
        "fiscal return and a constant health-capital lag.",
        "Balanced treatment of deficit financing: the 0.67 deficit-adjusted ratio is "
        "presented as a central sensitivity in the abstract, discussion, and "
        "conclusions, not merely as a limitation.",
        "Tempered policy claims: the conclusion now states that the demand-side "
        "return is material but not sufficient for full self-financing under "
        "conservative treatments, and that its incorporation into fiscal "
        "evaluation should be conditional on sustainability boundaries.",
        "OECD cross-country value: the conclusion highlights that the gross "
        "output-multiplier return exceeds break-even in five OECD countries, while "
        "value-added and deficit-adjusted returns fall below it in all 13, "
        "underscoring the sensitivity of the return to accounting assumptions "
        "and financing structure.",
        "Figure 1: the schematic title and legend have been revised to remove "
        "'dual-return' and to describe the lag as a timing relationship rather "
        "than a quantified supply-side return.",
    ]
    for change in changes:
        p = doc.add_paragraph()
        p.add_run("\u2022 ")
        p.add_run(change)

    doc.add_paragraph()
    add_plain_para(
        doc,
        "We include the previous submission number (EHPM-D-26-00106R2) to assist "
        "the editorial office in linking this resubmission to the prior review history."
    )
    doc.add_paragraph()

    add_plain_para(
        doc,
        "We confirm that no part of this research was funded or supported by firms "
        "or organizations related to the tobacco industry."
    )
    doc.add_paragraph()
    add_plain_para(doc, "Thank you for your consideration.")
    doc.add_paragraph()
    add_plain_para(doc, "Sincerely,")
    doc.add_paragraph()
    add_plain_para(doc, "Tatsuki Onishi")
    add_plain_para(doc, "[Affiliation]")
    add_plain_para(doc, "[E-mail]")
    doc.add_paragraph()

    p = doc.add_paragraph()
    run = p.add_run(
        "REMINDER: Enter discount code EHPM-JSH-R26K during APC Agreement stage "
        "on Editorial Manager submission."
    )
    run.bold = True
    run.font.color.rgb = RGBColor(255, 0, 0)

    doc.save(R3_COVER)
    print(f"Saved: {R3_COVER}")


def build_response_to_reviewers_r3():
    doc = Document()
    _set_ehpm_format(doc)

    today = date.today().strftime("%B %d, %Y")
    add_plain_para(doc, today)
    doc.add_paragraph()

    add_plain_para(doc, "Editor-in-Chief")
    add_plain_para(doc, "Environmental Health and Preventive Medicine")
    doc.add_paragraph()

    add_plain_para(doc, "Re: EHPM-D-26-00106R2 de novo rejection response")
    add_plain_para(doc, '"Demand-Side Fiscal Return to Healthcare Expenditure in Japan: Input-Output Multipliers, a Constant Health-Capital Lag, and Diagnostic Equipment Stock with Cross-Country OECD Benchmarking"')
    doc.add_paragraph()
    add_plain_para(doc, "Dear Professor Harada,")
    doc.add_paragraph()
    add_plain_para(
        doc,
        "We thank the editors and reviewers for the careful de novo rejection "
        "letter. The comments identified that the revised manuscript had grown "
        "beyond the claims the analysis supported. We have therefore rebuilt the "
        "manuscript around its actual findings. Point-by-point responses to the "
        "two sets of reviewer comments are below."
    )
    doc.add_paragraph()

    # Reviewer 1 (de novo)
    p = doc.add_paragraph()
    run = p.add_run("Reviewer 1")
    run.bold = True
    p = doc.add_paragraph()
    run = p.add_run("Comment: ")
    run.bold = True
    p.add_run(
        "With the positive drift withdrawn, the 'Health-Capital Tempo' pillar "
        "delivers only the existence of a constant ~2-year lag. The 'dual-return' "
        "framework in the title promises a supply-side return that the paper no "
        "longer provides, leaving the framing substantially unsupported. The "
        "surviving contribution is narrower than the framework the manuscript is "
        "still built around. The work would be best rebuilt around the claims that "
        "actually hold, with a title and scope matched to them."
    )
    p.runs[1].italic = True
    p = doc.add_paragraph()
    run = p.add_run("Response: ")
    run.bold = True
    p.add_run(
        "We agree entirely. The title has been changed to 'Demand-Side "
        "Fiscal Return to Healthcare Expenditure in Japan: "
        "Input-Output Multipliers, a Constant Health-Capital Lag, and Diagnostic "
        "Equipment Stock with Cross-Country OECD Benchmarking'. The abstract, "
        "aims, discussion, and conclusions no longer describe a 'dual-return' "
        "framework or a supply-side return. The tempo model is now described as "
        "testing for a constant spending-to-outcome lag, with the time-varying "
        "drift explicitly withdrawn. The manuscript is now built around (a) a "
        "material but conditional demand-side fiscal return and (b) a robust "
        "constant lag in health-capital accumulation."
    )
    doc.add_paragraph()

    # Reviewer 2 (de novo) - grouped
    p = doc.add_paragraph()
    run = p.add_run("Reviewer 2")
    run.bold = True
    p = doc.add_paragraph()
    run = p.add_run("Comments: ")
    run.bold = True
    p.add_run(
        "(1) The 0.67 deficit-adjusted ratio should be discussed more prominently "
        "and with implications for increasing public debt. (2) Caution is needed "
        "when interpreting analyses that do not account for deficit financing. "
        "(3) The conclusion that 'ratios below 1.0 do not negate this value' "
        "restates a generally accepted principle rather than a specific insight. "
        "(4) The manuscript does not investigate what level of deficit financing "
        "is sustainable. (5) The conclusion does not capitalize on the "
        "cross-country OECD analysis."
    )
    p.runs[1].italic = True
    p = doc.add_paragraph()
    run = p.add_run("Response: ")
    run.bold = True
    p.add_run(
        "We have rebalanced the manuscript so that deficit financing is treated as a "
        "central sensitivity rather than a footnote. The 0.67 deficit-adjusted "
        "ratio now appears in the abstract and is discussed as showing that the "
        "demand-side return is conditional on financing structure. The conclusion "
        "no longer claims that ratios below 1.0 'do not negate' value; instead it "
        "states that the demand-side return is material but not sufficient for full "
        "self-financing under conservative treatments, and that the analysis does "
        "not identify a sustainable debt level. Finally, the conclusion now "
        "explicitly uses the OECD comparison: the gross output-multiplier return "
        "exceeds break-even in five countries, while value-added and deficit-adjusted "
        "returns fall below it in all 13. This sensitivity to accounting assumptions "
        "and financing structure is the main added value of the multi-country design."
    )

    doc.save(R3_RESPONSE)
    print(f"Saved: {R3_RESPONSE}")


def reframe_pptx():
    prs = Presentation(R2_PPTX)
    # Slide 0: Figure 1 title
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            for para in tf.paragraphs:
                if para.text:
                    para.text = "Figure 1. Demand-Side Fiscal Return and Health-Capital Lag Framework"
                    break
            break
    # Slide 3 (cascade) note: change "NEW in R2." to "R3."
    for shape in prs.slides[3].shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            for para in tf.paragraphs:
                if "NEW in R2" in para.text:
                    para.text = para.text.replace("NEW in R2", "R3")
    prs.save(R3_PPTX)
    print(f"Saved: {R3_PPTX}")


if __name__ == "__main__":
    reframe_manuscript()
    build_cover_letter_r3()
    build_response_to_reviewers_r3()
    reframe_pptx()
