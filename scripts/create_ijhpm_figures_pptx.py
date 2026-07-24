#!/usr/bin/env python3
"""Create editable PowerPoint figures for IJHPM submission.

Reads output/ijhpm_results.json for caption placeholders and reuses existing
English figure PNGs in output/.
"""
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), '..', 'documents', 'IJHPM')
os.makedirs(OUTPUT_DIR, exist_ok=True)
REPO_ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
FIG_DIR = os.path.join(REPO_ROOT, 'output')

with open(os.path.join(REPO_ROOT, 'output', 'ijhpm_results.json')) as f:
    R = json.load(f)

meta = R['metadata']
n_areas = meta['n_areas']
n_univ = meta['n_univ_areas']
fiscal_year = meta.get('fiscal_year', 2022)

def build_deck(figures, out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for idx, (img, title, caption) in enumerate(figures, 1):
        slide = prs.slides.add_slide(blank_layout)
        title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2),
                                             Inches(12.5), Inches(0.6))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.bold = True
        run.font.size = Pt(22)
        run.font.name = 'Arial'

        img_path = os.path.join(FIG_DIR, img)
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(1.8), Inches(1.0),
                                     width=Inches(9.7))
        else:
            print(f"Warning: missing image {img_path}")

        cap_box = slide.shapes.add_textbox(Inches(0.4), Inches(6.4),
                                           Inches(12.5), Inches(1.0))
        cf = cap_box.text_frame
        cf.word_wrap = True
        cp = cf.paragraphs[0]
        crun = cp.add_run()
        crun.text = caption
        crun.font.size = Pt(12)
        crun.font.name = 'Arial'

    prs.save(out_path)
    print(f"Saved: {out_path}")


figures_en = [
    (
        'rapm_fig1_en.png',
        f'Figure 1. Geographic distribution of anaesthesia standardised claim ratios across {n_areas} secondary medical areas of Japan, fiscal year {fiscal_year}',
        '(A) General anaesthesia (L008). (B) Spinal anaesthesia (L004). '
        '(C) Epidural anaesthesia as main anaesthetic (L002). '
        '(D) Continuous epidural infusion (L003). Choropleth maps shaded '
        'by quintile of the standardised claim ratio (national average = 100). '
        'Red circles mark secondary medical areas containing at least one '
        'university hospital. Areas masked by the data provider owing to low '
        'volume are shown in grey.',
    ),
    (
        'rapm_fig2_en.png',
        'Figure 2. University hospital presence and the combined general-anaesthesia plus continuous-epidural measure.',
        f'(A) Distribution of secondary medical areas containing at least one '
        f'university hospital (n = {n_univ} of {n_areas}; red). (B) Choropleth map of the '
        f'combined general-anaesthesia plus continuous-epidural standardised '
        f'claim ratio (mean of L008 and L003 SCR; {R["combined"]["L008_L003"]["n"]} areas with data for '
        f'both codes), shaded by quintile. Red circles mark secondary medical '
        f'areas containing at least one university hospital. Areas masked by the '
        f'data provider for either code are shown in grey.',
    ),
]

build_deck(figures_en,
           os.path.join(OUTPUT_DIR, 'regional_anaesthesia_figures_IJHPM_EN.pptx'))
