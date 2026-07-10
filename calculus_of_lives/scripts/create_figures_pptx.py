#!/usr/bin/env python3
"""
Build an editable English PPTX of the manuscript figures (one per slide,
with title and caption). Outputs ../output/figures_en.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from create_manuscript_en import FIG_CAPTIONS

OUT = Path(__file__).resolve().parent.parent / "output"

SLIDES = [
    ("fig1_layers", "Figure 1. Two layers of the question", FIG_CAPTIONS[1]),
    ("fig2_quadrant", "Figure 2. Supply\u2013demand \u00d7 expansion\u2013contraction",
     FIG_CAPTIONS[2]),
    ("fig3_asymptote", "Figure 3. Asymptotic relief from the calculus",
     FIG_CAPTIONS[3]),
]


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for stem, title, caption in SLIDES:
        slide = prs.slides.add_slide(blank)

        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = title
        r.font.size = Pt(24); r.font.bold = True
        r.font.color.rgb = RGBColor(0x1a, 0x1a, 0x1a)
        p.alignment = PP_ALIGN.LEFT

        img = OUT / f"{stem}.png"
        from PIL import Image
        with Image.open(img) as im:
            w, h = im.size
        avail_w, avail_h = Inches(11.5), Inches(5.0)
        ratio = min(avail_w / (w * 9525), avail_h / (h * 9525))
        pic_w = int(w * 9525 * ratio)
        pic_h = int(h * 9525 * ratio)
        left = int((prs.slide_width - pic_w) / 2)
        slide.shapes.add_picture(str(img), left, Inches(1.15), width=pic_w, height=pic_h)

        cb = slide.shapes.add_textbox(Inches(0.6), Inches(6.35), Inches(12.1), Inches(1.0))
        cf = cb.text_frame; cf.word_wrap = True
        cp = cf.paragraphs[0]
        cr = cp.add_run(); cr.text = caption
        cr.font.size = Pt(12); cr.font.color.rgb = RGBColor(0x40, 0x40, 0x40)

    prs.save(OUT / "figures_en.pptx")
    print("figures_en.pptx written to", OUT)


if __name__ == "__main__":
    build()
