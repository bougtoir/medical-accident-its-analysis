#!/usr/bin/env python3
"""Generate editable PPTX with one figure per slide for EEH submission."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

FIG_DIR = Path(__file__).parent / "figures"
OUTPUT = Path(__file__).parent / "figures_pptx.pptx"

FIGURES = [
    {
        'file': 'Fig1.png',
        'title': 'Fig. 1',
        'caption': 'Conquest rates by closure classification across six sensitivity scenarios. '
                   'Bars show conquest rates for closed (including technically excluded) vs. open polities. '
                   'Significance thresholds: *p < 0.05, **p < 0.01.',
    },
    {
        'file': 'Fig2.png',
        'title': 'Fig. 2',
        'caption': 'Progressive strengthening of Fisher\'s exact test p-values as technically excluded '
                   'polities are reclassified from open to closed. The dashed line indicates p = 0.05.',
    },
    {
        'file': 'Fig3.png',
        'title': 'Fig. 3',
        'caption': 'Conquest rates disaggregated by closure type under the 7-country reclassification. '
                   'Technical exclusion (zero technology flow) produces the highest conquest rate, '
                   'followed by policy bans, sakoku, and bloc closures.',
    },
    {
        'file': 'Fig4.png',
        'title': 'Fig. 4',
        'caption': 'Forest plot of multivariate logistic regression coefficients (log-odds) under the '
                   '7-country reclassification with disrupted = overtaken. External threat and institutional '
                   'quality are the strongest predictors; the network closure indicator is not independently significant.',
    },
]


def main():
    prs = Presentation()
    # Widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for fig in FIGURES:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Title at top
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = fig['title']
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Image centered
        img_path = FIG_DIR / fig['file']
        if img_path.exists():
            # Calculate proportional sizing (max width 11", max height 5")
            pic = slide.shapes.add_picture(
                str(img_path),
                Inches(1.5), Inches(1.0),
                width=Inches(10), height=None
            )
            # Constrain height
            if pic.height > Inches(5):
                ratio = Inches(5) / pic.height
                pic.width = int(pic.width * ratio)
                pic.height = Inches(5)
            # Center horizontally
            pic.left = int((prs.slide_width - pic.width) / 2)

        # Caption at bottom
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12.3), Inches(1.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = fig['caption']
        p.font.size = Pt(11)
        p.alignment = PP_ALIGN.LEFT

    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")


if __name__ == '__main__':
    main()
