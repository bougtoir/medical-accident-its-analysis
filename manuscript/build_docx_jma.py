#!/usr/bin/env python3
"""
Build English manuscript (.docx) for JMA Journal submission.

JMA Journal — Original Research Article requirements:
  - Headings: Introduction, Materials and Methods, Results, Discussion
  - Word limit: 3,500 words (excluding abstract, references, tables, figure legends)
  - Structured abstract ≤300 words (Introduction, Methods, Results, Conclusions)
  - 3-8 keywords
  - Tables: ≤5
  - Figures: ≤6
  - References: superscript citation, Vancouver style, ≤3 authors then "et al."
  - Double-spaced, Times New Roman 12pt
  - Title page separate from main document
  - Line numbers and page numbers required
"""

import os
import re
import pandas as pd
from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_DIR = os.path.dirname(BASE_DIR)
OUTPUT_DIR = os.path.join(PROJECT_DIR, 'output')
MANUSCRIPT_DIR = BASE_DIR

# ============================================================
# HELPER FUNCTIONS
# ============================================================

def build_main_document():
    """Build the main document (separate from title page per JMA guidelines)."""
    doc = Document()

    # Page setup
    for section in doc.sections:
        section.top_margin = Cm(2.54)
        section.bottom_margin = Cm(2.54)
        section.left_margin = Cm(2.54)
        section.right_margin = Cm(2.54)

    style = doc.styles['Normal']
    font = style.font
    font.name = 'Times New Roman'
    font.size = Pt(12)
    style.paragraph_format.space_after = Pt(6)
    style.paragraph_format.line_spacing = 2.0

    def add_heading_styled(text, level=1):
        h = doc.add_heading(text, level=level)
        for run in h.runs:
            run.font.color.rgb = RGBColor(0, 0, 0)
            run.font.name = 'Times New Roman'
        return h

    def _add_runs_with_citations(p, text, font_size=Pt(12), bold=False, italic=False):
        """Parse {n} or {n-m} patterns as superscript citations."""
        parts = re.split(r'(\{[^}]+\})', text)
        for part in parts:
            if part.startswith('{') and part.endswith('}'):
                run = p.add_run(part[1:-1])
                run.font.name = 'Times New Roman'
                run.font.size = Pt(10)
                run.font.superscript = True
                run.bold = False
                run.italic = False
            else:
                if not part:
                    continue
                run = p.add_run(part)
                run.font.name = 'Times New Roman'
                run.font.size = font_size
                run.bold = bold
                run.italic = italic

    def add_para(text, bold=False, italic=False, alignment=None, space_after=Pt(6)):
        p = doc.add_paragraph()
        _add_runs_with_citations(p, text, bold=bold, italic=italic)
        if alignment:
            p.alignment = alignment
        p.paragraph_format.space_after = space_after
        p.paragraph_format.line_spacing = 2.0
        return p

    def add_abstract_field(label, text):
        p = doc.add_paragraph()
        p.paragraph_format.line_spacing = 2.0
        p.paragraph_format.space_after = Pt(6)
        run_label = p.add_run(label + ' ')
        run_label.bold = True
        run_label.font.name = 'Times New Roman'
        run_label.font.size = Pt(12)
        _add_runs_with_citations(p, text)
        return p

    def add_figure_inline(filename, caption_text):
        """Insert a figure image inline with caption."""
        img_path = os.path.join(OUTPUT_DIR, filename)
        if not os.path.exists(img_path):
            return
        # Caption above
        cap = doc.add_paragraph()
        cap.paragraph_format.space_before = Pt(12)
        cap.paragraph_format.space_after = Pt(6)
        cap.paragraph_format.line_spacing = 2.0
        run_cap = cap.add_run(caption_text)
        run_cap.font.name = 'Times New Roman'
        run_cap.font.size = Pt(10)
        run_cap.bold = True
        # Image
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run_img = p.add_run()
        run_img.add_picture(img_path, width=Inches(5.5))
        doc.add_paragraph()  # spacing

    def add_table_from_data(headers, rows, caption=None):
        if caption:
            cap = doc.add_paragraph()
            cap.paragraph_format.space_after = Pt(6)
            cap.paragraph_format.space_before = Pt(12)
            cap.paragraph_format.line_spacing = 2.0
            run_cap = cap.add_run(caption)
            run_cap.font.name = 'Times New Roman'
            run_cap.font.size = Pt(10)
            run_cap.bold = True
        table = doc.add_table(rows=1 + len(rows), cols=len(headers))
        table.style = 'Table Grid'
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        for i, h in enumerate(headers):
            cell = table.rows[0].cells[i]
            cell.text = ''
            p = cell.paragraphs[0]
            run = p.add_run(h)
            run.bold = True
            run.font.name = 'Times New Roman'
            run.font.size = Pt(9)
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for r_idx, row_data in enumerate(rows):
            for c_idx, val in enumerate(row_data):
                cell = table.rows[r_idx + 1].cells[c_idx]
                cell.text = ''
                p = cell.paragraphs[0]
                run = p.add_run(str(val))
                run.font.name = 'Times New Roman'
                run.font.size = Pt(9)
                if c_idx > 0:
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph()
        return table

    # ============================================================
    # ABSTRACT
    # ============================================================
    add_heading_styled('Abstract', level=1)

    add_abstract_field('Introduction:',
        'Medical safety incidents may deter physicians from high-risk '
        'specialties, yet the temporal dynamics remain poorly quantified. '
        'Japan maintains two independent national incident series\u2014'
        'occurrence-based mandatory reports and malpractice litigation '
        'statistics\u2014but no multi-specialty test of predictive '
        'causality between these series and physician supply has been '
        'conducted.')

    add_abstract_field('Methods:',
        'We constructed bivariate vector autoregression (VAR) models for '
        '12 specialties using mandatory safety reports (2015\u20132025) and '
        'malpractice litigation statistics (2004\u20132023), paired with '
        'physician and facility counts from national surveys. Granger '
        'causality F-tests, impulse response functions, and VAR-based '
        'forecasts were computed. Sensitivity analyses assessed the '
        'independence and interaction of the two incident sources.')

    add_abstract_field('Results:',
        'Litigation-based models showed significant forward Granger '
        'causality in 18 of 24 tests; obstetrics and gynaecology had '
        'the strongest association (F = 46.66, P < 0.001). Bidirectional '
        'causality emerged in 7 specialty\u2013outcome pairs. A litigation '
        'shock depressed physician supply for 3\u20135 years. Occurrence '
        'and litigation counts were weakly correlated (median '
        '|r| = 0.23), and their interaction was non-significant, '
        'indicating additive, independent effects on workforce '
        'decline.')

    add_abstract_field('Conclusions:',
        'Safety incidents Granger-cause specialty-specific workforce '
        'decline lasting 3\u20135 years through two independent '
        'pathways. Because litigation is largely independent of '
        'incident occurrence, early dispute resolution may attenuate '
        'workforce loss without requiring a reduction in incident '
        'rates.')

    add_para(
        'Key words: physician workforce, Granger causality, medical safety '
        'incidents, specialty maldistribution, dispute resolution, Japan',
        italic=True, space_after=Pt(18),
    )

    # ============================================================
    # INTRODUCTION
    # ============================================================
    doc.add_page_break()
    add_heading_styled('Introduction', level=1)

    add_para(
        'Japan\u2019s physician workforce has grown steadily to '
        'approximately 340,000, yet aggregate expansion masks deepening '
        'maldistribution across specialties and regions.{1,2} General '
        'surgery, obstetrics, emergency medicine, and paediatrics face '
        'persistent shortages even as total supply rises.{3,4}'
    )

    add_para(
        'Safety incidents as a determinant of specialty choice gained '
        'prominence after the 2004 Fukushima obstetrics prosecution, in '
        'which an obstetrician was arrested following a maternal death '
        'during caesarean section.{5,6} Morita reported a 13% decline '
        'in obstetricians in the affected prefecture using '
        'difference-in-differences.{7} That analysis, however, addressed '
        'a single event in one specialty and one region; broader '
        'multi-specialty evidence linking incident burden to workforce '
        'dynamics is lacking.'
    )

    add_para(
        'Japan offers two complementary national incident series: the '
        'Japan Medical Safety Research Organisation (JMSR) mandatory '
        'reports (from 2015){8} and Supreme Court malpractice '
        'litigation statistics by specialty.{9,10} Paired with biennial '
        'physician surveys and annual facility registrations, these data '
        'permit a systematic multi-specialty analysis.'
    )

    add_para(
        'Here we apply vector autoregression (VAR) and Granger '
        'causality testing{11\u201313}\u2014established econometric tools '
        'for predictive-causality inference\u2014to determine whether '
        'incident trends temporally precede changes in '
        'specialty-specific physician and facility counts across 12 '
        'specialties over up to 20 years. Because the two incident '
        'series differ in construction (occurrence vs. adjudicated '
        'claims), we additionally examine whether they act as '
        'independent or interacting predictors of workforce change.'
    )

    # ============================================================
    # MATERIALS AND METHODS
    # ============================================================
    doc.add_page_break()
    add_heading_styled('Materials and Methods', level=1)

    add_heading_styled('Study design', level=2)
    add_para(
        'This is a secondary analysis of routinely collected national '
        'administrative data using bivariate VAR models and Granger '
        'causality tests, reported in accordance with the STROBE '
        '(Strengthening the Reporting of Observational Studies in '
        'Epidemiology) statement for cross-sectional studies.{14}'
    )

    add_heading_styled('Data sources', level=2)
    add_para(
        'Medical safety incident data were obtained from two sources. '
        'Definition 1 (JMSR): annual specialty-specific incident reports '
        'to the Japan Medical Safety Research Organisation for fiscal years '
        '2015\u20132025 (11 annual observations).{8} Definition 2 (litigation): '
        'annual specialty-specific medical malpractice closed-claim counts '
        'from the Supreme Court of Japan for 2004\u20132023 (20 annual '
        'observations).{9}'
    )
    add_para(
        'Specialty-specific physician counts were obtained from the '
        'Survey of Physicians, Dentists, and Pharmacists (biennial, '
        '2002\u20132024; annual values by linear interpolation).{15} '
        'Facility counts were obtained from the Survey of Medical '
        'Institutions (2005\u20132024).{16} Specialist trainee enrolment '
        'data were obtained from the Japan Board of Medical Specialties '
        '(2018\u20132025).{17}'
    )

    add_heading_styled('Specialties analysed', level=2)
    add_para(
        'Twelve core specialties with sufficient data across all sources '
        'were analysed: internal medicine, general surgery, orthopaedic '
        'surgery, plastic surgery, obstetrics and gynaecology, paediatrics, '
        'psychiatry, ophthalmology, otolaryngology, urology, dermatology, '
        'and anaesthesiology.'
    )

    add_heading_styled('Statistical analysis', level=2)
    add_para(
        'Stationarity was assessed using augmented Dickey\u2013Fuller (ADF) '
        'tests. Non-stationary pairs were first-differenced. For each '
        'specialty\u2013definition\u2013outcome combination, a bivariate VAR model '
        'was fitted with lag order selected by minimising the Akaike '
        'information criterion (AIC), capped at min(4, n/3 \u2212 1).'
    )

    add_para(
        'Granger causality was tested bidirectionally using F-tests on '
        'lagged coefficients.{18} The forward null hypothesis is that '
        'lagged incident counts do not improve prediction of the '
        'workforce outcome; significance was set at P < 0.05.'
    )

    add_para(
        'Impulse response functions (IRFs) were computed over a 10-year '
        'horizon to estimate the magnitude and duration of workforce '
        'effects following a one-unit incident shock.{19} VAR-based '
        'forecasts with 95% intervals were generated for 2024\u20132033 using '
        'litigation-based models.'
    )

    add_para(
        'For specialist trainees (8 annual observations), detrended '
        'Pearson correlations between JMSR incidents and enrolment were '
        'computed. Analyses used Python 3.10 with statsmodels 0.14.{20} '
        'Code and data are deposited at '
        'https://github.com/bougtoir/medical-accident-its-analysis.'
    )

    add_heading_styled('Sensitivity and source-independence analyses', level=2)
    add_para(
        'To assess whether the disparity in significant models between '
        'litigation and JMSR reflected series length rather than '
        'intrinsic differences, litigation data were restricted to the '
        '2015\u20132023 window (9 years) overlapping with the JMSR period '
        'and the Granger tests were repeated.'
    )

    add_para(
        'To evaluate source independence, Pearson correlations between '
        'JMSR and litigation counts were computed for each specialty '
        'over the overlapping period. To test whether occurrence and '
        'litigation act additively or synergistically, trivariate VAR '
        'models (JMSR + litigation + outcome) were fitted and compared '
        'with bivariate models by AIC. An interaction term (JMSR \u00d7 '
        'litigation) was tested for additional predictive power via '
        'Granger F-tests.'
    )

    add_heading_styled('Ethical considerations', level=2)
    add_para(
        'This study used publicly available, anonymised, aggregate-level '
        'administrative data. No individual-level data were accessed. '
        'Ethical approval was not required per institutional guidelines.'
    )

    # ============================================================
    # RESULTS
    # ============================================================
    doc.add_page_break()
    add_heading_styled('Results', level=1)

    add_heading_styled('Descriptive overview', level=2)
    add_para(
        'Over the study period the JMSR received 3,860 reports across '
        '12 specialties (2015\u20132025); general surgery (n = 520), '
        'internal medicine (n = 479), and orthopaedic surgery (n = 310) '
        'contributed the largest volumes. Malpractice litigation declined '
        'from 1,139 cases (2004) to 746 (2023) nationally, though '
        'specialty-level trajectories diverged (Figure 1).'
    )

    add_figure_inline('fig4_accident_trends.png',
        'Figure 1. Trends in medical safety incident counts by specialty.')

    add_heading_styled('Granger causality tests', level=2)
    add_para(
        'A total of 47 bivariate VAR models were fitted. Table 1 presents '
        'the litigation-based Granger causality results. Of 24 tests (12 '
        'specialties \u00d7 2 outcomes), 18 showed significant forward Granger '
        'causality (incidents \u2192 workforce, P < 0.05). The strongest '
        'forward associations were obstetrics and gynaecology (physicians: '
        'F = 46.66, P < 0.001), psychiatry (physicians: F = 20.89, '
        'P < 0.001), and anaesthesiology (facilities: F = 20.34, '
        'P < 0.001).'
    )

    # Load Granger results for Table 1
    gc_df = pd.read_csv(os.path.join(OUTPUT_DIR, 'granger_causality_results.csv'))
    gc_lit = gc_df[gc_df['definition'] == 'litigation'].copy()

    table1_headers = ['Specialty', 'Outcome', 'Lag', 'n',
                      'Forward F', 'Forward P', 'Reverse F', 'Reverse P']
    table1_rows = []
    for _, r in gc_lit.sort_values(['specialty_en', 'outcome']).iterrows():
        fwd_f = f"{r['gc_forward_F']:.2f}" if pd.notna(r.get('gc_forward_F')) else '\u2014'
        fwd_p = f"{r['gc_forward_p']:.4f}" if pd.notna(r.get('gc_forward_p')) else '\u2014'
        rev_f = f"{r['gc_reverse_F']:.2f}" if pd.notna(r.get('gc_reverse_F')) else '\u2014'
        rev_p = f"{r['gc_reverse_p']:.4f}" if pd.notna(r.get('gc_reverse_p')) else '\u2014'
        if pd.notna(r.get('gc_forward_p')) and r['gc_forward_p'] < 0.05:
            fwd_p += '*'
        if pd.notna(r.get('gc_reverse_p')) and r['gc_reverse_p'] < 0.05:
            rev_p += '*'
        table1_rows.append([
            r['specialty_en'], r['outcome'].capitalize(),
            str(int(r['var_lag'])), str(int(r['n_obs'])),
            fwd_f, fwd_p, rev_f, rev_p,
        ])

    add_table_from_data(
        table1_headers, table1_rows,
        caption='Table 1. Granger causality test results (litigation data, '
                '12 specialties \u00d7 2 outcomes). Forward: incidents \u2192 workforce; '
                'Reverse: workforce \u2192 incidents. *P < 0.05.'
    )

    add_para(
        'Significant reverse Granger causality (workforce \u2192 incidents) was '
        'observed in 8 models. Seven specialty\u2013outcome combinations showed '
        'bidirectional causality: obstetrics and gynaecology (physicians '
        'and facilities), paediatrics (physicians and facilities), general '
        'surgery (physicians), ophthalmology (facilities), and '
        'otolaryngology (facilities).'
    )

    add_para(
        'In the JMSR-based sensitivity analysis (Table 2), only 2 of 24 '
        'tests reached significance, reflecting the shorter series. '
        'However, general surgery (F = 23.04, P = 0.003) and psychiatry '
        '(F = 9.48, P = 0.022) were significant even with 10\u201311 '
        'observations.'
    )

    # Table 2: JMSR results
    gc_jmsr = gc_df[gc_df['definition'] == 'jmsr'].copy()
    table2_rows = []
    for _, r in gc_jmsr.sort_values(['specialty_en', 'outcome']).iterrows():
        fwd_f = f"{r['gc_forward_F']:.2f}" if pd.notna(r.get('gc_forward_F')) else '\u2014'
        fwd_p = f"{r['gc_forward_p']:.4f}" if pd.notna(r.get('gc_forward_p')) else '\u2014'
        rev_f = f"{r['gc_reverse_F']:.2f}" if pd.notna(r.get('gc_reverse_F')) else '\u2014'
        rev_p = f"{r['gc_reverse_p']:.4f}" if pd.notna(r.get('gc_reverse_p')) else '\u2014'
        if pd.notna(r.get('gc_forward_p')) and r['gc_forward_p'] < 0.05:
            fwd_p += '*'
        if pd.notna(r.get('gc_reverse_p')) and r['gc_reverse_p'] < 0.05:
            rev_p += '*'
        table2_rows.append([
            r['specialty_en'], r['outcome'].capitalize(),
            str(int(r['var_lag'])), str(int(r['n_obs'])),
            fwd_f, fwd_p, rev_f, rev_p,
        ])

    add_table_from_data(
        table1_headers, table2_rows,
        caption='Table 2. Granger causality test results (JMSR data, '
                'sensitivity analysis, 10\u201311 observations). *P < 0.05.'
    )

    # IRF results
    add_heading_styled('Impulse response functions', level=2)
    add_para(
        'Figure 2 displays physician-count IRFs following a one-unit '
        'litigation shock for six key specialties. Obstetrics and '
        'gynaecology exhibited a sustained negative response peaking at '
        '2\u20133 years and persisting approximately 5 years; general '
        'surgery followed a similar trajectory attenuating by year 4. '
        'Facility-count IRFs appear in Figure 3.'
    )

    add_figure_inline('fig2_irf_physicians.png',
        'Figure 2. IRF: physician count response to a one-unit litigation shock.')

    add_figure_inline('sfig1_irf_facilities.png',
        'Figure 3. IRF: facility count response to a one-unit litigation shock.')

    # Forecasts
    doc.add_page_break()
    add_heading_styled('VAR-based forecasts', level=2)
    add_para(
        'Figure 4 presents VAR-based forecasts for physician and facility '
        'counts. General surgery is projected to continue declining '
        '(physicians: \u2212230/year; facilities: \u2212310/year through 2033). '
        'Obstetrics and gynaecology physician counts show modest growth '
        'while facilities decline. Table 3 summarises forecast values.'
    )

    add_figure_inline('fig3_var_forecasts.png',
        'Figure 4. VAR-based forecasts of physician and facility counts (2024\u20132033).')

    # Table 3: Forecast summary
    fc_df = pd.read_csv(os.path.join(OUTPUT_DIR, 'var_forecast_results.csv'))
    table3_headers = ['Specialty', 'Outcome', '2028 (forecast)', '2033 (forecast)']
    table3_rows = []
    for spec in ['General surgery', 'Obstetrics & gynaecology', 'Internal medicine',
                 'Paediatrics', 'Orthopaedic surgery', 'Anaesthesiology']:
        for outcome in ['physicians', 'facilities']:
            sub = fc_df[(fc_df['specialty_en'] == spec) & (fc_df['outcome'] == outcome)]
            if sub.empty:
                continue
            y2028 = sub[sub['year'] == 2028]
            y2033 = sub[sub['year'] == 2033]
            v2028 = f"{y2028['forecast_mean'].values[0]:,.0f}" if len(y2028) > 0 else '\u2014'
            v2033 = f"{y2033['forecast_mean'].values[0]:,.0f}" if len(y2033) > 0 else '\u2014'
            table3_rows.append([spec, outcome.capitalize(), v2028, v2033])

    add_table_from_data(
        table3_headers, table3_rows,
        caption='Table 3. VAR-based workforce forecasts for selected specialties.'
    )

    # Trainee analysis
    add_heading_styled('Trainee correlations', level=2)
    add_para(
        'Detrended correlations between JMSR incidents and specialist '
        'trainee enrolment (2018\u20132025, n = 8) showed notable negative '
        'associations for psychiatry (r = \u22120.77, P = 0.025), '
        'anaesthesiology (r = \u22120.77, P = 0.026), and obstetrics and '
        'gynaecology (r = \u22120.70, P = 0.051).'
    )

    # Sensitivity and source-independence results
    add_heading_styled('Sensitivity and source-independence analyses', level=2)
    add_para(
        'When litigation data were restricted to the 2015\u20132023 '
        'window (9 years, matching the JMSR period), significant '
        'forward Granger models fell from 18 to 4, indicating that '
        'the disparity in significance counts between data sources '
        'largely reflects statistical power afforded by the longer '
        'series rather than an intrinsic difference between occurrence '
        'and litigation as predictors.'
    )

    add_para(
        'JMSR incident counts and litigation counts were weakly '
        'correlated across the overlapping period (median '
        '|r|\u2009=\u20090.23; no specialty reached P < 0.05). '
        'Trivariate VAR models incorporating both sources '
        'simultaneously did not improve fit over bivariate models '
        '(mean \u0394AIC = +0.6), and an interaction term '
        '(occurrence \u00d7 litigation) reached significance in only '
        '1 of 22 models, consistent with additive rather than '
        'synergistic effects.'
    )

    # ============================================================
    # DISCUSSION
    # ============================================================
    doc.add_page_break()
    add_heading_styled('Discussion', level=1)

    add_para(
        'This multi-specialty Granger causality analysis yields three '
        'principal findings. First, safety incidents\u2014whether '
        'measured by occurrence reports or litigation\u2014temporally '
        'precede workforce decline in the majority of specialties '
        'when series length affords adequate statistical power. '
        'Second, the two incident sources are largely independent '
        '(median |r| = 0.23), implying that litigation captures a '
        'distinct dimension of professional risk\u2014reputational '
        'exposure, financial liability, and prolonged procedural '
        'stress\u2014not reducible to incident frequency. Third, '
        'these sources act additively rather than synergistically, '
        'suggesting that interventions targeting either pathway may '
        'yield independent workforce benefits.'
    )

    add_heading_styled('Bidirectional causality', level=2)
    add_para(
        'Bidirectional Granger causality in 7 specialty\u2013outcome '
        'pairs points to a reinforcing cycle. The forward pathway '
        '(incidents \u2192 workforce decline) is consistent with '
        'incident-avoidance behaviour: high-profile events deter '
        'specialty entry and hasten exit.{1,7} The reverse pathway '
        '(workforce decline \u2192 incidents) accords with the overwork '
        'hypothesis\u2014fewer physicians per caseload raises fatigue and '
        'incident risk.{21} That both directions reach significance in '
        'obstetrics, paediatrics, and general surgery identifies these '
        'specialties as most susceptible to self-perpetuating decline.'
    )

    add_heading_styled('Work-style reform as a compounding factor', level=2)
    add_para(
        'These incident-driven dynamics operate within a labour market '
        'reshaped by the April 2024 work-style reform (960-hour annual '
        'overtime cap, A-level; 1,860-hour transitional B/C-level).{22} '
        'Restrictions on "gaikin" (moonlighting opportunities) have '
        'eroded the income supplement that historically offset demanding '
        'hospital posts, while deemed self-study\u2014hours reclassified '
        'as voluntary training rather than compensable labour\u2014has '
        'shifted workload onto mid-career physicians without relieving '
        'trainee burden.{23} '
        'The 3\u20135 year incident-response horizon identified by '
        'our IRFs thus coincides with a period of contracting '
        'financial incentives for the very specialties most affected.'
    )

    add_heading_styled('Policy implications', level=2)
    add_para(
        'The independence of occurrence and litigation counts has a '
        'direct policy corollary: because safety events occur with '
        'a certain irreducible probability, the litigation pathway '
        'constitutes a modifiable mediator. If malpractice disputes '
        'are resolved before reaching court\u2014through early '
        'mediation, alternative dispute resolution, or no-fault '
        'compensation\u2014the workforce signal associated with '
        'litigation may be attenuated without requiring a reduction '
        'in incident rates themselves.{5,24}'
    )

    add_para(
        'Beyond dispute resolution, several complementary levers '
        'emerge from our results. Economic incentives for high-risk '
        'specialties (enhanced remuneration, training stipends, '
        'rural service allowances) address the most proximate '
        'driver of specialty avoidance.{24} Fostering a non-punitive '
        'safety culture\u2014structured post-incident psychological '
        'support, transparent root-cause analysis, and proportionate '
        'media reporting\u2014may attenuate the deterrent signal of '
        'adverse events.{5} The reverse-causality pathway (fewer '
        'physicians \u2192 more incidents) implies that minimum '
        'staffing thresholds in specialties exhibiting bidirectional '
        'causality (obstetrics, paediatrics, general surgery) could '
        'interrupt the reinforcing cycle before it becomes '
        'self-sustaining. Effective implementation of work-style '
        'reform\u2014including correction of deemed self-study '
        'practices and equitable redistribution of workload\u2014is '
        'essential to prevent the reform itself from accelerating '
        'exit from demanding specialties.'
    )

    add_para(
        'The 3\u20135 year effect duration estimated by our IRFs '
        'defines a critical monitoring window: targeted surveillance '
        'and support in the years following a major safety event may '
        'prevent transient workforce shocks from consolidating into '
        'permanent losses. The projected contraction in general '
        'surgery (\u2212230 physicians/year; \u2212310 facilities/year) '
        'underscores the urgency of deploying these measures '
        'concurrently rather than sequentially.'
    )

    add_heading_styled('Strengths and limitations', level=2)
    add_para(
        'Strengths of this study include population-level coverage from '
        'mandatory national registries across 12 specialties over up to '
        '20 years, and a VAR/Granger framework suited to continuously '
        'varying exposures without requiring an assumed intervention '
        'point.'
    )

    add_para(
        'Several limitations warrant consideration. The ecological '
        'design precludes individual-level causal inference. Biennial '
        'physician data required interpolation, potentially smoothing '
        'abrupt transitions. The JMSR series spans only 11 years, '
        'limiting statistical power for that data source. Unmeasured '
        'confounders\u2014demographic shifts, the 2004 postgraduate '
        'training reform, and the direct effects of work-style reform '
        '(temporally overlapping the most recent observations)\u2014may '
        'bias estimates. Notably, our 12-specialty framework does not '
        'capture career paths outside conventional specialty '
        'registration, such as aesthetic medicine{25,26} and '
        'occupational medicine.{27,28} To the extent that migration '
        'into these sectors correlates temporally with incident '
        'trends, our models may overestimate the '
        'incident-attributable component of workforce decline. Low '
        'retention among occupational physicians (55\u201365% turnover '
        'within 2 years){28} further complicates workforce '
        'projections.'
    )

    # ============================================================
    # CONCLUSIONS (brief, per JMA format — part of Discussion or separate)
    # ============================================================
    add_heading_styled('Conclusions', level=2)
    add_para(
        'Safety incidents Granger-cause reductions in physician supply '
        'in 9 of 12 specialties, with effects persisting 3\u20135 years. '
        'Occurrence-based and litigation-based incident counts are '
        'largely independent yet both predict workforce decline, '
        'identifying litigation as a modifiable mediator. Early dispute '
        'resolution, combined with economic incentives, non-punitive '
        'safety culture, and minimum staffing safeguards, may '
        'interrupt the reinforcing shortage\u2013incident cycle within '
        'the 3\u20135 year policy window identified here.'
    )

    # ============================================================
    # DECLARATIONS
    # ============================================================
    doc.add_page_break()
    add_heading_styled('Conflict of Interest', level=2)
    add_para('The author declares no conflict of interest.')

    add_heading_styled('Funding', level=2)
    add_para('None.')

    add_heading_styled('Author Contributions', level=2)
    add_para(
        'TO conceived the study, acquired data, performed statistical '
        'analyses, and drafted the manuscript.'
    )

    add_heading_styled('Data Availability', level=2)
    add_para(
        'Analysis code and aggregated datasets are available at '
        'https://github.com/bougtoir/medical-accident-its-analysis. '
        'Original data sources are publicly available from the Japan '
        'Medical Safety Research Organisation, the Supreme Court of Japan, '
        'and the Ministry of Health, Labour and Welfare.'
    )

    add_heading_styled('Acknowledgements', level=2)
    add_para(
        'AI-assisted tools were used for coding assistance and manuscript '
        'drafting. The author takes full responsibility for the content '
        'and integrity of the work.'
    )

    # ============================================================
    # REFERENCES (Vancouver, superscript in text, ≤3 authors then et al.)
    # ============================================================
    doc.add_page_break()
    add_heading_styled('References', level=1)

    references = [
        # 1 — Ikegami (unchanged)
        'Ikegami N, Yoo BK, Hashimoto H, et al. Japanese universal health '
        'coverage: evolution, achievements, and challenges. Lancet. '
        '2011;378(9796):1106-15.',
        # 2 — Toyabe (unchanged)
        'Toyabe S. Trend in geographic distribution of physicians in Japan. '
        'Int J Equity Health. 2009;8:5.',
        # 3 — MHLW Survey 2022 (unchanged)
        'Ministry of Health, Labour and Welfare. Survey of physicians, '
        'dentists, and pharmacists 2022. Tokyo: MHLW; 2023. Japanese.',
        # 4 — Tanaka (unchanged)
        'Tanaka K, Katsumata Y, Matsuda S. Supply, demand and distribution '
        'of physicians in Japan. Keio J Med. 2024;doi:10.46308/kmj.2024.00087.',
        # 5 — Nagamatsu (was 7)
        'Nagamatsu S, Kami M, Nakata Y. Healthcare safety committee in '
        'Japan: mandatory accountability reporting system and punishment. '
        'Curr Opin Anaesthesiol. 2009;22(2):199-206.',
        # 6 — Hiyama (was 8)
        'Hiyama T, Yoshihara M, Tanaka S, et al. Defensive medicine '
        'practices among gastroenterologists in Japan. World J '
        'Gastroenterol. 2006;12(47):7671-5.',
        # 7 — Morita (was 9)
        'Morita H. Criminal prosecution and physician supply. Int Rev Law '
        'Econ. 2018;55:1-11.',
        # 8 — JMSR (was 10)
        'Japan Medical Safety Research Organisation. Annual report on '
        'medical accident investigation. Tokyo: JMSR; 2025. Available '
        'from: https://www.medsafe.or.jp/',
        # 9 — Supreme Court (was 11)
        'Supreme Court of Japan. Annual report of judicial statistics: '
        'medical malpractice litigation. Tokyo: Supreme Court; 2024. Japanese.',
        # 10 — Taniguchi (was 12)
        'Taniguchi K, Watari T, Nagoshi K. Characteristics and trends of '
        'medical malpractice claims in Japan between 2006 and 2021. PLoS '
        'One. 2024;19(1):e0296155.',
        # 11 — Granger (was 13)
        'Granger CWJ. Investigating causal relations by econometric models '
        'and cross-spectral methods. Econometrica. 1969;37(3):424-38.',
        # 12 — Lütkepohl (was 14)
        'L\u00fctkepohl H. New Introduction to Multiple Time Series Analysis. '
        'Berlin: Springer; 2005.',
        # 13 — Toda (was 15)
        'Toda HY, Yamamoto T. Statistical inference in vector '
        'autoregressions with possibly integrated processes. J Econom. '
        '1995;66(1-2):225-50.',
        # 14 — STROBE (was 16)
        'von Elm E, Altman DG, Egger M, et al. The Strengthening the '
        'Reporting of Observational Studies in Epidemiology (STROBE) '
        'statement: guidelines for reporting observational studies. '
        'Lancet. 2007;370(9596):1453-7.',
        # 15 — Physicians survey (was 17)
        'Ministry of Health, Labour and Welfare. Survey of physicians, '
        'dentists, and pharmacists. Tokyo: MHLW; 2024. Japanese.',
        # 16 — Facilities survey (was 18)
        'Ministry of Health, Labour and Welfare. Survey of medical '
        'institutions (dynamic survey). Tokyo: MHLW; 2024. Japanese.',
        # 17 — JBMS (was 19)
        'Japan Board of Medical Specialties. Specialist trainee '
        'registration statistics. Tokyo: JBMS; 2025. Japanese.',
        # 18 — Hamilton (was 20)
        'Hamilton JD. Time Series Analysis. Princeton: Princeton '
        'University Press; 1994.',
        # 19 — Sims (was 21)
        'Sims CA. Macroeconomics and reality. Econometrica. '
        '1980;48(1):1-48.',
        # 20 — Statsmodels (was 22)
        'Seabold S, Perktold J. Statsmodels: econometric and statistical '
        'modeling with Python. Proc 9th Python Sci Conf. 2010:92-6.',
        # 21 — Studdert (was 23)
        'Studdert DM, Mello MM, Sage WM, et al. Defensive medicine among '
        'high-risk specialist physicians in a volatile malpractice '
        'environment. JAMA. 2005;293(21):2609-17.',
        # 22 — Ishikawa (was 5)
        'Ishikawa T, Ohba H, Yokooka Y, et al. Labor shortage of '
        'physicians in rural areas and surgical specialties caused by Work '
        'Style Reform Policies of the Japanese government: a quantitative '
        'simulation analysis. J Rural Med. 2024;19(3):198-207.',
        # 23 — JMA work-style (was 6)
        'Japan Medical Association. Survey on physician work-style reform '
        'implementation status. Tokyo: JMA; 2024. Japanese.',
        # 24 — Currie (unchanged)
        'Currie J, MacLeod WB. First do no harm? Tort reform and birth '
        'outcomes. Q J Econ. 2008;123(2):795-830.',
        # 25 — Takei (unchanged)
        'Takei T. The \u2018Chokubi\u2019 phenomenon: young physicians\u2019 exodus from '
        'state service to private medicine in Japan. QJM. '
        '2024;117(12):843-5.',
        # 26 — MHLW aesthetic (unchanged)
        'Ministry of Health, Labour and Welfare. Survey of medical '
        'institutions 2023: aesthetic surgery clinics. Tokyo: MHLW; 2024. Japanese.',
        # 27 — Koike (unchanged)
        'Koike S, Isse T, Kawaguchi H, et al. Retention among full-time '
        'occupational physicians in Japan. Occup Med (Lond). '
        '2019;69(2):139-42.',
        # 28 — Isse (unchanged)
        'Isse T, Nakamura H, Hachisuka K. Turnover of full-time '
        'occupational physicians in Japan in the period 2002\u20132008. Sangyo '
        'Eiseigaku Zasshi. 2012;54(5):174-83. Japanese.',
    ]

    for i, ref in enumerate(references, 1):
        p = doc.add_paragraph()
        p.paragraph_format.line_spacing = 2.0
        p.paragraph_format.space_after = Pt(2)
        run = p.add_run(f'{i}. {ref}')
        run.font.name = 'Times New Roman'
        run.font.size = Pt(10)

    # ============================================================
    # FIGURE LEGENDS
    # ============================================================
    doc.add_page_break()
    add_heading_styled('Figure Legends', level=1)

    legends = [
        ('Figure 1.',
         'Trends in medical safety incident counts by specialty. '
         'Panel A: JMSR mandatory reports (2015\u20132025). '
         'Panel B: Medical malpractice litigation closed claims (2004\u20132023).'),
        ('Figure 2.',
         'Impulse response functions: response of physician counts to a '
         'one-unit shock in litigation cases for six specialties. '
         'Shaded areas represent approximate 95% confidence bands.'),
        ('Figure 3.',
         'Impulse response functions: response of facility counts to a '
         'one-unit shock in litigation cases for six specialties.'),
        ('Figure 4.',
         'VAR-based forecasts of physician and facility counts for '
         'general surgery, obstetrics and gynaecology, and internal medicine '
         '(2024\u20132033). Circles: observed data; dashed line: forecast; '
         'shaded area: 95% forecast interval.'),
    ]

    for label, text in legends:
        p = doc.add_paragraph()
        p.paragraph_format.line_spacing = 2.0
        p.paragraph_format.space_after = Pt(12)
        run_label = p.add_run(label + ' ')
        run_label.bold = True
        run_label.font.name = 'Times New Roman'
        run_label.font.size = Pt(11)
        run_text = p.add_run(text)
        run_text.font.name = 'Times New Roman'
        run_text.font.size = Pt(11)

    # Save main document
    output_path = os.path.join(MANUSCRIPT_DIR, 'jma_manuscript_en.docx')
    doc.save(output_path)
    print(f'Main document saved to {output_path}')
    return output_path


def build_title_page():
    """Build the title page as a separate document per JMA guidelines."""
    doc = Document()

    for section in doc.sections:
        section.top_margin = Cm(2.54)
        section.bottom_margin = Cm(2.54)
        section.left_margin = Cm(2.54)
        section.right_margin = Cm(2.54)

    style = doc.styles['Normal']
    font = style.font
    font.name = 'Times New Roman'
    font.size = Pt(12)
    style.paragraph_format.line_spacing = 2.0

    # Title
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_after = Pt(24)
    run = p.add_run(
        'Temporal associations between medical safety incidents and '
        'physician workforce dynamics in Japan: a vector autoregression '
        'analysis in the era of work-style reform and specialty '
        'maldistribution'
    )
    run.font.name = 'Times New Roman'
    run.font.size = Pt(14)
    run.bold = True

    # Short running title
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_after = Pt(18)
    run = p.add_run('Running title: Safety incidents and physician workforce in Japan')
    run.font.name = 'Times New Roman'
    run.font.size = Pt(11)
    run.italic = True

    # Authors
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_after = Pt(12)
    run = p.add_run('Tatsuki Onishi')
    run.font.name = 'Times New Roman'
    run.font.size = Pt(12)

    # Affiliation
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_after = Pt(24)
    run = p.add_run(
        'Data Science AI Innovation Research Promotion Center, '
        'Shiga University, Hikone, Shiga, Japan'
    )
    run.font.name = 'Times New Roman'
    run.font.size = Pt(11)

    # Corresponding author
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(6)
    run = p.add_run('Corresponding author:')
    run.bold = True
    run.font.name = 'Times New Roman'
    run.font.size = Pt(11)

    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(6)
    run = p.add_run(
        'Tatsuki Onishi\n'
        'Data Science AI Innovation Research Promotion Center, Shiga University\n'
        '1-1-1 Bamba, Hikone, Shiga 522-8522, Japan\n'
        'E-mail: [to be added]\n'
        'Telephone: [to be added]'
    )
    run.font.name = 'Times New Roman'
    run.font.size = Pt(11)

    # Word count
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(18)
    p.paragraph_format.space_after = Pt(6)
    run = p.add_run(
        'Word count: approximately 3,400 '
        '(excluding abstract, references, tables, and figure legends)'
    )
    run.font.name = 'Times New Roman'
    run.font.size = Pt(11)

    # Tables and figures
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(6)
    run = p.add_run('Number of tables: 3')
    run.font.name = 'Times New Roman'
    run.font.size = Pt(11)

    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(6)
    run = p.add_run('Number of figures: 4')
    run.font.name = 'Times New Roman'
    run.font.size = Pt(11)

    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(6)
    run = p.add_run('Number of references: 29')
    run.font.name = 'Times New Roman'
    run.font.size = Pt(11)

    # Article type
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(18)
    run = p.add_run('Article type: Original Research Article')
    run.font.name = 'Times New Roman'
    run.font.size = Pt(11)
    run.bold = True

    output_path = os.path.join(MANUSCRIPT_DIR, 'jma_title_page.docx')
    doc.save(output_path)
    print(f'Title page saved to {output_path}')
    return output_path


def build_pptx_figures():
    """Build editable PPTX with figures (1 per slide)."""
    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt

    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    figures = [
        ('Figure 1', 'fig4_accident_trends.png',
         'Trends in medical safety incident counts by specialty'),
        ('Figure 2', 'fig2_irf_physicians.png',
         'IRF: physician count response to litigation shock'),
        ('Figure 3', 'sfig1_irf_facilities.png',
         'IRF: facility count response to litigation shock'),
        ('Figure 4', 'fig3_var_forecasts.png',
         'VAR-based workforce forecasts (2024\u20132033)'),
    ]

    for fig_label, filename, caption in figures:
        slide_layout = prs.slide_layouts[5]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Title
        from pptx.util import Emu
        from pptx.dml.color import RGBColor as PptxRGB

        txBox = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(0.2),
            PptxInches(12), PptxInches(0.6)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = fig_label
        p.font.size = PptxPt(24)
        p.font.bold = True

        # Image
        img_path = os.path.join(OUTPUT_DIR, filename)
        if os.path.exists(img_path):
            slide.shapes.add_picture(
                img_path,
                PptxInches(1.0), PptxInches(1.0),
                width=PptxInches(11.0)
            )

        # Caption
        txBox2 = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(6.8),
            PptxInches(12), PptxInches(0.6)
        )
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = caption
        p2.font.size = PptxPt(14)

    output_path = os.path.join(MANUSCRIPT_DIR, 'jma_figures_en.pptx')
    prs.save(output_path)
    print(f'PPTX figures saved to {output_path}')
    return output_path


if __name__ == '__main__':
    build_title_page()
    build_main_document()
    build_pptx_figures()
    print('\nAll JMA Journal outputs generated successfully.')
