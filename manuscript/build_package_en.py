#!/usr/bin/env python3
"""Supporting submission artifacts for the rate-based re-analysis, all derived
from the reproducible result files: editable figure PPTX, editable tables docx,
title page, cover letter, and STROBE checklist."""
import os, json
import pandas as pd
from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt

BASE = os.path.dirname(os.path.abspath(__file__))
PROJ = os.path.dirname(BASE)
OUT = os.path.join(PROJ, "output")
RES = json.load(open(os.path.join(PROJ, "results", "reanalysis_results.json")))

FIGS = [
    ("fig1_litigation_rate.png", "Figure 1",
     "Closed malpractice claims per 1,000 physicians by specialty, 2008–2024 (rates, not counts)."),
    ("fig2_physician_index.png", "Figure 2",
     "Physician workforce by specialty, indexed to 2008 (=100)."),
    ("fig3_equivalence.png", "Figure 3",
     "Equivalence (TOST) of the litigation-rate effect against ±1% and ±2% margins (90% CIs)."),
    ("fig4_counts_vs_rates.png", "Figure 4",
     "Biennial physician growth against lagged litigation exposure as counts (a) and rates (b)."),
]


def build_pptx():
    prs = Presentation()
    prs.slide_width = PInches(13.333); prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]
    for fn, num, cap in FIGS:
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(PInches(0.5), PInches(0.2), PInches(12.3), PInches(0.7))
        tf = tb.text_frame; tf.text = num
        tf.paragraphs[0].runs[0].font.size = PPt(24)
        tf.paragraphs[0].runs[0].font.bold = True
        img = os.path.join(OUT, fn)
        if os.path.exists(img):
            s.shapes.add_picture(img, PInches(1.2), PInches(1.1),
                                 height=PInches(5.2))
        cb = s.shapes.add_textbox(PInches(0.5), PInches(6.5), PInches(12.3), PInches(0.9))
        cf = cb.text_frame; cf.word_wrap = True; cf.text = cap
        cf.paragraphs[0].runs[0].font.size = PPt(14)
    p = os.path.join(BASE, "figures_en.pptx"); prs.save(p); print("wrote", p)


def _doc():
    doc = Document()
    for s in doc.sections:
        s.top_margin = s.bottom_margin = s.left_margin = s.right_margin = Cm(2.54)
    st = doc.styles["Normal"]; st.font.name = "Times New Roman"; st.font.size = Pt(11)
    return doc


def build_tables():
    doc = _doc()
    d = RES["descriptive"]["biennial_first_last"]
    by = d["by_specialty"]
    EN = list(by.keys())

    def tbl(title, headers, rows):
        h = doc.add_paragraph(); r = h.add_run(title); r.bold = True; r.font.size = Pt(11)
        t = doc.add_table(rows=1, cols=len(headers)); t.style = "Table Grid"
        for i, x in enumerate(headers):
            rr = t.rows[0].cells[i].paragraphs[0].add_run(str(x)); rr.bold = True
            rr.font.size = Pt(9)
        for row in rows:
            c = t.add_row().cells
            for i, v in enumerate(row):
                c[i].paragraphs[0].add_run(str(v)).font.size = Pt(9)
        doc.add_paragraph()

    tbl("Table 1. Primary data sources and their resolution.",
        ["Series", "Source", "Resolution", "Years", "Role"],
        [["Physicians by specialty", "MHLW Statistics of Physicians", "Biennial", "2004–2024", "Denominator & outcome"],
         ["Closed malpractice claims", "Supreme Court, by specialty", "Annual", "2008–2024", "Exposure"],
         ["Hospitals by specialty", "MHLW Survey of Medical Institutions", "Annual", "2008–2024", "Outcome"],
         ["Clinics by specialty", "MHLW Survey (static)", "Every 3 years", "2008–2023", "Descriptive"]])

    tbl(f"Table 2. Physicians, litigation rate (per 1,000 physicians) and hospitals by "
        f"specialty, {d['first']} vs {d['last']}.",
        ["Specialty", f"Phys {d['first']}", f"Phys {d['last']}",
         f"Rate {d['first']}", f"Rate {d['last']}", f"Hosp {d['first']}", f"Hosp {d['last']}"],
        [[k, by[k]["phys_first"], by[k]["phys_last"], f"{by[k]['litrate_first']:.2f}",
          f"{by[k]['litrate_last']:.2f}", by[k]["hosp_first"], by[k]["hosp_last"]] for k in EN])

    def g(coll, key):
        for x in RES[coll]:
            if key in x["label"]:
                return x
    P = g("primary", "physician growth ~ lagged litigation rate")
    H = g("primary", "hospital growth ~ lagged litigation rate")
    RV = g("primary", "Reverse")
    C = g("sensitivity", "COUNTS")
    A = g("sensitivity", "Annual hospital")
    I = g("sensitivity", "interpolated-annual")
    tbl("Table 3. Panel fixed-effects models and sensitivity analyses.",
        ["Model", "Coefficient", "95% CI", "p", "n"],
        [["Physician growth ~ lagged rate", f"{P['coef']:+.4f}", f"{P['ci_low']:+.4f}, {P['ci_high']:+.4f}", f"{P['p']:.2f}", P['n_obs']],
         ["Hospital growth ~ lagged rate", f"{H['coef']:+.4f}", f"{H['ci_low']:+.4f}, {H['ci_high']:+.4f}", f"{H['p']:.2f}", H['n_obs']],
         ["Counts contrast (physician)", f"{C['coef']:+.4f}", "—", f"{C['p']:.2f}", C['n_obs']],
         ["Annual hospital (sensitivity)", f"{A['coef']:+.4f}", "—", f"{A['p']:.2f}", A['n_obs']],
         ["Interpolated physician (sensitivity)", f"{I['coef']:+.4f}", "—", f"{I['p']:.2f}", I['n_obs']],
         ["Reverse (workforce→litigation)", f"{RV['coef']:+.3f}", "—", f"{RV['p']:.2f}", RV['n_obs']]])

    eqp, eqh = RES["equivalence"][0], RES["equivalence"][1]
    tbl("Table 4. Equivalence (TOST) results (effect per +1 SD litigation rate).",
        ["Outcome", "Coef per SD", "90% CI", "Margin", "TOST p", "Equivalent"],
        [[o, f"{e['coef_per_SD']*100:+.2f}%", f"{e['ci90_low']*100:+.2f}%, {e['ci90_high']*100:+.2f}%",
          f"±{int(t['margin']*100)}%", f"{t['p_tost']:.3f}", "Yes" if t['equivalent'] else "No"]
         for o, e in [("Physician growth", eqp), ("Hospital growth", eqh)] for t in e["tests"]])

    p = os.path.join(BASE, "tables_en.docx"); doc.save(p); print("wrote", p)


def build_title_page():
    doc = _doc()
    for _ in range(4):
        doc.add_paragraph()
    t = doc.add_paragraph(); t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = t.add_run("Litigation risk and specialty-level physician workforce in Japan, "
                  "2008–2024: a rate-based re-analysis with equivalence testing")
    r.bold = True; r.font.size = Pt(15)
    for lbl in ["Authors: [Author names]", "Affiliations: [Institutions]",
                "Corresponding author: [name, address, email]",
                "Word count (main text): [auto]", "Tables: 4  Figures: 4",
                "Conflicts of interest: none declared", "Funding: none"]:
        p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run(lbl).font.size = Pt(12)
    p = os.path.join(BASE, "title_page.docx"); doc.save(p); print("wrote", p)


def build_cover_letter():
    doc = _doc()
    for line in ["[Date]", "", "The Editor-in-Chief", "[Target Journal]", ""]:
        doc.add_paragraph(line)
    doc.add_paragraph("Dear Editor,")
    body = [
        "We submit an original research article, \"Litigation risk and specialty-level "
        "physician workforce in Japan, 2008–2024: a rate-based re-analysis with equivalence "
        "testing\", for consideration.",
        "Using national primary data for 12 specialties—biennial physician statistics, "
        "Supreme Court specialty-specific closed malpractice claims, and the annual medical-"
        "institution survey—we test whether litigation risk is associated with decline in the "
        "physician workforce and in hospitals offering each specialty. Crucially, we express "
        "exposure as a rate (claims per 1,000 physicians) to remove specialty-size confounding, "
        "analyse only measured biennial physician observations rather than interpolated annual "
        "values, and apply equivalence (TOST) testing.",
        "We find no association between litigation risk and workforce or hospital change; the "
        "effect is statistically equivalent to a null within a small, pre-specified margin. The "
        "common intuition that physicians avoid high-litigation specialties is not supported. "
        "We show that previously reported associations arise from size confounding and from "
        "treating interpolated years as independent observations, and we discuss how structural "
        "incentives—no-fault obstetric compensation (JOCS-CP) and procedure-based "
        "reimbursement—may sustain the specialty workforce.",
        "This manuscript is a substantially rebuilt analysis prepared after external peer "
        "review of an earlier version identified count-based modelling and interpolation as "
        "structural problems; we have addressed these by moving to rates, measured-only data, "
        "equivalence testing, explicit treatment of JOCS-CP, and correction of the Medical "
        "Accident Investigation System description. All data and code are openly available and "
        "every reported number is reproducible from raw files.",
        "The work is original, not under consideration elsewhere, and all authors approve the "
        "submission. We declare no conflicts of interest.",
    ]
    for b in body:
        p = doc.add_paragraph(); p.add_run(b).font.size = Pt(11)
        p.paragraph_format.space_after = Pt(8)
    doc.add_paragraph("Sincerely,"); doc.add_paragraph("[Corresponding author, on behalf of all authors]")
    p = os.path.join(BASE, "cover_letter.docx"); doc.save(p); print("wrote", p)


def build_strobe():
    doc = _doc()
    h = doc.add_paragraph(); r = h.add_run("STROBE Statement — checklist for observational studies")
    r.bold = True; r.font.size = Pt(13)
    items = [
        ("Title and abstract", "1", "Rate-based re-analysis design stated in title; structured abstract."),
        ("Background/rationale", "2", "Specialty maldistribution and the flight-from-risk hypothesis (Introduction)."),
        ("Objectives", "3", "Test association of litigation rate with workforce/hospital change; equivalence."),
        ("Study design", "4", "Ecological panel across 12 specialties, 2008–2024, biennial grid."),
        ("Setting", "5", "Japan; national administrative and judicial statistics."),
        ("Variables", "7", "Exposure: litigation rate per 1,000 physicians; outcomes: log-change in physicians/hospitals; JOCS-CP indicator."),
        ("Data sources/measurement", "8", "MHLW physician statistics, Supreme Court litigation, MHLW facility survey (Table 1)."),
        ("Bias", "9", "Size confounding removed by rates; interpolation avoided in primary analysis."),
        ("Study size", "10", "All 12 core specialties; 9 measured biennial waves."),
        ("Quantitative variables", "11", "Litigation rate; standardised exposure for TOST."),
        ("Statistical methods", "12", "Panel FE (specialty, wave), cluster-robust SE, TOST, sensitivity analyses."),
        ("Descriptive/main results", "13-16", "Trends (Figs 1–2, Table 2), null panel and equivalence (Table 3–4, Fig 3), counts vs rates (Fig 4)."),
        ("Limitations", "19", "Ecological; biennial resolution; litigation ≠ incident risk; 2004–2007 unavailable."),
        ("Interpretation", "20", "No support for flight-from-risk; structural incentives discussed."),
        ("Generalisability", "21", "Japan specialty-level; not individual careers."),
        ("Funding", "22", "None."),
    ]
    t = doc.add_table(rows=1, cols=3); t.style = "Table Grid"
    for i, x in enumerate(["Item", "No.", "Location/《addressed》"]):
        rr = t.rows[0].cells[i].paragraphs[0].add_run(x); rr.bold = True; rr.font.size = Pt(9)
    for a, b, c in items:
        cells = t.add_row().cells
        for i, v in enumerate([a, b, c]):
            cells[i].paragraphs[0].add_run(v).font.size = Pt(9)
    p = os.path.join(BASE, "strobe_checklist.docx"); doc.save(p); print("wrote", p)


if __name__ == "__main__":
    build_pptx(); build_tables(); build_title_page(); build_cover_letter(); build_strobe()
