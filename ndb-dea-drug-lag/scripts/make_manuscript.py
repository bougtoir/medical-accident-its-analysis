#!/usr/bin/env python3
"""
Generate the JA and EN manuscripts (docx, figures inline), separate editable
table docx, and editable figure pptx for the HCV DAA new-treatment-anticipation /
drug-lag study.

All result numbers are read from results/summary.json and data/*.csv and formatted
at runtime; no result value is hard-coded here. References are numbered in order of
first appearance (Vancouver). Figures/tables are inserted immediately after the
paragraph that first cites them.

Usage:
    python3 scripts/make_manuscript.py
"""
import json
import os

import pandas as pd
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt

BASE = os.path.join(os.path.dirname(__file__), "..")
DATA = os.path.join(BASE, "data")
RES = os.path.join(BASE, "results")
OUT = os.path.join(BASE, "output")

S = json.load(open(os.path.join(RES, "summary.json"), encoding="utf-8"))
TS = pd.read_csv(os.path.join(DATA, "hcv_timeseries.csv")).set_index("fy")
EV = pd.read_csv(os.path.join(DATA, "announcement_events.csv"))

Y0, Y1 = S["fiscal_year_range"]


def fmt(x, nd=0):
    return f"{x:,.{nd}f}"


# ---- derived display values (all traceable to summary/timeseries) -------------
peg_drop = abs(S["peginterferon"]["pct_change_first_to_last"])
rbv_zero_year = S["ribavirin"]["first_year_near_zero"]
daa_peak_fy = S["daa_total"]["peak_fy"]
daa_peak_val_m = S["daa_total"]["peak_value"] / 1e6
daa_rise = S["daa_total"]["pct_change_first_to_peak"]
daa_fall = abs(S["daa_total"]["pct_change_peak_to_last"])
daa_last_m = S["daa_total"]["fy_last"] / 1e6
n_daa = S["n_distinct_daa_products"]

# ------------------------------------------------------------------------------
# Reference text per source id (no fabricated citations; verifiable sources only).
# Citation NUMBERS are assigned dynamically in order of first appearance (Vancouver),
# tracked in CITE_ORDER during document build.
REF_TEXT = {
    "en": {
        "ndb": "Ministry of Health, Labour and Welfare (Japan). NDB Open Data "
               "(1st-10th editions). https://www.mhlw.go.jp/ndb/opendatasite/ "
               f"(accessed for fiscal years {Y0}-{Y1}).",
        "bms": "Bristol-Myers Squibb K.K. Press release: approval in Japan of the "
               "world's first all-oral, interferon- and ribavirin-free treatment for "
               "chronic hepatitis C (daclatasvir + asunaprevir). 4 July 2014.",
        "nhi": "Ministry of Health, Labour and Welfare / Central Social Insurance "
               "Medical Council (Chuikyo). NHI drug-price listings of direct-acting "
               "antivirals for hepatitis C (2014-2017).",
    },
    "ja": {
        "ndb": "厚生労働省. NDBオープンデータ（第1〜10回）. "
               "https://www.mhlw.go.jp/ndb/opendatasite/ "
               f"（{Y0}〜{Y1}年度分を使用）.",
        "bms": "ブリストル・マイヤーズ株式会社. プレスリリース：日本初の"
               "インターフェロンおよびリバビリンを必要としない経口薬のみによる"
               "C型慢性肝炎治療薬（ダクラタスビル＋アスナプレビル）の製造販売承認取得. "
               "2014年7月4日.",
        "nhi": "厚生労働省／中央社会保険医療協議会（中医協）. C型肝炎に対する"
               "直接作用型抗ウイルス薬の薬価基準収載（2014〜2017年）.",
    },
}

CITE_ORDER = []  # reset per document; source ids in order of first appearance


def cite(p, keys, lang):
    """Append a Vancouver superscript citation, numbering by first appearance."""
    for k in keys:
        if k not in CITE_ORDER:
            CITE_ORDER.append(k)
    nums = sorted(CITE_ORDER.index(k) + 1 for k in keys)
    run = p.add_run(",".join(str(n) for n in nums))
    run.font.superscript = True
    return p


def add_caption(doc, text):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(14)
    r = p.add_run(text)
    r.bold = True
    r.font.size = Pt(9)
    return p


def insert_fig(doc, path, width=6.3):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run().add_picture(path, width=Inches(width))


# ------------------------------------------------------------------------------
TXT = {
    "en": dict(
        title="Population-level replacement of interferon-based standard therapy for "
              "hepatitis C after interferon-free direct-acting antivirals: a "
              "new-treatment-anticipation / practical-drug-lag analysis of Japan's NDB "
              "Open Data",
        h_abs="Abstract", h_intro="Introduction", h_meth="Methods",
        h_res="Results", h_disc="Discussion", h_lim="Limitations", h_conc="Conclusion",
        h_ref="References", h_da="Data and code availability",
        abs_bg="Background: ", abs_me="Methods: ", abs_re="Results: ", abs_co="Conclusion: ",
    ),
    "ja": dict(
        title="IFNフリー直接作用型抗ウイルス薬導入後のC型肝炎インターフェロン標準治療の"
              "人口レベルでの置換：NDBオープンデータによる新治療法待望論／実用的"
              "ドラッグラグの検証",
        h_abs="要旨", h_intro="緒言", h_meth="方法",
        h_res="結果", h_disc="考察", h_lim="限界", h_conc="結論",
        h_ref="文献", h_da="データおよびコードの入手可能性",
        abs_bg="背景：", abs_me="方法：", abs_re="結果：", abs_co="結論：",
    ),
}


def build_manuscript(lang):
    T = TXT[lang]
    CITE_ORDER.clear()
    doc = Document()
    st = doc.styles["Normal"].font
    st.size = Pt(10.5)

    ttl = doc.add_paragraph()
    r = ttl.add_run(T["title"]); r.bold = True; r.font.size = Pt(14)

    # ---- Abstract ----
    doc.add_heading(T["h_abs"], level=1)
    ab = doc.add_paragraph()
    if lang == "en":
        ab.add_run(T["abs_bg"]).bold = True
        ab.add_run("If population-level anticipation of a newly reported/reimbursed "
                   "treatment is real, use of the prior standard therapy should fall "
                   "sharply once the awaited option arrives; if not, either the "
                   "anticipation/practical lag is absent or decision-makers are not "
                   "reacting to the news. ")
        ab.add_run(T["abs_me"]).bold = True
        ab.add_run(f"Using Japan's NDB Open Data ({Y0}-{Y1}), we tracked national "
                   "dispensed quantity of interferon (IFN)-based standard therapy "
                   "(peginterferon, ribavirin) and of interferon-free direct-acting "
                   "antivirals (DAAs) for hepatitis C, against official approval/"
                   "reimbursement milestones. ")
        ab.add_run(T["abs_re"]).bold = True
        ab.add_run(f"Peginterferon dispensing fell {fmt(peg_drop,1)}% from FY{Y0} to "
                   f"FY{Y1}; ribavirin reached near-zero by FY{rbv_zero_year}. Total DAA "
                   f"dispensing peaked in FY{daa_peak_fy} ({fmt(daa_peak_val_m,1)} million "
                   f"units, +{fmt(daa_rise,0)}% vs FY{Y0}) and then fell {fmt(daa_fall,0)}% "
                   f"by FY{Y1}. ")
        ab.add_run(T["abs_co"]).bold = True
        ab.add_run("The IFN-based standard therapy was replaced at the population level "
                   "within about two years, and the DAA surge-then-decay is consistent "
                   "with a finite pool of long-waiting patients being treated in a burst "
                   "(pent-up demand) rather than steady substitution.")
    else:
        ab.add_run(T["abs_bg"]).bold = True
        ab.add_run("新治療法の報道・収載に対する待望が人口レベルで実在するなら、待望された"
                   "治療が登場した時点で従来の標準治療の利用は急減するはずである。変化がなければ、"
                   "待望論・実用的ラグが存在しないか、治療選択の主体が報道に反応していないことになる。")
        ab.add_run(T["abs_me"]).bold = True
        ab.add_run(f"日本のNDBオープンデータ（{Y0}〜{Y1}年度）を用い、C型肝炎に対する"
                   "インターフェロン（IFN）ベース標準治療（ペグインターフェロン、リバビリン）と"
                   "IFNフリー直接作用型抗ウイルス薬（DAA）の全国処方数量を、公式の承認・"
                   "薬価収載イベントと対応づけて追跡した。")
        ab.add_run(T["abs_re"]).bold = True
        ab.add_run(f"ペグインターフェロンの処方数量はFY{Y0}からFY{Y1}で{fmt(peg_drop,1)}%減少し、"
                   f"リバビリンはFY{rbv_zero_year}までにほぼゼロとなった。DAA合計はFY{daa_peak_fy}に"
                   f"ピーク（{fmt(daa_peak_val_m,1)}百万単位、FY{Y0}比+{fmt(daa_rise,0)}%）を示し、"
                   f"その後FY{Y1}までに{fmt(daa_fall,0)}%減少した。")
        ab.add_run(T["abs_co"]).bold = True
        ab.add_run("IFNベース標準治療は約2年で人口レベルに置換され、DAAの急増→減衰は、"
                   "長く待機していた患者集団が一括して治療された（待望＝pent-up demand）ことと"
                   "整合的であり、定常的な置換フローではないことを示す。")

    # ---- Introduction ----
    doc.add_heading(T["h_intro"], level=1)
    p = doc.add_paragraph()
    if lang == "en":
        p.add_run("When a new therapy is announced through the media or scheduled for "
                  "reimbursement, advocates of the new treatment expect a population-wide "
                  "shift away from the current standard of care. Chronic hepatitis C offers "
                  "an unusually clean natural experiment: interferon-based therapy was the "
                  "standard until interferon-free DAAs arrived in Japan in 2014-2015, when "
                  "daclatasvir plus asunaprevir became the world's first all-oral, "
                  "interferon- and ribavirin-free regimen")
        cite(p, ["bms"], lang)
        p.add_run(". We use nationwide dispensing counts from NDB Open Data")
        cite(p, ["ndb"], lang)
        p.add_run(" to test whether the standard therapy was displaced at the population "
                  "level around these announcement/reimbursement events.")
    else:
        p.add_run("新しい治療法が報道され、あるいは保険収載が予定されると、新治療法待望の立場からは"
                  "従来の標準治療から人口規模での移行が期待される。C型慢性肝炎は例外的に明瞭な自然実験を"
                  "提供する。インターフェロンベース治療が長く標準であったが、2014〜2015年にIFNフリーDAAが"
                  "登場し、ダクラタスビル＋アスナプレビルが世界初の全経口・IFN／リバビリン不要療法となった")
        cite(p, ["bms"], lang)
        p.add_run("。本研究ではNDBオープンデータ")
        cite(p, ["ndb"], lang)
        p.add_run("の全国処方数量を用い、これらの承認・収載イベント前後で標準治療が人口レベルで"
                  "置換されたかを検証する。")

    # ---- Methods ----
    doc.add_heading(T["h_meth"], level=1)
    p = doc.add_paragraph()
    if lang == "en":
        p.add_run(f"We used NDB Open Data editions 1-10 (fiscal years {Y0}-{Y1})")
        cite(p, ["ndb"], lang)
        p.add_run(", extracting the national total dispensed quantity (総計/処方数量) for "
                  "each drug from the sex- and age-stratified prescription-drug tables "
                  "(oral, topical, injectable). Drugs were classified by their actual "
                  f"product names into IFN-based standard therapy (peginterferon, "
                  f"ribavirin) and {n_daa} interferon-free DAA products. Conventional "
                  "interferon is reported separately because it is not hepatitis-C-specific. "
                  "The metric is dispensed quantity (tablets/capsules for oral drugs, "
                  "syringes/vials for injections), not a patient count, and is compared "
                  "within a product over time. Official approval and NHI drug-price "
                  "listing milestones")
        cite(p, ["bms", "nhi"], lang)
        p.add_run(" were used as intervention markers at fiscal-year resolution "
                  "(Table 1). Analyses are descriptive; the full pipeline "
                  "(download -> build -> analyze -> figures) is openly reproducible.")
    else:
        p.add_run(f"NDBオープンデータ第1〜10回（{Y0}〜{Y1}年度）")
        cite(p, ["ndb"], lang)
        p.add_run("を用い、処方薬の性年齢別薬効分類別数量表（内服・外用・注射）から各薬剤の全国"
                  "総計（処方数量）を抽出した。薬剤は実際の製品名からIFNベース標準治療（ペグ"
                  f"インターフェロン、リバビリン）と{n_daa}製剤のIFNフリーDAAに分類した。従来型"
                  "インターフェロンはC型肝炎特異的でないため別掲した。指標は処方数量（内服は錠・"
                  "カプセル、注射はシリンジ・バイアル）であり患者数ではなく、製剤ごとに経時比較する。"
                  "公式の承認・薬価収載イベント")
        cite(p, ["bms", "nhi"], lang)
        p.add_run("を年度解像度の介入マーカーとした（表1）。解析は記述的であり、全パイプライン"
                  "（download→build→analyze→figures）は公開・再現可能である。")

    add_caption(doc, ("Table 1. Official approval / NHI drug-price listing milestones "
                      "used as intervention markers." if lang == "en"
                      else "表1．介入マーカーとして用いた公式の承認・薬価収載イベント。"))
    add_events_table(doc, lang)

    # ---- Results ----
    doc.add_heading(T["h_res"], level=1)
    p = doc.add_paragraph()
    if lang == "en":
        p.add_run(f"The interferon-based standard therapy collapsed after IFN-free DAAs "
                  f"became available (Fig. 1). Peginterferon dispensing fell "
                  f"{fmt(peg_drop,1)}% from FY{Y0} to FY{Y1}, and ribavirin reached "
                  f"near-zero by FY{rbv_zero_year}.")
    else:
        p.add_run(f"IFNフリーDAAの登場後、インターフェロンベース標準治療は消失した（図1）。"
                  f"ペグインターフェロンの処方数量はFY{Y0}からFY{Y1}で{fmt(peg_drop,1)}%減少し、"
                  f"リバビリンはFY{rbv_zero_year}までにほぼゼロとなった。")
    insert_fig(doc, os.path.join(OUT, f"fig1_ifn_collapse_{lang}.png"))
    add_caption(doc, ("Fig. 1. Collapse of interferon-based standard therapy "
                      "(indexed to FY%d = 100) with total DAA dispensed quantity and "
                      "dated announcement markers." % Y0 if lang == "en"
                      else "図1．インターフェロンベース標準治療の消失（FY%d=100指数）、"
                           "DAA合計処方数量、および収載イベント。" % Y0))

    p = doc.add_paragraph()
    if lang == "en":
        p.add_run(f"DAA use itself showed a surge-then-decay pattern (Fig. 2): total DAA "
                  f"dispensing peaked in FY{daa_peak_fy} at {fmt(daa_peak_val_m,1)} million "
                  f"units (+{fmt(daa_rise,0)}% vs FY{Y0}) and then declined {fmt(daa_fall,0)}% "
                  f"to {fmt(daa_last_m,1)} million units by FY{Y1}, with successive products "
                  "replacing earlier ones.")
    else:
        p.add_run(f"DAAの利用自体は急増→減衰パターンを示した（図2）。DAA合計はFY{daa_peak_fy}に"
                  f"{fmt(daa_peak_val_m,1)}百万単位（FY{Y0}比+{fmt(daa_rise,0)}%）でピークに達し、"
                  f"その後FY{Y1}までに{fmt(daa_fall,0)}%減少し{fmt(daa_last_m,1)}百万単位となった。"
                  "後発の製剤が先行製剤を置換した。")
    insert_fig(doc, os.path.join(OUT, f"fig2_daa_wave_{lang}.png"))
    add_caption(doc, ("Fig. 2. The DAA wave: dispensed quantity by product." if lang == "en"
                      else "図2．DAAの波：製剤別処方数量。"))

    # ---- Discussion ----
    doc.add_heading(T["h_disc"], level=1)
    p = doc.add_paragraph()
    if lang == "en":
        p.add_run("For hepatitis C DAAs there was no population-level practical lag: the "
                  "standard therapy was displaced within about two years of the "
                  "announcements/listings, indicating that decision-makers reacted rapidly. "
                  "The DAA surge-then-decay is consistent with a finite stock of "
                  "long-waiting patients being cured in a burst (pent-up demand) rather "
                  "than a steady replacement flow. This is descriptive evidence supporting "
                  "the anticipation hypothesis; it does not by itself establish that media "
                  "coverage caused individual treatment choices.")
    else:
        p.add_run("C型肝炎DAAでは人口レベルの実用的ラグは認められなかった。標準治療は報道・収載から"
                  "約2年で置換され、治療選択の主体が迅速に反応したことを示す。DAAの急増→減衰は、長く"
                  "待機していた患者ストックが一括して治癒された（pent-up demand）ことと整合的であり、"
                  "定常的置換フローではない。これは待望論仮説を支持する記述的証拠であり、報道が個々の"
                  "治療選択を引き起こしたことを単独で証明するものではない。")

    doc.add_heading(T["h_lim"], level=2)
    p = doc.add_paragraph()
    if lang == "en":
        p.add_run(f"NDB Open Data begins in FY{Y0}, the same period IFN-free DAAs launched, "
                  f"so there is no pre-DAA interferon baseline within NDB; the FY{Y0} value "
                  "already reflects decline from the pre-2014 peak. The metric is dispensed "
                  "quantity, not patient counts, and units differ across products, so the "
                  "summed DAA quantity is not a patient count. Data are annual, precluding "
                  "within-year interrupted time-series or formal causal estimation, and no "
                  "control condition or placebo event is included.")
    else:
        p.add_run(f"NDBオープンデータはFY{Y0}開始であり、これはIFNフリーDAA導入時期と重なるため、"
                  f"NDB内にDAA前のインターフェロン基準値は存在しない（FY{Y0}値は既に2014年以前の"
                  "ピークからの減少を反映）。指標は処方数量であり患者数ではなく、製剤間で単位が"
                  "異なるためDAA数量の合計は患者数ではない。データは年次であり、年内の中断時系列や"
                  "形式的因果推定はできず、対照条件・プラセボイベントも含まない。")

    # ---- Conclusion ----
    doc.add_heading(T["h_conc"], level=1)
    p = doc.add_paragraph()
    if lang == "en":
        p.add_run("In hepatitis C, the interferon-based standard therapy was replaced at "
                  "the population level within about two years of interferon-free DAAs, and "
                  "the surge-then-decay of DAA use is consistent with realized pent-up "
                  "demand. This supports the existence of population-level new-treatment "
                  "anticipation for this case, while cautioning that the same pattern need "
                  "not generalize to therapies with weaker anticipation or greater "
                  "practical constraints.")
    else:
        p.add_run("C型肝炎では、インターフェロンベース標準治療はIFNフリーDAA導入から約2年で人口"
                  "レベルに置換され、DAA利用の急増→減衰は待機需要の顕在化と整合的であった。これは"
                  "本事例における人口レベルの新治療法待望の存在を支持する一方、待望が弱い、あるいは"
                  "実用的制約の大きい治療には同じパターンが一般化するとは限らないことに注意を要する。")

    # ---- Data/code availability + references ----
    doc.add_heading(T["h_da"], level=1)
    p = doc.add_paragraph()
    p.add_run("Raw NDB Open Data workbooks are publicly available from MHLW and are "
              "re-downloadable via scripts/download_ndb.py; derived datasets, analysis "
              "code and figure/manuscript generators are in the project repository."
              if lang == "en" else
              "NDBオープンデータの元ファイルは厚生労働省より公開され、scripts/download_ndb.py で"
              "再取得できる。派生データ・解析コード・図表／原稿生成スクリプトはプロジェクトリポジトリ"
              "に含まれる。")

    doc.add_heading(T["h_ref"], level=1)
    for i, k in enumerate(CITE_ORDER, 1):
        rp = doc.add_paragraph()
        rp.paragraph_format.left_indent = Inches(0.3)
        rp.paragraph_format.first_line_indent = Inches(-0.3)
        rp.add_run(f"{i}. {REF_TEXT[lang][k]}")

    path = os.path.join(OUT, f"manuscript_{lang}.docx")
    doc.save(path)
    print("wrote", path)


def add_events_table(doc, lang):
    cols = (["FY", "Month", "Drug", "Milestone", "Source"] if lang == "en"
            else ["年度", "時期", "薬剤", "イベント", "出所"])
    t = doc.add_table(rows=1, cols=len(cols))
    t.style = "Table Grid"
    for j, c in enumerate(cols):
        r = t.rows[0].cells[j].paragraphs[0].add_run(c); r.bold = True
    drug_col = "drug_en" if lang == "en" else "drug_ja"
    mile = {"approval": "approval" if lang == "en" else "承認",
            "nhi_listing": "NHI listing" if lang == "en" else "薬価収載"}
    for _, e in EV.iterrows():
        cells = t.add_row().cells
        cells[0].text = str(int(e["fy"]))
        cells[1].text = str(e["event_month"])
        cells[2].text = str(e[drug_col])
        cells[3].text = mile.get(e["milestone"], e["milestone"])
        cells[4].text = str(e["source"])
    for row in t.rows:
        for cell in row.cells:
            for pph in cell.paragraphs:
                for rr in pph.runs:
                    rr.font.size = Pt(7.5)


def build_tables_doc(lang):
    doc = Document()
    doc.add_heading("Tables" if lang == "en" else "表", level=1)
    add_caption(doc, ("Table 1. Official approval / NHI drug-price listing milestones."
                      if lang == "en"
                      else "表1．公式の承認・薬価収載イベント。"))
    add_events_table(doc, lang)

    add_caption(doc, ("Table 2. National dispensed quantity by drug group and fiscal year "
                      "(NDB Open Data)." if lang == "en"
                      else "表2．薬効グループ別・年度別の全国処方数量（NDBオープンデータ）。"))
    groups = ["IFN_peg", "ribavirin", "IFN_conv", "DAA"]
    head = (["Fiscal year", "Peginterferon", "Ribavirin", "Conventional IFN", "DAA total"]
            if lang == "en"
            else ["年度", "ペグIFN", "リバビリン", "従来型IFN", "DAA合計"])
    t = doc.add_table(rows=1, cols=len(head)); t.style = "Table Grid"
    for j, c in enumerate(head):
        rr = t.rows[0].cells[j].paragraphs[0].add_run(c); rr.bold = True
    for fy, row in TS.iterrows():
        cells = t.add_row().cells
        cells[0].text = str(int(fy))
        for j, g in enumerate(groups, 1):
            cells[j].text = fmt(row[g], 0)
    path = os.path.join(OUT, f"tables_{lang}.docx")
    doc.save(path)
    print("wrote", path)


def build_pptx(lang):
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    figs = [
        (f"fig1_ifn_collapse_{lang}.png",
         "Fig. 1. Collapse of interferon-based standard therapy" if lang == "en"
         else "図1．インターフェロンベース標準治療の消失"),
        (f"fig2_daa_wave_{lang}.png",
         "Fig. 2. The DAA wave: dispensed quantity by product" if lang == "en"
         else "図2．DAAの波：製剤別処方数量"),
    ]
    for fn, cap in figs:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(PInches(0.5), PInches(0.2), PInches(12.3), PInches(0.7))
        tf = tb.text_frame; tf.text = cap
        tf.paragraphs[0].runs[0].font.size = PPt(20)
        tf.paragraphs[0].runs[0].font.bold = True
        slide.shapes.add_picture(os.path.join(OUT, fn), PInches(1.4), PInches(1.1),
                                 height=PInches(5.7))
    path = os.path.join(OUT, f"figures_{lang}.pptx")
    prs.save(path)
    print("wrote", path)


if __name__ == "__main__":
    for lang in ("en", "ja"):
        build_manuscript(lang)
        build_tables_doc(lang)
        build_pptx(lang)
