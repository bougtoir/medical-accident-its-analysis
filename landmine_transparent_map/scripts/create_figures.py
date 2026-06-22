"""
Generate figures for the transparent minefield paper.
Output: individual PNG files + combined PPTX.
Target: Journal on Baltic Security
"""

import numpy as np
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN

from model_simulation import (
    MineFieldParams, delay_map, delay_blind, delay_full_intel,
    postconflict_time_map, postconflict_time_blind, postconflict_residual_blind,
    entropy_with_map, entropy_no_map, entropy_reduction,
    positions_in_breach, clearance_efficiency, delay_ratio_vs_blind,
    normalised_delay, normalised_clearance_cost, welfare,
)

# Setup
OUTPUT_DIR = Path(__file__).parent.parent / "output"
OUTPUT_DIR.mkdir(exist_ok=True)
params = MineFieldParams()
r_values = np.linspace(0, 20, 201)

plt.rcParams.update({
    'font.size': 11,
    'axes.labelsize': 12,
    'axes.titlesize': 13,
    'legend.fontsize': 10,
    'figure.dpi': 300,
    'savefig.dpi': 300,
    'savefig.bbox': 'tight',
})


# ─── Figure 1: Conceptual Diagram ────────────────────────────────────────────

def fig1_concept():
    """Conceptual diagram showing the three regimes."""
    fig, axes = plt.subplots(1, 3, figsize=(12, 4))

    for ax in axes:
        ax.set_xlim(0, 10)
        ax.set_ylim(0, 6)
        ax.set_aspect('equal')
        ax.set_xticks([])
        ax.set_yticks([])
        # Draw minefield boundary
        rect = mpatches.FancyBboxPatch((0.5, 0.5), 9, 5, boxstyle="round,pad=0.1",
                                        edgecolor='grey', facecolor='#f5f5f5', linewidth=1.5)
        ax.add_patch(rect)

    np.random.seed(42)
    real_x = np.random.uniform(1, 9, 15)
    real_y = np.random.uniform(1, 5, 15)

    # Panel A: No map (current)
    ax = axes[0]
    ax.set_title('(a) No map (status quo)', fontweight='bold')
    # Mines are hidden - show as question marks
    for x, y in zip(real_x, real_y):
        ax.text(x, y, '?', fontsize=8, ha='center', va='center',
                color='black', alpha=0.4)
    ax.text(5, -0.3, 'Mines hidden; entire area uncertain',
            ha='center', fontsize=9, style='italic')

    # Panel B: Full intelligence
    ax = axes[1]
    ax.set_title('(b) Full intelligence (spy)', fontweight='bold')
    ax.scatter(real_x, real_y, c='red', s=60, marker='^', zorder=5,
              label='Real mines (known)')
    ax.legend(loc='upper right', fontsize=8)
    ax.text(5, -0.3, 'All real positions known; dummies irrelevant',
            ha='center', fontsize=9, style='italic')

    # Panel C: Transparent map with dummies
    ax = axes[2]
    ax.set_title('(c) Proposed: map with dummies (r=2)', fontweight='bold')
    ax.scatter(real_x, real_y, c='red', s=60, marker='^', zorder=5)
    # Add dummies (2× real)
    np.random.seed(123)
    dummy_x = np.random.uniform(1, 9, 30)
    dummy_y = np.random.uniform(1, 5, 30)
    ax.scatter(dummy_x, dummy_y, c='blue', s=40, marker='o', zorder=4, alpha=0.7)
    # Legend showing indistinguishable markers
    all_x = np.concatenate([real_x, dummy_x])
    all_y = np.concatenate([real_y, dummy_y])
    # Overlay with uniform markers to show indistinguishability
    ax.scatter(all_x, all_y, c='none', edgecolors='black', s=80, marker='D',
              zorder=6, linewidths=0.8)
    legend_elements = [
        mpatches.Patch(facecolor='none', edgecolor='black', label='Published positions (indistinguishable)'),
        plt.Line2D([0], [0], marker='^', color='w', markerfacecolor='red',
                   markersize=8, label='Real (hidden identity)'),
        plt.Line2D([0], [0], marker='o', color='w', markerfacecolor='blue',
                   markersize=8, label='Dummy (hidden identity)'),
    ]
    ax.legend(handles=legend_elements, loc='upper right', fontsize=7)
    ax.text(5, -0.3, 'All positions published; real vs dummy unknown',
            ha='center', fontsize=9, style='italic')

    plt.tight_layout()
    plt.savefig(OUTPUT_DIR / "fig1_concept.png")
    plt.close()
    print("  Fig 1: Conceptual diagram saved")


# ─── Figure 2: Breach Delay vs Dummy Ratio ───────────────────────────────────

def fig2_delay():
    """Breach delay as function of dummy ratio r."""
    fig, ax = plt.subplots(figsize=(8, 5))

    delays = [delay_map(params, r) for r in r_values]
    d_blind = delay_blind(params)
    d_intel = delay_full_intel(params)

    ax.plot(r_values, delays, 'b-', linewidth=2, label='Map regime (proposed)')
    ax.axhline(d_blind, color='red', linestyle='--', linewidth=1.5,
               label=f'Blind breach (no map): {d_blind:.0f} min')
    ax.axhline(d_intel, color='green', linestyle=':', linewidth=1.5,
               label=f'Full intelligence: {d_intel:.0f} min')

    # Annotate key points
    for r_ann in [3, 5, 10]:
        d = delay_map(params, r_ann)
        ax.annotate(f'r={r_ann}\n{d:.0f} min',
                    xy=(r_ann, d), xytext=(r_ann + 1.5, d + 30),
                    fontsize=9, ha='left',
                    arrowprops=dict(arrowstyle='->', color='grey'))

    ax.set_xlabel('Dummy ratio (r)')
    ax.set_ylabel('Breach delay (minutes)')
    ax.set_title('Attacker breach delay under transparent map regime')
    ax.legend(loc='lower right')
    ax.set_xlim(0, 20)
    ax.set_ylim(0, max(delays) * 1.1)
    ax.grid(True, alpha=0.3)

    # Secondary y-axis: hours
    ax2 = ax.twinx()
    ax2.set_ylim(0, max(delays) * 1.1 / 60)
    ax2.set_ylabel('Hours')

    plt.tight_layout()
    plt.savefig(OUTPUT_DIR / "fig2_breach_delay.png")
    plt.close()
    print("  Fig 2: Breach delay saved")


# ─── Figure 3: Post-Conflict Clearance ───────────────────────────────────────

def fig3_clearance():
    """Post-conflict clearance: time and residual risk."""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 5))

    # Panel A: Clearance time
    clear_times = [postconflict_time_map(params, r) for r in r_values]
    t_blind = postconflict_time_blind(params)

    ax1.plot(r_values, clear_times, 'b-', linewidth=2, label='Map-based clearance')
    ax1.axhline(t_blind, color='red', linestyle='--', linewidth=1.5,
                label=f'Blind sweep: {t_blind:.0f} hours')
    ax1.fill_between(r_values, clear_times, t_blind, alpha=0.1, color='green',
                     label='Time saved vs blind')
    ax1.set_xlabel('Dummy ratio (r)')
    ax1.set_ylabel('Clearance time (hours)')
    ax1.set_title('(a) Post-conflict clearance time')
    ax1.legend(loc='upper left')
    ax1.set_xlim(0, 20)
    ax1.grid(True, alpha=0.3)

    # Panel B: Residual risk comparison
    categories = ['No map\n(blind sweep)', 'Map, r=0', 'Map, r=3',
                  'Map, r=5', 'Map, r=10']
    residuals = [postconflict_residual_blind(params), 0, 0, 0, 0]
    colors = ['red', 'green', 'green', 'green', 'green']

    bars = ax2.bar(categories, residuals, color=colors, alpha=0.7, edgecolor='black')
    ax2.set_ylabel('Expected residual mines')
    ax2.set_title('(b) Residual mines after clearance')
    ax2.set_ylim(0, max(residuals) * 1.3)

    # Annotate
    for i, (cat, val) in enumerate(zip(categories, residuals)):
        ax2.text(i, val + 0.3, f'{val:.0f}', ha='center', fontweight='bold',
                fontsize=11)

    ax2.axhline(0, color='black', linewidth=0.5)
    ax2.text(2.5, max(residuals) * 0.8,
             'Map regime guarantees\nzero residual mines',
             ha='center', fontsize=10, style='italic',
             bbox=dict(boxstyle='round', facecolor='lightgreen', alpha=0.3))

    plt.tight_layout()
    plt.savefig(OUTPUT_DIR / "fig3_postconflict.png")
    plt.close()
    print("  Fig 3: Post-conflict clearance saved")


# ─── Figure 4: Information Entropy ────────────────────────────────────────────

def fig4_entropy():
    """Information entropy as function of r."""
    fig, ax = plt.subplots(figsize=(8, 5))

    entropies = [entropy_with_map(params, r) for r in r_values]
    H_no = entropy_no_map(params)

    ax.plot(r_values, entropies, 'purple', linewidth=2,
            label='Entropy with map H(r)')
    ax.axhline(H_no, color='red', linestyle='--', linewidth=1.5,
               label=f'No map (maximum uncertainty): {H_no:.0f} bits')
    ax.axhline(0, color='green', linestyle=':', linewidth=1.5,
               label='Full intelligence: 0 bits')

    # Shade the "controlled uncertainty" region
    ax.fill_between(r_values, 0, entropies, alpha=0.1, color='purple')

    ax.set_xlabel('Dummy ratio (r)')
    ax.set_ylabel('Positional entropy (bits)')
    ax.set_title('Attacker uncertainty under transparent map regime')
    ax.legend(loc='lower right')
    ax.set_xlim(0, 20)
    ax.grid(True, alpha=0.3)

    # Add right axis showing percentage
    ax2 = ax.twinx()
    ax2.set_ylim(0, max(entropies) * 1.05 / H_no * 100)
    ax2.set_ylabel('% of maximum uncertainty')

    plt.tight_layout()
    plt.savefig(OUTPUT_DIR / "fig4_entropy.png")
    plt.close()
    print("  Fig 4: Information entropy saved")


# ─── Figure 5: Policy Trade-off Summary ──────────────────────────────────────

def fig5_tradeoff():
    """Summary: delay vs clearance cost trade-off with Pareto frontier."""
    fig, ax = plt.subplots(figsize=(8, 6))

    delays_norm = [min(delay_map(params, r) / delay_blind(params), 1.5)
                   for r in r_values]
    clears_norm = [postconflict_time_map(params, r) / postconflict_time_blind(params)
                   for r in r_values]

    # Plot the trade-off curve
    scatter = ax.scatter(clears_norm, delays_norm, c=r_values, cmap='viridis',
                        s=20, zorder=3)
    plt.colorbar(scatter, ax=ax, label='Dummy ratio (r)')

    # Annotate key points
    for r_ann in [0, 1, 3, 5, 10, 15, 20]:
        d = min(delay_map(params, r_ann) / delay_blind(params), 1.5)
        c = postconflict_time_map(params, r_ann) / postconflict_time_blind(params)
        ax.annotate(f'r={r_ann}', xy=(c, d), fontsize=8, ha='left',
                    xytext=(c + 0.01, d + 0.02))

    # Reference points
    ax.plot(1.0, 1.0, 'rs', markersize=12, label='No map (status quo)', zorder=5)
    ax.plot(0, 0, 'g^', markersize=12, label='No mines (full ban)', zorder=5)

    ax.set_xlabel('Post-conflict clearance cost\n(fraction of blind sweep time)')
    ax.set_ylabel('Military utility\n(breach delay / blind sweep delay)')
    ax.set_title('Policy trade-off: military utility vs humanitarian cost')
    ax.legend(loc='upper left')
    ax.grid(True, alpha=0.3)

    # Ideal region annotation
    ax.annotate('Ideal region:\nhigh delay,\nlow clearance cost',
                xy=(0.1, 0.8), fontsize=10, style='italic',
                bbox=dict(boxstyle='round', facecolor='lightyellow', alpha=0.5))

    plt.tight_layout()
    plt.savefig(OUTPUT_DIR / "fig5_tradeoff.png")
    plt.close()
    print("  Fig 5: Policy trade-off saved")


# ─── Create PPTX with all figures ────────────────────────────────────────────

def create_pptx():
    """Create PPTX with one figure per slide."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    figures = [
        ("fig1_concept.png", "Figure 1",
         "Conceptual comparison of three information regimes for minefields"),
        ("fig2_breach_delay.png", "Figure 2",
         "Attacker breach delay as a function of dummy ratio r"),
        ("fig3_postconflict.png", "Figure 3",
         "Post-conflict clearance: time savings and residual risk elimination"),
        ("fig4_entropy.png", "Figure 4",
         "Positional entropy (attacker uncertainty) under transparent map regime"),
        ("fig5_tradeoff.png", "Figure 5",
         "Policy trade-off between military utility and humanitarian cost"),
    ]

    for fname, title, caption in figures:
        fpath = OUTPUT_DIR / fname
        if not fpath.exists():
            continue

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2),
                                         Inches(12), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Image
        slide.shapes.add_picture(str(fpath), Inches(0.5), Inches(1.0),
                                 width=Inches(12.3))

        # Caption
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(6.5),
                                          Inches(12), Inches(0.8))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = caption
        p2.font.size = Pt(14)
        p2.font.italic = True
        p2.alignment = PP_ALIGN.CENTER

    prs.save(str(OUTPUT_DIR / "figures.pptx"))
    print("  PPTX saved")


# ─── Main ────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    print("Generating figures...")
    fig1_concept()
    fig2_delay()
    fig3_clearance()
    fig4_entropy()
    fig5_tradeoff()
    create_pptx()
    print("Done! All figures saved to:", OUTPUT_DIR)
