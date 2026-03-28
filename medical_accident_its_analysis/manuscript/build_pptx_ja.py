#!/usr/bin/env python3
"""Build PPTX with one figure/table per slide (Japanese version).
Code-generated plots are embedded as images; tables are editable PowerPoint tables."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

OUTPUT_DIR = "/home/ubuntu/medical_analysis/output"
MANUSCRIPT_DIR = "/home/ubuntu/medical_analysis/manuscript"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

TITLE_FONT = "Times New Roman"
BODY_FONT = "Times New Roman"
DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0x52, 0x8A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
HEADER_BG = RGBColor(0x00, 0x52, 0x8A)
ALT_ROW = RGBColor(0xE8, 0xF0, 0xF8)


def add_image_slide(image_path, title_text, caption_text):
    """Add a slide with a code-generated figure (embedded as image)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Title bar
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.name = TITLE_FONT
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.LEFT

    # Image
    if os.path.exists(image_path):
        from PIL import Image
        img = Image.open(image_path)
        img_w, img_h = img.size
        aspect = img_w / img_h

        max_w = Inches(11.5)
        max_h = Inches(5.5)

        if aspect > (max_w / max_h):
            w = max_w
            h = int(w / aspect)
        else:
            h = max_h
            w = int(h * aspect)

        left = int((prs.slide_width - w) / 2)
        top = Inches(1.2)
        slide.shapes.add_picture(image_path, left, top, w, h)
    else:
        txBox2 = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f"[図が見つかりません: {image_path}]"
        p2.font.size = Pt(14)
        p2.font.italic = True

    # Caption
    txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.6))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = caption_text
    p3.font.size = Pt(11)
    p3.font.italic = True
    p3.font.name = BODY_FONT
    p3.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    p3.alignment = PP_ALIGN.LEFT


def add_table_slide(title_text, headers, rows, caption_text=None):
    """Add a slide with an editable PowerPoint table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.name = TITLE_FONT
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.LEFT

    # Table dimensions
    n_rows = 1 + len(rows)
    n_cols = len(headers)
    tbl_width = Inches(12)
    tbl_height = Inches(0.4) * n_rows
    left = Inches(0.667)
    top = Inches(1.3)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, tbl_width, tbl_height)
    table = table_shape.table

    # Column widths
    first_col_w = int(tbl_width * 0.22)
    other_col_w = int((tbl_width - first_col_w) / max(n_cols - 1, 1))
    table.columns[0].width = first_col_w
    for i in range(1, n_cols):
        table.columns[i].width = other_col_w

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.name = BODY_FONT
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_BG

    # Data rows
    for r_idx, row_data in enumerate(rows):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.text = str(val)
            p.font.size = Pt(11)
            p.font.name = BODY_FONT
            p.font.color.rgb = DARK
            p.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ALT_ROW

    # Caption
    if caption_text:
        y_cap = top + tbl_height + Inches(0.3)
        txBox2 = slide.shapes.add_textbox(Inches(0.5), y_cap, Inches(12.3), Inches(0.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = caption_text
        p2.font.size = Pt(11)
        p2.font.italic = True
        p2.font.name = BODY_FONT
        p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)


# ============================================================
# 表紙
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.3), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = (
    "医療安全事故が日本の12診療科における医師数および"
    "医療施設数に与える影響：分割時系列分析"
)
p.font.size = Pt(28)
p.font.bold = True
p.font.name = TITLE_FONT
p.font.color.rgb = DARK
p.alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(9.3), Inches(1))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "BMJ投稿用 図表一覧"
p2.font.size = Pt(18)
p2.font.name = BODY_FONT
p2.font.color.rgb = ACCENT
p2.alignment = PP_ALIGN.CENTER

p3 = tf2.add_paragraph()
p3.text = "[著者名：記入予定]"
p3.font.size = Pt(14)
p3.font.name = BODY_FONT
p3.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
p3.alignment = PP_ALIGN.CENTER

# ============================================================
# 図（コード出力、画像として埋め込み）
# ============================================================

add_image_slide(
    os.path.join(OUTPUT_DIR, "accident_trends.png"),
    "図1. 診療科別・定義別の医療安全事故件数の推移",
    "パネルA：医療事故調査制度（日本医療安全調査機構）による義務的報告（2015〜2025年）。"
    "パネルB：医事関係訴訟統計（最高裁判所、2004〜2023年）。事故件数上位5診療科をハイライト表示。"
)

add_image_slide(
    os.path.join(OUTPUT_DIR, "its_physicians_def1.png"),
    "図2. 分割時系列：医師数 vs 医療事故調査制度報告件数",
    "12主要診療科の分割回帰適合線。垂直破線は介入年（各診療科の事故報告ピーク年）を示す。"
)

add_image_slide(
    os.path.join(OUTPUT_DIR, "its_physicians_def2.png"),
    "図3. 分割時系列：医師数 vs 訴訟統計",
    "12主要診療科の分割回帰適合線（定義2：最高裁判所訴訟統計使用）。垂直破線は介入年を示す。"
)

add_image_slide(
    os.path.join(OUTPUT_DIR, "its_facilities_def1.png"),
    "図4. 分割時系列：施設数 vs 医療事故調査制度報告件数",
    "12主要診療科の医療施設数に対する分割回帰適合線。垂直破線は介入年を示す。"
)

add_image_slide(
    os.path.join(OUTPUT_DIR, "lead_time_heatmap.png"),
    "図5. 相互相関分析による推定リードタイム（年）のヒートマップ",
    "行：診療科、列：事故定義×アウトカムの組み合わせ。濃い色はより強い逆相関を示す。"
    "正のラグ値は事故件数の変化がアウトカムの変化に先行することを示す。"
)

add_image_slide(
    os.path.join(OUTPUT_DIR, "forecast_physicians.png"),
    "図6. 診療科別の予測医師数（2025〜2034年）",
    "直近10年間のデータに基づく線形トレンド外挿。95%予測区間付き。"
)

add_image_slide(
    os.path.join(OUTPUT_DIR, "forecast_facilities.png"),
    "図7. 診療科別の予測施設数（2025〜2034年）",
    "直近10年間のデータに基づく線形トレンド外挿。95%予測区間付き。"
)

add_image_slide(
    os.path.join(OUTPUT_DIR, "forecast_trainees.png"),
    "図8. 診療科別の予測専攻医登録数（2025〜2034年）",
    "利用可能データ（2018〜2025年）に基づく線形トレンド外挿。95%予測区間付き。"
)

# ============================================================
# 表（編集可能なPowerPoint表）
# ============================================================

add_table_slide(
    "表1. 線形トレンド除去後に強い逆相互相関（|r| > 0.9）を示した組み合わせ",
    ["診療科", "事故定義", "アウトカム", "ラグ（年）", "相関係数(r)"],
    [
        ["産婦人科", "医療事故調査制度", "医師数", "+3", "\u22120.946"],
        ["産婦人科", "複合指標", "医師数", "+3", "\u22120.947"],
        ["泌尿器科", "医療事故調査制度", "医師数", "+6", "\u22120.945"],
        ["麻酔科", "複合指標", "医師数", "\u22124", "\u22120.959"],
        ["耳鼻咽喉科", "訴訟統計", "医師数", "+6", "\u22120.973"],
        ["皮膚科", "訴訟統計", "医師数", "+7", "\u22120.988"],
        ["整形外科", "医療事故調査制度（専攻医）", "専攻医数", "+4", "\u22120.909"],
        ["耳鼻咽喉科", "医療事故調査制度（専攻医）", "専攻医数", "\u22123", "\u22120.941"],
        ["精神科", "訴訟統計", "施設数", "+8", "\u22120.981"],
        ["皮膚科", "訴訟統計", "施設数", "+8", "\u22120.969"],
    ],
    "トレンド除去後の事故系列とアウトカム系列のラグ\u22128〜+8年における相互相関。"
    "正のラグ＝事故件数の変化がアウトカムの変化に先行。"
)

add_table_slide(
    "表2. ITS分割回帰における統計的に有意（P<0.05）な事故効果係数",
    ["診療科", "定義", "アウトカム", "R\u00b2", "事故効果係数", "P値"],
    [
        ["外科", "訴訟統計", "医師数", "0.994", "+19.3/件", "<0.001"],
        ["産婦人科", "医療事故調査制度", "医師数", "0.995", "+17.7/報告", "0.021"],
        ["産婦人科", "訴訟統計", "医師数", "0.977", "+11.6/件", "<0.001"],
        ["産婦人科", "医療事故調査制度", "施設数", "0.999", "+4.3/報告", "0.042"],
        ["産婦人科", "訴訟統計", "施設数", "0.990", "+5.5/件", "<0.001"],
        ["眼科", "医療事故調査制度", "施設数", "1.000", "+4.1/報告", "0.039"],
        ["形成外科", "訴訟統計", "施設数", "0.999", "+4.1/件", "0.005"],
    ],
    "注：正の係数は事故件数とアウトカムの共有された長期トレンドを反映。考察の解釈を参照。"
)

add_table_slide(
    "表3. アウトカム別のリードタイムおよびウィンドウ期間推定の要約",
    ["アウトカム", "平均リードタイム（年）", "中央値リードタイム（年）",
     "平均ウィンドウ期間（年）", "中央値ウィンドウ期間（年）"],
    [
        ["医師数", "1.3", "2.5", "4.1", "4.0"],
        ["施設数", "\u22120.1", "0.0", "4.2", "5.0"],
    ],
    "リードタイム：事故変化からアウトカム変化までの遅延時間。"
    "ウィンドウ期間：AICに基づくモデル選択による効果持続期間の推定。"
)

add_table_slide(
    "表4. 診療科別の予測される医療人材変化（2025〜2034年）",
    ["診療科", "医師数（年間変化）", "施設数（年間変化）", "専攻医数（年間変化）"],
    [
        ["外科", "\u2212260", "\u2212350", "+3"],
        ["産婦人科", "+111", "\u221265", "+6"],
        ["内科", "+378", "\u2212256", "+30"],
        ["小児科", "+206", "\u2212188", "\u22125"],
        ["整形外科", "+190", "\u221234", "+30"],
        ["麻酔科", "+186", "+72", "0"],
        ["耳鼻咽喉科", "\u221216", "\u221297", "\u22125"],
        ["救急科", "\u2014", "\u2014", "+33"],
        ["総合診療", "\u2014", "\u2014", "+20"],
    ],
    "直近10年間の利用可能データに基づく線形トレンド外挿。"
)

# ============================================================
# 保存
# ============================================================
output_path = os.path.join(MANUSCRIPT_DIR, "bmj_figures_ja.pptx")
prs.save(output_path)
print(f"Saved to {output_path}")
