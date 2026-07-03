#!/usr/bin/env python3
"""Build English PPTX with one figure per slide for editorial use."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_DIR = os.path.dirname(BASE_DIR)
OUTPUT_DIR = os.path.join(PROJECT_DIR, 'output')

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

FIGURES = [
    ('fig4_accident_trends.png',
     'Figure 1. Medical Safety Incident Trends by Specialty',
     'JMSR mandatory reports (2015\u20132025) and litigation statistics (2004\u20132023) '
     'for the 12 core specialties.'),
    ('fig1_granger_heatmap.png',
     'Figure 2 (Supplementary). Granger Causality Heatmap',
     'Forward Granger causality p-values (-log10 scale) for incidents \u2192 '
     'workforce. Rows: specialties; columns: incident definitions.'),
    ('fig2_irf_physicians.png',
     'Figure 2. Impulse Response Functions: Litigation \u2192 Physicians',
     'Response of physician counts to a one-unit shock in litigation cases '
     'for six key specialties. Shaded areas: approximate 95% CI.'),
    ('fig3_var_forecasts.png',
     'Figure 3. VAR-Based Forecasts',
     'Physician and facility count forecasts (2024\u20132033) for general surgery, '
     'obstetrics & gynaecology, and internal medicine.'),
    ('sfig1_irf_facilities.png',
     'Supplementary Figure S1. IRF: Litigation \u2192 Facilities',
     'Response of facility counts to a one-unit shock in litigation cases.'),
]

for fname, title, caption in FIGURES:
    img_path = os.path.join(OUTPUT_DIR, fname)
    if not os.path.exists(img_path):
        print(f'  Skipping {fname} (not found)')
        continue
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # Image
    img_w = Inches(10)
    img_h = Inches(5.5)
    left = (prs.slide_width - img_w) // 2
    top = Inches(1.0)
    slide.shapes.add_picture(img_path, left, top, width=img_w)

    # Caption
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.7))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = caption
    p2.font.size = Pt(12)
    p2.font.italic = True
    p2.alignment = PP_ALIGN.LEFT

    print(f'  Added: {title}')

output_path = os.path.join(BASE_DIR, 'var_granger_figures_en.pptx')
prs.save(output_path)
print(f'Saved to {output_path}')
