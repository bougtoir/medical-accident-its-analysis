#!/usr/bin/env python3
"""Build a Healthcare Analytics (Elsevier) submission package for the rate-based
analysis of litigation risk and specialty physician workforce allocation (national data from Japan).

Outputs (all derived from results/reanalysis_results.json and data_primary/):
  - manuscript/ha_manuscript_en.docx   anonymised main manuscript
  - manuscript/ha_title_page.docx        title page (separate, with author info)
  - manuscript/ha_cover_letter.docx      cover letter addressed to Healthcare Analytics
  - manuscript/ha_highlights.docx      3-5 highlights (<=85 chars each)
  - manuscript/ha_supplementary.docx     supplementary figures & tables
  - output/ha_Figure_1.png .. Figure_3.png            main figure files
  - output/ha_Supplementary_Figure_1.png .. 3.png      supplementary figure files
  - manuscript/ha_figures.pptx           editable main figure slides
  - manuscript/ha_supplementary_figures.pptx editable supplementary figure slides

Main manuscript is double-anonymisation compliant: no author identifiers,
affiliations or acknowledgements in the body. Figures/tables in the main
manuscript are limited to 5 (3 figures + 2 tables); remaining figures/tables
are placed in the supplementary file.
"""
import os
import json
import re
import zipfile
import pandas as pd
from lxml import etree
from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt

BASE = os.path.dirname(os.path.abspath(__file__))
PROJ = os.path.dirname(BASE)
DP = os.path.join(PROJ, "data_primary")
OUT = os.path.join(PROJ, "output")
RES = json.load(open(os.path.join(PROJ, "results", "reanalysis_results.json")))

try:
    from latex2mathml.converter import convert as _latex_to_mathml
except ImportError:  # pragma: no cover
    _latex_to_mathml = None

_XSL_PATH = os.path.join(BASE, "MML2OMML.XSL")
if not os.path.exists(_XSL_PATH):
    _XSL_PATH = "/tmp/MML2OMML.XSL"
if os.path.exists(_XSL_PATH):
    _XSLT = etree.XSLT(etree.parse(_XSL_PATH))
else:
    _XSLT = None

CORE = ["内科", "外科", "整形外科", "形成外科", "産婦人科", "小児科", "精神科",
        "眼科", "耳鼻咽喉科", "泌尿器科", "皮膚科", "麻酔科"]
EN = {"内科": "Internal medicine", "外科": "Surgery", "整形外科": "Orthopaedics",
      "形成外科": "Plastic surgery", "産婦人科": "Obstetrics & gynaecology",
      "小児科": "Paediatrics", "精神科": "Psychiatry", "眼科": "Ophthalmology",
      "耳鼻咽喉科": "Otolaryngology", "泌尿器科": "Urology", "皮膚科": "Dermatology",
      "麻酔科": "Anaesthesiology"}


def load(name):
    df = pd.read_csv(os.path.join(DP, name)).set_index("specialty")
    df.columns = [int(c) for c in df.columns]
    return df.loc[CORE]


def prim(label_key):
    for r in RES["primary"]:
        if label_key in r["label"]:
            return r
    raise KeyError(label_key)


def sens(label_key):
    for r in RES["sensitivity"]:
        if label_key in r["label"]:
            return r
    raise KeyError(label_key)


def het(label_key):
    for r in RES.get("heterogeneity", []):
        if label_key in r["label"]:
            return r
    raise KeyError(label_key)


def trend(label_key):
    for r in RES.get("trend_sensitivity", []):
        if label_key in r["label"]:
            return r
    raise KeyError(label_key)


PHYS = prim("physician growth ~ lagged litigation rate")
HOSP = prim("hospital growth ~ lagged litigation rate")
REV = prim("Reverse")
CNT = sens("COUNTS contrast")
ANN = sens("Annual hospital")
INT = sens("interpolated-annual")
EQP, EQH = RES["equivalence"][0], RES["equivalence"][1]
# Bootstrap results for primary physician and hospital models (block resampling of specialties)
BS_PHYS = next((b for b in RES["bootstrap"] if b["outcome"] == "dlog_phys"), None)
BS_HOSP = next((b for b in RES["bootstrap"] if b["outcome"] == "dlog_hosp"), None)
SIM = RES.get("policy_simulation", {})
SURG_SIM = next((r for r in SIM.get("specialties", []) if r["specialty"] == "Surgery"), {}) if SIM else {}
TOTAL_SIM = SIM.get("totals", {}) if SIM else {}
SP = RES["descriptive"]["spearman_litrate_vs_physgrowth"]
n_pos = sum(1 for v in SP.values() if v["rho"] > 0)
n_sig = sum(1 for v in SP.values() if v["p"] < 0.05)
BIEN = RES["grid"]["biennial_years"]
YEARS = f"{BIEN[0]}\u2013{BIEN[-1]}"
N_SPECIALTIES = len(CORE)
PER = 1000  # rate scale (cases per 1,000 physicians)

def _simulation_table_rows():
    """Return rows for the counterfactual 2034 physician-count table."""
    rows = []
    for r in SIM.get("specialties", []):
        rows.append([
            EN.get(r["specialty"], r["specialty"]),
            f"{r['phys_2024']:,}",
            f"{r['projected_baseline']:,.0f}",
            f"{r['projected_litigation_zero_point']:,.0f}",
            f"{r['projected_litigation_zero_lower']:,.0f}",
            f"{r['projected_mde_lever']:,.0f}",
        ])
    t = SIM.get("totals", {})
    rows.append([
        "All specialties",
        f"{t.get('phys_2024', 0):,}",
        f"{t.get('base_2034', 0):,}",
        f"{t.get('lit_zero_point_2034', 0):,}",
        f"{t.get('lit_zero_lower_2034', 0):,}",
        f"{t.get('mde_2034', 0):,}",
    ])
    return rows

MARGIN1 = int(RES["equivalence"][0]["tests"][0]["margin"] * 100)
MARGIN2 = int(RES["equivalence"][1]["tests"][1]["margin"] * 100)

# JMSR/MAIS sensitivity (report counts added as a control)
JMSR = sens("JMSR")
JMSR_CORR = RES["jmsr_correlation"]
JMSR_START = JMSR_CORR["years"][0] + 1  # outcome years start one year after first JMSR lag

# Nikkei Telecom media coverage sensitivity (national annual article counts)
MEDIA = sens("Media-adjusted")
MEDIA_CORR = RES["media_correlation"]
MEDIA_START = MEDIA_CORR["years"][0] + 1   # outcome years start one year after first media lag
MEDIA_END = MEDIA_CORR["years"][-1]        # last outcome year of the media sensitivity

# Public repository and reproducibility metadata
PUBLIC_REPO = "https://github.com/bougtoir/medical-accident-its-analysis"

# Title used throughout the submission (global framing; country is stated as data provenance)
MANUSCRIPT_TITLE = (
    "Litigation risk and specialty-level physician workforce allocation: "
    "a transparent decision-analytics framework for evaluating healthcare workforce policy levers"
)

# Holm-adjusted p-value for the exploratory JOCS-CP indicator in the hospital model
JOCS_HOLM = next((
    t["holm_p"] for t in RES["multiple_comparison"]["tests"]
    if "hospital" in t["label"].lower() and "JOCS-CP" in t["label"]
), None)

# Load primary dataframes for year ranges/resolution (used in supplementary table and limitations)
PHYS_DF = load("physicians_by_specialty.csv")
HOSP_DF = load("facilities_hospital_by_specialty.csv")
CLINIC_DF = load("facilities_clinic_by_specialty.csv")
LIT_DF = load("litigation_by_specialty.csv")

# Specialist-trainee coverage context (optional reference table)
def _load_senkoi():
    path = os.path.join(DP, "senkoi_coverage.csv")
    if not os.path.exists(path):
        return None
    df = pd.read_csv(path)
    df = df.set_index("specialty_ja")
    df = df.reindex(CORE)
    return df.to_dict(orient="index")

SENKOI = _load_senkoi()

def _senkoi_table_rows():
    if SENKOI is None:
        return []
    rows = []
    for s in CORE:
        v = SENKOI[s]
        rows.append([
            v["specialty_en"],
            f"{v['senkoi_2018']:,}",
            f"{v['physicians_3_5_yr_2014']:,}",
            f"{v['yr3_2012_cohort']:,}",
            f"{v['yr5_2012_cohort']:,}",
            f"{v['coverage_pct']:.1f}%",
        ])
    return rows


def _resolution(df):
    diffs = [df.columns[i + 1] - df.columns[i] for i in range(len(df.columns) - 1)]
    return max(set(diffs), key=diffs.count)


def _year_label(years):
    return f"{years[0]}\u2013{years[-1]}"


PHYS_RES = _resolution(PHYS_DF)
HOSP_RES = _resolution(HOSP_DF)
CLINIC_RES = _resolution(CLINIC_DF)

# Descriptive counts used in abstract and results (computed once, not hard-coded)
DESCR = RES["descriptive"]["biennial_first_last"]["by_specialty"]
GREW = sum(1 for v in DESCR.values() if v["phys_last"] > v["phys_first"])
FELL = sum(1 for v in DESCR.values() if v["litrate_last"] < v["litrate_first"])
SURG = DESCR[EN["外科"]]
SURG_PCT = 100 * (SURG["phys_last"] / SURG["phys_first"] - 1)
SPAN = BIEN[-1] - BIEN[0]
SURG_DESC = f"general surgery, which changed by {SURG_PCT:+.1f}%"

# Useful helpers
_per = f"{PER:,}"
def _fmt_pct(x):
    return f"{x:+.2f}"


REFS = {
    "maldist": "Matsumoto M, Inoue K, Bowman R, Kajii E. Self-employment, specialty choice, and geographical distribution of physicians in Japan: a comparison with the United States. Health Policy. 2010;96(3):239-244. doi:10.1016/j.healthpol.2010.02.008.",
    "malprac": "Studdert DM, Mello MM, Brennan TA. Medical malpractice. N Engl J Med. 2004;350(3):283-292.",
    "defmed": "Hiyama T, Yoshihara M, Tanaka S, et al. Defensive medicine practices among gastroenterologists in Japan. World J Gastroenterol. 2006;12(47):7671-7675. doi:10.3748/wjg.v12.i47.7671.",
    "lakens": "Lakens D. Equivalence tests: a practical primer for t tests, correlations, and meta-analyses. Soc Psychol Personal Sci. 2017;8(4):355-362.",
    "schuir": "Schuirmann DJ. A comparison of the two one-sided tests procedure and the power approach for assessing the equivalence of average bioavailability. J Pharmacokinet Biopharm. 1987;15(6):657-680.",
    "jocscp": "Japan Council for Quality Health Care. Japan Obstetric Compensation System for Cerebral Palsy. Tokyo: Japan Council for Quality Health Care; 2009.",
    "phys": "Ministry of Health, Labour and Welfare. Statistics of Physicians, Dentists and Pharmacists. Tokyo: MHLW; 2024. Available from: https://www.mhlw.go.jp/toukei/list/33-20b.html (accessed 11 August 2026).",
    "court": "Supreme Court of Japan, Committee on Medical Litigation. Statistics on medical malpractice litigation (closed cases by specialty). Tokyo: Supreme Court of Japan; 2024. Available from: https://www.courts.go.jp/ (accessed 11 August 2026).",
    "mhlw_senkoi2018": "Ministry of Health, Labour and Welfare. Adoption status of specialist trainees under the new specialist training system (toward the FY2018 launch). Tokyo: MHLW; 2018. Available from: https://www.mhlw.go.jp/content/10803000/000452411.pdf (accessed 15 August 2026).",
    "mhlw_3_5yr": "Ministry of Health, Labour and Welfare. Number of physicians 3-5 years after medical registration, by principal specialty. Tokyo: MHLW; 2015. Available from: https://www.mhlw.go.jp/file/06-Seisakujouhou-10800000-Iseikyoku/323.pdf (accessed 15 August 2026).",
    "facil": "Ministry of Health, Labour and Welfare. Survey of Medical Institutions (Dynamic). Tokyo: MHLW; 2024. Available from: https://www.mhlw.go.jp/toukei/list/79-1a.html (accessed 11 August 2026).",
    "mais": "Act on the Promotion of Medical Safety; Medical Accident Investigation System (2015). Tokyo: MHLW; 2015.",
    "jmsr_data": "Japan Medical Safety Research Organisation. Annual reports of medical accident investigations (2015-2024). Tokyo: JMSR; 2025.",
    "nikkei": "Nikkei Inc. Nikkei Telecom 21. Tokyo: Nikkei Inc. Accessed 2024. Available from: https://telecom.nikkei.co.jp/",
    "angrist": "Angrist JD, Pischke JS. Mostly Harmless Econometrics. Princeton: Princeton University Press; 2009.",
    "cameron2015": "Cameron AC, Miller DL. A practitioner's guide to cluster-robust inference. J Hum Resour. 2015;50(2):317-372.",
    "strobe": "von Elm E, Altman DG, Egger M, et al. The STROBE statement. Lancet. 2007;370(9596):1453-1457.",
    "matsa2007": "Matsa DA. Does malpractice liability keep the doctor away? Evidence from tort reform damage caps. J Legal Stud. 2007;36(S2):S143-S182.",
    "hyman2015": "Hyman DA, Silver C, Black B, Paik M. Does tort reform affect physician supply? Evidence from Texas. Int Rev Law Econ. 2015;42:203-218.",
    "frakes2020": "Frakes MD, Frank MB, Seabury SA. The effect of malpractice law on physician supply: Evidence from negligence-standard reforms. J Health Econ. 2020;70:1-16.",
    "kessler1996": "Kessler DP, McClellan MB. Do doctors practice defensive medicine?. Q J Econ. 1996;111(2):353-390.",
    "sloan2008": "Sloan FA, Shadle JH. Is there empirical evidence for \"Defensive Medicine\"? A reassessment. J Health Econ. 2008;27(2):481-491.",
    "taniguchi2023": "Taniguchi K, Watari T, Nagoshi K. Characteristics and trends of medical malpractice claims in Japan between 2006 and 2021. PLoS One. 2023;18(12):e0296155.",
    "hasegawa2016": "Hasegawa J, Toyokawa S, Ikenoue T, et al. Relevant obstetric factors for cerebral palsy: from the Nationwide Obstetric Compensation System in Japan. PLoS One. 2016;11(1):e0148122.",
    "morita2018": "Morita H. Criminal prosecution and physician supply. Int Rev Law Econ. 2018;55:1-11.",
    "helland2015": "Helland E, Seabury SA. Tort reform and physician labor supply: a review of the evidence. Int Rev Law Econ. 2015;42:192-202.",
    "bismark2006": "Bismark M, Paterson R. No-fault compensation in New Zealand: harmonizing injury compensation, provider accountability, and patient safety. Health Aff. 2006;25(1):278-286.",
    "mello2011": "Mello MM, Kachalia A, Studdert DM. Administrative compensation for medical injuries: lessons from three foreign systems. New York: The Commonwealth Fund; 2011.",
    "kamijo2025": "Kamijo K, Wada Y, Ishida K, Warsof SL, Saade G, Kawakita T. Medical-legal claims in obstetrics and gynecology: Japan versus the United States. J Healthc Risk Manag. 2025;44(4):5-11. doi:10.1002/jhrm.70001.",
    "mcnamara2025": "McNamara C, Pineda-Torres M. Medical residency subsidies and physician shortages. J Public Econ. 2025;251:105494. doi:10.1016/j.jpubeco.2025.105494.",
    "lin2022": "Lin PL, Huang JP, Fujii T, Cho EH, Huang MC. A survey of specialty choice among obstetrics and gynecology residents in Japan, Korea, and Taiwan. J Obstet Gynaecol Res. 2022;48(7):1968-1977.",
    "tversky1973": "Tversky A, Kahneman D. Availability: a heuristic for judging frequency and probability. Cogn Psychol. 1973;5(2):207-232.",
    "kahneman1979": "Kahneman D, Tversky A. Prospect theory: an analysis of decision under risk. Econometrica. 1979;47(2):263-291.",
    "samuelson1988": "Samuelson W, Zeckhauser R. Status quo bias in decision making. J Risk Uncertain. 1988;1(1):7-59.",
}
_CITE_ORDER = []
BODY_TEXTS = []


def wc(text):
    return len(re.findall(r"\b[\w'-]+\b", text))


def add_math(doc, latex, inline=False, para=None):
    """Insert an OMML equation from LaTeX source.

    Falls back to italicised plain text if latex2mathml or the XSLT is not
    available, so the manuscript remains usable in any environment.
    """
    if _latex_to_mathml is None or _XSLT is None:
        if para is None:
            para = doc.add_paragraph()
        run = para.add_run(latex)
        run.italic = True
        run.font.name = "Cambria Math" if inline else "Times New Roman"
        return para

    mathml = _latex_to_mathml(latex)
    mathml_tree = etree.fromstring(mathml)
    omml_tree = _XSLT(mathml_tree)
    omml_root = omml_tree.getroot()

    if para is None:
        para = doc.add_paragraph()
        para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    para._element.append(omml_root)
    return para


def fmt(x, d=3):
    return f"{x:+.{d}f}" if isinstance(x, float) else str(x)


def p_tost_fmt(p, digits=3):
    """Format a TOST p-value; report '<0.001' when rounded to zero."""
    if p is None:
        return ""
    fmt_str = f"{{p:.{digits}f}}"
    s = fmt_str.format(p=p)
    return "<0.001" if s == "0." + "0" * digits else s


# Map common non-ASCII punctuation to ASCII equivalents. This keeps the
# manuscript free of double-byte or multi-byte textual characters while
# leaving Word OMML equations (which live under m:oMath, not w:r) untouched.
_ASCII_MAP = {
    "\u2013": "-",   # en dash
    "\u2014": "--", # em dash
    "\u00b1": "+/-", # plus-minus sign
    "\u2192": "->",  # rightwards arrow
    "\u2022": "*",   # bullet
    "\u2212": "-",   # minus sign
    "\u00d7": "x",   # multiplication sign
    "\u2018": "'",   # left single quotation mark
    "\u2019": "'",   # right single quotation mark
    "\u201c": '"',   # left double quotation mark
    "\u201d": '"',   # right double quotation mark
    "\u2026": "...", # horizontal ellipsis
}


def _is_in_math(run):
    """Return True if a run lives inside a Word OMML equation."""
    el = run._element
    while el is not None:
        if el.tag.endswith("oMath") or el.tag.endswith("oMathPara"):
            return True
        el = el.getparent()
    return False


def _normalize_text(text: str) -> str:
    for src, dst in _ASCII_MAP.items():
        text = text.replace(src, dst)
    return text


def normalize_docx(doc):
    """Replace non-ASCII punctuation in all paragraphs and tables of a docx."""
    for p in doc.paragraphs:
        for run in p.runs:
            if _is_in_math(run):
                continue
            if run.text:
                run.text = _normalize_text(run.text)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    for run in p.runs:
                        if _is_in_math(run):
                            continue
                        if run.text:
                            run.text = _normalize_text(run.text)


def normalize_pptx(prs):
    """Replace non-ASCII punctuation in all text frames of a pptx."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                tf = shape.text_frame
                for p in tf.paragraphs:
                    for run in p.runs:
                        if run.text:
                            run.text = _normalize_text(run.text)


_ZIP_ASCII_REPLACEMENTS = {
    # Greek letters and common math symbols are intentionally NOT replaced here.
    # sanitize_zip encodes any remaining non-ASCII characters as numeric XML
    # character references (&#...;), so equations stay visually correct while
    # the file itself contains only ASCII bytes.
    # East-Asian theme fonts -> ASCII font names
    "ＭＳ 明朝": "MS Mincho",
    "ＭＳ ゴシック": "MS Gothic",
    "ＭＳ Ｐ明朝": "MS PMincho",
    "ＭＳ Ｐゴシック": "MS PGothic",
    "맑은 고딕": "Malgun Gothic",
    "宋体": "SimSun",
    "新細明體": "PMingLiU",
    # Bullets and slide-number placeholders
    "\uf0b7": "*",
    "•": "*",
    "–": "-",
    "‹#›": "#",
}


def _sanitize_xml_text(text: str) -> str:
    # Longest first so multi-char font names are replaced before any subcomponent.
    for src, dst in sorted(_ZIP_ASCII_REPLACEMENTS.items(), key=lambda kv: -len(kv[0])):
        text = text.replace(src, dst)
    return text


def sanitize_zip(path: str) -> None:
    """Rewrite a .docx or .pptx so its XML contains only ASCII characters.

    Applies ASCII replacements for punctuation and East-Asian font names, then
    encodes any remaining non-ASCII characters as numeric XML character references
    (&#...;).  This keeps Word OMML equations (Greek letters, math symbols)
    visually intact while ensuring no multibyte bytes remain in the file.
    """
    import tempfile, shutil

    tmp = tempfile.mkstemp(suffix=os.path.splitext(path)[1])[1]
    shutil.move(path, tmp)
    try:
        with zipfile.ZipFile(tmp, "r") as zin, zipfile.ZipFile(
            path, "w", zipfile.ZIP_DEFLATED
        ) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename.endswith(".xml") or item.filename.endswith(".rels"):
                    text = data.decode("utf-8", errors="replace")
                    text = _sanitize_xml_text(text)
                    data = text.encode("ascii", "xmlcharrefreplace")
                zout.writestr(item, data)
    finally:
        os.remove(tmp)


def cite_number(keys):
    nums = []
    for k in keys:
        if k not in _CITE_ORDER:
            _CITE_ORDER.append(k)
        nums.append(_CITE_ORDER.index(k) + 1)
    nums = sorted(set(nums))
    out, i = [], 0
    while i < len(nums):
        j = i
        while j + 1 < len(nums) and nums[j + 1] == nums[j] + 1:
            j += 1
        out.append(str(nums[i]) if i == j else f"{nums[i]}-{nums[j]}")
        i = j + 1
    return ",".join(out)


def _setup_doc():
    doc = Document()
    for s in doc.sections:
        s.top_margin = Cm(2.54)
        s.bottom_margin = Cm(2.54)
        s.left_margin = Cm(2.54)
        s.right_margin = Cm(2.54)
    st = doc.styles["Normal"]
    st.font.name = "Times New Roman"
    st.font.size = Pt(12)
    st.paragraph_format.line_spacing = 2.0
    st.paragraph_format.space_after = Pt(6)
    return doc


def _add_runs(par, text, size=Pt(12), bold=False, italic=False):
    for part in re.split(r"(\{[^}]+\})", text):
        if part.startswith("{") and part.endswith("}"):
            keys = [k.strip() for k in part[1:-1].split(",")]
            r = par.add_run(f"[{cite_number(keys)}]")
            r.font.superscript = True
        elif part:
            r = par.add_run(part)
            r.font.size = size
            r.bold = bold
            r.italic = italic
        if part:
            par.runs[-1].font.name = "Times New Roman"


_head_counters = [0, 0, 0, 0]

def head(doc, text, level=1, numbered=True):
    global _head_counters
    if numbered:
        _head_counters[level - 1] += 1
        for i in range(level, len(_head_counters)):
            _head_counters[i] = 0
        num = ".".join(str(_head_counters[i]) for i in range(level))
        text = f"{num}. {text}"
    h = doc.add_heading(text, level=level)
    for r in h.runs:
        r.font.color.rgb = RGBColor(0, 0, 0)
        r.font.name = "Times New Roman"
    return h


def body(doc, text, **kw):
    p = doc.add_paragraph()
    _add_runs(p, text, **kw)
    p.paragraph_format.line_spacing = 2.0
    p.paragraph_format.space_after = Pt(6)
    BODY_TEXTS.append(p.text)
    return p


def para(doc, text, **kw):
    p = doc.add_paragraph()
    _add_runs(p, text, **kw)
    p.paragraph_format.line_spacing = 2.0
    p.paragraph_format.space_after = Pt(6)
    return p


def field(doc, label, text):
    p = doc.add_paragraph()
    r = p.add_run(label + " ")
    r.bold = True
    r.font.name = "Times New Roman"
    r.font.size = Pt(12)
    _add_runs(p, text)
    p.paragraph_format.line_spacing = 2.0
    p.paragraph_format.space_after = Pt(6)
    return p


def figure(doc, fn, caption, width=Inches(5.8)):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(14)
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    img = os.path.join(OUT, fn)
    if os.path.exists(img):
        p.add_run().add_picture(img, width=width)
    cap = doc.add_paragraph()
    cap.paragraph_format.space_after = Pt(6)
    rc = cap.add_run(caption)
    rc.bold = True
    rc.font.size = Pt(10)
    rc.font.name = "Times New Roman"
    doc.add_paragraph()


def table(doc, headers, rows, caption):
    cap = doc.add_paragraph()
    cap.paragraph_format.space_before = Pt(14)
    cap.paragraph_format.space_after = Pt(6)
    rc = cap.add_run(caption)
    rc.bold = True
    rc.font.size = Pt(10)
    rc.font.name = "Times New Roman"
    t = doc.add_table(rows=1, cols=len(headers))
    t.style = "Table Grid"
    for i, h in enumerate(headers):
        cell = t.rows[0].cells[i]
        p = cell.paragraphs[0]
        r = p.add_run(str(h))
        r.bold = True
        r.font.size = Pt(9)
        r.font.name = "Times New Roman"
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for row in rows:
        cells = t.add_row().cells
        for i, v in enumerate(row):
            p = cells[i].paragraphs[0]
            pr = p.add_run(str(v))
            pr.font.size = Pt(9)
            pr.font.name = "Times New Roman"
            if i > 0:
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph()


def build_manuscript():
    global _head_counters
    _head_counters = [0, 0, 0, 0]
    doc = _setup_doc()

    # Title
    t = doc.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    t.paragraph_format.space_after = Pt(18)
    rt = t.add_run(MANUSCRIPT_TITLE)
    rt.bold = True
    rt.font.size = Pt(14)
    rt.font.name = "Times New Roman"

    # Abstract (unstructured, <=250 words)
    abstract_text = (
        "Specialty maldistribution is a global healthcare workforce-allocation problem, "
        "and many health systems assume that malpractice-litigation risk pushes physicians "
        "away from high-risk specialties. We used this question as a test case for a "
        "transparent, reproducible decision-analytics and operations-research framework that evaluates a proposed "
        "workforce policy lever while avoiding size confounding and interpolation artifacts "
        f"in sparse administrative panels. Using national primary data for {N_SPECIALTIES} specialties "
        f"in Japan ({BIEN[0]}\u2013{BIEN[-1]}), we measured exposure as closed "
        f"malpractice claims per {_per} physicians (rate, not count) and regressed biennial "
        "log-change in physicians and hospitals on the lagged litigation rate in a panel "
        "with specialty and wave fixed effects, cluster-robust standard errors with a "
        "small-cluster t(G-1) correction, a cluster block-bootstrap robustness check, "
        "two one-sided equivalence (TOST) tests, and power diagnostics. The workforce grew "
        f"in {GREW} of {N_SPECIALTIES} specialties; {SURG_DESC}, was the exception. Litigation rate "
        f"was not associated with physician growth (coefficient {fmt(PHYS['coef'],4)}; "
        f"95% CI {fmt(PHYS['ci_low'],4)} to {fmt(PHYS['ci_high'],4)}; p={PHYS['p']:.2f}) or hospital "
        f"growth (p={HOSP['p']:.2f}). A one-SD higher rate changed physician growth by less than "
        f"\u00b1{MARGIN1}% (TOST p={p_tost_fmt(EQP['tests'][0]['p_tost'])}; less than \u00b1{MARGIN2}% "
        f"was p={p_tost_fmt(EQP['tests'][1]['p_tost'])}); the minimum detectable effect was "
        f"{EQP['mde_80pct']:.2f}% per SD at 80% power. Hospital growth was within \u00b1{MARGIN2}% "
        f"(p={p_tost_fmt(EQH['tests'][1]['p_tost'])}) but not the stricter \u00b1{MARGIN1}% margin "
        f"(p={p_tost_fmt(EQH['tests'][0]['p_tost'])}). Sensitivity analyses were unchanged. "
        f"Specialty-level litigation risk is not associated with "
        "workforce decline in these national data. The decision-analytics and operations-research framework is "
        "exportable to other healthcare workforce-policy levers."
    )
    abstract_wc = wc(abstract_text)
    if abstract_wc > 250:
        raise SystemExit(f"Abstract is {abstract_wc} words; must be <=250")

    head(doc, "Abstract", level=1, numbered=False)
    p = doc.add_paragraph()
    r = p.add_run(abstract_text)
    r.font.name = "Times New Roman"
    r.font.size = Pt(12)
    p.paragraph_format.line_spacing = 2.0
    p.paragraph_format.space_after = Pt(6)

    kw = doc.add_paragraph()
    kr = kw.add_run("Keywords: ")
    kr.bold = True
    kr.font.name = "Times New Roman"
    _add_runs(kw, "malpractice litigation; physician workforce; specialty maldistribution; "
                  "equivalence testing; healthcare decision analytics; operations research")
    kw.paragraph_format.space_after = Pt(18)

    # Introduction
    head(doc, "Introduction", level=1)
    body(doc,
         "Specialty maldistribution is a global healthcare workforce allocation challenge: "
         "high-acuity fields such as surgery, obstetrics and gynaecology, paediatrics and emergency care "
         "are widely perceived as understaffed across many health systems.{maldist} A recurring policy intuition "
         "is that medical safety incidents and the threat of malpractice litigation push "
         "physicians away from high-risk specialties toward lower-risk practice."
         "{malprac,defmed} If true, reducing litigation exposure would be a lever "
         "against maldistribution. Japan provides a well-documented national setting in which to test this "
         "intuition: it tracks closed malpractice claims, physician counts and hospital facility counts by specialty, "
         "and it faces the same fee-for-service and specialty-training pressures seen in other high-income countries.")
    body(doc,
         "Prior work in this area often related raw annual counts of incidents or "
         "lawsuits to raw counts of physicians or facilities. Two features make such "
         "designs prone to spurious association. "
         "First, counts are not adjusted for specialty size: a larger specialty "
         "mechanically accumulates more procedures, more claims and more physicians, so "
         "count-based associations can arise without any behavioural mechanism. Second, the physician census "
         "is collected only biennially; interpolating it to an annual series and "
         "analysing it as if each year were an independent observation inflates the "
         "degrees of freedom of any lag-based method. These pitfalls are not unique to "
         "malpractice research; they arise whenever administrative counts are used to "
         "infer behavioural responses in healthcare organisations. From a decision-analytics "
         "perspective, the relevant question is not whether litigation risk is correlated "
         "with raw counts, but whether it is a policy lever that workforce planners can use "
         "to shift specialty supply. We therefore treat the litigation-workforce question as "
         "a test case for a transparent, reproducible decision-analytics and operations-research framework for "
         "healthcare workforce allocation and resource planning. "
         "This is, to our knowledge, the first application of pre-specified equivalence testing "
         "and small-cluster inference diagnostics to the malpractice-litigation workforce literature.")
    body(doc,
         "We therefore examined the question using rates rather than counts, using "
         "only measured biennial physician observations, and using equivalence testing\u2014"
         "which can provide positive evidence for the absence of a meaningful effect "
         "rather than merely failing to reject a null.{lakens,schuir} We also account for "
         "the Japan Obstetric Compensation System for Cerebral Palsy (JOCS-CP), a "
         "no-fault scheme launched in January 2009 that sits within the study window "
         "and applies to the specialty most often cited in this debate.{jocscp}")

    # Materials and methods
    head(doc, "Materials and methods", level=1)
    head(doc, "Data sources", level=2)
    body(doc,
         "We report this observational study following the Strengthening the Reporting of "
         "Observational Studies in Epidemiology (STROBE) guidance.{strobe} We "
         f"studied {N_SPECIALTIES} core clinical specialties in Japan for which the Supreme Court reports "
         "specialty-specific litigation. Three official primary series drove the main "
         "analysis: physician counts by specialty from the biennial Statistics of "
         "Physicians, Dentists and Pharmacists{phys}; closed malpractice claims by "
         "specialty from the Supreme Court of Japan{court}; and hospital facility counts by "
         "specialty from the annual Survey of Medical Institutions{facil}. Two sensitivity "
         "series were also used: annual medical accident investigation reports by "
         "specialty from the Japan Medical Safety Research Organisation (JMSR, 2015-2024){jmsr_data} "
         f"and total national newspaper article counts from Nikkei Telecom 21 "
         "(2004-2018; the sensitivity analysis uses 2009-2018; keywords: medical error + medical malpractice).{nikkei} "
         f"The full extraction pipeline (with source identifiers and SHA-256 checksums) is "
         f"documented in the accompanying repository ({PUBLIC_REPO}).")
    body(doc,
         "As context for the specialty-training pipeline, first-year specialist-trainee (senkoi) counts "
         "in 2018 and counts of physicians 3\u20135 years after medical registration in 2014 were summarised "
         "by specialty; coverage rates range widely across the 12 primary fields (Supplementary Table 8).{mhlw_senkoi2018,mhlw_3_5yr}")
    body(doc,
         "Physician counts use the principal specialty classification; broad categories were matched to the Supreme Court's "
         "specialty labels, and subspecialties were aggregated in code. Because the "
         "Court assigns multi-specialty cases to a single principal specialty and states "
         "that the counts do not represent the intrinsic risk of each specialty, we "
         "treat litigation as an exposure signal rather than a measure of incident "
         "risk.{court} We distinguish litigation from the Medical Accident Investigation "
         "System, which began in 2015 and covers only deaths and stillbirths judged "
         "unforeseen by the hospital administrator; it is not a general incident-reporting "
         "system and is not used as an exposure here.{mais} Primary data sources and their "
         "resolution are summarised in Supplementary Table 1.")

    head(doc, "Statistical analysis", level=2)
    body(doc,
         "We formalised the evaluation as a sensitivity-analysis framework that varies the "
         "exposure definition (counts versus rates), panel frequency (measured biennial "
         "waves versus interpolated annual values), and potential confounders (JMSR reports, "
         "media coverage, JOCS-CP period) while holding the specialty-level panel "
         "structure constant. "
         f"The exposure was the litigation rate, defined as closed claims per {_per} "
         "physicians in each specialty-year, which removes specialty-size confounding "
         "because large specialties generate more claims for reasons unrelated to "
         f"per-physician risk. The primary analysis used the {len(BIEN)} measured biennial "
         f"physician waves ({BIEN[0]}\u2013{BIEN[-1]}). For each specialty we computed the "
         "biennial log-change in physicians (and, separately, in hospitals) and regressed "
         "it on the litigation rate at the start of the interval, in a panel with specialty "
         "and wave fixed effects and standard errors clustered by specialty.{angrist} "
         f"Clusters are defined by specialty, so G={N_SPECIALTIES} and the small-cluster correction "
         "uses a t-distribution with G-1 degrees of freedom for all cluster-robust inference. "
         "Supplementary Figure 1 summarises the sensitivity-analysis framework. The primary "
         "estimating equation, for specialty s and wave t, was as follows.")
    add_math(doc, r"\Delta \log(Y_{st}) = \alpha_s + \delta_t + \beta \cdot \text{litrate}_{s,t-1} + \epsilon_{st}")
    body(doc,
         "Here, Y is either physicians or hospitals, alpha_s are specialty fixed effects, "
         "delta_t are wave fixed effects, and standard errors are clustered by specialty. "
         "For the equivalence analysis we standardised litrate to a z-score, so beta "
         "gives the expected biennial log-change per one-SD increase in the litigation rate.")
    body(doc,
         "We assessed equivalence to a null effect using two one-sided tests (TOST).{lakens,schuir} "
         "For a pre-specified margin m, the two one-sided null hypotheses are")
    add_math(doc, r"H_0: \beta \leq -m \quad \text{and} \quad H_0: \beta \geq +m")
    body(doc,
         "Equivalence is declared when both one-sided tests yield p < alpha. We used "
         f"margins of {MARGIN1}% and {MARGIN2}% biennial workforce change because they are "
         "smaller than typical policy targets for specialty workforce rebalancing and "
         "represent changes that workforce planners would consider substantively small.")
    body(doc,
         "We complemented the analytical small-cluster inference with two diagnostic "
         "quantities. First, a cluster block-bootstrap with B = 1,999 replications "
         "resampled specialties with replacement, re-fitted the primary model, and "
         "reported percentile bootstrap 95% confidence intervals and a bootstrap p-value "
         "based on the distribution of absolute t-statistics. Second, we report the "
         "minimum detectable effect (MDE) at 80% power for the per-SD litigation-rate "
         "coefficient and the power to declare equivalence when the true effect is zero. "
         "The MDE is (t(0.975, df) + t(0.80, df)) \u00d7 SE; the equivalence power is "
         "2\u00b7F_t(m/SE) \u2212 1, where F_t is the cumulative distribution of the t(df) "
         "distribution. These diagnostics make the limited-panel information explicit to "
         "workforce planners evaluating this policy lever.")
    head(doc, "Policy lever simulation", level=2)
    body(doc,
         "To translate the regression results into a decision-analytics and operations-research "
         "output, we projected physician counts to 2034 under three stylised policy levers. The "
         "baseline used each specialty's observed mean biennial log-growth from 2004 to "
         "2024. The two litigation-reduction levers set the litigation rate to zero and "
         "applied the point estimate and the 95% lower-bound coefficient, respectively, so "
         "they show both the central projection and the most favourable effect consistent "
         "with the data. The third lever added the minimum detectable per-SD effect "
         f"({EQP['mde_80pct']:.2f}% per biennium) to baseline growth as a benchmark for the smallest "
         "policy effect this panel could detect with 80% power. The projections are "
         "deterministic counterfactuals, not forecasts, and are reported as the marginal "
         "percent change in 2034 relative to the baseline drift.")
    body(doc,
         "All reported confidence intervals and two-sided p-values use a "
         f"t-distribution with G-1 = {PHYS['df']} degrees of freedom, the small-cluster correction recommended "
         "by Cameron and Miller.{cameron2015} An indicator for obstetrics and gynaecology from 2009 onward captured the "
         "JOCS-CP period.{jocscp} Sensitivity analyses repeated the models on (i) the "
         "annual hospital series, (ii) a linearly interpolated annual physician series "
         "(with degrees of freedom governed by the measured waves, not the interpolated n), "
         "(iii) raw counts instead of rates, (iv) the annual hospital series 2016-2024 "
         f"additionally controlling for the JMSR report rate (reports per {_per} physicians), "
         f"and (v) the annual hospital series {MEDIA_START}\u2013{MEDIA_END} additionally "
         f"controlling for total Nikkei Telecom article counts. Because the article-count series "
         "is a national yearly variable, it is collinear with full wave fixed effects; this "
         "sensitivity therefore uses specialty fixed effects plus a linear time trend rather "
         "than wave dummies. The JOCS-CP indicator and all sensitivity models are exploratory; "
         "we report raw p-values and Holm step-down adjusted p-values for this family. "
         "Analyses used Python (statsmodels); code and data are openly available at "
         f"{PUBLIC_REPO}.")

    # Results
    head(doc, "Results", level=1)
    head(doc, "Workforce and litigation trends", level=2)
    body(doc,
         f"Litigation rates per {_per} physicians varied several-fold across "
         f"specialties and fell over time in {FELL} of {N_SPECIALTIES} fields (Supplementary Figure 2). "
         f"Over the same period the physician workforce grew in {GREW} of {N_SPECIALTIES} specialties "
         "(Supplementary Figure 3; Table 1); the only exception was general surgery, "
         f"which was essentially flat ({SURG_PCT:+.1f}% across {SPAN} years). Exposure and "
         "workforce therefore did not move in opposite directions as a flight-from-risk "
         "account would predict.")
    rows = []
    for s in CORE:
        v = DESCR[EN[s]]
        rows.append([EN[s], v["phys_first"], v["phys_last"],
                     f"{v['litrate_first']:.2f}", f"{v['litrate_last']:.2f}",
                     v["hosp_first"], v["hosp_last"]])
    table(doc,
          ["Specialty", f"Physicians {BIEN[0]}", f"Physicians {BIEN[-1]}",
           f"Lit. rate {BIEN[0]}", f"Lit. rate {BIEN[-1]}",
           f"Hospital facilities {BIEN[0]}", f"Hospital facilities {BIEN[-1]}"],
          rows,
          f"Table 1. Physicians, litigation rate (per {_per} physicians) and hospital facilities by "
          "specialty, first and last waves.")

    head(doc, "Primary association and equivalence", level=2)
    body(doc,
         f"The lagged litigation rate was not associated with biennial physician growth "
         f"(coefficient {fmt(PHYS['coef'],4)}; 95% CI {fmt(PHYS['ci_low'],4)} to "
         f"{fmt(PHYS['ci_high'],4)}; p={PHYS['p']:.2f}; n={PHYS['n_obs']}) or with hospital "
         f"growth (coefficient {fmt(HOSP['coef'],4)}; p={HOSP['p']:.2f}). Equivalence "
         f"testing (Figure 1; Table 2) showed that a 1-SD higher litigation rate "
         f"changed biennial physician growth by less than \u00b1{MARGIN1}% (TOST p={p_tost_fmt(EQP['tests'][0]['p_tost'])}; "
         f"point estimate {fmt(EQP['coef_per_SD']*100,2)}% with 90% CI "
         f"{fmt(EQP['ci90_low']*100,2)}% to {fmt(EQP['ci90_high']*100,2)}%). For hospital "
         f"growth the point estimate was {fmt(EQH['coef_per_SD']*100,2)}% (90% CI "
         f"{fmt(EQH['ci90_low']*100,2)}% to {fmt(EQH['ci90_high']*100,2)}%): it was within "
         f"the \u00b1{MARGIN2}% margin (p={p_tost_fmt(EQH['tests'][1]['p_tost'])}) but not the stricter "
         f"\u00b1{MARGIN1}% margin (p={p_tost_fmt(EQH['tests'][0]['p_tost'])}). Thus the data are "
         "consistent with the absence of a policy-relevant effect on physician growth, "
         "and with at most a small effect on hospital facility-count growth. Detailed TOST results by "
         "margin are reported in Supplementary Table 2.")
    figure(doc, "ha_Figure_1.png",
           f"Figure 1. Equivalence (TOST) of the litigation-rate effect against \u00b1{MARGIN1}% and "
           f"\u00b1{MARGIN2}% margins; horizontal bars are 90% confidence intervals.")
    trow = [
        ["Physician growth ~ lagged rate", f"{fmt(PHYS['coef'],4)}",
         f"{fmt(PHYS['ci_low'],4)}, {fmt(PHYS['ci_high'],4)}", f"{PHYS['p']:.2f}", PHYS['n_obs']],
        ["Hospital facility-count growth ~ lagged rate", f"{fmt(HOSP['coef'],4)}",
         f"{fmt(HOSP['ci_low'],4)}, {fmt(HOSP['ci_high'],4)}", f"{HOSP['p']:.2f}", HOSP['n_obs']],
        ["Counts contrast (physician)", f"{fmt(CNT['coef'],4)}", "\u2014", f"{CNT['p']:.2f}", CNT['n_obs']],
        ["Annual hospital facility-count growth (sensitivity)", f"{fmt(ANN['coef'],4)}", "\u2014", f"{ANN['p']:.2f}", ANN['n_obs']],
        ["Interpolated physician (sensitivity)", f"{fmt(INT['coef'],4)}", "\u2014", f"{INT['p']:.2f}", INT['n_obs']],
        ["Reverse (workforce\u2192litigation)", f"{fmt(REV['coef'],3)}", "\u2014", f"{REV['p']:.2f}", REV['n_obs']],
    ]
    table(doc,
          ["Model", "Coefficient", "95% CI", "p", "n"],
          trow,
          "Table 2. Panel fixed-effects models and sensitivity analyses.")

    head(doc, "Small-cluster robustness and power", level=2)
    body(doc,
         "Because inference is based on only 12 specialty clusters, we checked the "
         "primary results with a cluster block-bootstrap (B = 1,999). For physician "
         f"growth the bootstrap 95% CI for the lagged litigation-rate coefficient was "
         f"{fmt(BS_PHYS['coef_boot_ci_low'],4)} to {fmt(BS_PHYS['coef_boot_ci_high'],4)} and the "
         f"bootstrap p-value was {BS_PHYS['p_bootstrap']:.2f}; for hospital facility-count growth the "
         f"bootstrap 95% CI was {fmt(BS_HOSP['coef_boot_ci_low'],4)} to "
         f"{fmt(BS_HOSP['coef_boot_ci_high'],4)} and the bootstrap p-value was "
         f"{BS_HOSP['p_bootstrap']:.2f}. Both intervals comfortably contain zero. "
         "Power diagnostics make the panel information explicit. For physician growth, "
         f"the minimum detectable effect was {EQP['mde_80pct']:.2f}% per SD at 80% power, "
         f"and the power to declare equivalence within the \u00b1{MARGIN1}% margin if the true "
         f"effect were zero was {EQP['tests'][0]['power_if_null']*100:.1f}%. For hospital facility-count growth the minimum "
         f"detectable effect was {EQH['mde_80pct']:.2f}% per SD and the equivalent power "
         f"for the \u00b1{MARGIN1}% margin was {EQH['tests'][0]['power_if_null']*100:.1f}%. The panel is therefore "
         "informative enough to rule out policy-relevant effects for physicians, and to "
         "bound any hospital effect within a small margin.")

    head(doc, "Counts versus rates, and confounders", level=2)
    body(doc,
         f"Using raw litigation counts rather than rates did not reveal a negative "
         f"association in this measured-only design (p={CNT['p']:.2f}). Figure 2 "
         "shows the contrast: the count exposure is confounded by specialty size (panel a), "
         "whereas the rate-adjusted exposure is not (panel b); points are coloured and shaped "
         "by specialty so readers can identify which fields drive any apparent pattern. "
         f"The annual hospital and interpolated annual-physician sensitivity analyses were also null "
         f"(p={ANN['p']:.2f} and p={INT['p']:.2f}), confirming that the null result is robust to panel "
         "frequency and exposure definition. The JOCS-CP indicator was positive in sign in the "
         f"obstetric-hospital model (coefficient {fmt(HOSP['jocscp_coef'],3)}, raw p={HOSP['jocscp_p']:.3f}), "
         f"but it did not remain significant after the small-cluster correction and Holm "
         f"adjustment for the exploratory sensitivity family (Holm p={JOCS_HOLM:.3f}); we "
         "therefore treat it as exploratory and do not interpret it as a causal policy "
         "effect.")
    figure(doc, "ha_Figure_2.png",
           "Figure 2. Biennial physician growth against lagged litigation exposure measured as "
           "(a) counts and (b) rates. Points are coloured by specialty; the count panel shows "
           "the size confounding that the rate panel removes.")
    body(doc,
         "The same count-versus-rate contrast for biennial hospital facility-count growth is shown in "
         "Figure 3. As with physician growth, the count exposure creates a spurious size "
         "confound that disappears once the rate-adjusted exposure is used.")
    figure(doc, "ha_Figure_3.png",
           "Figure 3. Biennial hospital facility-count growth against lagged litigation exposure measured as "
           "(a) counts and (b) rates. Points are coloured by specialty; the rate-adjusted "
           "panel shows no systematic association.")
    body(doc,
         f"Descriptively, per-specialty rank correlations between the lagged litigation "
         f"rate and physician growth were positive in {n_pos} of {N_SPECIALTIES} specialties and "
         f"statistically significant in {n_sig}; the direction is therefore, if anything, "
         "opposite to a flight-from-risk hypothesis.")
    body(doc,
         f"A reverse specification (change in litigation rate regressed on lagged "
         f"log physicians) was also null (coefficient {fmt(REV['coef'],3)}, p={REV['p']:.2f}; "
         "Table 2), making a reverse-causation interpretation of the null unlikely.")
    body(doc,
         f"We also evaluated JMSR medical-accident investigation report counts as a "
         f"potential confounder or competing exposure.{{mais}} From {JMSR_CORR['years'][0]} to "
         f"{JMSR_CORR['years'][-1]}, raw litigation and JMSR report counts were strongly "
         f"correlated across specialties (Pearson r={JMSR_CORR['pooled_r']:.2f}), because large "
         "specialties generate more of both. After removing specialty-specific levels and "
         f"trends, however, the within-specialty correlation was negligible (r={JMSR_CORR['detrended_r']:.2f}). "
         f"A model of annual hospital facility-count growth for {JMSR_START}-2024 that included both the "
         "lagged litigation rate and the lagged JMSR report rate left the litigation "
         f"coefficient essentially unchanged ({fmt(JMSR['lit_coef'],4)}; p={JMSR['lit_p']:.2f}) "
         f"and the JMSR term was not associated with hospital facility-count growth (p={JMSR['med_p']:.2f}; "
         "Supplementary Table 3). The null litigation result is therefore neither explained nor "
         "masked by broader medical-accident reporting.")
    body(doc,
         f"Finally, we tested national newspaper coverage from Nikkei Telecom 21 as a "
         f"potential confounder.{{nikkei}} Total annual article counts (keywords: "
         f"medical error + medical malpractice) and total litigation counts were correlated "
         f"(Pearson r={MEDIA_CORR['total_r']:.2f}), consistent with greater public attention in "
         f"high-litigation years. Within the annual hospital panel, however, the lagged "
         "litigation rate and the media-count series were only weakly correlated. "
         f"A model of annual hospital facility-count growth for {MEDIA_START}-{MEDIA_END} that included both the "
         f"lagged litigation rate and the lagged article count (per 1,000 articles) left the "
         "litigation coefficient essentially unchanged and the media term was not associated "
         f"with hospital facility-count growth (p={MEDIA['media_p']:.2f}; Supplementary Table 4). Media coverage "
         "therefore does not explain the null litigation effect either. "
         "Holm step-down adjusted p-values for the exploratory sensitivity family are "
         "reported in Supplementary Table 5.")

    head(doc, "Policy lever simulation", level=2)
    body(doc,
         "The counterfactual projection made the practical implications of the null "
         "regression result explicit. Under the point estimate, eliminating all malpractice "
         "litigation would add only "
         f"{TOTAL_SIM.get('marginal_pct_lit_point', 0):.1f}% to the projected 2034 national physician "
         "stock relative to baseline drift. Even under the 95% lower-bound (most favourable) "
         f"coefficient it would add {TOTAL_SIM.get('marginal_pct_lit_lower', 0):.1f}%, comparable to the "
         f"{TOTAL_SIM.get('marginal_pct_mde', 0):.1f}% gain from a generic lever equal to the minimum detectable "
         "effect. General surgery, the only specialty with negative baseline drift, illustrates "
         "the break-even arithmetic: its projected 2024-2034 decline of "
         f"{SURG_SIM.get('pct_change_baseline', 0):.1f}% would be reduced to "
         f"{SURG_SIM.get('pct_change_lit_zero_point', 0):.1f}% under the point estimate and reversed to "
         f"{SURG_SIM.get('pct_change_lit_zero_lower', 0):.1f}% under the 95% lower bound. The latter "
         "requires eliminating every remaining closed claim and assumes the most adverse (most "
         "negative) coefficient compatible with the data; a more realistic policy would achieve "
         "far less. Full projected 2034 physician counts by specialty and lever are reported in "
         "Supplementary Table 7; Figure 4 summarises the same information as marginal percentage "
         "changes. Litigation reduction is therefore not a high-leverage instrument for workforce "
         "allocation in this setting.")
    figure(doc, "ha_Figure_4.png",
           "Figure 4. Counterfactual policy-lever simulation: marginal 10-year change in "
           "physician counts by specialty relative to the projected baseline drift. "
           "The MDE benchmark is the minimum detectable per-SD effect from the primary analysis.")

    # Discussion
    head(doc, "Discussion", level=1)
    body(doc,
         "Using national primary data, rates rather than counts, and only measured "
         "physician observations, we found no association between specialty-level "
         "malpractice-litigation risk and subsequent physician or hospital decline. "
         f"Equivalence testing showed that any effect of litigation risk on biennial physician "
         f"growth is smaller than {MARGIN1}% (90% CI within the {MARGIN1}% margin), and any "
         f"effect on hospital facility-count growth is smaller than {MARGIN2}% (but not confidently smaller than "
         f"{MARGIN1}%). These data therefore do not support the hypothesis that physicians "
         f"systematically abandon high-litigation specialties over {SPAN} years of official statistics.")
    body(doc,
         "A null result is not merely a failure to detect an effect. The narrow "
         "confidence intervals, pre-specified equivalence margins, and power diagnostics "
         "allow us to say that, if litigation risk does influence specialty-level workforce "
         "growth, the magnitude is too small to matter for workforce planning. Several "
         "interpretations are consistent with this finding. First, litigation risk may affect "
         "clinical behaviour on the intensive margin (defensive medicine) rather than the "
         "extensive margin (specialty exit). Second, the closed-claim rate is an objective "
         "measure, but the perceived risk that drives career decisions may be shaped more by "
         "rare, salient criminal prosecutions and media coverage than by routine civil claims. "
         "Third, high-risk specialties in Japan also carry substantial fee-for-service returns, "
         "so income effects, switching costs and status-quo bias can offset any deterrent. "
         "Finally, the negative finding is informative in its own right: studies that reject a "
         "commonly assumed policy mechanism are under-represented in the literature, and a well-"
         "powered equivalence result reduces the risk that policy is directed at an ineffective "
         "lever.{tversky1973,kahneman1979,samuelson1988}")
    body(doc,
         "The raw-count sensitivity did not reveal a negative association in these data, "
         "illustrating that the apparent count-litigation relationship does not translate into "
         "a behavioural effect once specialty size is accounted for (Figure 2). This is a "
         "cautionary example for workforce research that pairs administrative count series, "
         "and shows why rate-based, measured-only designs are preferable when testing "
         "litigation-workforce hypotheses.")
    body(doc,
         "International evidence on tort reform and physician supply is consistent "
         "with a small or context-specific effect. Matsa found that U.S. state damage "
         "caps increased the supply of frontier rural specialists by 10-12 percent, "
         "but did not affect physician supply for the average resident.{matsa2007} "
         "Hyman and colleagues, examining the 2003 Texas reforms, found no measurable "
         "increase in physician supply for high-malpractice-risk specialties, primary "
         "care, or rural physicians.{hyman2015} Frakes and colleagues showed that "
         "negligence-standard reforms could shift the composition of the physician "
         "workforce toward surgery in some regions, yet the effect was localized and "
         "modest.{frakes2020} Against this backdrop, a null effect of civil litigation "
         "risk on specialty supply is not surprising, especially in a system "
         "with comparatively low litigation volume and predictable damages.")
    body(doc,
         "Even when litigation risk is unrelated to the number of physicians, it may "
         "shape clinical behaviour through defensive medicine. Kessler and McClellan "
         "showed that U.S. malpractice reforms reduced medical expenditures for elderly "
         "heart-disease patients without increasing mortality or complications, "
         "suggesting that defensive practice is one margin of adjustment to liability "
         "pressure.{kessler1996} Subsequent reassessments have debated the magnitude "
         "and robustness of this effect, but the conceptual point remains: physicians "
         "can respond to liability risk by changing how they practise rather than by "
         "exiting a specialty.{sloan2008} Fee-for-service reimbursement in Japan "
         "rewards the high-acuity procedural work that also carries litigation "
         "exposure, so the financial return to remaining in surgery, obstetrics, or "
         "interventional specialties may dominate any deterrent from civil claims.")
    body(doc,
         "This pattern is consistent with well-documented behavioural-economics mechanisms. "
         "Media coverage of sensational malpractice or criminal prosecutions makes litigation "
         "risk highly available to physicians and trainees, and loss aversion can cause a rare "
         "but salient adverse outcome to be overweighted in career deliberations.{tversky1973,kahneman1979} "
         "Yet the actual decision to leave a specialty is governed by expected income, sunk "
         "training costs, switching costs and status-quo bias, all of which discourage exit "
         "even when perceived risk is high.{samuelson1988} The discrepancy between "
         "reported anxiety and measured supply is therefore not a contradiction; it is exactly "
         "what one would expect when a vivid, low-probability risk meets strong economic "
         "and institutional incentives to remain.")
    body(doc,
         "The civil litigation environment in Japan itself dampens the likelihood of a "
         "flight-from-risk response. Taniguchi and colleagues analysed all closed "
         "malpractice claims reported by the Supreme Court from 2006 to 2021 and "
         "found that more than half ended in settlement, plaintiffs won only about a "
         "quarter of judgments, and the number of claims has been declining, "
         "especially in obstetrics and gynaecology.{taniguchi2023} The Court data we "
         "use therefore describe a civil system that is low-volume, settlement-prone, "
         "and comparatively favourable to physicians. This context makes it unlikely "
         "that routine civil litigation risk alone would drive physicians out of "
         "high-risk fields.")
    body(doc,
         "No-fault obstetric compensation (the Japan Obstetric Compensation System for Cerebral Palsy, 2009) illustrates "
         "a different mechanism. It was introduced partly because of a shortage of "
         "young obstetricians and regional gaps in maternity care, and it combined "
         "no-fault compensation with investigation and prevention.{hasegawa2016} The "
         f"hospital-level JOCS-CP indicator was directionally positive (coefficient "
         f"{fmt(HOSP['jocscp_coef'],3)}, raw p={HOSP['jocscp_p']:.3f}), but it did not remain "
         f"significant after the small-cluster correction and Holm adjustment for the "
         f"exploratory sensitivity family (Holm p={JOCS_HOLM:.3f}). This suggests that, "
         "if the JOCS-CP did support obstetric hospital supply, the effect would be too small "
         "or too confounded by concurrent obstetric policies to be isolated here. "
         "Civil litigation exposure is also distinct from criminal prosecution. Morita "
         "studied the 2004 Fukushima obstetrician prosecution and found a 13 percent "
         "decline in obstetricians, with some switching to gynaecology.{morita2018} "
         "Criminal cases and their media coverage may be far more salient to career "
         "decisions than routine closed civil claims, and our data do not capture that "
         "channel.")
    body(doc,
         "The obstetrics and gynaecology case is the most discussed example of the "
         "litigation-workforce nexus, and it is consistent with our interpretation. "
         "A recent Japan\u2013U.S. comparison of medical-legal claims in obstetrics and "
         "gynaecology (OB/GYN) found that the proportion of malpractice claims in this specialty fell from "
         "15.1 percent in 2004 to 5.2 percent in 2022, and that claims per 100 OB/GYN "
         "physicians fell from 0.9 in 2007 to 0.4 in 2016, while maternal and neonatal "
         "mortality also declined.{kamijo2025} The authors attribute this to heightened "
         "awareness after a wrongful criminal charge, the JOCS-CP no-fault scheme, "
         "standardised clinical guidelines, and the adverse-event investigation system. "
         "This is not evidence that lowering litigation risk caused the workforce to "
         "grow; it is evidence that obstetric litigation, workforce support, and safety "
         "interventions moved together. Surveys of OB/GYN residents in Japan, Korea, "
         "and Taiwan likewise show that litigation is reported as a negative factor, "
         "but that its perceived importance is smaller where no-fault compensation exists "
         "and that workload, lifestyle, and professional interest remain dominant.{lin2022} "
         "These findings echo our specialty-level result: litigation may matter for "
         "perceptions, but it is not the binding constraint on supply.")
    body(doc,
         "What do these findings imply for policy? Reducing civil malpractice "
         "litigation is unlikely to be a powerful lever for correcting specialty "
         "maldistribution in this national setting. The 10-year counterfactual simulation "
         f"showed that eliminating all litigation would add only {TOTAL_SIM.get('marginal_pct_lit_point', 0):.1f}% "
         "to the projected national physician stock under the point estimate, and even "
         f"{TOTAL_SIM.get('marginal_pct_lit_lower', 0):.1f}% under the most favourable 95% lower-bound "
         "coefficient, before accounting for the implausibility of zero claims. Structural incentives are more promising: "
         "no-fault compensation can de-risk high-acuity specialties, and payment "
         "design can reward service in underserved settings and activities. The "
         "JOCS-CP experience supports the former; the country's fee-for-service schedule "
         "and rural/urban payment adjustments illustrate the latter. Malpractice reform "
         "may still matter for defensive medicine, patient compensation, and provider-patient "
         "trust. But our evidence does not support the claim, at least from these data, that lowering "
         "litigation risk will retain physicians in high-risk specialties. "
         "Policymakers should therefore target structural incentives before "
         "relying on litigation-avoidance messaging.")
    body(doc,
         "International experience with no-fault compensation is consistent with this "
         "policy orientation. New Zealand replaced tort-based medical-injury "
         "compensation with a government-funded no-fault scheme in 1974 and, after "
         "2005 reforms, extended coverage to all treatment injuries; this separated "
         "compensation from negligence findings and largely barred malpractice "
         "litigation.{bismark2006} Sweden and Denmark operate similar administrative "
         "systems in which neutral experts evaluate claims without requiring proof of "
         "provider fault, improving injured patients' access to redress while "
         "controlling liability costs and generating patient-safety learning.{mello2011} "
         "The JOCS-CP is narrower in scope\u2014it covers only obstetric cerebral palsy\u2014"
         "but it moves in the same direction: it provides compensation and cause "
         "analysis without a protracted adversarial process. Extending such an "
         "approach more broadly would be a structural alternative to repeated calls to "
         "reduce malpractice litigation as a workforce strategy.")
    body(doc,
         "Several questions remain for future research. We cannot observe individual "
         "physicians' risk perceptions, career intentions, or responses to media "
         "coverage of high-profile cases. The closed-claim rate is an objective "
         "exposure measure, but it may not capture the perceived risk that drives "
         "specialty choice, especially when criminal prosecutions or sensational media "
         "coverage shape beliefs. Helland and Seabury's review concludes that tort "
         "reform effects on physician supply are heterogeneous across states and "
         "specialties, and that more granular, state-specific designs are needed to "
         "settle the question in the U.S. context.{helland2015} Applying similar "
         "logic in the same national setting would require individual-level or prefecture-level career "
         "data linked to local litigation, media, and reimbursement environments. "
         "Until then, the present specialty-level rate analysis provides the most "
         "systematic evidence available on the central policy question: whether "
         "malpractice litigation risk drives physicians away from high-risk "
         f"specialties. The answer, in these data, is no—or at least not in a way that is "
         f"detectable or policy-relevant across {len(BIEN)} measured waves ({BIEN[0]}–{BIEN[-1]}).")

    body(doc,
         "The framework we used is intentionally general and falls within healthcare "
         "operations research and decision analytics. Specialty physician workforce "
         "allocation is a recurring healthcare decision problem, and the same "
         "sensitivity-analysis steps—rate adjustment to remove size confounding, measured-"
         "only panels to avoid interpolation, equivalence testing with policy-relevant "
         "margins, and small-cluster inference—can be applied to other proposed levers, "
         "such as fee schedules, regional quotas, or training subsidies. The analytic "
         "contribution is therefore not the malpractice finding itself but a transparent, "
         "reproducible decision-support workflow that helps planners distinguish meaningful "
         "workforce effects from spurious count- or interpolation-based associations.")

    head(doc, "Limitations", level=2)
    body(doc,
         f"This is an ecological, specialty-level analysis and cannot establish individual-level "
         f"causality. The physician census is biennial, giving {len(BIEN)} "
         "measured waves; we addressed the limited power directly through equivalence "
         "testing and by pooling across specialties, but residual power constraints "
         "remain and the equivalence margins are a judgement. Litigation rates may be endogenous "
         "to physician supply if a smaller workforce increases workload and hence incidents; "
         "the lagged exposure, fixed effects, and reverse specification make reverse causation "
         "unlikely, yet unobserved confounders at the specialty or prefecture level cannot be "
         "fully ruled out. Cluster block-bootstrap and power diagnostics are reported in "
         "Supplementary Table 6. Because clusters are "
         f"defined by the {N_SPECIALTIES} specialties, the small-cluster correction uses G-1={PHYS['df']} "
         "degrees of freedom; this is the minimum at which cluster-robust t inference is "
         "recommended and is inherent to the data. Specialty-specific "
         f"litigation counts could be recovered only from {BIEN[0]}; pre-{BIEN[0]} specialty tables were "
         "not retrievable from primary sources. Clinic counts by specialty are "
         f"published only every {CLINIC_RES} years and were used descriptively. JMSR "
         f"report counts are available only from {JMSR_CORR['years'][0]} and were used in a "
         f"{JMSR_START}-2024 sensitivity. Media article counts are available only for "
         f"{MEDIA_START}-{MEDIA_END} and are a national total, so they cannot be decomposed by "
         "specialty and are collinear with full wave fixed effects. Litigation counts are assigned to a principal specialty "
         "and, by the Court's own note, do not measure intrinsic specialty risk.{court} "
         "Finally, these findings are "
         "embedded in the country's particular legal, cultural and institutional "
         "context\u2014including its no-fault obstetric compensation scheme, its "
         "fee-for-service reimbursement structure and its comparatively low-volume "
         "malpractice-litigation culture\u2014so physician responses to litigation risk "
         "may differ in health systems with different liability regimes, compensation "
         "mechanisms or professional norms; the results should not be assumed to "
         "generalise across cultural spheres.")

    head(doc, "Conclusions", level=1)
    body(doc,
         f"Across {YEARS}, specialty-level malpractice-litigation risk was not associated "
         "with physician or hospital decline in these national data, and the "
         "physician effect was statistically equivalent to null within a small margin. "
         "From a decision-analytics standpoint, malpractice litigation is not a reliable "
         "policy lever for correcting specialty maldistribution in this setting. Policymakers "
         "may more productively target structural incentives, especially no-fault "
         "compensation, rather than rely on the assumption that reducing litigation will "
         "retain physicians in high-risk specialties. The transparent, reproducible "
         "sensitivity-analysis framework used here is exportable to other healthcare "
         "workforce-policy levers.")

    # Declaration of generative AI use (Elsevier requirement; place between Conclusions and References)
    head(doc, "Declaration of generative AI use", level=1, numbered=False)
    body(doc,
         "[Authors: insert the Elsevier AI declaration here before submission. "
         "Example wording: During the preparation of this work the author(s) used "
         "[TOOL NAME] in order to [PURPOSE]. After using this tool/service, the "
         "author(s) reviewed and edited the content as needed and take(s) full "
         "responsibility for the content of the publication. If no generative AI was "
         "used, state so.]")

    # Declarations
    head(doc, "Declarations", level=1, numbered=False)
    para(doc,
         "Funding: none. Competing interests: none declared. "
         "Ethics approval: this study used publicly available aggregated national "
         "statistics and did not involve human subjects, identifiable data or patient "
         "records; no ethics approval was required. "
         f"Data and code availability: all primary data files, extraction scripts and "
         "analysis code are openly available in the project repository ("
         f"{PUBLIC_REPO}), enabling full reproduction of every reported number.")

    # References
    head(doc, "References", level=1, numbered=False)
    missing = [k for k in REFS if k not in _CITE_ORDER]
    if missing:
        raise SystemExit(f"orphan references (in list, never cited): {missing}")
    for i, k in enumerate(_CITE_ORDER, 1):
        p = doc.add_paragraph()
        p.paragraph_format.line_spacing = 1.5
        r = p.add_run(f"{i}. {REFS[k]}")
        r.font.size = Pt(10)
        r.font.name = "Times New Roman"

    out = os.path.join(BASE, "ha_manuscript_en.docx")
    normalize_docx(doc)
    doc.save(out)
    main_wc = sum(wc(t) for t in BODY_TEXTS)
    print(f"wrote {out}; abstract {abstract_wc} words; main body ~{main_wc} words")
    return main_wc, abstract_wc


def build_title_page(main_word_count):
    doc = _setup_doc()
    for _ in range(4):
        doc.add_paragraph()
    t = doc.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    t.paragraph_format.space_after = Pt(18)
    rt = t.add_run(MANUSCRIPT_TITLE)
    rt.bold = True
    rt.font.size = Pt(15)
    rt.font.name = "Times New Roman"

    lines = [
        "Authors: Onishi Tatsuki",
        "Affiliation: Data Science AI Innovation Research Promotion Center, Shiga University, "
        "1-1-1 Bamba, Hikone, Shiga 522-8522, Japan",
        "Corresponding author: Onishi Tatsuki (email: [corresponding author email])",
        f"Word count (main text): approximately {main_word_count} words (excluding abstract, references, declarations, tables and figure legends)",
        "Article type: Original research article",
        "Target journal: Healthcare Analytics (Elsevier)",
        "Tables: 2  Figures: 4  Supplementary tables: 9  Supplementary figures: 3",
        "Conflicts of interest: none declared",
        "Funding: none",
        f"Data availability: all primary data and analysis code are openly available in the project repository ({PUBLIC_REPO}).",
    ]
    for line in lines:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p.add_run(line)
        r.font.size = Pt(12)
        r.font.name = "Times New Roman"

    out = os.path.join(BASE, "ha_title_page.docx")
    normalize_docx(doc)
    doc.save(out)
    print("wrote", out)


def build_highlights():
    highlights = [
        f"Litigation risk is unrelated to physician or hospital decline in {N_SPECIALTIES} specialties.",
        "A decision-analytics framework removes size confounding and sparse-panel bias.",
        "Equivalence, bootstrap, and power diagnostics support an informative null.",
        "Behavioural-economics mechanisms and structural incentives explain the null.",
        "Policy should target structural incentives, not litigation-avoidance messaging.",
    ]
    for h in highlights:
        if len(h) > 85:
            raise SystemExit(f"Highlight exceeds 85 characters ({len(h)}): {h}")

    doc = _setup_doc()
    h = doc.add_paragraph()
    r = h.add_run("Highlights")
    r.bold = True
    r.font.size = Pt(13)
    r.font.name = "Times New Roman"
    for item in highlights:
        p = doc.add_paragraph(style="List Bullet")
        p.clear()
        pr = p.add_run(item)
        pr.font.name = "Times New Roman"
        pr.font.size = Pt(12)
        p.paragraph_format.line_spacing = 1.5
        p.paragraph_format.space_after = Pt(4)
    out = os.path.join(BASE, "ha_highlights.docx")
    normalize_docx(doc)
    doc.save(out)
    print("wrote", out)


def build_cover_letter():
    doc = _setup_doc()
    for line in ["[Date]", "", "Prof. Dr. Madjid Tavana, PhD", "Editor-in-Chief", "Healthcare Analytics", ""]:
        p = doc.add_paragraph()
        if line:
            r = p.add_run(line)
            r.font.size = Pt(12)
            r.font.name = "Times New Roman"
    p = doc.add_paragraph()
    p.add_run("Dear Editor,").font.size = Pt(12)
    p.runs[0].font.name = "Times New Roman"

    paragraphs = [
        f'We submit an original research article, "{MANUSCRIPT_TITLE}", for '
        f'consideration by Healthcare Analytics.',
        "Healthcare Analytics advances data-driven analytics, operations research, and "
        "decision science for healthcare. Our study sits at this intersection: it treats "
        "specialty physician workforce planning as a healthcare resource-allocation problem "
        "and uses the litigation-workforce question as a test case for a transparent, "
        "reproducible decision-analytics and policy-lever evaluation framework. The "
        "analytical contribution is diagnostic: we show how two common observational fallacies\u2014"
        "size confounding and interpolation of sparse panel data\u2014can be identified and removed "
        "when a proposed policy lever (malpractice-litigation risk) is evaluated against "
        "workforce outcomes, and we add cluster block-bootstrap and power diagnostics to the "
        "standard fixed-effects / equivalence-testing toolkit. To our knowledge, this is the first "
        "application of pre-specified equivalence testing and small-cluster inference diagnostics to "
        "the malpractice-litigation workforce literature. We also include a counterfactual "
        "policy-lever simulation that projects 2034 physician counts under stylised litigation-"
        "reduction and structural-incentive scenarios, making the practical decision implications "
        "of the regression results explicit.",
        "The national administrative data we use come from Japan, a setting that provides a "
        f"complete, long-running test case. Using national primary data for {N_SPECIALTIES} clinical specialties, we "
        "measure exposure as a size-adjusted rate (closed malpractice claims per "
        f"{_per} physicians) and estimate panel fixed-effects models with small-cluster "
        "inference, equivalence (TOST) testing, and bootstrap robustness checks. The effect of "
        "litigation risk on physician and hospital facility-count growth is statistically equivalent to a null "
        "within policy-relevant margins. We show that count-based associations disappear once "
        "specialty-size confounding and interpolation are removed, and we discuss how "
        "behavioural-economics mechanisms (saliency, loss aversion, status-quo bias) and "
        "structural incentives\u2014notably no-fault obstetric compensation\u2014may sustain the "
        "specialty workforce. Because these behavioural patterns arise from real-world administrative "
        "and clinical data, the study offers a policy-relevant example of how perceived risk and economic "
        "incentives interact in medical labour markets.",
        f"All data and code are openly available in the project repository "
        f"({PUBLIC_REPO}) and every reported number is reproducible from the raw primary files.",
        "The work is original, not under consideration elsewhere, and all authors approve "
        "the submission. We declare no conflicts of interest.",
    ]
    for b in paragraphs:
        p = doc.add_paragraph()
        r = p.add_run(b)
        r.font.size = Pt(11)
        r.font.name = "Times New Roman"
        p.paragraph_format.space_after = Pt(8)
        p.paragraph_format.line_spacing = 1.5

    for line in ["Sincerely,", "", "[Corresponding author, on behalf of all authors]"]:
        p = doc.add_paragraph()
        if line:
            r = p.add_run(line)
            r.font.size = Pt(12)
            r.font.name = "Times New Roman"

    out = os.path.join(BASE, "ha_cover_letter.docx")
    normalize_docx(doc)
    doc.save(out)
    print("wrote", out)


def build_supplementary():
    doc = _setup_doc()
    head(doc, "Supplementary material", level=1, numbered=False)

    # Supplementary Note 1: detailed statistical methods (supports Methods section simplification)
    head(doc, "Supplementary Note 1. Detailed statistical methods", level=2, numbered=False)
    para(doc,
         "This note records the full model specification, equivalence-testing rationale, and power "
         "diagnostics summarised in the main text. It also describes the heterogeneity and trend-"
         "stability checks reported in Supplementary Table 9.")
    para(doc,
         "Primary estimating equation. Let s index specialty, t index the biennial wave, and Y be "
         "the count of physicians or hospitals. The baseline model is")
    add_math(doc, r"\Delta \log(Y_{st}) = \alpha_s + \delta_t + \beta \cdot \text{litrate}_{s,t-1} + \gamma \cdot \text{JOCS-CP}_{s,t-1} + \epsilon_{st}")
    para(doc,
         "where alpha_s are specialty fixed effects, delta_t are wave fixed effects, litrate is closed "
         "malpractice claims per 1,000 physicians, and JOCS-CP is an obstetrics-and-gynaecology-specific "
         "indicator from 2009 onward. Standard errors are clustered by specialty and inference uses a "
         "t-distribution with G-1 degrees of freedom, where G=12.")
    para(doc,
         "Equivalence testing. Two one-sided tests (TOST) evaluate whether the per-SD litigation-rate "
         "coefficient beta lies inside a pre-specified symmetric margin m. The null hypotheses are")
    add_math(doc, r"H_0: \beta \leq -m \quad \text{and} \quad H_0: \beta \geq +m")
    para(doc,
         "Equivalence is declared when both one-sided tests yield p < 0.05. Margins of 1% and 2% biennial "
         "workforce change were chosen because they are smaller than typical policy targets for specialty "
         "rebalancing.")
    para(doc,
         "Power diagnostics. The minimum detectable effect (MDE) at 80% power for the per-SD coefficient "
         "is (t(0.975, df) + t(0.80, df)) x SE. The power to declare equivalence within margin m when the "
         "true effect is zero is 2 x F_t(m / SE) - 1, where F_t is the cumulative distribution function of "
         "the t(df) distribution. These quantities make the limited information in a 12-cluster panel explicit.")
    para(doc,
         "Heterogeneity and trend stability. To test whether the association differed by baseline "
         "litigation risk or surgical orientation, the model was augmented with an interaction between "
         "litrate and a binary group indicator:")
    add_math(doc, r"\Delta \log(Y_{st}) = \alpha_s + \delta_t + \beta \cdot \text{litrate}_{s,t-1} + \delta_g \cdot (\text{litrate}_{s,t-1} \times \text{Group}_s) + \epsilon_{st}")
    para(doc,
         "Group_s is either high-litigation (above the median specialty mean lagged litigation rate) "
         "or surgical (surgery, orthopaedics, obstetrics and gynaecology, urology). The main effect of "
         "Group is absorbed by the specialty fixed effects and is omitted. A separate stability check "
         "added specialty-specific linear trends, C(specialty) x year, to the baseline specification.")

    para(doc, "Supplementary Figure 1. Sensitivity-analysis framework for evaluating "
              "malpractice-litigation risk as a healthcare workforce-allocation instrument.")
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    img = os.path.join(OUT, "ha_Supplementary_Figure_1.png")
    if os.path.exists(img):
        p.add_run().add_picture(img, width=Inches(5.8))
    doc.add_paragraph()

    para(doc, f"Supplementary Figure 2. Closed malpractice claims per {_per} physicians by "
              f"specialty, 2008\u20132024 (rates, not counts).")
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    img = os.path.join(OUT, "ha_Supplementary_Figure_2.png")
    if os.path.exists(img):
        p.add_run().add_picture(img, width=Inches(5.8))
    doc.add_paragraph()

    para(doc, "Supplementary Figure 3. Physician workforce by specialty, indexed to 2008 (=100).")
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    img = os.path.join(OUT, "ha_Supplementary_Figure_3.png")
    if os.path.exists(img):
        p.add_run().add_picture(img, width=Inches(5.8))
    doc.add_paragraph()

    # Supplementary Table 1: data sources
    res_words = {1: "Annual", 2: "Biennial", 3: "Every 3 years"}
    table(doc,
          ["Series", "Source", "Resolution", "Years", "Role"],
          [["Physicians by specialty", "MHLW Statistics of Physicians", res_words.get(PHYS_RES, f"Every {PHYS_RES} years"),
            _year_label(PHYS_DF.columns), "Denominator & outcome"],
           ["Closed malpractice claims", "Supreme Court, by specialty", res_words.get(_resolution(LIT_DF), f"Every {_resolution(LIT_DF)} years"),
            _year_label(LIT_DF.columns), "Exposure (numerator)"],
           ["Hospitals by specialty", "MHLW Survey of Medical Institutions", res_words.get(HOSP_RES, f"Every {HOSP_RES} years"),
            _year_label(HOSP_DF.columns), "Outcome"],
           ["Clinics by specialty", "MHLW Survey (static)", res_words.get(CLINIC_RES, f"Every {CLINIC_RES} years"),
            _year_label(CLINIC_DF.columns), "Descriptive only"],
           ["JMSR report counts", "JMSR / MAIS", "Annual",
            _year_label(JMSR_CORR['years']), "Sensitivity (2016-2024)"],
           ["Newspaper article counts", "Nikkei Telecom 21", "Annual",
            _year_label(MEDIA_CORR['years']), f"Sensitivity ({MEDIA_START}-{MEDIA_END})"]],
          "Supplementary Table 1. Primary data sources and their resolution.")

    # Supplementary Table 2: TOST details
    eqp, eqh = RES["equivalence"][0], RES["equivalence"][1]
    rows = []
    for outcome, eq in [("Physician growth", eqp), ("Hospital growth", eqh)]:
        for t in eq["tests"]:
            rows.append([
                outcome,
                f"{eq['coef_per_SD']*100:+.2f}%",
                f"{eq['ci90_low']*100:+.2f}%, {eq['ci90_high']*100:+.2f}%",
                f"\u00b1{int(t['margin']*100)}%",
                f"{p_tost_fmt(t['p_tost'])}",
                "Yes" if t["equivalent"] else "No"
            ])
    table(doc,
          ["Outcome", "Coef per SD", "90% CI", "Margin", "TOST p", "Equivalent"],
          rows,
          "Supplementary Table 2. Equivalence (TOST) results (effect per +1 SD litigation rate).")

    # Supplementary Table 3: JMSR-adjusted hospital model
    table(doc,
          ["Exposure", "Coefficient", "p", "n"],
          [["Lagged litigation rate", fmt(JMSR["lit_coef"], 4), f"{JMSR['lit_p']:.2f}", JMSR["n_obs"]],
           ["Lagged JMSR report rate", fmt(JMSR["med_coef"], 4), f"{JMSR['med_p']:.2f}", JMSR["n_obs"]]],
          f"Supplementary Table 3. JMSR-adjusted annual hospital facility-count growth model "
          f"({JMSR_START}-2024) with both exposures entered simultaneously.")

    # Supplementary Table 4: media-adjusted hospital model
    table(doc,
          ["Exposure", "Coefficient", "p", "n"],
          [["Lagged litigation rate", fmt(MEDIA["lit_coef"], 4), f"{MEDIA['lit_p']:.2f}", MEDIA["n_obs"]],
           ["Lagged media count (per 1,000 articles)", fmt(MEDIA["media_coef"], 4), f"{MEDIA['media_p']:.2f}", MEDIA["n_obs"]]],
          f"Supplementary Table 4. Media-adjusted annual hospital facility-count growth model "
          f"({MEDIA_START}-{MEDIA_END}) with a linear time trend; full year fixed effects are "
          "omitted because the national article-count series is collinear with them.")

    # Supplementary Table 5: multiple comparison adjustment
    mc_rows = []
    for t in RES["multiple_comparison"]["tests"]:
        mc_rows.append([t["label"], f"{t['raw_p']:.3f}", f"{t['holm_p']:.3f}"])
    table(doc,
          ["Test", "Raw p", "Holm-adjusted p"],
          mc_rows,
          "Supplementary Table 5. Holm step-down adjusted p-values for the exploratory "
          "sensitivity tests and the JOCS-CP indicator.")

    # Supplementary Table 6: bootstrap and power diagnostics
    bs_rows = []
    for label, b, eq in [("Physician growth", BS_PHYS, EQP),
                          ("Hospital growth", BS_HOSP, EQH)]:
        for t in eq["tests"]:
            bs_rows.append([
                label,
                f"{b['p_bootstrap']:.3f}" if b else "",
                f"{b['coef_boot_ci_low']:.5f}, {b['coef_boot_ci_high']:.5f}" if b else "",
                f"\u00b1{int(t['margin']*100)}%",
                f"{t['power_if_null']*100:.1f}%",
                f"{eq['mde_80pct']:.2f}%"
            ])
    table(doc,
          ["Outcome", "Bootstrap p", "Bootstrap 95% CI", "Equivalence margin",
           "Power if true null", "MDE (80% power)"],
          bs_rows,
          "Supplementary Table 6. Cluster block-bootstrap p-values, percentile 95% "
          "confidence intervals, equivalence power (if true effect is zero), and "
          "minimum detectable effect per 1-SD litigation-rate increase.")

    # Supplementary Table 7: counterfactual 2034 physician counts
    table(doc,
          ["Specialty", "2024 count", "Baseline 2034", "Litigation zero (point)",
           "Litigation zero (95% lower)", "MDE benchmark"],
          _simulation_table_rows(),
          "Supplementary Table 7. Counterfactual 2034 physician counts by specialty and policy instrument. "
          "Counts are projected from observed biennial baseline drift plus the indicated effect; "
          "marginal percentage changes are shown in Figure 4.")

    # Supplementary Table 8: specialist-trainee coverage context
    if SENKOI is not None:
        table(doc,
              ["Specialty", "2018 specialist trainees (senkoi)",
               "Physicians 3\u20135 yr after registration (2014)",
               "H24 cohort 3rd year (2012)", "H24 cohort 5th year (2014)", "Coverage"],
              _senkoi_table_rows(),
              "Supplementary Table 8. Specialist-trainee coverage by specialty. "
              "Coverage is the number of 2018 first-year specialist trainees divided by the number "
              "of physicians reported 3\u20135 years after medical registration in 2014, by primary specialty.")

    # Supplementary Table 9: heterogeneity and trend stability
    het_rows = []
    for r in RES.get("heterogeneity", []):
        het_rows.append([
            "Physician growth" if "physician" in r["label"].lower() else "Hospital growth",
            r["group"].capitalize(),
            f"{fmt(r['coef'], 4)} ({fmt(r['se'], 4)})",
            f"{r['p']:.2f}",
            f"{fmt(r['interact_coef'], 4)} ({fmt(r['interact_se'], 4)})",
            f"{r['interact_p']:.2f}",
            r["n_obs"],
        ])
    for r in RES.get("trend_sensitivity", []):
        het_rows.append([
            "Physician growth" if "physician" in r["label"].lower() else "Hospital growth",
            "Specialty trends",
            f"{fmt(r['coef'], 4)} ({fmt(r['se'], 4)})",
            f"{r['p']:.2f}",
            "\u2014", "\u2014",
            r["n_obs"],
        ])
    table(doc,
          ["Outcome", "Group / check", "Main coefficient (SE)", "p",
           "Interaction coefficient (SE)", "p", "n"],
          het_rows,
          "Supplementary Table 9. Heterogeneity and trend-stability checks. "
          "High-litigation and surgical models include an interaction between the lagged litigation rate "
          "and a binary group indicator (the main effect of the group is absorbed by specialty fixed effects). "
          "The trend-stability model adds specialty-specific linear time trends.")

    # STROBE checklist
    head(doc, "STROBE checklist", level=1, numbered=False)
    strobe_items = [
        ("1", "Title and abstract indicate the study design"),
        ("2", "Background and rationale"),
        ("3", "Specific objectives and hypotheses"),
        ("4", "Key elements of study design"),
        ("5", "Data sources, locations and dates"),
        ("6", "Eligibility and selection criteria"),
        ("7", "Definitions of outcomes, exposures and predictors"),
        ("8", "Data sources and measurement methods"),
        ("9", "Discussion of potential sources of bias"),
        ("10", "Explanation of study size"),
        ("11", "Quantitative variables: grouping and transformations"),
        ("12", "Description of all statistical methods"),
        ("13", "Numbers of participants at each stage"),
        ("14", "Descriptive characteristics of participants"),
        ("15", "Outcome data for each exposure category"),
        ("16", "Main results with measures of uncertainty"),
        ("17", "Other analyses (sensitivity and subgroup)"),
        ("18", "Summary of key results with reference to objectives"),
        ("19", "Limitations of bias and imprecision"),
        ("20", "Cautious overall interpretation"),
        ("21", "Generalisability to other populations"),
        ("22", "Registration, protocol, funding, data and code availability"),
    ]
    table(doc,
          ["Item", "Description", "Reported"],
          [[i, desc, "Yes"] for i, desc in strobe_items],
          "STROBE checklist for observational cohort studies.")

    out = os.path.join(BASE, "ha_supplementary.docx")
    normalize_docx(doc)
    doc.save(out)
    print("wrote", out)


def build_figure_pptx():
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]

    main_figs = [
        ("ha_Figure_1.png", "Figure 1",
         f"Equivalence (TOST) of the litigation-rate effect against \u00b1{MARGIN1}% and "
         f"\u00b1{MARGIN2}% margins; horizontal bars are 90% confidence intervals."),
        ("ha_Figure_2.png", "Figure 2",
         "Biennial physician growth against lagged litigation exposure measured as (a) counts and (b) rates. Points are coloured by specialty; the count panel shows the size confounding that the rate panel removes."),
        ("ha_Figure_3.png", "Figure 3",
         "Biennial hospital facility-count growth against lagged litigation exposure measured as (a) counts and (b) rates. Points are coloured by specialty; the rate-adjusted panel shows no systematic association."),
        ("ha_Figure_4.png", "Figure 4",
         "Counterfactual policy-instrument simulation: marginal 10-year change in physician counts by specialty relative to the projected baseline drift. The MDE benchmark is the minimum detectable per-SD effect from the primary analysis."),
    ]
    supp_figs = [
        ("ha_Supplementary_Figure_1.png", "Supplementary Figure 1",
         "Sensitivity-analysis framework for evaluating malpractice-litigation risk as a healthcare workforce-allocation instrument."),
        ("ha_Supplementary_Figure_2.png", "Supplementary Figure 2",
         f"Closed malpractice claims per {_per} physicians by specialty, 2008\u20132024 (rates, not counts)."),
        ("ha_Supplementary_Figure_3.png", "Supplementary Figure 3",
         "Physician workforce by specialty, indexed to 2008 (=100)."),
    ]

    def add_slide(prs, fn, num, cap):
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(PInches(0.5), PInches(0.2), PInches(12.3), PInches(0.7))
        tf = tb.text_frame
        tf.text = num
        tf.paragraphs[0].runs[0].font.size = PPt(24)
        tf.paragraphs[0].runs[0].font.bold = True
        img = os.path.join(OUT, fn)
        if os.path.exists(img):
            s.shapes.add_picture(img, PInches(1.2), PInches(1.1), height=PInches(5.2))
        cb = s.shapes.add_textbox(PInches(0.5), PInches(6.5), PInches(12.3), PInches(0.9))
        cf = cb.text_frame
        cf.word_wrap = True
        cf.text = cap
        cf.paragraphs[0].runs[0].font.size = PPt(14)

    for fn, num, cap in main_figs:
        add_slide(prs, fn, num, cap)
    out = os.path.join(BASE, "ha_figures.pptx")
    normalize_pptx(prs)
    prs.save(out)
    print("wrote", out)

    prs2 = Presentation()
    prs2.slide_width = PInches(13.333)
    prs2.slide_height = PInches(7.5)
    for fn, num, cap in supp_figs:
        add_slide(prs2, fn, num, cap)
    out2 = os.path.join(BASE, "ha_supplementary_figures.pptx")
    normalize_pptx(prs2)
    prs2.save(out2)
    print("wrote", out2)


def prepare_figures():
    """Verify that Healthcare Analytics figures have been generated."""
    required = [
        "ha_Figure_1.png", "ha_Figure_2.png", "ha_Figure_3.png", "ha_Figure_4.png",
        "ha_Supplementary_Figure_1.png", "ha_Supplementary_Figure_2.png",
        "ha_Supplementary_Figure_3.png",
    ]
    for fn in required:
        path = os.path.join(OUT, fn)
        if not os.path.exists(path):
            raise SystemExit(f"missing figure: {path}; run manuscript/build_figures_en.py")
        print("verified", path)


def create_submission_zip():
    """Bundle all generated Healthcare Analytics submission files into one archive."""
    zip_path = os.path.join(OUT, "ha_submission.zip")
    files = [
        os.path.join(BASE, "ha_manuscript_en.docx"),
        os.path.join(BASE, "ha_title_page.docx"),
        os.path.join(BASE, "ha_cover_letter.docx"),
        os.path.join(BASE, "ha_highlights.docx"),
        os.path.join(BASE, "ha_supplementary.docx"),
        os.path.join(BASE, "ha_figures.pptx"),
        os.path.join(BASE, "ha_supplementary_figures.pptx"),
        os.path.join(OUT, "ha_Figure_1.png"),
        os.path.join(OUT, "ha_Figure_2.png"),
        os.path.join(OUT, "ha_Figure_3.png"),
        os.path.join(OUT, "ha_Figure_4.png"),
        os.path.join(OUT, "ha_Supplementary_Figure_1.png"),
        os.path.join(OUT, "ha_Supplementary_Figure_2.png"),
        os.path.join(OUT, "ha_Supplementary_Figure_3.png"),
    ]
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as z:
        for path in files:
            if not os.path.exists(path):
                raise SystemExit(f"submission zip missing file: {path}")
            z.write(path, arcname=os.path.basename(path))
    print("wrote", zip_path)


def main():
    prepare_figures()
    main_wc, abs_wc = build_manuscript()
    build_title_page(main_wc)
    build_highlights()
    build_cover_letter()
    build_supplementary()
    build_figure_pptx()
    create_submission_zip()


if __name__ == "__main__":
    main()
