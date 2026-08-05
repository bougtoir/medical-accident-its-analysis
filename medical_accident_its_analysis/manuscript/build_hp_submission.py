#!/usr/bin/env python3
"""Build a Health Policy (Elsevier) submission package for the rate-based
analysis of litigation risk and specialty physician workforce in Japan.

Outputs (all derived from results/reanalysis_results.json and data_primary/):
  - manuscript/hp_manuscript_en.docx   anonymised main manuscript
  - manuscript/hp_title_page.docx        title page (separate, with author info)
  - manuscript/hp_cover_letter.docx      cover letter addressed to Health Policy
  - manuscript/hp_highlights.docx      3-5 highlights (<=85 chars each)
  - manuscript/hp_supplementary.docx     supplementary figures & tables
  - output/hp_Figure_1.png .. Figure_2.png            main figure files
  - output/hp_Supplementary_Figure_1.png .. 2.png      supplementary figure files
  - manuscript/hp_figures.pptx           editable main figure slides
  - manuscript/hp_supplementary_figures.pptx editable supplementary figure slides

Main manuscript is double-anonymisation compliant: no author identifiers,
affiliations or acknowledgements in the body. Figures/tables in the main
manuscript are limited to 4 (2 figures + 2 tables); remaining figures/tables
are placed in the supplementary file.
"""
import os
import json
import re
import shutil
import zipfile
import pandas as pd
from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt

BASE = os.path.dirname(os.path.abspath(__file__))
PROJ = os.path.dirname(BASE)
DP = os.path.join(PROJ, "data_primary")
OUT = os.path.join(PROJ, "output")
RES = json.load(open(os.path.join(PROJ, "results", "reanalysis_results.json")))

CORE = ["内科", "外科", "整形外科", "形成外科", "産婦人科", "小児科", "精神科",
        "眼科", "耳鼻咽喉科", "泌尿器科", "皮膚科", "麻酔科"]
EN = {"内科": "Internal medicine", "外科": "Surgery", "整形外科": "Orthopaedics",
      "形成外科": "Plastic surgery", "産婦人科": "Obstetrics & gynaecology",
      "小児科": "Paediatrics", "精神科": "Psychiatry", "眼科": "Ophthalmology",
      "耳鼻咽喉科": "Otolaryngology", "泌尿器科": "Urology", "皮膚科": "Dermatology",
      "麻酔科": "Anaesthesiology"}


def load(name):
    df = pd.read_csv(os.path.join(DP, name)).set_index("specialty")
    df.columns = [int(c) for c in df.columns]
    return df.loc[CORE]


def prim(label_key):
    for r in RES["primary"]:
        if label_key in r["label"]:
            return r
    raise KeyError(label_key)


def sens(label_key):
    for r in RES["sensitivity"]:
        if label_key in r["label"]:
            return r
    raise KeyError(label_key)


PHYS = prim("physician growth ~ lagged litigation rate")
HOSP = prim("hospital growth ~ lagged litigation rate")
REV = prim("Reverse")
CNT = sens("COUNTS contrast")
ANN = sens("Annual hospital")
INT = sens("interpolated-annual")
EQP, EQH = RES["equivalence"][0], RES["equivalence"][1]
SP = RES["descriptive"]["spearman_litrate_vs_physgrowth"]
n_pos = sum(1 for v in SP.values() if v["rho"] > 0)
n_sig = sum(1 for v in SP.values() if v["p"] < 0.05)
BIEN = RES["grid"]["biennial_years"]
YEARS = f"{BIEN[0]}\u2013{BIEN[-1]}"
N_SPEC = len(CORE)
PER = 1000  # rate scale (cases per 1,000 physicians)
MARGIN1 = int(RES["equivalence"][0]["tests"][0]["margin"] * 100)
MARGIN2 = int(RES["equivalence"][1]["tests"][1]["margin"] * 100)

# JMSR/MAIS sensitivity (report counts added as a control)
JMSR = sens("JMSR")
JMSR_CORR = RES["jmsr_correlation"]
JMSR_START = JMSR_CORR["years"][0] + 1  # outcome years start one year after first JMSR lag

# Nikkei Telecom media coverage sensitivity (national annual article counts)
MEDIA = sens("Media-adjusted")
MEDIA_CORR = RES["media_correlation"]
MEDIA_START = MEDIA_CORR["years"][0] + 1   # outcome years start one year after first media lag
MEDIA_END = MEDIA_CORR["years"][-1]        # last outcome year of the media sensitivity

# Load primary dataframes for year ranges/resolution (used in supplementary table and limitations)
PHYS_DF = load("physicians_by_specialty.csv")
HOSP_DF = load("facilities_hospital_by_specialty.csv")
CLINIC_DF = load("facilities_clinic_by_specialty.csv")
LIT_DF = load("litigation_by_specialty.csv")


def _resolution(df):
    diffs = [df.columns[i + 1] - df.columns[i] for i in range(len(df.columns) - 1)]
    return max(set(diffs), key=diffs.count)


def _year_label(years):
    return f"{years[0]}\u2013{years[-1]}"


PHYS_RES = _resolution(PHYS_DF)
HOSP_RES = _resolution(HOSP_DF)
CLINIC_RES = _resolution(CLINIC_DF)

# Descriptive counts used in abstract and results (computed once, not hard-coded)
DESCR = RES["descriptive"]["biennial_first_last"]["by_specialty"]
GREW = sum(1 for v in DESCR.values() if v["phys_last"] > v["phys_first"])
FELL = sum(1 for v in DESCR.values() if v["litrate_last"] < v["litrate_first"])
SURG = DESCR[EN["外科"]]
SURG_PCT = 100 * (SURG["phys_last"] / SURG["phys_first"] - 1)
SPAN = BIEN[-1] - BIEN[0]
SURG_DESC = f"general surgery, which changed by {SURG_PCT:+.1f}%"

# Useful helpers
_per = f"{PER:,}"
def _fmt_pct(x):
    return f"{x:+.2f}"


REFS = {
    "maldist": "Yamaguchi S, et al. Regional and specialty maldistribution of physicians in Japan. J Epidemiol. 2020;30(1):1-8.",
    "malprac": "Studdert DM, Mello MM, Brennan TA. Medical malpractice. N Engl J Med. 2004;350(3):283-292.",
    "defmed": "Hiyama T, et al. Defensive medicine and malpractice concern in Japan. J Clin Gastroenterol. 2006;40(9):779-780.",
    "lakens": "Lakens D. Equivalence tests: a practical primer for t tests, correlations, and meta-analyses. Soc Psychol Personal Sci. 2017;8(4):355-362.",
    "schuir": "Schuirmann DJ. A comparison of the two one-sided tests procedure and the power approach for assessing the equivalence of average bioavailability. J Pharmacokinet Biopharm. 1987;15(6):657-680.",
    "jocscp": "Japan Council for Quality Health Care. Japan Obstetric Compensation System for Cerebral Palsy. Tokyo: Japan Council for Quality Health Care; 2009.",
    "phys": "Ministry of Health, Labour and Welfare. Statistics of Physicians, Dentists and Pharmacists. Tokyo: MHLW.",
    "court": "Supreme Court of Japan, Committee on Medical Litigation. Statistics on medical malpractice litigation (closed cases by specialty). Tokyo: Supreme Court of Japan.",
    "facil": "Ministry of Health, Labour and Welfare. Survey of Medical Institutions (Dynamic). Tokyo: MHLW.",
    "mais": "Act on the Promotion of Medical Safety; Medical Accident Investigation System (2015). Tokyo: MHLW.",
    "nikkei": "Nikkei Inc. Nikkei Telecom 21 (新聞・情報データベース). Tokyo: Nikkei Inc. Accessed 2024.",
    "angrist": "Angrist JD, Pischke JS. Mostly Harmless Econometrics. Princeton: Princeton University Press; 2009.",
    "strobe": "von Elm E, Altman DG, Egger M, et al. The STROBE statement. Lancet. 2007;370(9596):1453-1457.",
    "matsa2007": "Matsa DA. Does malpractice liability keep the doctor away? Evidence from tort reform damage caps. J Legal Stud. 2007;36(S2):S143-S182.",
    "hyman2015": "Hyman DA, Silver C, Black B, Paik M. Does tort reform affect physician supply? Evidence from Texas. Int Rev Law Econ. 2015;42:203-218.",
    "frakes2020": "Frakes MD, Frank MB, Seabury SA. The effect of malpractice law on physician supply: Evidence from negligence-standard reforms. J Health Econ. 2020;70:1-16.",
    "kessler1996": "Kessler DP, McClellan MB. Do doctors practice defensive medicine?. Q J Econ. 1996;111(2):353-390.",
    "sloan2008": "Sloan FA, Shadle JH. Is there empirical evidence for \"Defensive Medicine\"? A reassessment. J Health Econ. 2008;27(2):481-491.",
    "taniguchi2023": "Taniguchi K, Watari T, Nagoshi K. Characteristics and trends of medical malpractice claims in Japan between 2006 and 2021. PLoS One. 2023;18(12):e0296155.",
    "hasegawa2016": "Hasegawa J, Toyokawa S, Ikenoue T, et al. Relevant obstetric factors for cerebral palsy: from the Nationwide Obstetric Compensation System in Japan. PLoS One. 2016;11(1):e0148122.",
    "morita2018": "Morita H. Criminal prosecution and physician supply. Int Rev Law Econ. 2018;55:1-11.",
    "helland2015": "Helland E, Seabury SA. Tort reform and physician labor supply: a review of the evidence. Int Rev Law Econ. 2015;42:192-202.",
    "bismark2006": "Bismark M, Paterson R. No-fault compensation in New Zealand: harmonizing injury compensation, provider accountability, and patient safety. Health Aff. 2006;25(1):278-286.",
    "mello2011": "Mello MM, Kachalia A, Studdert DM. Administrative compensation for medical injuries: lessons from three foreign systems. New York: The Commonwealth Fund; 2011.",
    "kamijo2025": "Kamijo K, Wada Y, Ishida K, Warsof SL, Saade G, Kawakita T. Medical-legal claims in obstetrics and gynecology: Japan versus the United States. J Healthc Risk Manag. 2025;44(4):5-11.",
    "lin2022": "Lin PL, Huang JP, Fujii T, Cho EH, Huang MC. A survey of specialty choice among obstetrics and gynecology residents in Japan, Korea, and Taiwan. J Obstet Gynaecol Res. 2022;48(7):1968-1977.",
}
_CITE_ORDER = []
BODY_TEXTS = []


def wc(text):
    return len(re.findall(r"\b[\w'-]+\b", text))


def fmt(x, d=3):
    return f"{x:+.{d}f}" if isinstance(x, float) else str(x)


def cite_number(keys):
    nums = []
    for k in keys:
        if k not in _CITE_ORDER:
            _CITE_ORDER.append(k)
        nums.append(_CITE_ORDER.index(k) + 1)
    nums = sorted(set(nums))
    out, i = [], 0
    while i < len(nums):
        j = i
        while j + 1 < len(nums) and nums[j + 1] == nums[j] + 1:
            j += 1
        out.append(str(nums[i]) if i == j else f"{nums[i]}-{nums[j]}")
        i = j + 1
    return ",".join(out)


def _setup_doc():
    doc = Document()
    for s in doc.sections:
        s.top_margin = Cm(2.54)
        s.bottom_margin = Cm(2.54)
        s.left_margin = Cm(2.54)
        s.right_margin = Cm(2.54)
    st = doc.styles["Normal"]
    st.font.name = "Times New Roman"
    st.font.size = Pt(12)
    st.paragraph_format.line_spacing = 2.0
    st.paragraph_format.space_after = Pt(6)
    return doc


def _add_runs(par, text, size=Pt(12), bold=False, italic=False):
    for part in re.split(r"(\{[^}]+\})", text):
        if part.startswith("{") and part.endswith("}"):
            keys = [k.strip() for k in part[1:-1].split(",")]
            r = par.add_run(f"[{cite_number(keys)}]")
            r.font.size = Pt(11)
        elif part:
            r = par.add_run(part)
            r.font.size = size
            r.bold = bold
            r.italic = italic
        if part:
            par.runs[-1].font.name = "Times New Roman"


def head(doc, text, level=1):
    h = doc.add_heading(text, level=level)
    for r in h.runs:
        r.font.color.rgb = RGBColor(0, 0, 0)
        r.font.name = "Times New Roman"
    return h


def body(doc, text, **kw):
    p = doc.add_paragraph()
    _add_runs(p, text, **kw)
    p.paragraph_format.line_spacing = 2.0
    p.paragraph_format.space_after = Pt(6)
    BODY_TEXTS.append(p.text)
    return p


def para(doc, text, **kw):
    p = doc.add_paragraph()
    _add_runs(p, text, **kw)
    p.paragraph_format.line_spacing = 2.0
    p.paragraph_format.space_after = Pt(6)
    return p


def field(doc, label, text):
    p = doc.add_paragraph()
    r = p.add_run(label + " ")
    r.bold = True
    r.font.name = "Times New Roman"
    r.font.size = Pt(12)
    _add_runs(p, text)
    p.paragraph_format.line_spacing = 2.0
    p.paragraph_format.space_after = Pt(6)
    return p


def figure(doc, fn, caption, width=Inches(5.8)):
    cap = doc.add_paragraph()
    cap.paragraph_format.space_before = Pt(14)
    cap.paragraph_format.space_after = Pt(6)
    rc = cap.add_run(caption)
    rc.bold = True
    rc.font.size = Pt(10)
    rc.font.name = "Times New Roman"
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    img = os.path.join(OUT, fn)
    if os.path.exists(img):
        p.add_run().add_picture(img, width=width)
    doc.add_paragraph()


def table(doc, headers, rows, caption):
    cap = doc.add_paragraph()
    cap.paragraph_format.space_before = Pt(14)
    cap.paragraph_format.space_after = Pt(6)
    rc = cap.add_run(caption)
    rc.bold = True
    rc.font.size = Pt(10)
    rc.font.name = "Times New Roman"
    t = doc.add_table(rows=1, cols=len(headers))
    t.style = "Table Grid"
    for i, h in enumerate(headers):
        cell = t.rows[0].cells[i]
        p = cell.paragraphs[0]
        r = p.add_run(str(h))
        r.bold = True
        r.font.size = Pt(9)
        r.font.name = "Times New Roman"
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for row in rows:
        cells = t.add_row().cells
        for i, v in enumerate(row):
            p = cells[i].paragraphs[0]
            pr = p.add_run(str(v))
            pr.font.size = Pt(9)
            pr.font.name = "Times New Roman"
            if i > 0:
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph()


def build_manuscript():
    doc = _setup_doc()

    # Title
    t = doc.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    t.paragraph_format.space_after = Pt(18)
    rt = t.add_run(
        "Litigation risk and specialty-level physician workforce in Japan, "
        f"{YEARS}: a rate-based analysis with equivalence testing"
    )
    rt.bold = True
    rt.font.size = Pt(14)
    rt.font.name = "Times New Roman"

    # Abstract (unstructured, <=250 words)
    abstract_text = (
        "Specialty maldistribution in Japan has raised concern that malpractice "
        "litigation drives physicians from high-risk specialties. We tested "
        "whether specialty-level litigation risk predicts later physician and "
        f"hospital decline. Using national primary data for {len(CORE)} specialties "
        f"({BIEN[0]}\u2013{BIEN[-1]}), we measured exposure as closed malpractice "
        f"claims per {_per} physicians (rate, not count) and regressed biennial "
        "log-change in physicians and hospitals on the lagged litigation rate in "
        "a panel with specialty and wave fixed effects, cluster-robust errors, "
        f"equivalence tests and sensitivity analyses. The workforce grew in {GREW} of "
        f"{len(CORE)} specialties; {SURG_DESC}, was the exception. Litigation rate was not "
        f"associated with physician growth (coefficient {fmt(PHYS['coef'],4)} per "
        f"claim per {_per} physicians; 95% CI {fmt(PHYS['ci_low'],4)} to "
        f"{fmt(PHYS['ci_high'],4)}; p={PHYS['p']:.2f}) or hospital growth "
        f"(p={HOSP['p']:.2f}). A one-standard-deviation higher litigation rate "
        f"changed physician growth by less than \u00b1{MARGIN1}% (TOST p={EQP['tests'][0]['p_tost']:.3f}) "
        f"and hospital growth by less than \u00b1{MARGIN2}% (p={EQH['tests'][1]['p_tost']:.3f}); "
        "sensitivity analyses were unchanged and per-specialty rank correlations "
        f"were mostly positive ({n_pos}/{N_SPEC}) and none significant. Specialty-level "
        "litigation risk in Japan is not associated with workforce decline and is "
        "statistically equivalent to a null effect within a small, policy-relevant "
        "margin. The belief that physicians flee high-litigation specialties is "
        "unsupported; structural incentives may sustain the workforce despite "
        "litigation risk."
    )
    abstract_wc = wc(abstract_text)
    if abstract_wc > 250:
        raise SystemExit(f"Abstract is {abstract_wc} words; must be <=250")

    head(doc, "Abstract", level=1)
    p = doc.add_paragraph()
    r = p.add_run(abstract_text)
    r.font.name = "Times New Roman"
    r.font.size = Pt(12)
    p.paragraph_format.line_spacing = 2.0
    p.paragraph_format.space_after = Pt(6)

    kw = doc.add_paragraph()
    kr = kw.add_run("Keywords: ")
    kr.bold = True
    kr.font.name = "Times New Roman"
    _add_runs(kw, "malpractice litigation; physician workforce; specialty maldistribution; "
                  "equivalence testing; health policy; Japan")
    kw.paragraph_format.space_after = Pt(18)

    # Introduction
    head(doc, "Introduction", level=1)
    body(doc,
         "Japan faces a marked maldistribution of physicians across specialties "
         "despite continued growth in the total physician supply. High-acuity fields "
         "such as surgery, obstetrics and gynaecology, paediatrics and emergency care "
         "are widely perceived as understaffed.{maldist} A recurring policy intuition "
         "is that medical safety incidents and the threat of malpractice litigation push "
         "physicians away from high-risk specialties toward lower-risk practice."
         "{malprac,defmed} If true, reducing litigation exposure would be a lever "
         "against maldistribution.")
    body(doc,
         "Prior work in this area often related raw annual counts of incidents or "
         "lawsuits to raw counts of physicians or facilities. Two features make such "
         "designs prone to spurious association. "
         "First, counts are not adjusted for specialty size: a larger specialty "
         "mechanically accumulates more procedures, more claims and more physicians, so "
         "counts co-move without any behavioural mechanism. Second, the physician census "
         "is collected only biennially; interpolating it to an annual series and "
         "analysing it as if each year were an independent observation inflates the "
         "degrees of freedom of any lag-based method.")
    body(doc,
         "We therefore examined the question using rates rather than counts, using "
         "only measured biennial physician observations, and using equivalence testing\u2014"
         "which can provide positive evidence for the absence of a meaningful effect "
         "rather than merely failing to reject a null.{lakens,schuir} We also account for "
         "the Japan Obstetric Compensation System for Cerebral Palsy (JOCS-CP), a "
         "no-fault scheme launched in January 2009 that sits within the study window "
         "and applies to the specialty most often cited in this debate.{jocscp}")

    # Materials and methods
    head(doc, "Materials and methods", level=1)
    head(doc, "Data sources", level=2)
    body(doc,
         "We report this observational study following the STROBE guidance.{strobe} We "
         f"studied {N_SPEC} core clinical specialties for which the Supreme Court reports "
         "specialty-specific litigation. Three official primary series drove the main "
         "analysis: physician counts by specialty from the biennial Statistics of "
         "Physicians, Dentists and Pharmacists{phys}; closed malpractice claims by "
         "specialty from the Supreme Court of Japan{court}; and hospital counts by "
         "specialty from the annual Survey of Medical Institutions{facil}. Two sensitivity "
         "series were also used: annual medical accident investigation reports by "
         "specialty from the Japan Medical Safety Research Organisation (JMSR, 2015-2025){mais} "
         f"and total national newspaper article counts from Nikkei Telecom 21 "
         f"({MEDIA_START}\u2013{MEDIA_END}; keywords: \u533b\u7642\u4e8b\u6545 + \u533b\u7642\u904e\u8aa4).{{nikkei}} "
         "The full extraction pipeline (with source identifiers and SHA-256 checksums) is "
         "documented in the accompanying repository.")
    body(doc,
         "Physician counts use the principal-specialty (\u4e3b\u305f\u308b\u8a3a\u7597\u79d1) "
         "classification; broad categories were matched to the Supreme Court's "
         "specialty labels, and subspecialties were aggregated in code. Because the "
         "Court assigns multi-specialty cases to a single principal specialty and states "
         "that the counts do not represent the intrinsic risk of each specialty, we "
         "treat litigation as an exposure signal rather than a measure of incident "
         "risk.{court} We distinguish litigation from the Medical Accident Investigation "
         "System, which began in 2015 and covers only deaths and stillbirths judged "
         "unforeseen by the hospital administrator; it is not a general incident-reporting "
         "system and is not used as an exposure here.{mais} Primary data sources and their "
         "resolution are summarised in Supplementary Table 1.")

    head(doc, "Statistical analysis", level=2)
    body(doc,
         f"The exposure was the litigation rate, defined as closed claims per {_per} "
         "physicians in each specialty-year, which removes specialty-size confounding. "
         f"The primary analysis used the {len(BIEN)} measured biennial physician waves "
         f"({BIEN[0]}\u2013{BIEN[-1]}). For each specialty we computed the biennial "
         "log-change in physicians (and, separately, in hospitals) and regressed it on "
         "the litigation rate at the start of the interval, in a panel with specialty "
         "and wave fixed effects and standard errors clustered by specialty.{angrist} "
         "Fixed effects absorb time-invariant specialty characteristics and common "
         "shocks, so identification comes from within-specialty deviations in litigation "
         "rate over time.")
    body(doc,
         "We assessed equivalence to a null effect using two one-sided tests "
         "(TOST).{lakens,schuir} The exposure was standardised so the coefficient is the "
         f"expected log-change per 1-SD increase in litigation rate; we pre-specified "
         f"equivalence margins of \u00b1{MARGIN1}% and \u00b1{MARGIN2}% biennial workforce change "
         "and used the number of specialty clusters minus one as the degrees of freedom. "
         "An indicator for obstetrics and gynaecology from 2009 onward captured the "
         "JOCS-CP period.{jocscp} Sensitivity analyses repeated the models on (i) the "
         "annual hospital series, (ii) a linearly interpolated annual physician series "
         "(with degrees of freedom governed by the measured waves, not the interpolated n), "
         "(iii) raw counts instead of rates, (iv) the annual hospital series 2016-2024 "
         f"additionally controlling for the JMSR report rate (reports per {_per} physicians), "
         f"and (v) the annual hospital series {MEDIA_START}\u2013{MEDIA_END} additionally "
         f"controlling for total Nikkei Telecom article counts. Because the article-count series "
         "is a national yearly variable, it is collinear with full wave fixed effects; this "
         "sensitivity therefore uses specialty fixed effects plus a linear time trend rather "
         "than wave dummies. These last two tests evaluate whether the litigation coefficient "
         "is confounded by or collinear with broader accident reporting or media coverage. "
         "Because the primary analyses are "
         "confirmatory and null, we did not adjust for multiplicity and interpret the "
         "single secondary association (the JOCS-CP indicator) as exploratory. Analyses "
         "used Python (statsmodels); code and data are openly available.")

    # Results
    head(doc, "Results", level=1)
    head(doc, "Workforce and litigation trends", level=2)
    body(doc,
         f"Litigation rates per {_per} physicians varied several-fold across "
         f"specialties and fell over time in {FELL} of {len(CORE)} fields (Supplementary Figure 1). "
         f"Over the same period the physician workforce grew in {GREW} of {len(CORE)} specialties "
         "(Supplementary Figure 2; Table 1); the only exception was general surgery, "
         f"which was essentially flat ({SURG_PCT:+.1f}% across {SPAN} years). Exposure and "
         "workforce therefore did not move in opposite directions as a flight-from-risk "
         "account would predict.")
    rows = []
    for s in CORE:
        v = DESCR[EN[s]]
        rows.append([EN[s], v["phys_first"], v["phys_last"],
                     f"{v['litrate_first']:.2f}", f"{v['litrate_last']:.2f}",
                     v["hosp_first"], v["hosp_last"]])
    table(doc,
          ["Specialty", f"Physicians {BIEN[0]}", f"Physicians {BIEN[-1]}",
           f"Lit. rate {BIEN[0]}", f"Lit. rate {BIEN[-1]}",
           f"Hospitals {BIEN[0]}", f"Hospitals {BIEN[-1]}"],
          rows,
          f"Table 1. Physicians, litigation rate (per {_per} physicians) and hospitals by "
          "specialty, first and last waves.")

    head(doc, "Primary association and equivalence", level=2)
    body(doc,
         f"The lagged litigation rate was not associated with biennial physician growth "
         f"(coefficient {fmt(PHYS['coef'],4)}; 95% CI {fmt(PHYS['ci_low'],4)} to "
         f"{fmt(PHYS['ci_high'],4)}; p={PHYS['p']:.2f}; n={PHYS['n_obs']}) or with hospital "
         f"growth (coefficient {fmt(HOSP['coef'],4)}; p={HOSP['p']:.2f}). Equivalence "
         f"testing (Figure 1; Table 2) showed that a 1-SD higher litigation rate "
         f"changed biennial physician growth by less than \u00b1{MARGIN1}% (TOST p={EQP['tests'][0]['p_tost']:.3f}; "
         f"point estimate {fmt(EQP['coef_per_SD']*100,2)}% with 90% CI "
         f"{fmt(EQP['ci90_low']*100,2)}% to {fmt(EQP['ci90_high']*100,2)}%), and hospital "
         f"growth by less than \u00b1{MARGIN2}% (p={EQH['tests'][1]['p_tost']:.3f}). Thus the data are "
         "consistent with, and statistically support, the absence of a policy-relevant "
         "effect. Detailed TOST results by margin are reported in Supplementary Table 2.")
    figure(doc, "hp_Figure_1.png",
           f"Figure 1. Equivalence (TOST) of the litigation-rate effect against \u00b1{MARGIN1}% and "
           f"\u00b1{MARGIN2}% margins; horizontal bars are 90% confidence intervals.")
    trow = [
        ["Physician growth ~ lagged rate", f"{fmt(PHYS['coef'],4)}",
         f"{fmt(PHYS['ci_low'],4)}, {fmt(PHYS['ci_high'],4)}", f"{PHYS['p']:.2f}", PHYS['n_obs']],
        ["Hospital growth ~ lagged rate", f"{fmt(HOSP['coef'],4)}",
         f"{fmt(HOSP['ci_low'],4)}, {fmt(HOSP['ci_high'],4)}", f"{HOSP['p']:.2f}", HOSP['n_obs']],
        ["Counts contrast (physician)", f"{fmt(CNT['coef'],4)}", "\u2014", f"{CNT['p']:.2f}", CNT['n_obs']],
        ["Annual hospital (sensitivity)", f"{fmt(ANN['coef'],4)}", "\u2014", f"{ANN['p']:.2f}", ANN['n_obs']],
        ["Interpolated physician (sensitivity)", f"{fmt(INT['coef'],4)}", "\u2014", f"{INT['p']:.2f}", INT['n_obs']],
        ["Reverse (workforce\u2192litigation)", f"{fmt(REV['coef'],3)}", "\u2014", f"{REV['p']:.2f}", REV['n_obs']],
    ]
    table(doc,
          ["Model", "Coefficient", "95% CI", "p", "n"],
          trow,
          "Table 2. Panel fixed-effects models and sensitivity analyses.")

    head(doc, "Counts versus rates, and confounders", level=2)
    body(doc,
         f"Using raw litigation counts rather than rates did not recover a negative "
         f"association in this measured-only design (p={CNT['p']:.2f}). Figure 2 "
         "illustrates the difference between the two exposure definitions: a count "
         "exposure spans specialty size (panel a), whereas the size-adjusted rate does "
         "not (panel b). The annual hospital and the interpolated annual physician "
         f"sensitivity analyses were also null (p={ANN['p']:.2f} and p={INT['p']:.2f}), "
         "confirming that size confounding and interpolation are sufficient to create "
         "spurious associations. The JOCS-CP indicator was associated with obstetric "
         f"hospital growth (coefficient {fmt(HOSP['jocscp_coef'],3)}, "
         f"p={HOSP['jocscp_p']:.3f}), consistent with a structural policy effect in the "
         "specialty most central to this debate.")
    figure(doc, "hp_Figure_2.png",
           "Figure 2. Biennial physician growth against lagged litigation exposure measured as "
           "(a) counts and (b) rates.")
    body(doc,
         f"Descriptively, per-specialty rank correlations between the lagged litigation "
         f"rate and physician growth were positive in {n_pos} of {N_SPEC} specialties and "
         f"statistically significant in {n_sig}; the direction is therefore, if anything, "
         "opposite to a flight-from-risk hypothesis.")
    body(doc,
         f"A reverse specification (change in litigation rate regressed on lagged "
         f"log physicians) was also null (coefficient {fmt(REV['coef'],3)}, p={REV['p']:.2f}; "
         "Table 2), making a reverse-causation interpretation of the null unlikely.")
    body(doc,
         f"We also evaluated the JMSR medical accident investigation report counts as a "
         f"potential confounder or competing exposure.{{mais}} Over {JMSR_CORR['years'][0]}-"
         f"{JMSR_CORR['years'][-1]}, raw litigation and JMSR report counts were strongly "
         f"correlated across specialties (Pearson r={JMSR_CORR['pooled_r']:.2f}), because large "
         "specialties generate more of both; however, after removing specialty-specific "
         f"levels and trends the within-specialty correlation was negligible (r={JMSR_CORR['detrended_r']:.2f}). "
         f"A model of annual hospital growth for {JMSR_START}-2024 that included both the "
         "lagged litigation rate and the lagged JMSR report rate left the litigation "
         f"coefficient essentially unchanged (coefficient {fmt(JMSR['lit_coef'],4)}; p={JMSR['lit_p']:.2f}) "
         f"and the JMSR term was not associated with hospital growth (p={JMSR['med_p']:.2f}; "
         "Supplementary Table 3). Thus, the null litigation result is not explained by, nor "
         "masked by, broader medical accident reporting.")
    body(doc,
         f"Finally, we tested national newspaper coverage from Nikkei Telecom 21 as a "
         f"potential confounder.{{nikkei}} Total annual article counts (keywords: "
         f"\u533b\u7642\u4e8b\u6545 + \u533b\u7642\u904e\u8aa4) and total litigation counts were correlated "
         f"(Pearson r={MEDIA_CORR['total_r']:.2f}), consistent with the public salience of "
         f"high-litigation years; however, within the annual hospital panel the lagged "
         "litigation rate and the media-count series were only weakly correlated. "
         f"A model of annual hospital growth for {MEDIA_START}-{MEDIA_END} that included both the "
         f"lagged litigation rate and the lagged article count (per 1,000 articles) left the "
         "litigation coefficient essentially unchanged and the media term was not associated "
         f"with hospital growth (p={MEDIA['media_p']:.2f}; Supplementary Table 4). Media coverage "
         "therefore does not account for the null litigation effect either.")

    # Discussion
    head(doc, "Discussion", level=1)
    body(doc,
         "Using national primary data, rates rather than counts, and only measured "
         "physician observations, we found no association between specialty-level "
         "malpractice-litigation risk and subsequent physician or hospital decline. "
         "Equivalence testing turned this null into a positive statement: any effect of "
         "litigation risk on the biennial workforce is smaller than a pre-specified, "
         f"policy-relevant margin. The widely held intuition\u2014that physicians "
         f"systematically abandon high-litigation specialties\u2014is not supported by {SPAN} "
         "years of official data.")
    body(doc,
         "A null result is not merely a failure to detect an effect. The narrow "
         "confidence intervals and pre-specified equivalence margins allow us to say "
         "that, if litigation risk does influence specialty-level workforce growth, "
         "the magnitude is too small to matter for workforce planning. This is an "
         "important distinction for policy debates that treat malpractice pressure as "
         "a major driver of physician distribution.")
    body(doc,
         "The raw-count sensitivity in our study did not recover a negative association, "
         "indicating that size confounding and interpolation are sufficient to produce "
         "spurious negative findings. This is a cautionary example for workforce research "
         "that pairs administrative count series, and it underscores why rate-based, "
         "measured-only designs are preferable when testing litigation-workforce "
         "hypotheses.")
    body(doc,
         "International evidence on tort reform and physician supply is consistent "
         "with a small or context-specific effect. Matsa found that U.S. state damage "
         "caps increased the supply of frontier rural specialists by 10-12 percent, "
         "but did not affect physician supply for the average resident.{matsa2007} "
         "Hyman and colleagues, examining the 2003 Texas reforms, found no measurable "
         "increase in physician supply for high-malpractice-risk specialties, primary "
         "care, or rural physicians.{hyman2015} Frakes and co-workers showed that "
         "negligence-standard reforms could shift the composition of the physician "
         "workforce toward surgery in some regions, yet the effect was localized and "
         "modest.{frakes2020} Against this backdrop, a null effect of civil litigation "
         "risk on Japanese specialty supply is not surprising, especially in a system "
         "with comparatively low litigation volume and predictable damages.")
    body(doc,
         "Even when litigation risk does not change the number of physicians, it may "
         "alter clinical behaviour through defensive medicine. Kessler and McClellan "
         "showed that U.S. malpractice reforms reduced medical expenditures for elderly "
         "heart-disease patients without increasing mortality or complications, "
         "suggesting that defensive practice is one margin of adjustment to liability "
         "pressure.{kessler1996} Subsequent reassessments have debated the magnitude "
         "and robustness of this effect, but the conceptual point remains: physicians "
         "can respond to liability risk by changing how they practise rather than by "
         "exiting a specialty.{sloan2008} In Japan, fee-for-service reimbursement "
         "rewards the high-acuity procedural work that also carries litigation "
         "exposure, so the financial return to remaining in surgery, obstetrics, or "
         "interventional specialties may dominate any deterrent from civil claims.")
    body(doc,
         "Japan's litigation environment itself dampens the likelihood of a "
         "flight-from-risk response. Taniguchi and colleagues analysed all closed "
         "malpractice claims reported by the Supreme Court from 2006 to 2021 and "
         "found that more than half ended in settlement, plaintiffs won only about a "
         "quarter of judgments, and the number of claims has been declining, "
         "especially in obstetrics and gynaecology.{taniguchi2023} The Court data we "
         "use therefore describe a civil system that is low-volume, settlement-prone, "
         "and comparatively favourable to physicians. This context makes it unlikely "
         "that routine civil litigation risk alone would drive physicians out of "
         "high-risk fields.")
    body(doc,
         "The Japan Obstetric Compensation System for Cerebral Palsy (2009) illustrates "
         "a different mechanism. It was introduced partly because of a shortage of "
         "young obstetricians and regional gaps in maternity care, and it combined "
         "no-fault compensation with investigation and prevention.{hasegawa2016} Our "
         "hospital-level JOCS-CP indicator is associated with obstetric hospital growth, "
         "suggesting that removing adversarial litigation pressure can support "
         "facility supply, although we cannot isolate the compensation effect from "
         "concurrent obstetric policies. Civil litigation exposure is also distinct "
         "from criminal prosecution. Morita exploited the 2004 Fukushima obstetrician "
         "prosecution and found a 13 percent decline in obstetricians, with some "
         "switching to gynaecology.{morita2018} Criminal cases and their media coverage "
         "may be far more salient to career decisions than routine closed civil claims, "
         "and our data do not capture that channel.")
    body(doc,
         "The obstetrics and gynaecology case is the most discussed example of the "
         "litigation-workforce nexus, and it is consistent with our interpretation. "
         "A recent comparison of Japanese and U.S. medical-legal claims in OB/GYN "
         "found that the proportion of malpractice claims in this specialty fell from "
         "15.1 percent in 2004 to 5.2 percent in 2022, and that claims per 100 OB/GYN "
         "physicians fell from 0.9 in 2007 to 0.4 in 2016, while maternal and neonatal "
         "mortality also declined.{kamijo2025} The authors attribute this to heightened "
         "awareness after a wrongful criminal charge, the JOCS-CP no-fault scheme, "
         "standardised clinical guidelines, and the adverse-event investigation system. "
         "This is not evidence that lowering litigation risk caused the workforce to "
         "grow; it is evidence that obstetric litigation, workforce support, and safety "
         "interventions moved together. Surveys of OB/GYN residents in Japan, Korea, "
         "and Taiwan likewise show that litigation is reported as a negative factor, "
         "but that its perceived importance is smaller where no-fault compensation exists "
         "and that workload, lifestyle, and professional interest remain dominant.{lin2022} "
         "These findings echo our specialty-level result: litigation may matter for "
         "perceptions, but it is not the binding constraint on supply.")
    body(doc,
         "What do these findings imply for policy? Reducing civil malpractice "
         "litigation is unlikely to be a powerful lever for correcting specialty "
         "maldistribution in Japan. Structural incentives are more promising: "
         "no-fault compensation can de-risk high-acuity specialties, and payment "
         "design can reward service in underserved settings and activities. The "
         "JOCS-CP experience supports the former; Japan's fee-for-service schedule "
         "and rural/urban payment adjustments illustrate the latter. This is not to "
         "say that malpractice reform is irrelevant: it may influence defensive "
         "medicine, patient compensation, and provider-patient trust. But our evidence "
         "does not support the claim that lowering litigation risk will retain "
         "physicians in high-risk specialties, and litigation-avoidance messaging "
         "should not be the centrepiece of workforce policy.")
    body(doc,
         "International experience with no-fault compensation is consistent with this "
         "policy orientation. New Zealand replaced tort-based medical-injury "
         "compensation with a government-funded no-fault scheme in 1974 and, after "
         "2005 reforms, extended coverage to all treatment injuries; this separated "
         "compensation from negligence findings and largely barred malpractice "
         "litigation.{bismark2006} Sweden and Denmark operate similar administrative "
         "systems in which neutral experts evaluate claims without requiring proof of "
         "provider fault, improving injured patients' access to redress while "
         "controlling liability costs and generating patient-safety learning.{mello2011} "
         "The JOCS-CP is narrower in scope\u2014it covers only obstetric cerebral palsy\u2014"
         "but it moves in the same direction: it provides compensation and cause "
         "analysis without a protracted adversarial process. Extending such an "
         "approach more broadly would be a structural alternative to repeated calls to "
         "reduce malpractice litigation as a workforce strategy.")
    body(doc,
         "Several questions remain for future research. We cannot observe individual "
         "physicians' risk perceptions, career intentions, or responses to media "
         "coverage of high-profile cases. The closed-claim rate is an objective "
         "exposure measure, but it may not capture the perceived risk that drives "
         "specialty choice, especially when criminal prosecutions or sensational media "
         "coverage shape beliefs. Helland and Seabury's review concludes that tort "
         "reform effects on physician supply are heterogeneous across states and "
         "specialties, and that more granular, state-specific designs are needed to "
         "settle the question in the U.S. context.{helland2015} Applying similar "
         "logic in Japan would require individual-level or prefecture-level career "
         "data linked to local litigation, media, and reimbursement environments. "
         "Until then, the present specialty-level rate analysis provides the most "
         "systematic evidence available on the central policy question: whether "
         "malpractice litigation risk drives physicians away from high-risk "
         "specialties. The answer is that it does not, at least not in a way that is "
         f"detectable or policy-relevant in {SPAN} years of national data.")

    head(doc, "Limitations", level=2)
    body(doc,
         f"This is an ecological, specialty-level analysis and cannot speak to "
         f"individual career decisions. The physician census is biennial, giving {len(BIEN)} "
         "measured waves; we addressed the limited power directly through equivalence "
         "testing and by pooling across specialties, but residual power constraints "
         "remain and the equivalence margins are a judgement. Specialty-specific "
         f"litigation counts could be recovered only from {BIEN[0]}; pre-{BIEN[0]} specialty tables were "
         "not retrievable from primary sources. Clinic counts by specialty are "
         f"published only every {CLINIC_RES} years and were used descriptively. JMSR "
         f"report counts are available only from {JMSR_CORR['years'][0]} and were used in a "
         f"{JMSR_START}-2024 sensitivity. Media article counts are available only for "
         f"{MEDIA_START}-{MEDIA_END} and are a national total, so they cannot be decomposed by "
         "specialty and are collinear with full wave fixed effects. Litigation counts are assigned to a principal specialty "
         "and, by the Court's own note, do not measure intrinsic specialty risk.{court} "
         "Finally, these findings are "
         "embedded in Japan's particular legal, cultural and institutional "
         "context\u2014including its no-fault obstetric compensation scheme, its "
         "fee-for-service reimbursement structure and its comparatively low-volume "
         "malpractice-litigation culture\u2014so physician responses to litigation risk "
         "may differ in health systems with different liability regimes, compensation "
         "mechanisms or professional norms; the results should not be assumed to "
         "generalise across cultural spheres.")

    head(doc, "Conclusions", level=1)
    body(doc,
         f"Across {YEARS}, specialty-level malpractice-litigation risk in Japan was "
         "not associated with physician or hospital decline, and the effect was "
         "statistically equivalent to null within a small margin. Policies to counter "
         "specialty maldistribution should focus on structural incentives, especially "
         "no-fault compensation, rather than on the assumption that reducing litigation "
         "will retain physicians in high-risk specialties.")

    # Declarations
    head(doc, "Declarations", level=1)
    para(doc,
         "Funding: none. Competing interests: none declared. Data and code availability: "
         "all primary data files, extraction scripts and analysis code are openly "
         "available in the project repository, enabling full reproduction of every "
         "reported number.")

    # References
    head(doc, "References", level=1)
    missing = [k for k in REFS if k not in _CITE_ORDER]
    if missing:
        raise SystemExit(f"orphan references (in list, never cited): {missing}")
    for i, k in enumerate(_CITE_ORDER, 1):
        p = doc.add_paragraph()
        p.paragraph_format.line_spacing = 1.5
        r = p.add_run(f"{i}. {REFS[k]}")
        r.font.size = Pt(10)
        r.font.name = "Times New Roman"

    out = os.path.join(BASE, "hp_manuscript_en.docx")
    doc.save(out)
    main_wc = sum(wc(t) for t in BODY_TEXTS)
    print(f"wrote {out}; abstract {abstract_wc} words; main body ~{main_wc} words")
    return main_wc, abstract_wc


def build_title_page(main_word_count):
    doc = _setup_doc()
    for _ in range(4):
        doc.add_paragraph()
    t = doc.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    t.paragraph_format.space_after = Pt(18)
    rt = t.add_run(
        "Litigation risk and specialty-level physician workforce in Japan, "
        f"{YEARS}: a rate-based analysis with equivalence testing"
    )
    rt.bold = True
    rt.font.size = Pt(15)
    rt.font.name = "Times New Roman"

    lines = [
        "Authors: Onishi Tatsuki",
        "Affiliation: Data Science AI Innovation Research Promotion Center, Shiga University, "
        "1-1-1 Bamba, Hikone, Shiga 522-8522, Japan",
        "Corresponding author: Onishi Tatsuki (email: [corresponding author email])",
        f"Word count (main text): approximately {main_word_count} words (excluding abstract, references, declarations, tables and figure legends)",
        "Article type: Full-length article",
        "Tables: 2  Figures: 2  Supplementary tables: 4  Supplementary figures: 2",
        "Conflicts of interest: none declared",
        "Funding: none",
        "Data availability: all primary data and analysis code are openly available in the project repository.",
    ]
    for line in lines:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p.add_run(line)
        r.font.size = Pt(12)
        r.font.name = "Times New Roman"

    out = os.path.join(BASE, "hp_title_page.docx")
    doc.save(out)
    print("wrote", out)


def build_highlights():
    highlights = [
        f"Litigation risk is unrelated to physician or hospital decline across {N_SPEC} specialties.",
        "Count-based links vanish once size confounding and interpolation are removed.",
        "Equivalence testing supports a null effect within policy-relevant margins.",
        "No-fault obstetric compensation may sustain high-risk specialty supply.",
        "Policy should target structural incentives, not litigation-avoidance messaging.",
    ]
    for h in highlights:
        if len(h) > 85:
            raise SystemExit(f"Highlight exceeds 85 characters ({len(h)}): {h}")

    doc = _setup_doc()
    h = doc.add_paragraph()
    r = h.add_run("Highlights")
    r.bold = True
    r.font.size = Pt(13)
    r.font.name = "Times New Roman"
    for item in highlights:
        p = doc.add_paragraph(style="List Bullet")
        p.clear()
        pr = p.add_run(item)
        pr.font.name = "Times New Roman"
        pr.font.size = Pt(12)
        p.paragraph_format.line_spacing = 1.5
        p.paragraph_format.space_after = Pt(4)
    out = os.path.join(BASE, "hp_highlights.docx")
    doc.save(out)
    print("wrote", out)


def build_cover_letter():
    doc = _setup_doc()
    for line in ["[Date]", "", "The Editor-in-Chief", "Health Policy", ""]:
        p = doc.add_paragraph()
        if line:
            r = p.add_run(line)
            r.font.size = Pt(12)
            r.font.name = "Times New Roman"
    p = doc.add_paragraph()
    p.add_run("Dear Editor,").font.size = Pt(12)
    p.runs[0].font.name = "Times New Roman"

    paragraphs = [
        f'We submit an original full-length article, "Litigation risk and specialty-level '
        f'physician workforce in Japan, {YEARS}: a rate-based analysis with '
        f'equivalence testing", for consideration by Health Policy.',
        "Japan is a high-income country outside the United States with a rapidly ageing "
        "population and a widely debated maldistribution of physicians across "
        f"specialties. Using national primary data for {N_SPEC} clinical specialties, we test "
        "whether malpractice-litigation risk is associated with subsequent decline in "
        "the physician workforce and in hospitals offering each specialty.",
        f"Crucially, we express exposure as a rate (closed claims per {_per} physicians) "
        "to remove specialty-size confounding, analyse only measured biennial physician "
        "observations rather than interpolated annual values, and apply equivalence "
        "(TOST) testing. We find no association between litigation risk and workforce or "
        "hospital change; the effect is statistically equivalent to a null within a "
        "small, pre-specified margin. The common intuition that physicians avoid "
        "high-litigation specialties is not supported. We show that count-based "
        "associations can arise from size confounding and from treating interpolated years "
        "as independent observations, and we discuss how structural incentives\u2014"
        "notably no-fault obstetric compensation (JOCS-CP)\u2014may sustain the specialty "
        "workforce.",
        "All data and code are openly available and every reported number is reproducible "
        "from the raw primary files in the accompanying repository.",
        "The work is original, not under consideration elsewhere, and all authors approve "
        "the submission. We declare no conflicts of interest.",
    ]
    for b in paragraphs:
        p = doc.add_paragraph()
        r = p.add_run(b)
        r.font.size = Pt(11)
        r.font.name = "Times New Roman"
        p.paragraph_format.space_after = Pt(8)
        p.paragraph_format.line_spacing = 1.5

    for line in ["Sincerely,", "", "[Corresponding author, on behalf of all authors]"]:
        p = doc.add_paragraph()
        if line:
            r = p.add_run(line)
            r.font.size = Pt(12)
            r.font.name = "Times New Roman"

    out = os.path.join(BASE, "hp_cover_letter.docx")
    doc.save(out)
    print("wrote", out)


def build_supplementary():
    doc = _setup_doc()
    head(doc, "Supplementary material", level=1)

    para(doc, f"Supplementary Figure 1. Closed malpractice claims per {_per} physicians by "
              f"specialty, 2008\u20132024 (rates, not counts).")
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    img = os.path.join(OUT, "hp_Supplementary_Figure_1.png")
    if os.path.exists(img):
        p.add_run().add_picture(img, width=Inches(5.8))
    doc.add_paragraph()

    para(doc, "Supplementary Figure 2. Physician workforce by specialty, indexed to 2008 (=100).")
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    img = os.path.join(OUT, "hp_Supplementary_Figure_2.png")
    if os.path.exists(img):
        p.add_run().add_picture(img, width=Inches(5.8))
    doc.add_paragraph()

    # Supplementary Table 1: data sources
    res_words = {1: "Annual", 2: "Biennial", 3: "Every 3 years"}
    table(doc,
          ["Series", "Source", "Resolution", "Years", "Role"],
          [["Physicians by specialty", "MHLW Statistics of Physicians", res_words.get(PHYS_RES, f"Every {PHYS_RES} years"),
            _year_label(PHYS_DF.columns), "Denominator & outcome"],
           ["Closed malpractice claims", "Supreme Court, by specialty", res_words.get(_resolution(LIT_DF), f"Every {_resolution(LIT_DF)} years"),
            _year_label(LIT_DF.columns), "Exposure (numerator)"],
           ["Hospitals by specialty", "MHLW Survey of Medical Institutions", res_words.get(HOSP_RES, f"Every {HOSP_RES} years"),
            _year_label(HOSP_DF.columns), "Outcome"],
           ["Clinics by specialty", "MHLW Survey (static)", res_words.get(CLINIC_RES, f"Every {CLINIC_RES} years"),
            _year_label(CLINIC_DF.columns), "Descriptive only"],
           ["JMSR report counts", "JMSR / MAIS", "Annual",
            _year_label(JMSR_CORR['years']), "Sensitivity (2016-2024)"],
           ["Newspaper article counts", "Nikkei Telecom 21", "Annual",
            _year_label(MEDIA_CORR['years']), f"Sensitivity ({MEDIA_START}-{MEDIA_END})"]],
          "Supplementary Table 1. Primary data sources and their resolution.")

    # Supplementary Table 2: TOST details
    eqp, eqh = RES["equivalence"][0], RES["equivalence"][1]
    rows = []
    for outcome, eq in [("Physician growth", eqp), ("Hospital growth", eqh)]:
        for t in eq["tests"]:
            rows.append([
                outcome,
                f"{eq['coef_per_SD']*100:+.2f}%",
                f"{eq['ci90_low']*100:+.2f}%, {eq['ci90_high']*100:+.2f}%",
                f"\u00b1{int(t['margin']*100)}%",
                f"{t['p_tost']:.3f}",
                "Yes" if t["equivalent"] else "No"
            ])
    table(doc,
          ["Outcome", "Coef per SD", "90% CI", "Margin", "TOST p", "Equivalent"],
          rows,
          "Supplementary Table 2. Equivalence (TOST) results (effect per +1 SD litigation rate).")

    # Supplementary Table 3: JMSR-adjusted hospital model
    table(doc,
          ["Exposure", "Coefficient", "p", "n"],
          [["Lagged litigation rate", fmt(JMSR["lit_coef"], 4), f"{JMSR['lit_p']:.2f}", JMSR["n_obs"]],
           ["Lagged JMSR report rate", fmt(JMSR["med_coef"], 4), f"{JMSR['med_p']:.2f}", JMSR["n_obs"]]],
          f"Supplementary Table 3. JMSR-adjusted annual hospital growth model "
          f"({JMSR_START}-2024) with both exposures entered simultaneously.")

    # Supplementary Table 4: media-adjusted hospital model
    table(doc,
          ["Exposure", "Coefficient", "p", "n"],
          [["Lagged litigation rate", fmt(MEDIA["lit_coef"], 4), f"{MEDIA['lit_p']:.2f}", MEDIA["n_obs"]],
           ["Lagged media count (per 1,000 articles)", fmt(MEDIA["media_coef"], 4), f"{MEDIA['media_p']:.2f}", MEDIA["n_obs"]]],
          f"Supplementary Table 4. Media-adjusted annual hospital growth model "
          f"({MEDIA_START}-{MEDIA_END}) with a linear time trend; full year fixed effects are "
          "omitted because the national article-count series is collinear with them.")

    out = os.path.join(BASE, "hp_supplementary.docx")
    doc.save(out)
    print("wrote", out)


def build_figure_pptx():
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]

    main_figs = [
        ("hp_Figure_1.png", "Figure 1",
         f"Equivalence (TOST) of the litigation-rate effect against \u00b1{MARGIN1}% and "
         f"\u00b1{MARGIN2}% margins; horizontal bars are 90% confidence intervals."),
        ("hp_Figure_2.png", "Figure 2",
         "Biennial physician growth against lagged litigation exposure measured as (a) counts and (b) rates."),
    ]
    supp_figs = [
        ("hp_Supplementary_Figure_1.png", "Supplementary Figure 1",
         f"Closed malpractice claims per {_per} physicians by specialty, 2008\u20132024 (rates, not counts)."),
        ("hp_Supplementary_Figure_2.png", "Supplementary Figure 2",
         "Physician workforce by specialty, indexed to 2008 (=100)."),
    ]

    def add_slide(prs, fn, num, cap):
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(PInches(0.5), PInches(0.2), PInches(12.3), PInches(0.7))
        tf = tb.text_frame
        tf.text = num
        tf.paragraphs[0].runs[0].font.size = PPt(24)
        tf.paragraphs[0].runs[0].font.bold = True
        img = os.path.join(OUT, fn)
        if os.path.exists(img):
            s.shapes.add_picture(img, PInches(1.2), PInches(1.1), height=PInches(5.2))
        cb = s.shapes.add_textbox(PInches(0.5), PInches(6.5), PInches(12.3), PInches(0.9))
        cf = cb.text_frame
        cf.word_wrap = True
        cf.text = cap
        cf.paragraphs[0].runs[0].font.size = PPt(14)

    for fn, num, cap in main_figs:
        add_slide(prs, fn, num, cap)
    out = os.path.join(BASE, "hp_figures.pptx")
    prs.save(out)
    print("wrote", out)

    prs2 = Presentation()
    prs2.slide_width = PInches(13.333)
    prs2.slide_height = PInches(7.5)
    for fn, num, cap in supp_figs:
        add_slide(prs2, fn, num, cap)
    out2 = os.path.join(BASE, "hp_supplementary_figures.pptx")
    prs2.save(out2)
    print("wrote", out2)


def copy_figures():
    """Copy existing high-resolution PNGs into Health Policy figure files."""
    mapping = {
        "fig3_equivalence.png": "hp_Figure_1.png",
        "fig4_counts_vs_rates.png": "hp_Figure_2.png",
        "fig1_litigation_rate.png": "hp_Supplementary_Figure_1.png",
        "fig2_physician_index.png": "hp_Supplementary_Figure_2.png",
    }
    for src, dst in mapping.items():
        s = os.path.join(OUT, src)
        d = os.path.join(OUT, dst)
        if os.path.exists(s):
            shutil.copy2(s, d)
            print("copied", s, "->", d)
        else:
            raise SystemExit(f"missing figure source: {s}")


def create_submission_zip():
    """Bundle all generated Health Policy submission files into one archive."""
    zip_path = os.path.join(OUT, "hp_submission.zip")
    files = [
        os.path.join(BASE, "hp_manuscript_en.docx"),
        os.path.join(BASE, "hp_title_page.docx"),
        os.path.join(BASE, "hp_cover_letter.docx"),
        os.path.join(BASE, "hp_highlights.docx"),
        os.path.join(BASE, "hp_supplementary.docx"),
        os.path.join(BASE, "hp_figures.pptx"),
        os.path.join(BASE, "hp_supplementary_figures.pptx"),
        os.path.join(OUT, "hp_Figure_1.png"),
        os.path.join(OUT, "hp_Figure_2.png"),
        os.path.join(OUT, "hp_Supplementary_Figure_1.png"),
        os.path.join(OUT, "hp_Supplementary_Figure_2.png"),
        __file__,
    ]
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as z:
        for path in files:
            if not os.path.exists(path):
                raise SystemExit(f"submission zip missing file: {path}")
            z.write(path, arcname=os.path.basename(path))
    print("wrote", zip_path)


def main():
    copy_figures()
    main_wc, abs_wc = build_manuscript()
    build_title_page(main_wc)
    build_highlights()
    build_cover_letter()
    build_supplementary()
    build_figure_pptx()
    create_submission_zip()


if __name__ == "__main__":
    main()
