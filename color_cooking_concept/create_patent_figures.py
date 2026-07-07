"""
TONGE Patent Figures - PowerPoint Generation
特許出願用図面7枚をPowerPointで生成
白黒のみ（JPO準拠）、全テキスト黒、日本語表記
"""

from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree


BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)


def set_font(run, size=10, bold=False):
    """Force font to black."""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = BLACK


def add_box(slide, left, top, width, height, text, font_size=10,
            bold=False, border_width=Pt(1.5), align=PP_ALIGN.CENTER):
    """Rounded rectangle, white fill, black border, black text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.color.rgb = BLACK
    shape.line.width = border_width
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        set_font(run, font_size, bold)
    return shape


def add_rect(slide, left, top, width, height, text="", font_size=9,
             align=PP_ALIGN.CENTER, fill_white=True):
    """Plain rectangle, white fill, black border, black text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.color.rgb = BLACK
    shape.line.width = Pt(1.5)
    if fill_white:
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = BLACK
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(3)
        tf.margin_right = Pt(3)
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, line in enumerate(text.split('\n')):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            run = p.add_run()
            run.text = line
            set_font(run, font_size)
    return shape


def add_diamond(slide, left, top, width, height, text, font_size=9):
    """Diamond shape, white fill, black border."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND, left, top, width, height)
    shape.line.color.rgb = BLACK
    shape.line.width = Pt(1.5)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_font(run, font_size, bold=True)
    return shape


def add_arrow(slide, sx, sy, ex, ey, width=Pt(1.5)):
    """Line with arrowhead, black."""
    c = slide.shapes.add_connector(1, sx, sy, ex, ey)
    c.line.color.rgb = BLACK
    c.line.width = width
    ln = c.line._ln
    tail = ln.find(qn('a:tailEnd'))
    if tail is None:
        tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('len', 'med')
    return c


def add_line(slide, sx, sy, ex, ey, dashed=False, width=Pt(1.0)):
    """Simple line, black."""
    c = slide.shapes.add_connector(1, sx, sy, ex, ey)
    c.line.color.rgb = BLACK
    c.line.width = width
    if dashed:
        c.line.dash_style = 4
    return c


def add_text(slide, left, top, width, height, text, font_size=9,
             bold=False, align=PP_ALIGN.LEFT):
    """Text box, black text, no border."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        set_font(run, font_size, bold)
    return txBox


def add_title(slide, text):
    add_text(slide, Cm(1), Cm(0.5), Cm(23), Cm(1.5), text,
             font_size=14, bold=True, align=PP_ALIGN.CENTER)


# ─── Presentation ────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = Cm(25.4)
prs.slide_height = Cm(19.05)
BLANK = prs.slide_layouts[6]


# ═══════════════════════════════════════════════════════════════════
# 【図1】全体構成ブロック図
# ═══════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK)
add_title(slide, "【図1】")

# Main frame
add_rect(slide, Cm(1.5), Cm(2), Cm(22), Cm(14.5), "")
add_text(slide, Cm(2), Cm(2.2), Cm(14), Cm(0.8),
         "調理タイミング判定装置 100", font_size=11, bold=True)

# 110 Color sensor
add_box(slide, Cm(2), Cm(3.5), Cm(5.5), Cm(4),
        "色彩センサー部 110\n\nRGBCセンサー 111\n(TCS3472)\n\n白色LED照明 112",
        font_size=9)

# 120 MCU
add_box(slide, Cm(8.5), Cm(3.5), Cm(7), Cm(4),
        "マイクロコントローラ部 120\n(ESP32-S3)\n\n"
        "- 色空間変換\n  (RGB -> XYZ -> L*a*b*)\n"
        "- 色差ΔE算出\n"
        "- 時系列追跡(環状バッファ)\n"
        "- 予測通知(変化率推定)",
        font_size=8, align=PP_ALIGN.LEFT)

# Input
add_box(slide, Cm(16.5), Cm(3.5), Cm(5.5), Cm(2.5),
        "入力手段\n(タッチパネル)", font_size=9)

# 130 Display
add_box(slide, Cm(2), Cm(8.5), Cm(5.5), Cm(3.5),
        "表示部 130\n(タッチスクリーン)\n\n"
        "- プリセット選択メニュー\n"
        "- ΔEプログレスバー\n"
        "- トレンドグラフ",
        font_size=8, align=PP_ALIGN.LEFT)

# 140 Notification
add_box(slide, Cm(8.5), Cm(8.5), Cm(4.5), Cm(3.5),
        "通知部 140\n\n"
        "- スピーカー\n  (880Hz / 1100Hz)\n"
        "- Wi-Fi / BLE",
        font_size=8, align=PP_ALIGN.LEFT)

# 150 Storage
add_box(slide, Cm(14), Cm(8.5), Cm(8), Cm(3.5),
        "記憶部 150\n(不揮発性メモリ)\n\n"
        "目標色プリセット辞書:\n"
        " きつね色  [68.5, 12.3, 42.1]\n"
        " 飴色      [73.2,  8.7, 38.5]\n"
        " ハシバミ色 [65.0, 10.1, 30.2]",
        font_size=8, align=PP_ALIGN.LEFT)

# External
add_box(slide, Cm(2), Cm(13.5), Cm(4.5), Cm(2),
        "食品表面\n(被計測対象)", font_size=9)

add_box(slide, Cm(7.5), Cm(13.5), Cm(4.5), Cm(2),
        "白色基準面\n(キャリブレーション用)", font_size=9)

# Arrows
add_text(slide, Cm(5.8), Cm(4.8), Cm(2.5), Cm(0.6), "I2C", font_size=8, bold=True)
add_arrow(slide, Cm(7.5), Cm(5.3), Cm(8.5), Cm(5.3))
add_arrow(slide, Cm(8.5), Cm(7.5), Cm(7.5), Cm(9))
add_arrow(slide, Cm(12), Cm(7.5), Cm(11), Cm(8.5))
add_arrow(slide, Cm(15.5), Cm(7.5), Cm(17), Cm(8.5))
add_arrow(slide, Cm(4.5), Cm(13.5), Cm(4.5), Cm(7.5))
add_arrow(slide, Cm(8.5), Cm(13.5), Cm(5.5), Cm(7.5))


# ═══════════════════════════════════════════════════════════════════
# 【図2】色空間変換の処理フロー
# ═══════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK)
add_title(slide, "【図2】")

steps = [
    "TCS3472 RGBC値取得\n(R, G, B, C)",
    "正規化\nr = R/C,  g = G/C,  b = B/C",
    "ホワイトバランス補正\nr' = r x coeff_R,  g' = g x coeff_G,  b' = b x coeff_B",
    "逆sRGBコンパンディング\nc > 0.04045 : ((c+0.055)/1.055)^2.4\nc <= 0.04045 : c / 12.92",
    "線形RGB -> CIE XYZ (D65)\n3x3 マトリクス変換",
    "CIE XYZ -> L*a*b* (D65白色点)\nL* = 116 f(Y/Yn) - 16\na* = 500 (f(X/Xn) - f(Y/Yn))\nb* = 200 (f(Y/Yn) - f(Z/Zn))",
    "出力: L*, a*, b*",
]

y = Cm(2.5)
for i, text in enumerate(steps):
    h = Cm(2.3) if i in [3, 5] else Cm(1.7)
    add_box(slide, Cm(5), y, Cm(15), h, text, font_size=9)
    if i < len(steps) - 1:
        add_arrow(slide, Cm(12.5), y + h, Cm(12.5), y + h + Cm(0.3))
    y += h + Cm(0.3)


# ═══════════════════════════════════════════════════════════════════
# 【図3】目標色プリセット辞書のデータ構造
# ═══════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK)
add_title(slide, "【図3】")

rows, cols = 8, 8
table_shape = slide.shapes.add_table(rows, cols, Cm(1.5), Cm(3), Cm(22), Cm(12))
table = table_shape.table

for i, w in enumerate([Cm(1.5), Cm(3.5), Cm(3.5), Cm(2), Cm(2), Cm(2), Cm(2.5), Cm(5)]):
    table.columns[i].width = w

# Header - white text on black background
headers = ["ID", "色名(JA)", "色名(EN)", "L*", "a*", "b*", "ΔE閾値", "カテゴリ"]
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = h
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = WHITE
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLACK

# Data rows
data = [
    ["1", "きつね色", "Golden Brown", "68.5", "12.3", "42.1", "6.0", "揚げ物・焼き物"],
    ["2", "飴色", "Caramel", "73.2", "8.7", "38.5", "5.0", "炒め物"],
    ["3", "ハシバミ色", "Hazelnut", "65.0", "10.1", "30.2", "5.0", "焼き菓子"],
    ["4", "こんがり", "Toasted", "60.0", "15.2", "35.8", "6.0", "パン"],
    ["5", "べっこう色", "Amber", "55.0", "20.3", "45.0", "4.0", "カラメル"],
    ["6", "焦がしバター色", "Beurre Noisette", "48.0", "12.5", "28.0", "4.0", "ソース"],
    ["N", "(ユーザー定義)", "(user-defined)", "...", "...", "...", "...", "..."],
]

for ri, row_data in enumerate(data):
    for ci, val in enumerate(row_data):
        cell = table.cell(ri + 1, ci)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = val
        run.font.size = Pt(8)
        run.font.color.rgb = BLACK
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE


# ═══════════════════════════════════════════════════════════════════
# 【図4】調理タイミング判定方法の処理フロー
# ═══════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK)
add_title(slide, "【図4】")

# START
add_box(slide, Cm(9), Cm(2.5), Cm(7), Cm(1.2),
        "開始", font_size=10, bold=True)
add_arrow(slide, Cm(12.5), Cm(3.7), Cm(12.5), Cm(4.2))

# S1
add_box(slide, Cm(7.5), Cm(4.2), Cm(10), Cm(1.3),
        "S1: プリセットから目標色を選択", font_size=9)
add_arrow(slide, Cm(12.5), Cm(5.5), Cm(12.5), Cm(5.9))

# S2
add_box(slide, Cm(7.5), Cm(5.9), Cm(10), Cm(1.3),
        "S2: 白色基準面の計測 -> 補正係数算出", font_size=9)
add_arrow(slide, Cm(12.5), Cm(7.2), Cm(12.5), Cm(7.6))

# S3
add_box(slide, Cm(7.5), Cm(7.6), Cm(10), Cm(1.3),
        "S3: 食品初期色を計測 -> 初期ΔE算出", font_size=9)
add_arrow(slide, Cm(12.5), Cm(8.9), Cm(12.5), Cm(9.3))

# S4
add_box(slide, Cm(7.5), Cm(9.3), Cm(10), Cm(1.5),
        "S4: 食品色の連続計測ループ\nRGBC -> L*a*b* -> ΔE算出", font_size=9)
add_arrow(slide, Cm(12.5), Cm(10.8), Cm(12.5), Cm(11.4))

# Decision
add_diamond(slide, Cm(9.5), Cm(11.4), Cm(6), Cm(2.5),
            "ΔE < 閾値?", font_size=10)

# No -> loop back
add_text(slide, Cm(16), Cm(12.2), Cm(2), Cm(0.6), "No", font_size=9, bold=True)
add_arrow(slide, Cm(15.5), Cm(12.7), Cm(18), Cm(12.7))
add_line(slide, Cm(18), Cm(12.7), Cm(18), Cm(10), width=Pt(1.5))
add_arrow(slide, Cm(18), Cm(10), Cm(17.5), Cm(10))

# Prediction branch
add_box(slide, Cm(1.5), Cm(11.7), Cm(6.5), Cm(2),
        "予測: 残りサイクル < 5 ?\n-> 予告通知出力", font_size=8)
add_arrow(slide, Cm(9.5), Cm(12.7), Cm(8), Cm(12.7))

# Yes
add_text(slide, Cm(11.5), Cm(13.8), Cm(2), Cm(0.6), "Yes", font_size=9, bold=True)
add_arrow(slide, Cm(12.5), Cm(13.9), Cm(12.5), Cm(14.8))

# S5
add_box(slide, Cm(7.5), Cm(14.8), Cm(10), Cm(1.5),
        "S5: 目標色到達通知\nアラーム音出力 / 到達画面表示", font_size=9)
add_arrow(slide, Cm(12.5), Cm(16.3), Cm(12.5), Cm(16.8))

# END
add_box(slide, Cm(9), Cm(16.8), Cm(7), Cm(1.2),
        "終了", font_size=10, bold=True)


# ═══════════════════════════════════════════════════════════════════
# 【図5】ΔE時系列追跡および予測通知
# ═══════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK)
add_title(slide, "【図5】")

# Axes
add_line(slide, Cm(3.5), Cm(3), Cm(3.5), Cm(15.5), width=Pt(1.5))
add_text(slide, Cm(1.5), Cm(2.5), Cm(2), Cm(0.8), "ΔE", font_size=11, bold=True)
add_line(slide, Cm(3.5), Cm(15.5), Cm(22), Cm(15.5), width=Pt(1.5))
add_text(slide, Cm(10), Cm(15.7), Cm(8), Cm(0.8),
         "時間 (計測サイクル)", font_size=9)

# Y labels
for val, y in [(50, Cm(3.5)), (40, Cm(5.3)), (30, Cm(7.1)),
               (20, Cm(9.0)), (10, Cm(11.8)), (5, Cm(13.2))]:
    add_text(slide, Cm(2), y, Cm(1.5), Cm(0.5), str(val), font_size=8)
    add_line(slide, Cm(3.3), y + Cm(0.2), Cm(3.5), y + Cm(0.2))

# X labels
for val, x in [(0, Cm(3.3)), (5, Cm(5.5)), (10, Cm(7.7)),
               (15, Cm(9.9)), (20, Cm(12.1)), (25, Cm(14.3)),
               (30, Cm(16.5)), (35, Cm(18.7))]:
    add_text(slide, x, Cm(15.5), Cm(1.2), Cm(0.5), str(val), font_size=7,
             align=PP_ALIGN.CENTER)

# Threshold line
thr_y = Cm(13.4)
add_line(slide, Cm(3.5), thr_y, Cm(22), thr_y, dashed=True, width=Pt(1.5))
add_text(slide, Cm(18.5), Cm(12.8), Cm(4), Cm(0.6),
         "ΔE閾値 = 5.0", font_size=8, bold=True)

# Data curve
pts = [
    (Cm(4.0), Cm(3.8)),
    (Cm(5.5), Cm(4.5)),
    (Cm(7.0), Cm(5.4)),
    (Cm(8.5), Cm(6.5)),
    (Cm(10.0), Cm(7.5)),
    (Cm(11.5), Cm(8.6)),
    (Cm(13.0), Cm(9.8)),
    (Cm(14.5), Cm(10.8)),   # (P) pre-alarm
    (Cm(16.0), Cm(11.8)),
    (Cm(17.5), Cm(12.7)),
    (Cm(19.0), Cm(13.3)),
    (Cm(20.0), Cm(13.4)),   # (R) reached
]

for i in range(len(pts) - 1):
    add_line(slide, pts[i][0], pts[i][1], pts[i+1][0], pts[i+1][1], width=Pt(2))

# Dots
dot_r = Cm(0.12)
for x, y in pts:
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL, x - dot_r, y - dot_r, dot_r*2, dot_r*2)
    d.fill.solid()
    d.fill.fore_color.rgb = BLACK
    d.line.color.rgb = BLACK
    d.line.width = Pt(0.5)

# Pre-alarm marker (open circle)
pax, pay = pts[7]
pa = slide.shapes.add_shape(MSO_SHAPE.OVAL, pax - Cm(0.25), pay - Cm(0.25), Cm(0.5), Cm(0.5))
pa.fill.solid()
pa.fill.fore_color.rgb = WHITE
pa.line.color.rgb = BLACK
pa.line.width = Pt(2)
add_text(slide, pax + Cm(0.5), pay - Cm(0.6), Cm(6), Cm(0.8),
         "(P) 予告通知発動 (残り約5サイクル)", font_size=8)

# Reached marker (double circle)
rx, ry = pts[-1]
r1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, rx - Cm(0.3), ry - Cm(0.3), Cm(0.6), Cm(0.6))
r1.fill.solid()
r1.fill.fore_color.rgb = BLACK
r1.line.color.rgb = BLACK
r1.line.width = Pt(2)
r2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, rx - Cm(0.15), ry - Cm(0.15), Cm(0.3), Cm(0.3))
r2.fill.solid()
r2.fill.fore_color.rgb = WHITE
r2.line.color.rgb = BLACK
r2.line.width = Pt(1)
add_text(slide, Cm(15), Cm(14.0), Cm(6), Cm(0.8),
         "(R) 目標到達 -> アラーム", font_size=8, bold=True)

# Legend
add_text(slide, Cm(3.5), Cm(16.5), Cm(20), Cm(0.8),
         "凡例:  --- ΔE計測値   (P) 予告通知発動点   "
         "(R) 目標到達 (ΔE < 閾値)   - - - ΔE閾値",
         font_size=8)


# ═══════════════════════════════════════════════════════════════════
# 【図6】外付けクリップオン型外観図
# ═══════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK)
add_title(slide, "【図6】")

# CoreS3 display
add_rect(slide, Cm(9), Cm(2.5), Cm(7), Cm(4.5),
         "表示部 130 / CoreS3\n\n"
         "  dE: 12.3\n"
         "  [========>        ] 62%\n"
         "  目標: きつね色",
         font_size=9, align=PP_ALIGN.LEFT)
add_text(slide, Cm(16.5), Cm(3.5), Cm(7), Cm(0.6),
         "<- タッチスクリーン (2インチ)", font_size=8)

# Cable
add_text(slide, Cm(10.5), Cm(7.2), Cm(5), Cm(0.6),
         "Groveケーブル (I2C)", font_size=8, align=PP_ALIGN.CENTER)
add_line(slide, Cm(12.5), Cm(7), Cm(12.5), Cm(8), width=Pt(1.5))

# Sensor
add_rect(slide, Cm(10), Cm(8), Cm(5), Cm(2.5),
         "センサー部 110\n\n"
         "(o) RGBCセンサー 111\n"
         "(o) 白色LED 112",
         font_size=9, align=PP_ALIGN.LEFT)

# Clip
add_rect(slide, Cm(8), Cm(11), Cm(9), Cm(0.6), "", fill_white=False)
add_text(slide, Cm(17.5), Cm(10.5), Cm(5), Cm(0.8),
         "<- クリップ機構", font_size=9)

# Pan rim (hatched look via gray)
add_rect(slide, Cm(7.5), Cm(11.6), Cm(10), Cm(1.2),
         "鍋 / フライパンの縁", font_size=9)

# Direction
add_text(slide, Cm(10), Cm(13.3), Cm(6), Cm(0.8),
         "照射・計測方向", font_size=9, bold=True, align=PP_ALIGN.CENTER)
add_arrow(slide, Cm(12.5), Cm(14), Cm(12.5), Cm(14.8))

# Food
add_rect(slide, Cm(7.5), Cm(14.8), Cm(10), Cm(2),
         "調理中の食品\n(被計測対象)", font_size=10)

# Annotation
add_text(slide, Cm(1), Cm(8.5), Cm(6.5), Cm(2.5),
         "計測距離:\n数cm - 10cm\n\n照明:\n白色LED内蔵\n(照明条件を統一)",
         font_size=8)


# ═══════════════════════════════════════════════════════════════════
# 【図7】UI画面遷移図
# ═══════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK)
add_title(slide, "【図7】")

# MENU
add_rect(slide, Cm(1), Cm(3), Cm(5.5), Cm(6),
         "MENU画面\n\nTONGE\n"
         "  > きつね色\n"
         "  > 飴色\n"
         "  > ハシバミ色\n"
         "  > こんがり\n"
         "  > べっこう色\n"
         "  > 焦がしバター",
         font_size=8, align=PP_ALIGN.LEFT)

# CALIBRATE
add_rect(slide, Cm(10), Cm(3), Cm(5.5), Cm(4.5),
         "CALIBRATE画面\n\n"
         "白い面を\n"
         "センサーに\n"
         "かざしてください\n\n"
         "[タッチで開始]",
         font_size=8, align=PP_ALIGN.LEFT)

# MONITORING
add_rect(slide, Cm(10), Cm(9), Cm(5.5), Cm(6),
         "MONITORING画面\n\n"
         "計測色:   目標色:\n"
         "dE: 15.2\n"
         "[========>        ] 62%\n\n"
         "L*=62  a*=11  b*=35\n"
         "[トレンドグラフ]\n\n"
         "[タッチ: 停止]",
         font_size=8, align=PP_ALIGN.LEFT)

# REACHED
add_rect(slide, Cm(19), Cm(9), Cm(5.5), Cm(5),
         "REACHED画面\n\n"
         "きつね色\n"
         "到達!\n\n"
         "Time: 05:23\n"
         "dE: 4.8\n\n"
         "アラーム鳴動",
         font_size=8, align=PP_ALIGN.LEFT)

# Arrows
add_text(slide, Cm(6.8), Cm(4.5), Cm(3), Cm(0.6),
         "長押し ->", font_size=8)
add_arrow(slide, Cm(6.5), Cm(5.2), Cm(10), Cm(5.2))

add_text(slide, Cm(8.2), Cm(7.8), Cm(3), Cm(0.6),
         "タッチ", font_size=8)
add_arrow(slide, Cm(12.5), Cm(7.5), Cm(12.5), Cm(9))

add_text(slide, Cm(15.7), Cm(10.8), Cm(3.5), Cm(0.6),
         "ΔE<閾値 ->", font_size=8)
add_arrow(slide, Cm(15.5), Cm(11.5), Cm(19), Cm(11.5))

add_text(slide, Cm(10), Cm(15.5), Cm(8), Cm(0.6),
         "<-- タッチ: メニューに戻る", font_size=8)
add_line(slide, Cm(21.5), Cm(14), Cm(21.5), Cm(16), width=Pt(1.5))
add_line(slide, Cm(21.5), Cm(16), Cm(3.5), Cm(16), width=Pt(1.5))
add_arrow(slide, Cm(3.5), Cm(16), Cm(3.5), Cm(9))

add_text(slide, Cm(5), Cm(11.5), Cm(4), Cm(0.6),
         "タッチ(停止)", font_size=7)
add_arrow(slide, Cm(10), Cm(12), Cm(6.5), Cm(9))


# ═══════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════

out = "/home/ubuntu/repos/wip/color_cooking_concept/filing_docs/patent_figures.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
