"""
Generate editable PPTX with all figures (1 per slide).
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pathlib import Path

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG_DIR = os.path.join(BASE_DIR, 'figures')
OUT_DIR = BASE_DIR

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

figures = [
    ('fig1_morphology_vs_introgression.png',
     'Figure 1',
     'Archaic introgression by language morphological type (Kruskal-Wallis test).'),
    ('fig2_nean_deni_scatter_morphology.png',
     'Figure 2',
     'Bivariate scatter: Neanderthal x Denisovan sharing colored by morphological type.'),
    ('fig3_tone_vs_introgression.png',
     'Figure 3',
     'Archaic introgression in tonal vs. non-tonal languages (Mann-Whitney U test).'),
    ('fig4_language_family_introgression.png',
     'Figure 4',
     'Archaic introgression by language family (mean +/- SD).'),
    ('fig5_world_map_typology.png',
     'Figure 5',
     'Global distribution: language typology x archaic introgression (marker size = Neanderthal).'),
]

for fname, title, caption in figures:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Image
    fig_path = os.path.join(FIG_DIR, fname)
    if Path(fig_path).exists():
        slide.shapes.add_picture(fig_path, Inches(1.0), Inches(1.0), width=Inches(11.0))

    # Caption
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.6))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = caption
    p2.font.size = Pt(14)
    p2.alignment = PP_ALIGN.CENTER

out_path = os.path.join(OUT_DIR, 'figures_archaic_language.pptx')
prs.save(out_path)
print(f"PPTX saved to: {out_path}")
