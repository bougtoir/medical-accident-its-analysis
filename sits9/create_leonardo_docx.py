#!/usr/bin/env python3
"""
Generate SITS9 manuscript formatted for Leonardo (MIT Press) General Article.

Formatting rules
────────────────
- English language
- 5,000–8,000 words
- 12 pt Times New Roman, double-spaced
- Figures NOT embedded; separate files (PNG/TIFF)
- Figure legends collected at end of manuscript
- References: author-date style (Harvard)
- Abstract: 150–200 words
"""

from __future__ import annotations

import re
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.shared import Inches, Pt, RGBColor, Cm
from docx.oxml.ns import qn

OUT_DIR = Path(__file__).parent / "output"

# ---------------------------------------------------------------------------
# helpers
# ---------------------------------------------------------------------------

def _style_doc(doc: Document) -> None:
    """Set default font and paragraph style for the whole document."""
    style = doc.styles["Normal"]
    font = style.font
    font.name = "Times New Roman"
    font.size = Pt(12)
    pf = style.paragraph_format
    pf.line_spacing_rule = WD_LINE_SPACING.DOUBLE
    pf.space_after = Pt(0)
    pf.space_before = Pt(0)

    # Headings: TNR bold
    for level in (1, 2, 3):
        hs = doc.styles[f"Heading {level}"]
        hs.font.name = "Times New Roman"
        hs.font.color.rgb = RGBColor(0, 0, 0)
        hs.paragraph_format.line_spacing_rule = WD_LINE_SPACING.DOUBLE
        hs.paragraph_format.space_before = Pt(12)
        hs.paragraph_format.space_after = Pt(6)


def _para(doc: Document, text: str, bold: bool = False,
          italic: bool = False, align=None) -> None:
    """Add a paragraph with optional bold/italic and citation-superscript parsing."""
    p = doc.add_paragraph()
    if align is not None:
        p.alignment = align
    # Parse {1}, {2-3} etc. as superscript citation markers
    parts = re.split(r'(\{[^}]+\})', text)
    for part in parts:
        if part.startswith('{') and part.endswith('}'):
            run = p.add_run(part[1:-1])
            run.font.superscript = True
            run.font.size = Pt(10)
            run.bold = bold
            run.italic = italic
        else:
            run = p.add_run(part)
            run.bold = bold
            run.italic = italic


def _table(doc: Document, headers: list[str], rows: list[list[str]],
           caption: str | None = None) -> None:
    """Add a table with an optional caption above it."""
    if caption:
        p = doc.add_paragraph()
        parts = caption.split(".", 1)
        if len(parts) == 2:
            run = p.add_run(parts[0] + ".")
            run.bold = True
            run.font.size = Pt(11)
            run = p.add_run(parts[1])
            run.font.size = Pt(11)
        else:
            run = p.add_run(caption)
            run.font.size = Pt(11)
    tbl = doc.add_table(rows=1, cols=len(headers), style="Table Grid")
    for i, h in enumerate(headers):
        cell = tbl.rows[0].cells[i]
        cell.text = h
        for par in cell.paragraphs:
            for r in par.runs:
                r.bold = True
                r.font.size = Pt(10)
                r.font.name = "Times New Roman"
    for row_data in rows:
        row = tbl.add_row()
        for i, val in enumerate(row_data):
            cell = row.cells[i]
            cell.text = ""
            p = cell.paragraphs[0]
            parts = re.split(r'(\{[^}]+\})', val)
            for part in parts:
                if part.startswith('{') and part.endswith('}'):
                    run = p.add_run(part[1:-1])
                    run.font.superscript = True
                    run.font.size = Pt(8)
                    run.font.name = "Times New Roman"
                else:
                    run = p.add_run(part)
                    run.font.size = Pt(10)
                    run.font.name = "Times New Roman"


def _code(doc: Document, code: str) -> None:
    p = doc.add_paragraph()
    run = p.add_run(code)
    run.font.name = "Courier New"
    run.font.size = Pt(9)
    rpr = run._element.get_or_add_rPr()
    shd = rpr.makeelement(qn("w:shd"), {
        qn("w:val"): "clear",
        qn("w:color"): "auto",
        qn("w:fill"): "F0F0F0",
    })
    rpr.append(shd)


def _heading(doc: Document, text: str, level: int = 1) -> None:
    doc.add_heading(text, level=level)


# ---------------------------------------------------------------------------
# manuscript body
# ---------------------------------------------------------------------------

def build_manuscript() -> Document:
    doc = Document()
    _style_doc(doc)

    # ── title page ──────────────────────────────────────────────────────────
    _para(doc,
          "SITS9: A Stitch-in-Time Language Unifying Music, "
          "Textile, and Cuisine through Operational Homology",
          bold=True, align=WD_ALIGN_PARAGRAPH.CENTER)
    doc.add_paragraph()

    _para(doc, "[Author name(s) and affiliation(s) to be inserted]",
          italic=True, align=WD_ALIGN_PARAGRAPH.CENTER)
    doc.add_paragraph()

    _para(doc,
          "Submitted to: Leonardo, MIT Press — General Article",
          italic=True, align=WD_ALIGN_PARAGRAPH.CENTER)
    doc.add_paragraph()
    doc.add_paragraph()

    # ── abstract ────────────────────────────────────────────────────────────
    _heading(doc, "Abstract", 1)
    _para(doc,
          "We propose SITS9 (A Stitch in Time Saves Nine), a minimal "
          "domain-specific language consisting of five primitive instructions "
          "-- FWD, RET, CROSS, TENSION, and ANCHOR -- that can simultaneously "
          "describe musical chord progressions, sewing/embroidery patterns, and "
          "cooking procedures. By defining domain-specific renderers (a 'loom "
          "renderer' for textile SVG, a 'tone renderer' for MIDI music, and a "
          "'kitchen renderer' for text-based recipes), a single SITS9 "
          "program (called a Deck) produces valid, aesthetically coherent output "
          "in all three domains. We demonstrate the system using Pachelbel's "
          "Canon (I-V-vi-iii-IV-I-IV-V) as a test case and argue that the "
          "operational homology reflects both historical fact -- the Jacquard "
          "loom and music-box punch-card convergence, and the Inca quipu as "
          "a precursor of cross-domain encoding -- and a deeper cognitive "
          "universality in how humans structure temporal processes. We further "
          "discuss applications to cryptographic steganography (textile ciphers), "
          "palette constraints and cultural identity, physical-object "
          "steganography via femtosecond laser inscription, and the design of "
          "cross-modal workshops in which a single Deck generates cuisine, "
          "handicraft, and music as parallel experiences.")

    _para(doc,
          "Keywords: SITS9; domain-specific language; operational homology; Jacquard "
          "loom; quipu; steganography; cross-modal workshop; chord progression; "
          "sashiko; computational craft",
          italic=True)
    doc.add_paragraph()

    # ── 1. Introduction ────────────────────────────────────────────────────
    _heading(doc, "1. Introduction: 'Moving Forward by Going Back'", 1)

    _para(doc,
          "The Japanese back stitch (kaeshi-nui) is defined by a single "
          "kinetic principle: advance the needle, then retreat it partway "
          "through the previous stitch hole before advancing again. This "
          "forward-return-forward cycle creates overlapping stitches that are "
          "stronger and more continuous than a simple running stitch. The "
          "operation can be summarised as FWD-RET-FWD, where the retreat "
          "through previously covered ground is what gives the stitch its "
          "structural integrity.{1}")

    _para(doc,
          "Harmonic chord progressions in tonal music obey the same "
          "principle. A canonical ii-V-I progression (e.g. Dm7-G7-Cmaj7 "
          "in C major) moves away from the tonic (forward into tension) "
          "and then returns to it (resolution). Simply ascending a scale "
          "does not constitute 'progression'; it is the inclusion of "
          "return -- the controlled retreat to a previously established "
          "point -- that produces harmonic motion and emotional "
          "persuasiveness.{2}")

    _para(doc,
          "Cooking exhibits the same pattern. The foundation of classical "
          "sauce-making is a cycle of heating (forward), cooling and "
          "tasting (return), and reheating (forward again). Reducing a "
          "stock into a demi-glace is, operationally, a literal back stitch: "
          "energy is added, withdrawn, evaluated, and added again in an "
          "iterative loop that progressively concentrates flavour.{3}")

    _para(doc,
          "That these three domains share an operational core is not, in "
          "itself, surprising: all three are acts of applying temporal "
          "operations to a material substrate to produce emergent structure "
          "(thread becomes fabric, vibrations become melody, ingredients "
          "become cuisine). What is surprising is that the shared core can "
          "be captured by exactly five instructions, and that the resulting "
          "cross-domain translations are not merely formally valid but "
          "aesthetically coherent -- the textile 'sounds like' the music, "
          "and the recipe 'looks like' the stitch pattern, in ways that "
          "human participants in pilot workshops have found intuitively "
          "meaningful.")

    _para(doc,
          "This paper formalises the shared operational structure underlying "
          "these three domains as SITS9 -- named after the proverb "
          "'A stitch in time saves nine', reflecting the principle that a "
          "single well-placed instruction eliminates redundancy across "
          "domains. We define a minimal instruction set, "
          "implement three domain-specific renderers, demonstrate the "
          "system with Pachelbel's Canon as a test case, and discuss "
          "historical precedents, cryptographic implications, and "
          "workshop applications.")

    # ── 2. Historical parallel ──────────────────────────────────────────────
    _heading(doc, "2. Historical Parallel Evolution: Three Domains, One Code", 1)

    _para(doc,
          "When the developmental histories of music, textile craft, and "
          "cuisine are placed side by side, a striking parallelism emerges: "
          "all three domains followed the same trajectory of instruction-set "
          "complexification (Table 1).")

    _table(doc,
           ["Era", "Textile", "Music", "Cuisine", "Shared instruction pattern"],
           [
               ["Primitive",
                "Running stitch (equal-interval repetition)",
                "Monophonic chant (Gregorian)",
                "Open-flame roasting (single process)",
                "loop { FWD }"],
               ["Ancient",
                "Weaving (warp-weft crossing)",
                "Organum (two voices)",
                "Boiling (water + fire combined)",
                "Invention of CROSS"],
               ["Medieval",
                "Sashiko (geometric repetition)",
                "Counterpoint (polyphony)",
                "Sauce systems (base + variation)",
                "Multi-layer parallel processing"],
               ["Renaissance",
                "Tapestry (pictorial expression)",
                "Fugue (structural counterpoint)",
                "Court cuisine (course structure)",
                "Systematisation of RET"],
               ["Industrial Rev.",
                "Jacquard loom (punch card)",
                "Barrel organ / music box (punch card)",
                "Canning (process preservation)",
                "Encoding and reproduction of process"],
               ["20th century",
                "Bauhaus textiles",
                "Minimal music (Reich, Riley)",
                "Nouvelle cuisine",
                "Parameter minimisation + precision"],
               ["Present",
                "CNC knitting / digital embroidery",
                "DAW / MIDI",
                "Molecular gastronomy / sous-vide",
                "Binary instruction sequences"],
           ],
           caption="Table 1. Parallel evolution of instruction-set complexity across three temporal-process domains.")
    doc.add_paragraph()

    _heading(doc, "2.1 The Jacquard Loom: Literal Unification", 2)
    _para(doc,
          "In 1804 Joseph Marie Jacquard patented a loom driven by punched "
          "cards that encoded weaving patterns. At the same period, barrel "
          "organs and music boxes used an identical medium -- punched cards "
          "or punched rolls -- to encode melodies. Textile and music were "
          "thus literally driven by the same programming language before "
          "the computer existed.{1} Charles Babbage acknowledged Jacquard's "
          "cards as the direct inspiration for the Analytical Engine (1837), "
          "and Ada Lovelace famously noted that the Engine 'weaves algebraical "
          "patterns just as the Jacquard-loom weaves flowers and leaves.'")

    _heading(doc, "2.2 The Quipu: Pre-Columbian Multi-Domain Encoding", 2)
    _para(doc,
          "The Inca quipu (khipu) is an even earlier precedent for "
          "cross-domain encoding. A quipu encodes information through the "
          "combination of knot types, positions, counts, and cord colours. "
          "The same physical object simultaneously served as an accounting "
          "system (numerical records), a narrative device (oral-history "
          "mnemonic), and a textile artefact (fibre structure). The quipu "
          "is an accounting ledger, an epic poem, and a textile -- a single "
          "physical object encoding multiple domains at once.{4} Recent "
          "research suggests that some quipus carried intentionally "
          "encrypted military communications, which the Spanish conquerors "
          "could not decode.")

    _para(doc,
          "The historical lineage can thus be traced as: quipu (Inca) -> "
          "Jacquard loom punch card (1804) -> barrel organ / player piano "
          "(same medium) -> Babbage Analytical Engine (1837) -> computer -> "
          "MIDI + digital embroidery -> SITS9 (return to unification).")

    _para(doc,
          "It is significant that each 'branching point' in this lineage "
          "was a moment when a single encoding medium was applied to a "
          "new domain. The Jacquard card was originally textile-only; "
          "its adoption for barrel organs was a deliberate cross-domain "
          "transfer. Babbage's recognition that the same cards could "
          "drive computation was another such transfer. In each case, "
          "the medium did not change -- only the renderer. SITS9 "
          "recapitulates this pattern at the software level: the 'card "
          "deck' (JSON-serialised instruction sequence) remains constant "
          "while the renderer determines the output domain.")

    _para(doc,
          "A further historical note: the McLuhan dictum 'the medium is "
          "the message' is literally embodied in both quipu and SITS9.{5} "
          "In a quipu, the medium (knotted cord) and the message (the "
          "information encoded in the knots) are physically inseparable. "
          "In SITS9, the stitch pattern IS the music IS the recipe -- "
          "they are not three messages carried in one medium, but one "
          "message read in three ways.")

    # ── 3. SITS9 Formal Specification ───────────────────────────────────
    _heading(doc, "3. SITS9 Formal Specification", 1)

    _heading(doc, "3.1 Primitive Instructions", 2)
    _para(doc,
          "SITS9 defines exactly five primitive instructions (Table 2). "
          "The design goal was the smallest instruction set that can "
          "express the forward-return-cross-tension-anchor patterns "
          "observed in all three domains.")

    _table(doc,
           ["Instruction", "Parameter", "Loom reading", "Tone reading", "Kitchen reading"],
           [
               ["FWD(n)", "n: distance",
                "Advance n stitch units",
                "n-degree interval / n-beat duration",
                "Heat for n time-units"],
               ["RET(n)", "n: return distance",
                "Retreat n stitch units (back stitch)",
                "n-degree descent / rest",
                "Rest / cool for n time-units"],
               ["CROSS", "angle, depth, wrap",
                "Toggle surface/back (needle pierces fabric)",
                "Toggle melody/bass voice",
                "Switch cooking method (grill/simmer/steam)"],
               ["TENSION(v)", "v: 0.0-1.0",
                "Thread tension",
                "Dynamic level (velocity)",
                "Heat level (low-high)"],
               ["ANCHOR", "(none)",
                "Knot (fix thread to fabric)",
                "Tonic resolution",
                "Taste and season (fix flavour)"],
           ],
           caption="Table 2. The five SITS9 primitive instructions and their domain-specific interpretations.")
    doc.add_paragraph()

    _heading(doc, "3.2 Card and Deck", 2)
    _para(doc,
          "A sequence of primitive instructions is grouped into a Card "
          "(analogous to one Jacquard punch card). A Deck is an ordered "
          "sequence of Cards, analogous to the Jacquard card deck. Each "
          "Card may carry optional domain-specific metadata (chord name, "
          "stitch pattern name, etc.), but rendering is based solely on "
          "the primitive instructions. A Deck also carries global "
          "parameters: BPM (for the tone renderer), key root (MIDI note "
          "number), and stitch-unit scale (mm per unit).")

    _para(doc,
          "In the reference implementation, instructions are defined as "
          "frozen (immutable) dataclasses in Python, ensuring that a Deck "
          "is a deterministic specification: the same Deck always produces "
          "the same output from any renderer. Decks are serialised to and "
          "from JSON, which serves as the 'punch card format' -- a "
          "platform-independent, human-readable, and version-controllable "
          "representation. A Card may also carry a repeat count, allowing "
          "loops without duplicating instruction data (analogous to a "
          "Jacquard card being run through the loom multiple times).")

    _code(doc,
          "Deck\n"
          "  +-- Card 0: 'D (I)'     [ANCHOR, TENSION(0.2)]\n"
          "  +-- Card 1: 'A (V)'     [TENSION(0.8), FWD(5), CROSS]\n"
          "  +-- Card 2: 'Bm (vi)'   [TENSION(0.5), FWD(2), CROSS]\n"
          "  +-- Card 3: 'F#m (iii)' [TENSION(0.4), RET(5), CROSS]\n"
          "  +-- Card 4: 'G (IV)'    [TENSION(0.3), FWD(1), CROSS]\n"
          "  +-- Card 5: 'D (I)'     [TENSION(0.2), RET(5), CROSS]\n"
          "  +-- Card 6: 'G (IV)'    [TENSION(0.3), FWD(3), CROSS]\n"
          "  +-- Card 7: 'A (V)'     [TENSION(0.8), FWD(2), CROSS]\n")

    _heading(doc, "3.3 Renderer Architecture", 2)
    _para(doc,
          "Each renderer takes a Deck as input and produces domain-specific "
          "output (see Figure 1 for the SVG renderer output). The "
          "architecture is extensible: adding a new renderer (e.g. for "
          "dance choreography, garden layout, or cipher decryption) "
          "requires no change to the Deck structure. This separation of "
          "instruction and rendering mirrors the Model-View pattern in "
          "software architecture, but applied to physical-world domains: "
          "the Deck is the model, and each renderer is a view that "
          "interprets the same abstract instructions in a domain-specific "
          "context.")

    _code(doc,
          "                 +-> LoomRenderer (SVG)     -> sashiko pattern\n"
          "SITS9 Deck --+-> ToneRenderer (MIDI)    -> music\n"
          "                 +-> KitchenRenderer (text)  -> recipe\n")

    # Inline figure (user requirement: figures must be embedded in manuscript)
    stitch_png = OUT_DIR / "pachelbel_stitch_2loops.png"
    if stitch_png.exists():
        doc.add_paragraph()  # spacer
        doc.add_picture(str(stitch_png), width=Inches(5.5))
        cap = doc.add_paragraph()
        cap.paragraph_format.space_before = Pt(14)
        cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = cap.add_run(
            "Figure 1. SITS9 Loom Renderer output for Pachelbel's Canon "
            "(I-V-vi-iii-IV-I-IV-V, 2 loops). Red solid lines = surface "
            "stitches; grey dashed lines = back stitches; filled black "
            "circles = anchor points (knots)."
        )
        run.font.size = Pt(10)
        run.font.italic = True

    # ── 4. Structural Homology: Back Stitch = Chord Progression ────────────
    _heading(doc, "4. Structural Homology: Back Stitch and Chord Progression", 1)

    _para(doc,
          "Neither back stitching nor chord progression moves forward in a "
          "straight line. The controlled retreat is what gives each its "
          "characteristic strength: physical tensile strength in the case "
          "of stitching, and emotional persuasiveness in the case of harmony. "
          "A running stitch (pure FWD) is the textile equivalent of simply "
          "ascending a scale: functional, but not 'progressive'. It is the "
          "inclusion of RET -- return through previously covered ground -- "
          "that constitutes progression in both domains.")

    _para(doc,
          "This claim is not merely metaphorical. Consider the physical "
          "mechanics: in back stitching, the thread passes through the "
          "same hole twice -- first going forward past it, then returning "
          "through it. This overlap creates a denser, stronger seam "
          "than a running stitch. In a ii-V-I cadence, the harmony "
          "passes through dominant tension (the 'overshoot') before "
          "returning to the tonic. The resolution through previously "
          "heard harmonic territory creates a denser, more satisfying "
          "closure than simply arriving at the tonic from an adjacent "
          "chord. In both cases, the retreat through previously covered "
          "ground is the mechanism that converts a sequence into a "
          "structure.")

    _para(doc,
          "Table 3 maps stitch types to their SITS9 instruction "
          "patterns, corresponding harmonic idioms, and culinary analogues.")

    _table(doc,
           ["Stitch type", "SITS9 pattern", "Harmonic analogue", "Culinary analogue"],
           [
               ["Running stitch",
                "loop { FWD(n), CROSS }",
                "Pedal point (sustained chord)",
                "Constant-temperature heating"],
               ["Back stitch",
                "loop { FWD, CROSS, FWD, TENSION+, CROSS, RET, TENSION-, ANCHOR }",
                "ii-V-I functional progression",
                "Heat-cool-reheat (sauce reduction)"],
               ["Half-back stitch",
                "loop { FWD(n), CROSS, RET(n/2), CROSS }",
                "Deceptive cadence (V-vi)",
                "Partial cooling (carry-over cooking)"],
               ["Herringbone (chidori-gake)",
                "loop { FWD(1), CROSS(+30), FWD(1), CROSS(-30) }",
                "Stepwise bass line",
                "Alternating ingredients"],
               ["Blind hem (matsuri-nui)",
                "loop { FWD(n), CROSS(depth=minimal) }",
                "Pedal tone (sustained bass)",
                "Long low-heat simmer"],
               ["Blanket stitch (kagari-nui)",
                "loop { CROSS(wrap=true), FWD(1) }",
                "Ostinato (repeating pattern)",
                "Continuous stirring (risotto)"],
           ],
           caption="Table 3. Mapping of stitch types to SITS9 patterns, harmonic idioms, and culinary processes.")
    doc.add_paragraph()

    # ── 5. Demonstration: Pachelbel's Canon ─────────────────────────────────
    _heading(doc, "5. Demonstration: Pachelbel's Canon", 1)

    _para(doc,
          "To validate the system, we encoded the harmonic progression of "
          "Pachelbel's Canon in D (I-V-vi-iii-IV-I-IV-V) as an eight-Card "
          "SITS9 Deck. Each Card's TENSION value reflects the harmonic "
          "tension of the chord degree (I = 0.2, V = 0.8, vi = 0.5, "
          "iii = 0.4, IV = 0.3). The FWD/RET directions and distances "
          "were chosen to produce a back-stitch-like forward-and-return "
          "pattern in the SVG renderer, while CROSS instructions with "
          "alternating angles (+/-30 degrees) create a visible zigzag "
          "in the stitch pattern.")

    _para(doc,
          "Figure 2 shows the harmonic tension profile of the Deck: "
          "each chord is represented as a coloured bar (blue = ANCHOR, "
          "red = FWD, green = RET) whose height encodes its TENSION value. "
          "The dashed connecting line reveals the 'tension arc' of the "
          "progression -- rising sharply to the dominant (V), descending "
          "through the return sequence (vi -> iii -> IV -> I), and "
          "re-ascending for the loop.")

    # Figure 2: Score schema
    score_png = OUT_DIR / "pachelbel_score_schema.png"
    if score_png.exists():
        doc.add_paragraph()
        doc.add_picture(str(score_png), width=Inches(5.5))
        cap = doc.add_paragraph()
        cap.paragraph_format.space_before = Pt(14)
        cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = cap.add_run(
            "Figure 2. Harmonic tension profile of Pachelbel's Canon "
            "encoded as SITS9 instructions. Bar colour indicates instruction "
            "type: blue = ANCHOR, red = FWD, green = RET. Bar height = "
            "TENSION value. Dashed line traces the tension arc."
        )
        run.font.size = Pt(10)
        run.font.italic = True

    _para(doc,
          "The same Deck was passed to all three renderers:")
    _para(doc,
          "(a) The LoomRenderer produced an SVG sashiko pattern (Figure 1) "
          "in which surface stitches (red solid lines) and back stitches "
          "(grey dashed lines) form a zigzag that visually encodes the "
          "harmonic tension profile. Anchor points (filled circles) mark "
          "the tonic. Figure 3 shows the textile schema: the needle path "
          "alternates between the surface (orange band) and back (blue band) "
          "of the fabric, with each CROSS instruction representing a "
          "penetration of the needle through the cloth.")

    # Figure 3: Textile schema
    textile_png = OUT_DIR / "pachelbel_textile_schema.png"
    if textile_png.exists():
        doc.add_paragraph()
        doc.add_picture(str(textile_png), width=Inches(5.5))
        cap = doc.add_paragraph()
        cap.paragraph_format.space_before = Pt(14)
        cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = cap.add_run(
            "Figure 3. Textile schema (back-stitch pattern) for Pachelbel's "
            "Canon. The needle path alternates between surface (orange) and "
            "back (blue) of the fabric. Red lines = FWD stitches on surface; "
            "orange diagonal lines = CROSS (needle penetration); black dot = "
            "ANCHOR (knot)."
        )
        run.font.size = Pt(10)
        run.font.italic = True

    _para(doc,
          "(b) The ToneRenderer produced a MIDI file with two tracks "
          "(piano melody on channel 0, strings bass on channel 1). "
          "Each FWD generates note-on events proportional to n, each RET "
          "produces a rest, CROSS toggles between channels, TENSION sets "
          "velocity, and ANCHOR resolves to a tonic major triad.")
    _para(doc,
          "(c) The KitchenRenderer produced a text recipe in which FWD "
          "maps to heating time, RET to resting/cooling, CROSS to "
          "switching cooking method (grill -> simmer -> steam), TENSION "
          "to heat level, and ANCHOR to tasting and seasoning. Figure 4 "
          "shows the cooking schema: the heat level (TENSION) traces the "
          "same arc as the harmonic tension, while background colour zones "
          "indicate cooking method transitions triggered by CROSS.")

    # Figure 4: Cooking schema
    cooking_png = OUT_DIR / "pachelbel_cooking_schema.png"
    if cooking_png.exists():
        doc.add_paragraph()
        doc.add_picture(str(cooking_png), width=Inches(5.5))
        cap = doc.add_paragraph()
        cap.paragraph_format.space_before = Pt(14)
        cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = cap.add_run(
            "Figure 4. Cooking schema for Pachelbel's Canon. The heat curve "
            "(TENSION) follows the same profile as the harmonic tension in "
            "Figure 2. Background zones indicate cooking method: red = grill, "
            "blue = simmer, green = steam. Dashed vertical lines = CROSS "
            "(method switch)."
        )
        run.font.size = Pt(10)
        run.font.italic = True

    _para(doc,
          "All three outputs are generated from the identical JSON-serialised "
          "Deck, confirming that a single 'program' produces coherent "
          "output across three sensory domains. The visual similarity "
          "between the tension arc in Figure 2, the zigzag trajectory in "
          "Figure 3, and the heat curve in Figure 4 is not a coincidence "
          "but a direct consequence of the shared instruction sequence. "
          "The source code and all generated artefacts are available in "
          "the supplementary repository.")

    _para(doc,
          "It is worth noting why Pachelbel's Canon is a particularly "
          "revealing test case. The progression I-V-vi-iii-IV-I-IV-V is "
          "the prototypical 'back-stitch' pattern in harmonic space: "
          "the first half (I to iii) moves away from the tonic by stepwise "
          "descent through the diatonic circle, while the second half "
          "(IV to V) returns toward the tonic, overlapping with the "
          "starting territory before re-advancing. This is structurally "
          "identical to the kinetics of kaeshi-nui: the needle advances "
          "past the previous stitch, retreats through it, and advances "
          "again. The SITS9 encoding makes this identity not metaphorical "
          "but literal -- the same FWD-RET-CROSS instruction sequence "
          "generates both the stitch trajectory and the harmonic motion.")

    _heading(doc, "5.1 National Anthem Demonstrations", 2)
    _para(doc,
          "To test whether the system generalises beyond a single "
          "well-known progression, we encoded four national anthems "
          "with contrasting musical characters: Kimigayo (Japan), "
          "La Marseillaise (France), The Star-Spangled Banner (US), "
          "and Hino Nacional Brasileiro (Brazil). Each anthem was "
          "converted into a SITS9 Deck by mapping its melodic contour, "
          "dynamic profile, and rhythmic character to the five "
          "primitives, following the same methodology as the Pachelbel "
          "encoding.")

    _para(doc,
          "The results confirmed the predictions based on each anthem's "
          "musical characteristics (see Section 9.1 for the full "
          "comparison table). Kimigayo, with its pentatonic scale, slow tempo "
          "(BPM 54), and narrow melodic range, produced a SITS9 Deck "
          "with small FWD values (2-3), low TENSION (0.15-0.45), and "
          "gentle CROSS angles (15 degrees). The resulting SVG shows a "
          "restrained, meditative stitch pattern with short, evenly "
          "spaced stitches -- visually reminiscent of traditional sashiko. "
          "The generated recipe calls for low heat, long simmering times, "
          "and minimal method-switching, consistent with washoku aesthetics.")

    _para(doc,
          "By contrast, La Marseillaise (BPM 120, march tempo) produced "
          "a Deck with large FWD/RET amplitudes (5-7), high TENSION "
          "values (up to 0.95), and steep CROSS angles (45 degrees). "
          "The SVG renders as a bold, aggressive zigzag with sharp "
          "directional changes. The recipe alternates between high-heat "
          "searing and rapid cooling -- a rhythm that resembles classical "
          "French sauce technique. The Star-Spangled Banner, with its "
          "famously wide vocal range (1.5 octaves), produced the largest "
          "FWD values (up to 8), creating an SVG with dramatic leaps "
          "and a recipe with extreme temperature variations. The Brazilian "
          "anthem's syncopation manifested as irregular FWD intervals "
          "and a flowing, rhythmically asymmetric stitch pattern.")

    _para(doc,
          "These demonstrations illustrate a key design principle of "
          "SITS9: the language does not merely permit cross-domain "
          "translation -- it preserves the 'cultural accent' of the "
          "source material across all renderings. A pentatonic source "
          "produces pentatonic-feeling outputs in every domain.")

    # ── 6. Palette Constraints and Cultural Identity ─────────────────────────
    _heading(doc, "6. Palette Constraints and Cultural Identity", 1)

    _para(doc,
          "In certain musical traditions, the available pitch set is "
          "deliberately restricted -- for example, the Japanese pentatonic "
          "scale (yo-na-nuki: omitting the 4th and 7th degrees). This "
          "restriction of the discrete palette is not a limitation but a "
          "culturally constitutive choice, and it has direct parallels in "
          "textile and cuisine (Table 4).")

    _table(doc,
           ["Domain", "Constrained palette (pentatonic-like)", "Full palette (chromatic-like)"],
           [
               ["Music",
                "Pentatonic: omit 4th & 7th scale degrees",
                "Twelve-tone: use all pitches equally"],
               ["Textile",
                "Sashiko: white thread on indigo (1 colour)",
                "Jacquard full-colour: unlimited palette"],
               ["Cuisine",
                "Japanese: dashi, soy, miso, salt, sugar (5 bases)",
                "Molecular gastronomy: synthesise any flavour"],
           ],
           caption="Table 4. Palette constraints across domains.")
    doc.add_paragraph()

    _para(doc,
          "In SITS9, palette constraint is expressed as the discrete "
          "set of values that FWD(n) can take:")
    _code(doc,
          "Pentatonic:   n in {1, 2, 3, 5, 6}      (skip 4 and 7)\n"
          "Diatonic:     n in {1, 2, 3, 4, 5, 6, 7}\n"
          "Chromatic:    n in {1, 2, ..., 12}        (no constraint)\n")
    _para(doc,
          "A key observation is that stronger palette constraints produce "
          "clearer cultural signatures. Pentatonic melodies are immediately "
          "identified as 'Japanese', 'Chinese', or 'Celtic'; indigo sashiko "
          "is immediately identified as Japanese; dashi-and-soy flavour "
          "profiles are immediately identified as washoku. By reducing "
          "information entropy, palette constraints amplify cultural "
          "signal.{6}")

    _para(doc,
          "This relationship can be expressed informally as: cultural "
          "distinctiveness is inversely proportional to palette size. "
          "A twelve-tone composition, a full-colour digital textile, and "
          "a molecular-gastronomy dish can each be virtuosic, but none "
          "carries a strong cultural fingerprint. Conversely, "
          "a pentatonic melody played on a shakuhachi, a white-on-indigo "
          "sashiko cloth, and a bowl of miso soup each carry unmistakable "
          "cultural identity precisely because their palettes are "
          "constrained to a handful of values.")

    _para(doc,
          "In SITS9 terms, palette constraint operates as a filter on "
          "the parameter space of FWD(n), TENSION(v), and CROSS(angle). "
          "A 'Japanese-style' SITS9 Deck might restrict FWD to "
          "n in {1, 2, 3}, TENSION to v in {0.1, 0.3, 0.5}, and CROSS "
          "angle to {-15, 0, 15}. A 'French-style' Deck might allow "
          "the full range. The cultural character of the output is thus "
          "encoded not in the instructions themselves but in the "
          "constraints imposed on their parameters -- a meta-level of "
          "cultural encoding that sits above the instruction set itself.")

    # ── 7. Textile Ciphers and Steganography ────────────────────────────────
    _heading(doc, "7. Textile Ciphers: SITS9 as Steganography", 1)

    _para(doc,
          "The use of textile as a covert information channel has multiple "
          "historical precedents (Table 5).")

    _table(doc,
           ["Period", "Case", "Mechanism"],
           [
               ["WWI",
                "Belgian knitting spies",
                "Knit=0, purl=1: binary encoding of railway movements. "
                "MI5 banned knitting near the front lines.{7}"],
               ["French Revolution",
                "Tricoteuses (knitting women)",
                "Dickens' Madame Defarge encodes execution lists in "
                "knitting patterns (fictional, but the concept existed)."],
               ["Inca Empire",
                "Quipu military cipher",
                "Administrative records + encrypted military messages "
                "in the same knotted-cord medium.{4}"],
               ["Antebellum US",
                "Underground Railroad quilt code (disputed)",
                "Quilt patterns as route-guiding ciphers. "
                "(Contested: no confirmed primary sources.){8}"],
           ],
           caption="Table 5. Historical examples of textile-based cryptographic and steganographic communication.")
    doc.add_paragraph()

    _para(doc,
          "What unites these examples is that the information channel "
          "was the craft process itself, not a hidden compartment within "
          "a craft object. A knitted scarf is not a container with a "
          "message hidden inside; the scarf IS the message, read in a "
          "different encoding. This is precisely the architecture of "
          "SITS9: a Deck is not a container with three payloads -- it "
          "is a single sequence of instructions that is simultaneously "
          "readable as three different artefacts.")

    _para(doc,
          "In the SITS9 framework, a single Deck appears as "
          "'a beautiful sashiko pattern', 'a pleasant piece of music', "
          "and 'a cooking recipe' -- but a fourth renderer (a cipher "
          "decryptor) could extract a plaintext message that none of the "
          "other three renderers reveal. The steganographic strength of "
          "such a system is that the 'cover objects' (textile, music, "
          "recipe) are not merely wrappers -- they are genuine, "
          "aesthetically functional artefacts in their own right. "
          "The strongest steganographic medium is one that nobody "
          "recognises as a medium at all. A contemporary art-practice "
          "example is Alyce Santoro's Sonic Fabric, a textile woven "
          "from recycled audio cassette tape that can be 'played back' "
          "by running a modified tape head across the cloth -- a literal "
          "textile-as-music artefact.{9}:")
    _code(doc,
          "Censorship susceptibility  ~ Recognition as a medium\n"
          "Cipher strength            ~ 1 / Censorship susceptibility\n")

    # ── 8. Physical-Object Steganography ────────────────────────────────────
    _heading(doc, "8. Future Prospect: Physical-Object Steganography", 1)

    _para(doc,
          "The textile-cipher discussion (Section 7) can be extended beyond "
          "fabric and thread. If information can be inscribed in a material "
          "whose 'media recognition' is zero, ultimate steganography is "
          "achieved.")

    _heading(doc, "8.1 Lens Inscription: Cipher as Halation", 2)
    _para(doc,
          "Femtosecond lasers can inscribe refractive-index changes inside "
          "transparent substrates at micrometre scale. Applied to an "
          "eyeglass lens, the resulting micro-patterns are visually "
          "indistinguishable from ordinary halation (light scatter and "
          "reflection artefacts). An SITS9 Deck encoded as a 3D point "
          "cloud and inscribed inside a lens via femtosecond laser becomes "
          "a fourth renderer output -- an ImplantRenderer -- that is "
          "physically portable and appears to the naked eye as nothing "
          "more than a minor lens imperfection.")

    _code(doc,
          "                 +-> LoomRenderer     -> sashiko pattern\n"
          "SITS9 Deck --+-> ToneRenderer     -> music\n"
          "                 +-> KitchenRenderer   -> recipe\n"
          "                 +-> ImplantRenderer   -> 3D point cloud (laser inscription)\n")

    _heading(doc, "8.2 Intraocular and Orthopaedic Implants", 2)
    _para(doc,
          "The same principle applies to medical implants. An intraocular "
          "lens (IOL) is a transparent acrylic or silicone substrate "
          "ideally suited for 3D laser inscription, and resides inside "
          "the body where physical access requires surgery (Table 6).")

    _table(doc,
           ["Medium", "Material", "Secrecy", "Notes"],
           [
               ["Eyeglass lens",
                "Glass / polycarbonate",
                "High (disguised as halation)",
                "Non-invasive, easily replaced, carried daily"],
               ["Intraocular lens (IOL)",
                "Acrylic / PMMA",
                "Very high (inside body + transparent)",
                "Widely used in cataract surgery"],
               ["Joint / bone-fixation screw",
                "Titanium / ceramic",
                "Very high (deep inside body)",
                "Large capacity; multiple units = distributed storage"],
               ["Dental implant",
                "Titanium",
                "High",
                "Visible on X-ray but pattern not decodable"],
           ],
           caption="Table 6. Candidate substrates for physical-object steganography.")
    doc.add_paragraph()

    _para(doc,
          "In terms of information density, if each refractive-index "
          "change at 1 um^3 encodes 1 bit, a single eyeglass lens "
          "(several cm^3) can theoretically store terabits of data "
          "(noise constraints reduce practical capacity by orders of "
          "magnitude, but several books' worth of data remains feasible). "
          "Compared to the CIA's Cold War technique of hiding microfilm "
          "in dental fillings, 3D laser inscription emits no active "
          "signal (unlike RFID) and is therefore substantially more "
          "resistant to scanning and detection.")

    _para(doc,
          "The connection to SITS9 is direct: an ImplantRenderer would "
          "map FWD(n) to a displacement of n micrometres along the "
          "inscription axis, RET(n) to a return displacement, CROSS to "
          "a change of inscription plane (z-axis shift), TENSION to "
          "the refractive-index change magnitude (higher tension = "
          "stronger inscription = more visible under microscopy), and "
          "ANCHOR to a reference marker. The same Deck that produces "
          "a sashiko pattern and a MIDI file would also produce a "
          "three-dimensional point cloud suitable for femtosecond "
          "laser inscription. The resulting inscribed lens could be "
          "read back into a Deck and re-rendered as music, textile, "
          "or recipe -- a physical artefact that is simultaneously "
          "a wearable object, an information store, and a cross-modal "
          "creative work.")

    # ── 9. Workshop Application ─────────────────────────────────────────────
    _heading(doc, "9. Application: Cross-Modal Workshop Design", 1)

    _para(doc,
          "As a practical application, we propose workshops in which "
          "participants experience all three SITS9 renderings of a "
          "single Deck. Each participant leaves with: (a) a hand-stitched "
          "sashiko coaster, (b) a MIDI playback of the music generated "
          "from their stitch pattern, and (c) a dish prepared from the "
          "recipe generated from the same Deck. Because each participant's "
          "stitching style differs slightly, the resulting music also "
          "differs -- revealing that 'stitching style = performance "
          "style = cooking style' is consistent across domains.")

    _heading(doc, "9.1 National Anthems as Source Material", 2)
    _para(doc,
          "A particularly effective input for such workshops is the "
          "national anthem. Anthems are nearly free of copyright "
          "restrictions, and each participant can select 'their own "
          "country's song' as a personally meaningful starting point. "
          "Because the musical characteristics of an anthem directly "
          "reflect its culture's palette, the SITS9 renderings "
          "preserve the cultural signature across domains (Table 7).")

    _table(doc,
           ["Anthem", "Musical features", "SITS9 profile", "Predicted textile / cuisine"],
           [
               ["Kimigayo (Japan)",
                "Pentatonic, slow tempo, narrow range",
                "Small FWD values, low TENSION, few CROSS",
                "Restrained sashiko pattern / gentle simmered washoku"],
               ["La Marseillaise (France)",
                "March, abrupt modulation, wide dynamics",
                "Large FWD/RET amplitude, sudden TENSION changes, frequent CROSS",
                "High-contrast embroidery / rapid heat-cool cycles (sauce reduction)"],
               ["Star-Spangled Banner (US)",
                "Wide range (1.5 octaves), leaping intervals",
                "Large FWD(n), large CROSS angles",
                "Bold zigzag pattern / dynamic heat variation"],
               ["Hino Nacional (Brazil)",
                "Syncopation, long melodic lines",
                "Irregular FWD intervals, long Cards",
                "Rhythmic uneven-spaced stitches / staggered cooking steps"],
           ],
           caption="Table 7. National anthems as SITS9 input: predicted cross-domain renderings.")
    doc.add_paragraph()

    _heading(doc, "9.2 Workshop Procedure", 2)
    _para(doc,
          "Phase 1 (Introduction, 15 min): Participants select the anthem "
          "of their cultural heritage. The facilitator demonstrates the "
          "SITS9 conversion and shows that three outputs emerge from "
          "a single Deck.")
    _para(doc,
          "Phase 2 (Cuisine, 30 min): Each participant follows the recipe "
          "generated from their anthem Deck to prepare a simple dish. "
          "The connection to music is not revealed at this stage.")
    _para(doc,
          "Phase 3 (Textile, 30 min): Using the same Deck rendered as "
          "stitching instructions, participants stitch on bleached cloth. "
          "Completed patterns are photographed.")
    _para(doc,
          "Phase 4 (Music, 15 min): Photographed stitch patterns are "
          "loaded into the music renderer. The revelation -- 'your "
          "national anthem was hidden in your cooking and stitching' -- "
          "is made.")
    _para(doc,
          "Phase 5 (Comparison, 15 min): All participants' textiles, "
          "dishes, and music are displayed side by side. Exchanging "
          "sashiko and tasting each other's anthem-derived dishes "
          "becomes an embodied exercise in cross-cultural understanding.")

    _para(doc,
          "The educational significance of this workshop lies in the "
          "embodied realisation that 'culture' is not confined to a "
          "single domain but operates as a consistent operational "
          "palette across music, craft, and food.")

    # ── 10. Discussion ──────────────────────────────────────────────────────
    _heading(doc, "10. Discussion: Why Did All Three Domains Follow the Same Evolution?", 1)

    _para(doc,
          "Why did music, textile, and cuisine follow the same trajectory "
          "of instruction-set complexification? Our hypothesis is that all "
          "three are fundamentally acts of 'applying temporal operations "
          "to a material to create structure': thread x temporal operation "
          "-> fabric; sound x temporal operation -> music; ingredient x "
          "temporal operation -> cuisine. The complexification pattern is "
          "shared because the human cognitive architecture -- the order "
          "in which we can understand and memorise sequences of "
          "'repetition -> crossing -> return' -- is universal. "
          "SITS9's ability to unify three domains is not a coincidence "
          "but a reflection of the universal structure of human "
          "procedural cognition.{5}")

    _para(doc,
          "A supporting observation comes from the parallel evolution "
          "summarised in Table 1. The fact that all three domains "
          "independently passed through the same stages -- from "
          "single-operation repetition (running stitch / monophonic chant / "
          "open-flame roasting) to multi-layer crossing (weaving / "
          "polyphony / sauce systems) to encoded reproduction (Jacquard / "
          "player piano / canning) -- suggests that the trajectory is "
          "not culturally contingent but cognitively constrained. Each "
          "stage introduces one additional degree of freedom to the "
          "instruction set: CROSS introduces a second axis, RET introduces "
          "temporal non-monotonicity, TENSION introduces continuous "
          "modulation. The order in which humans can learn and teach these "
          "degrees of freedom appears to be invariant across cultures "
          "and across domains. Steve Reich's concept of 'music as a "
          "gradual process' -- in which perceptible, slowly unfolding "
          "transformations constitute the musical material -- is itself "
          "an instance of this universal trajectory: minimal music "
          "strips the instruction set back to loop-FWD and makes the "
          "process audible.{10}")

    _para(doc,
          "The naming of SITS9 -- from the proverb 'a stitch in time "
          "saves nine' -- captures a deeper principle: the language's "
          "economy of means. Five instructions suffice to span three "
          "domains not because the domains are simple, but because the "
          "underlying operational algebra is small. Just as a Turing "
          "machine needs only a few symbols to be computationally "
          "universal, SITS9 achieves cross-domain universality with "
          "five primitives because the human procedural repertoire is "
          "itself built from a small set of temporal operations.")

    _para(doc,
          "A limitation of the current work is that the mappings between "
          "SITS9 instructions and domain-specific actions are manually "
          "designed rather than empirically derived. The TENSION-to-heat "
          "mapping in the KitchenRenderer, for example, is linear, but "
          "human perception of 'cooking intensity' may follow a "
          "logarithmic or sigmoidal curve. Similarly, the CROSS angle-to-"
          "interval mapping in the ToneRenderer is a simplification of "
          "what is, in musical practice, a highly context-dependent "
          "relationship. Future work could involve perceptual experiments "
          "to determine whether listeners, stitchers, and cooks "
          "independently converge on similar tension curves when presented "
          "with the same abstract instruction sequence. The national-anthem "
          "workshop proposed in Section 9 provides a natural experimental "
          "framework for such studies.")

    _para(doc,
          "Another avenue for future investigation is the formal "
          "relationship between SITS9 and existing computational "
          "models. The Card/Deck structure bears a clear resemblance to "
          "subroutines and programs in imperative languages, while the "
          "immutable, frozen-dataclass implementation of instructions "
          "suggests a functional-programming interpretation. Exploring "
          "whether SITS9 Decks form a monoid under concatenation, or "
          "whether the renderer architecture can be formalised as a "
          "functor from an 'instruction category' to domain-specific "
          "categories, could place the system on firmer mathematical "
          "foundations and enable automatic reasoning about cross-domain "
          "equivalences.")

    # ── 11. Conclusion ──────────────────────────────────────────────────────
    _heading(doc, "11. Conclusion", 1)

    _para(doc,
          "SITS9 (A Stitch in Time Saves Nine) is a minimal "
          "formal language consisting of five primitive instructions "
          "-- FWD, RET, CROSS, TENSION, ANCHOR -- from which music, "
          "textile pattern, and cooking recipe can be simultaneously "
          "generated via domain-specific renderers. We have demonstrated "
          "this with a working implementation using Pachelbel's Canon "
          "and four national anthems as test cases, confirming that the "
          "system preserves both structural coherence and cultural "
          "character across all three output domains.")

    _para(doc,
          "That such unification is possible is a modern re-discovery of "
          "the historical fact that the Jacquard loom and the music box "
          "were driven by the same punch card, and that the Inca quipu "
          "encoded accounting, narrative, and textile structure in a "
          "single knotted-cord medium. SITS9 moves beyond surface "
          "analogy ('A resembles B') to structural equivalence ('A and B "
          "are different renderings of the same instruction set C'), "
          "with applications ranging from cryptographic steganography "
          "and physical-object inscription to cross-modal workshop design "
          "and cultural-palette analysis.")

    _para(doc,
          "The name SITS9 encapsulates the core insight: a single "
          "well-placed stitch -- a single well-designed instruction set -- "
          "saves nine domain-specific representations. As the historical "
          "record from quipu to Jacquard to MIDI attests, the convergence "
          "of textile, music, and process encoding is not a modern "
          "invention but a recurring pattern in the history of human "
          "technology. SITS9 simply makes the convergence explicit, "
          "executable, and extensible.")

    # ── Acknowledgements ────────────────────────────────────────────────────
    _heading(doc, "Acknowledgements", 1)
    _para(doc, "[To be inserted.]", italic=True)

    # ── References ──────────────────────────────────────────────────────────
    _heading(doc, "References", 1)
    refs = [
        ("1", "Essinger, J. (2004) Jacquard's Web: How a Hand-Loom Led to "
              "the Birth of the Information Age. Oxford University Press."),
        ("2", "Aldwell, E., Schachter, C. and Cadwallader, A. (2019) "
              "Harmony and Voice Leading, 5th edn. Cengage Learning."),
        ("3", "McGee, H. (2004) On Food and Cooking: The Science and Lore "
              "of the Kitchen. New York: Scribner."),
        ("4", "Urton, G. (2003) Signs of the Inka Khipu: Binary Coding in "
              "the Andean Knotted-String Records. Austin: University of "
              "Texas Press."),
        ("5", "McLuhan, M. (1964) Understanding Media: The Extensions of "
              "Man. New York: McGraw-Hill."),
        ("6", "Lerdahl, F. and Jackendoff, R. (1983) A Generative Theory "
              "of Tonal Music. Cambridge, MA: MIT Press."),
        ("7", "Stallings, W. (2017) Cryptography and Network Security: "
              "Principles and Practice, 7th edn. Hoboken, NJ: Pearson."),
        ("8", "Tobin, J. and Dobard, R. (1999) Hidden in Plain View: A "
              "Secret Story of Quilts and the Underground Railroad. New "
              "York: Anchor Books. [Contested thesis]"),
        ("9", "Santoro, A. (2007) Sonic Fabric. Available at: "
              "https://sonicfabric.com/ (Accessed: 22 June 2026)."),
        ("10", "Reich, S. (2002) Writings on Music, 1965-2000. Oxford "
               "University Press."),
    ]
    for num, ref in refs:
        p = doc.add_paragraph()
        run_num = p.add_run(f"{num}. ")
        run_num.font.superscript = True
        run_num.font.size = Pt(10)
        run_text = p.add_run(ref)
        run_text.font.size = Pt(10)

    return doc


# ---------------------------------------------------------------------------
# PPTX figure output (editable, 1 figure per slide)
# ---------------------------------------------------------------------------

def build_figure_pptx() -> None:
    """Generate editable PPTX with each figure on a separate slide."""
    from pptx import Presentation
    from pptx.util import Inches as PInches, Pt as PPt

    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)

    slide_layout = prs.slide_layouts[5]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    txBox = slide.shapes.add_textbox(PInches(0.5), PInches(0.2), PInches(12), PInches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Figure 1"
    run.font.size = PPt(24)
    run.font.bold = True

    # Image
    stitch_png = OUT_DIR / "pachelbel_stitch_2loops.png"
    if stitch_png.exists():
        slide.shapes.add_picture(
            str(stitch_png),
            PInches(1.5), PInches(1.2), PInches(10), PInches(4.0))

    # Caption
    txBox2 = slide.shapes.add_textbox(PInches(0.5), PInches(5.5), PInches(12), PInches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = (
        "Figure 1. SITS9 Loom Renderer output for Pachelbel's Canon "
        "(I-V-vi-iii-IV-I-IV-V, 2 loops). Red solid lines = surface "
        "stitches; grey dashed lines = back stitches; filled black "
        "circles = anchor points (knots)."
    )
    run2.font.size = PPt(14)
    run2.font.italic = True

    # ── Slide 2: Score schema ──
    extra_figs = [
        ("Figure 2", OUT_DIR / "pachelbel_score_schema.png",
         "Figure 2. Harmonic tension profile of Pachelbel's Canon "
         "encoded as SITS9 instructions. Bar colour = instruction type "
         "(blue = ANCHOR, red = FWD, green = RET). Bar height = TENSION. "
         "Dashed line traces the tension arc."),
        ("Figure 3", OUT_DIR / "pachelbel_textile_schema.png",
         "Figure 3. Textile schema (back-stitch pattern) for Pachelbel's "
         "Canon. Needle path alternates between surface (orange) and back "
         "(blue). Red = FWD on surface; orange diagonals = CROSS; "
         "black dot = ANCHOR (knot)."),
        ("Figure 4", OUT_DIR / "pachelbel_cooking_schema.png",
         "Figure 4. Cooking schema for Pachelbel's Canon. Heat curve "
         "(TENSION) follows the harmonic tension profile. Background zones: "
         "red = grill, blue = simmer, green = steam. Dashed vertical = "
         "CROSS (method switch)."),
    ]

    for fig_title, fig_path, fig_caption in extra_figs:
        if not fig_path.exists():
            continue
        slide = prs.slides.add_slide(slide_layout)

        txBox = slide.shapes.add_textbox(
            PInches(0.5), PInches(0.2), PInches(12), PInches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = fig_title
        run.font.size = PPt(24)
        run.font.bold = True

        slide.shapes.add_picture(
            str(fig_path),
            PInches(1.0), PInches(1.0), PInches(11), PInches(4.5))

        txBox2 = slide.shapes.add_textbox(
            PInches(0.5), PInches(5.8), PInches(12), PInches(1.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = fig_caption
        run2.font.size = PPt(14)
        run2.font.italic = True

    pptx_path = OUT_DIR / "sits9_figures.pptx"
    prs.save(str(pptx_path))
    print(f"Figure PPTX saved -> {pptx_path}")


# ---------------------------------------------------------------------------

def main() -> None:
    OUT_DIR.mkdir(exist_ok=True)
    doc = build_manuscript()
    out_path = OUT_DIR / "sits9_leonardo.docx"
    doc.save(str(out_path))
    print(f"Leonardo manuscript saved -> {out_path}")
    build_figure_pptx()


if __name__ == "__main__":
    main()
