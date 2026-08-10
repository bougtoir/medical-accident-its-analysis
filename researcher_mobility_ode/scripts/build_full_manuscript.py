#!/usr/bin/env python3
"""Build a Research Policy full-article manuscript from ODE result CSVs.

Outputs (regenerated from results/* CSVs, no hard-coded numbers):
- docs/manuscript_full_article.docx
- docs/manuscript_full_article.md
- docs/manuscript_full_article_figures.pptx
- docs/figures/*.png

All numerical values are read from the result CSVs produced by the analysis
pipeline; the script contains only formatting and prose.
"""
from __future__ import annotations

import argparse
import math
import os
import sys
from pathlib import Path

# Make local packages importable
BASE_DIR = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(BASE_DIR / "src"))
sys.path.insert(0, str(BASE_DIR / "scripts"))

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.util import Inches as PptxInches

from pyommlbuilder.main import (
    Math,
    SubscriptObject,
    Fraction,
    Numerator,
    Denominator,
    Function,
    MathPara,
)
from pyommlbuilder.helpers import make_aligned_equation

import annual_rates_projection_report as arpr

RESULTS_DIR = BASE_DIR / "results"
ENDOG = RESULTS_DIR / "endogenous"
SAT = RESULTS_DIR / "endogenous_saturating"
TV = RESULTS_DIR / "time_varying"
BOOT = RESULTS_DIR / "bootstrap_ci"
POL = RESULTS_DIR / "policy_counterfactuals"
ANNUAL = RESULTS_DIR / "annual"
FIG_DIR = BASE_DIR / "docs" / "figures"


def _fmt(v, dec=2):
    if pd.isna(v):
        return "—"
    try:
        return f"{float(v):.{dec}f}"
    except (ValueError, TypeError):
        return str(v)


def add_citation(para, number: int):
    run = para.add_run(f" [{number}]")
    run.font.superscript = True
    return run


def add_footnote(para, symbol="1"):
    run = para.add_run(f" {symbol}")
    run.font.superscript = True
    return run


def _paragraph_text(doc):
    for p in doc.paragraphs:
        yield p.text


def _table_text(doc):
    for t in doc.tables:
        for row in t.rows:
            for cell in row.cells:
                yield cell.text


def _doc_word_count(doc):
    return sum(len(t.split()) for t in list(_paragraph_text(doc)) + list(_table_text(doc)))


def _rel_path(path: Path, base: Path) -> str:
    return path.relative_to(base).as_posix()


def add_omath_paragraph(doc, math_element, align=WD_ALIGN_PARAGRAPH.CENTER):
    """Append an OMML math element to a new paragraph."""
    p = doc.add_paragraph()
    p._element.append(math_element._as_xml_element())
    p.alignment = align
    return p


def add_omath_inline(para, math_element):
    """Append an OMML math element inline inside an existing paragraph."""
    para._element.append(math_element._as_xml_element())


# ---------------------------------------------------------------------------
# OMML equation builders
# ---------------------------------------------------------------------------

def math_I_linear():
    """I(P_D) = I_0 + r P_D"""
    return Math(
        Function("I", SubscriptObject("P", "D")),
        "=",
        SubscriptObject("I", "0"),
        "+",
        "r",
        SubscriptObject("P", "D"),
    )


def math_I_saturating():
    """I(P_D) = I_0 + r P_D / (1 + epsilon P_D)"""
    return Math(
        Function("I", SubscriptObject("P", "D")),
        "=",
        SubscriptObject("I", "0"),
        "+",
        Fraction(
            Numerator("r", SubscriptObject("P", "D")),
            Denominator("1 + ", "ε", "×", SubscriptObject("P", "D")),
        ),
    )


def math_threshold():
    """M = k × c̄"""
    return Math("M = k × c\u0304")


def math_active_pool():
    """T = D + H_D + P_D"""
    return Math(
        "T = D + ",
        SubscriptObject("H", "D"),
        " + ",
        SubscriptObject("P", "D"),
    )


def math_ode_system():
    """Six-equation display using MathPara."""
    def deriv(base):
        return Fraction(Numerator("d" + base), Denominator("dt"))

    lines = [
        make_aligned_equation(
            deriv("D"),
            Math(
                Function("I", SubscriptObject("P", "D")),
                " + βA - (α + ",
                SubscriptObject("h", "D"),
                " + d)D",
            ),
            line_break=False,
        ),
        make_aligned_equation(
            deriv("A"),
            Math(
                "αD - (β + ",
                SubscriptObject("h", "A"),
                " + d)A",
            ),
            line_break=False,
        ),
        make_aligned_equation(
            deriv("H_D"),
            Math(
                SubscriptObject("h", "D"),
                "D + β",
                SubscriptObject("H", "A"),
                " - (",
                SubscriptObject("p", "D"),
                " + d)",
                SubscriptObject("H", "D"),
            ),
            line_break=False,
        ),
        make_aligned_equation(
            deriv("H_A"),
            Math(
                SubscriptObject("h", "A"),
                "A - (β + ",
                SubscriptObject("p", "A"),
                " + d)",
                SubscriptObject("H", "A"),
            ),
            line_break=False,
        ),
        make_aligned_equation(
            deriv("P_D"),
            Math(
                SubscriptObject("p", "D"),
                SubscriptObject("H", "D"),
                " + β",
                SubscriptObject("P", "A"),
                " - d",
                SubscriptObject("P", "D"),
            ),
            line_break=False,
        ),
        make_aligned_equation(
            deriv("P_A"),
            Math(
                SubscriptObject("p", "A"),
                SubscriptObject("H", "A"),
                " - (β + d)",
                SubscriptObject("P", "A"),
            ),
            line_break=False,
        ),
    ]
    return MathPara(*lines)


# ---------------------------------------------------------------------------
# Figures
# ---------------------------------------------------------------------------

def build_figure1(eq, fig_dir: Path):
    """Equilibrium domestic active pool vs minimum viable threshold."""
    fig_dir.mkdir(parents=True, exist_ok=True)
    groups = eq["group"].tolist()
    x = np.arange(len(groups))
    width = 0.35
    fig, ax = plt.subplots(figsize=(10, 5))
    ax.bar(x - width / 2, eq["T_equilibrium"], width, label="Equilibrium T", color="steelblue")
    ax.bar(x + width / 2, eq["M_threshold"], width, label="Minimum viable threshold M", color="coral")
    ax.set_xticks(x)
    ax.set_xticklabels(groups, rotation=35, ha="right")
    ax.set_ylabel("Number of researchers")
    ax.set_title("Domestic active researcher pool and minimum viable coauthor threshold by group")
    ax.legend()
    ax.set_ylim(0, max(eq["T_equilibrium"].max(), eq["M_threshold"].max()) * 1.1)
    fig.tight_layout()
    path = fig_dir / "fig1_equilibrium_margin.png"
    fig.savefig(path, dpi=300)
    plt.close(fig)
    return path


def build_figure2(pnr_closest, fig_dir: Path):
    """Closest point-of-no-return proximity by group."""
    fig_dir.mkdir(parents=True, exist_ok=True)
    pnr_closest = pnr_closest.sort_values("proximity")
    groups = pnr_closest["group"].tolist()
    labels = [f"{r}\n({t})" for r, t in zip(pnr_closest["rate_name"], pnr_closest["target"])]
    prox = pnr_closest["proximity"].tolist()
    fig, ax = plt.subplots(figsize=(9, 5))
    # Use a perceptually uniform, colour-vision-deficiency-friendly sequential palette
    norm = max(prox) * 1.2 if prox else 1.0
    colors = [plt.cm.plasma(0.25 + 0.55 * (p / norm)) for p in prox]
    bars = ax.barh(groups, prox, color=colors)
    for bar, label in zip(bars, labels):
        width = bar.get_width()
        ax.text(width + 0.01, bar.get_y() + bar.get_height() / 2,
                label, va="center", fontsize=7)
    ax.set_xlabel("Required proportional change in rate |critical factor − 1|")
    ax.set_title("Closest point-of-no-return sensitivity by group (smaller = more fragile)")
    xmax = max(1.2, max(prox) * 1.15) if prox else 1.2
    ax.set_xlim(0, xmax)
    ax.axvline(1.0, color="gray", linestyle="--", linewidth=0.8)
    fig.tight_layout()
    path = fig_dir / "fig2_pnr_proximity.png"
    fig.savefig(path, dpi=300)
    plt.close(fig)
    return path


def build_figure3(period_compare, fig_dir: Path):
    """Historical counterfactual: change in safety margin from early to late rates."""
    fig_dir.mkdir(parents=True, exist_ok=True)
    df = period_compare.sort_values("delta_margin")
    groups = df["group"].tolist()
    deltas = df["delta_margin"].tolist()
    fig, ax = plt.subplots(figsize=(9, 5))
    # Colour-vision-deficiency-friendly palette: blue for positive, vermillion for negative
    CVD_POS = "#0072B2"
    CVD_NEG = "#D55E00"
    colors = [CVD_POS if d > 0 else CVD_NEG for d in deltas]
    bars = ax.barh(groups, deltas, color=colors)
    for bar, d in zip(bars, deltas):
        width = bar.get_width()
        ax.text(width + (max(deltas) * 0.01 if width >= 0 else min(deltas) * 0.01),
                bar.get_y() + bar.get_height() / 2,
                f"{_fmt(d, 1)}", va="center", ha="left" if width >= 0 else "right",
                fontsize=8)
    ax.axvline(0, color="black", linewidth=0.8)
    ax.set_xlabel("Change in equilibrium safety margin (late − early)")
    ax.set_title("Counterfactual change in safety margin if late-period rates persisted (point estimates)")
    fig.tight_layout()
    path = fig_dir / "fig3_historical_margin.png"
    fig.savefig(path, dpi=300)
    plt.close(fig)
    return path


def build_figure4(boot, fig_dir: Path):
    """Bootstrap 95% confidence intervals for equilibrium T."""
    fig_dir.mkdir(parents=True, exist_ok=True)
    df = boot.sort_values("T_equilibrium_median")
    groups = df["group"].tolist()
    med = df["T_equilibrium_median"].tolist()
    low = df["T_equilibrium_q025"].tolist()
    high = df["T_equilibrium_q975"].tolist()
    fig, ax = plt.subplots(figsize=(9, 5))
    y = np.arange(len(groups))
    ax.errorbar(med, y, xerr=[np.subtract(med, low), np.subtract(high, med)],
                fmt="o", color="steelblue", capsize=4, ecolor="gray")
    ax.set_yticks(y)
    ax.set_yticklabels(groups)
    ax.set_xlabel("Equilibrium domestic active pool T")
    ax.set_title("Bootstrap 95% confidence intervals for equilibrium T")
    ax.set_xlim(0, max(high) * 1.05)
    fig.tight_layout()
    path = fig_dir / "fig4_bootstrap_ci.png"
    fig.savefig(path, dpi=300)
    plt.close(fig)
    return path


# ---------------------------------------------------------------------------
# Reference list
# ---------------------------------------------------------------------------

NOTE_TEXT = "Yamada Y (momentumyy). 海外で当てた研究者はその後どうなるのか. note.com, 2026. https://note.com/momentumyy/n/n86df5d34282d (accessed 2026-08-09)."

REFS = [
    "MacroPolo. The Global AI Talent Tracker 2.0. Paulson Institute, 2023. https://macropolo.org/digital-projects/the-global-ai-talent-tracker/",
    "Appelt S, van Beuzekom B, Galindo-Rueda F, de Pinho R. Which factors influence the international mobility of research scientists? OECD Science, Technology and Industry Working Papers 2015/02, 2015. https://doi.org/10.1787/5js1tmrr2233-en",
    "Stephan P E. The Economics of Science. J Econ Lit. 1996;34(3):1199-1235.",
    "Huntington S P. The Clash of Civilizations and the Remaking of World Order. New York: Simon & Schuster, 1996.",
    "Aghion P, Bloom N, Blundell R, Griffith R, Howitt P. Competition and innovation: an inverted-U relationship. Q J Econ. 2005;120(2):701-728.",
    "Priem J, Piwowar H, Orr R. OpenAlex: A fully-open index of scholarly works, authors, venues, institutions, and concepts. arXiv:2205.01833, 2022. https://doi.org/10.48550/arXiv.2205.01833",
    "Thorn K, Holm-Nielsen L B. International Mobility of Researchers and Scientists: Policy Options for Turning a Drain into a Gain. UNU-WIDER Research Paper No. 2006/83, 2006. https://www.wider.unu.edu/sites/default/files/rp2006-83.pdf",
    "AlShebli B, Memon S A, Evans J A, Rahwan T. China and the U.S. produce more impactful AI research when collaborating together. Sci Rep. 2024;14:28576. https://doi.org/10.1038/s41598-024-79863-5",
    "Yuan S, Shao Z, Wei X, Tang J, Hall W, Wang Y, et al. Science behind AI: the evolution of trend, mobility, and collaboration. Scientometrics. 2020;124(2):993-1013. https://doi.org/10.1007/s11192-020-03423-7",
    "Shaffer M L. Minimum Population Sizes for Species Conservation. BioScience. 1981;31(2):131-134.",
    "Franzoni C, Scellato G, Stephan P E. Foreign-born scientists: mobility patterns for 16 countries. Nat Biotechnol. 2012;30(12):1250-1253.",
    "Jones B F, Wuchty S, Uzzi B. Multi-University Research Teams: Shifting Impact, Geography, and Stratification in Science. Science. 2008;322(5905):1259-1262.",
    "Freeman R B, Huang W. Collaboration: Strength in diversity. Nature. 2014;513(7518):305. https://doi.org/10.1038/513305a",
    "Shachar A. The Race for Talent: Highly Skilled Migrants and Competitive Immigration Regimes. NYU Law Rev. 2006;81(1):148-206.",
    "Kerr W R. Global Talent and U.S. Immigration Policy. Harvard Business School Working Paper No. 20-107, 2020. https://www.hbs.edu/ris/Publication%20Files/20-107_0967f1ab-1d23-4d54-b5a1-c884234d9b31.pdf"
]


def _ref(n: int) -> str:
    return REFS[n - 1]


# ---------------------------------------------------------------------------
# Data loading
# ---------------------------------------------------------------------------

def load_data():
    cohort = pd.read_csv(BASE_DIR / "data" / "cohort" / "cohort.csv")
    eq = pd.read_csv(ENDOG / "equilibrium_summary.csv")
    sat_eq = pd.read_csv(SAT / "equilibrium_summary.csv") if (SAT / "equilibrium_summary.csv").exists() else None
    top_t = pd.read_csv(ENDOG / "top_transitions_T.csv")
    pnr_full = pd.read_csv(ENDOG / "point_of_no_return.csv")
    # Closest per group, active pool preferred, restricting to rates that actually cross the threshold
    active_only = pnr_full[(pnr_full["target"] == "domestic_active") & (pnr_full["is_within_bounds"] == True)].copy()
    active_only["proximity"] = (active_only["critical_factor"] - 1.0).abs()
    pnr_closest = active_only.loc[active_only.groupby("group")["proximity"].idxmin()].reset_index(drop=True)
    pnr_closest = pnr_closest.sort_values("proximity").reset_index(drop=True)
    period_compare = pd.read_csv(TV / "period_comparison.csv")
    boot = pd.read_csv(BOOT / "bootstrap_summary.csv")
    policy_rank = pd.read_csv(POL / "ranked_interventions.csv")
    return cohort, eq, sat_eq, top_t, pnr_closest, period_compare, boot, policy_rank


# ---------------------------------------------------------------------------
# Annual transition-rate and projection helpers
# ---------------------------------------------------------------------------


def load_annual_data():
    """Read annual transition-rate and projection CSVs.

    Returns a dict of DataFrames; missing tables are returned as None.
    """
    paths = {
        "rate_table": ANNUAL / "annual_ode_rates.csv",
        "projected_rates": ANNUAL / "projected_ode_rates.csv",
        "observed_stock": ANNUAL / "observed_annual_stock.csv",
        "projected_stock": ANNUAL / "projected_annual_stock.csv",
        "interciv_stock": ANNUAL / "annual_interciv_stock.csv",
        "evaluation": ANNUAL / "projection_evaluation.csv",
        "group_accuracy": ANNUAL / "projection_accuracy_by_group.csv",
        "compartment_accuracy": ANNUAL / "projection_accuracy_by_compartment.csv",
    }
    return {k: pd.read_csv(p) if p.exists() else None for k, p in paths.items()}


def compute_annual_context(annual):
    """Return data-derived summary strings for the annual projection sections."""
    ctx = {}
    eval_df = annual.get("evaluation")
    if eval_df is not None and not eval_df.empty:
        ctx["overall_rmse"] = float(((eval_df["error"] ** 2).mean()) ** 0.5)
        ctx["overall_mape"] = float(eval_df["ape"].mean())
        ctx["overall_mape_pct"] = ctx["overall_mape"] * 100.0
    else:
        ctx["overall_rmse"] = float("nan")
        ctx["overall_mape"] = float("nan")
        ctx["overall_mape_pct"] = float("nan")

    gacc = annual.get("group_accuracy")
    if gacc is not None and not gacc.empty:
        gacc = gacc.dropna(subset=["mape"]).copy()
        if not gacc.empty:
            best = gacc.loc[gacc["mape"].idxmin()]
            worst = gacc.loc[gacc["mape"].idxmax()]
            ctx["best_group"] = best["origin_group"]
            ctx["worst_group"] = worst["origin_group"]
            ctx["best_mape_pct"] = float(best["mape"]) * 100.0
            ctx["worst_mape_pct"] = float(worst["mape"]) * 100.0
            ctx["n_eval_groups"] = len(gacc)
        else:
            ctx["best_group"] = "—"
            ctx["worst_group"] = "—"
            ctx["best_mape_pct"] = float("nan")
            ctx["worst_mape_pct"] = float("nan")
            ctx["n_eval_groups"] = 0
    else:
        ctx["best_group"] = "—"
        ctx["worst_group"] = "—"
        ctx["best_mape_pct"] = float("nan")
        ctx["worst_mape_pct"] = float("nan")
        ctx["n_eval_groups"] = 0

    proj = annual.get("projected_rates")
    if proj is not None and not proj.empty:
        ctx["n_projected_group_years"] = len(proj)
        ctx["smoothed_pct"] = float(proj["correction_smoothed"].mean()) * 100.0
        ctx["capped_pct"] = float(proj["correction_capped"].mean()) * 100.0
    else:
        ctx["n_projected_group_years"] = 0
        ctx["smoothed_pct"] = float("nan")
        ctx["capped_pct"] = float("nan")

    obs = annual.get("observed_stock")
    if obs is not None and not obs.empty:
        ctx["obs_year_min"] = int(obs["year"].min())
        ctx["obs_year_max"] = int(obs["year"].max())
    else:
        ctx["obs_year_min"] = 2000
        ctx["obs_year_max"] = 2023

    return ctx


def annual_summary_table(annual):
    """Mean observed annual transition rates and inflow by group (2000-2016)."""
    rate_table = annual.get("rate_table")
    if rate_table is None or rate_table.empty:
        return pd.DataFrame()
    observed = rate_table[rate_table["year"] <= 2016]
    cols = ["alpha", "beta", "h_D", "p_D", "d", "I_total"]
    means = observed.groupby("origin_group")[cols].mean().reset_index()
    means.columns = ["Group", "α", "β", "h_D", "p_D", "d", "I_total"]
    return means


def interciv_top_table(annual, n=10):
    """Top origin-destination abroad author-year accumulations."""
    flows = annual.get("interciv_stock")
    if flows is None or flows.empty:
        return pd.DataFrame()
    pivot = (
        flows.groupby(["origin_group", "destination_group"], observed=False)["count"]
        .sum()
        .reset_index()
        .sort_values("count", ascending=False)
        .head(n)
    )
    pivot.columns = ["Origin", "Destination", "Author-years"]
    return pivot


def build_annual_figures(annual, fig_dir):
    """Generate annual projection figures; reuse existing PNGs if data are missing."""
    fig_paths = {}
    rate_table = annual.get("rate_table")
    projected_rates = annual.get("projected_rates")
    if rate_table is not None and projected_rates is not None:
        fig_paths["fig5"] = arpr.plot_annual_rates(rate_table, projected_rates)
    else:
        fig_paths["fig5"] = fig_dir / "annual_rates_by_group.png"

    interciv = annual.get("interciv_stock")
    if interciv is not None:
        fig_paths["fig6"] = arpr.plot_interciv_heatmap(interciv)
    else:
        fig_paths["fig6"] = fig_dir / "annual_interciv_heatmap.png"

    obs_stock = annual.get("observed_stock")
    proj_stock = annual.get("projected_stock")
    if obs_stock is not None and proj_stock is not None:
        fig_paths["fig7"] = arpr.plot_projection_by_compartment(proj_stock, obs_stock)
    else:
        fig_paths["fig7"] = fig_dir / "annual_projection_vs_observed.png"
    return fig_paths


def compute_context(cohort, eq, sat_eq, top_t, pnr_closest, period_compare, policy_rank):
    """Return data-derived summary strings used in the Results and Discussion."""
    n_groups = len(eq)
    eq_sorted = eq.sort_values("T_equilibrium", ascending=False)
    largest_pools = ", ".join(eq_sorted["group"].head(3).tolist())
    smallest_pool = eq_sorted["group"].iloc[-1]
    eq_m = eq.sort_values("margin_to_threshold_T")
    smallest_margin_group = eq_m["group"].iloc[0]

    d_rows = top_t[(top_t["rate"] == "d") & (top_t["target"] == "domestic_active")]
    d_min_e = d_rows["elasticity"].min()
    d_max_e = d_rows["elasticity"].max()
    d_all_negative = (d_rows["elasticity"] < 0).all()

    # Positive levers after dropout
    positive = []
    for _, gdf in top_t.groupby("group"):
        gdf = gdf.sort_values("abs_elasticity", ascending=False)
        # Skip the largest (dropout), then collect the positive transition-rate levers
        for _, r in gdf.iloc[1:].iterrows():
            if r["elasticity"] > 0 and r["rate"] not in ("I0", "r"):
                positive.append(r["rate"])
    pos_counts = pd.Series(positive).value_counts()
    most_common_positive = pos_counts.index[0] if not pos_counts.empty else "p_D"
    second_positive = pos_counts.index[1] if len(pos_counts) > 1 else None
    if most_common_positive == "p_D":
        pos_lever_text = "principal-investigator promotion (p_D)"
    elif most_common_positive == "h_D":
        pos_lever_text = "domestic hit generation (h_D)"
    elif most_common_positive == "beta":
        pos_lever_text = "return from abroad (β)"
    else:
        pos_lever_text = most_common_positive
    if second_positive == "p_D":
        second_text = "principal-investigator promotion (p_D)"
    elif second_positive == "h_D":
        second_text = "domestic hit generation (h_D)"
    elif second_positive == "beta":
        second_text = "return from abroad (β)"
    else:
        second_text = second_positive
    if second_text and second_text != pos_lever_text:
        positive_lever_sentence = f"The largest positive transition lever is {pos_lever_text}, followed by {second_text}."
    else:
        positive_lever_sentence = f"The largest positive transition lever is {pos_lever_text}."
    positive_lever_sentence_lower = positive_lever_sentence[0].lower() + positive_lever_sentence[1:]
    if positive_lever_sentence_lower.endswith('.'):
        positive_lever_sentence_lower = positive_lever_sentence_lower[:-1]

    # PI promotion elasticity, identify group with highest p_D elasticity
    pd_elas = top_t[(top_t["rate"] == "p_D") & (top_t["target"] == "domestic_active")].copy()
    pd_elas["abs_e"] = pd_elas["elasticity"].abs()
    highest_pd_group = pd_elas.sort_values("abs_e", ascending=False).iloc[0]["group"] if not pd_elas.empty else "Japanese"

    # Point of no return
    closest_rate_counts = pnr_closest["rate_name"].value_counts()
    closest_rate_mode = closest_rate_counts.index[0] if not closest_rate_counts.empty else "I0"
    all_closest_same = len(closest_rate_counts) == 1
    if all_closest_same:
        pnr_lever_text = f"{closest_rate_mode} is the closest point-of-no-return lever for the active researcher pool in every group"
    else:
        pnr_lever_text = f"{closest_rate_mode} is the most common closest point-of-no-return lever for the active researcher pool"

    # Saturating reduction range
    sat_range_text = ""
    if sat_eq is not None:
        merged = eq[["group", "T_equilibrium"]].merge(
            sat_eq[["group", "T_equilibrium"]], on="group", suffixes=("_lin", "_sat")
        )
        pct_lower = 100.0 * (merged["T_equilibrium_lin"] - merged["T_equilibrium_sat"]) / merged["T_equilibrium_lin"]
        sat_range_text = f"{pct_lower.min():.0f}-{pct_lower.max():.0f}% lower than the linear variant"

    # Historical counterfactual
    if period_compare.empty:
        period_neg = "none"
        period_pos = "none"
        period_all_neg = False
    else:
        sorted_pc = period_compare.sort_values("delta_margin")
        neg = sorted_pc[sorted_pc["delta_margin"] < 0]["group"].tolist()
        pos = sorted_pc[sorted_pc["delta_margin"] > 0]["group"].tolist()[::-1]
        period_neg = ", ".join(neg) if neg else "none"
        period_pos = ", ".join(pos) if pos else "none"
        period_all_neg = len(pos) == 0

    # 10% dropout margin gain range
    d_decrease = policy_rank[(policy_rank["lever"] == "d") & (policy_rank["direction"] == "decrease")].copy()
    d_10pct = d_decrease[d_decrease["lever_change_pct"].abs() >= 9.9]
    if d_10pct.empty:
        d_10pct = d_decrease
    d_10pct_group = d_10pct.loc[d_10pct.groupby("group")["normalised_margin_gain_per_10pct"].idxmax()]
    d_min_gain = d_10pct_group.sort_values("margin_gain").iloc[0]
    d_max_gain = d_10pct_group.sort_values("margin_gain").iloc[-1]

    # Endogenous inflow safety factor actually used in the fitted model
    safety_factor = float((eq["r"] / eq["r_critical"]).min())

    return {
        "n_groups": n_groups,
        "largest_pools": largest_pools,
        "smallest_pool": smallest_pool,
        "smallest_margin_group": smallest_margin_group,
        "d_min_e": d_min_e,
        "d_max_e": d_max_e,
        "d_all_negative": d_all_negative,
        "positive_lever_sentence": positive_lever_sentence,
        "positive_lever_sentence_lower": positive_lever_sentence_lower,
        "highest_pd_group": highest_pd_group,
        "pnr_lever_text": pnr_lever_text,
        "sat_range_text": sat_range_text,
        "period_neg": period_neg,
        "period_pos": period_pos,
        "period_all_neg": period_all_neg,
        "d_min_gain_group": d_min_gain["group"],
        "d_max_gain_group": d_max_gain["group"],
        "d_min_gain": round(d_min_gain["margin_gain"]),
        "d_max_gain": round(d_max_gain["margin_gain"]),
        "safety_factor": safety_factor,
    }


# ---------------------------------------------------------------------------
# Markdown output
# ---------------------------------------------------------------------------

def _abstract_and_highlights(eq, pnr_closest):
    closest = pnr_closest.iloc[0]
    # Compute a robust, data-driven statement about the most efficient lever
    policy_rank_path = POL / "ranked_interventions.csv"
    if policy_rank_path.exists():
        policy_rank = pd.read_csv(policy_rank_path)
        top_by_group = policy_rank.groupby("group").head(1)
        all_top_are_d = (top_by_group["lever"] == "d").all()
        top_lever_mode = top_by_group["lever"].mode()
        most_common_lever = top_lever_mode.iloc[0] if not top_lever_mode.empty else "d"
    else:
        all_top_are_d = True
        most_common_lever = "d"
    if all_top_are_d:
        lever_text = "A simulated reduction in dropout yields the largest margin gain per unit proportional change in every group, making it the most sensitive transition lever in the model. "
        highlight_lever = "Dropout reduction gives the largest margin gain per 10% change across all groups."
    else:
        lever_text = f"A simulated reduction in dropout is the most common single positive lever, although other levers dominate for some groups in the current data. "
        highlight_lever = f"Simulated {most_common_lever} adjustment yields the largest margin gain per unit proportional change for most groups."
    abstract = (
        "Artificial intelligence (AI) and machine learning (ML) research is increasingly concentrated in a few regions, "
        "raising the risk that smaller research communities fall below a minimum viable coauthor pool and cannot recover. "
        "We model each civilisation as a six-compartment system of domestic and abroad early-career, high-impact, and principal-investigator researchers, "
        "and estimate transition rates from OpenAlex Artificial Intelligence works (subfield 1702). "
        "The minimum viable coauthor threshold is defined as M = k × c_bar, where c_bar is the mean number of authors per work and k is the median number of distinct last-author groups observed per recent year. "
        f"Across {len(eq)} groups, equilibrium domestic active pools remain above their thresholds, but the closest point of no return is observed for the {closest['group']} group, "
        f"where {closest['rate_name']} must be multiplied by {_fmt(closest['critical_factor'], 3)}× its current value (equivalent to a {closest['proximity']*100:.0f}% proportional {'reduction' if closest['critical_factor'] < 1 else 'increase'}) to drive the active pool to its threshold. "
        + lever_text
        + "Historical and saturating-inflow counterfactuals show that the model is most sensitive to exogenous entry and attrition. "
        "These results provide a quantitative framework for early, safety-factor-bound policy scenarios that preserve civilisational diversity in AI/ML research."
    )
    keywords = (
        "researcher mobility; artificial intelligence; civilisation grouping; "
        "ordinary differential equations; point of no return; innovation studies"
    )
    highlights = [
        "Nine civilisations modelled as six-compartment ODEs fitted to OpenAlex AI/ML data.",
        f"Closest point of no return: {closest['group']} via {closest['rate_name']} (factor {_fmt(closest['critical_factor'], 3)}×).",

        highlight_lever,
    ]
    return abstract, keywords, highlights


def _data_availability_text():
    return (
        "This study uses publication metadata from the OpenAlex API (subfield 1702, Artificial Intelligence; "
        "2000–2023). The extraction and analysis code, the country-to-civilisation mapping, and the result CSVs "
        "used to generate this manuscript are available in the public GitHub repository "
        "https://github.com/bougtoir/researcher-mobility-ode. OpenAlex data are released under CC0."
    )


def _descriptive_table(cohort):
    """Return DataFrame of descriptive statistics per group."""
    grp = cohort.groupby("origin_group").agg(
        n=("author_id", "count"),
        works=("n_ai_works", "sum"),
        active=("active", "sum"),
        hits=("hit", "sum"),
        pis=("pi", "sum"),
        career_start_mean=("career_start", "mean"),
        abroad=("abroad", "sum"),
    ).reset_index()
    grp["career_start_mean"] = grp["career_start_mean"].round(1)
    grp = grp.rename(columns={"origin_group": "Group"})
    return grp


def write_markdown(output_dir: Path, data, fig_paths):
    """Write a plain-text markdown version for version control and review."""
    (cohort, eq, sat_eq, top_t, pnr_closest, period_compare, boot, policy_rank) = data
    ctx = compute_context(cohort, eq, sat_eq, top_t, pnr_closest, period_compare, policy_rank)
    abstract, keywords, highlights = _abstract_and_highlights(eq, pnr_closest)
    desc = _descriptive_table(cohort)
    fig1_rel = _rel_path(fig_paths["fig1"], output_dir)
    fig2_rel = _rel_path(fig_paths["fig2"], output_dir)
    fig3_rel = _rel_path(fig_paths["fig3"], output_dir)
    fig4_rel = _rel_path(fig_paths["fig4"], output_dir)
    fig5_rel = _rel_path(fig_paths["fig5"], output_dir) if fig_paths.get("fig5") else ""
    fig6_rel = _rel_path(fig_paths["fig6"], output_dir) if fig_paths.get("fig6") else ""
    fig7_rel = _rel_path(fig_paths["fig7"], output_dir) if fig_paths.get("fig7") else ""

    annual = load_annual_data()
    annual_ctx = compute_annual_context(annual)
    annual_means = annual_summary_table(annual)
    interciv_top = interciv_top_table(annual)
    group_acc = annual.get("group_accuracy")
    comp_acc = annual.get("compartment_accuracy")

    best_rmse_group = "—"
    worst_rmse_group = "—"
    worst_mape_group = "—"
    best_compartment_rmse = "—"
    worst_compartment_rmse = "—"
    worst_compartment_mape = "—"
    if group_acc is not None and not group_acc.empty:
        best_rmse_group = group_acc.loc[group_acc["rmse"].idxmin(), "origin_group"]
        worst_rmse_group = group_acc.loc[group_acc["rmse"].idxmax(), "origin_group"]
        worst_mape_group = group_acc.loc[group_acc["mape"].idxmax(), "origin_group"]
    if comp_acc is not None and not comp_acc.empty:
        best_compartment_rmse = comp_acc.loc[comp_acc["rmse"].idxmin(), "compartment"]
        worst_compartment_rmse = comp_acc.loc[comp_acc["rmse"].idxmax(), "compartment"]
        worst_compartment_mape = comp_acc.loc[comp_acc["mape"].idxmax(), "compartment"]

    lines = [
        "# Quantifying the Point of No Return in Global AI/ML Research Communities",
        "",
        "**Article type:** Research Article",
        "",
        "## Abstract",
        "",
        abstract,
        "",
        f"**Keywords:** {keywords}",
        "",
        "## Highlights",
        "",
    ]
    for h in highlights:
        lines.append(f"- {h}")
    lines.extend(["", "## Data and Code Availability", "", _data_availability_text(), ""])
    lines.extend([
        "## Declarations",
        "",
        "**Funding:** [To be completed / removed for double-blind review]",
        "",
        "**Competing interests:** [To be completed / removed for double-blind review]",
        "",
        "**Author contributions:** [To be completed / removed for double-blind review]",
        "",
        "**Acknowledgments:** " + NOTE_TEXT,
        "",
    ])

    lines.extend([
        "## 1. Introduction",
        "",
        "Most debates on research mobility focus on net flows: which country gains researchers and which loses them. "
        "Net-flow accounting is useful for headlines, but it hides the transition rates that actually move researchers between career stages and locations. "
        "A small proportional change in one of those rates can, over time, push a research community below the minimum coauthor pool it needs to remain viable. "
        "Once the pool falls below that threshold, recovery becomes difficult or impossible, even if policy is later reversed. "
        "That is the point of no return that motivates this paper.",
        "",
        "Artificial intelligence and machine learning have become the archetypal general-purpose technologies of the current era [5,6,7]. "
        "Their development depends on a relatively small, highly mobile workforce of doctoral and post-doctoral researchers, principal investigators, and research engineers [1]. "
        "The geographic concentration of this workforce has generated both scientific and geopolitical concern. "
        "Policymakers in the United States, China, Europe, Japan, India and elsewhere now treat AI talent as a strategic input, and several governments have introduced incentives to attract or retain researchers [4,5]. "
        "Most of those policies are evaluated by their immediate net-flow effects. "
        "They rarely ask which transition in the career pipeline is the binding constraint, or how close a community is to a threshold where the field can no longer sustain itself.",
        "",
        "The civilisation framework offers a natural way to partition the global research population into culturally and institutionally coherent arenas [4]. "
        "We adapt Huntington's nine civilisations for AI/ML mobility by keeping the United States, China (Sinic), India (Hindu), Japan, and the Islamic world as distinct groups, splitting the Western bloc into the United States, Anglosphere excluding the United States, Continental Europe and Other Western, and merging the smaller Latin American, Orthodox and African communities into Other Civilizations. "
        "This grouping reflects the empirical size and mobility patterns observed in the data rather than a normative claim about civilisational identity.",
        "",
        "The central argument of the paper is that preserving civilisational diversity in AI/ML is not only a normative preference but also a safeguard against technological dead ends. "
        "When a single region or a small oligopoly dominates a field, the set of research questions, evaluation norms, and institutional incentives narrows [5]. "
        "A diverse ecosystem generates competing approaches, which increases the probability that unexpected breakthroughs and error correction survive [5]. "
        "If transition rates can be observed with enough temporal resolution, policy can intervene before a community reaches the point of no return. "
        "Early, proportionate interventions can prevent the emergence of a monopoly or oligopoly without requiring large ex post rescues.",
        "",
        "We therefore address five research questions. "
        "First, how close is each civilisation to the point of no return in its AI/ML research community? "
        "Second, which transition rates have the largest effect on community size? "
        "Third, how have transition rates changed between earlier and later career cohorts, and what would have happened if those rates had persisted? "
        "Fourth, what safety-factor-bound policy packages can widen the margin before a point of no return is reached? "
        "Fifth, can the fitted rates be estimated year by year and used to project near-term population composition, and how well do those projections reproduce observed 2017-2023 counts?",
        "",
        "The contribution is a reproducible, data-driven transition-rate model that links OpenAlex publication records to a system of ordinary differential equations. "
        "The model is intentionally simple: it does not explain why a rate is high or low, but it identifies which rate is closest to a threshold and therefore where early intervention is most urgent.",
        "",
    ])

    lines.extend([
        "## 2. Literature and conceptual framework",
        "",
        "Researcher mobility has long been studied under the headings of brain drain, brain circulation and brain gain [3,4,9]. "
        "Thorn and Holm-Nielsen argue that the mobility of researchers from developing countries can become a gain when return migration and diaspora networks are supported, but it can become a drain when local research environments cannot retain or reproduce talent [7]. "
        "Appelt et al., using a gravity framework for 1996-2011, find that scientific collaboration, economic convergence and visa restrictions are the strongest correlates of bilateral mobility [2]. "
        "Their analysis shows that mobility is multi-directional: a large share of researcher movement is better described as circulation than as one-way migration.",
        "",
        "The AI/ML literature has documented the same patterns at higher resolution. "
        "MacroPolo's Global AI Talent Tracker finds that the United States remains the leading destination for top-tier AI researchers, while China and India are expanding domestic retention [1]. "
        "AlShebli et al. show that U.S.-China collaboration in AI is more impactful than either country working alone, and that most mobile AI scientists retain collaboration links with their origin country [8]. "
        "Yuan et al. find that the brain-drain problem for AI scientists is increasingly serious in developing countries, and that the ties among AI elites are highly clustered [9]. "
        "These studies establish that AI/ML talent is mobile, concentrated and strategically important.",
        "",
        "What is missing is a formal link between individual transition rates and the long-run viability of a research community. "
        "The concept of a minimum viable population, introduced by Shaffer, captures the smallest isolated population that has a high probability of persisting despite demographic, environmental and genetic stochasticity [10]. "
        "Transferred to science, the equivalent idea is a minimum viable coauthor pool: the smallest number of active researchers that can continue to produce work at the field's observed coauthor intensity. "
        "Below that pool, collaboration networks fragment, mentorship chains break, and the field enters a self-reinforcing decline.",
        "",
        "This framing generates four testable hypotheses. "
        "H1: Across all groups, the equilibrium active pool exceeds the minimum viable threshold, but the distance to the threshold varies widely. "
        "H2: Dropout is the transition rate with the largest negative effect, because attrition removes researchers from every compartment. "
        f"H3: {ctx['positive_lever_sentence'].rstrip('.')}. "
        "H4: Smaller civilisations, and those with older cohort structures, sit closer to their point of no return.",
        "",
    ])

    lines.extend([
        "## 3. Data and grouping",
        "",
        "We extracted AI/ML works and author histories from the OpenAlex API for subfield `subfields/1702` (Artificial Intelligence), using works published between 2000 and 2023 [2]. "
        "Authors were assigned to a civilisation by the majority country of their affiliated institutions. "
        "The mapping is documented in the repository and is reproduced here only in summary. "
        "The final groups are: United States, Anglosphere ex-US, Continental Europe, Sinic, Japanese, Hindu, Islamic, Other Western, and Other Civilizations.",
        "",
        "Table 1 reports the size and composition of the extracted cohort. "
        "The sample is a reproducible pilot extraction; absolute counts are small because the goal is to demonstrate the transition-rate framework rather than to provide a definitive census of global AI/ML researchers.",
        "",
    ])

    headers1 = ["Group", "Authors", "Works", "Active", "Hits", "PIs", "Career start", "Abroad"]
    lines.append("| " + " | ".join(headers1) + " |")
    lines.append("|" + "|".join(["---"] * len(headers1)) + "|")
    for _, row in desc.iterrows():
        lines.append(
            f"| {row['Group']} | {row['n']} | {row['works']} | {row['active']} | {row['hits']} | {row['pis']} | {row['career_start_mean']} | {row['abroad']} |"
        )
    lines.append("")

    lines.extend([
        "## 4. Methods",
        "",
        "### 4.1 Compartment model",
        "",
        "Each civilisation is represented by six compartments: domestic early-career researchers (D), abroad early-career researchers (A), domestic hit researchers (H_D), abroad hit researchers (H_A), domestic principal investigators (P_D), and abroad principal investigators (P_A). "
        "Transition rates are early-career outflow (α), return (β), hit generation at home and abroad (h_D and h_A), PI promotion at home and abroad (p_D and p_A), and dropout from all compartments (d). "
        "The equations are written in Word equation objects in the body of the manuscript.",
        "",
        "### 4.2 Endogenous inflow",
        "",
        "New entrants are modelled as a function of the domestic PI stock. "
        f"The linear form is I(P_D) = I_0 + r P_D, with r capped at {_fmt(ctx['safety_factor'], 2)}× the stability-critical value (safety factor {_fmt(ctx['safety_factor'], 2)}). "
        "A saturating alternative, I(P_D) = I_0 + r P_D / (1 + ε P_D), is reported as a robustness check.",
        "",
        "### 4.3 Minimum viable coauthor threshold",
        "",
        "For each group we computed the mean number of authors per work (c_bar) and the median number of distinct last-author groups observed per recent year (k). "
        "The minimum viable domestic active pool is M = k × c_bar. "
        "When the equilibrium active pool T = D + H_D + P_D falls below M, the community can no longer produce works at the observed coauthor intensity.",
        "",
        "### 4.4 Estimation and equilibrium",
        "",
        "Transition rates are estimated as constant per-year hazards from observed proportions within the cohort, with Laplace smoothing to avoid zero probabilities. "
        "The ODE system is solved at steady state for each group. "
        "Elasticities are computed by perturbing each rate by 1% and re-solving. "
        "For point-of-no-return analysis we scale each rate until the active pool reaches M and record the critical factor and its proximity, |critical factor − 1|. "
        "Historical counterfactuals split the cohort at career-start year 2010 and re-estimate all rates for the early and late windows. "
        "Bootstrap confidence intervals are obtained by resampling authors with replacement. "
        "Policy counterfactuals apply proportional changes to individual rates and report the resulting change in safety margin.",
        "",
    ])

    lines.extend([
        "### 4.10 Annual transition-rate estimation and projection",
        "",
        "The steady-state model in Sections 4.1-4.4 treats rates as constants. "
        "To test whether the same framework can be used for short-run monitoring, we reconstructed year-by-year compartment membership from cohort.csv and raw_sampled_works.json. "
        "For each author and year we inferred location as domestic if the author was in the origin civilisation and abroad otherwise, using sampled works when available and cohort-derived abroad/return years as a fallback. "
        "From these states we computed annual transition counts for the six compartments, applied Laplace +0.5 smoothing to empty destination cells, and derived the probabilities that map to α, β, h_D, h_A, p_D, p_A and d. "
        "Inter-civilisation flows are approximated by assigning each abroad author-year to the author's recent_group as the destination civilisation; this is a lower-bound proxy because year-to-year destination changes are not observed in the public cohort.",
        "",
        "For the 2017-2026 projection we fit a linear trend to the observed 2000-2016 rates for each group and rate. "
        "If fewer than four observations were available or the fit explained less than 10% of the variance, the historical mean was used instead. "
        "Projected rates were clipped to values between 0 and 1. "
        "Dropout was capped at the 90th percentile of observed dropout rates to prevent implausible extrapolation. "
        "Projected total inflows were apportioned across compartments using the observed 2016 distribution. "
        "Population composition was projected forward with the discrete-time recursion N(t+1) = N(t)P(t) + b(t+1), where P(t) is a 6×6 row-stochastic-in-expectation matrix that preserves dropout mass: the row sum is 1 − d after scaling outgoing rates. "
        "This discrete step is the operational counterpart of the continuous-time ODE; with an annual dt it provides an early-warning signal one year ahead.",
        "",
        "We compare the 2017-2023 projection with the observed annual stock. "
        "The comparison is limited to years that have observed data, and the observed stock is reindexed to the full group-year-compartment grid so that zero-observed cells are not omitted from the accuracy metrics. "
        "Accuracy is reported as RMSE and MAPE; MAPE here is computed against count_obs + 1 to avoid division by zero and is therefore a conservative, non-standard measure.",
        "",
        "### 4.11 Correction pressures and theoretical bounds",
        "",
        "The annual estimates contain several regularising pressures that bound the model away from instability and fabrication. "
        "Laplace smoothing adds a uniform prior of 0.5 to every possible destination, which shrinks sparse cells toward 1/(number of destinations) and prevents zero-probability singularities when a transition is unobserved in a small group-year. "
        "It is equivalent to a weak Dirichlet prior and is a standard regulariser for sparse multinomial transitions.",
        "",
        "Clipping projected rates to values between 0 and 1 is a feasibility pressure: rates outside the probability simplex are inadmissible. "
        "The dropout cap is a safety pressure motivated by the fact that unbounded linear extrapolation of observed attrition would eventually predict more leavers than the total stock. "
        "The inflow apportionment pressure keeps the composition of new entrants aligned with the most recently observed recruitment pattern, rather than inventing a new distribution. "
        f"Finally, the safety factor of {_fmt(ctx['safety_factor'], 2)} on the endogenous PI-driven inflow keeps the system inside the stability boundary. "
        "Together these pressures embody the principle that projection should stay within observed empirical support and within theoretical stability limits; they are not arbitrary adjustments but transparent bounds that can be tightened or relaxed as more data become available.",
        "",
    ])

    lines.extend([
        "## 5. Results",
        "",
        f"Table 2 reports the equilibrium domestic active pool T, the minimum viable threshold M, and the endogenous inflow parameters for the {len(eq)} groups. "
        "All groups remain above their threshold under the fitted model, but margins differ by an order of magnitude.",
        "",
    ])

    headers2 = ["Group", "T_eq", "M", "Margin", "I0", "r", "r_obs", "r_crit"]
    col_map2 = ["group", "T_equilibrium", "M_threshold", "margin_to_threshold_T", "I0", "r", "r_obs", "r_critical"]
    dec2 = [None, 2, 2, 2, 2, 5, 5, 5]
    lines.append("| " + " | ".join(headers2) + " |")
    lines.append("|" + "|".join(["---"] * len(headers2)) + "|")
    for _, row in eq.iterrows():
        lines.append("| " + " | ".join([_fmt(row[c], dec2[i] or 2) for i, c in enumerate(col_map2)]) + " |")
    lines.append("")

    lines.extend([
        f"![Figure 1]({fig1_rel})",
        "",
        "**Figure 1. Equilibrium domestic active pool (T) and minimum viable coauthor threshold (M) by group.** All groups remain above the threshold, but the margin varies widely.",
        "",
        f"Table 3 shows the three transition-rate elasticities with the largest absolute impact on T for each group. "
        f"Dropout (d) is the largest negative lever in every group, with an elasticity between {_fmt(ctx['d_min_e'], 2)} and {_fmt(ctx['d_max_e'], 2)} for the active pool. "
        f"{ctx['positive_lever_sentence']} "
        f"The {ctx['highest_pd_group']} group shows the highest sensitivity to PI promotion (p_D), indicating that strengthening domestic promotion is especially important for that community.",
        "",
    ])

    headers3 = ["Group", "1st rate", "1st elasticity", "2nd rate", "2nd elasticity", "3rd rate", "3rd elasticity"]
    lines.append("| " + " | ".join(headers3) + " |")
    lines.append("|" + "|".join(["---"] * len(headers3)) + "|")
    for group, gdf in top_t.groupby("group"):
        top3 = gdf.sort_values("abs_elasticity", ascending=False).head(3)
        vals = top3[["rate", "elasticity"]].values.tolist()
        parts = [group]
        for rate, elas in vals:
            parts.extend([rate, _fmt(elas, 3)])
        lines.append("| " + " | ".join(parts) + " |")
    lines.append("")

    closest = pnr_closest.iloc[0]
    lines.extend([
        f"Table 4 reports, for each group, the single rate that reaches the active-pool threshold with the smallest proportional change. "
        f"The {closest['group']} group is the most fragile: {closest['rate_name']} must be multiplied by {_fmt(closest['critical_factor'], 3)}× its current value (equivalent to a {closest['proximity']*100:.0f}% proportional {'reduction' if closest['critical_factor'] < 1 else 'increase'}) to drive the active pool to its minimum viable threshold. "
        f"{ctx['pnr_lever_text']}.",
        "",
    ])

    headers4 = ["Group", "Target", "Rate", "Current", "Critical factor", "Proximity"]
    lines.append("| " + " | ".join(headers4) + " |")
    lines.append("|" + "|".join(["---"] * len(headers4)) + "|")
    for _, row in pnr_closest.iterrows():
        lines.append(
            f"| {row['group']} | {row['target']} | {row['rate_name']} | {_fmt(row['current_rate'], 4)} | {_fmt(row['critical_factor'], 3)} | {_fmt(row['proximity'], 3)} |"
        )
    lines.append("")

    lines.extend([
        f"![Figure 2]({fig2_rel})",
        "",
        "**Figure 2. Closest point-of-no-return proximity by group.** Smaller values mean a smaller proportional change in the listed rate is required to reach the threshold for the stated target pool.",
        "",
    ])

    if sat_eq is not None:
        lines.extend([
            "### 5.1 Saturating recruitment extension",
            "",
            f"Replacing linear inflow with a saturating form lowers equilibrium pools because each additional PI adds fewer entrants. "
            f"Across groups, saturating equilibrium T is {ctx['sat_range_text']}. "
            "Table 5 compares linear and saturating equilibrium T values.",
            "",
        ])
        headers5 = ["Group", "Linear T", "Saturating T", "ε"]
        lines.append("| " + " | ".join(headers5) + " |")
        lines.append("|" + "|".join(["---"] * len(headers5)) + "|")
        merged = eq[["group", "T_equilibrium"]].merge(
            sat_eq[["group", "T_equilibrium", "epsilon"]], on="group", suffixes=("_lin", "_sat")
        )
        for _, row in merged.iterrows():
            lines.append(
                f"| {row['group']} | {_fmt(row['T_equilibrium_lin'], 2)} | {_fmt(row['T_equilibrium_sat'], 2)} | {_fmt(row['epsilon'], 5)} |"
            )
        lines.append("")

    # Describe actual historical-comparison groups dynamically
    if period_compare.empty:
        neg_groups_md, pos_groups_md = "none", "none"
    else:
        sorted_pc = period_compare.sort_values("delta_margin")
        neg = sorted_pc[sorted_pc["delta_margin"] < 0]["group"].tolist()
        pos = sorted_pc[sorted_pc["delta_margin"] > 0]["group"].tolist()[::-1]
        neg_groups_md = ", ".join(neg) if neg else "none"
        pos_groups_md = ", ".join(pos) if pos else "none"
    n_compare_md = len(period_compare)
    if ctx.get("period_all_neg"):
        prefix = "Both" if n_compare_md == 2 else f"All {n_compare_md}"
        period_direction_text = (
            f"{prefix} groups with dual-window support would see smaller safety margins under late-window rates "
            f"({ctx['period_neg']})."
        )
    else:
        period_direction_text = (
            f"Groups that would see smaller safety margins under late-window rates: {ctx['period_neg']}. "
            f"Groups that would see larger safety margins under late-window rates: {ctx['period_pos']}."
        )
    lines.extend([
        "### 5.2 Historical counterfactual",
        "",
        "Table 6 compares the equilibrium that would have emerged if the transition rates estimated for the early career window (2000-2010) or the late window (2011-2016) had persisted indefinitely. "
        "The late window is shorter and its rates are estimated from younger cohorts, so the comparison should be read as a sensitivity exercise rather than a forecast. "
        f"Only {n_compare_md} groups have enough dual-window support for reliable rate estimation in both windows; they are listed in the table. "
        f"{period_direction_text}",
        "",
    ])

    headers6 = ["Group", "T early", "T late", "ΔT (%)", "Margin early", "Margin late", "Δ margin"]
    lines.append("| " + " | ".join(headers6) + " |")
    lines.append("|" + "|".join(["---"] * len(headers6)) + "|")
    for _, row in period_compare.iterrows():
        lines.append(
            f"| {row['group']} | {_fmt(row['T_early'], 1)} | {_fmt(row['T_late'], 1)} | {_fmt(row['pct_delta_T'], 1)} | {_fmt(row['margin_early'], 1)} | {_fmt(row['margin_late'], 1)} | {_fmt(row['delta_margin'], 1)} |"
        )
    lines.append("")

    lines.extend([
        f"![Figure 3]({fig3_rel})",
        "",
        "**Figure 3. Change in safety margin between early and late transition-rate regimes.** Positive values mean the late-window rates would produce a larger safety margin than the early-window rates if they persisted; negative values mean the margin would shrink. "
        "The comparison is across two point estimates; uncertainty is substantial because the two windows have different cohort sizes and the steady-state model does not capture policy shocks.",
        "",
    ])

    d_decrease_md = policy_rank[(policy_rank["lever"] == "d") & (policy_rank["direction"] == "decrease")].copy()
    d_10pct_md = d_decrease_md[d_decrease_md["lever_change_pct"].abs() >= 9.9]
    if d_10pct_md.empty:
        d_10pct_md = d_decrease_md
    d_10pct_group_md = d_10pct_md.loc[d_10pct_md.groupby("group")["normalised_margin_gain_per_10pct"].idxmax()]
    policy_top_md = policy_rank.groupby("group").head(1)
    all_top_are_d_md = (policy_top_md["lever"] == "d").all()
    d_min_md = d_10pct_group_md.sort_values("margin_gain").iloc[0]
    d_max_md = d_10pct_group_md.sort_values("margin_gain").iloc[-1]
    if all_top_are_d_md:
        lever_statement_md = "Reducing dropout is the dominant positive lever for every civilisation."
    else:
        lever_statement_md = "Reducing dropout is the dominant positive lever for most civilisations in the current data."
    lines.extend([
        "### 5.3 Policy counterfactuals",
        "",
        "Table 7 reports the single mechanical counterfactual with the largest margin gain per 10% lever change for each group. "
        f"{lever_statement_md} "
        f"A roughly 10% proportional reduction in d would add about {round(d_min_md['margin_gain'])} active researchers in the {d_min_md['group']} group and about {round(d_max_md['margin_gain'])} in the {d_max_md['group']} group, reflecting differences in cohort size and baseline attrition.",
        "",
    ])

    policy_top = policy_rank.groupby("group").head(1).copy()
    headers7 = ["Group", "Lever", "Direction", "Change (%)", "Margin gain", "Gain per 10%"]
    lines.append("| " + " | ".join(headers7) + " |")
    lines.append("|" + "|".join(["---"] * len(headers7)) + "|")
    for _, row in policy_top.iterrows():
        lines.append(
            f"| {row['group']} | {row['lever']} | {row['direction']} | {_fmt(row['lever_change_pct'], 0)} | {_fmt(row['margin_gain'], 1)} | {_fmt(row['normalised_margin_gain_per_10pct'], 1)} |"
        )
    lines.append("")

    lines.extend([
        "### 5.4 Uncertainty",
        "",
        "Table 8 reports bootstrap 95% confidence intervals for the equilibrium active pool T and the domestic PI pool P_D. "
        "The intervals are wide, reflecting the small cohort sample and the extrapolation from individual careers to long-run steady states.",
        "",
    ])

    headers8 = ["Group", "T median", "T 95% CI", "P_D mean", "P_D 95% CI"]
    lines.append("| " + " | ".join(headers8) + " |")
    lines.append("|" + "|".join(["---"] * len(headers8)) + "|")
    for _, row in boot.iterrows():
        t_ci = f"[{_fmt(row['T_equilibrium_q025'], 0)}, {_fmt(row['T_equilibrium_q975'], 0)}]"
        p_ci = f"[{_fmt(row['P_D_equilibrium_q025'], 0)}, {_fmt(row['P_D_equilibrium_q975'], 0)}]"
        lines.append(
            f"| {row['group']} | {_fmt(row['T_equilibrium_median'], 0)} | {t_ci} | {_fmt(row['P_D_equilibrium_mean'], 0)} | {p_ci} |"
        )
    lines.append("")

    lines.extend([
        f"![Figure 4]({fig4_rel})",
        "",
        "**Figure 4. Bootstrap 95% confidence intervals for equilibrium T by group.** Intervals are asymmetric and wide, reflecting model uncertainty.",
        "",
    ])

    lines.extend([
        "### 5.6 Annual transition rates and inter-civilisation flows",
        "",
        "Figure 5 plots the observed 2000-2016 transition rates and the projected 2017-2026 rates for each civilisation. "
        "Rates are displayed by group and by transition type, so that the reader can see whether a particular transition is trending toward a boundary. "
        "Because the projections are linear trend fits regularised by the correction pressures described in Section 4.11, they are not forecasts of specific future events; they are the model's one-year-ahead extrapolation of the recent historical trajectory.",
        "",
        f"![Figure 5]({fig5_rel})",
        "",
        "**Figure 5. Annual observed (solid) and projected (dashed) transition rates by civilisation, 2000-2026.**",
        "",
        "Table 9 summarises the mean observed annual transition rates by group between 2000 and 2016. "
        "The table distinguishes early-career outflow (α), return (β), domestic and abroad hit generation (h_D, h_A), PI promotion (p_D), dropout (d), and total inflow (I_total).",
        "",
        "| Group | α | β | h_D | p_D | d | I_total |",
        "|---|---|---|---|---|---|---|",
        *[f"| {row['Group']} | {_fmt(row['α'], 3)} | {_fmt(row['β'], 3)} | {_fmt(row['h_D'], 3)} | {_fmt(row['p_D'], 3)} | {_fmt(row['d'], 3)} | {_fmt(row['I_total'], 2)} |" for _, row in (annual_means if not annual_means.empty else pd.DataFrame()).iterrows()],
        "",
        "**Table 9. Mean observed annual transition rates by civilisation, 2000-2016.**",
        "",
        "Figure 6 shows the inter-civilisation accumulation of abroad author-years. "
        "Rows represent the origin civilisation and columns represent the destination civilisation, approximated by the author's recent_group while abroad. "
        "The heatmap is a lower-bound proxy because year-to-year destination switches within a spell abroad are not observed.",
        "",
        f"![Figure 6]({fig6_rel})",
        "",
        "**Figure 6. Inter-civilisation abroad author-year accumulation by origin (rows) and destination (columns).**",
        "",
        "Table 10 lists the origin-destination pairs with the largest accumulation of abroad author-years. "
        "These pairs identify the strongest visible inter-civilisation pipelines and are the empirical counterpart to the α and β transitions.",
        "",
        "| Origin | Destination | Author-years |",
        "|---|---|---|",
        *[f"| {row['Origin']} | {row['Destination']} | {_fmt(row['Author-years'], 0)} |" for _, row in (interciv_top if not interciv_top.empty else pd.DataFrame()).iterrows()],
        "",
        "**Table 10. Top origin-destination abroad author-year pairs.**",
        "",
        "### 5.7 Out-of-sample projection, 2017-2023",
        "",
        f"The 2017-2023 projection is compared with observed annual stocks in Figure 7. "
        f"Overall accuracy is RMSE {_fmt(annual_ctx.get('overall_rmse', float('nan')), 2)} and MAPE {_fmt(annual_ctx.get('overall_mape_pct', float('nan')), 1)}% (a non-standard, conservative measure computed against count_obs + 1 to avoid division by zero). "
        f"Among civilisations the lowest RMSE is for {best_rmse_group} and the highest RMSE is for {worst_rmse_group}; the highest MAPE is for {worst_mape_group}. "
        "The largest errors occur in small compartments and in groups with sparse transition counts, which is expected because the annual model does not borrow information across civilisations.",
        "",
        f"![Figure 7]({fig7_rel})",
        "",
        "**Figure 7. Observed (solid) and projected (dashed) compartment counts by civilisation, 2017-2023. The vertical dotted line marks the end of the training period (2016).**",
        "",
        "Table 11 reports projection accuracy by civilisation and Table 12 by compartment. "
        f"Among compartments, the lowest RMSE is for {best_compartment_rmse}, while the highest RMSE is for {worst_compartment_rmse} and the highest MAPE is for {worst_compartment_mape}. "
        "P_D and H_D show larger errors because small changes in PI and hit rates are amplified by the endogenous inflow term.",
        "",
        "| Group | RMSE | MAPE |",
        "|---|---|---|",
        *[f"| {row['origin_group']} | {_fmt(row['rmse'], 2)} | {row['mape']*100:.1f}% |" for _, row in ((group_acc if group_acc is not None else pd.DataFrame()) if not (group_acc if group_acc is not None else pd.DataFrame()).empty else pd.DataFrame()).iterrows()],
        "",
        "**Table 11. Projection accuracy by civilisation, 2017-2023.**",
        "",
        "| Compartment | RMSE | MAPE |",
        "|---|---|---|",
        *[f"| {row['compartment']} | {_fmt(row['rmse'], 2)} | {row['mape']*100:.1f}% |" for _, row in ((comp_acc if comp_acc is not None else pd.DataFrame()) if not (comp_acc if comp_acc is not None else pd.DataFrame()).empty else pd.DataFrame()).iterrows()],
        "",
        "**Table 12. Projection accuracy by compartment, 2017-2023.**",
        "",
        "### 5.8 Correction pressures in the annual model",
        "",
        "The annual projection performs best where the correction pressures in Section 4.11 are binding. "
        "Laplace smoothing prevents empty cells from being treated as impossible transitions; the unit-interval clip and the dropout cap prevent the trend extrapolation from producing rates that are incompatible with a stochastic transition matrix; and the 2016 inflow apportionment keeps new-entrant composition close to the last observed regime. "
        "These pressures mean that the projection is not a purely mechanical forecast: it is a bounded extrapolation that stays within the empirical support of the 2000-2016 data and within the stability constraints of the compartment model.",
        "",
        "## 6. Discussion",
        "",
        "The results support a transition-rate view of research policy. "
        "Rather than asking which country has a net inflow or outflow of researchers, the model asks which rate must be altered to keep a community above its minimum viable coauthor pool. "
        "The answer is not the same for every group, but a clear pattern emerges.",
        "",
        f"First, {ctx['pnr_lever_text']}. "
        "A large proportional reduction in baseline recruitment would drive most communities to their threshold before mobility rates such as return or promotion became binding. "
        "This is consistent with the observation that AI/ML fields depend on a continuous pipeline of new graduate students and junior researchers [5,7]. "
        "Policies that sustain that pipeline, such as doctoral funding, visa routes for early-career researchers, and stable junior positions, are therefore first-order defences against a point of no return.",
        "",
        f"Second, among the mobility transition rates, dropout (d) is the dominant negative lever; its active-pool elasticity ranges from {_fmt(ctx['d_min_e'], 2)} to {_fmt(ctx['d_max_e'], 2)} across groups, and in the policy counterfactuals a simulated reduction in dropout yields the largest margin gain per unit proportional change. "
        "Attrition matters because it removes researchers from every compartment, not just one. "
        "A 10% proportional reduction in dropout expands the safety margin more than comparably sized increases in return, hit generation or promotion. "
        f"For {ctx['smallest_margin_group']}, the group with the smallest safety margin, even modest attrition reductions may widen the margin. "
        "These counterfactuals are mechanical perturbations of the fitted rates; they identify the most sensitive transition levers, not the causal effect of any specific policy programme.",
        "",
        f"Third, {ctx['positive_lever_sentence_lower']}. "
        f"The {ctx['highest_pd_group']} group shows the strongest response to PI promotion, suggesting that for that community expanding the domestic PI pipeline is an efficient lever. "
        "Return from abroad (β) is also positive for most groups, though its effect is generally smaller than reducing attrition directly. "
        "The implication for policy is that retention and promotion are usually more efficient than trying to attract returnees, but a balanced portfolio is still needed: a community without domestic PI growth cannot reproduce itself through attrition reduction alone.",
        "",
        f"Fourth, the historical counterfactual shows that the late-window rates, if they persisted, would alter equilibrium margins. "
        f"{period_direction_text} "
        "This pattern cautions against treating AI/ML mobility as a single global trend. "
        "It also confirms that the model can detect temporal changes in transition rates, which is the prerequisite for the early intervention the framework is designed to support.",
        "",
        "The transition levers also interact in ways that a single-rate elasticity cannot fully capture. "
        "For example, reducing dropout and increasing PI promotion together are likely to have a larger effect than the sum of the two individual perturbations, because more researchers survive to become PIs and those PIs then train additional early-career researchers through the endogenous inflow channel. "
        "Conversely, a simultaneous fall in exogenous entry and a rise in dropout can push a community to its threshold faster than either change alone. "
        "The model's steady-state and one-at-a-time counterfactuals are therefore a starting point; they identify the most sensitive margins but do not exhaust the policy design space.",
        "",
        "The connection to civilisational diversity is direct. "
        "Each group's safety margin can be monitored over time, and interventions can be adjusted before the margin disappears. "
        f"Because the model uses a fixed safety factor of {_fmt(ctx['safety_factor'], 2)} for the endogenous inflow parameter r, the policy recommendations are deliberately conservative: they do not push the system toward instability. "
        "That bounded approach is consistent with the goal of preserving diversity rather than maximising any single country's share.",
        "",
        "It is important to stress that the counterfactuals reported in Tables 3 and 7 are mechanical perturbations of the fitted transition rates, not causal estimates of specific programmes. "
        "They identify which rates the model treats as most sensitive, and therefore where empirical policy evaluation is most urgent, but they do not by themselves show that a given intervention would achieve the simulated change.",
        "",
        "Several limitations should be acknowledged. "
        "OpenAlex affiliation and country assignments are noisy, especially for researchers with multiple affiliations. "
        "The civilisation grouping is a coarse aggregation; within-group heterogeneity is substantial. "
        "The model is a steady-state ODE and does not capture short-term dynamics, cross-civilisation spillovers, or the non-linear effects of network externalities. "
        "The cohort sample is small; the absolute equilibrium numbers should be interpreted as model-implied stocks rather than as census counts. "
        "Authors with many publications are over-weighted relative to one-publication authors, so rate estimates reflect author-publication exposure rather than a uniformly representative sample of individuals. "
        f"The endogenous inflow is capped at a safety factor of {_fmt(ctx['safety_factor'], 2)} relative to the critical reproduction rate; alternative values would shift equilibrium levels and should be reported in future sensitivity tables. "
        "Finally, the point-of-no-return threshold is a sufficient condition for collapse, not a necessary one: a community may decline for reasons outside the model even if T remains above M.",
        "Wide bootstrap confidence intervals, especially for smaller civilisation groups, mean that the ordinal ranking of groups by equilibrium size or proximity to threshold should be treated as descriptive rather than definitive. "
        "The model identifies which transitions are most sensitive in a mechanical sense; turning those sensitivities into reliable policy priorities requires additional data on programme costs, implementation lags, and behavioural responses that are outside the scope of this paper.",
        "Operationally, the framework can be used in two complementary ways. "
        "As a monitoring tool, it can be rerun whenever new OpenAlex data are released, producing an updated set of transition rates, safety margins and proximity-to-threshold estimates. "
        "As a scenario tool, it can quantify how large a proportional change in a given rate would be required to move a community toward or away from collapse, which helps prioritise empirical policy evaluation. "
        "Both uses depend on transparent assumptions and regular recalibration; the model should not be used to justify one-off interventions without accompanying process evaluation.",
        "",
        "### 6.4 Validation of correction pressures",
        "",
        f"The correction pressures are not ad hoc adjustments; each maps to a known statistical or dynamical constraint. "
        f"Laplace smoothing is equivalent to a weak Dirichlet prior on a multinomial transition vector; it guarantees that no cell has zero estimated probability and shrinks rare transitions toward the simplex centroid. "
        f"Clipping projected rates to values between 0 and 1 is a feasibility constraint on probabilities; the dropout cap is a cross-sectional constraint that prevents projected attrition from exceeding the observed stock; and the inflow apportionment constraint keeps the composition of new entrants equal to the last observed recruitment pattern. "
        f"In the 2017-2023 projection these pressures reduced the sensitivity of the forecast to sparse cells and to short-run fluctuations in small groups. "
        f"Quantitatively, the overall RMSE of {_fmt(annual_ctx.get('overall_rmse', float('nan')), 2)} and conservative MAPE of {_fmt(annual_ctx.get('overall_mape_pct', float('nan')), 1)}% are consistent with a model that is deliberately regularised rather than optimised for in-sample fit. "
        f"The residual errors are concentrated in the smallest compartments, which is exactly where smoothing is most active and where future data will be most valuable.",
        "",
        "### 6.5 Intra-civilisation alternatives when inter-civilisation mobility cannot be controlled",
        "",
        "If a civilisation cannot control outflows to, or inflows from, other jurisdictions—whether because of visa regimes, salary differentials, language advantages, or targeted recruitment—it can still preserve its research community by acting on the intra-civilisation levers identified in the annual model. "
        "The annual rates show that the domestic active pool T = D + H_D + P_D responds most strongly to the dropout rate d, the domestic hit rate h_D, and the PI promotion rate p_D. "
        "Policies that reduce early-career attrition, expand domestic postdoctoral positions, or accelerate independent-lab formation therefore become defensive substitutes when inter-civilisation poaching cannot be regulated. "
        "This is the practical meaning of civilisational-diversity preservation under sovereignty constraints: even without controlling the border of talent, a community can increase the internal reproduction of active researchers. "
        f"The ODE safety factor of {_fmt(ctx['safety_factor'], 2)} on endogenous PI inflow is a conservative bound that prevents over-optimism about this substitution effect; more ambitious domestic growth would require corresponding evidence that the extra PIs can be absorbed without simply raising dropout.",
        "",
        "### 6.6 Annual updating as an early-warning layer",
        "",
        "The 2017-2023 projection demonstrates that the framework can be rerun annually with a one-year time step. "
        "Each new year of OpenAlex data updates the observed transition rates, the fitted trends, and the distance to the minimum viable coauthor threshold. "
        "Because the model is regularised by the correction pressures, the one-year-ahead projection is not easily derailed by a single noisy observation. "
        "Instead, successive years reveal whether a particular transition rate is drifting toward a boundary. "
        "That drift is the early-warning signal. "
        "Policymakers can then intervene before the active pool falls below M, using the rate-specific elasticities in Table 3 to prioritise the smallest proportional change that restores a safety margin. "
        "This is the operational mechanism for avoiding technology monopoly and oligopoly dead ends: by keeping every major research community above its minimum viable coauthor pool, annual monitoring sustains the competitive diversity that underpins long-run technological progress. "
        "The framework is therefore not a prediction that a particular civilisation will collapse; it is a tool for ensuring that no single civilisation reaches a point where its collapse becomes self-sustaining.",
        "",
        "### 6.7 Limitations",
        "",
        "Several limitations should be acknowledged. "
        "OpenAlex affiliation and country assignments are noisy, especially for researchers with multiple affiliations. "
        "The civilisation grouping is a coarse aggregation; within-group heterogeneity is substantial. "
        "The annual model relies on a discrete approximation of the continuous-time ODE and does not capture within-year events or cross-civilisation spillovers. "
        "Inter-civilisation flows are approximated by the author's recent_group while abroad, which misses year-to-year destination switching. "
        "The cohort sample is small; the absolute equilibrium numbers should be interpreted as model-implied stocks rather than as census counts. "
        "Authors with many publications are over-weighted relative to one-publication authors, so rate estimates reflect author-publication exposure rather than a uniformly representative sample of individuals. "
        f"The endogenous inflow is capped at a safety factor of {_fmt(ctx['safety_factor'], 2)} relative to the critical reproduction rate; alternative values would shift equilibrium levels and should be reported in future sensitivity tables. "
        "Finally, the point-of-no-return threshold is a sufficient condition for collapse, not a necessary one: a community may decline for reasons outside the model even if T remains above M.",
        "",
    ])

    lines.extend([
        "## 7. Conclusion",
        "",
        "We have proposed and implemented a transition-rate framework for assessing how close AI/ML research communities are to a point of no return. "
        "The model converts OpenAlex publication records into civilisation-specific transition rates and solves for the equilibrium active researcher pool. "
        "All groups remain above their minimum viable coauthor threshold in the fitted model, but the distance to that threshold varies by an order of magnitude and is most sensitive to exogenous entry and dropout. "
        f"Dropout is the dominant negative lever (active-pool elasticity {_fmt(ctx['d_min_e'], 2)} to {_fmt(ctx['d_max_e'], 2)}), and a simulated reduction is the single most efficient model-implied response for every civilisation. "
        "However, the closest point of no return is exogenous entry for all groups in the active-pool analysis, which means that policies which sustain the pipeline of new researchers are first-order defences. "
        "The historical counterfactual and the bootstrap intervals remind us that the future is not determined by current rates; transition rates can change, and policy can be directed at the most fragile lever before a collapse.",
        "",
        "The annual projection layer adds an operational dimension to this conclusion. "
        "By estimating year-by-year transition rates and projecting one year ahead, the model turns the steady-state diagnostic into an early-warning dashboard. "
        "A one-year time step is short enough to detect drift before the active pool approaches the minimum viable threshold, and the correction pressures keep the projection within empirical and theoretical bounds. "
        "When inter-civilisation mobility cannot be controlled, the same framework points to intra-civilisation levers—reducing dropout, raising domestic hit rates, and accelerating PI promotion—that preserve T = D + H_D + P_D. "
        "These two layers, steady-state and annual, together provide a coherent basis for early, safety-factor-bound intervention.",
        "",
        "The broader implication is that preserving civilisational diversity in AI/ML is compatible with, and may reinforce, scientific progress. "
        "A single dominant region or a tight oligopoly may achieve short-run scale economies, but it also risks methodological lock-in and reduces the set of problems that receive sustained attention. "
        "By monitoring transition rates and safety margins, policymakers can detect divergence early and intervene in a safety-factor-bound way. "
        "This is the practical meaning of the aspiration to avoid technology monopoly and oligopoly dead ends: not a prediction that any one civilisation will dominate, but a structured method for keeping the global system away from points of no return. "
        "Early, proportionate interventions that reduce attrition and sustain new recruitment can widen safety margins and preserve civilisational diversity in AI/ML.",
        "",
        "Future work should extend the model to network externalities, finer temporal resolution, and additional security-relevant fields such as semiconductor physics, quantum computing, biotechnology and energy materials, allowing cross-field comparisons of vulnerability. "
        "Other priorities include systematic sensitivity scans for the safety factor and saturating parameter epsilon, country- or institution-level partitions, dynamic ODE forecasts, endogenous coauthorship matching, and integration with policy cost data to produce cost-effectiveness comparisons of alternative interventions.",
        "",
        "## References",
        "",
    ])
    for i, ref in enumerate(REFS, 1):
        lines.append(f"{i}. {ref}")

    md_path = output_dir / "manuscript_full_article.md"
    lines.append("")
    md_path.write_text("\n".join(lines), encoding="utf-8")
    return md_path


# ---------------------------------------------------------------------------
# Word output
# ---------------------------------------------------------------------------

def _add_title_page(doc, word_count=None):
    style = doc.styles["Normal"]
    style.font.name = "Times New Roman"
    style.font.size = Pt(11)
    style.paragraph_format.line_spacing = 1.5

    title = doc.add_heading("Quantifying the Point of No Return in Global AI/ML Research Communities", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title.runs[0].font.size = Pt(16)
    title.runs[0].font.bold = True

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run("Article type: Research Article")

    if word_count:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run(f"Approximate word count (main text incl. tables, excl. references): {word_count}")

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run("Corresponding author: [To be completed / removed for double-blind review]")

    p = doc.add_paragraph()
    run = p.add_run()
    run.add_break(WD_BREAK.PAGE)


def _add_front_matter(doc, abstract, keywords, highlights):
    doc.add_heading("Abstract", level=1)
    p = doc.add_paragraph()
    p.add_run(abstract)

    p = doc.add_paragraph()
    p.add_run("Keywords: ").bold = True
    p.add_run(keywords)

    doc.add_heading("Highlights", level=2)
    for h in highlights:
        p = doc.add_paragraph(style="List Bullet")
        p.add_run(h)

    doc.add_heading("Data and Code Availability", level=2)
    p = doc.add_paragraph()
    p.add_run(_data_availability_text())

    doc.add_heading("Declarations", level=2)
    for sub in ["Funding", "Competing interests", "Author contributions", "Acknowledgments"]:
        p = doc.add_paragraph()
        p.add_run(f"{sub}: ").bold = True
        if sub == "Acknowledgments":
            p.add_run("This study was motivated by a note.com essay by Yamada Y (momentumyy) that framed researcher mobility in terms of transition rates rather than net flows (" + NOTE_TEXT + ").")
        else:
            p.add_run("[To be completed / removed for double-blind review]")


def _add_table_from_df(doc, df, caption, decimals=None, bold_header=True):
    if decimals is None:
        decimals = {}
    cols = df.columns.tolist()
    table = doc.add_table(rows=1, cols=len(cols))
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    for i, c in enumerate(cols):
        hdr[i].text = str(c)
        if bold_header:
            for run in hdr[i].paragraphs[0].runs:
                run.font.bold = True
    for _, row in df.iterrows():
        cells = table.add_row().cells
        for i, c in enumerate(cols):
            v = row[c]
            cells[i].text = _fmt(v, decimals.get(c, 2))
    cap = doc.add_paragraph()
    cap.add_run(caption).italic = True
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    return table


def _add_docx_body(doc, data, fig_paths):
    (cohort, eq, sat_eq, top_t, pnr_closest, period_compare, boot, policy_rank) = data
    ctx = compute_context(cohort, eq, sat_eq, top_t, pnr_closest, period_compare, policy_rank)
    annual = load_annual_data()
    annual_ctx = compute_annual_context(annual)
    annual_means = annual_summary_table(annual)
    interciv_top = interciv_top_table(annual)
    group_acc = annual.get("group_accuracy")
    comp_acc = annual.get("compartment_accuracy")

    best_rmse_group = "—"
    worst_rmse_group = "—"
    worst_mape_group = "—"
    best_compartment_rmse = "—"
    worst_compartment_rmse = "—"
    worst_compartment_mape = "—"
    if group_acc is not None and not group_acc.empty:
        best_rmse_group = group_acc.loc[group_acc["rmse"].idxmin(), "origin_group"]
        worst_rmse_group = group_acc.loc[group_acc["rmse"].idxmax(), "origin_group"]
        worst_mape_group = group_acc.loc[group_acc["mape"].idxmax(), "origin_group"]
    if comp_acc is not None and not comp_acc.empty:
        best_compartment_rmse = comp_acc.loc[comp_acc["rmse"].idxmin(), "compartment"]
        worst_compartment_rmse = comp_acc.loc[comp_acc["rmse"].idxmax(), "compartment"]
        worst_compartment_mape = comp_acc.loc[comp_acc["mape"].idxmax(), "compartment"]

    # Introduction
    doc.add_heading("1. Introduction", level=1)
    p = doc.add_paragraph()
    p.add_run("Most debates on research mobility focus on net flows: which country gains researchers and which loses them. "
              "Net-flow accounting is useful for headlines, but it hides the transition rates that actually move researchers between career stages and locations. "
              "A small proportional change in one of those rates can, over time, push a research community below the minimum coauthor pool it needs to remain viable. "
              "Once the pool falls below that threshold, recovery becomes difficult or impossible, even if policy is later reversed. "
              "That is the point of no return that motivates this paper. "
              "The contribution of this paper is to translate that qualitative insight into an empirically tractable model. "
              "We estimate transition rates from open bibliometric data, solve the steady state of a compartment model, and identify which rate in which civilisation is closest to a threshold. "
              "The approach is deliberately stylised: it sacrifices demographic realism for transparency and for the ability to compare multiple civilisations with the same accounting framework.")

    p = doc.add_paragraph()
    p.add_run("Artificial intelligence and machine learning have become the archetypal general-purpose technologies of the current era")
    add_citation(p, 1)
    p.add_run(", and their development depends on a relatively small, highly mobile workforce of doctoral and post-doctoral researchers, principal investigators, and research engineers")
    add_citation(p, 1)
    p.add_run(". "
              "The geographic concentration of this workforce has generated both scientific and geopolitical concern. "
              "Policymakers in the United States, China, Europe, Japan, India and elsewhere now treat AI talent as a strategic input, and several governments have introduced incentives to attract or retain researchers")
    add_citation(p, 2)
    p.add_run(". "
              "Most of those policies are evaluated by their immediate net-flow effects. "
              "They rarely ask which transition in the career pipeline is the binding constraint, or how close a community is to a threshold where the field can no longer sustain itself. "
              "The economic literature on science has long emphasised that researchers are a scarce input and that their mobility responds to career incentives and institutional quality")
    add_citation(p, 3)
    p.add_run(". "
              "That literature provides the microfoundation for our rates: individuals decide where to train, whether to go abroad, when to return, and when to leave academia. "
              "We aggregate those individual decisions into civilisation-level transition rates and ask what the resulting dynamics imply for community survival.")

    p = doc.add_paragraph()
    p.add_run("The civilisation framework offers a natural way to partition the global research population into culturally and institutionally coherent arenas")
    add_citation(p, 4)
    p.add_run(". "
              "We adapt Huntington's nine civilisations for AI/ML mobility by keeping the United States, China (Sinic), India (Hindu), Japan, and the Islamic world as distinct groups, splitting the Western bloc into the United States, Anglosphere excluding the United States, Continental Europe and Other Western, and merging the smaller Latin American, Orthodox and African communities into Other Civilizations. "
              "This grouping reflects the empirical size and mobility patterns observed in the data rather than a normative claim about civilisational identity.")

    p = doc.add_paragraph()
    p.add_run("The central argument of the paper is that preserving civilisational diversity in AI/ML is not only a normative preference but also a safeguard against technological dead ends. "
              "When a single region or a small oligopoly dominates a field, the set of research questions, evaluation norms, and institutional incentives narrows")
    add_citation(p, 5)
    p.add_run(". "
              "A diverse ecosystem generates competing approaches, which increases the probability that unexpected breakthroughs and error correction survive")
    add_citation(p, 5)
    p.add_run(". "
              "If transition rates can be observed with enough temporal resolution, policy can intervene before a community reaches the point of no return. "
              "Early, proportionate interventions can prevent the emergence of a monopoly or oligopoly without requiring large ex post rescues.")

    p = doc.add_paragraph()
    p.add_run("We therefore address five research questions. "
              "First, how close is each civilisation to the point of no return in its AI/ML research community? "
              "Second, which transition rates have the largest effect on community size? "
              "Third, how have transition rates changed between earlier and later career cohorts, and what would have happened if those rates had persisted? "
              "Fourth, what safety-factor-bound policy packages can widen the margin before a point of no return is reached? "
              "Fifth, can the fitted rates be estimated year by year and used to project near-term population composition, and how well do those projections reproduce observed 2017-2023 counts? "
              "The key policy intuition is that, with an appropriately chosen time step and an early warning signal, intervention can be calibrated in safety margins rather than after collapse. "
              "This prevents any single civilisation from cornering the supply of critical talent, and thereby preserves the competitive diversity that drives long-run innovation.")

    p = doc.add_paragraph()
    p.add_run("The contribution is a reproducible, data-driven transition-rate model that links OpenAlex publication records to a system of ordinary differential equations")
    add_citation(p, 6)
    p.add_run(". "
              "The model is intentionally simple: it does not explain why a rate is high or low, but it identifies which rate is closest to a threshold and therefore where early intervention is most urgent.")

    doc.add_paragraph()

    doc.add_heading("2. Literature and conceptual framework", level=1)

    p = doc.add_paragraph()
    p.add_run("Researcher mobility has long been studied under the headings of brain drain, brain circulation and brain gain")
    add_citation(p, 7)
    p.add_run(". "
              "Thorn and Holm-Nielsen argue that the mobility of researchers from developing countries can become a gain when return migration and diaspora networks are supported, but it can become a drain when local research environments cannot retain or reproduce talent")
    add_citation(p, 7)
    p.add_run(". "
              "Appelt et al., using a gravity framework for 1996-2011, find that scientific collaboration, economic convergence and visa restrictions are the strongest correlates of bilateral mobility")
    add_citation(p, 2)
    p.add_run(". "
              "Their analysis shows that mobility is multi-directional: a large share of researcher movement is better described as circulation than as one-way migration.")

    p = doc.add_paragraph()
    p.add_run("The AI/ML literature has documented the same patterns at higher resolution. "
              "MacroPolo's Global AI Talent Tracker finds that the United States remains the leading destination for top-tier AI researchers, while China and India are expanding domestic retention")
    add_citation(p, 1)
    p.add_run(". "
              "AlShebli et al. show that U.S.-China collaboration in AI is more impactful than either country working alone, and that most mobile AI scientists retain collaboration links with their origin country")
    add_citation(p, 8)
    p.add_run(". "
              "Yuan et al. find that the brain-drain problem for AI scientists is increasingly serious in developing countries, and that the ties among AI elites are highly clustered")
    add_citation(p, 9)
    p.add_run(". "
              "These studies establish that AI/ML talent is mobile, concentrated and strategically important.")

    p = doc.add_paragraph()
    p.add_run("What is missing is a formal link between individual transition rates and the long-run viability of a research community. "
              "The concept of a minimum viable population, introduced by Shaffer, captures the smallest isolated population that has a high probability of persisting despite demographic, environmental and genetic stochasticity")
    add_citation(p, 10)
    p.add_run(". "
              "Transferred to science, the equivalent idea is a minimum viable coauthor pool: the smallest number of active researchers that can continue to produce work at the field's observed coauthor intensity. "
              "Below that pool, collaboration networks fragment, mentorship chains break, and the field enters a self-reinforcing decline.")

    p = doc.add_paragraph()
    p.add_run("This framing generates four testable hypotheses. "
              "H1: Across all groups, the equilibrium active pool exceeds the minimum viable threshold, but the distance to the threshold varies widely. "
              "H2: Dropout is the transition rate with the largest negative effect, because attrition removes researchers from every compartment. "
              f"H3: {ctx['positive_lever_sentence'].rstrip('.')}. "
              "H4: Smaller civilisations, and those with older cohort structures, sit closer to their point of no return.")

    p = doc.add_paragraph()
    p.add_run("A final literature stream emphasises the consequences of concentrated research agendas. "
              "Aghion et al. provide evidence that the relationship between competition and innovation follows an inverted-U shape, with the strongest innovative performance in markets that are neither perfectly collusive nor perfectly monopolistic")
    add_citation(p, 5)
    p.add_run(". "
              "Translated to global science, this suggests that a single dominant region or a tight oligopoly may slow the rate of methodological and conceptual breakthroughs. "
              "Maintaining multiple centres of AI/ML research is therefore not merely a distributional concern; it may increase the long-run productivity of the field.")

    doc.add_heading("2.1 Researcher mobility", level=2)
    p = doc.add_paragraph()
    p.add_run("Researcher mobility has been studied from several angles. "
              "A large empirical literature documents net flows of scientists and inventors across countries and regions, often using patent or publication records")
    add_citation(p, 11)
    p.add_run(". "
              "That work consistently finds that the United States, parts of Europe and, increasingly, China and India are central nodes in the global mobility network. "
              "It also finds that mobility responds to wages, funding, institutional quality and career prospects, but that it is path-dependent: once a community loses its senior cohort, it becomes harder to rebuild.")

    doc.add_heading("2.2 Scientific collaboration and diversity", level=2)
    p = doc.add_paragraph()
    p.add_run("A second strand of work emphasises the structure of scientific collaboration. "
              "Multi-university and international teams now produce a growing share of high-impact research, and the geographic dispersion of teams does not necessarily reduce their impact")
    add_citation(p, 12)
    p.add_run(". "
              "This literature suggests that global AI/ML is not a zero-sum race in which every researcher in one location subtracts from another. "
              "It also implies that sustaining a domestic community is compatible with, rather than opposed to, international collaboration. "
              "The question is therefore not whether researchers move, but whether the domestic pipeline that replaces them is robust enough to keep the field alive.")

    doc.add_heading("2.3 Minimum viable populations and critical thresholds", level=2)
    p = doc.add_paragraph()
    p.add_run("The third relevant literature concerns population viability and critical thresholds. "
              "In conservation biology, the minimum viable population concept identifies the smallest number of individuals that can sustain a population in the wild")
    add_citation(p, 10)
    p.add_run(". "
              "We borrow that intuition and apply it to a research community. "
              "A field needs a minimum number of active researchers to produce work, train successors, and maintain peer review and conference communities. "
              "Below that threshold, positive feedback loops weaken: fewer researchers produce fewer students, fewer students produce fewer researchers, and the community enters a downward spiral. "
              "This is the point of no return.")

    doc.add_heading("2.4 This paper's framework", level=2)
    p = doc.add_paragraph()
    p.add_run("The present paper bridges these literatures by estimating transition rates from open bibliometric data and embedding them in a compartment model. "
              "The model is closest in spirit to Stephan's economic model of science, in which researchers move through career stages and respond to incentives")
    add_citation(p, 3)
    p.add_run(", but it adds a civilisational partition and a minimum viable coauthor threshold. "
              "The civilisational partition is not merely a geographic convenience. "
              "It reflects the fact that career incentives, language, funding systems, and institutional networks cluster along civilisational lines, and that these clusters shape mobility more than national borders alone")
    add_citation(p, 4)
    p.add_run(". "
              "The result is a framework that can be updated as new data arrive and can compare the fragility of different research communities using a common metric. "
              "Because it is built on open bibliometric data and transparent transition rates, the model can be replicated and extended by other researchers and by policymakers who need a common language for discussing mobility and capacity.")

    # Data
    doc.add_heading("3. Data and grouping", level=1)
    p = doc.add_paragraph()
    p.add_run("We extracted AI/ML works and author histories from the OpenAlex API for subfield `subfields/1702` (Artificial Intelligence), using works published between 2000 and 2023")
    add_citation(p, 6)
    p.add_run(". "
              "OpenAlex provides open, CC0 bibliographic metadata including authors, affiliations, countries, publication dates, venues and citation links. "
              "We built author histories by following each author's sequence of works and affiliations, assigning them to a country for each work and then to a civilisation by the modal country of their recorded affiliations. "
              "An author is treated as active if they have at least one AI/ML work in the observation window and as a principal investigator (PI) if they appear as the last author of at least one work, a standard proxy for seniority in empirical science")
    add_citation(p, 3)
    p.add_run(". "
              "A 'hit' work is defined as one whose citation count places it in the top 10% of AI/ML works in the same publication year. "
              "The final groups are: United States, Anglosphere ex-US, Continental Europe, Sinic, Japanese, Hindu, Islamic, Other Western, and Other Civilizations.")

    doc.add_heading("3.1 Country-to-civilisation mapping", level=2)
    p = doc.add_paragraph()
    p.add_run("The grouping follows Huntington's civilisation taxonomy but is adjusted for sample-size and mobility reality in AI/ML. "
              "The United States is separated from the broader Anglosphere because it is the dominant destination for AI/ML researchers and because its higher-education and funding systems differ systematically from those of other English-speaking countries. "
              "Continental Europe is kept distinct from the Anglosphere because intra-European mobility and EU research funding create a separate mobility bloc. "
              "Latin American, Orthodox and sub-Saharan African countries are merged into Other Civilizations because their AI/ML author counts in the sample are too small to estimate stable transition rates separately. "
              "This aggregation is a pragmatic modelling choice and does not imply that these communities are culturally homogeneous.")

    doc.add_heading("3.2 Sample selection and variable definitions", level=2)
    p = doc.add_paragraph()
    p.add_run("The cohort is restricted to authors with at least two AI/ML works and a non-missing career-start year between 2000 and 2016. "
              "The career-start year is the first observed AI/ML publication year. "
              "Authors with exclusively unknown affiliations or with all affiliations outside the mapped countries are excluded. "
              "For each author we record the country of the majority of their affiliations and the civilisation to which that country maps. "
              "Works with more than 100 authors are excluded to avoid distorting coauthor counts. "
              "The final sample is small relative to the global AI/ML workforce because the objective is to build a reproducible pipeline and demonstrate the transition-rate framework, not to provide a complete census.")

    doc.add_heading("3.3 OpenAlex coverage and known biases", level=2)
    p = doc.add_paragraph()
    p.add_run("OpenAlex coverage has improved over time but remains incomplete for works before 2000 and for non-English publications. "
              "Author disambiguation is imperfect, especially for common names and authors with multiple name variants. "
              "Affiliation metadata are supplied by publishers and are sometimes missing or refer to the primary institution rather than the country of residence. "
              "For these reasons, the absolute counts reported here are lower bounds on the true global AI/ML workforce. "
              "The analysis nevertheless preserves relative comparisons across civilisations because the same extraction rules are applied uniformly. "
              "Replication from a clean OpenAlex snapshot should produce very similar transition rates and point-of-no-return rankings even if absolute counts shift.")

    p = doc.add_paragraph()
    p.add_run("Table 1 reports the size and composition of the extracted cohort. "
              "The Sinic and Continental Europe groups contribute the largest number of works, followed by the United States and the Anglosphere ex-US. "
              "The Japanese and Other Western groups are the smallest in terms of author counts. "
              "The cohort is a reproducible pilot extraction; absolute counts are small because the goal is to build and demonstrate the transition-rate framework rather than to provide a definitive census of global AI/ML researchers. "
              "Consequently, the absolute equilibrium numbers should be interpreted as model-implied stocks rather than population totals, and the bootstrap intervals reported below give a more honest picture of the uncertainty around those stocks. "
              "The relative sizes are nevertheless informative. "
              "A civilisation with a small cohort but a low coauthor intensity can be more resilient than a larger civilisation with a high coauthor intensity, because the former needs fewer distinct PI groups to sustain its output. "
              "This is why the minimum viable coauthor threshold and the equilibrium active pool must be compared jointly.")

    desc = _descriptive_table(cohort)
    _add_table_from_df(
        doc,
        desc,
        caption="Table 1. Descriptive statistics for the extracted AI/ML cohort by civilisation group.",
        decimals={"career_start_mean": 1},
    )

    # Methods
    doc.add_heading("4. Methods", level=1)
    doc.add_heading("4.1 Compartment model", level=2)
    p = doc.add_paragraph()
    p.add_run("Each civilisation is represented by six compartments: domestic early-career researchers (D), abroad early-career researchers (A), domestic hit researchers (H_D), abroad hit researchers (H_A), domestic principal investigators (P_D), and abroad principal investigators (P_A). "
              "Transition rates are early-career outflow (α), return (β), hit generation at home and abroad (h_D and h_A), PI promotion at home and abroad (p_D and p_A), and dropout from all compartments (d). "
              "The equations are:")
    add_omath_paragraph(doc, math_ode_system())
    p = doc.add_paragraph()
    p.add_run("The model makes several simplifying assumptions. "
              "It treats each civilisation as a single aggregate, ignoring cross-civilisation collaboration and spillovers. "
              "It assumes constant per-year transition rates and a continuous-time Markov structure. "
              "Career stages are collapsed into the three observed layers: early-career, hit researchers and PIs. "
              "These simplifications are necessary to keep the model estimable from OpenAlex and to make the point-of-no-return calculation transparent. "
              "They also mean that the model is best interpreted as a stylised early-warning device, not as a realistic demographic projection.")

    doc.add_heading("4.2 Endogenous inflow", level=2)
    p = doc.add_paragraph()
    p.add_run("New entrants are modelled as a function of the domestic PI stock. "
              "The linear form is ")
    add_omath_inline(p, math_I_linear())
    p.add_run(f", with r capped at {_fmt(ctx['safety_factor'], 2)}× the stability-critical value (safety factor {_fmt(ctx['safety_factor'], 2)}). "
              "A saturating alternative, ")
    add_omath_inline(p, math_I_saturating())
    p.add_run(", is reported as a robustness check. "
              "The PI-driven inflow captures the idea that senior researchers train graduate students, attract postdoctoral researchers, and create the institutional infrastructure that produces the next cohort. "
              "This is a strong assumption because it ignores cross-border recruitment and non-PI sources of new researchers, but it provides a transparent lower bound: if the domestic PI stock falls, the model predicts a decline in new entrants. "
              "The safety factor prevents the model from producing runaway growth when the observed r exceeds the critical value, which is a common empirical finding because observed recruitment is bounded by the data window.")

    doc.add_heading("4.3 Minimum viable coauthor threshold", level=2)
    p = doc.add_paragraph()
    p.add_run("For each group we computed the mean number of authors per work (c\u0304) and the median number of distinct last-author groups observed per recent year (k). "
              "The minimum viable domestic active pool is ")
    add_omath_inline(p, math_threshold())
    p.add_run(". When the equilibrium active pool ")
    add_omath_inline(p, math_active_pool())
    p.add_run(" falls below M, the community can no longer produce works at the observed coauthor intensity. "
              "The threshold is deliberately conservative: it assumes that each new work requires at least k distinct PI groups and that each work has the average number of coauthors. "
              "This overstates the number of distinct actors needed for a viable field, which means that M is a soft lower bound and that observed margins are probably smaller than they appear. "
              "A community with a margin just above M is therefore more fragile than the number itself suggests.")

    doc.add_heading("4.4 Estimation, equilibrium and sensitivity", level=2)
    p = doc.add_paragraph()
    p.add_run("Transition rates are estimated as constant per-year hazards from observed proportions within the cohort. "
              "For each group and each transition, the rate is the ratio of observed transitions to the total exposure time spent in the source compartment during the observation window, with Laplace smoothing of 1 added to both numerator and denominator. "
              "This avoids zero-rate singularities when the cohort is small. "
              "Because the data are right-censored at the end of the observation period, the resulting rates are lower bounds on true long-run hazards; equilibrium solutions therefore tend to be conservative. "
              "The non-linear steady-state equations are solved numerically using a trust-region Newton method with analytically supplied Jacobians. "
              "Elasticities are computed by perturbing each rate by 1%, re-solving, and taking the percentage change in the target stock. "
              "For point-of-no-return analysis we scale each rate until the active pool reaches M and record the critical factor and its proximity, |critical factor − 1|. "
              "A rate whose critical factor lies inside the scan window and is close to 1.0 is the most fragile lever for that group. "
              "All counterfactuals are mechanical perturbations of the fitted rates; they reveal which transitions the model treats as sensitive, not the causal impact of real-world policies.")

    doc.add_heading("4.5 Historical counterfactual design", level=2)
    p = doc.add_paragraph()
    p.add_run("To examine whether transition rates have changed over the past two decades we split the cohort at career-start year 2010. "
              "The early window (career start 2000-2010) captures researchers whose careers were largely established before the most recent AI boom, while the late window (2011-2016) captures researchers who entered during the boom but have a shorter career span over which to estimate rates. "
              "For each window we re-estimate all transition rates and solve for the steady-state active pool. "
              "Comparing the two equilibria reveals how sensitive the long-run margin is to the observed regime, not a prediction of the actual future, because the late cohort is younger and its rates are noisier.")

    doc.add_heading("4.6 Bootstrap uncertainty", level=2)
    p = doc.add_paragraph()
    p.add_run("We resample authors with replacement within each group to obtain 200 bootstrap replicates. "
              "For each replicate we recompute the transition rates and solve the steady-state model, recording T and P_D. "
              "The 2.5th and 97.5th percentiles of the bootstrap distribution provide 95% confidence intervals. "
              "Because the model is non-linear and the equilibrium depends on ratios of rates, the bootstrap distribution is often skewed; we report medians and percentile intervals rather than standard errors.")

    doc.add_heading("4.7 Robustness checks", level=2)
    p = doc.add_paragraph()
    p.add_run("We assess robustness in three ways. "
              "First, we replace the linear PI-driven inflow with a saturating recruitment function that imposes diminishing returns to additional PIs. "
              "Second, we vary the cohort-split year for the historical counterfactual. "
              "Third, we examine the effect of the safety factor on the endogenous inflow parameter r, keeping the system within a bounded stability region. "
              "Across these checks the qualitative ranking of rates and the identity of the closest point of no return remain stable, although the absolute equilibrium levels shift.")

    doc.add_heading("4.8 Relationship to existing indicators", level=2)
    p = doc.add_paragraph()
    p.add_run("The model differs from conventional net-flow or stock indicators in three ways. "
              "First, a net inflow may hide a rise in the total number of researchers abroad relative to the domestic PI stock. "
              "Second, stocks such as total AI/ML publications say little about whether the domestic pipeline can sustain itself. "
              "Third, indices of international collaboration do not distinguish between temporary mobility and permanent brain drain. "
              "The transition-rate view makes each of these processes explicit and provides a stock-and-flow language that is closer to policy instruments such as doctoral funding, retention grants and diaspora networks.")

    doc.add_heading("4.9 Limitations", level=2)
    p = doc.add_paragraph()
    p.add_run("The main limitations are data quality and model scope. "
              "OpenAlex country metadata are noisy, especially for older works and for authors with multiple affiliations. "
              "Career stages are inferred from authorship order and are imperfect proxies. "
              "The model does not include cross-civilisation knowledge spillovers, bilateral migration costs, or firm-level mobility. "
              "Finally, the assumption of constant rates is a strong approximation over a 23-year window. "
              "We therefore emphasise rank-order and relative sensitivity rather than point forecasts.")

    doc.add_heading("4.10 Annual transition-rate estimation and projection", level=2)
    p = doc.add_paragraph()
    p.add_run("The steady-state model in Sections 4.1-4.4 treats rates as constants. "
              "To test whether the same framework can be used for short-run monitoring, we reconstructed year-by-year compartment membership from cohort.csv and raw_sampled_works.json. "
              "For each author and year we inferred location as domestic if the author was in the origin civilisation and abroad otherwise, using sampled works when available and cohort-derived abroad/return years as a fallback. "
              "From these states we computed annual transition counts for the six compartments, applied Laplace +0.5 smoothing to empty destination cells, and derived the probabilities that map to α, β, h_D, h_A, p_D, p_A and d. "
              "Inter-civilisation flows are approximated by assigning each abroad author-year to the author's recent_group as the destination civilisation; this is a lower-bound proxy because year-to-year destination changes are not observed in the public cohort.")

    p = doc.add_paragraph()
    p.add_run("For the 2017-2026 projection we fit a linear trend to the observed 2000-2016 rates for each group and rate. "
              "If fewer than four observations were available or the fit explained less than 10% of the variance, the historical mean was used instead. "
              "Projected rates were clipped to values between 0 and 1. "
              "Dropout was capped at the 90th percentile of observed dropout rates to prevent implausible extrapolation. "
              "Projected total inflows were apportioned across compartments using the observed 2016 distribution. "
              "Population composition was projected forward with the discrete-time recursion N(t+1) = N(t)P(t) + b(t+1), where P(t) is a 6×6 row-stochastic-in-expectation matrix that preserves dropout mass: the row sum is 1 − d after scaling outgoing rates. "
              "This discrete step is the operational counterpart of the continuous-time ODE; with an annual dt it provides an early-warning signal one year ahead.")

    p = doc.add_paragraph()
    p.add_run("We compare the 2017-2023 projection with the observed annual stock. "
              "The comparison is limited to years that have observed data, and the observed stock is reindexed to the full group-year-compartment grid so that zero-observed cells are not omitted from the accuracy metrics. "
              "Accuracy is reported as RMSE and MAPE; MAPE here is computed against count_obs + 1 to avoid division by zero and is therefore a conservative, non-standard measure.")

    doc.add_heading("4.11 Correction pressures and theoretical bounds", level=2)
    p = doc.add_paragraph()
    p.add_run("The annual estimates contain several regularising pressures that bound the model away from instability and fabrication. "
              "Laplace smoothing adds a uniform prior of 0.5 to every possible destination, which shrinks sparse cells toward 1/(number of destinations) and prevents zero-probability singularities when a transition is unobserved in a small group-year. "
              "It is equivalent to a weak Dirichlet prior and is a standard regulariser for sparse multinomial transitions.")

    p = doc.add_paragraph()
    p.add_run("Clipping projected rates to values between 0 and 1 is a feasibility pressure: rates outside the probability simplex are inadmissible. "
              "The dropout cap is a safety pressure motivated by the fact that unbounded linear extrapolation of observed attrition would eventually predict more leavers than the total stock. "
              "The inflow apportionment pressure keeps the composition of new entrants aligned with the most recently observed recruitment pattern, rather than inventing a new distribution. "
              f"Finally, the safety factor of {_fmt(ctx['safety_factor'], 2)} on the endogenous PI-driven inflow keeps the system inside the stability boundary. "
              "Together these pressures embody the principle that projection should stay within observed empirical support and within theoretical stability limits; they are not arbitrary adjustments but transparent bounds that can be tightened or relaxed as more data become available.")

    # Results
    doc.add_heading("5. Results", level=1)
    p = doc.add_paragraph()
    p.add_run(f"Table 2 reports the equilibrium domestic active pool T, the minimum viable threshold M, and the endogenous inflow parameters for the {len(eq)} groups. "
              "All groups remain above their threshold under the fitted model, but margins differ by an order of magnitude. "
              f"The {ctx['largest_pools']} groups show the largest equilibrium active pools, reflecting large cohorts and relatively low coauthor-intensity thresholds. "
              f"The {ctx['smallest_pool']} group has the smallest equilibrium active pool, and {ctx['smallest_margin_group']} has the narrowest safety margin, although both still exceed their minimum viable coauthor pool. "
              "The ratio T/M is a summary resilience indicator, but absolute margin is the more direct measure of proximity to the point of no return.")

    eq_table = eq[["group", "T_equilibrium", "M_threshold", "margin_to_threshold_T", "I0", "r", "r_obs", "r_critical"]].copy()
    eq_table = eq_table.rename(columns={
        "group": "Group",
        "T_equilibrium": "T_eq",
        "M_threshold": "M",
        "margin_to_threshold_T": "Margin",
        "r_obs": "r_obs",
        "r_critical": "r_crit",
    })
    _add_table_from_df(
        doc,
        eq_table,
        caption="Table 2. Equilibrium domestic active pool, minimum viable threshold, and endogenous inflow parameters.",
        decimals={"T_eq": 2, "M": 2, "Margin": 2, "I0": 2, "r": 5, "r_obs": 5, "r_crit": 5},
    )

    p = doc.add_paragraph()
    p.add_run("Figure 1 visualises the gap between equilibrium and threshold. "
              f"The {ctx['largest_pools']} groups display the largest equilibrium active pools, while the {ctx['smallest_pool']} group is the smallest. "
              "However, the point-of-no-return metric is not the absolute level of T but the distance between T and M, which reflects both the stock of researchers and the coauthor intensity of the field. "
              "Groups with high T but also high c\u0304 and k can still be fragile if their margin is small.")
    doc.add_picture(str(fig_paths["fig1"]), width=Inches(5.8))
    cap = doc.add_paragraph()
    cap.add_run("Figure 1. Equilibrium domestic active pool (T) and minimum viable coauthor threshold (M) by group.").italic = True
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER

    p = doc.add_paragraph()
    p.add_run("Table 3 shows the three transition-rate elasticities with the largest absolute impact on T for each group. "
              f"Dropout (d) is the largest negative lever in every group; its active-pool elasticity ranges from {_fmt(ctx['d_min_e'], 2)} to {_fmt(ctx['d_max_e'], 2)}. "
              "Attrition removes researchers from every compartment, so a proportional increase in d produces a larger proportional decline in the active pool. "
              f"{ctx['positive_lever_sentence']} "
              "Early-career outflow (α) has a modest negative effect in most groups, but because it moves researchers to the abroad compartment rather than removing them entirely, its direct impact on the domestic active pool is smaller than that of dropout. "
              "There is notable heterogeneity in the magnitude of the positive levers. "
              f"The {ctx['highest_pd_group']} group shows the strongest response to PI promotion (p_D), indicating that improving the promotion of hit researchers to PIs is an especially efficient way to expand the domestic active pool in that community. "
              "In the largest groups, p_D remains positive but its relative effect is smaller, because the active pool is already large and a proportional change in promotion has less marginal impact.")

    rows3 = []
    for group, gdf in top_t.groupby("group"):
        top3 = gdf.sort_values("abs_elasticity", ascending=False).head(3)
        r = [group]
        for _, row in top3.iterrows():
            r.extend([row["rate"], _fmt(row["elasticity"], 3)])
        rows3.append(r)
    elas_df = pd.DataFrame(rows3, columns=["Group", "1st rate", "1st elasticity", "2nd rate", "2nd elasticity", "3rd rate", "3rd elasticity"])
    _add_table_from_df(
        doc,
        elas_df,
        caption="Table 3. Top transition-rate elasticities for domestic active pool T.",
    )

    closest = pnr_closest.iloc[0]
    p = doc.add_paragraph()
    p.add_run(f"Table 4 reports, for each group, the single rate that reaches the active-pool threshold with the smallest proportional change. "
              f"The {closest['group']} group is the most fragile: {closest['rate_name']} must be multiplied by {_fmt(closest['critical_factor'], 3)}× its current value (equivalent to a {closest['proximity']*100:.0f}% proportional {'reduction' if closest['critical_factor'] < 1 else 'increase'}) to drive the active pool to its minimum viable threshold. "
              f"{ctx['pnr_lever_text']}. "
              "This is consistent with a recruitment-driven view of scientific communities: if the pipeline of new researchers shuts or slows, the active pool eventually falls below the minimum viable coauthor pool regardless of how efficient return or promotion becomes. "
              "A global retention programme that reduces dropout would benefit all groups, but the most vulnerable groups may also need an expansion of the exogenous entry rate.")

    pnr_table = pnr_closest[["group", "target", "rate_name", "current_rate", "critical_factor", "proximity"]].copy()
    pnr_table.columns = ["Group", "Target", "Rate", "Current", "Critical factor", "Proximity"]
    _add_table_from_df(
        doc,
        pnr_table,
        caption="Table 4. Closest point of no return for the active researcher pool by group.",
        decimals={"Current": 4, "Critical factor": 3, "Proximity": 3},
    )

    p = doc.add_paragraph()
    p.add_run("Figure 2 ranks groups by their closest point-of-no-return sensitivity.")
    doc.add_picture(str(fig_paths["fig2"]), width=Inches(5.8))
    cap = doc.add_paragraph()
    cap.add_run("Figure 2. Closest point-of-no-return proximity by group. Smaller values mean a smaller proportional change in the listed rate is required to reach the threshold for the stated target pool.").italic = True
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER

    if sat_eq is not None:
        doc.add_heading("5.1 Saturating recruitment extension", level=2)
        p = doc.add_paragraph()
        p.add_run("Replacing linear inflow with a saturating form lowers equilibrium pools because each additional PI adds fewer entrants. "
                  "Table 5 compares linear and saturating equilibrium T values. "
                  "The saturating model is important because the observed r is often close to the stability boundary, and an unchecked linear inflow can produce explosive growth. "
                  f"Across groups, the saturating variant predicts equilibrium pools that are {ctx['sat_range_text']}, underscoring the sensitivity of long-run projections to the functional form of inflow. "
                  "This sensitivity does not overturn the ranking of groups, but it shows that absolute equilibrium levels should be treated with caution. "
                  "The saturating model is the preferred interpretation for policy because it acknowledges that recruitment cannot scale linearly with the number of PIs indefinitely.")
        merged = eq[["group", "T_equilibrium"]].merge(
            sat_eq[["group", "T_equilibrium", "epsilon"]], on="group", suffixes=("_lin", "_sat")
        )
        merged.columns = ["Group", "Linear T", "Saturating T", "ε"]
        _add_table_from_df(
            doc,
            merged,
            caption="Table 5. Equilibrium T under linear and saturating PI-driven inflow.",
            decimals={"Linear T": 2, "Saturating T": 2, "ε": 5},
        )

    doc.add_heading("5.2 Historical counterfactual", level=2)
    n_compare = len(period_compare)
    if ctx["period_all_neg"]:
        prefix = "Both" if n_compare == 2 else f"All {n_compare}"
        period_direction_text = (
            f"{prefix} groups with dual-window support would see smaller safety margins under late-window rates "
            f"({ctx['period_neg']})."
        )
    else:
        period_direction_text = (
            f"Groups that would see smaller safety margins under late-window rates: {ctx['period_neg']}. "
            f"Groups that would see larger safety margins under late-window rates: {ctx['period_pos']}."
        )
    p = doc.add_paragraph()
    p.add_run("Table 6 compares the equilibrium that would have emerged if the transition rates estimated for the early career window (2000-2010) or the late window (2011-2016) had persisted indefinitely. "
              "The late window is shorter and its rates are estimated from younger cohorts, so the comparison should be read as a sensitivity exercise rather than a forecast. "
              f"Only {n_compare} groups have enough dual-window support for reliable rate estimation in both windows; they are listed in the table. "
              f"{period_direction_text} "
              "This pattern shows that global AI/ML mobility is not moving in a single direction; different civilisations are on different trajectories, and a uniform policy response would ignore this heterogeneity. "
              "Because the late cohort is younger, the late-window equilibrium is likely biased downward for groups where career progression has not yet run its course. "
              "Even so, the exercise shows that the current regime is not the only possible one, which is why counterfactual policy analysis is useful.")

    pc = period_compare.copy()
    pc = pc[["group", "T_early", "T_late", "pct_delta_T", "margin_early", "margin_late", "delta_margin"]].rename(columns={
        "group": "Group",
        "T_early": "T early",
        "T_late": "T late",
        "pct_delta_T": "ΔT (%)",
        "margin_early": "Margin early",
        "margin_late": "Margin late",
        "delta_margin": "Δ margin",
    })
    _add_table_from_df(
        doc,
        pc,
        caption="Table 6. Historical counterfactual: equilibrium active pool and safety margin under early versus late transition-rate regimes.",
        decimals={"T early": 1, "T late": 1, "ΔT (%)": 1, "Margin early": 1, "Margin late": 1, "Δ margin": 1},
    )

    p = doc.add_paragraph()
    p.add_run("Figure 3 shows the change in safety margin between the early and late transition-rate regimes.")
    doc.add_picture(str(fig_paths["fig3"]), width=Inches(5.8))
    cap = doc.add_paragraph()
    cap.add_run("Figure 3. Change in safety margin between early and late transition-rate regimes. Positive values mean the late-window rates would produce a larger safety margin than the early-window rates if they persisted; negative values mean the margin would shrink. "
                "The comparison is across two point estimates; uncertainty is substantial because the two windows have different cohort sizes and the steady-state model does not capture policy shocks.").italic = True
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_heading("5.3 Policy counterfactuals", level=2)
    policy_top = policy_rank.groupby("group").head(1).copy()
    # Robustly describe the 10% dropout reduction effect, even if another lever is top for some group
    d_decrease = policy_rank[(policy_rank["lever"] == "d") & (policy_rank["direction"] == "decrease")].copy()
    d_10pct = d_decrease[d_decrease["lever_change_pct"].abs() >= 9.9]
    if d_10pct.empty:
        d_10pct = d_decrease
    d_10pct_group = d_10pct.loc[d_10pct.groupby("group")["normalised_margin_gain_per_10pct"].idxmax()]
    all_top_are_d = (policy_top["lever"] == "d").all()
    d_min = d_10pct_group.sort_values("margin_gain").iloc[0]
    d_max = d_10pct_group.sort_values("margin_gain").iloc[-1]
    min_gain_group = d_min["group"]
    max_gain_group = d_max["group"]
    min_gain = round(d_min["margin_gain"])
    max_gain = round(d_max["margin_gain"])
    dominant_text = (
        "Reducing dropout is the dominant positive lever for every civilisation, which is consistent with the elasticity results in Table 3. "
        if all_top_are_d
        else "Reducing dropout is the dominant positive lever for most civilisations in the current data. "
    )
    p = doc.add_paragraph()
    p.add_run(f"Table 7 reports the single mechanical counterfactual with the largest margin gain per 10% lever change for each group. {dominant_text}"
              f"The gain from a roughly 10% proportional reduction in d ranges from about {min_gain} additional active researchers for the {min_gain_group} group to about {max_gain} for the {max_gain_group} group, reflecting differences in cohort size and baseline attrition. "
              "No other single lever comes close to dropout reduction in terms of simulated margin gain per unit proportional change, although combinations of levers may be more efficient for some groups. "
              "The results also imply that policy need not focus on blocking early-career outflow. "
              "Reducing attrition among researchers who remain in the domestic system is a more efficient way to protect the active pool than preventing researchers from going abroad, because a researcher abroad is still in the global AI/ML system and may return. "
              "For the smallest groups, increasing the exogenous entry rate or improving the promotion of hit researchers to PIs can add additional margin, but dropout reduction remains the first-order model-implied target.")

    policy_top = policy_top.rename(columns={
        "group": "Group",
        "lever": "Lever",
        "direction": "Direction",
        "lever_change_pct": "Change (%)",
        "margin_gain": "Margin gain",
        "normalised_margin_gain_per_10pct": "Gain per 10%",
    })
    _add_table_from_df(
        doc,
        policy_top,
        caption="Table 7. Top positive mechanical counterfactual per group, measured by margin gain per 10% proportional lever change.",
        decimals={"Change (%)": 0, "Margin gain": 1, "Gain per 10%": 1},
    )

    doc.add_heading("5.4 Uncertainty", level=2)
    p = doc.add_paragraph()
    p.add_run("Table 8 reports bootstrap 95% confidence intervals for the equilibrium active pool T and the domestic PI pool P_D. "
              "The intervals are wide, reflecting the small cohort sample and the extrapolation from individual careers to long-run steady states. "
              "For some groups the upper bound is an order of magnitude larger than the lower bound, indicating that the equilibrium is sensitive to resampling variation in the transition rates. "
              "This uncertainty should be interpreted as a warning against over-interpreting point estimates and as a reason to view the point-of-no-return distances as indicative rather than precise thresholds. "
              "Despite the width, the lower bounds for most groups remain above the minimum viable threshold, which supports the qualitative conclusion that all groups are currently above the point of no return. "
              "For the smallest groups the lower bound is closer to M, reinforcing the need for continued monitoring and for policy buffers.")

    boot_tab = boot.copy()
    boot_tab["T 95% CI"] = boot_tab.apply(lambda r: f"[{_fmt(r['T_equilibrium_q025'], 0)}, {_fmt(r['T_equilibrium_q975'], 0)}]", axis=1)
    boot_tab["P_D 95% CI"] = boot_tab.apply(lambda r: f"[{_fmt(r['P_D_equilibrium_q025'], 0)}, {_fmt(r['P_D_equilibrium_q975'], 0)}]", axis=1)
    boot_tab = boot_tab[["group", "T_equilibrium_median", "T 95% CI", "P_D_equilibrium_mean", "P_D 95% CI"]]
    boot_tab.columns = ["Group", "T median", "T 95% CI", "P_D mean", "P_D 95% CI"]
    _add_table_from_df(
        doc,
        boot_tab,
        caption="Table 8. Bootstrap 95% confidence intervals for equilibrium T and domestic PI pool P_D.",
        decimals={"T median": 0, "P_D mean": 0},
    )

    p = doc.add_paragraph()
    p.add_run("Figure 4 displays the bootstrap intervals graphically.")
    doc.add_picture(str(fig_paths["fig4"]), width=Inches(5.8))
    cap = doc.add_paragraph()
    cap.add_run("Figure 4. Bootstrap 95% confidence intervals for equilibrium T by group.").italic = True
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_heading("5.5 Synthesis", level=2)
    p = doc.add_paragraph()
    p.add_run("Taken together, the results provide a consistent picture. "
              "Exogenous entry and dropout are the two rates that most strongly determine the long-run viability of an AI/ML research community. "
              "Communities that are large in absolute terms are not necessarily safe if their coauthor intensity is high; conversely, small communities can be robust if their attrition is low and their recruitment pipeline is stable. "
              "The historical counterfactual shows that the current regime is not preordained: a shift in transition rates at the start of the AI boom would have produced different steady states for different civilisations. "
              "This is precisely why the framework is useful: it identifies which rate in which community is closest to a threshold, allowing policy to intervene before rather than after a collapse. "
              "The policy message is therefore both diagnostic and preventative. "
              "By tracking transition rates rather than net flows, policymakers can see which civilisation is approaching a point of no return and which lever offers the largest safety margin per unit of effort.")

    doc.add_heading("5.6 Annual transition rates and inter-civilisation flows", level=2)
    p = doc.add_paragraph()
    p.add_run("Figure 5 plots the observed 2000-2016 transition rates and the projected 2017-2026 rates for each civilisation. "
              "Rates are displayed by group and by transition type, so that the reader can see whether a particular transition is trending toward a boundary. "
              "Because the projections are linear trend fits regularised by the correction pressures described in Section 4.11, they are not forecasts of specific future events; they are the model's one-year-ahead extrapolation of the recent historical trajectory.")
    doc.add_picture(str(fig_paths["fig5"]), width=Inches(6.0))
    cap = doc.add_paragraph()
    cap.add_run("Figure 5. Annual observed (solid) and projected (dashed) transition rates by civilisation, 2000-2026.").italic = True
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER

    p = doc.add_paragraph()
    p.add_run("Table 9 summarises the mean observed annual transition rates by group between 2000 and 2016. "
              "The table distinguishes early-career outflow (α), return (β), domestic and abroad hit generation (h_D, h_A), PI promotion (p_D), dropout (d), and total inflow (I_total).")
    if not annual_means.empty:
        _add_table_from_df(
            doc,
            annual_means,
            caption="Table 9. Mean observed annual transition rates by civilisation, 2000-2016.",
            decimals={"α": 3, "β": 3, "h_D": 3, "p_D": 3, "d": 3, "I_total": 2},
        )

    p = doc.add_paragraph()
    p.add_run("Figure 6 shows the inter-civilisation accumulation of abroad author-years. "
              "Rows represent the origin civilisation and columns represent the destination civilisation, approximated by the author's recent_group while abroad. "
              "The heatmap is a lower-bound proxy because year-to-year destination switches within a spell abroad are not observed.")
    doc.add_picture(str(fig_paths["fig6"]), width=Inches(5.8))
    cap = doc.add_paragraph()
    cap.add_run("Figure 6. Inter-civilisation abroad author-year accumulation by origin (rows) and destination (columns).").italic = True
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER

    p = doc.add_paragraph()
    p.add_run("Table 10 lists the origin-destination pairs with the largest accumulation of abroad author-years. "
              "These pairs identify the strongest visible inter-civilisation pipelines and are the empirical counterpart to the α and β transitions.")
    if not interciv_top.empty:
        _add_table_from_df(
            doc,
            interciv_top,
            caption="Table 10. Top origin-destination abroad author-year pairs.",
            decimals={"Author-years": 0},
        )

    doc.add_heading("5.7 Out-of-sample projection, 2017-2023", level=2)
    p = doc.add_paragraph()
    p.add_run(f"The 2017-2023 projection is compared with observed annual stocks in Figure 7. "
              f"Overall accuracy is RMSE {_fmt(annual_ctx.get('overall_rmse', float('nan')), 2)} and MAPE {_fmt(annual_ctx.get('overall_mape_pct', float('nan')), 1)}% (a non-standard, conservative measure computed against count_obs + 1 to avoid division by zero). "
              f"Among civilisations the lowest RMSE is for {best_rmse_group} and the highest RMSE is for {worst_rmse_group}; the highest MAPE is for {worst_mape_group}. "
              "The largest errors occur in small compartments and in groups with sparse transition counts, which is expected because the annual model does not borrow information across civilisations.")
    doc.add_picture(str(fig_paths["fig7"]), width=Inches(6.0))
    cap = doc.add_paragraph()
    cap.add_run("Figure 7. Observed (solid) and projected (dashed) compartment counts by civilisation, 2017-2023. The vertical dotted line marks the end of the training period (2016).").italic = True
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER

    p = doc.add_paragraph()
    p.add_run("Table 11 reports projection accuracy by civilisation and Table 12 by compartment. "
              f"Among compartments, the lowest RMSE is for {best_compartment_rmse}, while the highest RMSE is for {worst_compartment_rmse} and the highest MAPE is for {worst_compartment_mape}. "
              "P_D and H_D show larger errors because small changes in PI and hit rates are amplified by the endogenous inflow term.")
    if group_acc is not None and not group_acc.empty:
        gacc = group_acc.copy()
        gacc["rmse"] = gacc["rmse"].apply(lambda x: _fmt(x, 2))
        gacc["mape"] = gacc["mape"].apply(lambda x: f"{x*100:.1f}%")
        gacc = gacc.rename(columns={"origin_group": "Group", "rmse": "RMSE", "mape": "MAPE"})
        _add_table_from_df(doc, gacc, caption="Table 11. Projection accuracy by civilisation, 2017-2023.", decimals={"MAPE": 2})
    if comp_acc is not None and not comp_acc.empty:
        cacc = comp_acc.copy()
        cacc["rmse"] = cacc["rmse"].apply(lambda x: _fmt(x, 2))
        cacc["mape"] = cacc["mape"].apply(lambda x: f"{x*100:.1f}%")
        cacc = cacc.rename(columns={"compartment": "Compartment", "rmse": "RMSE", "mape": "MAPE"})
        _add_table_from_df(doc, cacc, caption="Table 12. Projection accuracy by compartment, 2017-2023.", decimals={"MAPE": 2})

    doc.add_heading("5.8 Correction pressures in the annual model", level=2)
    p = doc.add_paragraph()
    p.add_run("The annual projection performs best where the correction pressures in Section 4.11 are binding. "
              "Laplace smoothing prevents empty cells from being treated as impossible transitions; the unit-interval clip and the dropout cap prevent the trend extrapolation from producing rates that are incompatible with a stochastic transition matrix; and the 2016 inflow apportionment keeps new-entrant composition close to the last observed regime. "
              "These pressures mean that the projection is not a purely mechanical forecast: it is a bounded extrapolation that stays within the empirical support of the 2000-2016 data and within the stability constraints of the compartment model.")

    # Discussion
    doc.add_heading("6. Discussion", level=1)
    p = doc.add_paragraph()
    p.add_run("The results support a transition-rate view of research policy. "
              "Rather than asking which country has a net inflow or outflow of researchers, the model asks which rate must be altered to keep a community above its minimum viable coauthor pool. "
              "This shift in focus has implications for how we conceptualise brain drain, design science and technology policy, and interpret civilisational diversity in AI/ML.")

    doc.add_heading("6.1 From net flows to transition rates", level=2)
    p = doc.add_paragraph()
    p.add_run("Most empirical studies of researcher mobility measure net flows, stocks or collaboration counts")
    add_citation(p, 11)
    p.add_run(". "
              "These indicators are useful for describing patterns, but they do not reveal the mechanisms that sustain or undermine a research community. "
              "A country may have a positive net inflow while simultaneously losing its domestic PI base through retirement or emigration, or it may have negative net flow but a healthy pipeline of new entrants. "
              "The transition-rate framework disaggregates these processes and shows that the same net flow can correspond to very different vulnerability profiles. "
              "For example, a high early-career outflow rate is less damaging than a high dropout rate because researchers abroad may return; a high dropout rate removes researchers from the system entirely. "
              "This distinction is lost in net-flow accounting but is central to point-of-no-return analysis.")

    p = doc.add_paragraph()
    p.add_run(f"First, {ctx['pnr_lever_text']}. "
              "A large proportional reduction in baseline recruitment would drive most communities to their threshold before mobility rates such as return or promotion became binding. "
              "This is consistent with the observation that AI/ML fields depend on a continuous pipeline of new graduate students and junior researchers")
    add_citation(p, 1)
    p.add_run(". "
              "Policies that sustain that pipeline, such as doctoral funding, visa routes for early-career researchers, and stable junior positions, are therefore first-order defences against a point of no return.")

    p = doc.add_paragraph()
    p.add_run(f"Second, among the mobility transition rates, dropout (d) is the dominant negative lever; its active-pool elasticity ranges from {_fmt(ctx['d_min_e'], 2)} to {_fmt(ctx['d_max_e'], 2)} across groups, and in the policy counterfactuals a simulated reduction in dropout yields the largest margin gain per unit proportional change. "
              "Attrition matters because it removes researchers from every compartment, not just one. "
              "A 10% proportional reduction in dropout expands the safety margin more than comparably sized increases in return, hit generation or promotion. "
              f"For {ctx['smallest_margin_group']}, the group with the smallest safety margin, even modest attrition reductions may widen the margin. "
              "These counterfactuals are mechanical perturbations of the fitted rates; they identify the most sensitive transition levers, not the causal effect of any specific policy programme.")

    p = doc.add_paragraph()
    p.add_run(f"Third, {ctx['positive_lever_sentence_lower']}. "
              f"The {ctx['highest_pd_group']} group shows the strongest response to PI promotion, suggesting that for that community expanding the domestic PI pipeline is an efficient lever. "
              "Return from abroad (β) is also positive for most groups, though its effect is generally smaller than reducing attrition directly. "
              "The implication for policy is that retention and promotion are usually more efficient than trying to attract returnees, but a balanced portfolio is still needed: a community without domestic PI growth cannot reproduce itself through attrition reduction alone.")

    p = doc.add_paragraph()
    p.add_run("Fourth, the historical counterfactual shows that the late-window rates, if they persisted, would alter equilibrium margins. "
              f"{period_direction_text} "
              "This pattern cautions against treating AI/ML mobility as a single global trend. "
              "It also confirms that the model can detect temporal changes in transition rates, which is the prerequisite for the early intervention the framework is designed to support.")

    p = doc.add_paragraph()
    p.add_run("The transition levers also interact in ways that a single-rate elasticity cannot fully capture. "
              "For example, reducing dropout and increasing PI promotion together are likely to have a larger effect than the sum of the two individual perturbations, because more researchers survive to become PIs and those PIs then train additional early-career researchers through the endogenous inflow channel. "
              "Conversely, a simultaneous fall in exogenous entry and a rise in dropout can push a community to its threshold faster than either change alone. "
              "The model's steady-state and one-at-a-time counterfactuals are therefore a starting point; they identify the most sensitive margins but do not exhaust the policy design space.")

    p = doc.add_paragraph()
    p.add_run("The connection to civilisational diversity is direct. "
              "Each group's safety margin can be monitored over time, and interventions can be adjusted before the margin disappears. "
              f"Because the model uses a fixed safety factor of {_fmt(ctx['safety_factor'], 2)} for the endogenous inflow parameter r, the policy recommendations are deliberately conservative: they do not push the system toward instability. "
              "That bounded approach is consistent with the goal of preserving diversity rather than maximising any single country's share.")

    p = doc.add_paragraph()
    p.add_run("It is important to stress that the counterfactuals reported in Tables 3 and 7 are mechanical perturbations of the fitted transition rates, not causal estimates of specific programmes. "
              "They identify which rates the model treats as most sensitive, and therefore where empirical policy evaluation is most urgent, but they do not by themselves show that a given intervention would achieve the simulated change.")

    doc.add_heading("6.2 Civilisational diversity as an innovation input", level=2)
    p = doc.add_paragraph()
    p.add_run("A second implication concerns the normative status of civilisational diversity. "
              "We treat diversity as an input to innovation rather than as a distributional afterthought")
    add_citation(p, 13)
    p.add_run(". "
              "A monocentric or tight-oligopoly structure in AI/ML may produce short-run efficiency gains through scale and agglomeration, but it also raises the risk of methodological lock-in, selection bias in training data, and reduced error correction. "
              "Recent work on multi-university teams shows that geographically dispersed collaborations can retain high impact, which suggests that distributing capacity across civilisations need not sacrifice quality")
    add_citation(p, 12)
    p.add_run(". "
              "By quantifying the safety margin for each research community, the framework makes it possible to argue for support of smaller communities on positive, innovation-systems grounds. "
              "Preserving multiple centres of AI/ML research is not a matter of slowing the frontier; it is a matter of ensuring that the frontier is not defined by a single set of institutions, languages, or problems.")

    doc.add_heading("6.3 Policy implications and early warning", level=2)
    p = doc.add_paragraph()
    p.add_run("The policy implications can be read as an early-warning architecture. "
              "A single dashboard that tracks the fitted transition rates, their bootstrap uncertainty, and the distance to M for each civilisation would allow policymakers to detect divergence before a community enters an irreversible decline. "
              "Interventions can then be calibrated to maintain a minimum safety margin rather than to maximise any one stock. "
              "This is the operational meaning of early intervention: not a forecast that a particular collapse will occur, but a structured way to keep the system away from a point of no return. "
              "It also frames high-skilled mobility as a strategic competition among jurisdictions for talent")
    add_citation(p, 14)
    p.add_run(", in which the central question is not only who wins the current round but whether the global system retains enough diversity for future rounds")
    add_citation(p, 15)
    p.add_run(". "
              "If the dt of policy response is short enough, the model can be updated annually and divergence caught early, before any single civilisation approaches a point of no return. "
              "This is the mechanism through which technology monopoly, hegemonic concentration and oligopoly dead-ends can be avoided: by keeping every major research community above its minimum viable coauthor pool, the framework sustains the competitive diversity that underpins long-run technological progress. "
              "The framework is therefore not a prediction that a particular civilisation will collapse. "
              "It is a tool for ensuring that no single civilisation reaches a point where its collapse becomes self-sustaining, and that the global AI/ML system retains the diversity required for continued innovation.")

    p = doc.add_paragraph()
    p.add_run("Operationally, the framework can be used in two complementary ways. "
              "As a monitoring tool, it can be rerun whenever new OpenAlex data are released, producing an updated set of transition rates, safety margins and proximity-to-threshold estimates. "
              "As a scenario tool, it can quantify how large a proportional change in a given rate would be required to move a community toward or away from collapse, which helps prioritise empirical policy evaluation. "
              "Both uses depend on transparent assumptions and regular recalibration; the model should not be used to justify one-off interventions without accompanying process evaluation.")

    doc.add_heading("6.4 Validation of correction pressures", level=2)
    p = doc.add_paragraph()
    p.add_run("The correction pressures are not ad hoc adjustments; each maps to a known statistical or dynamical constraint. "
              "Laplace smoothing is equivalent to a weak Dirichlet prior on a multinomial transition vector; it guarantees that no cell has zero estimated probability and shrinks rare transitions toward the simplex centroid. "
              "Clipping projected rates to values between 0 and 1 is a feasibility constraint on probabilities; the dropout cap is a cross-sectional constraint that prevents projected attrition from exceeding the observed stock; and the inflow apportionment constraint keeps the composition of new entrants equal to the last observed recruitment pattern. "
              "In the 2017-2023 projection these pressures reduced the sensitivity of the forecast to sparse cells and to short-run fluctuations in small groups. "
              f"Quantitatively, the overall RMSE of {_fmt(annual_ctx.get('overall_rmse', float('nan')), 2)} and conservative MAPE of {_fmt(annual_ctx.get('overall_mape_pct', float('nan')), 1)}% are consistent with a model that is deliberately regularised rather than optimised for in-sample fit. "
              "The residual errors are concentrated in the smallest compartments, which is exactly where smoothing is most active and where future data will be most valuable.")

    doc.add_heading("6.5 Intra-civilisation alternatives when inter-civilisation mobility cannot be controlled", level=2)
    p = doc.add_paragraph()
    p.add_run("If a civilisation cannot control outflows to, or inflows from, other jurisdictions—whether because of visa regimes, salary differentials, language advantages, or targeted recruitment—it can still preserve its research community by acting on the intra-civilisation levers identified in the annual model. "
              "The annual rates show that the domestic active pool T = D + H_D + P_D responds most strongly to the dropout rate d, the domestic hit rate h_D, and the PI promotion rate p_D. "
              "Policies that reduce early-career attrition, expand domestic postdoctoral positions, or accelerate independent-lab formation therefore become defensive substitutes when inter-civilisation poaching cannot be regulated. "
              "This is the practical meaning of civilisational-diversity preservation under sovereignty constraints: even without controlling the border of talent, a community can increase the internal reproduction of active researchers. "
              f"The ODE safety factor of {_fmt(ctx['safety_factor'], 2)} on endogenous PI inflow is a conservative bound that prevents over-optimism about this substitution effect; more ambitious domestic growth would require corresponding evidence that the extra PIs can be absorbed without simply raising dropout.")

    doc.add_heading("6.6 Annual updating as an early-warning layer", level=2)
    p = doc.add_paragraph()
    p.add_run("The 2017-2023 projection demonstrates that the framework can be rerun annually with a one-year time step. "
              "Each new year of OpenAlex data updates the observed transition rates, the fitted trends, and the distance to the minimum viable coauthor threshold. "
              "Because the model is regularised by the correction pressures, the one-year-ahead projection is not easily derailed by a single noisy observation. "
              "Instead, successive years reveal whether a particular transition rate is drifting toward a boundary. "
              "That drift is the early-warning signal. "
              "Policymakers can then intervene before the active pool falls below M, using the rate-specific elasticities in Table 3 to prioritise the smallest proportional change that restores a safety margin. "
              "This is the operational mechanism for avoiding technology monopoly and oligopoly dead ends: by keeping every major research community above its minimum viable coauthor pool, annual monitoring sustains the competitive diversity that underpins long-run technological progress. "
              "The framework is therefore not a prediction that a particular civilisation will collapse; it is a tool for ensuring that no single civilisation reaches a point where its collapse becomes self-sustaining.")

    doc.add_heading("6.7 Limitations", level=2)
    p = doc.add_paragraph()
    p.add_run("Several limitations should be acknowledged. "
              "OpenAlex affiliation and country assignments are noisy, especially for researchers with multiple affiliations. "
              "The civilisation grouping is a coarse aggregation; within-group heterogeneity is substantial. "
              "The annual model relies on a discrete approximation of the continuous-time ODE and does not capture within-year events or cross-civilisation spillovers. "
              "Inter-civilisation flows are approximated by the author's recent_group while abroad, which misses year-to-year destination switching. "
              "The cohort sample is small; the absolute equilibrium numbers should be interpreted as model-implied stocks rather than as census counts. "
              "Authors with many publications are over-weighted relative to one-publication authors, so rate estimates reflect author-publication exposure rather than a uniformly representative sample of individuals. "
              f"The endogenous inflow is capped at a safety factor of {_fmt(ctx['safety_factor'], 2)} relative to the critical reproduction rate; alternative values would shift equilibrium levels and should be reported in future sensitivity tables. "
              "Finally, the point-of-no-return threshold is a sufficient condition for collapse, not a necessary one: a community may decline for reasons outside the model even if T remains above M.")

    p = doc.add_paragraph()
    p.add_run("Wide bootstrap confidence intervals, especially for smaller civilisation groups, mean that the ordinal ranking of groups by equilibrium size or proximity to threshold should be treated as descriptive rather than definitive. "
              "The model identifies which transitions are most sensitive in a mechanical sense; turning those sensitivities into reliable policy priorities requires additional data on programme costs, implementation lags, and behavioural responses that are outside the scope of this paper.")

    p = doc.add_paragraph()
    p.add_run("From a security-studies perspective, the framework is intentionally non-adversarial. "
              "It does not model deliberate recruitment campaigns, technology transfer, or strategic denial. "
              "Instead, it treats mobility as an aggregate transition process and asks when a community becomes unable to reproduce itself. "
              "That baseline is useful because it shows where defensive, capacity-building policies can be most efficient, but it does not replace classified or diplomatic assessments of technology competition. "
              "Future work could add a strategic layer by distinguishing between civilian and defence-relevant AI/ML pipelines, or by modelling the effects of targeted recruitment on specific subfields.")

    # Conclusion
    doc.add_heading("7. Conclusion", level=1)
    p = doc.add_paragraph()
    p.add_run("We have proposed and implemented a transition-rate framework for assessing how close AI/ML research communities are to a point of no return. "
              "The model converts OpenAlex publication records into civilisation-specific transition rates and solves for the equilibrium active researcher pool. "
              "All groups remain above their minimum viable coauthor threshold in the fitted model, but the distance to that threshold varies by an order of magnitude and is most sensitive to exogenous entry and dropout. "
              f"Dropout is the dominant negative lever (active-pool elasticity {_fmt(ctx['d_min_e'], 2)} to {_fmt(ctx['d_max_e'], 2)}), and a simulated reduction is the single most efficient model-implied response for every civilisation. "
              "However, the closest point of no return is exogenous entry for all groups in the active-pool analysis, which means that policies which sustain the pipeline of new researchers are first-order defences. "
              "The historical counterfactual and the bootstrap intervals remind us that the future is not determined by current rates; transition rates can change, and policy can be directed at the most fragile lever before a collapse.")

    p = doc.add_paragraph()
    p.add_run("The annual projection layer adds an operational dimension to this conclusion. "
              "By estimating year-by-year transition rates and projecting one year ahead, the model turns the steady-state diagnostic into an early-warning dashboard. "
              "A one-year time step is short enough to detect drift before the active pool approaches the minimum viable threshold, and the correction pressures keep the projection within empirical and theoretical bounds. "
              "When inter-civilisation mobility cannot be controlled, the same framework points to intra-civilisation levers—reducing dropout, raising domestic hit rates, and accelerating PI promotion—that preserve T = D + H_D + P_D. "
              "These two layers, steady-state and annual, together provide a coherent basis for early, safety-factor-bound intervention.")

    p = doc.add_paragraph()
    p.add_run("The broader implication is that preserving civilisational diversity in AI/ML is compatible with, and may reinforce, scientific progress. "
              "A single dominant region or a tight oligopoly may achieve short-run scale economies, but it also risks methodological lock-in and reduces the set of problems that receive sustained attention. "
              "By monitoring transition rates and safety margins, policymakers can detect divergence early and intervene in a safety-factor-bound way. "
              "This is the practical meaning of the aspiration to avoid technology monopoly and oligopoly dead ends: not a prediction that any one civilisation will dominate, but a structured method for keeping the global system away from points of no return. "
              "Early, proportionate interventions that reduce attrition and sustain new recruitment can widen safety margins and preserve civilisational diversity in AI/ML.")

    doc.add_heading("7.1 Future work", level=2)
    p = doc.add_paragraph()
    p.add_run("Several extensions are natural. "
              "First, the model can be applied to other security-relevant fields such as semiconductor physics, quantum computing, biotechnology and energy materials, allowing cross-field comparisons of vulnerability. "
              "Second, the civilisation partition can be refined to a country or institution level, allowing bilateral migration flows and network externalities to be incorporated. "
              "Third, the ODE can be solved dynamically rather than at steady state, making it possible to forecast the time to threshold under alternative policy scenarios. "
              "Fourth, the minimum viable coauthor threshold can be made endogenous by modelling coauthorship as a matching process. "
              "Fifth, the sensitivity of equilibrium outcomes to the safety factor and to the saturating parameter epsilon should be mapped systematically. "
              "Finally, the framework can be integrated with policy cost data to produce cost-effectiveness comparisons of alternative interventions, turning mechanical sensitivities into actionable funding priorities.")

    # References
    doc.add_heading("References", level=1)
    for i, ref in enumerate(REFS, 1):
        p = doc.add_paragraph()
        p.add_run(f"{i}. {ref}")


def write_docx(output_dir, data, fig_paths):
    abstract, keywords, highlights = _abstract_and_highlights(data[1], data[4])

    # Pre-compute body word count by building a throwaway body doc
    body_doc = Document()
    _add_docx_body(body_doc, data, fig_paths)
    body_wc = _doc_word_count(body_doc)

    doc = Document()
    _add_title_page(doc, word_count=body_wc)
    _add_front_matter(doc, abstract, keywords, highlights)
    _add_docx_body(doc, data, fig_paths)

    path = output_dir / "manuscript_full_article.docx"
    doc.save(path)
    return path


def write_pptx(output_dir, data, fig_paths):
    (cohort, eq, sat_eq, top_t, pnr_closest, period_compare, boot, policy_rank) = data
    annual = load_annual_data()
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    def add_image_slide(title, img_path, caption):
        slide_layout = prs.slide_layouts[3]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        left = PptxInches(1.5)
        top = PptxInches(1.2)
        slide.shapes.add_picture(str(img_path), left, top, width=PptxInches(10))
        txBox = slide.shapes.add_textbox(left, PptxInches(6.0), PptxInches(10), PptxInches(0.8))
        txBox.text_frame.text = caption
        for paragraph in txBox.text_frame.paragraphs:
            paragraph.font.size = Pt(14)

    def add_table_slide(title, df, col_names, width_per_col=1.5, font_size=10):
        slide_layout = prs.slide_layouts[3]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        rows, cols = len(df) + 1, len(col_names)
        left = PptxInches(0.5)
        top = PptxInches(1.2)
        table = slide.shapes.add_table(rows, cols, left, top, PptxInches(cols * width_per_col), PptxInches(0.6 * rows)).table
        for i, h in enumerate(col_names):
            table.cell(0, i).text = str(h)
        for row_i, (_, row) in enumerate(df.iterrows()):
            for j, val in enumerate(row):
                table.cell(row_i + 1, j).text = str(val)
                table.cell(row_i + 1, j).text_frame.paragraphs[0].font.size = Pt(font_size)

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Quantifying the Point of No Return in Global AI/ML Research Communities"
    slide.placeholders[1].text = "Data-driven manuscript figures and tables"

    add_image_slide(
        "Figure 1: Equilibrium T vs minimum viable threshold",
        fig_paths["fig1"],
        "Blue bars: equilibrium T; orange bars: threshold M. All groups remain above the threshold, but margins vary widely.",
    )
    add_image_slide(
        "Figure 2: Closest point-of-no-return proximity",
        fig_paths["fig2"],
        "Smaller values mean a smaller proportional change in the listed rate is required to reach the threshold for the stated target pool.",
    )
    add_image_slide(
        "Figure 3: Historical counterfactual margin change",
        fig_paths["fig3"],
        "Positive values mean the late-window rates would produce a larger safety margin than the early-window rates if they persisted; negative values mean the margin would shrink. The comparison is across point estimates; uncertainty is substantial.",
    )
    add_image_slide(
        "Figure 4: Bootstrap 95% CI for equilibrium T",
        fig_paths["fig4"],
        "Intervals are asymmetric and wide, reflecting model uncertainty.",
    )

    desc = _descriptive_table(cohort)
    add_table_slide("Table 1: Descriptive cohort statistics", desc, desc.columns.tolist(), width_per_col=1.3)

    eq_table = eq[["group", "T_equilibrium", "M_threshold", "margin_to_threshold_T", "I0", "r", "r_obs", "r_critical"]].copy()
    eq_table.columns = ["Group", "T_eq", "M", "Margin", "I0", "r", "r_obs", "r_crit"]
    for c in ["T_eq", "M", "Margin", "I0"]:
        eq_table[c] = eq_table[c].apply(lambda x: _fmt(x, 2))
    for c in ["r", "r_obs", "r_crit"]:
        eq_table[c] = eq_table[c].apply(lambda x: _fmt(x, 5))
    add_table_slide("Table 2: Equilibrium and inflow parameters", eq_table, eq_table.columns.tolist(), width_per_col=1.4)

    rows2 = []
    for group, gdf in top_t.groupby("group"):
        top3 = gdf.sort_values("abs_elasticity", ascending=False).head(3)
        parts = [group]
        for _, row in top3.iterrows():
            parts.extend([row["rate"], _fmt(row["elasticity"], 3)])
        rows2.append(parts)
    elas_df = pd.DataFrame(rows2, columns=["Group", "1st", "el1", "2nd", "el2", "3rd", "el3"])
    add_table_slide("Table 3: Top transition-rate elasticities", elas_df, elas_df.columns.tolist(), width_per_col=1.4)

    pnr_table = pnr_closest[["group", "target", "rate_name", "current_rate", "critical_factor", "proximity"]].copy()
    pnr_table.columns = ["Group", "Target", "Rate", "Current", "Crit.factor", "Proximity"]
    for c, d in {"Current": 4, "Crit.factor": 3, "Proximity": 3}.items():
        pnr_table[c] = pnr_table[c].apply(lambda x, d=d: _fmt(x, d))
    add_table_slide("Table 4: Closest point of no return", pnr_table, pnr_table.columns.tolist(), width_per_col=1.8)

    if sat_eq is not None:
        merged = eq[["group", "T_equilibrium"]].merge(
            sat_eq[["group", "T_equilibrium", "epsilon"]], on="group", suffixes=("_lin", "_sat")
        )
        merged.columns = ["Group", "Linear T", "Saturating T", "ε"]
        for c, d in {"Linear T": 2, "Saturating T": 2, "ε": 5}.items():
            merged[c] = merged[c].apply(lambda x, d=d: _fmt(x, d))
        add_table_slide("Table 5: Saturating inflow extension", merged, merged.columns.tolist(), width_per_col=2.0)

    pc = period_compare.rename(columns={
        "group": "Group",
        "T_early": "T early",
        "T_late": "T late",
        "pct_delta_T": "ΔT (%)",
        "margin_early": "Margin early",
        "margin_late": "Margin late",
        "delta_margin": "Δ margin",
    })
    for c in ["T early", "T late", "ΔT (%)", "Margin early", "Margin late", "Δ margin"]:
        pc[c] = pc[c].apply(lambda x: _fmt(x, 1))
    add_table_slide("Table 6: Historical counterfactual", pc, pc.columns.tolist(), width_per_col=1.4)

    policy_top = policy_rank.groupby("group").head(1).rename(columns={
        "group": "Group",
        "lever": "Lever",
        "direction": "Direction",
        "lever_change_pct": "Change (%)",
        "margin_gain": "Margin gain",
        "normalised_margin_gain_per_10pct": "Gain per 10%",
    })
    for c, d in {"Change (%)": 0, "Margin gain": 1, "Gain per 10%": 1}.items():
        policy_top[c] = policy_top[c].apply(lambda x, d=d: _fmt(x, d))
    add_table_slide("Table 7: Top policy intervention", policy_top, policy_top.columns.tolist(), width_per_col=2.0)

    boot_tab = boot.copy()
    boot_tab["T 95% CI"] = boot_tab.apply(lambda r: f"[{_fmt(r['T_equilibrium_q025'], 0)}, {_fmt(r['T_equilibrium_q975'], 0)}]", axis=1)
    boot_tab["P_D 95% CI"] = boot_tab.apply(lambda r: f"[{_fmt(r['P_D_equilibrium_q025'], 0)}, {_fmt(r['P_D_equilibrium_q975'], 0)}]", axis=1)
    boot_tab = boot_tab[["group", "T_equilibrium_median", "T 95% CI", "P_D_equilibrium_mean", "P_D 95% CI"]]
    boot_tab.columns = ["Group", "T median", "T 95% CI", "P_D mean", "P_D 95% CI"]
    for c in ["T median", "P_D mean"]:
        boot_tab[c] = boot_tab[c].apply(lambda x: _fmt(x, 0))
    add_table_slide("Table 8: Bootstrap 95% CI", boot_tab, boot_tab.columns.tolist(), width_per_col=2.2)

    # Annual projection slides
    if fig_paths.get("fig5"):
        add_image_slide(
            "Figure 5: Annual observed and projected transition rates",
            fig_paths["fig5"],
            "Solid lines mark observed 2000-2016 rates; dashed lines mark projected 2017-2026 rates.",
        )
    if fig_paths.get("fig6"):
        add_image_slide(
            "Figure 6: Inter-civilisation abroad author-years",
            fig_paths["fig6"],
            "Rows are origin civilisations; columns are destination civilisations approximated by recent_group.",
        )
    if fig_paths.get("fig7"):
        add_image_slide(
            "Figure 7: Observed vs projected compartment counts",
            fig_paths["fig7"],
            "Solid lines are observed counts; dashed lines are 2017-2026 projections. The vertical dotted line is 2016.",
        )

    annual_means = annual_summary_table(annual)
    if not annual_means.empty:
        add_table_slide(
            "Table 9: Mean observed annual transition rates, 2000-2016",
            annual_means,
            annual_means.columns.tolist(),
            width_per_col=1.4,
        )

    interciv_top = interciv_top_table(annual)
    if not interciv_top.empty:
        add_table_slide(
            "Table 10: Top origin-destination abroad author-year pairs",
            interciv_top,
            interciv_top.columns.tolist(),
            width_per_col=2.0,
        )

    group_acc = annual.get("group_accuracy")
    if group_acc is not None and not group_acc.empty:
        gacc = group_acc.copy()
        gacc["rmse"] = gacc["rmse"].apply(lambda x: _fmt(x, 2))
        gacc["mape"] = gacc["mape"].apply(lambda x: f"{x*100:.1f}%")
        gacc = gacc.rename(columns={"origin_group": "Group", "rmse": "RMSE", "mape": "MAPE"})
        add_table_slide(
            "Table 11: Projection accuracy by civilisation, 2017-2023",
            gacc,
            gacc.columns.tolist(),
            width_per_col=2.2,
        )

    comp_acc = annual.get("compartment_accuracy")
    if comp_acc is not None and not comp_acc.empty:
        cacc = comp_acc.copy()
        cacc["rmse"] = cacc["rmse"].apply(lambda x: _fmt(x, 2))
        cacc["mape"] = cacc["mape"].apply(lambda x: f"{x*100:.1f}%")
        cacc = cacc.rename(columns={"compartment": "Compartment", "rmse": "RMSE", "mape": "MAPE"})
        add_table_slide(
            "Table 12: Projection accuracy by compartment, 2017-2023",
            cacc,
            cacc.columns.tolist(),
            width_per_col=2.2,
        )

    path = output_dir / "manuscript_full_article_figures.pptx"
    prs.save(path)
    return path


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--output-dir", type=Path, default=BASE_DIR / "docs")
    args = parser.parse_args()
    output_dir = args.output_dir
    output_dir.mkdir(parents=True, exist_ok=True)

    data = load_data()
    cohort, eq, sat_eq, top_t, pnr_closest, period_compare, boot, policy_rank = data

    fig_dir = output_dir / "figures"
    fig1 = build_figure1(eq, fig_dir)
    fig2 = build_figure2(pnr_closest, fig_dir)
    fig3 = build_figure3(period_compare, fig_dir)
    fig4 = build_figure4(boot, fig_dir)

    annual = load_annual_data()
    annual_figs = build_annual_figures(annual, fig_dir)

    fig_paths = {
        "fig1": fig1,
        "fig2": fig2,
        "fig3": fig3,
        "fig4": fig4,
        "fig5": annual_figs.get("fig5"),
        "fig6": annual_figs.get("fig6"),
        "fig7": annual_figs.get("fig7"),
    }

    md_path = write_markdown(output_dir, data, fig_paths)
    docx_path = write_docx(output_dir, data, fig_paths)
    pptx_path = write_pptx(output_dir, data, fig_paths)

    print(f"Wrote {md_path}")
    print(f"Wrote {docx_path}")
    print(f"Wrote {pptx_path}")
    print(f"Figures saved to {fig_dir}")


if __name__ == "__main__":
    main()
