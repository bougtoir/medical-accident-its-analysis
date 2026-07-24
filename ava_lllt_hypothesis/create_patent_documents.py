#!/usr/bin/env python3
"""Generate formal Japanese patent application documents (JPO format).

Outputs:
  patent/meisaisho.docx            -- Specification
  patent/claims.docx               -- Claims
  patent/abstract.docx             -- Abstract
  patent/rejection_strategy.docx   -- Rejection strategy (internal)
  patent/fig1_mechanism.png        -- Figure 1: Mechanism
  patent/fig2_device.png           -- Figure 2: Device config
  patent/fig3_control.png          -- Figure 3: Control flowchart
  patent/fig4_pulse.png            -- Figure 4: Pulse timing
  patent/fig5_comparison.png       -- Figure 5: Prior art comparison
  patent/fig6_staged.png           -- Figure 6: Staged upgrade
  patent/drawings.pptx             -- All figures editable
"""
from pathlib import Path
from docx import Document
from docx.shared import Pt, Mm, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

OUT = Path(__file__).parent / "patent"
OUT.mkdir(exist_ok=True)
FONT = "MS Mincho"
SZ = Pt(11)


def _sf(run, name=FONT, size=SZ, bold=False):
    run.font.name = name
    run.font.size = size
    run.font.bold = bold
    rpr = run._element.find(qn("w:rPr"))
    if rpr is None:
        rpr = OxmlElement("w:rPr")
        run._element.append(rpr)
    rf = rpr.find(qn("w:rFonts"))
    if rf is None:
        rf = OxmlElement("w:rFonts")
        rpr.append(rf)
    rf.set(qn("w:eastAsia"), name)


def _doc():
    d = Document()
    for s in d.sections:
        s.page_width, s.page_height = Mm(210), Mm(297)
        s.top_margin = s.bottom_margin = s.left_margin = s.right_margin = Mm(25)
    st = d.styles["Normal"]
    st.font.name = FONT
    st.font.size = SZ
    return d


def _h(doc, text, level=1):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    p.paragraph_format.space_before = Pt(14 if level == 1 else 10)
    p.paragraph_format.space_after = Pt(6)
    r = p.add_run(text)
    _sf(r, bold=True, size=Pt(13 if level == 1 else 11))
    return p


def _np(doc, n, text):
    """Numbered paragraph like 【0001】."""
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    p.paragraph_format.space_after = Pt(4)
    p.paragraph_format.line_spacing = Pt(18)
    r1 = p.add_run(f"\u3010{n:04d}\u3011")
    _sf(r1, bold=True)
    r2 = p.add_run(f"\n{text}")
    _sf(r2)
    return p


def _p(doc, text, bold=False):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    p.paragraph_format.space_after = Pt(4)
    p.paragraph_format.line_spacing = Pt(18)
    r = p.add_run(text)
    _sf(r, bold=bold)
    return p


def _tbl(doc, headers, rows):
    t = doc.add_table(rows=1 + len(rows), cols=len(headers))
    t.alignment = WD_TABLE_ALIGNMENT.CENTER
    t.style = "Table Grid"
    for i, h in enumerate(headers):
        c = t.rows[0].cells[i]
        r = c.paragraphs[0].add_run(h)
        _sf(r, bold=True, size=Pt(9))
        c.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
    for ri, row in enumerate(rows):
        for ci, v in enumerate(row):
            c = t.rows[ri + 1].cells[ci]
            r = c.paragraphs[0].add_run(str(v))
            _sf(r, size=Pt(9))
    return t


# ==========================================================
# 1. MEISAISHO
# ==========================================================
def build_meisaisho():
    doc = _doc()
    _h(doc, "\u3010\u66f8\u985e\u540d\u3011\u660e\u7d30\u66f8")
    _h(doc, "\u3010\u767a\u660e\u306e\u540d\u79f0\u3011")
    _p(doc, "\u5149\u751f\u4f53\u8abf\u7bc0\u306b\u3088\u308b\u52d5\u9759\u8108\u543e\u5408\u9078\u629e\u7684\u62e1\u5f35\u3092\u5229\u7528\u3057\u305f"
       "\u30a6\u30a7\u30a2\u30e9\u30d6\u30eb\u672b\u68a2\u8840\u6d41\u4fc3\u9032\u30c7\u30d0\u30a4\u30b9\u304a\u3088\u3073\u305d\u306e\u5236\u5fa1\u65b9\u6cd5")

    _h(doc, "\u3010\u6280\u8853\u5206\u91ce\u3011")
    _np(doc, 1,
        "\u672c\u767a\u660e\u306f\u3001\u624b\u638c\u304a\u3088\u3073/\u307e\u305f\u306f\u8db3\u5e95\u306e\u7121\u6bdb\u76ae\u819a\uff08glabrous skin\uff09\u306b\u5b58\u5728\u3059\u308b"
        "\u52d5\u9759\u8108\u543e\u5408\uff08AVA: Arteriovenous Anastomoses\uff09\u3092\u5149\u751f\u4f53\u8abf\u7bc0"
        "\uff08PBM: Photobiomodulation\uff09\u306b\u3088\u308a\u9078\u629e\u7684\u306b\u62e1\u5f35\u3055\u305b\u3001\u672b\u68a2\u8840\u6d41\u3092\u5897\u52a0\u3055\u305b\u308b"
        "\u30a6\u30a7\u30a2\u30e9\u30d6\u30eb\u30c7\u30d0\u30a4\u30b9\u304a\u3088\u3073\u305d\u306e\u5236\u5fa1\u65b9\u6cd5\u306b\u95a2\u3059\u308b\u3002\u3088\u308a\u5177\u4f53\u7684\u306b\u306f\u3001"
        "630\u301c850 nm\u5e2f\u306e\u975e\u71b1\u7684\u5149\u7167\u5c04\u306b\u3088\u308a\u8840\u7ba1\u5185\u76ae\u7d30\u80de\u306b\u84c4\u7a4d\u3055\u308c\u305f"
        "NO\uff08\u4e00\u9178\u5316\u7a92\u7d20\uff09\u3092\u5149\u5316\u5b66\u7684\u306b\u904a\u96e2\u3055\u305b\u3001AVA\u8840\u7ba1\u5e73\u6ed1\u7b4b\u3092\u5f1b\u7de9\u3055\u305b\u308b\u3053\u3068\u3067\u3001"
        "\u672b\u68a2\u8840\u7ba1\u62b5\u6297\uff08TPR: Total Peripheral Resistance\uff09\u3092\u4f4e\u4e0b\u3055\u305b\u308b\u30c7\u30d0\u30a4\u30b9\u306b\u95a2\u3059\u308b\u3002")

    _h(doc, "\u3010\u80cc\u666f\u6280\u8853\u3011")

    bg_paras = [
        (2, "\u9ad8\u8840\u5727\u306e\u85ac\u7269\u6cbb\u7642\u306f\u3001\u964d\u5727\u85ac\uff08\u5229\u5c3f\u85ac\u3001Ca\u62ee\u6297\u85ac\u3001ACE\u963b\u5bb3\u85ac\u3001ARB\u7b49\uff09\u306e"
            "\u9577\u671f\u670d\u7528\u3092\u8981\u3059\u308b\u3002\u7279\u306b\u7b2c3\u9078\u629e\u85ac\u3067\u3042\u308b\u30b5\u30a4\u30a2\u30b6\u30a4\u30c9\u7cfb\u5229\u5c3f\u85ac\u306f\u4f4e\u30ab\u30ea\u30a6\u30e0\u8840\u75c7\u3001"
            "\u9ad8\u5c3f\u9178\u8840\u75c7\u3001\u8010\u7cd6\u80fd\u969c\u5bb3\u7b49\u306e\u4ee3\u8b1d\u6027\u526f\u4f5c\u7528\u3092\u4f34\u3044\u3001\u670d\u85ac\u30a2\u30c9\u30d2\u30a2\u30e9\u30f3\u30b9\u306e\u4f4e\u4e0b\u3092\u62db\u304f"
            "\uff08\u975e\u7279\u8a31\u6587\u732e6\u53c2\u7167\uff09\u3002\u4e16\u754c\u7684\u306b\u9ad8\u8840\u5727\u60a3\u8005\u306e\u7d0450%\u304c\u51e6\u65b9\u85ac\u30921\u5e74\u4ee5\u5185\u306b\u4e2d\u65ad\u3059\u308b"
            "\u3068\u3055\u308c\u3001\u975e\u85ac\u7269\u7684\u4ecb\u5165\u3078\u306e\u9700\u8981\u306f\u5927\u304d\u3044\u3002"),
        (3, "\u7279\u8a31\u6587\u732e2\uff08US 8,679,170 B2, Muehlbauer et al., 2014\uff09\u306b\u958b\u793a\u3055\u308c\u308b"
            "AVACEN\u30c7\u30d0\u30a4\u30b9\u306f\u3001\u624b\u638c\u3092\u525b\u6027\u30c1\u30e3\u30f3\u30d0\u30fc\u306b\u633f\u5165\u3057\u3001\u9670\u5727\uff08\u221230 mmHg\uff09\u3068"
            "\u4f1d\u5c0e\u52a0\u6e29\uff0842\u00b0C\uff09\u3092\u540c\u6642\u306b\u5370\u52a0\u3057\u3066AVA\u8840\u6d41\u3092\u5897\u52a0\u3055\u305b\u308b\u3002\u636e\u3048\u7f6e\u304d\u578b\u88c5\u7f6e"
            "\uff08\u7d045 kg\uff09\u3067\u3042\u308a\u3001\u6cbb\u7642\u4e2d\u306f\u4e00\u65b9\u306e\u624b\u304c\u30c1\u30e3\u30f3\u30d0\u30fc\u5185\u306b\u62d8\u675f\u3055\u308c\u308b\u305f\u3081\u3001"
            "\u65e5\u4e2d\u306e\u30a6\u30a7\u30a2\u30e9\u30d6\u30eb\u4f7f\u7528\u306f\u4e0d\u53ef\u80fd\u3067\u3042\u308b\u3002\u9670\u5727\u6a5f\u69cb\uff08\u30c0\u30a4\u30e4\u30d5\u30e9\u30e0\u30dd\u30f3\u30d7\u3001"
            "\u6c17\u5bc6\u30b7\u30fc\u30eb\u3001\u5727\u529b\u30bb\u30f3\u30b5\uff09\u306e\u5b58\u5728\u304c\u5c0f\u578b\u5316\u30fb\u643a\u5e2f\u6027\u30fb\u7121\u97f3\u5316\u3092\u6839\u672c\u7684\u306b\u5236\u9650\u3057\u3066\u3044\u308b\u3002"),
        (4, "AVACEN\u306e\u964d\u5727\u30e1\u30ab\u30cb\u30ba\u30e0\u306f\u3001(i) \u9670\u5727\u306b\u3088\u308b\u7d4c\u58c1\u5727\u4f4e\u4e0b\u2192\u53d7\u52d5\u7684\u8840\u7ba1\u62e1\u5f35\u3001"
            "\u304a\u3088\u3073 (ii) \u4f1d\u5c0e\u52a0\u71b1\u306b\u3088\u308bTRPV4\u30c1\u30e3\u30cd\u30eb\u6d3b\u6027\u5316\u2192\u5e73\u6ed1\u7b4b\u5f1b\u7de9\u3001\u306e2\u3064\u306b\u4f9d\u5b58\u3059\u308b\u3002"
            "\u3044\u305a\u308c\u3082\u300c\u71b1\u30a8\u30cd\u30eb\u30ae\u30fc\uff0b\u6a5f\u68b0\u30a8\u30cd\u30eb\u30ae\u30fc\u300d\u306b\u3088\u308b\u7269\u7406\u7684\u30a2\u30d7\u30ed\u30fc\u30c1\u3067\u3042\u308a\u3001"
            "\u5149\u5316\u5b66\u7684NO\u653e\u51fa\u6a5f\u5e8f\u306f\u542b\u307e\u308c\u306a\u3044\u3002"),
        (5, "\u7279\u8a31\u6587\u732e1\uff08US 11,229,548 B2, Diller et al., 2022\uff09\u306f\u3001\u9838\u690e/\u8170\u690e\u9818\u57df\u3078\u306e"
            "\u6e29\u71b1\u523a\u6fc0\uff08Selective Thermal Stimulation: STS\uff09\u306b\u3088\u308a\u3001\u8996\u5e8a\u4e0b\u90e8\u3092\u4ecb\u3057\u305f"
            "\u53cd\u5c04\u7684AVA\u62e1\u5f35\u3092\u60f9\u8d77\u3057\u3001\u8840\u5727\u3092\u30d5\u30a3\u30fc\u30c9\u30d0\u30c3\u30af\u5236\u5fa1\u3067\u7ba1\u7406\u3059\u308b\u30b7\u30b9\u30c6\u30e0\u3092\u958b\u793a\u3059\u308b\u3002"
            "\u30d9\u30c3\u30c9\u7d71\u5408\u578b\u30de\u30c3\u30c8\u30ec\u30b9\u5f62\u614b\u3067\u3042\u308a\u3001\u4f7f\u7528\u306f\u7761\u7720\u4e2d\u306b\u9650\u5b9a\u3055\u308c\u308b\u3002"
            "\u964d\u5727\u6a5f\u5e8f\u306f\u4e2d\u67a2\u6027\u53cd\u5c04\u306b\u4f9d\u5b58\u3057\u3001\u672b\u68a2AVA\u3078\u306e\u76f4\u63a5\u7684\u4f5c\u7528\u3067\u306f\u306a\u304f\u3001"
            "\u52b9\u679c\u767a\u73fe\u307e\u3067\u306e\u6642\u9593\u9045\u5ef6\uff0810\u301c30\u5206\uff09\u3092\u4f34\u3046\u3002"),
        (6, "\u975e\u7279\u8a31\u6587\u732e1\uff08Lobo MD et al., Lancet 2015\uff09\u306f\u3001\u8178\u9aa8\u52d5\u9759\u8108\u9593\u306b\u30cb\u30c1\u30ce\u30fc\u30eb\u88fd"
            "\u30ab\u30d7\u30e9\u30fc\u3092\u7d4c\u76ae\u7684\u306b\u7559\u7f6e\u3057\u3001\u56fa\u5b9a\u7684AV\u543e\u5408\u3092\u4f5c\u6210\u3057\u3066\u964d\u5727\u3059\u308b\u624b\u6cd5\u3092\u5831\u544a\u3059\u308b"
            "\uff08ROX CONTROL HTN\u8a66\u9a13\u3001SBP \u221226.9 mmHg\uff09\u3002\u4fb5\u8972\u7684\u624b\u6280\u3092\u8981\u3057\u3001\u9759\u8108\u72ed\u7a84"
            "\uff0829%\uff09\u306e\u5408\u4f75\u75c7\u3092\u4f34\u3046\u304c\u3001\u300cAVA/AV shunt\u958b\u901a\u2192TPR\u4f4e\u4e0b\u2192\u8840\u5727\u964d\u4e0b\u300d\u3068\u3044\u3046"
            "\u30e1\u30ab\u30cb\u30ba\u30e0\u306e\u6709\u52b9\u6027\u3092\u9ad8\u3044\u30ec\u30d9\u30eb\u306e\u30a8\u30d3\u30c7\u30f3\u30b9\uff08RCT\uff09\u3067\u8a3c\u660e\u3057\u305f\u3002"),
        (7, "\u5149\u751f\u4f53\u8abf\u7bc0\uff08PBM\uff09\u30c7\u30d0\u30a4\u30b9\u306f\u3001\u5275\u50b7\u6cbb\u7652\u4fc3\u9032\u3001\u75bc\u75db\u7de9\u548c\u3001\u708e\u75c7\u6291\u5236\u3092\u76ee\u7684\u3068\u3057\u3066"
            "\u591a\u6570\u304c\u516c\u77e5\u3067\u3042\u308b\u3002\u3053\u308c\u3089\u306e\u65e2\u5b58PBM\u30c7\u30d0\u30a4\u30b9\u306b\u304a\u3044\u3066\u3001(i) \u7167\u5c04\u90e8\u4f4d\u306fglabrous skin\u306b"
            "\u9650\u5b9a\u3055\u308c\u305a\u5168\u8eab\u306e\u4efb\u610f\u90e8\u4f4d\u3092\u5bfe\u8c61\u3068\u3057\u3001(ii) \u6cbb\u7642\u76ee\u7684\u306f\u5c40\u6240\u7684\u7d44\u7e54\u4fee\u5fa9\u3067\u3042\u308a"
            "\u5168\u8eab\u7684\u8840\u884c\u52d5\u614b\u3078\u306e\u5f71\u97ff\u306f\u610f\u56f3\u3055\u308c\u3066\u304a\u3089\u305a\u3001(iii) AVA\u306e\u9078\u629e\u7684\u62e1\u5f35\u3092\u4ecb\u3057\u305f"
            "TPR\u4f4e\u4e0b\u3068\u3044\u3046\u4f5c\u7528\u6a5f\u5e8f\u306f\u958b\u793a\u3082\u793a\u5506\u3082\u3055\u308c\u3066\u3044\u306a\u3044\u3002"),
        (8, "\u975e\u7279\u8a31\u6587\u732e3\uff08Ribeiro BG et al., Lasers Med Sci 2022\uff09\u306f\u3001\u30e9\u30c3\u30c8\u306b\u304a\u3044\u3066"
            "660 nm PBM\u304c\u8840\u5727\u3092\u4f4e\u4e0b\u3055\u305b\u308b\u3053\u3068\u3092\u5831\u544a\u3057\u3001NOS\u963b\u5bb3\u5264\uff08L-NAME\uff09\u6295\u4e0e\u3067"
            "\u964d\u5727\u52b9\u679c\u304c\u6d88\u5931\u3059\u308b\u3053\u3068\u304b\u3089NO\u4f9d\u5b58\u6027\u3092\u8a3c\u660e\u3057\u305f\u3002\u3057\u304b\u3057\u3001\u3053\u306e\u7814\u7a76\u306f "
            "(i) \u7167\u5c04\u90e8\u4f4d\u3092\u5c3e\u52d5\u8108\u3068\u3057\u3066\u304a\u308aAVA\u5bc6\u96c6\u90e8\u4f4d\u3067\u306f\u306a\u3044\u3001"
            "(ii) \u30e9\u30c3\u30c8\u306e\u5168\u8eab\u306b\u7167\u5c04\u3057\u3066\u304a\u308a\u90e8\u4f4d\u9078\u629e\u6027\u304c\u306a\u3044\u3001"
            "(iii) \u30a6\u30a7\u30a2\u30e9\u30d6\u30eb\u30c7\u30d0\u30a4\u30b9\u3068\u3057\u3066\u306e\u5b9f\u88c5\u306f\u4e00\u5207\u691c\u8a0e\u3055\u308c\u3066\u3044\u306a\u3044\u3002"),
        (9, "\u4e0a\u8a18\u5f93\u6765\u6280\u8853\u3092\u7dcf\u62ec\u3059\u308b\u3068\u3001\u4ee5\u4e0b\u306e\u5168\u3066\u3092\u540c\u6642\u306b\u6e80\u305f\u3059\u30c7\u30d0\u30a4\u30b9\u306f\u5b58\u5728\u3057\u306a\u3044\uff1a\n"
            "(A) \u975e\u4fb5\u8972\u3067\u3042\u308b\u3053\u3068\n"
            "(B) \u30a6\u30a7\u30a2\u30e9\u30d6\u30eb\u5f62\u614b\u3067\u65e5\u4e2d\u88c5\u7740\u53ef\u80fd\u3067\u3042\u308b\u3053\u3068\n"
            "(C) \u9670\u5727\u6a5f\u69cb\u304c\u4e0d\u8981\u3067\u3042\u308b\u3053\u3068\n"
            "(D) \u4f4e\u6d88\u8cbb\u96fb\u529b\uff08\u30d0\u30c3\u30c6\u30ea\u30fc\u99c6\u52d54\u6642\u9593\u4ee5\u4e0a\uff09\u3067\u3042\u308b\u3053\u3068\n"
            "(E) glabrous skin\u306eAVA\u3092\u300c\u9078\u629e\u7684\u306b\u300d\u62e1\u5f35\u3059\u308b\u6a5f\u5e8f\u3092\u6709\u3059\u308b\u3053\u3068\n"
            "(F) NO\u306e\u5149\u5316\u5b66\u7684\u653e\u51fa\u3068\u3044\u3046\u975e\u71b1\u7684\u6a5f\u5e8f\u3092\u5229\u7528\u3059\u308b\u3053\u3068\n"
            "(G) \u30d5\u30a3\u30fc\u30c9\u30d0\u30c3\u30af\u5236\u5fa1\u306b\u3088\u308a\u904e\u5ea6\u306e\u52b9\u679c\u3092\u9632\u6b62\u3059\u308b\u3053\u3068"),
        (10, "\u7279\u306b\u3001\u300c\u5149\u5316\u5b66\u7684NO\u653e\u51fa\u300d\u306b\u3088\u308bAVA\u9078\u629e\u7684\u62e1\u5f35\u3068\u3044\u3046\u6280\u8853\u601d\u60f3\u306f\u3001\u4e0a\u8a18\u3044\u305a\u308c\u306e"
             "\u5f93\u6765\u6280\u8853\u306b\u3082\u958b\u793a\u3055\u308c\u3066\u304a\u3089\u305a\u3001\u307e\u305f\u3001(i) PBM\u5206\u91ce\u306e\u5f53\u696d\u8005\u304cAVA\u5206\u5e03\u306e"
             "\u89e3\u5256\u5b66\u7684\u7279\u7570\u6027\uff08glabrous skin\u5c40\u5728\uff09\u306b\u7740\u76ee\u3057\u3066\u30c7\u30d0\u30a4\u30b9\u306e\u7167\u5c04\u90e8\u4f4d\u3092\u9650\u5b9a\u3059\u308b"
             "\u52d5\u6a5f\u4ed8\u3051\u306f\u5b58\u5728\u305b\u305a\u3001(ii) \u5faa\u74b0\u5668\u5206\u91ce\u306e\u5f53\u696d\u8005\u304cPBM\u306b\u3088\u308bNO\u653e\u51fa\u3092\u964d\u5727\u624b\u6bb5\u3068\u3057\u3066"
             "\u9069\u7528\u3059\u308b\u52d5\u6a5f\u4ed8\u3051\u3082\u5b58\u5728\u3057\u306a\u3044\u3002\u3059\u306a\u308f\u3061\u3001\u5f93\u6765\u6280\u8853\u306e\u7d44\u307f\u5408\u308f\u305b\u304b\u3089\u672c\u767a\u660e\u306b\u60f3\u5230\u3059\u308b"
             "\u3053\u3068\u306b\u306f\u3001\u6280\u8853\u5206\u91ce\u6a2a\u65ad\u7684\u306a\u77e5\u898b\u306e\u7d71\u5408\u3092\u8981\u3059\u308b\u3068\u3044\u3046\u9855\u8457\u306a\u963b\u5bb3\u8981\u56e0\u304c\u5b58\u5728\u3059\u308b\u3002"),
    ]
    for n, t in bg_paras:
        _np(doc, n, t)

    # Prior art refs
    _h(doc, "\u3010\u5148\u884c\u6280\u8853\u6587\u732e\u3011")
    _h(doc, "\u3010\u7279\u8a31\u6587\u732e\u3011", 2)
    _np(doc, 11,
        "\u7279\u8a31\u6587\u732e1: US 11,229,548 B2 (Diller et al., 2022)\n"
        "\u7279\u8a31\u6587\u732e2: US 8,679,170 B2 (Muehlbauer et al., 2014)\n"
        "\u7279\u8a31\u6587\u732e3: US 6,974,442 B2 (Grahn & Heller, 2005)\n"
        "\u7279\u8a31\u6587\u732e4: US 8,551,032 B2 (Stanford, 2013)\n"
        "\u7279\u8a31\u6587\u732e5: US 10,596,388 B2 (Enwemeka et al., 2020)")
    _h(doc, "\u3010\u975e\u7279\u8a31\u6587\u732e\u3011", 2)
    refs = [
        "Lobo MD et al. Lancet. 2015;385(9978):1634-1641.",
        "Walløe L. Temperature (Austin). 2016;3(1):92-103.",
        "Ribeiro BG et al. Lasers Med Sci. 2022;37(6):2551-2559.",
        "Kräuchi K et al. Chronobiol Int. 2012;29(9):1189-1197.",
        "Gani F et al. J Am Heart Assoc. 2021;10:e018979.",
        "Tomiyama H et al. J Hypertens. 2022;40(5):839-848.",
        "Granger DN et al. Am J Physiol Heart Circ Physiol. 2019;316(2):H375-H392.",
        "Karu TI et al. Lasers Surg Med. 2005;36(4):307-314.",
        "Chung H et al. Ann Biomed Eng. 2012;40(2):516-533.",
        "Hale GM et al. Appl Optics. 1973;12(3):555-563.",
    ]
    _np(doc, 12, "\n".join(f"\u975e\u7279\u8a31\u6587\u732e{i+1}: {r}" for i, r in enumerate(refs)))

    # Summary
    _h(doc, "\u3010\u767a\u660e\u306e\u6982\u8981\u3011")
    _h(doc, "\u3010\u767a\u660e\u304c\u89e3\u6c7a\u3057\u3088\u3046\u3068\u3059\u308b\u8ab2\u984c\u3011", 2)
    _np(doc, 13,
        "\u672c\u767a\u660e\u306f\u3001\u4ee5\u4e0b\u306e\u6280\u8853\u7684\u8ab2\u984c\u3092\u89e3\u6c7a\u3059\u308b\u3082\u306e\u3067\u3042\u308b\uff1a\n"
        "1. \u85ac\u7269\u3092\u7528\u3044\u305a\u306b\u3001glabrous skin\u306b\u5c40\u5728\u3059\u308bAVA\u306e\u8840\u7ba1\u5e73\u6ed1\u7b4b\u3092\u5149\u5316\u5b66\u7684NO\u653e\u51fa\u306b\u3088\u308a"
        "\u80fd\u52d5\u7684\u306b\u5f1b\u7de9\u3055\u305b\u3001\u672b\u68a2\u8840\u7ba1\u62b5\u6297\u3092\u4f4e\u4e0b\u3055\u305b\u308b\u3053\u3068\n"
        "2. \u9670\u5727\u6a5f\u69cb\u304a\u3088\u3073\u4f1d\u5c0e\u52a0\u71b1\u6a5f\u69cb\u3092\u6392\u9664\u3057\u3001\u5149\u7167\u5c04\u306e\u307f\u3067AVA\u62e1\u5f35\u3092\u9054\u6210\u3059\u308b"
        "\u30a6\u30a7\u30a2\u30e9\u30d6\u30eb\u5f62\u614b\u3092\u5b9f\u73fe\u3059\u308b\u3053\u3068\n"
        "3. AVA\u5468\u56f2\u306eNO\u84c4\u7a4d\u91cf\u306e\u67af\u6e07\u3092\u9632\u6b62\u3059\u308b\u30d1\u30eb\u30b9\u7167\u5c04\u5236\u5fa1\u306b\u3088\u308a\u3001\u6301\u7d9a\u7684\uff084\u6642\u9593\u4ee5\u4e0a\uff09\u306a"
        "AVA\u62e1\u5f35\u7dad\u6301\u3092\u53ef\u80fd\u306b\u3059\u308b\u3053\u3068\n"
        "4. \u672b\u68a2\u76ae\u819a\u6e29\u304a\u3088\u3073/\u307e\u305f\u306f\u8840\u5727\u306e\u30d5\u30a3\u30fc\u30c9\u30d0\u30c3\u30af\u60c5\u5831\u306b\u57fa\u3065\u304d\u3001\u5149\u7167\u5c04\u30d1\u30e9\u30e1\u30fc\u30bf\u3092"
        "\u81ea\u5f8b\u7684\u306b\u6700\u9069\u5316\u3059\u308b\u3053\u3068\n"
        "5. \u5bd2\u51b7\u74b0\u5883\u4e0b\u306b\u304a\u3051\u308bAVA\u653e\u7e2e\u3092\u4e88\u9632\u7684\u306b\u6291\u5236\u3059\u308b\u3053\u3068")
    _h(doc, "\u3010\u8ab2\u984c\u3092\u89e3\u6c7a\u3059\u308b\u305f\u3081\u306e\u624b\u6bb5\u3011", 2)
    _np(doc, 14, "\u4e0a\u8a18\u8ab2\u984c\u3092\u89e3\u6c7a\u3059\u308b\u305f\u3081\u3001\u672c\u767a\u660e\u306f\u4ee5\u4e0b\u306e\u6280\u8853\u7684\u624b\u6bb5\u3092\u63a1\u7528\u3059\u308b\u3002")
    _np(doc, 15,
        "\u7b2c\u4e00\u306b\u3001\u5149\u7167\u5c04\u624b\u6bb5\u3068\u3057\u3066630\u301c850 nm\u5e2f\u306e\u767a\u5149\u30c0\u30a4\u30aa\u30fc\u30c9\uff08LED\uff09\u30a2\u30ec\u30a4\u3092"
        "\u30d5\u30ec\u30ad\u30b7\u30d6\u30eb\u57fa\u677f\u4e0a\u306b\u914d\u7f6e\u3057\u3001\u624b\u638c\u304a\u3088\u3073/\u307e\u305f\u306f\u8db3\u5e95\u306eglabrous skin\u306b\u5bc6\u7740\u3055\u305b\u3066"
        "\u88c5\u7740\u3059\u308b\u3002\u3053\u306e\u6ce2\u9577\u5e2f\u306f\u3001(i) \u30d8\u30e2\u30b0\u30ed\u30d3\u30f3\u7d50\u5408NO\u306e\u5149\u89e3\u96e2\u5438\u53ce\u5e2f\u3001"
        "(ii) \u30b7\u30c8\u30af\u30ed\u30e0c\u30aa\u30ad\u30b7\u30c0\u30fc\u30bc\uff08CcO\uff09\u306eCu_A/Cu_B\u4e2d\u5fc3\u306e\u5438\u53ce\u5e2f\u3001"
        "\u304a\u3088\u3073 (iii) S-\u30cb\u30c8\u30ed\u30bd\u30c1\u30aa\u30fc\u30eb\u306e\u5149\u5206\u89e3\u6ce2\u9577\u5e2f\u306b\u5bfe\u5fdc\u3059\u308b\uff08\u975e\u7279\u8a31\u6587\u732e8, 9\u53c2\u7167\uff09\u3002")
    _np(doc, 16,
        "\u7b2c\u4e8c\u306b\u3001AVA\u306f\u771f\u76ae\u6df1\u5c64\u304b\u3089\u76ae\u4e0b\u6d45\u5c64\uff08\u6df1\u30551\u301c3 mm\uff09\u306b\u4f4d\u7f6e\u3057\uff08\u975e\u7279\u8a31\u6587\u732e2\uff09\u3001"
        "660 nm\u5149\u306e\u76ae\u819a\u900f\u904e\u6df1\u5ea6\uff0814\u301c21 mm\uff09\u306f\u5341\u5206\u306b\u3053\u306e\u6df1\u5ea6\u306b\u5230\u9054\u3059\u308b\u3002"
        "\u3057\u305f\u304c\u3063\u3066\u3001\u4f53\u8868\u304b\u3089\u306e\u5149\u7167\u5c04\u306e\u307f\u3067AVA\u58c1\u306e\u8840\u7ba1\u5185\u76ae\u7d30\u80de\u306b\u76f4\u63a5\u4f5c\u7528\u3057\u5f97\u308b\u3002")
    _np(doc, 17,
        "\u7b2c\u4e09\u306b\u3001NO\u653e\u51fa\u5f8c\u306eNO\u84c4\u7a4d\u91cf\u56de\u5fa9\u306e\u305f\u3081\u3001\u9593\u6b47\u7684\u30d1\u30eb\u30b9\u7167\u5c04\u5236\u5fa1\u3092\u63a1\u7528\u3059\u308b\u3002"
        "\u9023\u7d9a\u7167\u5c04\u3067\u306fS-\u30cb\u30c8\u30ed\u30bd\u30c1\u30aa\u30fc\u30eb\u7b49\u306eNO\u30b9\u30c8\u30a2\u304c\u67af\u6e07\u3057\u3001\u52b9\u679c\u304c\u6e1b\u5f31\u3059\u308b"
        "\uff08\u975e\u7279\u8a31\u6587\u732e3\uff09\u3002\u672c\u767a\u660e\u3067\u306f\u3001\u7167\u5c04\u671f\u9593\uff08T_on\uff09\u3068\u975e\u7167\u5c04\u671f\u9593\uff08T_off\uff09\u306e"
        "\u30c7\u30e5\u30fc\u30c6\u30a3\u30b5\u30a4\u30af\u30eb\u3092\u52d5\u7684\u306b\u5236\u5fa1\u3057\u3001T_off\u671f\u9593\u4e2d\u306eNO\u30b9\u30c8\u30a2\u56de\u5fa9\u3092\u78ba\u4fdd\u3059\u308b\u3002")
    _np(doc, 18,
        "\u7b2c\u56db\u306b\u3001\u672b\u68a2\u76ae\u819a\u6e29\u30bb\u30f3\u30b5\u304a\u3088\u3073/\u307e\u305f\u306f\u8840\u5727\u63a8\u5b9a\u624b\u6bb5\u304b\u3089\u306e\u30d5\u30a3\u30fc\u30c9\u30d0\u30c3\u30af\u4fe1\u53f7\u306b\u57fa\u3065\u304d\u3001"
        "\u5149\u7167\u5c04\u306e\u30d1\u30ef\u30fc\u5bc6\u5ea6\u3001\u30d1\u30eb\u30b9\u5468\u6ce2\u6570\u3001\u30c7\u30e5\u30fc\u30c6\u30a3\u30b5\u30a4\u30af\u30eb\u3092\u9589\u30eb\u30fc\u30d7\u5236\u5fa1\u3059\u308b\u3002"
        "\u3053\u308c\u306b\u3088\u308a\u3001\u500b\u4f53\u5dee\uff08AVA\u5bc6\u5ea6\u3001NO\u84c4\u7a4d\u91cf\u3001\u76ae\u819a\u8272\u7d20\u91cf\uff09\u306b\u5bfe\u3059\u308b\u81ea\u52d5\u9069\u5fdc\u304c\u5b9f\u73fe\u3055\u308c\u308b\u3002")

    # Detailed description
    _h(doc, "\u3010\u767a\u660e\u3092\u5b9f\u65bd\u3059\u308b\u305f\u3081\u306e\u5f62\u614b\u3011")
    _np(doc, 19, "\u672c\u660e\u7d30\u66f8\u306b\u304a\u3044\u3066\u4f7f\u7528\u3059\u308b\u7528\u8a9e\u3092\u4ee5\u4e0b\u306e\u3088\u3046\u306b\u5b9a\u7fa9\u3059\u308b\u3002")
    _np(doc, 20,
        "\u300c\u52d5\u9759\u8108\u543e\u5408\uff08AVA\uff09\u300d\u3068\u306f\u3001\u7d30\u52d5\u8108\u3068\u7d30\u9759\u8108\u3092\u6bdb\u7d30\u8840\u7ba1\u3092\u4ecb\u3055\u305a\u306b\u76f4\u63a5\u63a5\u7d9a\u3059\u308b"
        "\u8840\u7ba1\u69cb\u9020\u3067\u3042\u308a\u3001\u53e3\u5f8420\u301c150 \u00b5m\u306e\u30b0\u30ed\u30e0\u30b9\u4f53\uff08glomus body\uff09\u3092\u542b\u3080\u3002"
        "\u30d2\u30c8\u306b\u304a\u3044\u3066AVA\u306f\u4e3b\u3068\u3057\u3066glabrous skin\uff08\u624b\u638c\u3001\u8db3\u5e95\u3001\u6307\u8dbe\u5148\u7aef\uff09\u306b\u5c40\u5728\u3057\u3001"
        "\u63a8\u5b9a\u5bc6\u5ea6100\u301c600\u500b/cm\u00b2\u3067\u3042\u308b\uff08\u975e\u7279\u8a31\u6587\u732e2\uff09\u3002")
    _np(doc, 21,
        "\u300c\u5149\u751f\u4f53\u8abf\u7bc0\uff08PBM\uff09\u300d\u3068\u306f\u3001\u975e\u30a4\u30aa\u30f3\u5316\u5149\uff08600\u301c1100 nm\u5e2f\uff09\u3092\u751f\u4f53\u7d44\u7e54\u306b\u7167\u5c04\u3057\u3001"
        "\u5149\u5316\u5b66\u53cd\u5fdc\u3092\u8a98\u5c0e\u3057\u3066\u751f\u7406\u5b66\u7684\u5fdc\u7b54\u3092\u60f9\u8d77\u3059\u308b\u6280\u8853\u3092\u3044\u3046\u3002"
        "\u672c\u767a\u660e\u306b\u304a\u3044\u3066PBM\u306f\u300c\u975e\u71b1\u7684\u300d\u4f5c\u7528\u3092\u6307\u3057\u3001\u7167\u5c04\u306b\u3088\u308b\u7d44\u7e54\u6e29\u5ea6\u4e0a\u6607\u304c"
        "1\u00b0C\u672a\u6e80\u306e\u6761\u4ef6\u4e0b\u3067\u306e\u5149\u5316\u5b66\u7684\u52b9\u679c\u3092\u3044\u3046\u3002")
    _np(doc, 22,
        "\u300c\u9078\u629e\u7684\u306bAVA\u3092\u62e1\u5f35\u3055\u305b\u308b\u300d\u3068\u306f\u3001AVA\u58c1\u306e\u8840\u7ba1\u5e73\u6ed1\u7b4b\u3092\u5f1b\u7de9\u3055\u305b\u308b\u3053\u3068\u306b\u3088\u308a\u3001"
        "glabrous skin\u5185\u306eAVA\u3092\u4ecb\u3057\u305f\u8840\u6d41\u3092\u512a\u5148\u7684\u306b\u5897\u52a0\u3055\u305b\u308b\u3053\u3068\u3092\u3044\u3046\u3002"
        "\u3053\u3053\u3067\u300c\u9078\u629e\u7684\u300d\u3068\u306f\u3001\u7167\u5c04\u90e8\u4f4d\u3092\u89e3\u5256\u5b66\u7684\u306bAVA\u9ad8\u5bc6\u5ea6\u9818\u57df\u306b\u9650\u5b9a\u3059\u308b\u3053\u3068\u306b\u3088\u308b"
        "\u89e3\u5256\u5b66\u7684\u9078\u629e\u6027\u3001\u304a\u3088\u3073NO\u653e\u51fa\u2192cGMP\u2192\u5e73\u6ed1\u7b4b\u5f1b\u7de9\u3068\u3044\u3046AVA\u58c1\u306b\u76f4\u63a5\u4f5c\u7528\u3059\u308b"
        "\u85ac\u7406\u5b66\u7684\u9078\u629e\u6027\u306e\u4e21\u65b9\u3092\u5305\u542b\u3059\u308b\u3002\u306a\u304a\u3001\u4e00\u822c\u6bdb\u7d30\u8840\u7ba1\u3078\u306e\u8efd\u5ea6\u306e\u8840\u6d41\u5897\u52a0\u306f\u6392\u9664\u3055\u308c\u306a\u3044\u3002")
    _np(doc, 23, "\u300c\u30d1\u30ef\u30fc\u5bc6\u5ea6\u300d\u3068\u306f\u3001\u7167\u5c04\u9762\u306b\u304a\u3051\u308b\u5358\u4f4d\u9762\u7a4d\u3042\u305f\u308a\u306e\u5149\u51fa\u529b\uff08mW/cm\u00b2\uff09\u3092\u3044\u3046\u3002")
    _np(doc, 24, "\u300c\u30c7\u30e5\u30fc\u30c6\u30a3\u30b5\u30a4\u30af\u30eb\u300d\u3068\u306f\u3001\u30d1\u30eb\u30b9\u7167\u5c04\u306b\u304a\u3051\u308b\u7167\u5c04\u671f\u9593\uff08T_on\uff09\u306e\u3001"
        "\u7167\u5c041\u5468\u671f\uff08T_on + T_off\uff09\u306b\u5bfe\u3059\u308b\u6bd4\u7387\uff08%\uff09\u3092\u3044\u3046\u3002")

    # Mechanism
    _np(doc, 25,
        "\u672c\u767a\u660e\u306e\u4f5c\u7528\u539f\u7406\u3092\u4ee5\u4e0b\u306b\u8a73\u8ff0\u3059\u308b\u3002\u672c\u9805\u306e\u8a18\u8f09\u306f\u767a\u660e\u306e\u6280\u8853\u7684\u7406\u89e3\u3092\u52a9\u3051\u308b\u305f\u3081\u306e"
        "\u3082\u306e\u3067\u3042\u308a\u3001\u672c\u767a\u660e\u306e\u6a29\u5229\u7bc4\u56f2\u306f\u7279\u8a31\u8acb\u6c42\u306e\u7bc4\u56f2\u306e\u8a18\u8f09\u306b\u57fa\u3065\u3044\u3066\u5b9a\u3081\u3089\u308c\u308b\u3002")
    _np(doc, 26,
        "630\u301c850 nm\u5e2f\u306e\u5149\u306f\u3001\u8840\u7ba1\u5185\u76ae\u7d30\u80de\u5185\u306e\u4ee5\u4e0b\u306eNO\u84c4\u7a4d\u6e90\u306b\u4f5c\u7528\u3057\u3066NO\u3092\u904a\u96e2\u3055\u305b\u308b\uff1a\n"
        "(i) S-\u30cb\u30c8\u30ed\u30bd\u30c1\u30aa\u30fc\u30eb\uff08RSNO\uff09: \u5185\u76ae\u7d30\u80de\u5185\u306e\u30b0\u30eb\u30bf\u30c1\u30aa\u30f3-SNO\u3001\u30a2\u30eb\u30d6\u30df\u30f3-SNO\u7b49\u306e"
        "\u786b\u9ec4-NO\u7d50\u5408\u3092\u5149\u5206\u89e3\u3059\u308b\uff08\u975e\u7279\u8a31\u6587\u732e8\uff09\u3002\n"
        "(ii) \u30b7\u30c8\u30af\u30ed\u30e0c\u30aa\u30ad\u30b7\u30c0\u30fc\u30bc\uff08CcO\uff09\u7d50\u5408NO: \u30df\u30c8\u30b3\u30f3\u30c9\u30ea\u30a2\u96fb\u5b50\u4f1d\u9054\u7cfb"
        "\u8907\u5408\u4f53IV\u306eCu_B\u4e2d\u5fc3\u306b\u53ef\u9006\u7684\u306b\u7d50\u5408\u3057\u305fNO\u304c\u5149\u5438\u53ce\u306b\u3088\u308a\u89e3\u96e2\u3059\u308b\uff08\u975e\u7279\u8a31\u6587\u732e9\uff09\u3002\n"
        "(iii) \u30d8\u30e2\u30b0\u30ed\u30d3\u30f3\u7d50\u5408NO\uff08Hb-NO\uff09: \u8d64\u8840\u7403\u5185\u30d8\u30e2\u30b0\u30ed\u30d3\u30f3\u306e\u30d8\u30e0\u9244\u306b\u7d50\u5408\u3057\u305f"
        "NO\u304cR/T\u69cb\u9020\u5909\u5316\u306b\u4f34\u3044\u653e\u51fa\u3055\u308c\u308b\u3002")
    _np(doc, 27,
        "\u904a\u96e2\u3057\u305fNO\u306f\u4ee5\u4e0b\u306e\u7d4c\u8def\u3067AVA\u8840\u7ba1\u5e73\u6ed1\u7b4b\u3092\u5f1b\u7de9\u3055\u305b\u308b\uff1a\n"
        "NO \u2192 \u53ef\u6eb6\u6027\u30b0\u30a2\u30cb\u30eb\u9178\u30b7\u30af\u30e9\u30fc\u30bc\uff08sGC\uff09\u6d3b\u6027\u5316 \u2192 GTP\u2192cGMP\u5909\u63db\u4fc3\u9032 \u2192 "
        "cGMP\u4f9d\u5b58\u6027\u30d7\u30ed\u30c6\u30a4\u30f3\u30ad\u30ca\u30fc\u30bc\uff08PKG\uff09\u6d3b\u6027\u5316 \u2192 "
        "\u30df\u30aa\u30b7\u30f3\u8efd\u9396\u30db\u30b9\u30d5\u30a1\u30bf\u30fc\u30bc\u6d3b\u6027\u5316 + Ca\u00b2\u207a\u30c1\u30e3\u30cd\u30eb\u963b\u5bb3 \u2192 "
        "\u8840\u7ba1\u5e73\u6ed1\u7b4b\u5f1b\u7de9 \u2192 AVA\u62e1\u5f35\n"
        "\u3053\u306e\u7d4c\u8def\u306f\u85ac\u7406\u5b66\u7684\u306b\u5341\u5206\u306b\u78ba\u7acb\u3055\u308c\u3066\u3044\u308b\uff08\u975e\u7279\u8a31\u6587\u732e7\uff09\u3002")
    _np(doc, 28,
        "AVA\u306f\u53e3\u5f8420\u301c150 \u00b5m\u3067\u3042\u308a\u3001\u4e00\u822c\u6bdb\u7d30\u8840\u7ba1\uff085\u301c8 \u00b5m\uff09\u306e3\u301c30\u500d\u306e\u592a\u3055\u3092\u6709\u3059\u308b\u3002"
        "Poiseuille's law\u306b\u3088\u308a\u3001AVA 1\u672c\u306e\u958b\u901a\u306f\u6841\u9055\u3044\u306e\u8840\u6d41\u5897\u52a0\u3092\u3082\u305f\u3089\u3059\u3002"
        "\u624b\u638c\u306eAVA\u7dcf\u8840\u6d41\u5bb9\u91cf\u306f0\u301c1 L/min\u3067\u3042\u308a\u3001\u5fc3\u62cd\u51fa\u91cf\u306e\u6700\u592720%\u306b\u76f8\u5f53\u3059\u308b\u3002"
        "AVA\u8840\u6d41\u304c0.5 L/min\u5897\u52a0\u3057\u305f\u5834\u5408\u3001TPR\u306f\u7406\u8ad6\u4e0a10%\u4f4e\u4e0b\u3057\u3001"
        "MAP 100 mmHg\u306e\u60a3\u8005\u30675\u301c10 mmHg\u306e\u964d\u5727\u306b\u76f8\u5f53\u3059\u308b\u3002")
    _np(doc, 29,
        "\u975e\u7279\u8a31\u6587\u732e3\u306b\u304a\u3044\u3066\u3001NOS\u963b\u5bb3\u5264\u6295\u4e0e\u30e9\u30c3\u30c8\u3067\u306fPBM\u306b\u3088\u308b\u964d\u5727\u304c\u6d88\u5931\u3057\u305f\u3002"
        "\u672c\u767a\u660e\u306e\u30d1\u30eb\u30b9\u7167\u5c04\u5236\u5fa1\u306f\u3001\u3053\u306e\u77e5\u898b\u306b\u57fa\u3065\u304f\u6280\u8853\u7684\u89e3\u6c7a\u7b56\u3067\u3042\u308b\u3002"
        "T_off\u671f\u9593\u4e2d\u306beNOS\u304cL-\u30a2\u30eb\u30ae\u30cb\u30f3\u304b\u3089NO\u3092de novo\u5408\u6210\u3057\u3001"
        "S-\u30cb\u30c8\u30ed\u30bd\u30c1\u30aa\u30fc\u30eb\u7b49\u306e\u30b9\u30c8\u30a2\u304c\u518d\u5145\u586b\u3055\u308c\u308b\u3002"
        "\u672c\u767a\u660e\u30673\u301c30\u79d2\u306eT_on\u30681\u301c15\u79d2\u306eT_off\u3092\u7d44\u307f\u5408\u308f\u305b\u3001"
        "4\u6642\u9593\u4ee5\u4e0a\u306e\u6301\u7d9a\u7684AVA\u62e1\u5f35\u52b9\u679c\u3092\u9054\u6210\u3059\u308b\u3002")

    # Prior art comparison
    _np(doc, 30, "\u5f93\u6765\u6280\u8853\u3068\u306e\u5dee\u7570\u3092\u4ee5\u4e0b\u306e\u8868\u306b\u793a\u3059\uff08\u56f35\u53c2\u7167\uff09\u3002")
    _tbl(doc,
        ["\u8981\u7d20", "AVACEN\n(\u7279\u8a31\u6587\u732e2)", "Diller\n(\u7279\u8a31\u6587\u732e1)", "\u672c\u767a\u660e"],
        [
            ["AVA\u62e1\u5f35\u6a5f\u5e8f", "\u9670\u5727\uff0b\u4f1d\u5c0e\u52a0\u71b1", "\u4e2d\u67a2\u53cd\u5c04", "\u5149\u5316\u5b66\u7684NO\u653e\u51fa"],
            ["\u7167\u5c04\u30a8\u30cd\u30eb\u30ae\u30fc", "\u6a5f\u68b0\uff0b\u71b1", "\u71b1", "\u5149\uff08\u975e\u71b1\u7684\uff09"],
            ["\u9670\u5727\u6a5f\u69cb", "\u5fc5\u9808 (\u221230 mmHg)", "\u306a\u3057", "\u4e0d\u8981"],
            ["\u5f62\u614b", "\u636e\u3048\u7f6e\u304d (5 kg)", "\u30d9\u30c3\u30c9\u578b", "\u30a6\u30a7\u30a2\u30e9\u30d6\u30eb (<80 g)"],
            ["\u4f7f\u7528\u5834\u9762", "\u9759\u6b62\u72b6\u614b\u306e\u307f", "\u7761\u7720\u4e2d\u306e\u307f", "\u65e5\u5e38\u52d5\u4f5c\u4e2d"],
            ["\u6d88\u8cbb\u96fb\u529b", ">50 W", "\u4e0d\u660e", "<2 W"],
            ["\u30d1\u30eb\u30b9\u5236\u5fa1", "\u306a\u3057", "\u306a\u3057", "\u3042\u308a (NO\u30b9\u30c8\u30a2\u7ba1\u7406)"],
        ])

    # Inhibiting factors
    _np(doc, 31,
        "\u65e2\u5b58PBM\u30c7\u30d0\u30a4\u30b9\u304b\u3089\u672c\u767a\u660e\u306b\u60f3\u5230\u3059\u308b\u306b\u306f\u3001\u4ee5\u4e0b\u306e\u8907\u5408\u7684\u6280\u8853\u7684\u969c\u58c1\u304c\u5b58\u5728\u3059\u308b\uff1a\n"
        "(i) \u7167\u5c04\u90e8\u4f4d\u306e\u9650\u5b9a: \u65e2\u5b58PBM\u306f\u5168\u8eab\u4efb\u610f\u90e8\u4f4d\u3092\u5bfe\u8c61\u3068\u3057\u3001glabrous skin\u306b\u7167\u5c04\u90e8\u4f4d\u3092"
        "\u9650\u5b9a\u3059\u308b\u52d5\u6a5f\u4ed8\u3051\u304c\u306a\u3044\u3002\n"
        "(ii) \u52b9\u679c\u306e\u4e88\u6e2c\u4e0d\u53ef\u80fd\u6027: \u300c\u5c40\u6240\u7167\u5c04\u2192\u5168\u8eab\u8840\u884c\u52d5\u614b\u5909\u5316\u300d\u306f\u4e88\u6e2c\u56f0\u96e3\u3002\n"
        "(iii) \u6280\u8853\u5206\u91ce\u306e\u9694\u7d76: PBM\u5149\u7269\u7406\u5b66\u3001\u8840\u7ba1\u89e3\u5256\u5b66\u3001\u5faa\u74b0\u751f\u7406\u5b66\u3001\u5236\u5fa1\u5de5\u5b66\u306e4\u5206\u91ce\u306e\u7d71\u5408\u304c\u5fc5\u8981\u3002\n"
        "(iv) \u30d1\u30eb\u30b9\u5236\u5fa1\u306e\u975e\u81ea\u660e\u6027: NO\u30b9\u30c8\u30a2\u7ba1\u7406\u306e\u305f\u3081\u306e\u52d5\u7684\u30c7\u30e5\u30fc\u30c6\u30a3\u30b5\u30a4\u30af\u30eb\u5236\u5fa1\u306f"
        "2022\u5e74\u306b\u521d\u3081\u3066\u793a\u5506\u3055\u308c\u305f\u3002")
    _np(doc, 32,
        "\u672c\u767a\u660e\u306e\u9855\u8457\u306a\u52b9\u679c\uff1a\n"
        "(i) \u9670\u5727\u6a5f\u69cb\u306e\u5b8c\u5168\u6392\u9664\u306b\u3088\u308b\u300c\u30a6\u30a7\u30a2\u30e9\u30d6\u30ebAVA\u30c7\u30d0\u30a4\u30b9\u300d\u3068\u3044\u3046\u65b0\u30ab\u30c6\u30b4\u30ea\u306e\u5275\u51fa\n"
        "(ii) \u30d1\u30eb\u30b9\u7167\u5c04\u5236\u5fa1\u306b\u3088\u308b4\u6642\u9593\u4ee5\u4e0a\u306e\u6301\u7d9a\u7684\u52b9\u679c\u7dad\u6301\n"
        "(iii) \u76ae\u819a\u6e29\u30d5\u30a3\u30fc\u30c9\u30d0\u30c3\u30af\u306b\u3088\u308b\u5bd2\u51b7\u6607\u5727\u306e\u4e88\u9632\u7684\u6291\u5236")

    # Embodiments
    _np(doc, 33,
        "\u3014\u5b9f\u65bd\u4f8b1: \u30b0\u30ed\u30fc\u30d6\u578b\u30c7\u30d0\u30a4\u30b9\uff08\u56f32\u53c2\u7167\uff09\u3015\n"
        "\u30d5\u30ec\u30ad\u30b7\u30d6\u30ebPCB\u57fa\u677f\uff08\u30dd\u30ea\u30a4\u30df\u30c9\u30014\u5c64\u30010.2 mm\u539a\uff09\u4e0a\u306b660 nm LED\u7d20\u5b50"
        "\uff08\u5149\u51fa\u529b50 mW/\u7d20\u5b50@350 mA\uff09\u30924\u00d78\u30a2\u30ec\u30a4\uff08\u8a0832\u500b\uff09\u3067\u914d\u7f6e\u3057\u3001"
        "\u624b\u638c\u9762\uff08\u7d0440 cm\u00b2\uff09\u3092\u5747\u4e00\u306b\u7167\u5c04\u3059\u308b\u30b0\u30ed\u30fc\u30d6\u578b\u30c7\u30d0\u30a4\u30b9\u3092\u69cb\u6210\u3059\u308b\u3002"
        "LED\u7d20\u5b50\u306f\u30b7\u30ea\u30b3\u30fc\u30f3\u30b4\u30e0\u88fd\u30b0\u30ed\u30fc\u30d6\u306e\u5185\u9762\u306b\u57cb\u8a2d\u3055\u308c\u308b\u3002"
        "\u30ea\u30b9\u30c8\u30d0\u30f3\u30c9\u90e8\u306b\u5236\u5fa1\u57fa\u677f\uff08ESP32-S3\u3001TLC5940\u00d72\u3001MLX90614\u3001MAX30102\u3001BLE 5.0\uff09"
        "\u304a\u3088\u3073LiPo\u30d0\u30c3\u30c6\u30ea\u30fc\uff083.7V, 800 mAh\uff09\u3092\u5185\u8535\u3059\u308b\u3002")
    _np(doc, 34,
        "\u5236\u5fa1\u30a2\u30eb\u30b4\u30ea\u30ba\u30e0\uff08\u56f33\u53c2\u7167\uff09\uff1a\n"
        "\u30fb\u521d\u671f\u7167\u5c04: \u30d1\u30ef\u30fc\u5bc6\u5ea640 mW/cm\u00b2\u3001\u30c7\u30e5\u30fc\u30c6\u30a3\u30b5\u30a4\u30af\u30eb67% (T_on=10\u79d2, T_off=5\u79d2)\n"
        "\u30fb\u6e29\u5ea6FB: \u76ae\u819a\u6e29<30\u00b0C\u219260 mW/cm\u00b2\u3078\u5897\u5f37; >38\u00b0C\u219220 mW/cm\u00b2\u3078\u6e1b\u5f31; >42\u00b0C\u2192\u5373\u6642\u505c\u6b62\n"
        "\u30fbBP FB (PPG\u6709\u52b9\u6642): BP>\u76ee\u6a19+5\u2192\u30c7\u30e5\u30fc\u30c6\u30a3\u5897\u52a0; BP<\u76ee\u6a19-5\u2192\u6e1b\u5c11; BP\u2264\u76ee\u6a19-15\u2192\u505c\u6b62")
    _np(doc, 35,
        "\u3014\u5b9f\u65bd\u4f8b2: \u30a4\u30f3\u30bd\u30fc\u30eb\u578b\u30c7\u30d0\u30a4\u30b9\u3015\n"
        "\u9774\u306e\u30a4\u30f3\u30bd\u30fc\u30eb\u5f62\u72b6\u306e\u57fa\u677f\u4e0a\u306b830 nm LED\u7d20\u5b50\u3092\u8db3\u5e95\u30a2\u30fc\u30c1\u90e8\u306b\u96c6\u4e2d\u914d\u7f6e"
        "\uff083\u00d76\u30a2\u30ec\u30a4\u3001\u8a0818\u500b\uff09\u3002\u5727\u96fb\u30bb\u30f3\u30b5\uff08FSR\uff09\u3092\u914d\u7f6e\u3057\u3001"
        "\u5ea7\u4f4d\u5b89\u9759\u6642\u306b\u306e\u307f\u7167\u5c04\u3092\u884c\u3046\u30e2\u30fc\u30c9\u3092\u5099\u3048\u308b\u3002")
    _np(doc, 36,
        "\u3014\u5b9f\u65bd\u4f8b3: \u6c11\u751f\u54c1\u7248\uff08\u672b\u68a2\u4fdd\u6e29\u30b0\u30ed\u30fc\u30d6\uff09\u3015\n"
        "\u8acb\u6c42\u984210\u301c12\u306b\u5bfe\u5fdc\u3002\u533b\u7642\u6a5f\u5668\u3068\u3057\u3066\u306e\u85ac\u4e8b\u627f\u8a8d\u3092\u8981\u3057\u306a\u3044\u4e00\u822c\u6d88\u8cbb\u8005\u5411\u3051\u88fd\u54c1\u3002\n"
        "\u30cf\u30fc\u30c9\u30a6\u30a7\u30a2\u306f\u5b9f\u65bd\u4f8b1\u3068\u540c\u4e00\u3060\u304c\u3001MAX30102 PPG\u30bb\u30f3\u30b5\u306f\u30d5\u30a1\u30fc\u30e0\u30a6\u30a7\u30a2\u3067\u7121\u52b9\u5316\u3002\n"
        "\u88fd\u54c1\u8a34\u6c42: \u300c\u5149\u306b\u3088\u308b\u672b\u68a2\u8840\u6d41\u4fc3\u9032\u3067\u624b\u5148\u307d\u304b\u307d\u304b\u300d\u3002\u60f3\u5b9a\u4fa1\u683c: \u00a59,800\u301c14,800\u3002\n"
        "\u6bb5\u968e\u7684\u6a5f\u80fd\u6607\u683c\uff08\u56f36\uff09: Phase 1\uff08\u6c11\u751f\u54c1\uff09\u2192Phase 2\uff08OTA\u30a2\u30c3\u30d7\u30c7\u30fc\u30c8\u3067PPG\u6709\u52b9\u5316\u3001"
        "\u30af\u30e9\u30b9II\u533b\u7642\u6a5f\u5668\u3068\u3057\u3066\u8a8d\u8a3c\uff09\u3002\n"
        "\u91cf\u7523\u52b9\u679c: \u5e74\u959310,000\u301c50,000\u53f0\u8ca9\u58f2\u2192BOM \u00a52,000\u4ee5\u4e0b\u3002"
        "\u533b\u7642\u6a5f\u5668\u7248\u3082\u00a530,000\u4ee5\u4e0b\u3067\u63d0\u4f9b\u53ef\u80fd\u3002")

    # Effects
    _h(doc, "\u3010\u767a\u660e\u306e\u52b9\u679c\u3011")
    _np(doc, 37,
        "\u672c\u767a\u660e\u306b\u3088\u308c\u3070\u3001\u4ee5\u4e0b\u306e\u52b9\u679c\u304c\u5f97\u3089\u308c\u308b\uff1a\n"
        "1. \u5149\u5316\u5b66\u7684AVA\u62e1\u5f35: \u5f93\u6765\u30c7\u30d0\u30a4\u30b9\u306b\u5b58\u5728\u3057\u306a\u3044\u4f5c\u7528\u6a5f\u5e8f\u3002\n"
        "2. \u30a6\u30a7\u30a2\u30e9\u30d6\u30eb\u5316: \u30c7\u30d0\u30a4\u30b9\u7dcf\u91cd\u91cf<80 g\u3002\n"
        "3. \u4f4e\u6d88\u8cbb\u96fb\u529b: \u7dcf\u6d88\u8cbb\u96fb\u529b2.1W\u4ee5\u4e0b\u3001\u30d0\u30c3\u30c6\u30ea\u30fc\u30673\u301c6\u6642\u9593\u3002\n"
        "4. \u6301\u7d9a\u7684\u52b9\u679c: NO\u30b9\u30c8\u30a2\u7ba1\u7406\u30674\u6642\u9593\u4ee5\u4e0a\u7dad\u6301\u3002\n"
        "5. \u5bd2\u51b7\u6607\u5727\u4e88\u9632: \u76ae\u819a\u6e29\u4f4e\u4e0b\u691c\u77e5\u306b\u3088\u308b\u5148\u884c\u7684\u51fa\u529b\u5897\u5f37\u3002\n"
        "6. \u9069\u5fdc\u62e1\u5927: \u30ec\u30a4\u30ce\u30fc\u73fe\u8c61\u3001\u672b\u68a2\u4fdd\u6e29\u3001\u5bd2\u51b7\u8a98\u767a\u6027\u9ad8\u8840\u5727\u3002\n"
        "7. \u6c11\u751f\u54c1\u2192\u533b\u7642\u6a5f\u5668\u306e\u6bb5\u968e\u7684\u5c55\u958b\u3002")

    # Drawing descriptions
    _h(doc, "\u3010\u56f3\u9762\u306e\u7c21\u5358\u306a\u8aac\u660e\u3011")
    _np(doc, 38,
        "\u3010\u56f31\u3011\u672c\u767a\u660e\u306e\u4f5c\u7528\u6a5f\u5e8f\u3092\u793a\u3059\u6a21\u5f0f\u56f3\n"
        "\u3010\u56f32\u3011\u30b0\u30ed\u30fc\u30d6\u578b\u30c7\u30d0\u30a4\u30b9\u306e\u69cb\u6210\u3092\u793a\u3059\u5e73\u9762\u56f3\u304a\u3088\u3073\u65ad\u9762\u56f3\n"
        "\u3010\u56f33\u3011\u9589\u30eb\u30fc\u30d7\u5236\u5fa1\u30a2\u30eb\u30b4\u30ea\u30ba\u30e0\u306e\u30d5\u30ed\u30fc\u30c1\u30e3\u30fc\u30c8\n"
        "\u3010\u56f34\u3011\u30d1\u30eb\u30b9\u7167\u5c04\u5236\u5fa1\u306e\u30bf\u30a4\u30df\u30f3\u30b0\u30c1\u30e3\u30fc\u30c8\n"
        "\u3010\u56f35\u3011\u5f93\u6765\u6280\u8853\u3068\u306e\u6bd4\u8f03\u8868\n"
        "\u3010\u56f36\u3011\u6bb5\u968e\u7684\u6a5f\u80fd\u6607\u683c\u306e\u6982\u5ff5\u56f3")

    # Symbol description
    _h(doc, "\u3010\u7b26\u53f7\u306e\u8aac\u660e\u3011")
    _np(doc, 39,
        "1: \u30a6\u30a7\u30a2\u30e9\u30d6\u30eb\u672b\u68a2\u8840\u6d41\u4fc3\u9032\u30c7\u30d0\u30a4\u30b9\n"
        "10: \u30d5\u30ec\u30ad\u30b7\u30d6\u30eb\u57fa\u677f\uff08FPC\uff09  11: LED\u7d20\u5b50  12: LED\u30a2\u30ec\u30a4\n"
        "20: \u5236\u5fa1\u57fa\u677f  21: \u30de\u30a4\u30af\u30ed\u30b3\u30f3\u30c8\u30ed\u30fc\u30e9 (ESP32-S3)\n"
        "22: LED\u30c9\u30e9\u30a4\u30d0 (TLC5940)  23: BLE\u30a2\u30f3\u30c6\u30ca\n"
        "30: \u6e29\u5ea6\u30bb\u30f3\u30b5 (MLX90614)  31: PPG\u30bb\u30f3\u30b5 (MAX30102)\n"
        "40: \u30d0\u30c3\u30c6\u30ea\u30fc (LiPo)\n"
        "50: \u30b7\u30ea\u30b3\u30fc\u30f3\u30b4\u30e0\u88fd\u30b0\u30ed\u30fc\u30d6  51: \u30ea\u30b9\u30c8\u30d0\u30f3\u30c9\u90e8\n"
        "60: \u5916\u90e8\u7aef\u672b  61: BLE\u7121\u7dda\u901a\u4fe1\n"
        "A: \u624b\u638c\u306eglabrous skin  B: \u52d5\u9759\u8108\u543e\u5408 (AVA)")

    doc.save(str(OUT / "meisaisho.docx"))
    print(f"  Saved: {OUT / 'meisaisho.docx'}")


# ==========================================================
# 2. CLAIMS
# ==========================================================
def build_claims():
    doc = _doc()
    _h(doc, "\u3010\u66f8\u985e\u540d\u3011\u7279\u8a31\u8acb\u6c42\u306e\u7bc4\u56f2")

    claims_data = {
        1: ("手掌および/または足底の無毛皮膚（glabrous skin）の表面に装着可能な"
            "ウェアラブルデバイスであって、\n"
            "(a) フレキシブル基板上に配列された複数の発光ダイオード（LED）素子を含み、"
            "630〜850 nmの波長帯の非コヒーレント光を前記無毛皮膚の表面に向けて出射する"
            "光照射手段であって、前記無毛皮膚の真皮深層から皮下浅層に存在する"
            "動静脈吻合（AVA）に到達可能なパワー密度で光を照射するもの、\n"
            "(b) 前記光照射手段を間欠的に駆動するパルス照射モードを有し、照射期間（T_on）"
            "と非照射期間（T_off）のデューティサイクルを制御する制御手段であって、"
            "前記非照射期間において前記AVA周囲の血管内皮細胞における一酸化窒素（NO）"
            "蓄積量の回復を許容するように前記デューティサイクルを設定するもの、および\n"
            "(c) 装着者の末梢皮膚温を検知する温度センサ\n"
            "を備え、\n"
            "前記制御手段は、前記温度センサにより検知された皮膚温に基づき前記光照射手段の"
            "出力パラメータをフィードバック制御し、\n"
            "前記光照射手段から照射された光により、前記AVA周囲の血管内皮細胞において"
            "NOが光化学的に遊離し、遊離したNOが前記AVAの血管平滑筋を弛緩させることにより、"
            "前記AVAを拡張させて末梢血流を増加させる、\n"
            "ウェアラブル末梢血流促進デバイス。"),
        2: ("請求項1に記載のデバイスにおいて、前記光照射手段は、前記無毛皮膚表面において"
            "20〜100 mW/cm²のパワー密度で光を照射する、デバイス。"),
        3: ("請求項1または2に記載のデバイスにおいて、前記LED素子は、660±10 nm帯の"
            "第一波長群と830±10 nm帯の第二波長群を含み、前記第一波長群は"
            "S-ニトロソチオールの光分解に適し、前記第二波長群はシトクロムcオキシダーゼの"
            "Cu_B中心からのNO解離に適する、デバイス。"),
        4: ("請求項1〜3のいずれかに記載のデバイスにおいて、前記制御手段は、"
            "前記照射期間（T_on）を3〜30秒、前記非照射期間（T_off）を1〜15秒に設定し、"
            "デューティサイクルを50〜80%の範囲で前記フィードバック制御により動的に"
            "調節する、デバイス。"),
        5: ("請求項1〜4のいずれかに記載のデバイスにおいて、前記制御手段は、\n"
            "(i) 前記温度センサにより検知された末梢皮膚温が所定の下限閾値"
            "（28〜32°Cの範囲で設定可能）を下回った場合に、前記光照射手段の出力"
            "パワー密度を段階的に増強し、\n"
            "(ii) 前記末梢皮膚温が所定の上限閾値（38〜42°Cの範囲で設定可能）を超えた場合に、"
            "前記光照射手段を停止する安全遮断機能を有する、\n"
            "デバイス。"),
        6: ("請求項1〜5のいずれかに記載のデバイスにおいて、\n"
            "光電容積脈波法（PPG）センサをさらに備え、前記PPGセンサの出力信号から"
            "パルストランジットタイム（PTT）またはパルス波形解析に基づき装着者の血圧を"
            "推定する血圧推定手段をさらに有し、\n"
            "前記制御手段は、前記血圧推定手段により推定された血圧値が所定の目標血圧範囲を"
            "超えている場合に前記光照射手段の出力を増強し、前記目標血圧範囲内にある場合に"
            "前記出力を維持または減弱する閉ループ血圧フィードバック制御を実行する、\n"
            "デバイス。"),
        7: ("請求項1〜6のいずれかに記載のデバイスにおいて、前記デバイスは、\n"
            "(A) 手掌面を覆うグローブ型、\n"
            "(B) 足底面に配置されるインソール型、または\n"
            "(C) 手首から手掌にかけて巻き付けるバンド型\n"
            "のいずれかのウェアラブル形態を有する、デバイス。"),
        8: ("請求項1〜7のいずれかに記載のウェアラブル末梢血流促進デバイスの制御方法であって、\n"
            "(i) 前記光照射手段により、装着者の無毛皮膚に630〜850 nm帯の光をパルス照射する"
            "照射工程、\n"
            "(ii) 前記温度センサにより、装着者の末梢皮膚温を所定間隔で測定する測定工程、\n"
            "(iii) 測定された前記皮膚温に基づき、前記光照射手段のパワー密度および"
            "デューティサイクルの少なくとも一方を調節するフィードバック制御工程、および\n"
            "(iv) 測定された前記皮膚温が安全上限閾値を超えた場合に前記光照射手段を停止する"
            "安全遮断工程\n"
            "を含む、制御方法。"),
        9: ("請求項8に記載の制御方法において、前記フィードバック制御工程は、\n"
            "前記皮膚温の時間変化率（dT/dt）が負の所定閾値を下回った場合に、"
            "寒冷環境への移行と判断して前記光照射手段の出力を先行的に増強する"
            "予測制御を含む、制御方法。"),
        10: ("手掌および/または足底の無毛皮膚の表面に装着可能なウェアラブル"
             "末梢保温デバイスであって、\n"
             "(a) フレキシブル基板上に配列された複数のLED素子を含み、630〜850 nmの"
             "波長帯の光を前記無毛皮膚に照射する光照射手段、\n"
             "(b) 前記光照射手段を間欠的に駆動し、照射期間と非照射期間のデューティサイクルを"
             "制御する制御手段、\n"
             "(c) 装着者の末梢皮膚温を検知する温度センサ、および\n"
             "(d) 前記制御手段と無線通信により接続され、照射スケジュールの設定、使用履歴の記録、"
             "および末梢皮膚温の表示を提供する外部端末用アプリケーション\n"
             "を備え、\n"
             "前記制御手段は、前記温度センサにより検知された皮膚温に基づき前記光照射手段の"
             "出力をフィードバック制御し、前記光照射により末梢血流が増加して手足の保温効果が"
             "提供される、\n"
             "ウェアラブル末梢保温デバイス。"),
        11: ("請求項10に記載の末梢保温デバイスにおいて、\n"
             "光電容積脈波法（PPG）センサのハードウェアが前記フレキシブル基板上に実装されており、"
             "出荷時には前記PPGセンサの機能がファームウェアにより無効化されており、\n"
             "ソフトウェアアップデートの適用により前記PPGセンサの機能が有効化され、"
             "前記PPGセンサの出力に基づく血圧推定機能および血圧フィードバック制御機能が"
             "追加されるように構成された、デバイス。"),
        12: ("請求項10または11に記載のデバイスにおいて、前記外部端末用アプリケーションは、"
             "前記温度センサから取得した皮膚温データ、照射パラメータの履歴、および使用時間を"
             "匿名化して蓄積するデータ収集機能を有し、蓄積されたデータが前記デバイスの"
             "効果検証に利用可能である、デバイス。"),
    }

    for n, text in claims_data.items():
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        p.paragraph_format.space_before = Pt(8)
        p.paragraph_format.space_after = Pt(4)
        p.paragraph_format.line_spacing = Pt(18)
        r1 = p.add_run(f"\u3010\u8acb\u6c42\u9805{n}\u3011")
        _sf(r1, bold=True)
        r2 = p.add_run(f"\n{text}")
        _sf(r2)

    doc.save(str(OUT / "claims.docx"))
    print(f"  Saved: {OUT / 'claims.docx'}")


# ==========================================================
# 3. ABSTRACT
# ==========================================================
def build_abstract():
    doc = _doc()
    _h(doc, "\u3010\u66f8\u985e\u540d\u3011\u8981\u7d04\u66f8")
    _h(doc, "\u3010\u8981\u7d04\u3011")
    _h(doc, "\u3010\u8ab2\u984c\u3011", 2)
    _p(doc, "薬物を用いずにウェアラブル形態で末梢血管抵抗を低下させ、"
       "血圧を降下させるデバイスを提供する。")
    _h(doc, "\u3010\u89e3\u6c7a\u624b\u6bb5\u3011", 2)
    _p(doc, "手掌/足底のglabrous skinに装着可能なウェアラブルデバイスであって、"
       "フレキシブル基板上に配列されたLEDアレイ（630〜850 nm）と、"
       "パルス照射を制御する制御手段と、末梢皮膚温センサを備える。"
       "LEDから照射された光により、AVA周囲の血管内皮細胞からNOが光化学的に遊離し、"
       "cGMP経路を介してAVA平滑筋が弛緩してAVAが拡張し、末梢血流が増加する。"
       "パルス照射制御（T_on/T_off）によりNOストア枯渇を防止し、"
       "持続的なAVA拡張を維持する。"
       "温度フィードバックにより個体差への自動適応と安全遮断を実現する。")
    _h(doc, "\u3010\u9078\u629e\u56f3\u3011", 2)
    _p(doc, "\u56f31")

    doc.save(str(OUT / "abstract.docx"))
    print(f"  Saved: {OUT / 'abstract.docx'}")


# ==========================================================
# 4. FIGURES
# ==========================================================
def build_figures():
    plt.rcParams["font.family"] = ["Noto Sans CJK JP", "DejaVu Sans", "sans-serif"]

    # --- Figure 1: Mechanism ---
    fig, ax = plt.subplots(figsize=(12, 4))
    ax.set_xlim(0, 12); ax.set_ylim(0, 4)
    ax.axis("off")
    fig.suptitle("【図1】作用機序模式図", fontsize=14, fontweight="bold", y=0.98)

    boxes = [
        (0.3, 1.5, "PBM Light\n630-850 nm"),
        (2.3, 1.5, "NO Store\n(RSNO, CcO,\nHb-NO)"),
        (4.3, 1.5, "NO Release\n(photolysis)"),
        (6.3, 1.5, "sGC/cGMP\nActivation"),
        (8.3, 1.5, "AVA Smooth\nMuscle\nRelaxation"),
        (10.3, 1.5, "AVA Dilation\n→ TPR↓\n→ MAP↓"),
    ]
    for x, y, txt in boxes:
        rect = mpatches.FancyBboxPatch((x, y), 1.6, 1.2,
            boxstyle="round,pad=0.1", facecolor="#E3F2FD", edgecolor="#1565C0", linewidth=1.5)
        ax.add_patch(rect)
        ax.text(x + 0.8, y + 0.6, txt, ha="center", va="center", fontsize=8, fontweight="bold")

    for i in range(5):
        ax.annotate("", xy=(boxes[i+1][0], 2.1), xytext=(boxes[i][0]+1.6, 2.1),
            arrowprops=dict(arrowstyle="->", color="#1565C0", lw=2))

    # Pulse control feedback arrow
    ax.annotate("Pulse Control\n(T_on/T_off)\nNO Store Management",
        xy=(2.3, 1.5), xytext=(4.3, 0.3),
        arrowprops=dict(arrowstyle="->", color="#E65100", lw=1.5, connectionstyle="arc3,rad=0.3"),
        fontsize=7, color="#E65100", ha="center")

    fig.tight_layout()
    fig.savefig(str(OUT / "fig1_mechanism.png"), dpi=300, bbox_inches="tight")
    plt.close(fig)
    print(f"  Saved: {OUT / 'fig1_mechanism.png'}")

    # --- Figure 2: Device config ---
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 5))
    fig.suptitle("【図2】グローブ型デバイスの構成", fontsize=14, fontweight="bold")

    # Palm view
    ax1.set_xlim(0, 10); ax1.set_ylim(0, 12)
    ax1.set_title("(a) 平面図（手掌側）", fontsize=11)
    ax1.axis("off")

    # Glove outline
    glove_x = [2, 1.5, 1, 0.5, 0.8, 1.5, 1, 0.5, 0.8, 1.5, 1.5, 0.8, 1.2, 2,
               2.5, 2, 1.5, 2, 3, 3.5, 4.5, 5, 5.5, 6, 7, 8, 8, 7.5, 7, 6, 5, 4, 3, 2]
    glove_y = [2, 2, 3, 5, 5.5, 4, 5.5, 7.5, 8, 6, 7, 9.5, 10, 7.5,
               8, 10, 10.5, 9, 7.5, 10, 10.5, 10, 9, 7.5, 6, 5, 3, 2.5, 2.2, 2, 1.5, 1.2, 1.5, 2]
    ax1.fill(glove_x, glove_y, alpha=0.1, color="#90CAF9")
    ax1.plot(glove_x, glove_y, color="#1565C0", linewidth=2)

    # LED array
    for r in range(4):
        for c in range(8):
            x = 2.8 + c * 0.45
            y = 2.5 + r * 0.9
            ax1.plot(x, y, "rs", markersize=5, markerfacecolor="#F44336")
    ax1.text(5, 1.5, "12: LEDアレイ (4×8)\n    660nm × 32", fontsize=8, ha="center", color="#D32F2F")

    # Control unit at wrist
    ctrl = mpatches.FancyBboxPatch((3, 0.3), 4, 1,
        boxstyle="round,pad=0.1", facecolor="#FFF9C4", edgecolor="#F57F17", linewidth=1.5)
    ax1.add_patch(ctrl)
    ax1.text(5, 0.8, "20: 制御基板 + 40: バッテリー\n(51: リストバンド部)", fontsize=7, ha="center")

    # Cross section
    ax2.set_xlim(0, 10); ax2.set_ylim(0, 8)
    ax2.set_title("(b) 断面図（A-A'線）", fontsize=11)
    ax2.axis("off")

    layers = [
        (0.5, 0.5, 9, 1.0, "#FFECB3", "50: シリコーンゴム製グローブ"),
        (0.5, 1.5, 9, 0.5, "#FFCDD2", "10: FPC + 11: LED素子"),
        (0.5, 2.0, 9, 0.3, "#E0E0E0", "接着層"),
        (0.5, 2.3, 9, 1.0, "#FFCCBC", "A: glabrous skin"),
        (0.5, 3.3, 9, 0.8, "#FFAB91", "真皮 (AVA: 深さ1-3mm)"),
        (0.5, 4.1, 9, 1.0, "#FF8A65", "皮下組織"),
    ]
    for x, y, w, h, color, label in layers:
        rect = mpatches.FancyBboxPatch((x, y), w, h,
            boxstyle="square,pad=0", facecolor=color, edgecolor="black", linewidth=0.5)
        ax2.add_patch(rect)
        ax2.text(x + w + 0.1, y + h/2, label, fontsize=8, va="center")

    # LED arrows
    for i in range(5):
        ax2.annotate("", xy=(2 + i*1.5, 2.3), xytext=(2 + i*1.5, 1.5),
            arrowprops=dict(arrowstyle="->", color="#F44336", lw=1))
    ax2.text(5, 1.2, "↓ 660nm light", fontsize=8, ha="center", color="#F44336")

    # AVA dots
    for i in range(4):
        ax2.plot(2.5 + i*1.5, 3.7, "o", color="#D32F2F", markersize=8)
    ax2.text(5, 3.0, "B: AVA", fontsize=8, ha="center", color="#D32F2F", fontweight="bold")

    fig.tight_layout()
    fig.savefig(str(OUT / "fig2_device.png"), dpi=300, bbox_inches="tight")
    plt.close(fig)
    print(f"  Saved: {OUT / 'fig2_device.png'}")

    # --- Figure 3: Control flowchart ---
    fig, ax = plt.subplots(figsize=(10, 12))
    ax.set_xlim(0, 10); ax.set_ylim(0, 16)
    ax.axis("off")
    fig.suptitle("【図3】閉ループ制御アルゴリズム", fontsize=14, fontweight="bold")

    def _box(x, y, w, h, text, color="#E3F2FD", ec="#1565C0"):
        r = mpatches.FancyBboxPatch((x, y), w, h,
            boxstyle="round,pad=0.1", facecolor=color, edgecolor=ec, linewidth=1.5)
        ax.add_patch(r)
        ax.text(x+w/2, y+h/2, text, ha="center", va="center", fontsize=8, fontweight="bold")

    def _diamond(cx, cy, w, h, text, color="#FFF9C4", ec="#F57F17"):
        pts = np.array([[cx, cy+h/2], [cx+w/2, cy], [cx, cy-h/2], [cx-w/2, cy], [cx, cy+h/2]])
        ax.fill(pts[:,0], pts[:,1], color=color)
        ax.plot(pts[:,0], pts[:,1], color=ec, linewidth=1.5)
        ax.text(cx, cy, text, ha="center", va="center", fontsize=7, fontweight="bold")

    def _arr(x1, y1, x2, y2, label=""):
        ax.annotate("", xy=(x2, y2), xytext=(x1, y1),
            arrowprops=dict(arrowstyle="->", color="#333", lw=1.5))
        if label:
            mx, my = (x1+x2)/2, (y1+y2)/2
            ax.text(mx+0.15, my, label, fontsize=7, color="#333")

    _box(3, 14.5, 4, 0.8, "START: デバイス装着")
    _arr(5, 14.5, 5, 14)
    _box(3, 13.2, 4, 0.8, "初期設定:\n40 mW/cm², DC=67%")
    _arr(5, 13.2, 5, 12.7)
    _box(3, 11.9, 4, 0.8, "パルス照射開始\n(T_on / T_off)")
    _arr(5, 11.9, 5, 11.4)
    _box(3, 10.6, 4, 0.8, "皮膚温 T_skin 測定\n(MLX90614, 0.5秒間隔)")
    _arr(5, 10.6, 5, 10.1)

    _diamond(5, 9.3, 3.5, 1.2, "T_skin > 42°C ?")
    _arr(6.75, 9.3, 8, 9.3, "YES")
    _box(7.5, 8.9, 2, 0.8, "即時停止\n(安全遮断)", "#FFCDD2", "#D32F2F")

    _arr(5, 8.7, 5, 8.2, "NO")
    _diamond(5, 7.4, 3.5, 1.2, "T_skin > 38°C ?")
    _arr(6.75, 7.4, 8, 7.4, "YES")
    _box(7.5, 7.0, 2, 0.8, "出力減弱\n→20 mW/cm²", "#E8F5E9", "#2E7D32")

    _arr(5, 6.8, 5, 6.3, "NO")
    _diamond(5, 5.5, 3.5, 1.2, "T_skin < 30°C ?")
    _arr(6.75, 5.5, 8, 5.5, "YES")
    _box(7.5, 5.1, 2, 0.8, "出力増強\n→60 mW/cm²", "#FFF3E0", "#E65100")

    _arr(5, 4.9, 5, 4.4, "NO")
    _diamond(5, 3.6, 3.5, 1.2, "PPG有効 &\nBP > 目標+5 ?")
    _arr(6.75, 3.6, 8, 3.6, "YES")
    _box(7.5, 3.2, 2, 0.8, "DC増加\n(+5%/min)", "#E8EAF6", "#283593")

    _arr(5, 3.0, 5, 2.5, "NO")
    _diamond(5, 1.7, 3.5, 1.2, "セッション\n終了条件 ?")
    _arr(6.75, 1.7, 8, 1.7, "YES")
    _box(7.5, 1.3, 2, 0.8, "照射停止\nEND", "#EFEBE9", "#4E342E")

    _arr(3.25, 1.7, 1, 1.7)
    ax.annotate("", xy=(1, 10.6), xytext=(1, 1.7),
        arrowprops=dict(arrowstyle="->", color="#333", lw=1.5))
    ax.annotate("", xy=(3, 11.0), xytext=(1, 11.0),
        arrowprops=dict(arrowstyle="->", color="#333", lw=1.5))
    ax.text(0.5, 6, "LOOP", fontsize=8, rotation=90, va="center", color="#666")

    fig.tight_layout()
    fig.savefig(str(OUT / "fig3_control.png"), dpi=300, bbox_inches="tight")
    plt.close(fig)
    print(f"  Saved: {OUT / 'fig3_control.png'}")

    # --- Figure 4: Pulse timing ---
    fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(10, 5), sharex=True)
    fig.suptitle("【図4】パルス照射制御のタイミングチャート", fontsize=14, fontweight="bold")

    t = np.arange(0, 60, 0.01)
    t_on, t_off = 10, 5
    cycle = t_on + t_off
    led = np.array([1 if (ti % cycle) < t_on else 0 for ti in t])

    ax1.fill_between(t, led * 40, alpha=0.3, color="#F44336")
    ax1.plot(t, led * 40, color="#F44336", linewidth=1)
    ax1.set_ylabel("Irradiance\n(mW/cm²)", fontsize=9)
    ax1.set_ylim(-5, 50)
    ax1.axhline(y=0, color="black", linewidth=0.5)
    for i in range(4):
        x_start = i * cycle
        ax1.annotate("", xy=(x_start + t_on, 45), xytext=(x_start, 45),
            arrowprops=dict(arrowstyle="<->", color="#333", lw=1))
        ax1.text(x_start + t_on/2, 47, f"T_on={t_on}s", ha="center", fontsize=7)
        ax1.annotate("", xy=(x_start + cycle, 45), xytext=(x_start + t_on, 45),
            arrowprops=dict(arrowstyle="<->", color="#333", lw=1))
        ax1.text(x_start + t_on + t_off/2, 47, f"T_off={t_off}s", ha="center", fontsize=7)
    ax1.set_title("LED出力", fontsize=10, loc="left")

    # NO store level (estimated)
    no_store = np.zeros_like(t)
    no_store[0] = 100
    for i in range(1, len(t)):
        if led[i] > 0:
            no_store[i] = max(20, no_store[i-1] - 0.8)
        else:
            no_store[i] = min(100, no_store[i-1] + 1.5)

    ax2.plot(t, no_store, color="#1565C0", linewidth=2)
    ax2.fill_between(t, no_store, alpha=0.2, color="#1565C0")
    ax2.set_ylabel("NO Store\nLevel (%)", fontsize=9)
    ax2.set_xlabel("Time (s)", fontsize=9)
    ax2.set_ylim(0, 110)
    ax2.axhline(y=20, color="#D32F2F", linestyle="--", linewidth=1, label="Depletion threshold")
    ax2.legend(fontsize=8, loc="lower right")
    ax2.set_title("推定NOストアレベル", fontsize=10, loc="left")

    fig.tight_layout()
    fig.savefig(str(OUT / "fig4_pulse.png"), dpi=300, bbox_inches="tight")
    plt.close(fig)
    print(f"  Saved: {OUT / 'fig4_pulse.png'}")

    # --- Figure 5: Comparison table ---
    fig, ax = plt.subplots(figsize=(12, 5))
    ax.axis("off")
    fig.suptitle("【図5】従来技術との比較表", fontsize=14, fontweight="bold")

    headers = ["要素", "AVACEN\n(特許文献2)", "Diller\n(特許文献1)", "ROX Coupler\n(非特許文献1)", "本発明"]
    data = [
        ["AVA拡張機序", "陰圧+加温\n(TRPV4)", "中枢反射\n(視床下部)", "物理的AV shunt\n(侵襲的)", "光化学的\nNO放出"],
        ["刺激部位", "手掌\n(チャンバー内)", "脊椎\n(頸椎/腰椎)", "腸骨\n(経皮的留置)", "手掌/足底\n(glabrous skin)"],
        ["形態", "据え置き型\n(5 kg)", "ベッド統合型", "埋込型", "ウェアラブル\n(<80 g)"],
        ["侵襲性", "非侵襲", "非侵襲", "侵襲的", "非侵襲"],
        ["使用場面", "静止状態のみ", "睡眠中のみ", "常時", "日常動作中"],
        ["消費電力", ">50 W", "不明", "N/A", "<2 W"],
        ["パルス制御", "なし", "なし", "N/A", "あり\n(NOストア管理)"],
        ["降圧効果", "要検証", "SBP -3~-8", "SBP -26.9", "理論値\nSBP -5~-10"],
    ]

    table = ax.table(cellText=data, colLabels=headers, loc="center",
                     cellLoc="center", colColours=["#E3F2FD"]*5)
    table.auto_set_font_size(False)
    table.set_fontsize(8)
    table.scale(1.0, 1.6)

    # Highlight "本発明" column
    for i in range(len(data) + 1):
        table[i, 4].set_facecolor("#E8F5E9")
        table[i, 4].set_text_props(fontweight="bold")

    fig.tight_layout()
    fig.savefig(str(OUT / "fig5_comparison.png"), dpi=300, bbox_inches="tight")
    plt.close(fig)
    print(f"  Saved: {OUT / 'fig5_comparison.png'}")

    # --- Figure 6: Staged upgrade ---
    fig, ax = plt.subplots(figsize=(12, 5))
    ax.set_xlim(0, 12); ax.set_ylim(0, 6)
    ax.axis("off")
    fig.suptitle("【図6】段階的機能昇格（民生品→医療機器）", fontsize=14, fontweight="bold")

    # Phase 1
    r1 = mpatches.FancyBboxPatch((0.5, 1), 3.5, 3.5,
        boxstyle="round,pad=0.2", facecolor="#E8F5E9", edgecolor="#2E7D32", linewidth=2)
    ax.add_patch(r1)
    ax.text(2.25, 4.0, "Phase 1: 民生品", fontsize=11, fontweight="bold",
            ha="center", color="#2E7D32")
    ax.text(2.25, 3.2, "末梢保温グローブ\n¥9,800〜14,800\n\n"
            "・LED照射 + 温度FB ✓\n・PPGセンサ: 搭載済み\n  (FW無効化)\n・薬事: 不要 (雑貨)",
            fontsize=8, ha="center", va="top")

    # Arrow
    ax.annotate("", xy=(5.5, 2.75), xytext=(4.2, 2.75),
        arrowprops=dict(arrowstyle="-|>", color="#333", lw=2.5))
    ax.text(4.85, 3.5, "OTAアップデート\n(FW更新のみ)\nHW変更なし",
            fontsize=8, ha="center", va="center", color="#666")

    # Phase 2
    r2 = mpatches.FancyBboxPatch((5.5, 1), 3.5, 3.5,
        boxstyle="round,pad=0.2", facecolor="#E3F2FD", edgecolor="#1565C0", linewidth=2)
    ax.add_patch(r2)
    ax.text(7.25, 4.0, "Phase 2: 医療機器", fontsize=11, fontweight="bold",
            ha="center", color="#1565C0")
    ax.text(7.25, 3.2, "血圧管理デバイス\n¥30,000\n\n"
            "・LED照射 + 温度FB ✓\n・PPGセンサ: 有効化 ✓\n・BP推定 + BPフィードバック ✓\n"
            "・薬事: クラスII認証",
            fontsize=8, ha="center", va="top")

    # RWD arrow
    ax.annotate("", xy=(9.5, 2.75), xytext=(10.5, 2.75),
        arrowprops=dict(arrowstyle="<-", color="#E65100", lw=1.5))
    r3 = mpatches.FancyBboxPatch((10, 1.5), 1.8, 2.5,
        boxstyle="round,pad=0.1", facecolor="#FFF3E0", edgecolor="#E65100", linewidth=1.5)
    ax.add_patch(r3)
    ax.text(10.9, 3.5, "Phase 1 RWD\n(リアルワールド\nデータ)", fontsize=8,
            ha="center", va="center", color="#E65100", fontweight="bold")
    ax.text(10.9, 2.0, "N=10,000+の\n使用データで\n臨床試験簡略化\n(PMDA 2023)",
            fontsize=7, ha="center", va="center", color="#E65100")

    # Volume effect
    ax.text(2.25, 0.5, "量産: 10,000〜50,000台/年\n→ BOM ¥2,000以下",
            fontsize=8, ha="center", color="#2E7D32")

    fig.tight_layout()
    fig.savefig(str(OUT / "fig6_staged.png"), dpi=300, bbox_inches="tight")
    plt.close(fig)
    print(f"  Saved: {OUT / 'fig6_staged.png'}")

    # --- PPTX with all figures ---
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for fname, title in [
        ("fig1_mechanism.png", "【図1】作用機序模式図"),
        ("fig2_device.png", "【図2】グローブ型デバイスの構成"),
        ("fig3_control.png", "【図3】閉ループ制御フローチャート"),
        ("fig4_pulse.png", "【図4】パルス照射タイミングチャート"),
        ("fig5_comparison.png", "【図5】従来技術との比較表"),
        ("fig6_staged.png", "【図6】段階的機能昇格"),
    ]:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        fpath = str(OUT / fname)
        slide.shapes.add_picture(fpath, Inches(0.5), Inches(0.5), Inches(12), Inches(6.5))

    prs.save(str(OUT / "drawings.pptx"))
    print(f"  Saved: {OUT / 'drawings.pptx'}")


# ==========================================================
# 5. REJECTION STRATEGY (出願人控え)
# ==========================================================
def build_rejection_strategy():
    doc = _doc()
    _h(doc, "拒絶理由通知への対応方針（出願人控え）")
    _p(doc, "本書類は出願書類には含めない。審査官対応の際の方針として記録する。", bold=True)

    sections = [
        ("§1. 進歩性拒絶（特許法29条2項）への対応",
         "想定される引用例の組み合わせ:\n"
         "・引用文献1: 特許文献2（AVACEN）— AVA拡張による降圧の概念\n"
         "・引用文献2: 非特許文献3（Ribeiro 2022）— PBMによるNO放出→降圧\n"
         "・審査官の論理: 「引用文献1のAVA拡張を、引用文献2のPBMで実現することは"
         "当業者が容易に想到し得る」\n\n"
         "反論方針:\n"
         "(i) 動機付けの欠如: Ribeiroはラットの尾動脈に全身照射した実験であり、"
         "「手掌glabrous skinのAVAにPBMを適用する」という技術思想は記載されていない。"
         "AVACENはPBMを一切使用しておらず、光化学的NO放出への言及がない。\n"
         "(ii) 阻害要因: AVACENは陰圧が「AVA拡張に必須」としており、"
         "陰圧を排除して光照射のみでAVA拡張を達成するという技術思想は、"
         "引用文献1の教示に反する（阻害要因）。\n"
         "(iii) 予測できない顕著な効果: パルス照射制御による持続4時間の効果維持、"
         "寒冷昇圧の予防的抑制、消費電力25分の1以下でのウェアラブル化。"),

        ("§2. 新規性拒絶（特許法29条1項）への対応",
         "想定引用: 特許文献5（Enwemeka, US 10,596,388）\n\n"
         "反論: Enwemekaは創傷治癒・感染症治療目的。AVA拡張によるTPR低下/血圧降下は"
         "開示されていない。照射対象は創傷部位であり、glabrous skinのAVAへの適用は"
         "開示されていない。パルス照射によるNOストア管理の概念がない。"
         "温度/血圧フィードバック制御がない。"),

        ("§3. 記載不備（特許法36条）への対応",
         "想定指摘: 「NOが光化学的に遊離し」は物理現象の記述であり構成要件ではない\n\n"
         "反論: whereby節はデバイスの「使用状態における作用」を特定する機能的限定として"
         "許容される（知財高裁平成22年(行ケ)第10340号参照）。\n\n"
         "予備的補正案: whereby節を削除し、構成(a)(b)(c)のみの構成クレームに補正。"),

        ("§4. サポート要件（特許法36条6項1号）への対応",
         "想定指摘: 「20〜100 mW/cm²」の数値範囲のサポート\n\n"
         "反論:\n"
         "・下限20 mW/cm²: 非特許文献9のbiphasic dose response閾値から逆算\n"
         "・上限100 mW/cm²: IEC 62471免除グループ上限値。「非熱的PBM」の定義上限。"),

        ("§5. 医療行為除外への対応",
         "想定指摘: 請求項8が「人間を治療する方法」に該当\n\n"
         "反論: 請求項8は「デバイスの制御方法」であり、医療行為を規定していない。\n"
         "予備的補正案: 「コンピュータに実行させるプログラム」形式に変更。"),

        ("§6. 単一性（特許法37条）への対応",
         "想定指摘: 請求項1-9と請求項10-12が単一の発明を構成しない\n\n"
         "反論: 共通する特別な技術的特徴: 「フレキシブル基板上のLEDアレイによる"
         "glabrous skin AVAへの630-850nm光照射＋パルス制御＋温度フィードバック」。"
         "請求項10は請求項1の構成(a)(b)(c)の全てを含む。"),

        ("§7. 補正シナリオ",
         "シナリオA: 進歩性拒絶に対し請求項1に限定を加える場合\n"
         "・A1: パワー密度20-100 mW/cm²を統合\n"
         "・A2: 二波長660nm+830nmを必須に\n"
         "・A3: デューティサイクル具体値を必須に\n"
         "・A4: A1+A3の組み合わせ\n\n"
         "シナリオB: 新規性拒絶に対し機能的限定を追加\n"
         "・B1: whereby節に「末梢血管抵抗を低下させて血圧を降下させる」を明記\n"
         "・B2: AVA口径20-150µmを限定\n\n"
         "シナリオC: 記載不備に対しwhereby節を削除\n"
         "・構成(a)(b)(c)のみで請求項構成"),

        ("§8. 分割出願の検討",
         "・分割案1: 「寒冷昇圧予防モード」（請求項9）を独立出願に\n"
         "・分割案2: 「ソフトウェアアップデートによる民生品→医療機器昇格」"
         "（請求項11-12）を独立出願に\n"
         "・分割案3: 「PPGベースBPフィードバック制御」（請求項6）を独立出願に"),
    ]

    for title, body in sections:
        _h(doc, title, 2)
        _p(doc, body)

    doc.save(str(OUT / "rejection_strategy.docx"))
    print(f"  Saved: {OUT / 'rejection_strategy.docx'}")


# ==========================================================
if __name__ == "__main__":
    print("Generating patent documents...")
    build_meisaisho()
    build_claims()
    build_abstract()
    build_figures()
    build_rejection_strategy()
    print(f"\nDone. All files in: {OUT}")
