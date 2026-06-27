"""
Generate Medical Hypotheses manuscript:
"Photobiomodulation-induced vasodilation of arteriovenous anastomoses
as a non-pharmacological strategy for blood pressure reduction:
A hypothesis"

Format: Research Article--Hypotheses
- Unstructured abstract: max 250 words
- Body: max 3,000 words
- Max 3 tables/figures
- Max 6 keywords
- Max 50 references
- Vancouver numbered citation style
"""

from docx import Document
from docx.shared import Pt, Inches, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN
import re
from pathlib import Path

OUTPUT_DIR = Path(__file__).parent


# ─────────────────────────────────────────────
# Manuscript content
# ─────────────────────────────────────────────

TITLE = (
    "Photobiomodulation-induced vasodilation of glabrous skin "
    "arteriovenous anastomoses as a wearable, non-pharmacological "
    "strategy for blood pressure reduction: A hypothesis"
)

AUTHORS = "Tatsuki Onishi"
AFFILIATION = ""  # To be filled by author

KEYWORDS = [
    "Arteriovenous anastomoses",
    "Photobiomodulation",
    "Blood pressure",
    "Low-level laser therapy",
    "Peripheral vascular resistance",
    "Wearable device",
]

ABSTRACT = (
    "Hypertension remains the leading modifiable risk factor for "
    "cardiovascular disease, yet medication non-adherence and adverse "
    "effects of antihypertensive drugs, particularly diuretics, limit "
    "treatment efficacy. Arteriovenous anastomoses (AVAs) are direct "
    "vascular connections between arterioles and venules concentrated "
    "in glabrous skin of the palms, soles, and digits. When dilated, "
    "AVAs create low-resistance parallel pathways that reduce total "
    "peripheral resistance (TPR) and consequently lower arterial blood "
    "pressure. We hypothesize that targeted photobiomodulation (PBM) "
    "at 630\u2013850 nm, delivered via a wearable LED array to palmar and "
    "plantar surfaces, can selectively dilate AVAs through nitric oxide "
    "(NO) release from endothelial stores, thereby producing clinically "
    "significant blood pressure reduction without pharmacological "
    "intervention. This hypothesis is supported by three convergent "
    "lines of evidence: (1) surgically created arteriovenous fistulae "
    "reduce systolic blood pressure by 27 mmHg in randomized trials; "
    "(2) PBM at 660 nm produces NO-dependent vasodilation and "
    "hypotensive effects in experimental models; and (3) AVA blood "
    "flow fluctuations correlate inversely with mean arterial pressure "
    "(r = \u22120.8). A wearable PBM device targeting AVA-rich glabrous "
    "skin could offer a drug-free, side-effect-free alternative to "
    "third-line antihypertensives, with the added benefit of peripheral "
    "warming in cold environments."
)

INTRODUCTION = [
    (
        "Hypertension affects over 1.3 billion adults globally and remains "
        "the single largest contributor to cardiovascular mortality.{1} "
        "Despite the availability of multiple drug classes, approximately "
        "50% of treated patients fail to achieve target blood pressure, "
        "and medication non-adherence rates reach 50\u201380% within one year "
        "of initiation.{2} Third-line agents, particularly thiazide "
        "diuretics, carry clinically significant adverse effects including "
        "hypokalemia, hyperuricemia, glucose intolerance, and dehydration, "
        "which further compromise adherence.{3}"
    ),
    (
        "The fundamental hemodynamic equation, mean arterial pressure "
        "(MAP) = cardiac output (CO) \u00d7 total peripheral resistance (TPR), "
        "dictates that any sustained reduction in TPR, without compensatory "
        "increases in CO, will lower blood pressure.{4} Arteriovenous "
        "anastomoses (AVAs) are direct vascular connections between "
        "arterioles and venules, concentrated predominantly in glabrous "
        "(non-hairy) skin of the palms, soles, lips, and ears.{5} With "
        "luminal diameters of 20\u2013150 \u00b5m (5\u201310 times larger than "
        "capillaries), AVAs present dramatically lower resistance to flow "
        "when dilated.{6} Their primary physiological role is "
        "thermoregulation: sympathetic withdrawal dilates AVAs, increasing "
        "skin blood flow up to 1 L/min and facilitating radiative heat "
        "loss.{7}"
    ),
    (
        "Crucially, AVA dilation adds low-resistance parallel pathways to "
        "the systemic circulation, mechanically reducing TPR. This "
        "relationship is not merely theoretical: Wall\u00f8e demonstrated that "
        "spontaneous AVA blood flow oscillations correlate inversely with "
        "MAP fluctuations (r = \u22120.3 to \u22120.8), with peak-to-trough MAP "
        "variations of 15 mmHg coinciding with AVA opening and closing "
        "cycles.{5} Furthermore, surgically created central arteriovenous "
        "fistulae (the ROX Coupler) reduce office systolic blood pressure "
        "by 26.9 mmHg (p < 0.0001) in patients with resistant "
        "hypertension, providing definitive proof that adding parallel "
        "low-resistance vascular pathways lowers blood pressure.{8}"
    ),
    (
        "Photobiomodulation (PBM), formerly termed low-level laser therapy "
        "(LLLT), employs non-ionizing light at 600\u20131000 nm to modulate "
        "cellular function.{9} A well-characterized mechanism of PBM is "
        "the release of nitric oxide (NO) from intracellular stores "
        "including cytochrome c oxidase (Complex IV) and "
        "S-nitrosothiols.{10} NO is a potent vasodilator that activates "
        "soluble guanylate cyclase, elevates cyclic GMP, and relaxes "
        "vascular smooth muscle.{11} Recent studies demonstrate that PBM "
        "at 660 nm produces NO-dependent hypotensive effects in "
        "hypertensive animal models, with the effect abolished by NOS "
        "inhibition and restored by NO donors.{12}"
    ),
    (
        "We propose that these two established phenomena\u2014AVA-mediated TPR "
        "reduction and PBM-induced NO release\u2014can be integrated into a "
        "single wearable device that non-invasively lowers blood pressure "
        "by photobiomodulating AVA-rich glabrous skin."
    ),
]

HYPOTHESIS = [
    (
        "We hypothesize that targeted photobiomodulation (630\u2013850 nm) "
        "delivered to palmar and/or plantar glabrous skin via a wearable "
        "LED array will produce clinically significant reductions in "
        "arterial blood pressure (\u22655 mmHg systolic) through selective "
        "vasodilation of arteriovenous anastomoses mediated by "
        "photochemical release of nitric oxide from endothelial stores."
    ),
    (
        "The proposed mechanism operates through the following causal "
        "chain (Fig. 1): PBM irradiation (20\u2013100 mW/cm\u00b2) penetrates the "
        "dermis to a depth of 14\u201326 mm,{13} reaching AVAs located at "
        "1\u20133 mm depth in glabrous skin.{5} Photon absorption by "
        "cytochrome c oxidase and hemoglobin-bound nitrosyl groups "
        "releases NO into the vascular smooth muscle layer "
        "surrounding AVAs.{10,12} NO activates the sGC/cGMP/PKG pathway, "
        "causing smooth muscle relaxation and AVA dilation.{11} The "
        "resulting increase in AVA conductance adds parallel low-resistance "
        "pathways to the systemic circulation, reducing TPR.{5} With "
        "MAP = CO \u00d7 TPR, a 5\u201310% reduction in TPR yields a 4.5\u20139 mmHg "
        "reduction in MAP (assuming baseline MAP of 90 mmHg), which exceeds "
        "the threshold for clinical significance.{14}"
    ),
    (
        "This mechanism is fundamentally distinct from the thermal "
        "vasodilation employed by existing devices (e.g., AVACEN, which "
        "uses conductive heating plus negative pressure,{15} and the "
        "selective thermal stimulation system of Diller et al., which "
        "targets spinal thermoregulatory centers during sleep).{29} PBM offers "
        "a photochemical pathway to vasodilation that does not require "
        "tissue heating, operates at milliwatt-level power consumption "
        "suitable for battery-powered wearables, and avoids the "
        "engineering complexity of vacuum-sealed chambers."
    ),
]

EVALUATION = [
    (
        "Three independent lines of evidence support the plausibility "
        "of this hypothesis:"
    ),
    (
        "First, the principle that adding arteriovenous shunt flow "
        "reduces systemic blood pressure is established at the highest "
        "level of clinical evidence. The ROX CONTROL HTN randomized "
        "controlled trial (n = 83) demonstrated that a surgically placed "
        "iliac arteriovenous coupler reduced office systolic blood "
        "pressure by 26.9 mmHg versus 3.7 mmHg in controls (p < 0.0001), "
        "with 24-hour ambulatory reductions of 13.5 mmHg sustained at "
        "12 months.{8,16} Similarly, Faul et al. showed that iliac "
        "arteriovenous shunts in hypertensive COPD patients reduced "
        "systolic blood pressure by 13 mmHg at 12 months (p < 0.0001).{17} "
        "These trials confirm that TPR reduction via parallel "
        "low-resistance pathways is an effective antihypertensive strategy."
    ),
    (
        "Second, PBM-induced vasodilation and blood pressure reduction "
        "have been demonstrated experimentally. Ribeiro et al. showed "
        "that 660 nm laser irradiation reduced MAP in hypertensive rats, "
        "and that this effect was abolished by the NOS inhibitor L-NAME, "
        "confirming NO-dependence.{12} A 2025 systematic review and "
        "meta-analysis confirmed that PBM reduces systolic blood pressure, "
        "diastolic blood pressure, and MAP while increasing NO levels, "
        "though with very low certainty of evidence requiring further "
        "high-quality trials.{18} Importantly, PBM at clinically relevant "
        "power densities (50\u2013100 mW/cm\u00b2) increases dermal blood flow by "
        "25\u201330% as measured by laser Doppler flowmetry.{19}"
    ),
    (
        "Third, the physiological link between AVA blood flow and blood "
        "pressure is well-established. Wall\u00f8e's comprehensive review "
        "documented inverse correlations between AVA blood flow and MAP "
        "(r = \u22120.3 to \u22120.8) during thermoneutral-zone oscillations, "
        "with MAP fluctuations of approximately 15 mmHg.{5} Kr\u00e4uchi et al. "
        "demonstrated in 51 subjects that the distal-to-proximal skin "
        "temperature gradient (a surrogate for AVA blood flow) predicts "
        "blood pressure dipping status (r = \u22120.436, p = 0.0014), with "
        "78.6% concordance between wrist skin temperature rhythms and "
        "blood pressure dipping patterns.{20,21} Furthermore, Gani et al. "
        "showed that local passive heat application (40\u201342\u00b0C), which "
        "activates AVA vasodilation, reduced systolic blood pressure by "
        "19 mmHg acutely and 28 mmHg overnight in patients with autonomic "
        "failure (p < 0.001 versus sham).{22}"
    ),
    (
        "A potential counterargument is baroreflex compensation: acute "
        "TPR reduction may trigger sympathetic activation that increases "
        "heart rate and cardiac output, partially offsetting the blood "
        "pressure decrease.{4} However, several observations mitigate this "
        "concern. The Vanderbilt heat study demonstrated sustained "
        "overnight blood pressure reductions of 28 mmHg, suggesting that "
        "baroreflex compensation is incomplete during prolonged "
        "stimulation.{22} Additionally, PBM has been reported to enhance "
        "parasympathetic tone,{23} which may blunt the sympathetic "
        "reflex. Gradual-onset stimulation (ramping PBM intensity over "
        "minutes rather than seconds) could further minimize baroreflex "
        "triggering. Finally, elderly patients with resistant "
        "hypertension\u2014the primary target population\u2014typically exhibit "
        "reduced baroreflex sensitivity, making sustained TPR-mediated "
        "blood pressure reduction more achievable.{24}"
    ),
]

TESTING = [
    (
        "The hypothesis can be tested through the following approaches:"
    ),
    (
        "A proof-of-concept study in healthy volunteers (n = 20\u201330) would "
        "apply PBM (660\u2013830 nm LED array, 50\u2013100 mW/cm\u00b2, 20\u201330 minutes) "
        "to palmar skin while continuously monitoring blood pressure via "
        "finger photoplethysmography (e.g., Finapres) and skin blood flow "
        "via laser Doppler flowmetry. The primary endpoint would be change "
        "in MAP from baseline. A sham-controlled crossover design with "
        "an identical but non-emitting device would control for placebo "
        "effects. Simultaneous measurement of distal skin temperature "
        "would confirm AVA engagement, while plasma nitrite/nitrate levels "
        "would verify the NO-dependent mechanism."
    ),
    (
        "If the proof-of-concept yields positive results, a Phase II "
        "randomized controlled trial in patients with Stage 1\u20132 "
        "hypertension (n = 60\u201380) would compare active PBM (worn for "
        "4\u20138 hours daily for 8 weeks) versus sham, with 24-hour "
        "ambulatory blood pressure as the primary endpoint. Subgroup "
        "analyses by season (winter vs. summer) would test the additional "
        "hypothesis that cold-weather AVA closure contributes to seasonal "
        "blood pressure elevation and that PBM can prevent this increase."
    ),
    (
        "Dose-response relationships should be characterized across "
        "wavelengths (630, 660, 810, 850 nm), power densities "
        "(20\u2013200 mW/cm\u00b2), and irradiation durations to identify "
        "optimal parameters for sustained AVA dilation without tachyphylaxis."
    ),
]

CONSEQUENCES = [
    (
        "If confirmed, this hypothesis has several important implications. "
        "First, a wearable PBM device targeting AVA-rich skin would "
        "represent a fundamentally new class of non-pharmacological "
        "antihypertensive intervention\u2014one that operates through physical "
        "vascular mechanics rather than neurohormonal modulation. Unlike "
        "renal denervation or baroreflex activation therapy, this approach "
        "is entirely non-invasive and reversible."
    ),
    (
        "Second, the device concept uniquely addresses the problem of "
        "cold-induced hypertension.{25} In temperate climates, blood "
        "pressure rises by 5\u201310 mmHg in winter due to sympathetically "
        "mediated AVA closure and resultant TPR elevation.{26} A wearable "
        "PBM device that maintains AVA patency could simultaneously "
        "prevent seasonal blood pressure elevation and provide peripheral "
        "warming\u2014transforming an antihypertensive device into one that "
        "patients actively desire to use, potentially resolving the "
        "longstanding challenge of treatment adherence."
    ),
    (
        "Third, this approach may complement or substitute for third-line "
        "diuretic therapy. For patients intolerant of thiazides due to "
        "metabolic adverse effects (hypokalemia, hyperuricemia, impaired "
        "glucose tolerance), a PBM-AVA device could provide equivalent "
        "blood pressure reduction (5\u201310 mmHg) without systemic metabolic "
        "perturbation.{3}"
    ),
    (
        "Fourth, the same device platform could address Raynaud's "
        "phenomenon, a condition characterized by episodic AVA closure "
        "in digits, with clear endpoints (attack frequency, duration) and "
        "smaller required sample sizes\u2014offering a pragmatic initial "
        "regulatory pathway.{27}"
    ),
    (
        "Limitations of the proposed approach include inter-individual "
        "variability in PBM response (responders vs. non-responders),{19} "
        "potential NO store depletion during prolonged use, and the "
        "possibility that baroreflex compensation may attenuate chronic "
        "blood pressure reduction in some patient populations. These "
        "concerns can be addressed through closed-loop feedback control "
        "(adjusting PBM parameters in response to continuous blood "
        "pressure monitoring) and pulsed irradiation protocols that allow "
        "NO store replenishment between active periods."
    ),
    (
        "In conclusion, the convergence of evidence from invasive "
        "arteriovenous shunt trials, PBM-induced vasodilation studies, "
        "and AVA physiology research supports the hypothesis that "
        "wearable photobiomodulation targeting glabrous skin can produce "
        "clinically meaningful blood pressure reduction. This approach "
        "merits experimental validation as a potential drug-free "
        "alternative for hypertension management."
    ),
]

# Table 1: Comparison of approaches
TABLE1_TITLE = (
    "Table 1. Comparison of arteriovenous shunt-based blood pressure "
    "reduction strategies"
)
TABLE1_HEADERS = [
    "Approach", "Mechanism", "BP reduction\n(SBP, mmHg)",
    "Invasiveness", "Evidence level"
]
TABLE1_DATA = [
    [
        "ROX Coupler\n(iliac AV fistula)",
        "Surgically created\nfixed AV shunt",
        "\u221226.9",
        "Invasive\n(catheter)",
        "RCT (n=83){8}",
    ],
    [
        "AVACEN 100\n(palm heating +\nneg. pressure)",
        "Conductive heat +\nvacuum \u2192 AVA dilation",
        "Under\ninvestigation",
        "Non-invasive\n(stationary)",
        "Pilot studies{15}",
    ],
    [
        "Selective thermal\nstimulation\n(Diller, 2022)",
        "Spinal heating \u2192\nhypothalamic reflex\n\u2192 AVA dilation",
        "Not yet\nreported",
        "Non-invasive\n(bed-integrated)",
        "POC (n=10){28}",
    ],
    [
        "PBM-AVA device\n(proposed)",
        "PBM \u2192 NO release\n\u2192 AVA dilation\n\u2192 TPR\u2193",
        "Predicted:\n\u22125 to \u221210",
        "Non-invasive\n(wearable)",
        "Hypothesis",
    ],
]

REFERENCES = [
    "GBD 2019 Risk Factors Collaborators. Global burden of 87 risk factors in 204 countries and territories, 1990-2019. Lancet. 2020;396(10258):1223-1249.",
    "Burnier M, Egan BM. Adherence in hypertension. Circ Res. 2019;124(7):1124-1140.",
    "Roush GC, Ernst ME, Kostis JB, et al. Head-to-head comparisons of hydrochlorothiazide with indapamide and chlorthalidone: antihypertensive and metabolic effects. Hypertension. 2015;65(5):1041-1046.",
    "Hall JE. Guyton and Hall Textbook of Medical Physiology. 14th ed. Philadelphia: Elsevier; 2021.",
    "Walløe L. Arterio-venous anastomoses in the human skin and their role in temperature control. Temperature (Austin). 2016;3(1):92-103.",
    "Midttun M, Sejrsen P, Paaske WP. Blood flow rate in arteriovenous anastomoses from the finger in control subjects and in patients with Raynaud's phenomenon. Clin Physiol. 1996;16(3):275-284.",
    "Johnson JM, Minson CT, Kellogg DL Jr. Cutaneous vasodilator and vasoconstrictor mechanisms in temperature regulation. Compr Physiol. 2014;4(1):33-89.",
    "Lobo MD, Sobotka PA, Stanton A, et al. Central arteriovenous anastomosis for the treatment of patients with uncontrolled hypertension (the ROX CONTROL HTN study): a randomised controlled trial. Lancet. 2015;385(9978):1634-1641.",
    "Anders JJ, Lanzafame RJ, Arany PR. Low-level light/laser therapy versus photobiomodulation therapy. Photomed Laser Surg. 2015;33(4):183-184.",
    "Karu TI, Pyatibrat LV, Afanasyeva NI. Cellular effects of low power laser therapy can be mediated by nitric oxide. Lasers Surg Med. 2005;36(4):307-314.",
    "Moncada S, Higgs EA. The discovery of nitric oxide and its role in vascular biology. Br J Pharmacol. 2006;147(S1):S193-S201.",
    "Ribeiro BG, Alves AN, Dos Santos LA, et al. Nitric oxide storage levels modulate vasodilation and the hypotensive effect induced by photobiomodulation using an AlGaAs diode laser (660 nm). Lasers Med Sci. 2022;37(6):2551-2559.",
    "De Freitas Rodrigues A, Sordillo LA, Sordillo PP, et al. In vivo attenuation profile of 660 nm and 830 nm wavelengths on human skin. Lasers Med Sci. 2024;39(1):28.",
    "Lewington S, Clarke R, Qizilbash N, et al. Age-specific relevance of usual blood pressure to vascular mortality. Lancet. 2002;360(9349):1903-1913.",
    "AVACEN Medical. AVACEN 100 Treatment System. US Patent 8,679,170 B2. 2014.",
    "Lobo MD, Ott C, Gane EJ, et al. Central iliac arteriovenous anastomosis for uncontrolled hypertension: one-year results from the ROX CONTROL HTN trial. Hypertension. 2017;70(6):1099-1105.",
    "Faul JL, Galindo J, Engelsgjerd JS, et al. Creation of an iliac arteriovenous shunt lowers blood pressure in chronic obstructive pulmonary disease patients with hypertension. J Vasc Surg. 2014;59(4):1078-1083.",
    "Damasceno GM, Pfaff AMCS, Gomes RSA, et al. Photobiomodulation therapy in hypertension management\u2014Evidence from a systematic review and meta-analysis. J Clin Med. 2025;14(19):6716.",
    "Rogatkin DA, Dunaev AV. Stimulation of blood microcirculation at low level laser irradiation. Proc SPIE. 2014;9129:912922.",
    "Kräuchi K, Gompper B, Hauenstein D, et al. Diurnal blood pressure variations are associated with changes in distal-proximal skin temperature gradient. Chronobiol Int. 2012;29(9):1189-1197.",
    "Martinez-Nicolas A, Ortiz-Tudela E, Rol MA, et al. Wrist skin temperature, motor activity, and body position as determinants of the circadian pattern of blood pressure. Chronobiol Int. 2013;30(5):668-678.",
    "Gani F, Guzman RC, Ahmed SS, et al. Local passive heat for the treatment of hypertension in autonomic failure. J Am Heart Assoc. 2021;10(13):e018979.",
    "Ferraresi C, Huang YY, Hamblin MR. Photobiomodulation in human muscle tissue: an advantage in sports performance? J Biophotonics. 2016;9(11-12):1273-1299.",
    "Grassi G, Ram VS. Evidence for a critical role of the sympathetic nervous system in hypertension. J Am Soc Hypertens. 2016;10(5):457-466.",
    "Modesti PA. Season, temperature and blood pressure: a complex interaction. Eur J Intern Med. 2013;24(7):604-607.",
    "Barnett AG, Sans S, Salomaa V, et al. The effect of temperature on systolic blood pressure. Blood Press Monit. 2007;12(3):195-203.",
    "Flavahan NA. A vascular mechanistic approach to understanding Raynaud phenomenon. Nat Rev Rheumatol. 2015;11(3):146-158.",
    "Haghayegh S, Khoshnevis S, Smolensky MH, et al. Novel temperature-controlled sleep system to improve sleep: a proof-of-concept study. J Sleep Res. 2022;31(5):e13662.",
    "Diller KR, Khoshnevis S, Hemmen L. Thermoregulatory manipulation of systemic blood pressure. US Patent 11,229,548 B2. 2022.",
]


# ─────────────────────────────────────────────
# Helper functions
# ─────────────────────────────────────────────

def add_paragraph_with_refs(doc, text, style="Normal"):
    """Add paragraph with superscript citation numbers from {n} markers."""
    para = doc.add_paragraph(style=style)
    parts = re.split(r"(\{[^}]+\})", text)
    for part in parts:
        if part.startswith("{") and part.endswith("}"):
            run = para.add_run(part[1:-1])
            run.font.superscript = True
            run.font.size = Pt(8)
        else:
            run = para.add_run(part)
            run.font.size = Pt(11)
    return para


def set_table_style(table):
    """Apply consistent formatting to a table."""
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for row in table.rows:
        for cell in row.cells:
            for para in cell.paragraphs:
                para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                for run in para.runs:
                    run.font.size = Pt(9)


# ─────────────────────────────────────────────
# Build DOCX
# ─────────────────────────────────────────────

def create_docx():
    doc = Document()

    # Set default font
    style = doc.styles["Normal"]
    font = style.font
    font.name = "Times New Roman"
    font.size = Pt(11)

    # Title
    title_para = doc.add_paragraph()
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title_para.add_run(TITLE)
    title_run.bold = True
    title_run.font.size = Pt(14)

    # Authors
    author_para = doc.add_paragraph()
    author_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    author_run = author_para.add_run(AUTHORS)
    author_run.font.size = Pt(12)

    # Affiliation placeholder
    aff_para = doc.add_paragraph()
    aff_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    aff_run = aff_para.add_run("[Affiliation to be inserted]")
    aff_run.font.size = Pt(10)
    aff_run.italic = True

    # Keywords
    kw_para = doc.add_paragraph()
    kw_run = kw_para.add_run("Keywords: ")
    kw_run.bold = True
    kw_run.font.size = Pt(10)
    kw_text = kw_para.add_run("; ".join(KEYWORDS))
    kw_text.font.size = Pt(10)

    doc.add_paragraph()

    # Abstract
    abs_heading = doc.add_paragraph()
    abs_run = abs_heading.add_run("Abstract")
    abs_run.bold = True
    abs_run.font.size = Pt(12)
    abs_para = doc.add_paragraph()
    abs_text = abs_para.add_run(ABSTRACT)
    abs_text.font.size = Pt(11)

    doc.add_paragraph()

    # Introduction
    intro_heading = doc.add_paragraph()
    intro_run = intro_heading.add_run("1. Introduction")
    intro_run.bold = True
    intro_run.font.size = Pt(12)

    for para_text in INTRODUCTION:
        add_paragraph_with_refs(doc, para_text)

    # The Hypothesis
    hyp_heading = doc.add_paragraph()
    hyp_run = hyp_heading.add_run("2. The Hypothesis")
    hyp_run.bold = True
    hyp_run.font.size = Pt(12)

    for para_text in HYPOTHESIS:
        add_paragraph_with_refs(doc, para_text)

    # Figure 1 (inline)
    fig1_para = doc.add_paragraph()
    fig1_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    fig1_para.space_before = Pt(18)
    fig1_path = OUTPUT_DIR / "figure1_mechanism.png"
    if fig1_path.exists():
        fig1_para.add_run().add_picture(str(fig1_path), width=Inches(5.5))
    else:
        fig1_run = fig1_para.add_run("[Figure 1: Run create_figure1.py first]")
        fig1_run.italic = True
        fig1_run.font.size = Pt(10)

    fig1_cap = doc.add_paragraph()
    fig1_cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    fig1_cap_run = fig1_cap.add_run(
        "Fig. 1. Schematic of proposed causal chain: PBM irradiation "
        "(630\u2013850 nm) \u2192 photon absorption by CcO/nitrosyl-hemoglobin "
        "\u2192 NO release \u2192 sGC/cGMP activation \u2192 AVA smooth muscle "
        "relaxation \u2192 AVA dilation \u2192 TPR reduction \u2192 MAP decrease."
    )
    fig1_cap_run.font.size = Pt(9)
    fig1_cap_run.italic = True

    # Evaluation of the hypothesis
    eval_heading = doc.add_paragraph()
    eval_run = eval_heading.add_run("3. Evaluation of the Hypothesis")
    eval_run.bold = True
    eval_run.font.size = Pt(12)

    for para_text in EVALUATION:
        add_paragraph_with_refs(doc, para_text)

    # Table 1
    tbl_cap = doc.add_paragraph()
    tbl_cap.space_before = Pt(18)
    tbl_cap_run = tbl_cap.add_run(TABLE1_TITLE)
    tbl_cap_run.bold = True
    tbl_cap_run.font.size = Pt(10)

    table = doc.add_table(rows=1 + len(TABLE1_DATA), cols=len(TABLE1_HEADERS))
    table.style = "Table Grid"

    # Headers
    for i, header in enumerate(TABLE1_HEADERS):
        cell = table.rows[0].cells[i]
        cell.text = header
        for para in cell.paragraphs:
            for run in para.runs:
                run.bold = True

    # Data rows
    for row_idx, row_data in enumerate(TABLE1_DATA):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.rows[row_idx + 1].cells[col_idx]
            # Handle superscript refs in table cells
            parts = re.split(r"(\{[^}]+\})", cell_text)
            cell.text = ""
            para = cell.paragraphs[0]
            for part in parts:
                if part.startswith("{") and part.endswith("}"):
                    run = para.add_run(part[1:-1])
                    run.font.superscript = True
                    run.font.size = Pt(8)
                else:
                    run = para.add_run(part)
                    run.font.size = Pt(9)

    set_table_style(table)

    # Hypothesis testing
    test_heading = doc.add_paragraph()
    test_heading.space_before = Pt(12)
    test_run = test_heading.add_run("4. Hypothesis Testing")
    test_run.bold = True
    test_run.font.size = Pt(12)

    for para_text in TESTING:
        add_paragraph_with_refs(doc, para_text)

    # Consequences and discussion
    cons_heading = doc.add_paragraph()
    cons_run = cons_heading.add_run("5. Consequences of the Hypothesis and Discussion")
    cons_run.bold = True
    cons_run.font.size = Pt(12)

    for para_text in CONSEQUENCES:
        add_paragraph_with_refs(doc, para_text)

    # References
    doc.add_page_break()
    ref_heading = doc.add_paragraph()
    ref_run = ref_heading.add_run("References")
    ref_run.bold = True
    ref_run.font.size = Pt(12)

    for i, ref in enumerate(REFERENCES, 1):
        ref_para = doc.add_paragraph()
        num_run = ref_para.add_run(f"{i}. ")
        num_run.font.size = Pt(9)
        ref_text_run = ref_para.add_run(ref)
        ref_text_run.font.size = Pt(9)

    # Save
    output_path = OUTPUT_DIR / "manuscript_medical_hypotheses.docx"
    doc.save(str(output_path))
    print(f"Manuscript saved: {output_path}")
    return output_path


# ─────────────────────────────────────────────
# Build PPTX (Figure 1 + Table 1)
# ─────────────────────────────────────────────

def create_pptx():
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    # Slide 1: Figure 1 (mechanism diagram - text-based schematic)
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(0.8)
    )
    tf = title_box.text_frame
    tf.text = "Figure 1. Proposed mechanism of PBM-induced AVA vasodilation"
    tf.paragraphs[0].font.size = PptxPt(18)
    tf.paragraphs[0].font.bold = True

    # Mechanism flow (text-based)
    flow_box = slide.shapes.add_textbox(
        PptxInches(1), PptxInches(1.5), PptxInches(11), PptxInches(4.5)
    )
    tf = flow_box.text_frame
    tf.word_wrap = True
    lines = [
        "PBM Irradiation (630\u2013850 nm, 20\u2013100 mW/cm\u00b2)",
        "        \u2193",
        "Photon absorption by CcO / nitrosyl-Hb in dermal vessels",
        "        \u2193",
        "NO release from endothelial stores (S-nitrosothiols)",
        "        \u2193",
        "sGC activation \u2192 cGMP\u2191 \u2192 PKG activation",
        "        \u2193",
        "AVA smooth muscle relaxation \u2192 AVA DILATION",
        "        \u2193",
        "Parallel low-resistance pathways added to systemic circulation",
        "        \u2193",
        "TPR reduction (estimated 5\u201310%)",
        "        \u2193",
        "\u0394MAP = \u22124.5 to \u22129 mmHg (clinically significant)",
    ]
    for i, line in enumerate(lines):
        if i == 0:
            tf.paragraphs[0].text = line
            tf.paragraphs[0].font.size = PptxPt(14)
        else:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = PptxPt(14)
            if "\u2193" in line:
                from pptx.dml.color import RGBColor as PptxRGB
                p.font.color.rgb = PptxRGB(0x00, 0x70, 0xC0)
            if "DILATION" in line or "clinically significant" in line:
                p.font.bold = True

    # Caption
    cap_box = slide.shapes.add_textbox(
        PptxInches(0.5), PptxInches(6.3), PptxInches(12), PptxInches(1)
    )
    tf = cap_box.text_frame
    tf.word_wrap = True
    tf.text = (
        "Fig. 1. Schematic of proposed causal chain from photobiomodulation "
        "to blood pressure reduction via AVA vasodilation. CcO = cytochrome c "
        "oxidase; Hb = hemoglobin; NO = nitric oxide; sGC = soluble guanylate "
        "cyclase; cGMP = cyclic guanosine monophosphate; PKG = protein kinase G; "
        "AVA = arteriovenous anastomosis; TPR = total peripheral resistance; "
        "MAP = mean arterial pressure."
    )
    tf.paragraphs[0].font.size = PptxPt(10)
    tf.paragraphs[0].font.italic = True

    # Slide 2: Table 1
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(
        PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(0.8)
    )
    tf = title_box.text_frame
    tf.text = TABLE1_TITLE
    tf.paragraphs[0].font.size = PptxPt(16)
    tf.paragraphs[0].font.bold = True

    # Table
    rows = len(TABLE1_DATA) + 1
    cols = len(TABLE1_HEADERS)
    tbl = slide.shapes.add_table(
        rows, cols, PptxInches(0.5), PptxInches(1.2),
        PptxInches(12.3), PptxInches(5)
    ).table

    # Headers
    for i, h in enumerate(TABLE1_HEADERS):
        cell = tbl.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = PptxPt(11)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

    # Data
    for r_idx, row_data in enumerate(TABLE1_DATA):
        for c_idx, val in enumerate(row_data):
            cell = tbl.cell(r_idx + 1, c_idx)
            # Strip citation markers for pptx
            clean = re.sub(r"\{[^}]+\}", "", val)
            cell.text = clean
            for p in cell.text_frame.paragraphs:
                p.font.size = PptxPt(10)
                p.alignment = PP_ALIGN.CENTER

    output_path = OUTPUT_DIR / "figures_tables.pptx"
    prs.save(str(output_path))
    print(f"Figures/Tables PPTX saved: {output_path}")
    return output_path


if __name__ == "__main__":
    create_docx()
    create_pptx()
    print("Done.")
