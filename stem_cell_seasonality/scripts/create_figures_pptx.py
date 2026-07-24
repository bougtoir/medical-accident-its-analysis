#!/usr/bin/env python3
"""
Generate editable PPTX figures for Cell Stem Cell Perspective:
"The Invisible Variables"

Figure 1: The Environmental Iceberg of Stem Cell Culture
Figure 2: Research Roadmap (Evidence matrix + Three-phase strategy)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "output")
os.makedirs(OUTPUT_DIR, exist_ok=True)

# Widescreen slide dimensions
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_textbox(slide, left, top, width, height, text, font_size=12,
                bold=False, italic=False, color=RGBColor(0, 0, 0),
                alignment=PP_ALIGN.LEFT, font_name='Arial'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=10, font_color=RGBColor(0, 0, 0), bold=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = 'Arial'
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(2)
        tf.paragraphs[0].space_after = Pt(2)
    return shape


def add_rect(slide, left, top, width, height, fill_color, text="",
             font_size=10, font_color=RGBColor(255, 255, 255), bold=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(4)
        tf.margin_right = Pt(4)
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = 'Arial'
        p.alignment = PP_ALIGN.CENTER
    return shape


def create_figure1(prs):
    """Figure 1: The Environmental Iceberg of Stem Cell Culture."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                "Figure 1. The Environmental Iceberg of Stem Cell Culture",
                font_size=20, bold=True)

    # ── "Waterline" concept ──
    # Above waterline: controlled variables
    waterline_y = Inches(3.2)

    # Sky / controlled area background
    add_rect(slide, Inches(1.0), Inches(1.0), Inches(5.0), Inches(2.2),
             RGBColor(220, 240, 255))  # light blue sky

    add_textbox(slide, Inches(1.2), Inches(1.1), Inches(4.6), Inches(0.4),
                "CONTROLLED (standard practice)", font_size=14, bold=True,
                color=RGBColor(0, 100, 0))

    # Controlled variable boxes
    ctrl_vars = [
        ("Temperature\n37°C ± 0.1°C", RGBColor(144, 238, 144)),
        ("CO₂\n5% ± 0.1%", RGBColor(144, 238, 144)),
    ]
    for i, (txt, col) in enumerate(ctrl_vars):
        add_rounded_rect(slide, Inches(1.5 + i * 2.2), Inches(1.7),
                        Inches(2.0), Inches(1.0), col, txt,
                        font_size=14, bold=True)

    # Waterline
    add_rect(slide, Inches(0.8), waterline_y, Inches(5.4), Inches(0.06),
             RGBColor(0, 80, 160))
    add_textbox(slide, Inches(0.5), Inches(3.0), Inches(1.5), Inches(0.3),
                "— waterline —", font_size=9, italic=True,
                color=RGBColor(0, 80, 160))

    # Below waterline: uncontrolled area
    add_rect(slide, Inches(1.0), Inches(3.3), Inches(5.0), Inches(3.8),
             RGBColor(200, 220, 240))  # deeper blue

    add_textbox(slide, Inches(1.2), Inches(3.35), Inches(4.6), Inches(0.4),
                "UNCONTROLLED (invisible variables)", font_size=14, bold=True,
                color=RGBColor(180, 0, 0))

    # Uncontrolled variable boxes
    unctrl_vars = [
        ("Humidity\n(seasonal: summer ↑)", RGBColor(255, 200, 200)),
        ("Ambient Light\n(photoperiod)", RGBColor(255, 220, 180)),
        ("VOCs\n(T/RH-dependent)", RGBColor(255, 200, 200)),
        ("ELF-EMF\n(50/60 Hz, HVAC)", RGBColor(255, 220, 180)),
        ("Barometric Pressure\n(weather systems)", RGBColor(255, 200, 200)),
        ("Background Radiation\n(cosmic rays)", RGBColor(255, 220, 180)),
        ("Water Quality\n(endotoxin, TOC)", RGBColor(255, 200, 200)),
        ("Vibration\n(building, traffic)", RGBColor(255, 220, 180)),
    ]

    for i, (txt, col) in enumerate(unctrl_vars):
        row = i // 2
        col_idx = i % 2
        add_rounded_rect(
            slide,
            Inches(1.3 + col_idx * 2.4),
            Inches(3.9 + row * 0.85),
            Inches(2.2), Inches(0.75),
            col, txt, font_size=10, bold=False
        )

    # ── Right panel: Seasonal patterns by hemisphere ──
    add_textbox(slide, Inches(7.0), Inches(0.2), Inches(5.5), Inches(0.5),
                "Seasonal Variation Patterns", font_size=18, bold=True)

    # Table header
    headers = ["Variable", "NH Summer\n(Jun–Aug)", "SH Summer\n(Dec–Feb)", "Phase"]
    col_widths = [Inches(2.2), Inches(1.4), Inches(1.4), Inches(1.2)]
    col_starts = [Inches(7.0)]
    for w in col_widths[:-1]:
        col_starts.append(col_starts[-1] + w)

    row_h = Inches(0.4)
    table_top = Inches(1.0)

    # Header row
    for j, (hdr, w) in enumerate(zip(headers, col_widths)):
        add_rect(slide, col_starts[j], table_top, w, row_h,
                RGBColor(0, 80, 130), hdr, font_size=9, bold=True)

    # Data rows
    data_rows = [
        ["Humidity", "↑ High", "↑ High", "Inverted"],
        ["Photoperiod", "Long days", "Long days", "Inverted"],
        ["VOCs (off-gassing)", "↑ (warm)", "↑ (warm)", "Inverted"],
        ["ELF-EMF (HVAC)", "↑ cooling", "↑ heating*", "Inverted"],
        ["Barometric press.", "Stable", "Stable", "Inverted"],
        ["IVF outcome", "↑ Better", "↑ Better", "Inverted"],
        ["Melatonin", "↓ Less", "↓ Less", "Inverted"],
        ["Contamination risk", "↑ Higher", "↑ Higher", "Inverted"],
        ["Geomagnetic storms", "Spring/Fall", "Spring/Fall", "Synchronous"],
        ["Cosmic ray flux", "Winter ↑", "Winter ↑", "Inverted"],
        ["Solar UV (11yr)", "Max at solar max", "Same", "Synchronous"],
    ]

    colors_alt = [RGBColor(235, 245, 255), RGBColor(255, 255, 255)]
    phase_colors = {
        "Inverted": RGBColor(200, 230, 200),
        "Synchronous": RGBColor(255, 230, 200),
    }

    for i, row_data in enumerate(data_rows):
        y = table_top + row_h * (i + 1)
        bg = colors_alt[i % 2]
        for j, (cell, w) in enumerate(zip(row_data, col_widths)):
            cell_bg = bg
            if j == 3:  # Phase column
                cell_bg = phase_colors.get(cell, bg)
            add_rect(slide, col_starts[j], y, w, row_h, cell_bg, cell,
                    font_size=8, font_color=RGBColor(0, 0, 0))

    # Legend
    add_textbox(slide, Inches(7.0), Inches(6.0), Inches(6.0), Inches(0.4),
                "Green = hemisphere-inverted (solar-driven); "
                "Orange = hemisphere-synchronous (geomagnetic)",
                font_size=9, italic=True, color=RGBColor(80, 80, 80))

    # Caption
    add_textbox(slide, Inches(0.5), Inches(7.1), Inches(12.5), Inches(0.4),
                "Figure 1. Left: Controlled vs. uncontrolled environmental parameters in standard PSC facilities. "
                "Right: Expected seasonal patterns by hemisphere.",
                font_size=10, italic=True, color=RGBColor(60, 60, 60))

    # Note
    add_textbox(slide, Inches(7.0), Inches(6.3), Inches(6.0), Inches(0.3),
                "*SH HVAC patterns depend on climate zone",
                font_size=8, italic=True, color=RGBColor(120, 120, 120))


def create_figure2(prs):
    """Figure 2: Research Roadmap — Evidence matrix + Three-phase strategy."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
                "Figure 2. Research Roadmap for Environmental Profiling of PSC Facilities",
                font_size=20, bold=True)

    # ── Panel A: Evidence Matrix ──
    add_textbox(slide, Inches(0.5), Inches(0.8), Inches(6.0), Inches(0.4),
                "A  Evidence Matrix", font_size=16, bold=True,
                color=RGBColor(0, 80, 130))

    # Y axis: Evidence level
    evidence_levels = ["Direct PSC\nevidence", "Indirect\n(IVF/culture)", "Theoretical\nonly"]
    ev_colors = [RGBColor(0, 130, 80), RGBColor(200, 150, 0), RGBColor(180, 80, 80)]
    ev_y_starts = [Inches(1.6), Inches(3.0), Inches(4.4)]

    for i, (lbl, col) in enumerate(zip(evidence_levels, ev_colors)):
        add_textbox(slide, Inches(0.3), ev_y_starts[i], Inches(1.0), Inches(1.0),
                    lbl, font_size=9, bold=True, color=col,
                    alignment=PP_ALIGN.CENTER)

    # X axis: Hemisphere dependence
    add_textbox(slide, Inches(1.8), Inches(5.5), Inches(2.5), Inches(0.3),
                "Hemisphere-inverted\n(solar-driven)", font_size=9, bold=True,
                color=RGBColor(0, 80, 130), alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(4.3), Inches(5.5), Inches(2.5), Inches(0.3),
                "Hemisphere-synchronous\n(geomagnetic)", font_size=9, bold=True,
                color=RGBColor(180, 80, 0), alignment=PP_ALIGN.CENTER)

    # Variable bubbles positioned in the matrix
    # (left, top, label, evidence_level, hemisphere_type)
    bubbles = [
        # Direct PSC evidence, hemisphere-inverted
        (Inches(2.0), Inches(1.7), "CRY1 →\nreprogramming", RGBColor(0, 150, 80)),
        (Inches(3.5), Inches(1.8), "50Hz EMF →\nES cell genes", RGBColor(0, 150, 80)),

        # Indirect (IVF/culture), hemisphere-inverted
        (Inches(1.5), Inches(3.1), "Humidity →\ncontamination", RGBColor(220, 170, 0)),
        (Inches(3.0), Inches(3.2), "IVF summer\nadvantage", RGBColor(220, 170, 0)),
        (Inches(4.5), Inches(3.1), "VOC →\nembryo quality", RGBColor(220, 170, 0)),

        # Indirect, hemisphere-synchronous
        (Inches(5.2), Inches(3.3), "Cosmic ray →\nproliferation", RGBColor(220, 170, 0)),

        # Theoretical, hemisphere-inverted
        (Inches(1.8), Inches(4.5), "Barometric\npressure", RGBColor(200, 100, 100)),
        (Inches(3.3), Inches(4.6), "Melatonin →\ndonor cells", RGBColor(200, 100, 100)),

        # Theoretical, hemisphere-synchronous
        (Inches(5.0), Inches(4.5), "Geomagnetic →\nCRY/NF-κB", RGBColor(200, 100, 100)),
    ]

    for (x, y, txt, col) in bubbles:
        add_rounded_rect(slide, x, y, Inches(1.4), Inches(0.8), col, txt,
                        font_size=8, font_color=RGBColor(255, 255, 255), bold=True)

    # ── Panel B: Three-phase strategy ──
    add_textbox(slide, Inches(7.2), Inches(0.8), Inches(5.5), Inches(0.4),
                "B  Three-Phase Investigation Strategy", font_size=16, bold=True,
                color=RGBColor(0, 80, 130))

    phases = [
        {
            "title": "Phase I: Passive Monitoring",
            "color": RGBColor(0, 130, 180),
            "items": [
                "Deploy IoT environmental sensors",
                "Record: humidity, lux, EMF, pressure,\nVOC, vibration (1-min resolution)",
                "Log differentiation outcomes in parallel",
                "Duration: ≥12 months (full seasonal cycle)",
                "Cost: ~$5,000 per laboratory",
            ]
        },
        {
            "title": "Phase II: Retrospective Mining",
            "color": RGBColor(0, 150, 100),
            "items": [
                "IVF registries: HFEA, ANZARD, CDC ART",
                "GEO/ArrayExpress: submission date × QC",
                "iPSC banks: RIKEN BRC, Coriell, EBiSC",
                "Model: outcome ~ month + latitude + weather",
                "Purely computational — no wet lab needed",
            ]
        },
        {
            "title": "Phase III: Controlled Intervention",
            "color": RGBColor(180, 100, 0),
            "items": [
                "Humidity-controlled vs standard incubators",
                "Light-tight hoods vs standard BSC",
                "Mu-metal EMF shielding around incubators",
                "Test each variable × seasonal baseline",
                "Multi-center, cross-hemisphere design",
            ]
        },
    ]

    phase_top = Inches(1.4)
    phase_w = Inches(5.5)
    phase_h = Inches(1.7)
    phase_gap = Inches(0.15)

    for i, phase in enumerate(phases):
        y = phase_top + i * (phase_h + phase_gap)

        # Phase header
        add_rect(slide, Inches(7.2), y, phase_w, Inches(0.35),
                phase["color"], phase["title"],
                font_size=12, bold=True)

        # Phase items
        item_y = y + Inches(0.4)
        for j, item in enumerate(phase["items"]):
            add_textbox(slide, Inches(7.5), item_y + Inches(j * 0.25),
                       Inches(5.0), Inches(0.25),
                       f"• {item}", font_size=9)

    # Arrow between phases
    add_textbox(slide, Inches(9.5), Inches(3.15), Inches(0.5), Inches(0.3),
                "↓", font_size=18, bold=True, color=RGBColor(100, 100, 100),
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(9.5), Inches(5.0), Inches(0.5), Inches(0.3),
                "↓", font_size=18, bold=True, color=RGBColor(100, 100, 100),
                alignment=PP_ALIGN.CENTER)

    # Caption
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.5), Inches(0.5),
                "Figure 2. (A) Evidence matrix classifying environmental variables by hemisphere-dependence "
                "and evidence level. (B) Three-phase strategy from passive monitoring to controlled intervention.",
                font_size=10, italic=True, color=RGBColor(60, 60, 60))


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    create_figure1(prs)
    create_figure2(prs)

    out_path = os.path.join(OUTPUT_DIR, "CellStemCell_Figures_InvisibleVariables.pptx")
    prs.save(out_path)
    print(f"Figures PPTX saved to: {out_path}")


if __name__ == "__main__":
    main()
