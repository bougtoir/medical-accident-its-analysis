#!/usr/bin/env python3
"""
Generate Methods (Statistical Analysis) and Results sections
for the cesarean section bleeding & transfusion manuscript.

Also creates an editable PPTX with all figures.
"""

import json, re
from pathlib import Path

import pandas as pd
import numpy as np
from docx import Document
from docx.shared import Inches, Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt, Emu
from pptx.enum.text import PP_ALIGN

BASE = Path(__file__).resolve().parent
FIG_DIR = BASE / "figures"
TBL_DIR = BASE / "tables"

# Load summary stats
with open(BASE / "summary_stats.json") as f:
    S = json.load(f)

# Load regression results
bl_reg = pd.read_csv(TBL_DIR / "regression_blood_loss.csv")
or_tx = pd.read_csv(TBL_DIR / "logistic_regression_transfusion.csv")
or_mh = pd.read_csv(TBL_DIR / "logistic_regression_massive_hemorrhage.csv")
table1_cont = pd.read_csv(TBL_DIR / "table1_continuous.csv")
table1_cat = pd.read_csv(TBL_DIR / "table1_categorical.csv")
yearly = pd.read_csv(TBL_DIR / "yearly_trends.csv")

# Variable label mapping
VAR_LABELS = {
    "年齢(歳)": "age",
    "BMI": "BMI",
    "GA_weeks": "gestational age",
    "emergency": "emergency cesarean section",
    "prior_cs": "prior cesarean delivery",
    "HDP": "hypertensive disorders of pregnancy",
    "placenta_previa": "placenta previa or low-lying placenta",
    "epidural": "epidural anesthesia",
    "手術時間_min": "surgical duration",
    "hypotension_count": "number of hypotension episodes",
}

# ============================================================
# MANUSCRIPT DOCX
# ============================================================

doc = Document()

# Styles
style = doc.styles["Normal"]
font = style.font
font.name = "Times New Roman"
font.size = Pt(12)
style.paragraph_format.line_spacing = 1.5
style.paragraph_format.space_after = Pt(6)

def add_heading(text, level=1):
    h = doc.add_heading(text, level=level)
    for run in h.runs:
        run.font.name = "Times New Roman"
        run.font.color.rgb = RGBColor(0, 0, 0)
    return h

def add_para(text, bold=False, italic=False):
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.bold = bold
    run.italic = italic
    run.font.name = "Times New Roman"
    run.font.size = Pt(12)
    return p

def add_superscript(paragraph, text):
    run = paragraph.add_run(text)
    run.font.superscript = True
    run.font.name = "Times New Roman"
    run.font.size = Pt(12)
    return run

def fmt_p(val, digits=3):
    if val < 0.001:
        return "P < 0.001"
    return f"P = {val:.{digits}f}"

def fmt_or(row):
    return f"{row['OR']:.2f} (95% CI: {row['95% CI lower']:.2f}\u2013{row['95% CI upper']:.2f})"

def fmt_pct_change(row):
    return f"{row['% change']:.1f}% (95% CI: {(np.exp(row['95% CI lower']) - 1)*100:.1f}\u2013{(np.exp(row['95% CI upper']) - 1)*100:.1f}%)"

# ============================================================
# METHODS - Statistical Analysis
# ============================================================
add_heading("Methods", level=1)
add_heading("Statistical Analysis", level=2)

methods_text = (
    "Continuous variables were described as mean \u00b1 standard deviation (SD) or "
    "median with interquartile range (IQR), depending on the distribution. "
    "Categorical variables were expressed as frequencies and percentages. "
    "The Shapiro\u2013Wilk test was used to evaluate the normality of continuous variables."
)
add_para(methods_text)

methods_text2 = (
    "Univariate comparisons between groups were performed using the "
    "Mann\u2013Whitney U test for continuous variables and the chi-squared test "
    "for categorical variables. "
    "Spearman rank correlation coefficients were calculated to assess monotonic "
    "associations between continuous predictors and estimated blood loss (EBL). "
    "The Kruskal\u2013Wallis test was used for comparisons across three or more groups."
)
add_para(methods_text2)

methods_text3 = (
    "To identify independent predictors of EBL, multivariable linear regression "
    "was performed using log-transformed EBL as the dependent variable. "
    "Coefficients were back-transformed and expressed as percentage changes in EBL "
    "with 95% confidence intervals (CIs). "
    "Multivariable logistic regression was used to identify independent risk factors "
    "for two binary outcomes: (1) intraoperative transfusion and (2) massive "
    "hemorrhage, defined as EBL \u2265 1500 mL. Results were expressed as odds ratios (ORs) "
    "with 95% CIs. The following covariates were included in all multivariable models: "
    "age, body mass index (BMI), gestational age, emergency vs. elective surgery, "
    "prior cesarean delivery, hypertensive disorders of pregnancy (HDP), "
    "placenta previa or low-lying placenta, epidural anesthesia, and surgical duration. "
    "The number of intraoperative hypotension episodes was additionally included in "
    "the linear regression model for EBL."
)
add_para(methods_text3)

methods_text4 = (
    "Model discrimination was evaluated using the area under the receiver operating "
    "characteristic curve (AUC). "
    "Temporal trends in blood loss, transfusion rate, and massive hemorrhage rate "
    "were assessed using Pearson correlation with year and compared between "
    "the early (2014\u20132018) and late (2019\u20132024) periods using the "
    "Mann\u2013Whitney U test for continuous outcomes and the chi-squared test "
    "for proportions."
)
add_para(methods_text4)

methods_text5 = (
    "All tests were two-sided, and P < 0.05 was considered statistically significant. "
    "Statistical analyses were performed using Python 3.12 with SciPy 1.17, "
    "statsmodels 0.14, and scikit-learn 1.8."
)
add_para(methods_text5)

# ============================================================
# RESULTS
# ============================================================
add_heading("Results", level=1)

# --- Patient characteristics ---
add_heading("Patient Characteristics", level=2)

n_final = int(S["n_final"])
n_raw = int(S["total_raw"])
n_excluded = int(S["n_excluded"])

results_p1 = (
    f"A total of {n_raw:,} singleton cesarean sections were identified during the "
    f"study period ({S['year_range']}). After excluding {n_excluded} cases "
    "(130 for general anesthesia, 10 for systolic blood pressure < 90 mmHg at admission, "
    "5 for triplet pregnancy, 3 for intrauterine fetal death, 3 for vaginal delivery or "
    "forceps delivery, 3 for cardiac arrest or death, and 4 for other reasons including "
    "incomplete records), {n_final:,} patients were included in the final analysis (Table 1)."
)
add_para(results_p1)

results_p2 = (
    f"The mean patient age was {S['age_mean']:.1f} \u00b1 {S['age_sd']:.1f} years. "
    f"The median gestational age was {S['ga_median']:.1f} weeks "
    f"(IQR: {S['ga_iqr_low']:.1f}\u2013{S['ga_iqr_high']:.1f}). "
    f"Emergency cesarean sections accounted for {float(S['pct_emergency']):.1f}% (n = {int(S['n_emergency']):,}) "
    f"of the cohort."
)
add_para(results_p2)

# Table 1 - Patient characteristics
add_para("Table 1. Patient characteristics", bold=True)
table = doc.add_table(rows=1, cols=5)
table.style = "Table Grid"
table.alignment = WD_TABLE_ALIGNMENT.CENTER
hdr = table.rows[0].cells
hdr[0].text = "Variable"
hdr[1].text = "n"
hdr[2].text = "Mean \u00b1 SD or n (%)"
hdr[3].text = "Median [IQR]"
hdr[4].text = "Range"
for _, row in table1_cont.iterrows():
    r = table.add_row().cells
    r[0].text = str(row["Variable"])
    r[1].text = str(row["n"])
    r[2].text = f"{row['Mean']} \u00b1 {row['SD']}"
    r[3].text = f"{row['Median']} [{row['IQR']}]"
    r[4].text = f"{row['Min']}\u2013{row['Max']}"
for _, row in table1_cat.iterrows():
    r = table.add_row().cells
    r[0].text = str(row["Variable"])
    r[1].text = str(row["n"])
    r[2].text = f"{row['n']} ({row['Percentage']})"
    r[3].text = ""
    r[4].text = ""

# Format table fonts
for r in table.rows:
    for cell in r.cells:
        for para in cell.paragraphs:
            for run in para.runs:
                run.font.name = "Times New Roman"
                run.font.size = Pt(10)

doc.add_paragraph()  # spacing

# --- Blood loss ---
add_heading("Estimated Blood Loss", level=2)

bl_mean = float(S["bl_mean"])
bl_sd = float(S["bl_sd"])
bl_median = float(S["bl_median"])
bl_iqr_low = float(S["bl_iqr_low"])
bl_iqr_high = float(S["bl_iqr_high"])
n_bl = int(S["n_bl_available"])

results_bl1 = (
    f"Blood loss data were available for {n_bl:,} patients ({100*n_bl/n_final:.1f}%). "
    f"The mean EBL was {bl_mean:.0f} \u00b1 {bl_sd:.0f} mL, with a median of "
    f"{bl_median:.0f} mL (IQR: {bl_iqr_low:.0f}\u2013{bl_iqr_high:.0f} mL; "
    f"range: {float(S['bl_range_min']):.0f}\u2013{float(S['bl_range_max']):.0f} mL). "
    f"The distribution was right-skewed (Fig. 1A)."
)
add_para(results_bl1)

# Insert Figure 1
doc.add_paragraph()
doc.add_picture(str(FIG_DIR / "fig1_blood_loss_distribution.png"), width=Inches(6.0))
p_cap = doc.add_paragraph()
p_cap.alignment = WD_ALIGN_PARAGRAPH.LEFT
run = p_cap.add_run(
    "Figure 1. Distribution of estimated blood loss during cesarean section. "
    "(A) Histogram with median (dashed red) and massive hemorrhage threshold "
    "(dotted orange, 1500 mL). (B) Box plots by surgical urgency."
)
run.font.name = "Times New Roman"
run.font.size = Pt(10)
run.italic = True
doc.add_paragraph()

# Subgroup comparisons
results_bl2 = (
    "Elective cesarean sections had significantly higher median EBL compared with "
    "emergency cases (730 mL vs. 570 mL; P < 0.001) (Fig. 1B). "
    "Patients with placenta previa or low-lying placenta had substantially higher "
    "median EBL (1188 mL vs. 621 mL; P < 0.001). "
    "Epidural anesthesia was associated with higher median EBL (730 mL vs. 620 mL; P < 0.001). "
    "No significant difference was observed between patients with and without "
    "prior cesarean delivery (P = 0.95) (Fig. 5)."
)
add_para(results_bl2)

# Insert Figure 5
doc.add_paragraph()
doc.add_picture(str(FIG_DIR / "fig5_subgroup_violins.png"), width=Inches(6.0))
p_cap = doc.add_paragraph()
p_cap.alignment = WD_ALIGN_PARAGRAPH.LEFT
run = p_cap.add_run(
    "Figure 5. Estimated blood loss by clinical subgroups (violin plots). "
    "Red dotted line indicates the massive hemorrhage threshold (1500 mL). "
    "Horizontal lines within violins represent the median and interquartile range."
)
run.font.name = "Times New Roman"
run.font.size = Pt(10)
run.italic = True
doc.add_paragraph()

# Multivariable - blood loss
add_heading("Multivariable Analysis of Blood Loss", level=2)

# Sort by significance
bl_sig = bl_reg[bl_reg["P-value"] < 0.05].sort_values("P-value")
bl_nonsig = bl_reg[bl_reg["P-value"] >= 0.05]

results_mv1 = (
    f"In multivariable linear regression (n = {int(S['ols_n']):,}; "
    f"adjusted R\u00b2 = {float(S['ols_adj_r2']):.3f}), "
    "the following factors were independently associated with increased EBL: "
)

sig_parts = []
for _, row in bl_sig.iterrows():
    label = VAR_LABELS.get(row["Variable"], row["Variable"])
    pct = row["% change"]
    direction = "increased" if pct > 0 else "decreased"
    sig_parts.append(
        f"{label} ({direction} EBL by {abs(pct):.1f}%; {fmt_p(row['P-value'])})"
    )

results_mv1 += "; ".join(sig_parts) + "."

if len(bl_nonsig) > 0:
    ns_labels = [VAR_LABELS.get(r["Variable"], r["Variable"]) for _, r in bl_nonsig.iterrows()]
    results_mv1 += f" {', '.join(ns_labels).capitalize()} were not significantly associated with EBL."

add_para(results_mv1)

add_para(
    "Table 2 and Figure 3 show the full regression results.",
)

# Table 2 - Regression
add_para("Table 2. Multivariable linear regression: Factors associated with estimated blood loss (log-transformed)", bold=True)
table2 = doc.add_table(rows=1, cols=4)
table2.style = "Table Grid"
table2.alignment = WD_TABLE_ALIGNMENT.CENTER
hdr2 = table2.rows[0].cells
hdr2[0].text = "Variable"
hdr2[1].text = "% Change in EBL"
hdr2[2].text = "95% CI"
hdr2[3].text = "P-value"
for _, row in bl_reg.sort_values("P-value").iterrows():
    r = table2.add_row().cells
    r[0].text = VAR_LABELS.get(row["Variable"], row["Variable"])
    r[1].text = f"{row['% change']:.1f}%"
    ci_lo = (np.exp(row["95% CI lower"]) - 1) * 100
    ci_hi = (np.exp(row["95% CI upper"]) - 1) * 100
    r[2].text = f"{ci_lo:.1f} to {ci_hi:.1f}"
    r[3].text = fmt_p(row["P-value"])

for r in table2.rows:
    for cell in r.cells:
        for para in cell.paragraphs:
            for run in para.runs:
                run.font.name = "Times New Roman"
                run.font.size = Pt(10)

doc.add_paragraph()

# Insert Figure 3
doc.add_picture(str(FIG_DIR / "fig3_forest_blood_loss.png"), width=Inches(6.0))
p_cap = doc.add_paragraph()
p_cap.alignment = WD_ALIGN_PARAGRAPH.LEFT
run = p_cap.add_run(
    "Figure 3. Forest plot showing the percentage change in estimated blood loss "
    "and 95% confidence intervals from multivariable linear regression (log-transformed EBL)."
)
run.font.name = "Times New Roman"
run.font.size = Pt(10)
run.italic = True
doc.add_paragraph()

# --- Transfusion ---
add_heading("Transfusion", level=2)

n_tx = int(S["n_transfusion"])
pct_tx = float(S["pct_transfusion"])

results_tx1 = (
    f"Intraoperative transfusion was administered in {n_tx} patients ({pct_tx:.1f}%). "
    "Patients who received transfusion were older (median 35 vs. 34 years; P < 0.001), "
    "had earlier gestational age (median 37.1 vs. 38.1 weeks; P < 0.001), "
    "received more intravenous fluids (median 1425 vs. 990 mL; P < 0.001), "
    "and had higher oxytocin doses (median 15 vs. 10 U; P < 0.001) (Table 1)."
)
add_para(results_tx1)

# Logistic regression transfusion
results_tx2 = (
    f"Multivariable logistic regression (n = {int(S['logit_tx_n']):,}; "
    f"{int(S['logit_tx_events'])} events) identified placenta previa as the strongest "
    "independent risk factor for transfusion "
)

pp_row = or_tx[or_tx["Variable"] == "placenta_previa"].iloc[0]
results_tx2 += (
    f"(OR {fmt_or(pp_row)}; {fmt_p(pp_row['P-value'])}). "
    "Additional significant predictors included "
)

tx_sig = or_tx[(or_tx["P-value"] < 0.05) & (or_tx["Variable"] != "placenta_previa")].sort_values("P-value")
tx_sig_parts = []
for _, row in tx_sig.iterrows():
    label = VAR_LABELS.get(row["Variable"], row["Variable"])
    tx_sig_parts.append(f"{label} (OR {fmt_or(row)}; {fmt_p(row['P-value'])})")

results_tx2 += ", ".join(tx_sig_parts) + ". "
results_tx2 += f"The model achieved excellent discrimination (AUC = {float(S['auc_transfusion']):.3f}) (Fig. 6A)."

add_para(results_tx2)

# Table 3 - Logistic regression transfusion
add_para("Table 3. Multivariable logistic regression: Factors associated with intraoperative transfusion", bold=True)
table3 = doc.add_table(rows=1, cols=4)
table3.style = "Table Grid"
table3.alignment = WD_TABLE_ALIGNMENT.CENTER
hdr3 = table3.rows[0].cells
hdr3[0].text = "Variable"
hdr3[1].text = "Odds Ratio"
hdr3[2].text = "95% CI"
hdr3[3].text = "P-value"
for _, row in or_tx.sort_values("P-value").iterrows():
    r = table3.add_row().cells
    r[0].text = VAR_LABELS.get(row["Variable"], row["Variable"])
    r[1].text = f"{row['OR']:.2f}"
    r[2].text = f"{row['95% CI lower']:.2f}\u2013{row['95% CI upper']:.2f}"
    r[3].text = fmt_p(row["P-value"])

for r in table3.rows:
    for cell in r.cells:
        for para in cell.paragraphs:
            for run in para.runs:
                run.font.name = "Times New Roman"
                run.font.size = Pt(10)

doc.add_paragraph()

# Insert Figure 4
doc.add_picture(str(FIG_DIR / "fig4_forest_transfusion.png"), width=Inches(6.0))
p_cap = doc.add_paragraph()
p_cap.alignment = WD_ALIGN_PARAGRAPH.LEFT
run = p_cap.add_run(
    "Figure 4. Forest plot showing odds ratios and 95% confidence intervals from "
    "multivariable logistic regression for intraoperative transfusion (log scale)."
)
run.font.name = "Times New Roman"
run.font.size = Pt(10)
run.italic = True
doc.add_paragraph()

# --- Massive hemorrhage ---
add_heading("Massive Hemorrhage", level=2)

n_mh = int(S["n_massive"])
pct_mh = float(S["pct_massive"])

results_mh1 = (
    f"Massive hemorrhage (EBL \u2265 1500 mL) occurred in {n_mh} patients ({pct_mh:.1f}%). "
    "In multivariable logistic regression "
    f"(n = {int(S['logit_mh_n']):,}; {int(S['logit_mh_events'])} events), "
    "the strongest independent risk factor was placenta previa "
)

mh_pp = or_mh[or_mh["Variable"] == "placenta_previa"].iloc[0]
results_mh1 += f"(OR {fmt_or(mh_pp)}; {fmt_p(mh_pp['P-value'])}). "

mh_sig = or_mh[(or_mh["P-value"] < 0.05) & (or_mh["Variable"] != "placenta_previa")].sort_values("P-value")
mh_parts = []
for _, row in mh_sig.iterrows():
    label = VAR_LABELS.get(row["Variable"], row["Variable"])
    direction = "risk factor" if row["OR"] > 1 else "protective factor"
    mh_parts.append(f"{label} was a {direction} (OR {fmt_or(row)}; {fmt_p(row['P-value'])})")
results_mh1 += ". ".join(mh_parts) + ". "

results_mh1 += f"Model discrimination was good (AUC = {float(S['auc_massive']):.3f}) (Fig. 6B)."
add_para(results_mh1)

# Table 4 - Massive hemorrhage
add_para("Table 4. Multivariable logistic regression: Factors associated with massive hemorrhage (EBL \u2265 1500 mL)", bold=True)
table4 = doc.add_table(rows=1, cols=4)
table4.style = "Table Grid"
table4.alignment = WD_TABLE_ALIGNMENT.CENTER
hdr4 = table4.rows[0].cells
hdr4[0].text = "Variable"
hdr4[1].text = "Odds Ratio"
hdr4[2].text = "95% CI"
hdr4[3].text = "P-value"
for _, row in or_mh.sort_values("P-value").iterrows():
    r = table4.add_row().cells
    r[0].text = VAR_LABELS.get(row["Variable"], row["Variable"])
    r[1].text = f"{row['OR']:.2f}"
    r[2].text = f"{row['95% CI lower']:.2f}\u2013{row['95% CI upper']:.2f}"
    r[3].text = fmt_p(row["P-value"])

for r in table4.rows:
    for cell in r.cells:
        for para in cell.paragraphs:
            for run in para.runs:
                run.font.name = "Times New Roman"
                run.font.size = Pt(10)

doc.add_paragraph()

# Insert Figure 6
doc.add_picture(str(FIG_DIR / "fig6_roc_curves.png"), width=Inches(6.0))
p_cap = doc.add_paragraph()
p_cap.alignment = WD_ALIGN_PARAGRAPH.LEFT
run = p_cap.add_run(
    "Figure 6. Receiver operating characteristic (ROC) curves for the multivariable "
    "logistic regression models. (A) Transfusion prediction. (B) Massive hemorrhage "
    "(\u22651500 mL) prediction."
)
run.font.name = "Times New Roman"
run.font.size = Pt(10)
run.italic = True
doc.add_paragraph()

# --- Temporal trends ---
add_heading("Temporal Trends", level=2)

results_trend = (
    "Significant temporal improvements were observed over the study period (Fig. 2). "
    "Median EBL decreased from 765 mL in 2014 to 560 mL in 2024 "
    "(Pearson r = \u22120.893; P < 0.001). "
    "The transfusion rate declined from 6.1% in 2014 to 1.4% in 2024 "
    "(r = \u22120.733; P = 0.010), and the massive hemorrhage rate declined from "
    "9.2% in 2014 to 2.8% in 2024 (r = \u22120.742; P = 0.009)."
)
add_para(results_trend)

results_trend2 = (
    "When comparing the early (2014\u20132018; n = 1,468) and late (2019\u20132024; n = 1,487) "
    "periods, median EBL was significantly lower in the late period (620 vs. 700 mL; "
    "P < 0.001). The transfusion rate (2.6% vs. 6.1%; P < 0.001) and massive hemorrhage "
    "rate (5.1% vs. 7.4%; P = 0.012) were also significantly lower in the late period."
)
add_para(results_trend2)

# Insert Figure 2
doc.add_paragraph()
doc.add_picture(str(FIG_DIR / "fig2_temporal_trends.png"), width=Inches(6.0))
p_cap = doc.add_paragraph()
p_cap.alignment = WD_ALIGN_PARAGRAPH.LEFT
run = p_cap.add_run(
    "Figure 2. Temporal trends in cesarean section outcomes (2014\u20132024). "
    "(A) Median estimated blood loss with IQR shading. (B) Transfusion rate. "
    "(C) Massive hemorrhage rate. (D) Annual case volume."
)
run.font.name = "Times New Roman"
run.font.size = Pt(10)
run.italic = True

# Save
doc.save(str(BASE / "manuscript_methods_results.docx"))
print("Manuscript saved: manuscript_methods_results.docx")

# ============================================================
# PPTX - Editable figures
# ============================================================
prs = Presentation()
prs.slide_width = PptxInches(13.333)
prs.slide_height = PptxInches(7.5)

figures_info = [
    ("fig1_blood_loss_distribution.png",
     "Figure 1",
     "Distribution of estimated blood loss during cesarean section"),
    ("fig2_temporal_trends.png",
     "Figure 2",
     "Temporal trends in cesarean section outcomes (2014\u20132024)"),
    ("fig3_forest_blood_loss.png",
     "Figure 3",
     "Forest plot: Factors associated with estimated blood loss"),
    ("fig4_forest_transfusion.png",
     "Figure 4",
     "Forest plot: Factors associated with transfusion"),
    ("fig5_subgroup_violins.png",
     "Figure 5",
     "Estimated blood loss by clinical subgroups"),
    ("fig6_roc_curves.png",
     "Figure 6",
     "ROC curves for transfusion and massive hemorrhage prediction"),
]

for fname, title, caption in figures_info:
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.2),
                                      PptxInches(12.333), PptxInches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = PptxPt(24)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Image
    img_path = FIG_DIR / fname
    slide.shapes.add_picture(str(img_path),
                             PptxInches(1.0), PptxInches(0.9),
                             PptxInches(11.333), PptxInches(5.5))

    # Caption
    txBox2 = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(6.6),
                                       PptxInches(12.333), PptxInches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = caption
    p2.font.size = PptxPt(14)
    p2.font.italic = True
    p2.alignment = PP_ALIGN.CENTER

prs.save(str(BASE / "figures.pptx"))
print("PPTX saved: figures.pptx")
print("Done!")
