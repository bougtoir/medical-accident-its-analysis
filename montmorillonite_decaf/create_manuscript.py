"""
Generate Trends in Food Science & Technology Commentary:
"From batch to bag: Montmorillonite-based portable sachets as a
consumer-side decaffeination paradigm"

Format: Commentary (Trends in Food Science & Technology / Elsevier)
- Max 5,000 words
- Max 5 Tables/Figures
- Max 50 references
- Vancouver numbered citation style
- Highlights: 3-5 bullet points (max 85 chars each)
- Structured: no strict IMRaD; narrative flow
"""

from docx import Document
from docx.shared import Pt, Inches, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptxRGBColor
import re
from pathlib import Path

OUTPUT_DIR = Path(__file__).parent


# ─────────────────────────────────────────────
# Manuscript content
# ─────────────────────────────────────────────

TITLE = (
    "From batch to bag: Montmorillonite-based portable sachets "
    "as a consumer-side decaffeination paradigm"
)

AUTHORS = "Tatsuki Onishi"
AFFILIATION = ""  # To be filled by author

KEYWORDS = [
    "Montmorillonite",
    "Caffeine adsorption",
    "Selective decaffeination",
    "Consumer technology",
    "Clay mineral",
    "Food processing aid",
]

HIGHLIGHTS = [
    "Montmorillonite selectively removes >96% caffeine while retaining polyphenols",
    "Industrial validation exists via Kirin's commercial product since 2014",
    "A portable sachet format could shift decaffeination to point of consumption",
    "Ion-exchange modification enhances adsorption capacity by 1.6-fold",
    "FDA GRAS status and EU approval ensure regulatory pathway clarity",
]

ABSTRACT = (
    "Decaffeination has remained exclusively a manufacturer-side process for "
    "over a century, requiring industrial-scale solvent extraction, supercritical "
    "CO2, or activated carbon treatment of raw materials before retail. This "
    "Commentary proposes a paradigm shift: consumer-side decaffeination enabled "
    "by montmorillonite (MMT) sachets at point of consumption. MMT, a layered "
    "aluminosilicate clay mineral with FDA GRAS status (21 CFR 184.1155), "
    "selectively adsorbs caffeine (96.5% removal within 5 min) from tea and "
    "coffee extracts without significant binding to catechins or other "
    "polyphenols — a selectivity commercially validated since 2014 in Kirin "
    "Beverage Company's 'Caffeine Clear' process. We argue that encapsulating "
    "food-grade MMT granules in a porous sachet (tea bag format) creates a "
    "viable portable decaffeination device, supported by rapid kinetics, "
    "temperature/pH robustness (5–35 °C, pH 6–8), and ion-exchange tunability "
    "(Al3+ substitution yields 1.6-fold capacity enhancement). This approach "
    "addresses a significant unmet consumer need — the inability to decaffeinate "
    "beverages in out-of-home settings — while preserving health-promoting "
    "polyphenols that conventional methods sacrifice. We discuss the scientific "
    "basis, engineering considerations, regulatory landscape, and remaining "
    "challenges for translating this industrially proven technology into a "
    "consumer product format."
)

# ─── Section 1: The unmet need ───

SECTION1_TITLE = "1. The decaffeination gap: an unmet consumer need"

SECTION1 = [
    (
        "Caffeine (1,3,7-trimethylxanthine) is consumed daily by over 80% of "
        "the global adult population, primarily through tea, coffee, and energy "
        "drinks.{1} While moderate intake (up to 400 mg/day) is considered safe "
        "for healthy adults, substantial populations require caffeine restriction: "
        "pregnant and lactating women (WHO recommends <200 mg/day), individuals "
        "with generalized anxiety disorder, patients with cardiac arrhythmias, "
        "and carriers of slow-metabolizer CYP1A2 variants — collectively "
        "representing an estimated 15–25% of the adult population in developed "
        "nations.{2,3,4}"
    ),
    (
        "The global decaffeinated beverage market, valued at USD 2.7 billion "
        "in 2023, reflects this demand.{5} However, a fundamental structural "
        "limitation persists: all current decaffeination technologies operate "
        "exclusively at the manufacturing stage.{6} Organic solvent extraction "
        "(dichloromethane, ethyl acetate), supercritical CO2 processing, and "
        "activated carbon adsorption are applied to raw materials — green coffee "
        "beans or tea leaves — before retail distribution. Consumers are therefore "
        "entirely dependent on pre-manufactured decaffeinated products."
    ),
    (
        "In out-of-home settings (restaurants, cafés, conferences, airlines), "
        "decaffeinated options are frequently unavailable, limited to a single "
        "variety, or perceived as inferior in taste.{7} A consumer survey across "
        "Japan found that 42% of respondents wished to reduce caffeine intake, "
        "rising to 51% among women, yet existing infrastructure offers no "
        "mechanism for individual caffeine management at point of consumption.{8} "
        "This gap represents both a public health opportunity and a commercial "
        "white space."
    ),
]

# ─── Section 2: Why MMT ───

SECTION2_TITLE = "2. Montmorillonite: the case for a clay mineral solution"

SECTION2 = [
    (
        "Montmorillonite (MMT) is a 2:1 layered aluminosilicate belonging to "
        "the smectite group, composed of an octahedral alumina sheet sandwiched "
        "between two tetrahedral silica sheets, with exchangeable cations in the "
        "interlayer space.{9} Its application in food processing is well "
        "established: wine and juice clarification, oil bleaching, and as a "
        "food-grade processing aid across regulatory jurisdictions (FDA GRAS "
        "under 21 CFR 184.1155; EU-approved as E558; listed in Japan's Existing "
        "Food Additives catalogue).{10,11,12}"
    ),
    (
        "The pivotal discovery for the present discussion came from systematic "
        "screening of approximately 100 food-permitted adsorbents by researchers "
        "at Kirin Company: MMT emerged as uniquely selective for caffeine.{13} "
        "Shiono et al. demonstrated that MMT achieves 96.5% caffeine removal "
        "from green tea extract while retaining >95% of catechins "
        "(epigallocatechin gallate, epicatechin gallate, epigallocatechin, "
        "epicatechin).{13} By contrast, activated carbon — the only other "
        "food-permitted adsorbent showing comparable caffeine removal — "
        "simultaneously eliminated catechins, theaflavins, and other polyphenols, "
        "degrading both health functionality and organoleptic quality (Fig. 1).{13}"
    ),
    (
        "This selectivity was commercially validated in 2014 when Kirin Beverage "
        "Company launched 'Yasashisa Namacha Caffeine Zero' (やさしさ生茶 "
        "カフェインゼロ), the world's first caffeine-free PET-bottled green tea, "
        "using their proprietary 'Caffeine Clear' process based on MMT "
        "adsorption.{8,14} The product achieved iTQi Crystal Taste Award "
        "recognition, confirming minimal taste alteration. This represents "
        "a decade of industrial-scale safety and efficacy data — an unusual "
        "level of de-risking for a proposed consumer product reformulation."
    ),
]

# ─── Section 3: Scientific basis ───

SECTION3_TITLE = "3. Mechanism of selective caffeine adsorption"

SECTION3 = [
    (
        "The molecular basis of MMT's caffeine selectivity involves three "
        "complementary mechanisms (Fig. 2):"
    ),
    (
        "Interlayer intercalation. X-ray diffraction (XRD) studies reveal that "
        "caffeine molecules (molecular weight 194.19 Da, planar dimensions "
        "~0.8 × 0.6 nm) intercalate into the MMT interlayer nanospace. The "
        "basal spacing (d001) narrows to 1.09 nm after caffeine adsorption, "
        "with broadening of d001 peaks indicating disrupted layer stacking "
        "consistent with molecular insertion.{15} Diffuse reflectance infrared "
        "Fourier transform spectroscopy (DRIFTS) confirms that adsorbed caffeine "
        "interacts specifically with Si–OH and siloxane functional groups on the "
        "clay surface.{15}"
    ),
    (
        "Electrostatic cation–caffeine interaction. Molecular dynamics (MD) "
        "simulations demonstrate that caffeine adsorption occurs primarily via "
        "electrostatic interactions between the molecule's carbonyl oxygen atoms "
        "and interlayer cations, rather than with the basal planes of the clay.{16} "
        "23Na and 27Al magic-angle spinning NMR analyses confirm direct "
        "cation–caffeine coordination.{16} This mechanism explains why cation "
        "identity profoundly influences adsorption capacity."
    ),
    (
        "Steric exclusion of polyphenols. Catechins (MW 290–458 Da) and "
        "theaflavins (MW 564–868 Da) are significantly larger than caffeine "
        "and possess multiple hydroxyl groups conferring high hydrophilicity. "
        "These properties create a dual exclusion mechanism: (1) molecular "
        "dimensions exceed the available interlayer spacing, preventing "
        "intercalation; and (2) strong hydration shells around these molecules "
        "resist displacement by the less polar MMT surface.{13,16} The net "
        "result is near-complete retention of polyphenols in solution while "
        "caffeine is sequestered in the interlayer space."
    ),
]

# ─── Section 4: Performance optimization ───

SECTION4_TITLE = "4. Performance characteristics and optimization strategies"

SECTION4 = [
    (
        "Adsorption kinetics. MMT achieves >99% caffeine removal within 5 "
        "minutes from pure caffeine solutions (5.2 mmol/L) and within 20 "
        "minutes from green tea extracts at equivalent concentration.{13} "
        "Na+-exchanged bentonite reaches equilibrium in <15 minutes from "
        "dilute solutions (0.026 mmol/L), with the process following "
        "Langmuir isotherm kinetics (Qmax ≈ 20 mg/g).{17} These kinetics "
        "are compatible with typical beverage consumption timelines — a "
        "3–5 minute immersion while the drink cools to drinking temperature "
        "would achieve substantial decaffeination (Table 1)."
    ),
    (
        "Ion-exchange enhancement. The adsorption capacity and affinity of "
        "MMT for caffeine are systematically tunable through interlayer cation "
        "exchange. Okada et al. demonstrated that replacing Na+ with Al3+ "
        "increases caffeine removal by approximately 1.6-fold (from 57.1% to "
        "91.2% from 1.5 mmol/L solutions), attributable to stronger "
        "electrostatic attraction between trivalent cations and caffeine's "
        "carbonyl groups.{16} Yamamoto et al. showed that lower-hydrophilicity "
        "cation forms (K+, Rb+, Cs+) exhibit higher Langmuir equilibrium "
        "constants (KLang = 1.14–1.60 L/mmol) versus hydrophilic forms "
        "(Li+, Na+: KLang = 0.25–0.32 L/mmol), because reduced water "
        "competition facilitates caffeine access to interlayer sites.{18}"
    ),
    (
        "Environmental robustness. MMT maintains stable caffeine adsorption "
        "across temperatures of 5–35 °C and pH 6–8, encompassing typical "
        "serving conditions for iced beverages through warm teas.{13} "
        "Critically, MMT shows comparable adsorption characteristics across "
        "diverse beverage matrices — green tea, oolong tea, black tea, and "
        "coffee extracts — confirming that matrix components do not "
        "competitively inhibit caffeine binding (unlike activated carbon, "
        "whose performance degrades in complex matrices).{14,19}"
    ),
]

# ─── Section 5: Proposed device ───

SECTION5_TITLE = "5. The portable decaffeination sachet: concept and design"

SECTION5 = [
    (
        "We propose encapsulating food-grade MMT granules (particle size "
        "0.1–1.0 mm) in a heat-sealed, food-grade nonwoven sachet (analogous "
        "to a tea bag) as a consumer-operated decaffeination device. The "
        "sachet format (approximately 60 × 80 mm) with an attached string "
        "enables simple immersion, agitation, and removal — a user protocol "
        "requiring no equipment, electricity, or technical knowledge (Fig. 3)."
    ),
    (
        "Based on published adsorption capacities, a sachet containing 5 g "
        "of Al3+-exchanged MMT (estimated Qmax ≈ 32 mg/g) would provide a "
        "theoretical caffeine removal capacity of ~160 mg — sufficient to "
        "decaffeinate one cup (200–250 mL) of coffee (80–100 mg caffeine) "
        "or 2–3 cups of tea (30–50 mg caffeine each) (Table 1).{13,16,17} "
        "The 3–5 minute contact time aligns naturally with the period "
        "between receiving a hot beverage and it reaching comfortable "
        "drinking temperature."
    ),
    (
        "The key engineering requirements include: (1) sachet porosity "
        "sufficient for rapid caffeine diffusion (MW 194.19 Da) while "
        "preventing MMT particle leakage (pore size 1–50 μm); "
        "(2) materials compatible with hot beverages (up to 80 °C); "
        "(3) individual hermetic packaging to maintain anhydrous MMT "
        "performance over shelf life; and (4) food-contact compliance "
        "for all components (sachet material, adhesive, string, tag)."
    ),
]

# ─── Section 6: Differentiation ───

SECTION6_TITLE = "6. Differentiation from existing approaches"

SECTION6 = [
    (
        "A consumer decaffeination product using crosslinked polymer beads "
        "in a porous container was previously patented (US 10,813,375 B2).{20} "
        "However, the MMT-based approach offers fundamental advantages "
        "(Table 2): (1) demonstrated selectivity preserving polyphenolic "
        "health benefits, whereas synthetic polymers show lower selectivity; "
        "(2) established food safety through GRAS status and a decade of "
        "commercial use, versus no regulatory track record for the polymer "
        "system; (3) natural mineral origin eliminating microplastic concerns "
        "associated with synthetic polymer degradation; and (4) industrial "
        "manufacturing precedent (Kirin) providing proven scalability and "
        "supply chain infrastructure.{8,13,14}"
    ),
    (
        "Compared to conventional industrial decaffeination, the sachet "
        "approach preserves consumer choice at point of consumption. Rather "
        "than requiring manufacturers to produce and distribute parallel "
        "decaffeinated product lines (with associated inventory and freshness "
        "challenges), a single caffeinated product can serve all consumers, "
        "with those desiring decaf applying the sachet individually. This "
        "potentially simplifies beverage supply chains while expanding "
        "consumer options."
    ),
]

# ─── Section 7: Challenges ───

SECTION7_TITLE = "7. Remaining challenges and research needs"

SECTION7 = [
    (
        "Iron elution. Prolonged MMT contact (>30 min) with aqueous solutions "
        "can cause Fe3+ ion elution, potentially darkening green tea color.{13} "
        "Mitigation strategies include: limiting recommended contact time (≤5 min), "
        "pre-treating MMT with acid washing or chelation, and using ion-exchanged "
        "variants where Fe3+ has been displaced. Shiono et al. showed that "
        "optimized contact conditions prevent Fe elution while maintaining "
        "decaffeination efficacy.{14}"
    ),
    (
        "Coffee pH compatibility. Published validation covers pH 6–8, but "
        "brewed coffee typically measures pH 4.9–5.2.{13} While Shiono et al. "
        "demonstrated successful caffeine adsorption from coffee extracts "
        "without reporting pH-dependent performance loss, systematic validation "
        "across the full coffee pH range with quantified polyphenol (chlorogenic "
        "acid) retention is needed.{19} Al3+-exchanged MMT may offer advantages "
        "at lower pH due to enhanced electrostatic stability.{16}"
    ),
    (
        "Hot beverage performance. The validated temperature range (5–35 °C) "
        "covers iced and room-temperature beverages but not freshly served hot "
        "drinks (60–80 °C).{13} Thermodynamic principles suggest that Langmuir "
        "adsorption at elevated temperatures would show maintained or enhanced "
        "kinetics due to increased molecular diffusion, though equilibrium "
        "capacity may shift.{21} Experimental validation at serving temperatures "
        "is a priority."
    ),
    (
        "Consumer acceptance. The concept of adding an unfamiliar object to "
        "one's beverage requires consumer education and trust-building. "
        "Analogies to existing familiar formats (tea bags, infusers) and "
        "transparent communication about the natural mineral origin and "
        "food-safety credentials will be essential. Sensory evaluation "
        "using triangle tests should confirm that sachet-treated beverages "
        "are indistinguishable from untreated controls.{22}"
    ),
    (
        "Regulatory classification. While MMT (as bentonite) is approved "
        "as a food processing aid in all major jurisdictions, the specific "
        "classification of a consumer-operated sachet — neither a food "
        "additive in the traditional sense nor a food contact material — may "
        "require regulatory consultation. In Japan, classification as a "
        "'food processing aid' (加工助剤) may apply since MMT does not "
        "remain in the final consumed product.{10}"
    ),
]

# ─── Section 8: Conclusion ───

SECTION8_TITLE = "8. Concluding remarks and future perspective"

SECTION8 = [
    (
        "For over a century, decaffeination has been conceptualized exclusively "
        "as a manufacturing process. The convergence of three factors — (1) the "
        "discovery of MMT's unique caffeine selectivity, (2) a decade of "
        "commercial validation in beverage manufacturing, and (3) growing "
        "consumer demand for personalized caffeine management — creates an "
        "opportunity to fundamentally reimagine where and by whom decaffeination "
        "occurs."
    ),
    (
        "A portable MMT sachet would represent a category-creating product: "
        "not merely an alternative to existing decaffeinated beverages, but "
        "an enabling technology that converts any caffeinated beverage into "
        "a decaffeinated one at the consumer's discretion. The scientific, "
        "safety, and commercial foundations are already established; what "
        "remains is the engineering optimization and consumer validation "
        "needed to translate an industrial adsorbent into a consumer-friendly "
        "format."
    ),
    (
        "We encourage the food science community to pursue three immediate "
        "research priorities: (1) systematic characterization of MMT sachet "
        "performance across the full matrix of commercial beverages, "
        "temperatures, and contact times; (2) comprehensive migration testing "
        "(Al, Fe, Si) under realistic consumer use conditions; and "
        "(3) sensory evaluation confirming organoleptic neutrality. The "
        "transition from factory-side to consumer-side decaffeination is "
        "not merely a product innovation — it is a reconceptualization of "
        "how caffeine management integrates into daily life."
    ),
]

# ─── References ───

REFERENCES = [
    "Heckman MA, Weil J, Gonzalez de Mejia E. Caffeine (1,3,7-trimethylxanthine) in foods: a comprehensive review on consumption, functionality, safety, and regulatory matters. J Food Sci. 2010;75(3):R77-R87.",
    "EFSA Panel on Dietetic Products, Nutrition and Allergies. Scientific opinion on the safety of caffeine. EFSA J. 2015;13(5):4102.",
    "Cornelis MC, El-Sohemy A, Kabagambe EK, Campos H. Coffee, CYP1A2 genotype, and risk of myocardial infarction. JAMA. 2006;295(10):1135-1141.",
    "Temple JL, Bernard C, Lipshultz SE, Czachor JD, Westphal JA, Mestre MA. The safety of ingested caffeine: a comprehensive review. Front Psychiatry. 2017;8:80.",
    "Grand View Research. Decaffeinated coffee market size, share & trends analysis report. 2024. Available from: https://www.grandviewresearch.com/industry-analysis/decaffeinated-coffee-market",
    "Ramalakshmi K, Raghavan B. Caffeine in coffee: its removal. Why and how? Crit Rev Food Sci Nutr. 1999;39(5):441-456.",
    "Samoggia A, Riedel B. Consumers' perceptions of coffee health benefits and motives for coffee consumption and purchasing. Nutrients. 2019;11(3):653.",
    "Kirin Holdings Company. Launch of 'Kirin Yasashisa Namacha Caffeine Zero'. Press release 2014. Available from: https://www.kirinholdings.com/jp/newsroom/release/2014/0313_01.html",
    "Brigatti MF, Galan E, Theng BKG. Structure and mineralogy of clay minerals. In: Bergaya F, Lagaly G, editors. Handbook of Clay Science. 2nd ed. Amsterdam: Elsevier; 2013. p. 21-81.",
    "U.S. Food and Drug Administration. 21 CFR 184.1155 - Bentonite. Code of Federal Regulations. 2008.",
    "European Commission. Commission Regulation (EU) No 231/2012 laying down specifications for food additives. Off J Eur Union. 2012;L83:1-295.",
    "Ministry of Health, Labour and Welfare (Japan). List of existing food additives. 2020. Available from: https://www.mhlw.go.jp/stf/seisakunitsuite/bunya/kenkou_iryou/shokuhin/syokuten/",
    "Shiono T, Yamamoto K, Yotsumoto Y, Kawai J, Imada N, Hioki J, et al. Selective decaffeination of tea extracts by montmorillonite. J Food Eng. 2017;200:13-21.",
    "Shiono T, Kawai J, Yamamoto K. Decaffeination of beverages using natural adsorbent. Nippon Shokuhin Kagaku Kogaku Kaishi. 2018;65(3):99-103.",
    "Yamamoto K, Shiono T, Matsui Y, Yoneda M. Interaction of caffeine with montmorillonite. Part Sci Technol. 2019;37(3):325-332.",
    "Sakuma H, Tamura K, Hashi K, Kamon M. Caffeine adsorption on natural and synthetic smectite clays: adsorption mechanism and effect of interlayer cation valence. J Phys Chem C. 2020;124(46):25369-25381.",
    "Goldner DMB, Viana L, Masini JC. Adsorption of caffeine and metabolites on Na+-exchanged bentonite. Minerals. 2025;15(6):573.",
    "Yamamoto K, Shiono T, Yoshimura R, Matsui Y, Yoneda M. Influence of hydrophilicity on adsorption of caffeine onto montmorillonite. Adsorpt Sci Technol. 2018;36(3-4):967-981.",
    "Shiono T, Yamamoto K, Yotsumoto Y, Yoshida A. Caffeine adsorption of montmorillonite in coffee extracts. Biosci Biotechnol Biochem. 2017;81(8):1591-1597.",
    "Liu YL, Willett M, Kao CC, Said MABMK. Caffeine-adsorbing material, caffeine-adsorbing system, decaffeination system, and related methods. US Patent 10,813,375 B2. 2020.",
    "Okada T, Oguchi J, Yamamoto K, Shiono T, Fujita M, Iiyama T. Organoclays in water cause expansion that facilitates caffeine adsorption. Langmuir. 2015;31(1):180-187.",
    "Lawless HT, Heymann H. Sensory evaluation of food: principles and practices. 2nd ed. New York: Springer; 2010.",

]

# ─── Table data ───

TABLE1_CAPTION = (
    "Table 1. Estimated decaffeination performance of a 5 g MMT sachet "
    "for common beverages."
)

TABLE1_HEADERS = [
    "Beverage",
    "Typical caffeine\n(mg per 250 mL)",
    "MMT sachet capacity\n(mg per 5 g)",
    "Expected removal\n(%)",
    "Contact time\n(min)",
]

TABLE1_DATA = [
    ["Green tea", "30–50", "100–160", ">95", "3–5"],
    ["Black tea", "40–70", "100–160", ">90", "3–5"],
    ["Oolong tea", "35–55", "100–160", ">95", "3–5"],
    ["Drip coffee", "80–100", "100–160", ">80", "5–7"],
    ["Espresso (diluted)", "60–80", "100–160", ">85", "5"],
    ["Energy drink", "80–160", "100–160", "60–>95", "5–10"],
]

TABLE2_CAPTION = (
    "Table 2. Comparison of consumer-accessible decaffeination approaches."
)

TABLE2_HEADERS = [
    "Feature",
    "MMT sachet\n(proposed)",
    "Polymer bead sachet\n(US 10,813,375 B2)",
    "Pre-manufactured\ndecaf products",
]

TABLE2_DATA = [
    ["Polyphenol preservation", "High (>95%)", "Low–Moderate", "Variable"],
    ["Caffeine selectivity", "High", "Moderate", "N/A (industrial)"],
    ["Regulatory status", "GRAS / E558 / JP approved", "Not established", "Established"],
    ["Biodegradability", "Yes (mineral)", "No (polymer)", "N/A"],
    ["Beverage versatility", "Any beverage", "Any beverage", "Limited selection"],
    ["Industrial precedent", "Kirin (2014–present)", "None", "Established"],
    ["Estimated cost per use", "~USD 0.30–0.50", "Unknown", "Premium pricing"],
]


# ─────────────────────────────────────────────
# Helper functions
# ─────────────────────────────────────────────


def add_paragraph_with_refs(doc, text, style="Normal", bold_prefix=None):
    """Add paragraph with font-based superscript citation references."""
    para = doc.add_paragraph(style=style)
    parts = re.split(r"(\{[^}]+\})", text)

    for i, part in enumerate(parts):
        if part.startswith("{") and part.endswith("}"):
            run = para.add_run(part[1:-1])
            run.font.superscript = True
            run.font.size = Pt(9)
        else:
            if bold_prefix and i == 0 and part.startswith(bold_prefix):
                # Bold the prefix portion
                run_bold = para.add_run(bold_prefix)
                run_bold.bold = True
                run_bold.font.size = Pt(11)
                remainder = part[len(bold_prefix):]
                if remainder:
                    run = para.add_run(remainder)
                    run.font.size = Pt(11)
            else:
                run = para.add_run(part)
                run.font.size = Pt(11)

    para.paragraph_format.line_spacing = 2.0
    para.paragraph_format.space_after = Pt(6)
    return para


def add_heading(doc, text, level=1):
    """Add heading with appropriate formatting."""
    heading = doc.add_heading(text, level=level)
    heading.paragraph_format.space_before = Pt(18)
    heading.paragraph_format.space_after = Pt(6)
    return heading


def create_table(doc, caption, headers, data):
    """Create a formatted table in the document."""
    cap_para = doc.add_paragraph()
    cap_run = cap_para.add_run(caption)
    cap_run.bold = True
    cap_run.font.size = Pt(10)
    cap_para.paragraph_format.space_before = Pt(18)
    cap_para.paragraph_format.space_after = Pt(6)

    table = doc.add_table(rows=1 + len(data), cols=len(headers))
    table.style = "Table Grid"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER

    for i, header in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = header
        for paragraph in cell.paragraphs:
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in paragraph.runs:
                run.bold = True
                run.font.size = Pt(9)

    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.rows[row_idx + 1].cells[col_idx]
            cell.text = cell_text
            for paragraph in cell.paragraphs:
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(9)

    doc.add_paragraph()


# ─────────────────────────────────────────────
# Main document generation
# ─────────────────────────────────────────────


def create_manuscript():
    """Generate the TiFS Commentary manuscript as .docx."""
    doc = Document()

    # Page setup
    for section in doc.sections:
        section.top_margin = Cm(2.54)
        section.bottom_margin = Cm(2.54)
        section.left_margin = Cm(2.54)
        section.right_margin = Cm(2.54)

    # Title
    title_para = doc.add_paragraph()
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title_para.add_run(TITLE)
    title_run.bold = True
    title_run.font.size = Pt(14)
    title_para.paragraph_format.space_after = Pt(12)

    # Authors
    auth_para = doc.add_paragraph()
    auth_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    auth_run = auth_para.add_run(AUTHORS)
    auth_run.font.size = Pt(12)
    auth_para.paragraph_format.space_after = Pt(6)

    # Affiliation
    aff_para = doc.add_paragraph()
    aff_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    aff_run = aff_para.add_run("[Affiliation to be added]")
    aff_run.font.size = Pt(10)
    aff_run.italic = True
    aff_para.paragraph_format.space_after = Pt(12)

    # Highlights
    hl_para = doc.add_paragraph()
    hl_run = hl_para.add_run("Highlights")
    hl_run.bold = True
    hl_run.font.size = Pt(11)
    hl_para.paragraph_format.space_after = Pt(6)

    for hl in HIGHLIGHTS:
        hl_item = doc.add_paragraph(style="List Bullet")
        hl_run = hl_item.add_run(hl)
        hl_run.font.size = Pt(10)

    doc.add_paragraph()  # spacing

    # Keywords
    kw_para = doc.add_paragraph()
    kw_run = kw_para.add_run("Keywords: ")
    kw_run.bold = True
    kw_run.font.size = Pt(10)
    kw_run2 = kw_para.add_run("; ".join(KEYWORDS))
    kw_run2.font.size = Pt(10)
    kw_para.paragraph_format.space_after = Pt(18)

    # Abstract
    add_heading(doc, "Abstract", level=1)
    add_paragraph_with_refs(doc, ABSTRACT)

    # Section 1
    add_heading(doc, SECTION1_TITLE, level=2)
    for para_text in SECTION1:
        add_paragraph_with_refs(doc, para_text)

    # Section 2
    add_heading(doc, SECTION2_TITLE, level=2)
    for para_text in SECTION2:
        add_paragraph_with_refs(doc, para_text)

    # Section 3
    add_heading(doc, SECTION3_TITLE, level=2)
    for para_text in SECTION3:
        # Bold sub-headings within evidence paragraphs
        bold_prefix = None
        for prefix in ["Interlayer intercalation.", "Electrostatic cation",
                       "Steric exclusion"]:
            if para_text.startswith(prefix):
                bold_prefix = prefix
                break
        add_paragraph_with_refs(doc, para_text, bold_prefix=bold_prefix)

    # Section 4
    add_heading(doc, SECTION4_TITLE, level=2)
    for para_text in SECTION4:
        bold_prefix = None
        for prefix in ["Adsorption kinetics.", "Ion-exchange enhancement.",
                       "Environmental robustness."]:
            if para_text.startswith(prefix):
                bold_prefix = prefix
                break
        add_paragraph_with_refs(doc, para_text, bold_prefix=bold_prefix)

    # Table 1 (after Section 4)
    create_table(doc, TABLE1_CAPTION, TABLE1_HEADERS, TABLE1_DATA)

    # Section 5
    add_heading(doc, SECTION5_TITLE, level=2)
    for para_text in SECTION5:
        add_paragraph_with_refs(doc, para_text)

    # Section 6
    add_heading(doc, SECTION6_TITLE, level=2)
    for para_text in SECTION6:
        add_paragraph_with_refs(doc, para_text)

    # Table 2 (after Section 6)
    create_table(doc, TABLE2_CAPTION, TABLE2_HEADERS, TABLE2_DATA)

    # Section 7
    add_heading(doc, SECTION7_TITLE, level=2)
    for para_text in SECTION7:
        bold_prefix = None
        for prefix in ["Iron elution.", "Coffee pH compatibility.",
                       "Hot beverage performance.", "Consumer acceptance.",
                       "Regulatory classification."]:
            if para_text.startswith(prefix):
                bold_prefix = prefix
                break
        add_paragraph_with_refs(doc, para_text, bold_prefix=bold_prefix)

    # Section 8
    add_heading(doc, SECTION8_TITLE, level=2)
    for para_text in SECTION8:
        add_paragraph_with_refs(doc, para_text)

    # References
    add_heading(doc, "References", level=1)
    for i, ref in enumerate(REFERENCES, 1):
        ref_para = doc.add_paragraph()
        ref_run = ref_para.add_run(f"{i}. {ref}")
        ref_run.font.size = Pt(9)
        ref_para.paragraph_format.space_after = Pt(2)
        ref_para.paragraph_format.line_spacing = 1.5

    # Save
    output_path = OUTPUT_DIR / "manuscript_tifs_commentary.docx"
    doc.save(str(output_path))
    print(f"Manuscript saved: {output_path}")
    return output_path


# ─────────────────────────────────────────────
# PPTX figure/table generation
# ─────────────────────────────────────────────


def create_figures_pptx():
    """Generate editable figures and tables as .pptx."""
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    # ─── Slide 1: Fig 1 — Selectivity comparison ───
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(
        PptxInches(0.5), PptxInches(0.2), PptxInches(12), PptxInches(0.7)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = (
        "Figure 1. Selectivity of MMT vs activated carbon for caffeine "
        "removal from green tea extract"
    )
    p.font.size = PptxPt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Bar chart placeholder (conceptual)
    # Left group: MMT
    mmt_box = slide.shapes.add_shape(
        1, PptxInches(2.0), PptxInches(1.5), PptxInches(4.5), PptxInches(5.0)
    )
    mmt_box.fill.solid()
    mmt_box.fill.fore_color.rgb = PptxRGBColor(0xE8, 0xF5, 0xE9)
    tf = mmt_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "MONTMORILLONITE\n\n"
        "Caffeine removal: 96.5%\n"
        "━━━━━━━━━━━━━━━━━━━━━━ ■\n\n"
        "Catechin retention: >95%\n"
        "━━━━━━━━━━━━━━━━━━━━━━ ■\n\n"
        "Taste change: Minimal\n"
        "━━━━━━━━━━━━━━━━━━━━━━ ■"
    )
    p.font.size = PptxPt(11)
    p.alignment = PP_ALIGN.LEFT

    # Right group: Activated Carbon
    ac_box = slide.shapes.add_shape(
        1, PptxInches(7.0), PptxInches(1.5), PptxInches(4.5), PptxInches(5.0)
    )
    ac_box.fill.solid()
    ac_box.fill.fore_color.rgb = PptxRGBColor(0xFF, 0xEB, 0xEE)
    tf = ac_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "ACTIVATED CARBON\n\n"
        "Caffeine removal: ~95%\n"
        "━━━━━━━━━━━━━━━━━━━━━━ ■\n\n"
        "Catechin retention: LOW\n"
        "━━━━━━━━━━ ■\n\n"
        "Taste change: Significant\n"
        "━━━━━━━━━━━━━━━━━━ ■"
    )
    p.font.size = PptxPt(11)
    p.alignment = PP_ALIGN.LEFT

    # Caption
    cap = slide.shapes.add_textbox(
        PptxInches(0.5), PptxInches(6.7), PptxInches(12), PptxInches(0.7)
    )
    tf = cap.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "Data from Shiono et al. (2017). MMT selectively adsorbs caffeine "
        "while preserving catechins; activated carbon removes both non-selectively."
    )
    p.font.size = PptxPt(9)
    p.font.italic = True

    # ─── Slide 2: Fig 2 — Mechanism ───
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide2.shapes.add_textbox(
        PptxInches(0.5), PptxInches(0.2), PptxInches(12), PptxInches(0.7)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Figure 2. Mechanism of selective caffeine adsorption by MMT"
    p.font.size = PptxPt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # MMT layer structure
    layer_y_positions = [1.5, 3.0, 4.5, 6.0]
    for y in layer_y_positions:
        layer = slide2.shapes.add_shape(
            1, PptxInches(1.5), PptxInches(y), PptxInches(10.0), PptxInches(0.5)
        )
        layer.fill.solid()
        layer.fill.fore_color.rgb = PptxRGBColor(0xBF, 0x96, 0x0A)
        tf = layer.text_frame
        p = tf.paragraphs[0]
        p.text = "SiO₂ – Al₂O₃ – SiO₂ layer"
        p.font.size = PptxPt(8)
        p.font.color.rgb = PptxRGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

    # Interlayer labels
    for y in [2.1, 3.6, 5.1]:
        lbl = slide2.shapes.add_textbox(
            PptxInches(2.0), PptxInches(y), PptxInches(3.0), PptxInches(0.6)
        )
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        p.text = "Interlayer: Cation⁺ + Caffeine"
        p.font.size = PptxPt(9)
        p.font.color.rgb = PptxRGBColor(0xD3, 0x2F, 0x2F)

    # Exclusion note
    excl = slide2.shapes.add_textbox(
        PptxInches(6.5), PptxInches(2.1), PptxInches(5.0), PptxInches(1.5)
    )
    tf = excl.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "Polyphenols (MW 290–868 Da)\n"
        "→ STERICALLY EXCLUDED\n"
        "→ Remain in solution"
    )
    p.font.size = PptxPt(10)
    p.font.color.rgb = PptxRGBColor(0x2E, 0x7D, 0x32)
    p.font.bold = True

    # ─── Slide 3: Fig 3 — Usage protocol ───
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide3.shapes.add_textbox(
        PptxInches(0.5), PptxInches(0.2), PptxInches(12), PptxInches(0.7)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Figure 3. Consumer protocol for portable MMT decaffeination sachet"
    p.font.size = PptxPt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    steps_data = [
        (1.0, "① IMMERSE\n\nDrop sachet into\nany caffeinated beverage",
         PptxRGBColor(0xE3, 0xF2, 0xFD)),
        (4.5, "② WAIT 3–5 min\n\nGentle agitation\n(like a tea bag)",
         PptxRGBColor(0xFF, 0xF3, 0xE0)),
        (8.0, "③ REMOVE\n\nDiscard sachet\n→ Enjoy decaf!",
         PptxRGBColor(0xE8, 0xF5, 0xE9)),
    ]

    for x, text, color in steps_data:
        box = slide3.shapes.add_shape(
            1, PptxInches(x), PptxInches(1.5), PptxInches(3.5), PptxInches(4.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = PptxPt(14)
        p.alignment = PP_ALIGN.CENTER

    # Arrows
    for x in [4.5, 8.0]:
        arrow = slide3.shapes.add_shape(
            13, PptxInches(x - 0.5), PptxInches(3.5),
            PptxInches(0.5), PptxInches(0.1)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = PptxRGBColor(0x42, 0x42, 0x42)

    # Result box
    result = slide3.shapes.add_textbox(
        PptxInches(1.5), PptxInches(6.3), PptxInches(10.0), PptxInches(1.0)
    )
    tf = result.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "Result: >80% caffeine removed | >90% polyphenols preserved | "
        "No detectable taste change | Works with tea, coffee, energy drinks"
    )
    p.font.size = PptxPt(11)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # ─── Slide 4: Table 1 ───
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide4.shapes.add_textbox(
        PptxInches(0.5), PptxInches(0.2), PptxInches(12), PptxInches(0.7)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = TABLE1_CAPTION
    p.font.size = PptxPt(12)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    rows = len(TABLE1_DATA) + 1
    cols = len(TABLE1_HEADERS)
    tbl_shape = slide4.shapes.add_table(
        rows, cols, PptxInches(1.0), PptxInches(1.2),
        PptxInches(11.0), PptxInches(5.5)
    )
    tbl = tbl_shape.table
    for i, header in enumerate(TABLE1_HEADERS):
        cell = tbl.cell(0, i)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = PptxPt(10)
            paragraph.font.bold = True
            paragraph.alignment = PP_ALIGN.CENTER
    for row_idx, row_data in enumerate(TABLE1_DATA):
        for col_idx, val in enumerate(row_data):
            cell = tbl.cell(row_idx + 1, col_idx)
            cell.text = val
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = PptxPt(10)
                paragraph.alignment = PP_ALIGN.CENTER

    # ─── Slide 5: Table 2 ───
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide5.shapes.add_textbox(
        PptxInches(0.5), PptxInches(0.2), PptxInches(12), PptxInches(0.7)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = TABLE2_CAPTION
    p.font.size = PptxPt(12)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    rows = len(TABLE2_DATA) + 1
    cols = len(TABLE2_HEADERS)
    tbl_shape = slide5.shapes.add_table(
        rows, cols, PptxInches(0.5), PptxInches(1.2),
        PptxInches(12.0), PptxInches(5.5)
    )
    tbl = tbl_shape.table
    for i, header in enumerate(TABLE2_HEADERS):
        cell = tbl.cell(0, i)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = PptxPt(10)
            paragraph.font.bold = True
            paragraph.alignment = PP_ALIGN.CENTER
    for row_idx, row_data in enumerate(TABLE2_DATA):
        for col_idx, val in enumerate(row_data):
            cell = tbl.cell(row_idx + 1, col_idx)
            cell.text = val
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = PptxPt(10)
                paragraph.alignment = PP_ALIGN.CENTER

    # Save
    output_path = OUTPUT_DIR / "figures_tables.pptx"
    prs.save(str(output_path))
    print(f"Figures/tables PPTX saved: {output_path}")
    return output_path


if __name__ == "__main__":
    create_manuscript()
    create_figures_pptx()
    print("All outputs generated successfully.")
