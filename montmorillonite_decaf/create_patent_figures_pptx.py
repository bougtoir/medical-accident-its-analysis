"""
特許出願用 図面（図1〜図7）を編集可能なPowerPoint（PPTX）として作成する。

画像貼り付けではなく、PowerPointのネイティブ図形（オートシェイプ・テキストボックス）
およびネイティブグラフ（棒・折れ線・散布図）で作図するため、PowerPoint上で
自由に編集できる。1スライド＝1図（図1〜図7）。

出力: patent_figures.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_MARKER_STYLE
from pptx.enum.dml import MSO_PATTERN

# JPO図面の慣行に合わせ、無彩色（白黒）のみで作図する。
# 色の代わりにハッチング・破線・マーカー形状で系列を区別する。
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "IPAGothic"


def _pattern(fmt, pattern):
    """系列の塗りを白地に黒のハッチングとする。"""
    fmt.fill.patterned()
    fmt.fill.pattern = pattern
    fmt.fill.fore_color.rgb = BLACK
    fmt.fill.back_color.rgb = WHITE
    fmt.line.color.rgb = BLACK


def _bw_line(series, dash=False, marker=XL_MARKER_STYLE.CIRCLE):
    """折れ線系列を黒線（必要に応じ破線）＋白塗りマーカーとする。"""
    series.format.line.color.rgb = BLACK
    series.format.line.width = Pt(2.0 if not dash else 1.75)
    if dash:
        series.format.line._get_or_add_ln().append(
            series.format.line._get_or_add_ln().makeelement(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash",
                {"val": "dash"}))
    series.marker.style = marker
    series.marker.size = 7
    series.marker.format.fill.solid()
    series.marker.format.fill.fore_color.rgb = WHITE
    series.marker.format.line.color.rgb = BLACK

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide(title):
    slide = prs.slides.add_slide(BLANK)
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.name = FONT
    r.font.color.rgb = BLACK
    return slide


def textbox(slide, x, y, w, h, text, size=12, align=PP_ALIGN.LEFT, bold=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = FONT
        r.font.color.rgb = BLACK
    return tb


def shape(slide, kind, x, y, w, h, fill=None, line=BLACK, line_w=1.5, dash=None):
    sp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = line
    sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if sp.has_text_frame:
        sp.text_frame.word_wrap = True
    return sp


def connector(slide, x1, y1, x2, y2, w=1.5, dash=False, arrow=False):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = BLACK
    cn.line.width = Pt(w)
    cn.shadow.inherit = False
    if dash:
        _set_dash(cn)
    if arrow:
        _set_arrow_end(cn)
    return cn


def _set_dash(cn):
    ln = cn.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    d = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
    ln.append(d)


def _set_arrow_end(cn):
    from pptx.oxml.ns import qn
    ln = cn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle"})
    ln.append(tail)


# ─────────────────────────── 図1 ───────────────────────────
def fig1():
    s = add_slide("図1　携帯型カフェイン吸着サシェの外観図および断面図")
    textbox(s, 1.5, 1.0, 4, 0.4, "(A) 外観図", 14, PP_ALIGN.CENTER, True)
    textbox(s, 8.0, 1.0, 4, 0.4, "(B) 断面図", 14, PP_ALIGN.CENTER, True)

    # (A) タグ・紐・テトラ本体
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 2.9, 1.6, 1.4, 0.6, fill=WHITE)
    textbox(s, 2.9, 1.68, 1.4, 0.4, "4", 12, PP_ALIGN.CENTER)
    connector(s, 3.6, 2.2, 3.6, 3.1)
    tri = shape(s, MSO_SHAPE.ISOSCELES_TRIANGLE, 2.1, 3.1, 3.0, 2.8, fill=WHITE)
    textbox(s, 5.2, 4.2, 1.5, 0.4, "← 1（サシェ）", 11)
    textbox(s, 0.4, 4.6, 1.6, 0.4, "2（包材）→", 11, PP_ALIGN.RIGHT)

    # (B) 断面：包材枠＋粒子
    shape(s, MSO_SHAPE.RECTANGLE, 8.3, 2.4, 3.4, 3.6, fill=WHITE)
    import random
    random.seed(3)
    for _ in range(22):
        cx = 8.5 + random.uniform(0, 3.0)
        cy = 2.6 + random.uniform(0, 3.2)
        shape(s, MSO_SHAPE.OVAL, cx, cy, 0.22, 0.22, fill=WHITE, line_w=1.0)
    textbox(s, 7.2, 1.6, 3.0, 0.7, "3（吸着材粒子\n粒径0.05〜1.0mm）", 10)
    textbox(s, 11.4, 3.4, 1.9, 0.7, "2（包材の孔\n孔径1〜50μm）", 10)


# ─────────────────────────── 図2 ───────────────────────────
def fig2():
    s = add_slide("図2　モンモリロナイト層間の選択的吸着機構（選択図）")
    ys = [5.6, 4.0, 2.4]
    for y in ys:
        rect = shape(s, MSO_SHAPE.RECTANGLE, 2.5, y, 8.3, 0.5)
        rect.fill.patterned()
        rect.fill.pattern = MSO_PATTERN.PERCENT_20
        rect.fill.fore_color.rgb = BLACK
        rect.fill.back_color.rgb = WHITE
        textbox(s, 0.4, y + 0.05, 2.0, 0.4, "シリケート層", 10, PP_ALIGN.RIGHT)
    # 層間カチオン
    for y in [4.7, 3.1]:
        for i in range(6):
            textbox(s, 2.9 + i * 1.3, y, 0.5, 0.4, "＋", 12, PP_ALIGN.CENTER)
    # 吸着済みカフェイン（下層間）
    for cx in [4.0, 6.0, 8.0]:
        shape(s, MSO_SHAPE.OVAL, cx, 3.0, 0.32, 0.32, fill=WHITE)
    # 取り込まれるカフェイン
    shape(s, MSO_SHAPE.OVAL, 4.2, 1.0, 0.34, 0.34, fill=WHITE)
    connector(s, 4.37, 1.4, 4.37, 2.35, arrow=True)
    textbox(s, 0.5, 0.9, 3.5, 0.6, "カフェイン（MW194）\n→層間に吸着", 11)
    # 排除されるポリフェノール
    shape(s, MSO_SHAPE.HEXAGON, 9.3, 0.9, 0.7, 0.6, fill=WHITE)
    connector(s, 9.65, 1.5, 9.65, 2.3, dash=True, arrow=True)
    textbox(s, 10.2, 0.8, 3.0, 0.9, "ポリフェノール\n（MW290〜869）\n→立体排除", 11)
    # ×印（排除）
    connector(s, 9.4, 2.5, 9.9, 2.9)
    connector(s, 9.4, 2.9, 9.9, 2.5)


# ─────────────────────────── 図3〜6（グラフ） ───────────────────────────
def _style_chart(chart, legend=True):
    chart.has_title = False
    if legend:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(11)
        chart.legend.font.name = FONT
    else:
        chart.has_legend = False
    for plot in chart.plots:
        try:
            plot.categories
        except Exception:
            pass


def fig3():
    s = add_slide("図3　各種飲料における除去率・保持率（実施例6）")
    cd = CategoryChartData()
    cd.categories = ["緑茶", "紅茶", "烏龍茶", "コーヒー", "エスプレッソ", "エナジードリンク"]
    cd.add_series("カフェイン除去率", (93, 90, 92, 88, 85, 72))
    cd.add_series("ポリフェノール保持率", (96, 93, 95, 92, 90, 0))
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                            Inches(1.0), Inches(1.2), Inches(11.3), Inches(5.8), cd)
    _style_chart(gf.chart)
    plot = gf.chart.plots[0]
    _pattern(plot.series[0].format, MSO_PATTERN.WIDE_UPWARD_DIAGONAL)
    _pattern(plot.series[1].format, MSO_PATTERN.PERCENT_25)


def fig4():
    s = add_slide("図4　粒径とカフェイン吸着効率の関係（実施例4）")
    cd = CategoryChartData()
    cd.categories = ["0.05", "0.1", "0.3", "0.5", "1.0", "2.0"]
    cd.add_series("3分後カフェイン除去率(%)", (75, 89, 91, 85, 68, 42))
    gf = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS,
                            Inches(1.2), Inches(1.3), Inches(10.9), Inches(5.6), cd)
    _style_chart(gf.chart, legend=False)
    _bw_line(gf.chart.plots[0].series[0])
    textbox(s, 1.5, 6.9, 8, 0.4, "横軸: MMT粒径 (mm)　／　最適範囲 0.1〜0.5mm", 11)


def fig5():
    s = add_slide("図5　接触時間と除去率・Fe溶出量の関係（実施例5）")
    cd = CategoryChartData()
    cd.categories = ["1", "3", "5", "10", "20", "30", "60"]
    cd.add_series("カフェイン除去率(%)", (62, 88, 93, 96, 97, 97, 98))
    cd.add_series("Fe溶出量(mg/L)×100", (0.5, 2, 5, 12, 28, 45, 120))
    gf = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS,
                            Inches(1.2), Inches(1.3), Inches(10.9), Inches(5.4), cd)
    _style_chart(gf.chart)
    _bw_line(gf.chart.plots[0].series[0], marker=XL_MARKER_STYLE.CIRCLE)
    _bw_line(gf.chart.plots[0].series[1], dash=True, marker=XL_MARKER_STYLE.SQUARE)
    textbox(s, 1.5, 6.8, 11, 0.5,
            "横軸: 接触時間 (分)。Fe溶出量は視認性のため×100表示（飲料水基準0.3mg/L=30）。", 10)


def fig6():
    s = add_slide("図6　吸着材の選択性比較")
    cd = XyChartData()
    ser = cd.add_series("吸着材")
    pts = [("MMT/Al交換MMT(本発明)", 93, 96), ("活性炭", 89, 31),
           ("ゼオライト", 28, 94), ("架橋ポリマー", 78, 62)]
    for _, xx, yy in pts:
        ser.add_data_point(xx, yy)
    gf = s.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER,
                            Inches(1.5), Inches(1.3), Inches(9.0), Inches(5.4), cd)
    _style_chart(gf.chart, legend=False)
    sc = gf.chart.plots[0].series[0]
    sc.marker.style = XL_MARKER_STYLE.CIRCLE
    sc.marker.size = 9
    sc.marker.format.fill.solid()
    sc.marker.format.fill.fore_color.rgb = BLACK
    sc.marker.format.line.color.rgb = BLACK
    textbox(s, 1.6, 6.8, 11, 0.6,
            "横軸: カフェイン除去率(%)、縦軸: カテキン保持率(%)。"
            "本発明(93,96)／活性炭(89,31)／ゼオライト(28,94)／架橋ポリマー(78,62)。", 10)


# ─────────────────────────── 図7 ───────────────────────────
def _cup(slide, x, y):
    """台形のカップ（フリーフォーム）を描く。x,yは左上基準の概略位置。"""
    from pptx.util import Inches as I
    fb = slide.shapes.build_freeform(Emu(I(x + 0.3).emu), Emu(I(y).emu), scale=1)
    pts = [(x + 0.3, y), (x + 2.9, y), (x + 3.3, y + 3.0), (x - 0.1, y + 3.0)]
    fb.add_line_segments([(Emu(I(px).emu), Emu(I(py).emu)) for px, py in pts[1:]],
                         close=True)
    shp = fb.convert_to_shape()
    shp.fill.background()
    shp.line.color.rgb = BLACK
    shp.line.width = Pt(1.6)
    shp.shadow.inherit = False
    return shp


def fig7():
    s = add_slide("図7　平型サシェとテトラ型サシェのティーカップ内での収まり比較（実施例7）")
    textbox(s, 1.0, 1.0, 4, 0.4, "(A) 平型（比較）", 14, PP_ALIGN.CENTER, True)
    textbox(s, 7.6, 1.0, 4, 0.4, "(B) テトラ型（本発明）", 14, PP_ALIGN.CENTER, True)
    # (A) カップ + はみ出す平型
    _cup(s, 2.0, 3.2)
    shape(s, MSO_SHAPE.RECTANGLE, 3.1, 1.8, 0.9, 4.0, fill=WHITE)
    textbox(s, 4.1, 1.7, 3.0, 0.7, "1（60×80mm）\nはみ出す", 10)
    textbox(s, 1.4, 6.4, 3.5, 0.4, "7（内径75mm）", 11, PP_ALIGN.CENTER)
    # (B) カップ + 収まるテトラ
    _cup(s, 8.6, 3.2)
    shape(s, MSO_SHAPE.ISOSCELES_TRIANGLE, 8.6, 4.0, 1.6, 1.8, fill=WHITE)
    textbox(s, 10.3, 3.9, 3.0, 0.7, "1'（底面35mm角）\nカップ内に収まる", 10)
    textbox(s, 8.0, 6.4, 3.5, 0.4, "7（内径75mm）", 11, PP_ALIGN.CENTER)


if __name__ == "__main__":
    fig1()
    fig2()
    fig3()
    fig4()
    fig5()
    fig6()
    fig7()
    prs.save("patent_figures.pptx")
    print("saved: patent_figures.pptx")
