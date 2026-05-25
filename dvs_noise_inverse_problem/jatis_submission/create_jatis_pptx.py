#!/usr/bin/env python3
"""
Generate editable PPTX with one figure per slide for JATIS submission.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pathlib import Path

OUT_DIR = Path(__file__).resolve().parent

FIGURES = [
    {
        'file': 'fig2_gapmap.png',
        'title': 'Fig. 1 — Gap Map',
        'caption': 'Gap map showing the four surveyed domains and identified research gaps.',
    },
    {
        'file': 'fig1_pipeline.png',
        'title': 'Fig. 2 — PI-DC-DVS Pipeline',
        'caption': 'System architecture: noise forward model → Bayesian inverse → residual → calibration.',
    },
    {
        'file': 'fig3_evaluation.png',
        'title': 'Fig. 3 — Systematic Evaluation',
        'caption': 'Boxplot comparison (NRR, SPR, F1, AUC) of three methods on 20 EBSSA recordings.',
    },
    {
        'file': 'fig5_per_recording.png',
        'title': 'Fig. 4 — Per-Recording Comparison',
        'caption': 'Per-recording noise removal rate for Fano filter, PI-DC-DVS NN, and temporal filter.',
    },
    {
        'file': 'fig4_a5_simulation.png',
        'title': 'Fig. 5 — A5 Model Simulation',
        'caption': 'Noise rate and SNR improvement across temperature–illuminance parameter space.',
    },
    {
        'file': 'fig6_demo.png',
        'title': 'Fig. 6 — Noise Inverse Demo',
        'caption': 'Raw events → noise map → P_noise distribution → residual with satellite trail.',
    },
    {
        'file': 'fig7_snr.png',
        'title': 'Fig. 7 — S/N Improvement',
        'caption': 'Fano factor map, temporal dynamics, and per-pixel SNR distribution.',
    },
]


def build_pptx():
    prs = Presentation()
    # Widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for fig_info in FIGURES:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2),
                                         Inches(12), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = fig_info['title']
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Image
        img_path = OUT_DIR / fig_info['file']
        if img_path.exists():
            # Scale to fit
            max_w = Inches(11)
            max_h = Inches(5.5)
            pic = slide.shapes.add_picture(str(img_path),
                                           Inches(1.2), Inches(0.9),
                                           width=max_w)
            # Constrain height
            if pic.height > max_h:
                ratio = max_h / pic.height
                pic.width = int(pic.width * ratio)
                pic.height = max_h
            # Center horizontally
            pic.left = int((prs.slide_width - pic.width) / 2)

        # Caption
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(6.7),
                                          Inches(12), Inches(0.7))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = fig_info['caption']
        p2.font.size = Pt(14)
        p2.alignment = PP_ALIGN.CENTER

    out_path = OUT_DIR / 'figures_jatis.pptx'
    prs.save(str(out_path))
    print(f"PPTX saved: {out_path}")


if __name__ == '__main__':
    build_pptx()
