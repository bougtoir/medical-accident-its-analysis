#!/usr/bin/env python3
"""Generate TiFS Commentary manuscript for MMT oil restoration sachet."""

import re
from pathlib import Path
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN

OUT_DIR = Path(__file__).parent


def add_superscript_text(paragraph, text):
    """Parse text with {n} markers and render as superscript runs."""
    parts = re.split(r"(\{[^}]+\})", text)
    for part in parts:
        if part.startswith("{") and part.endswith("}"):
            run = paragraph.add_run(part[1:-1])
            run.font.superscript = True
            run.font.size = Pt(8)
        else:
            run = paragraph.add_run(part)
            run.font.size = Pt(11)
    return paragraph


def set_cell_text(cell, text, bold=False, size=9):
    cell.text = ""
    p = cell.paragraphs[0]
    run = p.add_run(text)
    run.font.size = Pt(size)
    run.font.bold = bold


def add_table(doc, headers, rows, caption_text, caption_num):
    """Add a table with caption to the document."""
    cap = doc.add_paragraph()
    cap.paragraph_format.space_before = Pt(18)
    run = cap.add_run(f"Table {caption_num}. ")
    run.font.bold = True
    run.font.size = Pt(10)
    run = cap.add_run(caption_text)
    run.font.size = Pt(10)

    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.style = "Table Grid"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER

    for i, h in enumerate(headers):
        set_cell_text(table.rows[0].cells[i], h, bold=True)

    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            set_cell_text(table.rows[r_idx + 1].cells[c_idx], str(val))

    doc.add_paragraph()


TITLE = (
    "From plant to pot: Acid-activated montmorillonite sachets "
    "as a consumer-side cooking oil restoration paradigm"
)

# Each highlight must be <=85 characters
HIGHLIGHTS = [
    "Bleaching earth is the industrial standard for edible oil refining",          # 66
    "Consumer oil pots lack chemical purification capability",                      # 56
    "MMT sachets selectively adsorb polar degradation products",                    # 58
    "Projected acid value reduction of 30-50% per immersion",                       # 55
    "Existing sachet manufacturing enables immediate adoption",                     # 57
]

KEYWORDS = [
    "Montmorillonite",
    "Bleaching earth",
    "Cooking oil restoration",
    "Acid-activated clay",
    "Sachet",
    "Selective adsorption",
]

ABSTRACT = (
    "The global practice of reusing cooking oil after deep-frying generates chemically "
    "degraded oil containing elevated levels of free fatty acids (FFA), polar compounds, "
    "peroxides, and polymeric triglycerides. While industrial edible oil refining has relied "
    "on acid-activated montmorillonite (bleaching earth) for selective removal of polar "
    "degradation products for over a century, this proven technology has never been "
    "translated to the consumer setting. Current household oil pots employ paper or "
    "activated carbon filters that provide only physical filtration of particulates, "
    "with negligible chemical purification. This Commentary proposes a paradigm shift: "
    "sachets (tea-bag format) containing granulated acid-activated montmorillonite that "
    "are immersed in used cooking oil to selectively adsorb polar degradation products "
    "while retaining triglycerides. Drawing on industrial bleaching earth performance data "
    "and the established food-safety profile of bentonite (FDA GRAS 21 CFR 184.1155; "
    "EU E558), we project acid value reductions of 30-50% and total polar compound "
    "reductions of 20-40% per 20-minute immersion at recommended temperatures (50-80 C). "
    "This approach could extend cooking oil lifespan 2-3 fold, reducing household food "
    "waste and oil purchasing costs. Manufacturing infrastructure already exists in both "
    "the bleaching earth and tea-bag/filter sachet industries, making near-term "
    "commercialization feasible."
)

SECTIONS = [
    (
        "1. Introduction: The cooking oil reuse problem",
        [
            (
                "Deep-frying is one of the most widespread cooking methods globally, "
                "with an estimated 20 million tonnes of cooking oil consumed annually "
                "for frying applications.{1} During frying at 160-190 C, cooking oil "
                "undergoes thermal oxidation, hydrolysis, and polymerization, producing "
                "a complex mixture of degradation products including free fatty acids "
                "(FFA), peroxides, aldehydes, polymeric triglycerides, and polar "
                "compounds.{2,3} These degradation products progressively accumulate "
                "with repeated use, deteriorating oil quality and raising food safety "
                "concerns.{4}"
            ),
            (
                "Regulatory frameworks worldwide set limits on cooking oil degradation. "
                "The European Union and many Asian countries enforce a total polar "
                "compound (TPC) limit of 25% for frying oil, while Japan's Food "
                "Sanitation Act specifies an acid value (AV) limit of 2.5 mg KOH/g.{5,6} "
                "When oil exceeds these thresholds, it must be discarded, generating "
                "significant volumes of waste cooking oil (WCO). In Japan alone, "
                "household WCO generation is estimated at 100,000 tonnes annually.{7}"
            ),
            (
                "Despite these challenges, the practice of reusing cooking oil is "
                "deeply embedded in food culture. In Japan, dedicated oil pot "
                "(abura-potto) products are ubiquitous household items, used to filter "
                "and store oil between frying sessions. However, as we discuss below, "
                "current oil pot filtration technology addresses only the symptom "
                "(visible particulates) while ignoring the underlying cause (chemical "
                "degradation). This Commentary proposes a solution that bridges a "
                "century-old industrial technology with consumer convenience."
            ),
        ],
    ),
    (
        "2. Industrial precedent: Bleaching earth in oil refining",
        [
            (
                "Acid-activated montmorillonite, commercially known as bleaching earth "
                "or activated bleaching earth, has been the global standard for edible "
                "oil decolorization and purification since the early 20th century.{8,9} "
                "The annual global consumption of bleaching earth exceeds 2 million "
                "tonnes, serving the refining of virtually all commercially produced "
                "vegetable oils.{10}"
            ),
            (
                "The effectiveness of bleaching earth derives from its unique "
                "physicochemical properties. Acid activation of montmorillonite with "
                "mineral acids (HCl or H2SO4, 2-6 M, 2-6 h) partially dissolves the "
                "octahedral sheet, creating amorphous silica domains and expanding the "
                "interlayer space.{11} This treatment increases the specific surface "
                "area from 80-120 m2/g to 200-400 m2/g and generates abundant surface "
                "silanol groups that serve as adsorption sites for polar molecules.{12} "
                "The resulting material selectively adsorbs polar components -- "
                "carotenoids, chlorophyll degradation products, FFA, phospholipids, "
                "peroxides, and trace metals -- while non-polar triglycerides pass "
                "through with minimal loss (oil retention typically >96%).{8}"
            ),
            (
                "In Japan, bleaching earth manufacturers such as Mizusawa Chemical "
                "Industries and Chuo Silica supply acid-activated montmorillonite for "
                "JAS-standard refined oil production.{13} The safety of bentonite "
                "(the parent material of montmorillonite) as a food processing aid is "
                "well established: FDA GRAS status (21 CFR 184.1155), EU approval as "
                "E558, and listing in Japan's existing food additives registry.{14}"
            ),
            (
                "Critically, the industrial bleaching process operates as a batch or "
                "continuous slurry process: powdered bleaching earth (<0.1 mm) is "
                "dispersed directly into oil, mechanically agitated for 20-60 minutes, "
                "and then separated by filtration or centrifugation.{8} This format is "
                "entirely unsuitable for household use. The question, therefore, is "
                "whether bleaching earth technology can be reformatted for consumer "
                "deployment -- and if so, what technical adaptations are required."
            ),
        ],
    ),
    (
        "3. Current state of consumer oil filtration",
        [
            (
                "The Japanese oil pot market represents a mature consumer product "
                "category with established manufacturers (KAI Corporation, Pearl Metal, "
                "Duskin) offering products ranging from simple mesh strainers to "
                "cartridge-type filter systems (Table 1).{15} These products share a "
                "common design: hot or warm oil is poured through a filter element into "
                "a storage container, where it is held until the next frying session."
            ),
            "TABLE_1",
            (
                "Despite this product diversity, a fundamental limitation persists: "
                "all current consumer oil pot filters operate on physical filtration "
                "principles. Paper filters remove suspended particulates (batter "
                "fragments, food debris) but have no capacity for dissolved polar "
                "degradation products.{16} Activated carbon cartridges provide some "
                "odor reduction but show minimal effect on acid value, polar compound "
                "content, or peroxide value. Moreover, activated carbon lacks polarity "
                "selectivity and can non-selectively adsorb triglycerides, reducing "
                "oil yield.{17}"
            ),
            (
                "This represents a striking gap: the industrial oil refining sector "
                "has perfected chemical purification of edible oils using bleaching "
                "earth for over a century, yet no consumer product has translated this "
                "capability to the household setting. The technology exists; the "
                "consumer format does not."
            ),
        ],
    ),
    (
        "4. The paradigm shift: From physical filtration to chemical restoration",
        [
            (
                "We propose that acid-activated montmorillonite can be formatted as "
                "immersion sachets -- analogous to tea bags -- that are placed "
                "directly into used cooking oil to selectively adsorb polar degradation "
                "products. This represents a shift from physical filtration (particle "
                "removal) to chemical restoration (selective adsorption) (Fig. 1)."
            ),
            "FIGURE_1",
            (
                "The key technical adaptations required for this translation are: "
                "(i) granulation of bleaching earth to 0.3-0.5 mm particles (vs. "
                "industrial powder <0.1 mm) to prevent particle shedding through the "
                "sachet membrane while maintaining adequate surface area; "
                "(ii) heat-resistant sachet construction (PP or PTFE nonwoven, "
                "heat-sealed) tolerant of oil temperatures up to 180 C; "
                "(iii) extended immersion time (15-30 min vs. beverage sachets' 3-5 min) "
                "to compensate for the high viscosity of cooking oil (40-80x water); "
                "and (iv) consumer-friendly quality indicators for assessing oil "
                "improvement."
            ),
            (
                "The polarity-based selectivity of bleaching earth is ideally suited "
                "for this application.{18} Degradation products (FFA, peroxides, "
                "aldehyde decomposition products) are significantly more polar than "
                "intact triglycerides, and acid-activated montmorillonite's surface "
                "silanol groups preferentially interact with these polar species. "
                "Industrial data consistently show oil retention rates >96% during "
                "bleaching, indicating minimal triglyceride loss.{8,19}"
            ),
            (
                "Notably, the sachet format mirrors our parallel proposal for "
                "montmorillonite-based portable decaffeination sachets, where the same "
                "mineral's selective adsorption properties are exploited to remove "
                "caffeine from beverages while retaining polyphenols. The oil "
                "restoration sachet extends this consumer-side paradigm to a second "
                "major food processing application."
            ),
        ],
    ),
    (
        "5. Performance projections and quality metrics",
        [
            (
                "Based on industrial bleaching earth performance data and bench-scale "
                "adsorption studies,{20,21} we project the following performance "
                "for a 25 g acid-activated montmorillonite sachet immersed in 500 mL "
                "of used cooking oil at 80 C for 20 minutes (Table 2)."
            ),
            "TABLE_2",
            (
                "Several parameters merit discussion. First, the granulated format "
                "(0.3-0.5 mm) necessarily sacrifices some adsorption kinetics compared "
                "to industrial powder (<0.1 mm), but the 15-30 minute immersion time "
                "compensates for the diffusion-limited mass transfer through the sachet "
                "membrane.{20} Second, the recommended immersion temperature of "
                "50-80 C balances oil viscosity (affecting diffusion rate) with "
                "practical safety considerations.{22} Third, sachet reuse for 2-3 "
                "sessions is projected based on industrial bleaching earth capacity "
                "data scaled to the consumer oil volume range.{8}"
            ),
            (
                "The practical significance of these projections is considerable. "
                "For a household performing weekly deep-frying, extending oil lifespan "
                "from 3-4 uses to 8-10 uses would reduce annual oil purchases by "
                "approximately 50%, with corresponding reductions in waste cooking "
                "oil generation.{7}"
            ),
        ],
    ),
    (
        "6. Regulatory and safety considerations",
        [
            (
                "The regulatory pathway for montmorillonite-based oil sachets is "
                "straightforward, as bentonite is already approved as a food-contact "
                "material and processing aid across major jurisdictions.{14} In Japan, "
                "bentonite is listed in the existing food additives registry and is "
                "routinely used in edible oil refining under JAS standards.{13} The "
                "FDA classifies bentonite as GRAS (21 CFR 184.1155) with no limitation "
                "other than current good manufacturing practice, and the EU approves "
                "bentonite as E558 for food processing.{14}"
            ),
            (
                "A key safety consideration is the potential for residual acid (from "
                "the activation process) or trace metals (inherent to the clay mineral) "
                "to leach into the oil during immersion. Industrial practice addresses "
                "this through thorough post-activation washing until the rinse water "
                "reaches neutral pH, followed by drying and quality testing (heavy metal "
                "content, acid residue).{8,12} The same quality controls would apply to "
                "consumer sachets, with the additional advantage that the granulated "
                "format (0.3-0.5 mm) is retained within the sachet membrane, preventing "
                "direct particle contact with the end product."
            ),
            (
                "From a product classification standpoint, the sachet would likely "
                "be classified as a food-contact article (kitchenware) rather than a "
                "food additive, as the montmorillonite remains within the sachet and "
                "does not become a component of the food.{23} This classification "
                "simplifies the regulatory approval pathway."
            ),
        ],
    ),
    (
        "7. Manufacturing feasibility and market considerations",
        [
            (
                "A distinguishing feature of the proposed concept is that the "
                "manufacturing infrastructure already exists on both sides of the "
                "value chain. Bleaching earth production is a mature industry with "
                "established suppliers (Mizusawa Chemical Industries, Clariant, BASF, "
                "Oil-Dri) offering food-grade products.{10,13} Sachet and tea-bag "
                "manufacturing is similarly well-established, with existing production "
                "lines for heat-resistant nonwoven sachets that could be adapted for "
                "montmorillonite filling with minimal retooling. Oil pot filter "
                "manufacturers already produce heat-resistant nonwoven components "
                "that provide the necessary membrane technology.{15}"
            ),
            (
                "The economic proposition is favorable from both manufacturer and "
                "consumer perspectives. Montmorillonite raw material costs are "
                "approximately 300-500 JPY/kg for food-grade acid-activated "
                "product,{10} translating to a material cost of 8-13 JPY per sachet "
                "(25 g fill). Including packaging, assembly, and distribution, a "
                "consumer price of 100-200 JPY per sachet (2-3 uses each) would "
                "yield margins comparable to existing oil pot filter products while "
                "offering substantially greater functionality."
            ),
            (
                "Consumer adoption may be facilitated by the familiarity of the "
                "tea-bag format and the deeply embedded practice of oil reuse in "
                "Japanese households.{15} The sachet's simplicity -- immerse, wait, "
                "remove -- requires no new equipment and minimal behavioral change."
            ),
        ],
    ),
    (
        "8. Future directions and conclusions",
        [
            (
                "Several research directions could advance this concept toward "
                "commercialization. First, systematic bench-scale studies are needed "
                "to optimize acid activation conditions specifically for the granulated "
                "sachet format (balancing surface area, mechanical strength, and particle "
                "integrity during oil immersion). Second, consumer-oriented quality "
                "indicators -- such as colorimetric test strips for acid value or "
                "polar compound content -- would enhance the user experience and "
                "ensure safe oil reuse practices.{24} Third, life-cycle assessment "
                "comparing the environmental footprint of extended oil reuse (with "
                "sachet waste) versus more frequent oil replacement would quantify "
                "the sustainability benefits."
            ),
            (
                "In conclusion, the translation of bleaching earth technology from "
                "industrial oil refining to consumer-format sachets represents a "
                "compelling opportunity. The science is established, the safety "
                "profile is documented, the manufacturing infrastructure exists, "
                "and the consumer need is clear. What remains is the engineering "
                "challenge of reformatting a century-old industrial technology into "
                "a sachet that fits into the kitchen counter -- a challenge that, "
                "given the precedents in water filtration (activated carbon), tea "
                "preparation (tea bags), and our parallel proposal for decaffeination "
                "sachets, appears eminently solvable."
            ),
        ],
    ),
]

REFERENCES = [
    "1. Choe E, Min DB. Chemistry of deep-fat frying oils. J Food Sci. 2007;72(5):R77-R86.",
    "2. Zhang Q, Saleh ASM, Chen J, Shen Q. Chemical alterations taken place during deep-fat frying based on certain reaction products: a review. Chem Phys Lipids. 2012;165(6):662-681.",
    "3. Choe E, Min DB. Chemistry and reactions of reactive oxygen species in foods. Crit Rev Food Sci Nutr. 2006;46(1):1-22.",
    "4. Paul S, Mittal GS. Regulating the use of degraded oil/fat in deep-fat/oil food frying. Crit Rev Food Sci Nutr. 1997;37(7):635-662.",
    "5. Firestone D. Regulation of frying fats and oils. In: Erickson MD, editor. Deep Frying: Chemistry, Nutrition, and Practical Applications. 2nd ed. Champaign: AOCS Press; 2007. p. 373-385.",
    "6. Ministry of Health, Labour and Welfare, Japan. Standards and Criteria for Food Additives, etc. Under the Food Sanitation Act. Tokyo: MHLW; 2020.",
    "7. Ministry of the Environment, Japan. Food Waste Statistics 2022. Tokyo: MOE; 2023.",
    "8. Zschau W. Bleaching of fats and oils. Eur J Lipid Sci Technol. 2001;103(8):505-551.",
    "9. Mag T. Bleaching -- Theory and Practice. In: Wan PJ, editor. Introduction to Fats and Oils Technology. Champaign: AOCS Press; 1991. p. 124-159.",
    "10. IMARC Group. Bleaching Earth Market: Global Industry Trends, Share, Size, Growth, Opportunity and Forecast 2024-2032. 2024.",
    "11. Christidis GE, Scott PW, Dunham AC. Acid activation and bleaching capacity of bentonites from the islands of Milos and Chios, Aegean, Greece. Appl Clay Sci. 1997;12(4):329-347.",
    "12. Rossi M, Gianazza M, Alamprese C, Stanga F. The role of bleaching clays and synthetic silica in palm oil physical refining. Food Chem. 2003;82(2):291-296.",
    "13. Mizusawa Chemical Industries. Technical Data: Activated Bleaching Earth for Edible Oil Refining. Tokyo: Mizusawa; 2022.",
    "14. U.S. Food and Drug Administration. 21 CFR 184.1155 -- Bentonite. Washington: FDA; 2023.",
    "15. KAI Corporation. Oil Pot Product Catalog 2024. Tokyo: KAI; 2024.",
    "16. Takeoka GR, Full GH, Dao LT. Effect of heating on the characteristics and chemical composition of selected frying oils and fats. J Agric Food Chem. 1997;45(8):3244-3249.",
    "17. Lin S, Akoh CC, Reynolds AE. Recovery of used frying oils with adsorbent combinations: refrying and physicochemical studies. J Food Lipids. 2001;8(1):1-14.",
    "18. Proctor A, Toro-Vazquez JF. The Freundlich isotherm in studying adsorption in oil processing. J Am Oil Chem Soc. 1996;73(12):1627-1633.",
    "19. Miyagi A, Nakajima M. Regeneration of used frying oils using adsorption processing. J Am Oil Chem Soc. 2003;80(1):91-96.",
    "20. Maskan M, Bagci HI. Effect of different adsorbents on purification of used sunflower seed oil utilized for frying. Eur Food Res Technol. 2003;217(3):215-218.",
    "21. Yates RA, Caldwell JD. Adsorptive capacity of active filter aids for used cooking oil. J Am Oil Chem Soc. 1993;70(5):507-511.",
    "22. Sabah E. Decolorization of vegetable oils: chlorophyll-a adsorption by acid-activated sepiolite. J Colloid Interface Sci. 2007;310(1):1-7.",
    "23. Ministry of Health, Labour and Welfare, Japan. Food Contact Materials Regulations under the Food Sanitation Act. Tokyo: MHLW; 2020.",
    "24. Osawa CC, Goncalves LAG, Ragazzi S. Determination of hydroperoxides in oils and fats using kits. J Sci Food Agric. 2007;87(9):1659-1666.",
]


def build_docx():
    doc = Document()

    style = doc.styles["Normal"]
    style.font.name = "Times New Roman"
    style.font.size = Pt(11)
    style.paragraph_format.line_spacing = 2.0

    # Title
    t = doc.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = t.add_run(TITLE)
    run.font.size = Pt(14)
    run.font.bold = True

    doc.add_paragraph()

    # Highlights
    h = doc.add_paragraph()
    run = h.add_run("Highlights")
    run.font.bold = True
    run.font.size = Pt(12)
    for hl in HIGHLIGHTS:
        p = doc.add_paragraph(style="List Bullet")
        run = p.add_run(hl)
        run.font.size = Pt(10)

    doc.add_paragraph()

    # Keywords
    kw = doc.add_paragraph()
    run = kw.add_run("Keywords: ")
    run.font.bold = True
    run = kw.add_run("; ".join(KEYWORDS))

    doc.add_paragraph()

    # Abstract
    ab = doc.add_paragraph()
    run = ab.add_run("Abstract")
    run.font.bold = True
    run.font.size = Pt(12)
    doc.add_paragraph(ABSTRACT)

    doc.add_paragraph()

    # Sections
    for sec_title, paragraphs in SECTIONS:
        sh = doc.add_paragraph()
        run = sh.add_run(sec_title)
        run.font.bold = True
        run.font.size = Pt(12)

        for para_text in paragraphs:
            if para_text == "TABLE_1":
                add_table(
                    doc,
                    ["Filter type", "Mechanism", "AV reduction", "TPC reduction", "Oil recovery"],
                    [
                        ["Paper filter", "Particle filtration", "Negligible", "Negligible", ">99%"],
                        ["Activated carbon", "Non-selective adsorption", "~20%", "~15%", "~88%"],
                        ["Metal mesh", "Coarse filtration", "None", "None", ">99%"],
                        ["MMT sachet (proposed)", "Polar-selective adsorption", "30-50%", "20-40%", ">95%"],
                    ],
                    "Comparison of consumer oil treatment technologies.",
                    1,
                )
            elif para_text == "TABLE_2":
                add_table(
                    doc,
                    ["Parameter", "Projected value"],
                    [
                        ["Sachet fill", "25 g acid-activated MMT (0.3-0.5 mm)"],
                        ["Oil volume", "500 mL"],
                        ["Immersion temperature", "50-80 C (recommended 80 C)"],
                        ["Immersion time", "15-30 min"],
                        ["AV reduction", "30-50% per immersion"],
                        ["TPC reduction", "20-40% per immersion"],
                        ["PV reduction", "40-70% per immersion"],
                        ["Oil recovery", ">95%"],
                        ["Sachet lifetime", "2-3 uses"],
                        ["Cost per sachet", "100-200 JPY"],
                    ],
                    "Projected performance specifications for a montmorillonite oil restoration sachet.",
                    2,
                )
            elif para_text == "FIGURE_1":
                fig_cap = doc.add_paragraph()
                fig_cap.paragraph_format.space_before = Pt(18)
                run = fig_cap.add_run("Fig. 1. ")
                run.font.bold = True
                run.font.size = Pt(10)
                run = fig_cap.add_run(
                    "Schematic comparison of industrial bleaching earth process "
                    "(left: powder slurry, mechanical agitation, centrifugal separation) "
                    "and proposed consumer sachet format (right: granulated MMT in "
                    "tea-bag sachet, immersed in used oil, removed after treatment). "
                    "The polarity-based selectivity mechanism is conserved across both "
                    "formats."
                )
                run.font.size = Pt(10)
                run.font.italic = True
                doc.add_paragraph("[Figure placeholder -- to be prepared for submission]")
            else:
                p = doc.add_paragraph()
                add_superscript_text(p, para_text)

    # References
    doc.add_paragraph()
    rh = doc.add_paragraph()
    run = rh.add_run("References")
    run.font.bold = True
    run.font.size = Pt(12)

    for ref in REFERENCES:
        p = doc.add_paragraph()
        run = p.add_run(ref)
        run.font.size = Pt(10)

    out = OUT_DIR / "manuscript_tifs_commentary.docx"
    doc.save(str(out))
    print(f"Saved {out}")
    return out


def build_pptx():
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    # Slide 1: Table 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(
        PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(0.8)
    )
    tf = txBox.text_frame
    tf.text = "Table 1. Comparison of consumer oil treatment technologies"
    tf.paragraphs[0].font.size = PptxPt(18)
    tf.paragraphs[0].font.bold = True

    headers = ["Filter type", "Mechanism", "AV reduction", "TPC reduction", "Oil recovery"]
    rows_data = [
        ["Paper filter", "Particle filtration", "Negligible", "Negligible", ">99%"],
        ["Activated carbon", "Non-selective adsorption", "~20%", "~15%", "~88%"],
        ["Metal mesh", "Coarse filtration", "None", "None", ">99%"],
        ["MMT sachet (proposed)", "Polar-selective adsorption", "30-50%", "20-40%", ">95%"],
    ]
    tbl = slide.shapes.add_table(
        len(rows_data) + 1, len(headers),
        PptxInches(0.5), PptxInches(1.5), PptxInches(12), PptxInches(3)
    ).table
    for i, h in enumerate(headers):
        tbl.cell(0, i).text = h
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            tbl.cell(r + 1, c).text = val

    # Slide 2: Table 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(
        PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(0.8)
    )
    tf = txBox.text_frame
    tf.text = "Table 2. Projected performance for MMT oil restoration sachet"
    tf.paragraphs[0].font.size = PptxPt(18)
    tf.paragraphs[0].font.bold = True

    headers2 = ["Parameter", "Projected value"]
    rows2 = [
        ["Sachet fill", "25 g acid-activated MMT (0.3-0.5 mm)"],
        ["Oil volume", "500 mL"],
        ["Immersion temperature", "50-80 C (recommended 80 C)"],
        ["Immersion time", "15-30 min"],
        ["AV reduction", "30-50% per immersion"],
        ["TPC reduction", "20-40% per immersion"],
        ["PV reduction", "40-70% per immersion"],
        ["Oil recovery", ">95%"],
        ["Sachet lifetime", "2-3 uses"],
        ["Cost per sachet", "100-200 JPY"],
    ]
    tbl2 = slide.shapes.add_table(
        len(rows2) + 1, len(headers2),
        PptxInches(0.5), PptxInches(1.5), PptxInches(12), PptxInches(5)
    ).table
    for i, h in enumerate(headers2):
        tbl2.cell(0, i).text = h
    for r, row in enumerate(rows2):
        for c, val in enumerate(row):
            tbl2.cell(r + 1, c).text = val

    # Slide 3: Fig 1 placeholder
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(
        PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(0.8)
    )
    tf = txBox.text_frame
    tf.text = "Fig. 1. Industrial bleaching earth process vs. consumer sachet format"
    tf.paragraphs[0].font.size = PptxPt(18)
    tf.paragraphs[0].font.bold = True

    txBox2 = slide.shapes.add_textbox(
        PptxInches(1), PptxInches(2), PptxInches(11), PptxInches(4)
    )
    tf2 = txBox2.text_frame
    lines = [
        "[Figure placeholder]",
        "",
        "Left panel: Industrial process",
        "- Powdered bleaching earth (<0.1 mm) dispersed in oil",
        "- Mechanical agitation 20-60 min",
        "- Centrifugal separation",
        "",
        "Right panel: Consumer sachet",
        "- Granulated MMT (0.3-0.5 mm) in tea-bag sachet",
        "- Immersed in used oil 15-30 min",
        "- Pull out and discard",
    ]
    tf2.text = lines[0]
    for line in lines[1:]:
        p = tf2.add_paragraph()
        p.text = line
    for p in tf2.paragraphs:
        p.font.size = PptxPt(14)

    out = OUT_DIR / "figures_tables.pptx"
    prs.save(str(out))
    print(f"Saved {out}")


if __name__ == "__main__":
    build_docx()
    build_pptx()
    print("Done.")
