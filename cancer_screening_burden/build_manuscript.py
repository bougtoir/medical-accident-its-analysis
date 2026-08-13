"""Build reproducible manuscript materials from the simulation outputs.

This script reads parameters.yaml and the generated CSV/PNG files and writes:
  - manuscript.md (Markdown draft with inline figure/table references)
  - manuscript/manuscript_figures.pptx (one slide per figure)
  - manuscript/manuscript_tables.docx (editable tables)
  - manuscript/manuscript.docx (full draft with inline figures and tables)

No numeric results are hard-coded; all numbers are read from output files.
"""

from __future__ import annotations

import argparse
import re
from pathlib import Path
from typing import Any, Dict, List, Tuple

import pandas as pd
import yaml
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt

from simulate import choose_summary_rate, find_capacity_threshold


def _ppv(tp: float, fp: float) -> float:
    total = tp + fp
    return tp / total if total > 0 else 0.0


def _fp_to_tp_ratio(tp: float, fp: float) -> float:
    return fp / tp if tp > 0 else float("inf")


def load_aggregate(output_dir: Path) -> pd.DataFrame:
    return pd.read_csv(output_dir / "aggregate_by_followup.csv")


def load_by_cancer(output_dir: Path) -> pd.DataFrame:
    return pd.read_csv(output_dir / "by_cancer_and_followup.csv")


def load_weighted_ppv(output_dir: Path) -> pd.DataFrame:
    return pd.read_csv(output_dir / "weighted_ppv_by_distribution.csv")


def load_capacity_summary(output_dir: Path) -> Dict[str, Any]:
    path = output_dir / "specialist_capacity_summary.yaml"
    if path.exists():
        with open(path, "r", encoding="utf-8") as f:
            return yaml.safe_load(f) or {}
    return {}


def parse_capacity_threshold(agg: pd.DataFrame) -> Tuple[float, str]:
    """Return the first follow-up rate at which capacity exceeds 100% and the bottleneck resource."""
    s = find_capacity_threshold(agg, threshold=100.0)
    # find_capacity_threshold returns a string like "40% (Specialist Total Visits)"
    # or "not reached in the scanned range".
    if s.startswith("not"):
        return float("nan"), "not reached"
    m = re.search(r"([0-9]+)%\s*\(([^)]+)\)", s)
    if not m:
        return float("nan"), "unknown"
    rate = int(m.group(1)) / 100.0
    resource = m.group(2).strip().lower()
    return rate, resource


def compute_per_cancer_at_followup(by_cancer: pd.DataFrame, follow_up: float) -> pd.DataFrame:
    sub = by_cancer[by_cancer["follow_up_rate"] == follow_up].copy()
    sub["ppv_pct"] = (
        sub["true_positives"] / (sub["true_positives"] + sub["false_positives"]) * 100.0
    )
    return sub


def make_table_1(params: Dict[str, Any]) -> List[List[str]]:
    """Table 1: Data sources and key scenario assumptions."""
    cap = params["capacity"]
    table = [
        ["Parameter", "Value", "Source / assumption"],
        ["Screened population", f"{params['simulation']['screened_population']:,}", "Model cohort"],
        ["Base sensitivity", f"{params['cancers'][0]['sensitivity']:.2f}", params["data_sources"]["test_performance"]["source"][0]],
        ["Base specificity", f"{params['cancers'][0]['specificity']:.3f}", params["data_sources"]["test_performance"]["source"][0]],
        ["CT capacity per 100k per year", f"{cap['ct_exams_per_year']:.0f}", params["data_sources"]["diagnostic_capacity"]["source"]],
        ["MRI capacity per 100k per year", f"{cap['mri_exams_per_year']:.0f}", params["data_sources"]["diagnostic_capacity"]["source"]],
        ["Endoscopy capacity per 100k per year", f"{cap['endoscopy_exams_per_year']:.0f}", params["data_sources"]["diagnostic_capacity"]["source"]],
        ["Specialist capacity per 100k per year", f"{cap['specialist_visits_per_year']:.0f}", cap.get("specialist_capacity_note", "Illustrative scenario assumption")],
        ["Primary care capacity per 100k per year", f"{cap.get('primary_care_visits_per_year', 0.0):.0f}", cap.get("primary_care_capacity_note", "Illustrative scenario assumption")],
    ]
    for cancer in params["cancers"]:
        table.append(
            [
                f"{cancer['name']} prevalence (per 100k)",
                f"{cancer['prevalence_per_100k']:.1f}",
                params["data_sources"]["cancer_incidence"]["source"],
            ]
        )
    return table


def make_table_2(by_cancer_at_50: pd.DataFrame, weighted_ppv: pd.DataFrame) -> List[List[str]]:
    """Table 2: Per-cancer outcomes at 50% follow-up."""
    table = [
        [
            "Cancer",
            "Prevalence per 100k",
            "True positives",
            "False positives",
            "PPV (%)",
            "FP/TP ratio",
            "Primary care visits",
            "Total visits",
            "Max resource utilisation (%)",
        ]
    ]
    weighted_counts = weighted_ppv.set_index("cancer")[["true_positives", "false_positives"]]
    for _, row in by_cancer_at_50.iterrows():
        cancer = row["cancer"]
        if cancer in weighted_counts.index:
            tp = float(weighted_counts.loc[cancer, "true_positives"])
            fp = float(weighted_counts.loc[cancer, "false_positives"])
        else:
            tp = float(row["true_positives"])
            fp = float(row["false_positives"])
        ppv = _ppv(tp, fp)
        fp_tp = _fp_to_tp_ratio(tp, fp)
        table.append(
            [
                row["cancer"],
                f"{row['prevalence_per_100k']:.1f}",
                f"{row['true_positives']:.1f}",
                f"{row['false_positives']:.1f}",
                f"{ppv * 100:.2f}",
                f"{fp_tp:.1f}",
                f"{row.get('primary_care_visits', 0.0):.1f}",
                f"{row['total_visits']:.1f}",
                f"{row['max_capacity_utilization_pct']:.1f}",
            ]
        )
    return table


def make_table_3(agg: pd.DataFrame) -> List[List[str]]:
    """Table 3: Aggregate outcomes by follow-up rate."""
    table = [
        [
            "Follow-up rate",
            "True positives",
            "False positives",
            "Total positives",
            "PPV (%)",
            "FP/TP ratio",
            "Primary care visits",
            "Total visits",
            "Max capacity utilisation (%)",
        ]
    ]
    for _, row in agg.iterrows():
        ppv = _ppv(row["true_positives"], row["false_positives"])
        fp_tp = _fp_to_tp_ratio(row["true_positives"], row["false_positives"])
        table.append(
            [
                f"{row['follow_up_rate']:.0%}",
                f"{row['true_positives']:.1f}",
                f"{row['false_positives']:.1f}",
                f"{row['total_positives']:.1f}",
                f"{ppv * 100:.2f}",
                f"{fp_tp:.1f}",
                f"{row.get('primary_care_visits', 0.0):.1f}",
                f"{row['total_visits']:.1f}",
                f"{row['max_capacity_utilization_pct']:.1f}",
            ]
        )
    return table


def make_table_4(capacity_impact: pd.DataFrame) -> List[List[str]]:
    """Table 4: Baseline cases per specialist and the incremental FP burden at 50% follow-up."""
    table = [
        [
            "Cancer",
            "Relevant specialty",
            "Baseline cases per specialist",
            "False-positive specialist visits per specialist",
            "Cases per specialist with FP",
            "Increase (%)",
        ]
    ]
    for _, row in capacity_impact.iterrows():
        table.append(
            [
                str(row["cancer"]),
                str(row["specialty"]),
                f"{row['baseline_cases_per_specialist']:.1f}",
                f"{row['fp_visits_per_specialist']:.2f}",
                f"{row['cases_per_specialist_with_fp']:.1f}",
                f"{row['percent_change_in_cases_per_specialist']:.1f}",
            ]
        )
    return table


def build_markdown(
    params: Dict[str, Any],
    agg: pd.DataFrame,
    by_cancer_at_50: pd.DataFrame,
    weighted_ppv: pd.DataFrame,
    capacity_impact: pd.DataFrame,
    capacity_summary: Dict[str, Any],
    output_dir: Path,
) -> str:
    """Generate the Markdown manuscript body."""
    follow_up_rates = agg["follow_up_rate"].tolist()
    row_50 = agg[agg["follow_up_rate"] == choose_summary_rate(follow_up_rates, 0.5)].iloc[0]

    threshold_rate, threshold_resource = parse_capacity_threshold(agg)

    # Recompute PPV and FP/TP ratio from counts so prose matches tables exactly.
    row_50_ppv = _ppv(row_50["true_positives"], row_50["false_positives"]) * 100
    row_50_fp_tp = _fp_to_tp_ratio(row_50["true_positives"], row_50["false_positives"])

    weighted_ppv = weighted_ppv.copy()
    weighted_ppv["ppv"] = weighted_ppv.apply(
        lambda r: _ppv(r["true_positives"], r["false_positives"]), axis=1
    )
    weighted_ppv["fp_to_tp_ratio"] = weighted_ppv.apply(
        lambda r: _fp_to_tp_ratio(r["true_positives"], r["false_positives"]), axis=1
    )
    if pd.isna(threshold_rate):
        threshold_str = "not reached in the scanned follow-up range"
    else:
        threshold_str = f"{threshold_rate:.0%} ({threshold_resource})"

    # Specificity sweep PPVs at 95% and 99%.
    sweep = pd.read_csv(output_dir / "specificity_sweep.csv")
    sweep_agg = (
        sweep.groupby("sweep_specificity")[["true_positives", "false_positives"]]
        .sum()
        .reset_index()
    )
    sweep_agg["ppv"] = sweep_agg["true_positives"] / (
        sweep_agg["true_positives"] + sweep_agg["false_positives"]
    )
    spec_values = sweep_agg["sweep_specificity"].tolist()
    spec_95 = choose_summary_rate(spec_values, 0.95)
    spec_99 = choose_summary_rate(spec_values, 0.99)
    ppv_95 = sweep_agg[sweep_agg["sweep_specificity"] == spec_95]["ppv"].values[0] * 100
    ppv_99 = sweep_agg[sweep_agg["sweep_specificity"] == spec_99]["ppv"].values[0] * 100
    ppv_ratio = ppv_95 / ppv_99 if ppv_99 > 0 else 0.0

    lowest_ppv = weighted_ppv.loc[weighted_ppv["ppv"].idxmin()]
    highest_ppv = weighted_ppv.loc[weighted_ppv["ppv"].idxmax()]
    weighted_ppv_min_fp_tp = weighted_ppv["fp_to_tp_ratio"].min()
    weighted_ppv_max_fp_tp = weighted_ppv["fp_to_tp_ratio"].max()

    cap = params["capacity"]

    baseline_cases_per_spec = capacity_summary.get("baseline_cases_per_specialist", 0.0)
    fp_visits_per_spec = capacity_summary.get("fp_visits_per_specialist", 0.0)
    percent_change = capacity_summary.get("percent_change", 0.0)

    references = [
        "1. National Cancer Center of Japan. Cancer Statistics in Japan 2016-2023. https://ganjoho.jp/reg_stat/statistics/data/dl/en.html",
        "2. Ministry of Health, Labour and Welfare. 2023 Medical Facility Survey (Static/Dynamic). https://www.mhlw.go.jp/toukei/saikin/hw/iryosd/23/",
        "3. Kahwati LC, et al. Blood-Based Tests for Multiple Cancer Screening: A Systematic Review. AHRQ; 2025. https://www.ncbi.nlm.nih.gov/books/NBK618307/",
        "4. Schnabel JL, et al. Predictive Performance of Cell-Free Nucleic Acid-Based Multi-Cancer Early Detection Tests: A Systematic Review. PubMed; 2024. https://pubmed.ncbi.nlm.nih.gov/37791504/",
        "5. Nagamachi S, et al. Nationwide PET/CT facility survey on N-NOSE-triggered examinations. J Nucl Med Technol (Japanese report), 2024. https://jcpet.jp/2024/10/senchu-chosa.html",
        "6. Ministry of Health, Labour and Welfare. Patient Survey 2023 (r05syobyo.pdf). https://www.mhlw.go.jp/toukei/saikin/hw/kanja/10syoubyo/",
        "7. Ministry of Health, Labour and Welfare. NDB Open Data 11th release (April 2024–March 2025). https://www.mhlw.go.jp/stf/seisakunitsuite/bunya/0000177179.html",
        "8. Japanese Board of Medical Specialties. Summary of the Japanese specialist system, as tabulated by senkoi.net. https://senkoi.net/650/ and https://senkoi.net/778/",
    ]

    md = f"""# False-positive cascade and healthcare capacity burden of direct-to-consumer multi-cancer early detection blood tests: a scenario modelling study

## Abstract

**Background:** Direct-to-consumer (DTC) blood-based multi-cancer early detection (MCED) tests are marketed as a simple alternative to organised screening, but their positive predictive value (PPV) is low in asymptomatic populations and each screen-positive person may trigger multiple confirmatory examinations.

**Methods:** We built a deterministic expected-value model parameterised with 2023 adult cancer incidence rates, 2023 population counts, and 2023 national medical-facility diagnostic volumes. Specialist capacity was additionally derived from NDB Open Data outpatient patient counts and JMSB specialist counts. We estimated true positives, false positives, downstream visits, capacity utilisation, and the per-specialist case load across follow-up rates of 0–100% and specificities of 95.0–99.9%.

**Results:** At 50% follow-up, a screening wave of 100,000 persons would generate {row_50['true_positives']:.1f} true positives and {row_50['false_positives']:.1f} false positives (overall PPV {row_50_ppv:.2f}%; FP/TP ratio {row_50_fp_tp:.1f}). Total downstream visits would reach {row_50['total_visits']:.1f}, with maximum capacity utilisation {row_50['max_capacity_utilization_pct']:.1f}% (specialist visits). The first illustrative capacity ceiling was exceeded at a follow-up rate of {threshold_str}. PPV ranged from {lowest_ppv['ppv']*100:.2f}% ({lowest_ppv['cancer']}) to {highest_ppv['ppv']*100:.2f}% ({highest_ppv['cancer']}) across cancer types and was strongly age-dependent. Relative to the MHLW Patient Survey 2023 baseline case load, a 100,000-person DTC wave at 50% follow-up would add {fp_visits_per_spec:.1f} false-positive specialist visits per relevant specialist, raising the effective cases per specialist from {baseline_cases_per_spec:.1f} to {baseline_cases_per_spec + fp_visits_per_spec:.1f} ({percent_change:.0f}% increase).

**Conclusions:** Even with optimistic 99% specificity, a DTC MCED wave can trigger a false-positive cascade that exceeds outpatient and endoscopic capacity. Regulatory guardrails on performance claims, follow-up obligations, and reporting are needed before routine adoption.

**Keywords:** multi-cancer early detection, false positive, healthcare capacity, direct-to-consumer testing, scenario model

---

## Introduction

Blood-based multi-cancer early detection (MCED) tests are increasingly advertised directly to consumers as a convenient “single blood draw” cancer screen [^1^]. Because most positive results in asymptomatic populations are false positives, each abnormal test generates a cascade of confirmatory imaging and specialist visits [^3^][^4^]. In this setting, where endoscopy and specialist visits are already constrained [^2^], widespread DTC use could displace routine care. We quantified this burden as a function of follow-up behaviour, test specificity, and age structure.

## Methods

### Data sources

Cancer incidence by site, age, sex, and calendar year (2023) and the corresponding 2023 population by age and sex were taken from the National Cancer Center of Japan [^1^]. Annual volumes of CT, MRI, and upper/lower gastrointestinal endoscopies were derived from the 2023 Ministry of Health, Labour and Welfare Medical Facility Survey [^2^]. Test sensitivity and specificity ranges were informed by two recent systematic reviews of blood-based MCED tests [^3^][^4^], and real-world PPV evidence came from a nationwide PET/CT facility survey of N-NOSE-triggered examinations [^5^].

Specialist capacity was defined using NDB Open Data first/revisit outpatient patient counts and Japanese Board of Medical Specialties specialist counts [^7^][^8^]. Disease-specific baseline patient numbers for the case-per-specialist ratio were taken from the MHLW 2023 Patient Survey [^6^].

### Model

We used a deterministic expected-value cohort model. For each cancer \(c\):

- Actual cases = screened population × prevalence per 100,000 / 100,000.
- True positives = actual cases × sensitivity.
- False positives = (screened population − actual cases) × (1 − specificity).

Each positive individual who followed up (follow-up rate, 0–100%) generated a primary care visit plus visits to CT, MRI, endoscopy, and specialist care according to cancer-specific pathway probabilities. Additional visits per true and false positive were added. Capacity utilisation for each resource was calculated as total visits divided by the annual capacity per 100,000 population. A full description of equations is available in `simulate.py`.

Prevalence was approximated by adult (20+) incidence because point prevalence of undiagnosed, screen-detectable cancers is not publicly reported. This is a conservative lower-bound for true positives and therefore an upper-bound for PPV and FP/TP ratios. Pathway probabilities and the share of facility capacity available for a new DTC-related wave were scenario assumptions, documented in `parameters.yaml`.

### Specialist and primary care capacity definition

Baseline specialist capacity is defined as the annual outpatient caseload per cancer-relevant specialist. NDB Open Data unique first/revisit outpatient patient counts were used (April 2024–March 2025) divided by the total number of basic JMSB specialists, giving an average annual caseload per specialist [^7^][^8^]. This value was then multiplied by the number of cancer-relevant specialists per 100,000 population and applied the same 20% share assumed available for a DTC wave. The resulting `specialist_visits_per_year` is an illustrative capacity ceiling for a 100,000-person cohort.

The same approach was applied to primary care (internal medicine and general practice) specialists to derive `primary_care_visits_per_year`.

### Scenarios

Base-case sensitivity and specificity were {params['cancers'][0]['sensitivity']:.2f} and {params['cancers'][0]['specificity']:.3f}. We varied follow-up rate from 0 to 100% and specificity from 0.950 to 0.999 in a sensitivity sweep. The available-for-cancer-workup share of national diagnostic capacity was set to {params['assumptions']['available_for_cancer_share']:.0%}.

## Results

### Scenario parameters

Table 1 summarises the data sources and base-case parameter values.

**Table 1. Data sources and scenario parameters.**

{format_markdown_table(make_table_1(params))}

### Per-cancer burden at 50% follow-up

At a 50% follow-up rate, the model estimated {row_50['true_positives']:.1f} true positives and {row_50['false_positives']:.1f} false positives across all eight cancers (Table 2). {highest_ppv['cancer']} had the highest age-adjusted PPV ({highest_ppv['ppv']*100:.2f}%) and {lowest_ppv['cancer']} the lowest ({lowest_ppv['ppv']*100:.2f}%). The FP/TP ratio ranged from {weighted_ppv_min_fp_tp:.1f} to {weighted_ppv_max_fp_tp:.1f}.

**Table 2. Per-cancer outcomes at 50% follow-up (per 100,000 screened).**

{format_markdown_table(make_table_2(by_cancer_at_50, weighted_ppv))}

### Capacity impact

Total downstream visits rose from {agg[agg['follow_up_rate']==0.0]['total_visits'].iloc[0]:.1f} at 0% follow-up to {agg[agg['follow_up_rate']==1.0]['total_visits'].iloc[0]:.1f} at 100% follow-up (Fig. 1). At 50% follow-up this included {row_50['primary_care_visits']:.1f} primary care visits ({row_50['primary_care_visits_utilization_pct']:.1f}% of the illustrative primary-care capacity). Resource utilisation by modality is shown in Fig. 2. The first illustrative capacity ceiling was exceeded at {threshold_str}. At 50% follow-up maximum utilisation was {row_50['max_capacity_utilization_pct']:.1f}%.

![Figure 1: Total downstream visits by follow-up rate](output/total_visits_by_followup.png)
**Fig. 1.** Total downstream diagnostic, primary care, and specialist visits generated by a blood-based MCED screening wave of 100,000 persons, by follow-up rate.

![Figure 2: Diagnostic capacity utilisation by follow-up rate](output/capacity_utilization.png)
**Fig. 2.** Capacity utilisation (%) for CT, MRI, endoscopy, specialist, and primary care visits as follow-up rate increases. Values above 100% indicate demand exceeding the illustrative annual capacity available for a DTC screening wave.

### Specialist capacity and the false-positive cascade

Table 4 compares the MHLW Patient Survey 2023 baseline cancer case load per specialist with the additional false-positive specialist visits generated by a 100,000-person DTC wave at 50% follow-up. Across all cancer-relevant specialties, the baseline case load is about {baseline_cases_per_spec:.1f} patients per specialist; the DTC wave adds about {fp_visits_per_spec:.1f} false-positive specialist visits per specialist, an increase of {percent_change:.0f}%.

**Table 4. Baseline cases per specialist and incremental false-positive burden at 50% follow-up.**

{format_markdown_table(make_table_4(capacity_impact))}

### Age-specific positive predictive value

PPV was strongly age-dependent (Fig. 3). In younger age groups PPV fell below 1% for several cancers, rising above 20% only in the oldest groups. This implies that if DTC MCED users are younger than the general screening population, aggregate PPV would be lower and the false-positive burden larger than the base-case estimate.

![Figure 3: Age-specific PPV by cancer type](output/ppv_by_age.png)
**Fig. 3.** Age-specific positive predictive value for each cancer, assuming sensitivity 0.70 and specificity 0.99.

### Sensitivity to test specificity

Lowering specificity from 99% to 95% reduced aggregate PPV to roughly {ppv_ratio:.1f} of the 99% value (from {ppv_99:.2f}% to {ppv_95:.2f}%) and dramatically increased total visits and capacity pressure (Fig. 4; Table 3).

**Table 3. Aggregate outcomes by follow-up rate (base-case specificity).**

{format_markdown_table(make_table_3(agg))}

![Figure 4: PPV and total positives across specificity values](output/specificity_sweep.png)
**Fig. 4.** Aggregate positive predictive value (%) and total positive results per 100,000 screened across specificity values at a 70% follow-up rate.

## Discussion

Our scenario model shows that, even under optimistic assumptions for test accuracy, a DTC blood-based MCED screening wave can produce roughly {row_50_fp_tp:.0f} false-positive workups for every true cancer detected. At 50% follow-up the illustrative specialist capacity ceiling is exceeded by {row_50['max_capacity_utilization_pct']-100:.1f} percentage points, and the effective case load per cancer-relevant specialist rises by about {percent_change:.0f}% due to false-positive follow-up visits. The burden is not uniform: cancers with low prevalence (ovarian, liver, pancreatic) had the lowest PPV and the highest FP/TP ratios, while age strongly modulates PPV.

These findings align with real-world data from the N-NOSE PET/CT survey, in which the cancer discovery rate after a high-risk result was low and well below the company's advertised PPV [^5^].

### Limitations

Our analysis intentionally uses scenario assumptions for test performance, diagnostic pathways, and the age distribution of DTC users, because these data are not publicly reported. Prevalence was approximated by adult incidence; true point prevalence of undiagnosed cancers may differ. Diagnostic capacity was annualised from a one-month facility survey and specialist capacity from NDB outpatient patient counts; both were reduced by an arbitrary available-for-cancer-workup share. The model is deterministic and does not capture stochastic variation, geographic maldistribution, or queueing effects.

### Conclusion

Without regulatory guardrails—clear performance thresholds, transparent PPV reporting, and follow-up obligations—direct-to-consumer MCED tests risk converting a marketing promise into a large-scale false-positive cascade that stresses diagnostic capacity.

## References

""" + "\n\n".join(references) + "\n"
    return md


def format_markdown_table(table: List[List[str]]) -> str:
    lines = ["| " + " | ".join(row) + " |" for row in table]
    lines.insert(1, "|" + "|".join(["---" for _ in table[0]]) + "|")
    return "\n".join(lines)


def _add_formatted_text(paragraph, text: str, font_size: int = 11) -> None:
    """Add text to a paragraph, applying bold and superscript citation formatting."""
    parts = re.split(r"(\[\^\d+\^\]|\*\*[^*]+\*\*)", text)
    for part in parts:
        if not part:
            continue
        run = paragraph.add_run()
        if re.fullmatch(r"\*\*[^*]+\*\*", part):
            run.text = part[2:-2]
            run.font.bold = True
        elif re.fullmatch(r"\[\^\d+\^\]", part):
            run.text = part[2:-2].strip("[^]")
            run.font.superscript = True
        else:
            run.text = part
        run.font.size = Pt(font_size)


def build_tables_docx(tables: Dict[str, List[List[str]]], docx_path: Path) -> None:
    doc = Document()
    for title, table_data in tables.items():
        doc.add_heading(title, level=2)
        t = doc.add_table(rows=1, cols=len(table_data[0]))
        t.style = "Table Grid"
        hdr_cells = t.rows[0].cells
        for i, val in enumerate(table_data[0]):
            hdr_cells[i].text = val
            for para in hdr_cells[i].paragraphs:
                for r in para.runs:
                    r.font.bold = True
                    r.font.size = Pt(10)
        for row in table_data[1:]:
            row_cells = t.add_row().cells
            for i, val in enumerate(row):
                row_cells[i].text = val
                for para in row_cells[i].paragraphs:
                    for r in para.runs:
                        r.font.size = Pt(10)
        doc.add_paragraph()
    doc.save(docx_path)


def build_figures_pptx(output_dir: Path, pptx_path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    figures = [
        ("Figure 1. Total downstream visits by follow-up rate", output_dir / "total_visits_by_followup.png"),
        ("Figure 2. Diagnostic capacity utilisation by follow-up rate", output_dir / "capacity_utilization.png"),
        ("Figure 3. Age-specific PPV by cancer type", output_dir / "ppv_by_age.png"),
        ("Figure 4. PPV and total positives across specificity values", output_dir / "specificity_sweep.png"),
        ("Supplementary Figure S1. One-way sensitivity of maximum capacity utilisation", output_dir / "tornado_max_capacity.png"),
        ("Supplementary Figure S2. One-way sensitivity of aggregate PPV", output_dir / "tornado_ppv.png"),
        ("Supplementary Figure S3. Aggregate PPV under alternative age-distribution scenarios", output_dir / "age_scenario_ppv.png"),
    ]

    for title, img_path in figures:
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
        tf = txBox.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        if img_path.exists():
            slide.shapes.add_picture(str(img_path), Inches(1.0), Inches(1.3), width=Inches(11.0))
    prs.save(pptx_path)


def markdown_to_docx(md_text: str, docx_path: Path, output_dir: Path) -> None:
    """Convert the Markdown manuscript into a single editable .docx with inline figures and tables."""
    doc = Document()
    # Set a default font for the document.
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    lines = md_text.splitlines()
    i = 0
    table_pattern = re.compile(r"^\|(.*)\|\s*$")

    while i < len(lines):
        line = lines[i].rstrip()

        if line.startswith("# ") and not line.startswith("## "):
            p = doc.add_heading(line[2:], level=0)
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            i += 1
            continue

        if line.startswith("### "):
            doc.add_heading(line[4:], level=3)
            i += 1
            continue

        if line.startswith("## "):
            doc.add_heading(line[3:], level=2)
            i += 1
            continue

        if line.startswith("!"):
            m = re.match(r"!\[([^\]]*)\]\(([^)]+)\)", line)
            if m:
                img_path = output_dir / Path(m.group(2)).name
                if img_path.exists():
                    doc.add_picture(str(img_path), width=Inches(5.5))
                    last_paragraph = doc.paragraphs[-1]
                    last_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                # The next line is typically the caption.
                if i + 1 < len(lines) and lines[i + 1].startswith("**Fig."):
                    i += 1
                    cap = doc.add_paragraph()
                    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    _add_formatted_text(cap, lines[i].strip(), font_size=10)
                    cap.paragraph_format.space_before = Pt(6)
            i += 1
            continue

        if line.startswith("**Table"):
            p = doc.add_paragraph()
            _add_formatted_text(p, line.strip(), font_size=11)
            i += 1
            continue

        if table_pattern.match(line):
            # Collect table lines.
            table_lines = []
            while i < len(lines) and table_pattern.match(lines[i]):
                table_lines.append(lines[i])
                i += 1
            # Remove header separator line.
            filtered = [l for l in table_lines if not re.match(r"^\|\s*[-:]+", l.strip())]
            if filtered:
                rows = []
                for l in filtered:
                    cells = [c.strip() for c in l.strip("|").split("|")]
                    rows.append(cells)
                if rows:
                    t = doc.add_table(rows=1, cols=len(rows[0]))
                    t.style = "Table Grid"
                    for col_idx, val in enumerate(rows[0]):
                        t.rows[0].cells[col_idx].text = val
                        for para in t.rows[0].cells[col_idx].paragraphs:
                            for r in para.runs:
                                r.font.bold = True
                                r.font.size = Pt(10)
                    for row in rows[1:]:
                        cells = t.add_row().cells
                        for col_idx, val in enumerate(row):
                            cells[col_idx].text = val
                            for para in cells[col_idx].paragraphs:
                                for r in para.runs:
                                    r.font.size = Pt(10)
                    doc.add_paragraph()
            continue

        if line == "---":
            doc.add_paragraph("\u2501" * 60)
            i += 1
            continue

        if line.startswith("- "):
            p = doc.add_paragraph(style="List Bullet")
            _add_formatted_text(p, line[2:], font_size=11)
            i += 1
            continue

        if line.strip():
            p = doc.add_paragraph()
            _add_formatted_text(p, line, font_size=11)
            i += 1
            continue

        # Blank line.
        i += 1

    doc.save(docx_path)


def main() -> None:
    parser = argparse.ArgumentParser(description="Build manuscript materials")
    parser.add_argument("--params", type=Path, default=Path("parameters.yaml"))
    parser.add_argument("--output", type=Path, default=Path("output"))
    parser.add_argument("--manuscript", type=Path, default=Path("manuscript"))
    args = parser.parse_args()

    args.manuscript.mkdir(parents=True, exist_ok=True)

    with open(args.params, "r", encoding="utf-8") as f:
        params = yaml.safe_load(f)

    agg = load_aggregate(args.output)
    by_cancer = load_by_cancer(args.output)
    weighted_ppv = load_weighted_ppv(args.output)
    by_cancer_at_50 = compute_per_cancer_at_followup(by_cancer, 0.5)
    capacity_impact = pd.read_csv(args.output / "specialist_capacity_impact.csv")
    capacity_summary = load_capacity_summary(args.output)

    md = build_markdown(
        params,
        agg,
        by_cancer_at_50,
        weighted_ppv,
        capacity_impact,
        capacity_summary,
        args.output,
    )
    (args.manuscript / "manuscript.md").write_text(md, encoding="utf-8")

    tables = {
        "Table 1. Data sources and scenario parameters": make_table_1(params),
        "Table 2. Per-cancer outcomes at 50% follow-up": make_table_2(by_cancer_at_50, weighted_ppv),
        "Table 3. Aggregate outcomes by follow-up rate": make_table_3(agg),
        "Table 4. Baseline cases per specialist and incremental FP burden": make_table_4(capacity_impact),
    }
    build_tables_docx(tables, args.manuscript / "manuscript_tables.docx")

    build_figures_pptx(args.output, args.manuscript / "manuscript_figures.pptx")

    markdown_to_docx(md, args.manuscript / "manuscript.docx", args.output)

    print(f"Manuscript materials written to {args.manuscript.resolve()}")


if __name__ == "__main__":
    main()
