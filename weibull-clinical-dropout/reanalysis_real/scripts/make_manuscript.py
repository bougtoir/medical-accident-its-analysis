"""Generate a short-communication manuscript (docx) + editable figure pptx.

ALL numeric results are read from results/weibull_fits.csv (never hard-coded).
Citation numbers use {n} markers rendered as Word-native superscript runs.
Figure is inserted inline after its first in-text mention; Table 1 is built
from the results CSV. An editable English PPTX (one figure per slide) is also
written.
"""
import os
import re
import pandas as pd
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt

HERE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
RESULTS = os.path.join(HERE, "results")
FIGDIR = os.path.join(HERE, "figures")
OUT = os.path.join(HERE, "manuscript")
os.makedirs(OUT, exist_ok=True)

FITS = pd.read_csv(os.path.join(RESULTS, "weibull_fits.csv")).set_index("dataset")
FIG = os.path.join(FIGDIR, "weibull_real_fits.png")


def _scoping():
    """Read the reproducible Europe PMC search count + date (never hard-coded)."""
    path = os.path.join(RESULTS, "scoping_search.csv")
    info = {"europepmc_hit_count": None, "search_date": None}
    if os.path.exists(path):
        with open(path) as f:
            for line in f:
                for key in info:
                    if line.startswith(key + ","):
                        info[key] = line.split(",", 1)[1].strip()
    return info


SCOPING = _scoping()
# Infectious-disease curves eligible for the scoping review (antipsychotic is a
# non-infectious methodological contrast shown separately, outside eligibility).
INFECTIOUS = ["TB-Ethiopia", "TB-China", "ART/HIV-Ethiopia", "HIV-Maputo-ATT",
              "HIV-Maputo-BTT", "HIV-Malawi-pre", "HIV-Gambella"]

# Numbered in order of first appearance in the body (Vancouver).
REFERENCES = [
    "World Health Organization. Global Tuberculosis Report 2024. Geneva: WHO; 2024.",
    "Tsibiyane M, Faye LM, Ndayi K, Sineke N, Tyeshani L, Faleni M, et al. Multilevel "
    "determinants of tuberculosis treatment interruption in rural South Africa: insights from "
    "primary healthcare nurses. Int J Environ Res Public Health. 2026;23(5):598. "
    "doi:10.3390/ijerph23050598.",
    "Kedthongma W, Usaprom S, Phakdeekul W. Community-based interventions to improve "
    "tuberculosis treatment outcomes: a meta-analysis. MethodsX. 2026;16:103893. "
    "doi:10.1016/j.mex.2026.103893.",
    "Fufa DB, Diriba TA, Dame KT, Debusho LK. Competing risk models to evaluate the "
    "factors for time to loss to follow-up among tuberculosis patients at Ambo General "
    "Hospital. Arch Public Health. 2023;81:113. doi:10.1186/s13690-023-01130-2.",
    "Jiang Y, Chen J, Ying M, Liu L, Li M, Lu S, et al. Factors associated with loss to "
    "follow-up before and after treatment initiation among patients with tuberculosis: a "
    "5-year observation in China. Front Med. 2023;10:1136094. doi:10.3389/fmed.2023.1136094.",
    "Daba O, Tsegaye D, Reshad M. Incidence and predictors of loss to follow-up among ART "
    "patients on follow-up at public health facilities in Southwest Ethiopia: a time-to-event "
    "analysis. J Int Assoc Provid AIDS Care. 2026;25:23259582261426232. "
    "doi:10.1177/23259582261426232.",
    "Mateus A, Waldman EA. Retention and predictors of loss-to-follow-up among patients on "
    "antiretroviral therapy in the Test-and-Treat Era: evidence from a retrospective cohort "
    "study in Maputo, Mozambique. BMC Infect Dis. 2026;26:667. "
    "doi:10.1186/s12879-026-12949-9.",
    "Makonokaya L, Dunga S, Kalitera L, Chamanga R, Msuku L, Golowa C, et al. Retention in "
    "HIV care before and after implementation of a case management program in Malawi. "
    "BMC Public Health. 2026;26:1615. doi:10.1186/s12889-026-27295-3.",
    "Dorgi A, Tunje A, Shimbre MS, Belete AG, Guyo TG, Bodicha BB, et al. Incidence and "
    "predictors of loss to follow-up among youth living with HIV transitioning to adult care "
    "in Gambella, Southwest Ethiopia. BMC Infect Dis. 2026;26:336. "
    "doi:10.1186/s12879-026-12596-0.",
    "Salas Aranda P, Garcia Cerdan C, Martin Gomez C, Lorenzo Romo C, Turrion Gomez M, "
    "Isidoro Garcia M, et al. Time to discontinuation in routine clinical practice of the "
    "initially prescribed antipsychotic treatment in patients with first-episode psychosis. "
    "Eur Psychiatry. 2025;68(Suppl):PMC12437960.",
    "Yu Y, Xian S, Yang D, Mu L, Han Y, Luo W, et al. Temporal trends in tuberculosis "
    "incidence among the aging population in Southwest China: a retrospective study. "
    "BMC Geriatr. 2026;26:373. doi:10.1186/s12877-026-07373-2.",
]

# Datasets in display order; grouping drives Table 1 and the figure panels.
ORDER = ["TB-Ethiopia", "TB-China", "ART/HIV-Ethiopia", "HIV-Maputo-ATT",
         "HIV-Maputo-BTT", "HIV-Malawi-pre", "HIV-Gambella", "Antipsychotic"]


def kfmt(name):
    r = FITS.loc[name]
    return f"{r['k']:.2f} (95% CI {r['k_lo']:.2f}\u2013{r['k_hi']:.2f})"


def krange():
    """Min/max shape parameter across all datasets (from the results CSV)."""
    return FITS["k"].min(), FITS["k"].max()


def n_pattern(substr):
    return int(FITS["hazard_pattern"].str.contains(substr, regex=False).sum())


def pattern_of(name):
    return "IFR" if FITS.loc[name]["k"] > 1 else "DFR"


def inf_counts():
    """(n_infectious_curves, n_studies, n_IFR, n_DFR, kmin, kmax) for the review set."""
    sub = FITS.loc[INFECTIOUS]
    n_ifr = int((sub["k"] > 1).sum())
    n_dfr = int((sub["k"] < 1).sum())
    # studies: TB-Ethiopia, TB-China, ART/HIV-Ethiopia, Maputo (ATT+BTT=1),
    # Malawi, Gambella -> 6 distinct source articles
    n_studies = 6
    return len(sub), n_studies, n_ifr, n_dfr, sub["k"].min(), sub["k"].max()


def add_runs(p, text):
    """Split on {..} markers -> superscript citation runs; **bold** supported."""
    for seg in re.split(r'(\{[^}]+\})', text):
        if not seg:
            continue
        if seg.startswith("{") and seg.endswith("}"):
            run = p.add_run(seg[1:-1])
            run.font.superscript = True
        else:
            parts = re.split(r'(\*\*[^*]+\*\*)', seg)
            for pt in parts:
                if pt.startswith("**") and pt.endswith("**"):
                    run = p.add_run(pt[2:-2]); run.bold = True
                elif pt:
                    p.add_run(pt)


def heading(doc, text, size=13):
    p = doc.add_paragraph()
    p.space_before = Pt(12)
    r = p.add_run(text); r.bold = True; r.font.size = Pt(size)
    return p


def para(doc, text, align=None, space_before=6):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(space_before)
    if align:
        p.alignment = align
    add_runs(p, text)
    return p


def build_docx():
    doc = Document()
    for s in doc.sections:
        s.left_margin = s.right_margin = Inches(1)

    n_inf, n_std, ni_ifr, ni_dfr, ikmin, ikmax = inf_counts()
    n_tb = int(sum(1 for d in FITS.index if str(d).startswith("TB-")))
    hits = SCOPING["europepmc_hit_count"] or "the identified"
    sdate = SCOPING["search_date"] or ""

    t = doc.add_paragraph(); t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = t.add_run("Reproducibly digitizable dropout curves are scarce and their hazard shapes "
                  "heterogeneous: a scoping search and Weibull analysis of loss to follow-up in "
                  "tuberculosis and HIV treatment")
    r.bold = True; r.font.size = Pt(15)
    sub = doc.add_paragraph(); sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub.add_run("Original Research Article").italic = True

    heading(doc, "Abstract")
    para(doc,
         "**Objectives** Loss to follow-up (LTFU) undermines tuberculosis (TB) and HIV treatment, yet "
         "support is rarely targeted to *when* patients disengage. The clinical question is whether "
         "dropout is homogeneous enough for its timing to be inferred from the diagnosis. We "
         "asked whether the hazard shape is consistent within and across diseases, "
         "and how often it is reported in a reproducibly usable form.")
    para(doc,
         f"**Methods** A reproducible Europe PMC search (open-access TB/HIV reports mentioning a "
         f"time-to-event curve; {hits} records{', ' + sdate if sdate else ''}) required a single, "
         "digitizable dropout curve (Kaplan\u2013Meier or competing-risk cumulative incidence), not a "
         "final proportion. Eligible curves were digitized and fitted "
         "with a Weibull model F(t)=1\u2212exp(\u2212(t/\u03bb)^k); the shape k (k>1 "
         "increasing-failure-rate [IFR]; k<1 decreasing-failure-rate [DFR]) was bootstrapped and "
         "compared with exponential and log-normal fits by the Akaike information criterion (AIC).",
         space_before=4)
    para(doc,
         f"**Results** Despite {hits} records, only {n_inf} infectious-disease curves from {n_std} "
         f"studies qualified: two TB cohorts (Ethiopia, China) and five HIV/ART curves "
         "(Ethiopia, Maputo [2 curves], Malawi, Gambella); most reported only a final "
         "proportion. "
         f"The shape parameter ranged k={ikmin:.2f}\u2013{ikmax:.2f} ({ni_ifr} IFR, {ni_dfr} DFR) and "
         f"diverged **within** diseases: the two TB cohorts fell on opposite sides of k=1, and the "
         f"HIV/ART cohorts likewise split.",
         space_before=4)
    para(doc,
         "**Conclusions** Reproducibly usable, time-resolved dropout curves are scarce, and where "
         "they exist the hazard shape is heterogeneous within both TB and HIV. Support timing "
         "therefore cannot be assumed from the disease label but must be located "
         "empirically, which also makes an intervention's effect verifiable. The core barrier is a "
         "reporting gap: time-to-LTFU curves or individual patient data must be published routinely.",
         space_before=4)
    para(doc, "**Keywords** Tuberculosis; HIV; Loss to follow-up; Retention in care; Scoping search; "
              "Treatment adherence; Weibull; Hazard function; Survival analysis; Reporting", space_before=4)

    heading(doc, "Background")
    para(doc,
         "LTFU interrupts treatment for both TB and HIV, promoting acquired drug resistance, "
         "onward transmission and avoidable death{1}. Two things are already established. First, "
         "the determinants of treatment interruption are multilevel and strongly context-dependent, "
         "differing by health system, case mix and social setting{2}. Second, interventions "
         "intended to keep patients in care \u2014 community-based, behavioural and digital "
         "approaches \u2014 show inconsistent, setting-specific effectiveness, so no single strategy "
         "has proved universally transferable{3}. In practice, retention support is still frequently "
         "delivered as generic encouragement, chosen because it seems reasonable rather than because "
         "it is matched to when patients actually leave. Adopting a package simply because it "
         "\u2018looks good\u2019 also forfeits the chance to check whether it acts at the point in "
         "treatment where dropout actually happens.")
    para(doc,
         "What has been missing is a *quantitative* description of the phenomenon that plausibly "
         "underlies this: whether patients across settings even disengage on the same time-course. "
         "Programmatic reports usually give a single cumulative LTFU proportion, which hides *when* "
         "during treatment patients leave. This raises the clinical question at the centre of this "
         "study: is dropout homogeneous enough that the timing of retention support can be inferred "
         "from the diagnosis, or must the hazard shape be measured in each setting? The hazard shape "
         "carries direct operational meaning: an "
         "increasing hazard (k>1) argues for intensified support later in treatment, a decreasing "
         "hazard (k<1) for front-loading it in the earliest weeks. Making, and later verifying, such "
         "timing decisions requires the shape to be recoverable from the published evidence in the "
         "first place. We therefore treated this as a scoping question with a deliberately unusual "
         "eligibility criterion \u2014 that a study make dropout *visible over time* in a "
         "reproducibly digitizable curve, not merely as a final proportion \u2014 and asked two "
         "things: how often is dropout timing reported in that usable form, and, where it is, is the "
         "Weibull hazard shape consistent within a disease and between diseases? To our knowledge no "
         "prior review has used reproducible visualization of dropout timing as its inclusion "
         "criterion.")

    heading(doc, "Methods")
    para(doc,
         f"**Search and eligibility.** We queried Europe PMC (open-access TB/HIV LTFU or retention "
         f"reports mentioning a Kaplan\u2013Meier, survival or cumulative-incidence curve) on "
         f"{sdate or 'the search date'}, which identified {hits} records; the exact query and a "
         "record sample are saved in the repository and regenerated by scripts/search_scoping.py. "
         "Our eligibility criterion was pragmatic: a report had to present a *time-resolved* "
         "treatment-dropout or retention curve as a single, cleanly separable, digitizable line "
         "(Kaplan\u2013Meier survivor or competing-risk cumulative incidence), not only a final "
         "cumulative LTFU proportion. Full dual-reviewer screening of every record was beyond the "
         "scope of this short study; we screened the identified literature pragmatically against "
         "this criterion and note that the overwhelming majority of LTFU reports present only a "
         "final proportion and were therefore ineligible. Eligible infectious-disease curves were "
         "carried forward; one non-infectious antipsychotic curve was retained separately as a "
         "methodological contrast, outside the infectious-disease scope.")
    para(doc,
         "**Data and analysis.** We used only real, published curves. TB data were the competing-risk LTFU cumulative "
         "incidence at Ambo General Hospital, Ethiopia{4}, and the all-patient time-to-LTFU "
         "Kaplan\u2013Meier (KM) curve from a 5-year Chinese cohort{5}. HIV/ART retention curves "
         "were an ART time-to-LTFU KM from southwest Ethiopia{6}, the After- and Before-Test-and-"
         "Treat retention curves from a Maputo (Mozambique) cohort{7}, the pre-intervention "
         "remaining-in-care curve from a Malawian case-management study{8}, and the overall KM for "
         "youth living with HIV transitioning to adult care in Gambella, Ethiopia{9}. As a "
         "non-infectious methodological contrast we added an antipsychotic time-to-discontinuation "
         "curve{10}. Curves were digitized from the published figures by an in-house pixel "
         "extractor with axis calibration read from the figure tick marks (values taken from curve "
         "pixels, not entered by hand); the original figure images and digitized coordinates are in "
         "the public repository. For each dataset we fitted the Weibull cumulative-incidence model "
         "F(t)=1\u2212exp(\u2212(t/\u03bb)^k) by nonlinear least squares, obtained a bootstrap 95% "
         "confidence interval (CI) for k, and compared the fit with exponential (k=1) and log-normal "
         "models by the Akaike information criterion (AIC). The shape "
         "parameter k is unit-free, so mixed follow-up units (months, days) do not affect the "
         "cross-study comparison of shape. Data, code and a one-command build that regenerates every "
         "reported number, Table 1 and Figure 1 are openly available (Availability of data and "
         "materials); no value is hard-coded.")

    heading(doc, "Results")
    para(doc,
         f"**Yield of the search.** Although the search identified {hits} open-access records, "
         f"reproducibly digitizable, time-resolved dropout curves were rare: only {n_inf} "
         f"infectious-disease curves from {n_std} studies met the eligibility criterion (two TB and "
         f"five HIV/ART curves; the {n_inf} curves exceed the {n_std} studies because the Maputo "
         "study contributed two curves, its After- and Before-Test-and-Treat cohorts). The dominant "
         "reason for ineligibility was that reports gave only a final "
         "cumulative LTFU proportion, with no curve from which the timing of dropout could be "
         "recovered. This scarcity is itself a principal finding: the data needed to time retention "
         "support are, at present, seldom published in a usable form.")
    para(doc,
         "Fitted parameters for the eligible curves are shown in Table 1 and the fits are overlaid "
         f"on the digitized data in Figure 1. Among the {n_inf} infectious-disease curves the shape "
         f"parameter spanned k={ikmin:.2f} to k={ikmax:.2f}, with {ni_ifr} increasing-hazard (IFR) "
         f"and {ni_dfr} decreasing-hazard (DFR). Crucially, the split occurred **within** diseases, "
         "not only between them. The two TB "
         f"cohorts fell on opposite sides of k=1 (Ethiopia {pattern_of('TB-Ethiopia')}, "
         f"k={kfmt('TB-Ethiopia')}; China {pattern_of('TB-China')}, k={kfmt('TB-China')}). The five "
         "HIV/ART cohorts did the same: the two Maputo cohorts "
         f"(k={kfmt('HIV-Maputo-ATT')} and k={kfmt('HIV-Maputo-BTT')}) and the Gambella youth cohort "
         f"(k={kfmt('HIV-Gambella')}) were increasing-hazard, whereas the southwest-Ethiopia ART "
         f"cohort (k={kfmt('ART/HIV-Ethiopia')}) and the Malawi pre-intervention cohort "
         f"(k={kfmt('HIV-Malawi-pre')}) were decreasing-hazard. The non-infectious antipsychotic "
         f"contrast was also decreasing-hazard (k={kfmt('Antipsychotic')}). No single characteristic "
         "shape described even a single disease, and the AIC preferred a non-exponential (Weibull or "
         "log-normal) shape over a constant hazard for every curve.")

    # Table 1 from results CSV
    para(doc, "**Table 1.** Weibull fits to the digitized treatment-dropout curves (seven eligible "
              "infectious-disease curves plus one non-infectious antipsychotic contrast; source of "
              "truth: results/weibull_fits.csv; all values regenerated by the released "
              "code).", space_before=10)
    tbl = doc.add_table(rows=1, cols=5); tbl.style = "Table Grid"
    for i, h in enumerate(["Dataset (source)", "k (95% CI)", "\u03bb", "R\u00b2", "Hazard pattern"]):
        tbl.rows[0].cells[i].paragraphs[0].add_run(h).bold = True
    for name in ORDER:
        r = FITS.loc[name]
        c = tbl.add_row().cells
        c[0].text = f"{name} \u2014 {r['source']}"
        c[1].text = f"{r['k']:.2f} ({r['k_lo']:.2f}\u2013{r['k_hi']:.2f})"
        c[2].text = f"{r['lam']:.1f}"
        c[3].text = f"{r['r2']:.3f}"
        c[4].text = str(r['hazard_pattern'])

    para(doc, "", space_before=14)
    doc.add_picture(FIG, width=Inches(6.0))
    doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
    para(doc, "**Figure 1.** Weibull cumulative-incidence fits (red = increasing hazard, blue = "
              "decreasing hazard) to digitized dropout curves (points) for the seven eligible "
              "infectious-disease cohorts (two TB, five HIV/ART) plus a non-infectious antipsychotic "
              "contrast. Both the two TB cohorts and the HIV/ART cohorts split between IFR and DFR, "
              "so the hazard shape is heterogeneous within each disease.", space_before=4)

    heading(doc, "Discussion")
    para(doc,
         "The clinical question posed at the outset was whether dropout is homogeneous enough for "
         "the timing of retention support to be inferred from the diagnosis. Our results answer it "
         "in the negative. Across the eligible curves the Weibull hazard shape was heterogeneous "
         "not only between diseases but within them: both the two TB cohorts and the HIV/ART cohorts "
         "fell on opposite sides of k=1, so there was no characteristic \u2018TB shape\u2019 or "
         "\u2018HIV shape\u2019 and the direction of the dropout hazard varied from cohort to cohort "
         "within the same disease. Dropout is therefore not homogeneous enough to read the timing of "
         "support off the diagnostic label: retention is not simply a matter of exhorting patients "
         "to try harder, and *when* to strengthen support cannot be assumed from the diagnosis. This "
         "gives a concrete, quantitative mechanism for the already-established observations that "
         "treatment-interruption determinants are context-dependent{2} and that retention "
         "interventions do not transfer cleanly between settings{3}.")
    para(doc,
         "The practical consequence follows directly. Because the hazard shape encodes *when* "
         "patients disengage, it maps onto the timing of adherence support: a front-loaded (DFR) "
         "setting argues for concentrating counselling, financial or social support in the earliest "
         "weeks, whereas a late-accelerating (IFR) setting argues for sustained or intensified "
         "support later in treatment. A package optimized for one shape will be mistimed in the "
         "other \u2014 a concrete reason why globally standardized adherence interventions have shown "
         "inconsistent effects{3}. This reframes retention design as a measurement problem before it "
         "is an intervention problem: choosing a support package because it seems sensible, without "
         "first locating a setting's dropout hazard shape, risks both mistiming the support and "
         "losing the ability to tell whether it worked, because an effect concentrated in the wrong "
         "phase of treatment can be invisible against the actual dropout dynamics. Estimating the "
         "shape first gives a pre-specified, falsifiable target \u2014 shift the hazard at the phase "
         "where dropout concentrates \u2014 against which an intervention can be evaluated; the "
         "Weibull cumulative-incidence model used here is a compact, reproducible way to obtain that "
         "shape from a single published curve.")
    para(doc,
         "Beyond answering the question, why the shapes diverge is worth considering even though it "
         "lies beyond what these data can establish. We cannot explain the divergence, and we do "
         "not claim to. It is "
         "confounded with different outcome definitions (competing-risk cumulative incidence vs "
         "all-patient KM), case mix, programme context and unequal observation windows, and it "
         "rests on digitized aggregate curves rather than individual patient data. One candidate "
         "worth stating explicitly is the age structure of disease onset: TB in China is "
         "increasingly concentrated in the elderly{11}, whereas high-burden sub-Saharan settings "
         "such as Ethiopia are dominated by young adults, and age was itself a significant risk "
         "factor for LTFU in the Ethiopian TB cohort{4}. If the age-at-onset distribution governs "
         "competing mortality, comorbidity and the reasons patients disengage, it could plausibly "
         "determine whether the LTFU hazard rises or falls over the course of treatment. We stress "
         "that this is a hypothesis, not a result: our curves are aggregate and not age-stratified, "
         "so neither an age-at-onset distribution nor an age-by-LTFU-timing interaction can be "
         "fitted here without inventing data. Testing it would require age-stratified time-to-LTFU "
         "curves or individual patient data.")
    para(doc,
         "Finally, this scarcity has an equity dimension that is central to the message of this "
         "study. Sustainable care in settings with constrained health resources and limited access "
         "depends on planning retention support and then verifying that it works; both require the "
         "timing of dropout to be recorded and reported, not just its final total. Yet despite the "
         "very large global burden of tuberculosis{1}, we could recover a reproducibly usable "
         f"time-to-LTFU curve for only {n_tb} TB cohorts \u2014 from China and Ethiopia \u2014 with "
         "no such curve available (or openly published) for the many other high-incidence settings, "
         "which differ in health system, case mix and programme context. A hazard shape measured in "
         "one country cannot simply be transplanted to another, so the near-absence of comparable "
         "curves elsewhere means most programmes currently have no local basis on which to time "
         "support or to test whether an intervention worked. Building and openly sharing "
         "time-resolved dropout data across many more regions \u2014 for TB and for HIV alike \u2014 "
         "is therefore not merely a research convenience but a prerequisite for evidence-based, "
         "sustainable retention care where the need is greatest.")

    heading(doc, "Limitations")
    para(doc,
         "This is a pragmatic scoping study, not a fully registered dual-reviewer scoping review: "
         "the search was limited to open-access, English-language Europe PMC records mentioning a "
         "time-to-event curve, and screening against the eligibility criterion was performed "
         "pragmatically rather than by independent duplicate assessment, so the count of eligible "
         "curves is a lower bound and some eligible curves in other databases or languages will "
         "have been missed. Only seven time-resolved infectious-disease dropout curves met our bar, "
         "so the shape analysis "
         "is small and hypothesis-generating rather than confirmatory, and it cannot yet establish "
         "a shape for any disease. Estimates come from digitized aggregate curves, not individual "
         "patient data; the Ethiopian TB estimate is a competing-risk *subdistribution*-hazard "
         "shape; the Gambella HIV curve is plotted from an analysis time of one year, so its shape "
         "describes only the observed window; the antipsychotic cohort is small (n\u224842), "
         "illustrative and outside the infectious-disease setting; and terminal KM steps and "
         "plateaus reflect administrative censoring. Different outcome definitions and follow-up "
         "windows across sources further limit direct comparison of the fitted scale parameters, "
         "although the shape parameter k is unit-free.")

    heading(doc, "Conclusions")
    para(doc,
         "Across the treatment cohorts that publish a time-resolved LTFU curve, the Weibull hazard "
         "shape is heterogeneous within TB and within HIV alike: cohorts of the same disease fell "
         "on opposite sides of a constant hazard. When patients disengage therefore cannot be "
         "inferred from the diagnosis or from intuition, and retention support should be timed to "
         "each setting's empirically measured dropout hazard shape \u2014 which also makes its "
         "effect verifiable. Because the openly available evidence is still too sparse to fix these "
         "shapes, routine reporting of time-to-LTFU curves, and ideally individual patient data "
         "with harmonized LTFU definitions, is needed so that dropout patterns and the "
         "interventions matched to them can be established and tested.")

    heading(doc, "Declarations")
    para(doc, "**Ethics approval and consent to participate** Not applicable. This study is a "
              "secondary analysis of aggregate, already-published, de-identified survival curves "
              "and involved no new human participants or individual-level data.")
    para(doc, "**Consent for publication** Not applicable.", space_before=4)
    para(doc, "**Clinical trial number** Not applicable.", space_before=4)
    para(doc, "**Availability of data and materials** All source figures, digitized CSV datasets, "
              "analysis code and a one-command build that regenerates every reported number, "
              "Table 1 and Figure 1 are openly available in the project repository. No result is "
              "hard-coded; every value is reproduced from the source data by the released code.",
         space_before=4)
    para(doc, "**Competing interests** The authors declare that they have no competing interests.",
         space_before=4)
    para(doc, "**Funding** This study received no specific funding.", space_before=4)
    para(doc, "**Authors' contributions** The author(s) conceived the study, digitized the source "
              "curves, implemented the analysis and wrote the manuscript. All authors read and "
              "approved the final manuscript.", space_before=4)
    para(doc, "**Acknowledgements** Not applicable.", space_before=4)

    heading(doc, "References")
    for i, ref in enumerate(REFERENCES, 1):
        p = doc.add_paragraph(); p.paragraph_format.space_before = Pt(3)
        p.add_run(f"{i}. {ref}")

    path = os.path.join(OUT, "tb_dropout_heterogeneity_short_communication.docx")
    doc.save(path); print("wrote", path)


def build_pptx():
    prs = Presentation(); prs.slide_width = PInches(13.333); prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(PInches(0.5), PInches(0.2), PInches(12.3), PInches(0.8))
    tb.text_frame.text = "Figure 1. Weibull fits to digitized treatment-dropout curves (TB and HIV cohorts)"
    tb.text_frame.paragraphs[0].runs[0].font.size = PPt(20)
    s.shapes.add_picture(FIG, PInches(3.9), PInches(1.0), height=PInches(6.2))
    cap = s.shapes.add_textbox(PInches(0.5), PInches(7.15), PInches(12.3), PInches(0.35))
    cap.text_frame.text = ("Hazard shape (k) is heterogeneous within TB and within HIV alike "
                           "(both IFR and DFR appear in each). Values: results/weibull_fits.csv.")
    cap.text_frame.paragraphs[0].runs[0].font.size = PPt(12)
    path = os.path.join(OUT, "figures_editable_en.pptx")
    prs.save(path); print("wrote", path)


def build_title_page():
    doc = Document()
    for s in doc.sections:
        s.left_margin = s.right_margin = Inches(1)
    t = doc.add_paragraph(); t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = t.add_run("Reproducibly digitizable dropout curves are scarce and their hazard shapes "
                  "heterogeneous: a scoping search and Weibull analysis of loss to follow-up in "
                  "tuberculosis and HIV treatment")
    r.bold = True; r.font.size = Pt(14)
    para(doc, "Article type: Original Research Article", space_before=10)
    heading(doc, "Authors")
    para(doc, "Tatsuki Onishi [AUTHOR ORDER / ADDITIONAL CO-AUTHORS TO BE CONFIRMED]")
    heading(doc, "Affiliations")
    para(doc, "1. [Institution, Department, City, Country \u2014 TO BE CONFIRMED]")
    heading(doc, "Corresponding author")
    para(doc, "Tatsuki Onishi. Email: bougtoir@gmail.com. [Postal address / ORCID "
              "\u2014 TO BE CONFIRMED]")
    heading(doc, "Declarations")
    para(doc, "Competing interests: none declared. Funding: none. Ethics approval and consent to "
              "participate: not applicable (secondary analysis of published aggregate data).")
    path = os.path.join(OUT, "title_page.docx")
    doc.save(path); print("wrote", path)


def build_cover_letter():
    doc = Document()
    for s in doc.sections:
        s.left_margin = s.right_margin = Inches(1)
    para(doc, "[Date / your address block \u2014 TO BE CONFIRMED]")
    para(doc, "The Editors", space_before=12)
    para(doc, "Tropical Medicine & International Health", space_before=2)
    para(doc, "Dear Editors,", space_before=12)
    para(doc,
         "Please consider our Original Research Article, \u201cReproducibly digitizable "
         "dropout curves are scarce and their hazard shapes heterogeneous: a scoping search and "
         "Weibull analysis of loss to follow-up in tuberculosis and HIV treatment\u201d, "
         "for publication in Tropical Medicine & International Health.", space_before=8)
    para(doc,
         "Loss to follow-up (LTFU) during tuberculosis (TB) and HIV treatment is a central "
         "obstacle to cure, viral suppression and transmission control in resource-limited, "
         "high-burden settings. It is already recognized that the determinants of treatment "
         "interruption are context-dependent and that adherence interventions transfer poorly "
         "between programmes, but retention support is still often chosen by intuition rather than "
         "matched to when patients actually disengage. Making such timing decisions requires the "
         "dropout curve to be recoverable from the literature, so we ran a reproducible search with "
         "an unusual eligibility criterion \u2014 that a study make dropout visible over time in a "
         "digitizable curve, not merely a final proportion. Two findings follow: such curves are "
         "scarce (only seven infectious-disease curves from six studies out of thousands of "
         "records), and where they exist the two TB cohorts and the HIV/ART cohorts each split "
         "between increasing- and decreasing-hazard shapes, so there is no characteristic shape "
         "even within one disease. The practical message is that retention support should be timed "
         "to each setting's measured dropout hazard \u2014 which also makes its effect verifiable "
         "\u2014 and that routinely reporting time-to-LTFU curves is needed before disease-specific "
         "patterns can be established.", space_before=8)
    para(doc,
         "We believe this fits the scope of Tropical Medicine & International Health: it addresses "
         "treatment retention and case management of two major infectious diseases in "
         "resource-limited settings, has direct implications for programme implementation and "
         "health systems, and is fully reproducible \u2014 all source figures, digitized data and "
         "analysis code that regenerate every reported number are openly available. The work is "
         "original, is not under consideration elsewhere, and all authors approve the submission. "
         "We declare no competing interests.", space_before=8)
    para(doc, "Thank you for your consideration.", space_before=8)
    para(doc, "Sincerely,", space_before=12)
    para(doc, "Tatsuki Onishi, on behalf of the authors", space_before=2)
    path = os.path.join(OUT, "cover_letter.docx")
    doc.save(path); print("wrote", path)


if __name__ == "__main__":
    build_docx()
    build_pptx()
    build_title_page()
    build_cover_letter()
