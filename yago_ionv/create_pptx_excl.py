"""Create editable English PPTX with figures from exclusion sensitivity analysis."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pathlib import Path

BASE = Path(__file__).resolve().parent
FIG = BASE / "figures_excl"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

figures = [
    ("fig1_rates_comparison.png",
     "Fig. 1: IONV Rates — Elective Low-Risk Subgroup"),
    ("fig2_forest_narrow_primary.png",
     "Fig. 2: Logistic Regression — 5-HT3 Antagonist Use (Elective Low-Risk)"),
    ("fig3_covariate_sensitivity.png",
     "Fig. 3: Covariate Sensitivity Analysis (Elective Low-Risk)"),
    ("fig4_broad_vs_narrow.png",
     "Fig. 4: Twin Effect — Broad vs Narrow Definition (Elective Low-Risk)"),
]

for fname, title in figures:
    fpath = FIG / fname
    if not fpath.exists():
        continue
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    slide.shapes.add_picture(str(fpath), Inches(0.5), Inches(1.0), Inches(12), Inches(6))

out = BASE / "figures_excl_sensitivity.pptx"
prs.save(str(out))
print(f"PPTX saved: {out}")
