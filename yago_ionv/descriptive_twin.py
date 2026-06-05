#!/usr/bin/env python3
"""
Descriptive Statistics for Twin Cesarean Section Dataset
========================================================
Twin-only subset from the IONV study (2014-04-01 to 2024-10-23).
Produces comprehensive descriptive statistics tables, figures, and a summary docx.

Outputs
-------
- tables_twin/   : CSV tables
- figures_twin/  : PNG figures
- figures_twin.pptx : Editable figures (1 per slide)
- descriptive_twin.docx : Summary document with inline tables and figures
- descriptive_twin_stats.json : Machine-readable summary
"""

import warnings
warnings.filterwarnings("ignore")

import json
import os
import re
from pathlib import Path

import numpy as np
import pandas as pd
from scipy import stats
from statsmodels.stats.proportion import proportion_confint
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import seaborn as sns

plt.rcParams.update({
    "font.size": 10,
    "figure.dpi": 300,
    "savefig.dpi": 300,
    "savefig.bbox": "tight",
    "font.family": "DejaVu Sans",
})
sns.set_style("whitegrid")
sns.set_context("paper", font_scale=1.2)

BASE = Path(__file__).resolve().parent
FIG = BASE / "figures_twin"
TAB = BASE / "tables_twin"
FIG.mkdir(exist_ok=True)
TAB.mkdir(exist_ok=True)

# ============================================================
# 1. LOAD TWIN DATA
# ============================================================
print("=" * 60)
print("1. LOADING TWIN DATA")
print("=" * 60)

twin_file = [f for f in os.listdir(BASE) if f.startswith("twin") and f.endswith(".xlsx")][0]
raw_twin = pd.read_excel(BASE / twin_file, engine="openpyxl")
print(f"Twin raw: {len(raw_twin)} rows, {len(raw_twin.columns)} columns")

# ============================================================
# 2. RENAME & HARMONIZE COLUMNS
# ============================================================
print("\n" + "=" * 60)
print("2. HARMONIZING COLUMNS")
print("=" * 60)

twin = raw_twin.copy()
twin.rename(columns={
    "手術時間(min)": "手術時間_min",
    "麻酔時間(min)": "麻酔時間_min",
    "出血量(ml)": "出血量_ml",
    "輸液量(ml)": "輸液量_ml",
    "制吐薬投与の有無": "antiemetic_any_str",
    "入室〜麻酔開始": "ae_pre_anesthesia_str",
    "麻酔開始〜胎児娩出": "ae_to_delivery_str",
    "胎児娩出〜退室": "ae_post_delivery_str",
    "低血圧(回)": "hypotension_count",
    "メトクロプラミド(mg)": "metoclopramide_mg",
    "ドロペリドール(mg)": "droperidol_mg",
    "オンダンセトロン(mg)": "ondansetron_mg",
    "グラニセトロン(mg)": "granisetron_mg",
    "ノバミン(mg)": "novamin_mg",
    "アタラックスP(mg)": "atarax_p_mg",
    "デカドロン(mg)": "dexamethasone_mg",
    "術前24時間以内の制吐薬投与": "preop_antiemetic",
    "昇圧持続投与": "vasopressor_continuous",
    "エフェドリン(mg)": "ephedrine_mg",
    "フェニレフリン(mg)": "phenylephrine_mg",
    "脊髄くも膜下への高比重ブピバカイン(mg)": "bupivacaine_mg",
    "脊髄くも膜下へのフェンタニル(μg )": "fentanyl_ug",
    "脊髄くも膜下へのモルヒネ(mg)": "morphine_mg",
    "ディプリバン(mg)": "propofol_mg",
    "ミダゾラム(mg)": "midazolam_mg",
    "デクスメデトミジン": "dexmedetomidine_ug",
    "メモ": "exclusion_note",
    "術前1週間以内のステロイド投与": "preop_steroid",
}, inplace=True)

# Convert 有/無 text to numeric
you_mu_map = {"有": 1, "無": 0}
for src, dst in [("antiemetic_any_str", "antiemetic_any"),
                 ("ae_pre_anesthesia_str", "ae_pre_anesthesia"),
                 ("ae_to_delivery_str", "ae_to_delivery"),
                 ("ae_post_delivery_str", "ae_post_delivery")]:
    if src in twin.columns:
        twin[dst] = twin[src].map(you_mu_map).astype(float)

for col in ["帝王切開の既往", "vasopressor_continuous", "高血圧合併妊娠", "妊娠高血圧症候群",
            "術前24時間以内の降圧薬使用", "preop_antiemetic", "preop_steroid"]:
    if col in twin.columns and not pd.api.types.is_numeric_dtype(twin[col]):
        twin[col] = twin[col].map(you_mu_map).astype(float)

# Emergency from 緊急適応疾患
if "緊急適応疾患" in twin.columns:
    twin["emergency"] = twin["緊急適応疾患"].apply(
        lambda x: 0 if pd.isna(x) or str(x).strip() == "無" else 1
    )
else:
    twin["emergency"] = 0

# Parse gestational age "36w2d" → float weeks
def parse_ga_weeks(val):
    if pd.isna(val):
        return np.nan
    s = str(val).strip()
    m = re.match(r"(\d+)w(\d+)d?", s)
    if m:
        return int(m.group(1)) + int(m.group(2)) / 7
    try:
        return float(s)
    except ValueError:
        return np.nan

twin["GA_weeks"] = twin["妊娠週数(週)"].apply(parse_ga_weeks)

# Numeric conversions
twin["手術時間_min"] = pd.to_numeric(twin["手術時間_min"], errors="coerce")
twin["麻酔時間_min"] = pd.to_numeric(twin["麻酔時間_min"], errors="coerce")
twin["出血量_ml"] = pd.to_numeric(twin["出血量_ml"], errors="coerce")

for col in ["年齢(歳)", "身長(cm)", "体重(kg)", "hypotension_count",
            "全身麻酔", "硬膜外麻酔", "脊髄くも膜下麻酔",
            "帝王切開の既往", "高血圧合併妊娠", "妊娠高血圧症候群",
            "antiemetic_any", "ae_pre_anesthesia", "ae_to_delivery", "ae_post_delivery",
            "preop_antiemetic", "vasopressor_continuous"]:
    if col in twin.columns:
        twin[col] = pd.to_numeric(twin[col], errors="coerce")

# BMI
h = pd.to_numeric(twin["身長(cm)"], errors="coerce") / 100
w = pd.to_numeric(twin["体重(kg)"], errors="coerce")
twin["BMI"] = w / (h ** 2)

# Surgery date & year
twin["手術日"] = pd.to_datetime(twin["手術日"], format="mixed", errors="coerce")
twin["year"] = twin["手術日"].dt.year

# Derived variables
twin["HDP"] = ((twin["高血圧合併妊娠"] == 1) | (twin["妊娠高血圧症候群"] == 1)).astype(int)
twin["hypotension"] = (twin["hypotension_count"] >= 1).astype(int).where(
    twin["hypotension_count"].notna(), np.nan
)
twin["epidural"] = pd.to_numeric(twin["硬膜外麻酔"], errors="coerce")
twin["prior_cs"] = pd.to_numeric(twin["帝王切開の既往"], errors="coerce")
twin["bmi_ge35"] = (twin["BMI"] >= 35).astype(int).where(twin["BMI"].notna(), np.nan)

# Uterine exteriorization: 子宮脱転 (0,1,2=unknown)
if "子宮脱転" in twin.columns:
    twin["uterine_exteriorization"] = twin["子宮脱転"].replace(2, np.nan)

# IONV outcomes
twin["ionv_primary"] = ((twin["ae_to_delivery"] == 1) | (twin["ae_post_delivery"] == 1)).astype(int)
twin["ionv_secondary"] = (twin["ae_to_delivery"] == 1).astype(int)

# 5-HT3 antagonists (Definition E)
for col in ["ondansetron_mg", "granisetron_mg"]:
    twin[col] = pd.to_numeric(twin[col], errors="coerce").fillna(0)
twin["serotonin_antagonist"] = ((twin["ondansetron_mg"] > 0) | (twin["granisetron_mg"] > 0)).astype(int)

print(f"Columns after harmonization: {len(twin.columns)}")

# ============================================================
# 3. APPLY EXCLUSION CRITERIA
# ============================================================
print("\n" + "=" * 60)
print("3. APPLYING EXCLUSION CRITERIA")
print("=" * 60)

n_raw = len(twin)
exclusions = []

ga_mask = twin["全身麻酔"] == 1
exclusions.append(("General anesthesia", ga_mask.sum()))

sbp_mask = twin["exclusion_note"].str.contains("SBP90", na=False) | \
           twin["exclusion_note"].str.contains("入室時SBP", na=False)
exclusions.append(("Pre-anesthesia SBP < 90 mmHg", sbp_mask.sum()))

iufd_mask = twin["exclusion_note"].str.contains("胎児死亡|死亡", na=False) & \
            ~twin["exclusion_note"].str.contains("全身麻酔", na=False)
exclusions.append(("Intrauterine fetal death", iufd_mask.sum()))

vt_mask = twin["exclusion_note"].str.contains("vanishing", case=False, na=False)
exclusions.append(("Vanishing twin", vt_mask.sum()))

triplet_mask = twin["exclusion_note"].str.contains("品胎", na=False)
exclusions.append(("Triplet pregnancy", triplet_mask.sum()))

generic_exclude_mask = twin["exclusion_note"].str.contains("除外", na=False)
all_exclude = generic_exclude_mask | ga_mask | sbp_mask | iufd_mask | vt_mask | triplet_mask

no_data_mask = twin["antiemetic_any"].isna() & ~all_exclude
exclusions.append(("No anesthesia data available", no_data_mask.sum()))

all_exclude = all_exclude | no_data_mask

n_excluded = all_exclude.sum()
print(f"Total twin before exclusion: {n_raw}")
for reason, n in exclusions:
    print(f"  {reason}: {n}")
print(f"  Total excluded: {n_excluded}")

df = twin[~all_exclude].copy()

# Further exclude pre-anesthesia antiemetic users (per protocol)
pre_ae_count = (df["ae_pre_anesthesia"] == 1).sum()
print(f"Pre-anesthesia antiemetic users: {pre_ae_count}")
df = df[df["ae_pre_anesthesia"] != 1].copy()
print(f"Analysis cohort: {len(df)}")

# ============================================================
# 4. DESCRIPTIVE STATISTICS — CONTINUOUS VARIABLES
# ============================================================
print("\n" + "=" * 60)
print("4. DESCRIPTIVE STATISTICS — CONTINUOUS VARIABLES")
print("=" * 60)

continuous_vars = [
    ("年齢(歳)", "Age (years)"),
    ("BMI", "BMI (kg/m²)"),
    ("GA_weeks", "Gestational age (weeks)"),
    ("手術時間_min", "Surgery time (min)"),
    ("麻酔時間_min", "Anesthesia time (min)"),
    ("出血量_ml", "Estimated blood loss (mL)"),
    ("輸液量_ml", "Infusion volume (mL)"),
    ("hypotension_count", "Hypotensive episodes (count)"),
    ("bupivacaine_mg", "Intrathecal bupivacaine (mg)"),
    ("fentanyl_ug", "Intrathecal fentanyl (μg)"),
    ("morphine_mg", "Intrathecal morphine (mg)"),
    ("ephedrine_mg", "Ephedrine (mg)"),
    ("phenylephrine_mg", "Phenylephrine (mg)"),
]

cont_rows = []
for var, label in continuous_vars:
    if var not in df.columns:
        continue
    s = pd.to_numeric(df[var], errors="coerce").dropna()
    n = len(s)
    if n == 0:
        continue
    row = {
        "Variable": label,
        "n": n,
        "Missing": len(df) - n,
        "Mean": s.mean(),
        "SD": s.std(),
        "Median": s.median(),
        "Q1": s.quantile(0.25),
        "Q3": s.quantile(0.75),
        "Min": s.min(),
        "Max": s.max(),
    }
    cont_rows.append(row)
    print(f"  {label}: n={n}, median={row['Median']:.1f} [{row['Q1']:.1f}–{row['Q3']:.1f}], "
          f"mean={row['Mean']:.1f} ± {row['SD']:.1f}")

cont_df = pd.DataFrame(cont_rows)
cont_df.to_csv(TAB / "table1_continuous.csv", index=False, float_format="%.2f")

# ============================================================
# 5. DESCRIPTIVE STATISTICS — CATEGORICAL VARIABLES
# ============================================================
print("\n" + "=" * 60)
print("5. DESCRIPTIVE STATISTICS — CATEGORICAL VARIABLES")
print("=" * 60)

categorical_vars = [
    ("emergency", "Emergency CS"),
    ("prior_cs", "Prior CS"),
    ("HDP", "Hypertensive disorders of pregnancy"),
    ("preop_steroid", "Preoperative steroid (≤1 week)"),
    ("epidural", "Epidural anesthesia"),
    ("vasopressor_continuous", "Continuous vasopressor infusion"),
    ("hypotension", "Hypotension (SBP < 90 mmHg)"),
    ("uterine_exteriorization", "Uterine exteriorization"),
    ("bmi_ge35", "BMI ≥ 35 kg/m²"),
    ("preop_antiemetic", "Preoperative antiemetic (≤24 h)"),
]

cat_rows = []
for var, label in categorical_vars:
    if var not in df.columns:
        continue
    s = df[var].dropna()
    n_valid = len(s)
    n_pos = int(s.sum())
    pct = 100 * n_pos / n_valid if n_valid > 0 else 0
    ci_lo, ci_hi = proportion_confint(n_pos, n_valid, method="wilson") if n_valid > 0 else (0, 0)
    row = {
        "Variable": label,
        "n/N": f"{n_pos}/{n_valid}",
        "Percentage": pct,
        "95% CI lower": 100 * ci_lo,
        "95% CI upper": 100 * ci_hi,
        "Missing": len(df) - n_valid,
    }
    cat_rows.append(row)
    print(f"  {label}: {n_pos}/{n_valid} ({pct:.1f}%) [95%CI {100*ci_lo:.1f}–{100*ci_hi:.1f}%]")

cat_df = pd.DataFrame(cat_rows)
cat_df.to_csv(TAB / "table2_categorical.csv", index=False, float_format="%.2f")

# ============================================================
# 6. IONV OUTCOMES
# ============================================================
print("\n" + "=" * 60)
print("6. IONV OUTCOMES (TWIN ONLY)")
print("=" * 60)

ionv_rows = []
for var, label in [("ionv_primary", "Primary: any IONV (anesthesia→exit)"),
                    ("ionv_secondary", "Secondary: IONV before delivery only"),
                    ("ae_post_delivery", "Post-delivery IONV only"),
                    ("serotonin_antagonist", "5-HT3 antagonist use (Definition E)")]:
    s = df[var].dropna()
    n_valid = len(s)
    n_pos = int(s.sum())
    pct = 100 * n_pos / n_valid if n_valid > 0 else 0
    ci_lo, ci_hi = proportion_confint(n_pos, n_valid, method="wilson") if n_valid > 0 else (0, 0)
    ionv_rows.append({
        "Outcome": label,
        "Events": n_pos,
        "N": n_valid,
        "Rate (%)": pct,
        "95% CI lower": 100 * ci_lo,
        "95% CI upper": 100 * ci_hi,
    })
    print(f"  {label}: {n_pos}/{n_valid} ({pct:.1f}%) [95%CI {100*ci_lo:.1f}–{100*ci_hi:.1f}%]")

ionv_df = pd.DataFrame(ionv_rows)
ionv_df.to_csv(TAB / "table3_ionv_outcomes.csv", index=False, float_format="%.2f")

# ============================================================
# 7. ANTIEMETIC DRUG DETAILS
# ============================================================
print("\n" + "=" * 60)
print("7. ANTIEMETIC DRUG USAGE")
print("=" * 60)

drug_vars = [
    ("metoclopramide_mg", "Metoclopramide"),
    ("droperidol_mg", "Droperidol"),
    ("ondansetron_mg", "Ondansetron"),
    ("granisetron_mg", "Granisetron"),
    ("novamin_mg", "Prochlorperazine (Novamin)"),
    ("atarax_p_mg", "Hydroxyzine (Atarax-P)"),
    ("dexamethasone_mg", "Dexamethasone"),
]

drug_rows = []
for var, label in drug_vars:
    if var not in df.columns:
        continue
    vals = pd.to_numeric(df[var], errors="coerce").fillna(0)
    n_any = int((vals > 0).sum())
    pct = 100 * n_any / len(df) if len(df) > 0 else 0
    doses = vals[vals > 0]
    row = {
        "Drug": label,
        "n (%)": f"{n_any} ({pct:.1f}%)",
        "n_users": n_any,
        "Percentage": pct,
    }
    if len(doses) > 0:
        row["Dose median"] = doses.median()
        row["Dose IQR"] = f"{doses.quantile(0.25):.1f}–{doses.quantile(0.75):.1f}"
    else:
        row["Dose median"] = np.nan
        row["Dose IQR"] = "–"
    drug_rows.append(row)
    dose_str = f", median dose {row['Dose median']:.1f} [{row['Dose IQR']}]" if n_any > 0 else ""
    print(f"  {label}: {n_any} ({pct:.1f}%){dose_str}")

drug_df = pd.DataFrame(drug_rows)
drug_df.to_csv(TAB / "table4_drug_details.csv", index=False, float_format="%.2f")

# ============================================================
# 8. SUBGROUP: EMERGENCY vs ELECTIVE
# ============================================================
print("\n" + "=" * 60)
print("8. SUBGROUP ANALYSIS — EMERGENCY vs ELECTIVE")
print("=" * 60)

subgroup_rows = []
for grp_name, mask in [("Emergency", df["emergency"] == 1),
                        ("Elective", df["emergency"] == 0)]:
    sub = df[mask]
    n = len(sub)
    if n == 0:
        continue
    row = {"Subgroup": grp_name, "n": n}
    for var, label in [("年齢(歳)", "Age"), ("BMI", "BMI"), ("GA_weeks", "GA (weeks)"),
                        ("出血量_ml", "Blood loss (mL)"), ("手術時間_min", "Surgery time (min)"),
                        ("麻酔時間_min", "Anesthesia time (min)")]:
        if var in sub.columns:
            s = sub[var].dropna()
            row[f"{label} median [IQR]"] = f"{s.median():.1f} [{s.quantile(0.25):.1f}–{s.quantile(0.75):.1f}]" if len(s) > 0 else "–"
    ionv_p = sub["ionv_primary"].sum()
    row["IONV primary n (%)"] = f"{ionv_p} ({100*ionv_p/n:.1f}%)" if n > 0 else "–"
    ionv_s = sub["ionv_secondary"].sum()
    row["IONV secondary n (%)"] = f"{ionv_s} ({100*ionv_s/n:.1f}%)" if n > 0 else "–"
    subgroup_rows.append(row)
    print(f"\n  {grp_name} (n={n}):")
    for k, v in row.items():
        if k not in ("Subgroup", "n"):
            print(f"    {k}: {v}")

subgroup_df = pd.DataFrame(subgroup_rows)
subgroup_df.to_csv(TAB / "table5_emergency_vs_elective.csv", index=False)

# Emergency vs Elective comparison (p-values)
emerg = df[df["emergency"] == 1]
elect = df[df["emergency"] == 0]
comparison_rows = []
for var, label in continuous_vars:
    if var not in df.columns:
        continue
    e_vals = pd.to_numeric(emerg[var], errors="coerce").dropna()
    l_vals = pd.to_numeric(elect[var], errors="coerce").dropna()
    if len(e_vals) > 0 and len(l_vals) > 0:
        stat, p = stats.mannwhitneyu(e_vals, l_vals, alternative="two-sided")
    else:
        p = np.nan
    comparison_rows.append({
        "Variable": label,
        "Emergency median [IQR]": f"{e_vals.median():.1f} [{e_vals.quantile(0.25):.1f}–{e_vals.quantile(0.75):.1f}]" if len(e_vals) > 0 else "–",
        "Elective median [IQR]": f"{l_vals.median():.1f} [{l_vals.quantile(0.25):.1f}–{l_vals.quantile(0.75):.1f}]" if len(l_vals) > 0 else "–",
        "P-value": p,
    })

comp_df = pd.DataFrame(comparison_rows)
comp_df.to_csv(TAB / "table6_emerg_vs_elect_continuous.csv", index=False, float_format="%.4f")

# ============================================================
# 9. TEMPORAL TRENDS
# ============================================================
print("\n" + "=" * 60)
print("9. TEMPORAL TRENDS")
print("=" * 60)

yearly = df.groupby("year").agg(
    n=("year", "size"),
    age_median=("年齢(歳)", "median"),
    ga_median=("GA_weeks", "median"),
    blood_loss_median=("出血量_ml", "median"),
    ionv_primary_n=("ionv_primary", "sum"),
    ionv_secondary_n=("ionv_secondary", "sum"),
    serotonin_n=("serotonin_antagonist", "sum"),
).reset_index()
yearly["ionv_primary_pct"] = 100 * yearly["ionv_primary_n"] / yearly["n"]
yearly["ionv_secondary_pct"] = 100 * yearly["ionv_secondary_n"] / yearly["n"]
yearly["serotonin_pct"] = 100 * yearly["serotonin_n"] / yearly["n"]
yearly.to_csv(TAB / "table7_yearly_trends.csv", index=False, float_format="%.2f")

for _, row in yearly.iterrows():
    print(f"  {int(row['year'])}: n={int(row['n'])}, IONV={int(row['ionv_primary_n'])} ({row['ionv_primary_pct']:.1f}%), "
          f"EBL median={row['blood_loss_median']:.0f} mL")

# ============================================================
# 10. FIGURES
# ============================================================
print("\n" + "=" * 60)
print("10. GENERATING FIGURES")
print("=" * 60)

# --- Fig 1: Distribution of continuous variables ---
fig, axes = plt.subplots(2, 3, figsize=(14, 9))
plot_vars = [("年齢(歳)", "Age (years)"), ("BMI", "BMI (kg/m²)"),
             ("GA_weeks", "Gestational age (weeks)"),
             ("出血量_ml", "Estimated blood loss (mL)"),
             ("手術時間_min", "Surgery time (min)"),
             ("麻酔時間_min", "Anesthesia time (min)")]

for ax, (var, label) in zip(axes.flatten(), plot_vars):
    vals = df[var].dropna()
    ax.hist(vals, bins=25, color="#5B9BD5", edgecolor="white", alpha=0.85)
    med = vals.median()
    ax.axvline(med, color="#C0504D", linestyle="--", linewidth=1.5, label=f"Median: {med:.1f}")
    ax.set_xlabel(label)
    ax.set_ylabel("Count")
    ax.legend(fontsize=8)

fig.suptitle(f"Distribution of Key Variables — Twin CS (n={len(df)})", fontsize=13, fontweight="bold")
plt.tight_layout(rect=[0, 0, 1, 0.96])
fig.savefig(FIG / "fig1_distributions.png")
plt.close(fig)
print("  Fig 1: distributions saved")

# --- Fig 2: IONV rates with 95% CI ---
fig, ax = plt.subplots(figsize=(8, 5))
outcome_labels = ["Primary\n(any IONV)", "Secondary\n(before delivery)", "Post-delivery\nonly", "5-HT3\nantagonist"]
rates = [float(r["Rate (%)"]) for r in ionv_rows]
ci_lo_list = [float(r["95% CI lower"]) for r in ionv_rows]
ci_hi_list = [float(r["95% CI upper"]) for r in ionv_rows]
yerr_lo = [r - lo for r, lo in zip(rates, ci_lo_list)]
yerr_hi = [hi - r for r, hi in zip(rates, ci_hi_list)]

colors = ["#4C72B0", "#55A868", "#C44E52", "#8172B2"]
bars = ax.bar(outcome_labels, rates, color=colors, width=0.5,
              yerr=[yerr_lo, yerr_hi], capsize=5, edgecolor="black", linewidth=0.5)
for bar, rate in zip(bars, rates):
    ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 1,
            f"{rate:.1f}%", ha="center", fontsize=10, fontweight="bold")
ax.set_ylabel("Rate (%)")
ax.set_title(f"IONV Outcome Rates — Twin CS (n={len(df)})")
ax.set_ylim(0, max(rates) * 1.4 + 5)
fig.tight_layout()
fig.savefig(FIG / "fig2_ionv_rates.png")
plt.close(fig)
print("  Fig 2: IONV rates saved")

# --- Fig 3: Temporal trends ---
if len(yearly) > 1:
    fig, ax1 = plt.subplots(figsize=(10, 5))
    ax1.bar(yearly["year"], yearly["n"], color="#B0C4DE", alpha=0.6, label="n cases")
    ax1.set_xlabel("Year")
    ax1.set_ylabel("Number of cases", color="#4169E1")
    ax1.tick_params(axis="y", labelcolor="#4169E1")
    ax1.xaxis.set_major_locator(mticker.MaxNLocator(integer=True))

    ax2 = ax1.twinx()
    ax2.plot(yearly["year"], yearly["ionv_primary_pct"], "o-", color="#C0504D", linewidth=2, label="IONV rate (%)")
    ax2.set_ylabel("IONV rate (%)", color="#C0504D")
    ax2.tick_params(axis="y", labelcolor="#C0504D")
    ax2.set_ylim(0, max(yearly["ionv_primary_pct"]) * 1.5 + 5)

    lines1, labels1 = ax1.get_legend_handles_labels()
    lines2, labels2 = ax2.get_legend_handles_labels()
    ax1.legend(lines1 + lines2, labels1 + labels2, loc="upper right")
    fig.suptitle(f"Temporal Trends — Twin CS (n={len(df)})", fontsize=13, fontweight="bold")
    fig.tight_layout(rect=[0, 0, 1, 0.96])
    fig.savefig(FIG / "fig3_temporal_trends.png")
    plt.close(fig)
    print("  Fig 3: temporal trends saved")

# --- Fig 4: Drug usage breakdown ---
fig, ax = plt.subplots(figsize=(9, 5))
drug_names = [r["Drug"] for r in drug_rows]
drug_pcts = [r["Percentage"] for r in drug_rows]
drug_colors = sns.color_palette("Set2", len(drug_names))
bars = ax.barh(drug_names[::-1], drug_pcts[::-1], color=drug_colors[::-1], edgecolor="black", linewidth=0.5)
for bar, pct in zip(bars, drug_pcts[::-1]):
    if pct > 0:
        ax.text(bar.get_width() + 0.3, bar.get_y() + bar.get_height() / 2,
                f"{pct:.1f}%", va="center", fontsize=9)
ax.set_xlabel("Usage rate (%)")
ax.set_title(f"Antiemetic Drug Usage — Twin CS (n={len(df)})")
fig.tight_layout()
fig.savefig(FIG / "fig4_drug_usage.png")
plt.close(fig)
print("  Fig 4: drug usage saved")

# --- Fig 5: Blood loss distribution (box + strip) ---
fig, ax = plt.subplots(figsize=(8, 5))
emerg_bl = df[df["emergency"] == 1]["出血量_ml"].dropna()
elect_bl = df[df["emergency"] == 0]["出血量_ml"].dropna()
data_for_box = [elect_bl, emerg_bl]
bp = ax.boxplot(data_for_box, labels=["Elective", "Emergency"], patch_artist=True,
                widths=0.4, showfliers=False)
bp["boxes"][0].set_facecolor("#5B9BD5")
bp["boxes"][1].set_facecolor("#ED7D31")
for i, d in enumerate(data_for_box, 1):
    jitter = np.random.normal(i, 0.04, size=len(d))
    ax.scatter(jitter, d, alpha=0.3, s=10, color="black", zorder=3)
ax.set_ylabel("Estimated blood loss (mL)")
ax.set_title(f"Blood Loss by Emergency Status — Twin CS (n={len(df)})")
fig.tight_layout()
fig.savefig(FIG / "fig5_blood_loss_by_emergency.png")
plt.close(fig)
print("  Fig 5: blood loss by emergency saved")

# ============================================================
# 11. SUMMARY JSON
# ============================================================
print("\n" + "=" * 60)
print("11. SAVING SUMMARY JSON")
print("=" * 60)

summary = {
    "n_raw": int(n_raw),
    "n_excluded": int(n_excluded),
    "n_pre_anesthesia_excluded": int(pre_ae_count),
    "n_analysis": int(len(df)),
    "exclusions": {reason: int(n) for reason, n in exclusions},
    "continuous_summary": {
        r["Variable"]: {
            "n": int(r["n"]),
            "mean": round(float(r["Mean"]), 2),
            "sd": round(float(r["SD"]), 2),
            "median": round(float(r["Median"]), 2),
            "q1": round(float(r["Q1"]), 2),
            "q3": round(float(r["Q3"]), 2),
            "min": round(float(r["Min"]), 2),
            "max": round(float(r["Max"]), 2),
        }
        for r in cont_rows
    },
    "categorical_summary": {
        r["Variable"]: {
            "n_N": r["n/N"],
            "pct": round(r["Percentage"], 2),
            "ci_lower": round(r["95% CI lower"], 2),
            "ci_upper": round(r["95% CI upper"], 2),
        }
        for r in cat_rows
    },
    "ionv_outcomes": {
        r["Outcome"]: {
            "events": int(r["Events"]),
            "N": int(r["N"]),
            "rate_pct": round(float(r["Rate (%)"]), 2),
        }
        for r in ionv_rows
    },
}

with open(BASE / "descriptive_twin_stats.json", "w") as f:
    json.dump(summary, f, indent=2, ensure_ascii=False)
print("  descriptive_twin_stats.json saved")

# ============================================================
# 12. GENERATE DOCX (descriptive_twin.docx)
# ============================================================
print("\n" + "=" * 60)
print("12. GENERATING DOCX")
print("=" * 60)

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT

doc = Document()
style = doc.styles["Normal"]
style.font.name = "Times New Roman"
style.font.size = Pt(11)
style.paragraph_format.line_spacing = 1.5

def add_heading(d, text, level=2):
    h = d.add_heading(text, level=level)
    for run in h.runs:
        run.font.color.rgb = RGBColor(0, 0, 0)
    return h

# Title
add_heading(doc, "Descriptive Statistics: Twin Cesarean Section Cohort", level=1)
p = doc.add_paragraph()
p.add_run(f"Analysis date: {pd.Timestamp.now().strftime('%Y-%m-%d')}\n").bold = False
p.add_run(f"Study period: 2014-04-01 to 2024-10-23\n")
p.add_run(f"Cohort: twin pregnancies undergoing cesarean section under spinal anesthesia\n")
p.add_run(f"N (raw): {n_raw} → N (analysis): {len(df)} (after exclusions and pre-anesthesia antiemetic removal)")

# Flow
add_heading(doc, "Patient Flow")
p = doc.add_paragraph()
p.add_run(f"A total of {n_raw} twin cesarean sections were identified. ")
for reason, n in exclusions:
    if n > 0:
        p.add_run(f"{reason} ({n}), ")
p.add_run(f"were excluded ({n_excluded} total). ")
p.add_run(f"An additional {pre_ae_count} patients who received antiemetics before anesthesia induction were excluded per protocol. ")
p.add_run(f"The final analysis cohort comprised {len(df)} twin cesarean sections.")

# Table 1: Continuous
add_heading(doc, "Table 1. Patient Characteristics — Continuous Variables")
tbl = doc.add_table(rows=1, cols=6, style="Table Grid")
tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
headers = ["Variable", "n", "Mean ± SD", "Median [IQR]", "Min", "Max"]
for i, h_text in enumerate(headers):
    cell = tbl.rows[0].cells[i]
    cell.text = h_text
    cell.paragraphs[0].runs[0].bold = True
    cell.paragraphs[0].runs[0].font.size = Pt(9)

for r in cont_rows:
    row_cells = tbl.add_row().cells
    vals = [
        r["Variable"],
        str(int(r["n"])),
        f"{r['Mean']:.1f} ± {r['SD']:.1f}",
        f"{r['Median']:.1f} [{r['Q1']:.1f}–{r['Q3']:.1f}]",
        f"{r['Min']:.1f}",
        f"{r['Max']:.1f}",
    ]
    for i, v in enumerate(vals):
        row_cells[i].text = v
        row_cells[i].paragraphs[0].runs[0].font.size = Pt(9)

# Table 2: Categorical
add_heading(doc, "Table 2. Patient Characteristics — Categorical Variables")
tbl2 = doc.add_table(rows=1, cols=4, style="Table Grid")
tbl2.alignment = WD_TABLE_ALIGNMENT.CENTER
for i, h_text in enumerate(["Variable", "n/N (%)", "95% CI", "Missing"]):
    cell = tbl2.rows[0].cells[i]
    cell.text = h_text
    cell.paragraphs[0].runs[0].bold = True
    cell.paragraphs[0].runs[0].font.size = Pt(9)

for r in cat_rows:
    row_cells = tbl2.add_row().cells
    vals = [
        r["Variable"],
        f"{r['n/N']} ({r['Percentage']:.1f}%)",
        f"{r['95% CI lower']:.1f}–{r['95% CI upper']:.1f}%",
        str(int(r["Missing"])),
    ]
    for i, v in enumerate(vals):
        row_cells[i].text = v
        row_cells[i].paragraphs[0].runs[0].font.size = Pt(9)

# Fig 1 inline
add_heading(doc, "Figure 1. Distribution of Key Variables")
fig1_path = FIG / "fig1_distributions.png"
if fig1_path.exists():
    doc.add_picture(str(fig1_path), width=Inches(6))
    p_cap = doc.add_paragraph()
    p_cap.add_run("Figure 1. ").bold = True
    p_cap.add_run(f"Distribution of age, BMI, gestational age, blood loss, surgery time, and "
                  f"anesthesia time in the twin CS cohort (n={len(df)}). "
                  "Red dashed lines indicate median values.")
    p_cap.runs[0].font.size = Pt(9)
    p_cap.paragraph_format.space_before = Pt(12)

# Table 3: IONV outcomes
add_heading(doc, "Table 3. IONV Outcomes")
tbl3 = doc.add_table(rows=1, cols=4, style="Table Grid")
tbl3.alignment = WD_TABLE_ALIGNMENT.CENTER
for i, h_text in enumerate(["Outcome", "Events/N", "Rate (%)", "95% CI"]):
    cell = tbl3.rows[0].cells[i]
    cell.text = h_text
    cell.paragraphs[0].runs[0].bold = True
    cell.paragraphs[0].runs[0].font.size = Pt(9)

for r in ionv_rows:
    row_cells = tbl3.add_row().cells
    vals = [
        r["Outcome"],
        f"{int(r['Events'])}/{int(r['N'])}",
        f"{float(r['Rate (%)']):.1f}",
        f"{float(r['95% CI lower']):.1f}–{float(r['95% CI upper']):.1f}%",
    ]
    for i, v in enumerate(vals):
        row_cells[i].text = v
        row_cells[i].paragraphs[0].runs[0].font.size = Pt(9)

# Fig 2 inline
add_heading(doc, "Figure 2. IONV Outcome Rates")
fig2_path = FIG / "fig2_ionv_rates.png"
if fig2_path.exists():
    doc.add_picture(str(fig2_path), width=Inches(5))
    p_cap = doc.add_paragraph()
    p_cap.add_run("Figure 2. ").bold = True
    p_cap.add_run(f"IONV outcome rates with 95% Wilson confidence intervals in twin CS (n={len(df)}).")
    p_cap.runs[0].font.size = Pt(9)
    p_cap.paragraph_format.space_before = Pt(12)

# Table 4: Drug details
add_heading(doc, "Table 4. Antiemetic Drug Usage")
tbl4 = doc.add_table(rows=1, cols=4, style="Table Grid")
tbl4.alignment = WD_TABLE_ALIGNMENT.CENTER
for i, h_text in enumerate(["Drug", "n (%)", "Median dose", "Dose IQR"]):
    cell = tbl4.rows[0].cells[i]
    cell.text = h_text
    cell.paragraphs[0].runs[0].bold = True
    cell.paragraphs[0].runs[0].font.size = Pt(9)

for r in drug_rows:
    row_cells = tbl4.add_row().cells
    dose_med = f"{r['Dose median']:.1f}" if not pd.isna(r["Dose median"]) else "–"
    vals = [r["Drug"], r["n (%)"], dose_med, r["Dose IQR"]]
    for i, v in enumerate(vals):
        row_cells[i].text = str(v)
        row_cells[i].paragraphs[0].runs[0].font.size = Pt(9)

# Fig 4 inline
add_heading(doc, "Figure 3. Antiemetic Drug Usage")
fig4_path = FIG / "fig4_drug_usage.png"
if fig4_path.exists():
    doc.add_picture(str(fig4_path), width=Inches(5.5))
    p_cap = doc.add_paragraph()
    p_cap.add_run("Figure 3. ").bold = True
    p_cap.add_run(f"Frequency of individual antiemetic drug usage in twin CS (n={len(df)}).")
    p_cap.runs[0].font.size = Pt(9)
    p_cap.paragraph_format.space_before = Pt(12)

# Table 5: Emergency vs Elective
add_heading(doc, "Table 5. Emergency vs Elective Subgroup")
if len(subgroup_rows) > 0:
    cols_sub = list(subgroup_rows[0].keys())
    tbl5 = doc.add_table(rows=1, cols=len(cols_sub), style="Table Grid")
    tbl5.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, h_text in enumerate(cols_sub):
        cell = tbl5.rows[0].cells[i]
        cell.text = h_text
        cell.paragraphs[0].runs[0].bold = True
        cell.paragraphs[0].runs[0].font.size = Pt(8)
    for r in subgroup_rows:
        row_cells = tbl5.add_row().cells
        for i, k in enumerate(cols_sub):
            row_cells[i].text = str(r[k])
            row_cells[i].paragraphs[0].runs[0].font.size = Pt(8)

# Fig 5 inline
add_heading(doc, "Figure 4. Blood Loss by Emergency Status")
fig5_path = FIG / "fig5_blood_loss_by_emergency.png"
if fig5_path.exists():
    doc.add_picture(str(fig5_path), width=Inches(5))
    p_cap = doc.add_paragraph()
    p_cap.add_run("Figure 4. ").bold = True
    p_cap.add_run(f"Estimated blood loss distribution in emergency vs elective twin CS (n={len(df)}). "
                  "Individual data points shown with jitter.")
    p_cap.runs[0].font.size = Pt(9)
    p_cap.paragraph_format.space_before = Pt(12)

# Fig 3 (temporal) inline
fig3_path = FIG / "fig3_temporal_trends.png"
if fig3_path.exists():
    add_heading(doc, "Figure 5. Temporal Trends")
    doc.add_picture(str(fig3_path), width=Inches(5.5))
    p_cap = doc.add_paragraph()
    p_cap.add_run("Figure 5. ").bold = True
    p_cap.add_run(f"Annual case volume and IONV rates in twin CS (n={len(df)}).")
    p_cap.runs[0].font.size = Pt(9)
    p_cap.paragraph_format.space_before = Pt(12)

# Table 7: Yearly trends
add_heading(doc, "Table 6. Yearly Trends")
if len(yearly) > 0:
    tbl7 = doc.add_table(rows=1, cols=5, style="Table Grid")
    tbl7.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, h_text in enumerate(["Year", "n", "IONV primary n (%)", "EBL median (mL)", "5-HT3 n (%)"]):
        cell = tbl7.rows[0].cells[i]
        cell.text = h_text
        cell.paragraphs[0].runs[0].bold = True
        cell.paragraphs[0].runs[0].font.size = Pt(9)
    for _, yr in yearly.iterrows():
        row_cells = tbl7.add_row().cells
        vals = [
            str(int(yr["year"])),
            str(int(yr["n"])),
            f"{int(yr['ionv_primary_n'])} ({yr['ionv_primary_pct']:.1f}%)",
            f"{yr['blood_loss_median']:.0f}",
            f"{int(yr['serotonin_n'])} ({yr['serotonin_pct']:.1f}%)",
        ]
        for i, v in enumerate(vals):
            row_cells[i].text = v
            row_cells[i].paragraphs[0].runs[0].font.size = Pt(9)

doc.save(BASE / "descriptive_twin.docx")
print("  descriptive_twin.docx saved")

# ============================================================
# 13. GENERATE PPTX (editable figures)
# ============================================================
print("\n" + "=" * 60)
print("13. GENERATING PPTX")
print("=" * 60)

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = PptxInches(13.333)
prs.slide_height = PptxInches(7.5)

fig_items = [
    ("fig1_distributions.png", "Figure 1",
     f"Distribution of key variables — Twin CS (n={len(df)})"),
    ("fig2_ionv_rates.png", "Figure 2",
     f"IONV outcome rates with 95% CI — Twin CS (n={len(df)})"),
    ("fig3_temporal_trends.png", "Figure 3",
     f"Temporal trends — Twin CS (n={len(df)})"),
    ("fig4_drug_usage.png", "Figure 4",
     f"Antiemetic drug usage — Twin CS (n={len(df)})"),
    ("fig5_blood_loss_by_emergency.png", "Figure 5",
     f"Blood loss by emergency status — Twin CS (n={len(df)})"),
]

for fname, title, caption in fig_items:
    fpath = FIG / fname
    if not fpath.exists():
        continue
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title
    txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.2), PptxInches(12), PptxInches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = PptxPt(24)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Image
    from PIL import Image
    img = Image.open(fpath)
    img_w, img_h = img.size
    max_w, max_h = 11.5, 5.2
    scale = min(max_w / (img_w / 96), max_h / (img_h / 96))
    w_in = (img_w / 96) * scale
    h_in = (img_h / 96) * scale
    left = (13.333 - w_in) / 2
    slide.shapes.add_picture(str(fpath), PptxInches(left), PptxInches(1.0),
                             PptxInches(w_in), PptxInches(h_in))

    # Caption
    txBox2 = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(6.5), PptxInches(12), PptxInches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = caption
    p2.font.size = PptxPt(12)
    p2.alignment = PP_ALIGN.CENTER

prs.save(BASE / "figures_twin.pptx")
print("  figures_twin.pptx saved")

# ============================================================
# DONE
# ============================================================
print("\n" + "=" * 60)
print("COMPLETE")
print("=" * 60)
print(f"Analysis cohort: n={len(df)} twin CS")
print(f"Outputs:")
print(f"  - tables_twin/ ({len(list(TAB.glob('*.csv')))} CSV files)")
print(f"  - figures_twin/ ({len(list(FIG.glob('*.png')))} PNG files)")
print(f"  - descriptive_twin.docx")
print(f"  - figures_twin.pptx")
print(f"  - descriptive_twin_stats.json")
