"""Create editable English PPTX with all figures from Definition E analysis."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pathlib import Path

BASE = Path(__file__).resolve().parent
FIG = BASE / "figures_e"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

figures = [
    ("fig1_rates_comparison.png", "Fig. 1: IONV Rates — Protocol (A) vs 5-HT3 (E)"),
    ("fig2_forest_E_primary.png", "Fig. 2: Logistic Regression — 5-HT3 Use (Def E Primary)"),
    ("fig3_covariate_sensitivity.png", "Fig. 3: Covariate Sensitivity Analysis (Def E Primary)"),
    ("fig4_protocol_vs_defE.png", "Fig. 4: Twin Effect — Protocol vs 5-HT3 Definition"),
]

for fname, title in figures:
    fpath = FIG / fname
    if not fpath.exists():
        continue
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    slide.shapes.add_picture(str(fpath), Inches(0.5), Inches(1.0), Inches(12), Inches(6))

out = BASE / "figures_ionv_e.pptx"
prs.save(str(out))
print(f"PPTX saved: {out}")
