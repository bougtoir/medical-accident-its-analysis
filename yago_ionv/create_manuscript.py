"""
Generate Methods (Statistical Analysis) and Results sections
for the IONV during cesarean section study.

Outputs:
  - manuscript_methods_results.docx  (inline figures + tables)
  - figures.pptx                     (editable, one figure per slide)
"""

import json, re
from pathlib import Path
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
import pandas as pd

BASE = Path(__file__).resolve().parent
FIG = BASE / "figures"

with open(BASE / "summary_stats.json") as f:
    S = json.load(f)

# ============================================================
# Helper: superscript citations using font.superscript
# ============================================================
def add_text_with_refs(para, text):
    parts = re.split(r"(\{[^}]+\})", text)
    for part in parts:
        if part.startswith("{") and part.endswith("}"):
            run = para.add_run(part[1:-1])
            run.font.superscript = True
            run.font.size = Pt(8)
        else:
            run = para.add_run(part)
            run.font.size = Pt(10)


def add_heading(doc, text, level=2):
    h = doc.add_heading(text, level=level)
    for run in h.runs:
        run.font.color.rgb = RGBColor(0, 0, 0)
    return h


# ============================================================
# Build document
# ============================================================
doc = Document()

style = doc.styles["Normal"]
style.font.name = "Times New Roman"
style.font.size = Pt(10)
style.paragraph_format.line_spacing = 1.5
style.paragraph_format.space_after = Pt(6)

# ========================= METHODS ========================= #
add_heading(doc, "Methods", level=1)
add_heading(doc, "Study Design and Setting")

p = doc.add_paragraph()
add_text_with_refs(p,
    "This retrospective, single-center, observational study was conducted at "
    "Juntendo University Shizuoka Hospital, Japan, using medical records from "
    "April 1, 2014, to October 23, 2024. The study was approved by the institutional "
    "review board and conducted in accordance with the Declaration of Helsinki (2013 revision) "
    "and the Japanese Ethical Guidelines for Medical and Biological Research Involving Human "
    "Subjects (2021). Informed consent was waived via an opt-out mechanism."
)

add_heading(doc, "Participants")
p = doc.add_paragraph()
add_text_with_refs(p,
    "We included women aged ≥18 years who underwent cesarean section under regional "
    "anesthesia (spinal, epidural, or combined spinal-epidural). "
    "Exclusion criteria were: general anesthesia, cardiac arrest in the operating room, "
    "pre-anesthesia systolic blood pressure <90 mmHg, intrauterine fetal death, "
    "vanishing twin, and triplet pregnancy."
)

add_heading(doc, "Outcome Measures")
p = doc.add_paragraph()
add_text_with_refs(p,
    "Because this was a retrospective study and nausea and vomiting could not be directly "
    "ascertained from medical records, intraoperative antiemetic use was employed as a "
    "surrogate marker for IONV.{2,3} "
    "The primary outcome was antiemetic use during the period from anesthesia induction to "
    "delivery or from delivery to operating room exit (i.e., any intraoperative antiemetic use). "
    "The secondary outcome was antiemetic use from anesthesia induction to delivery only. "
    "Antiemetics included ondansetron, granisetron, metoclopramide, droperidol, "
    "prochlorperazine (Novamin), hydroxyzine (Atarax-P), and dexamethasone."
)

p = doc.add_paragraph()
add_text_with_refs(p,
    "Patients who received antiemetics before anesthesia induction (during the admission-to-anesthesia "
    "period) were excluded from the outcome analysis, as pre-anesthesia antiemetic administration "
    "was considered likely prophylactic.{12-14}"
)

add_heading(doc, "Statistical Analysis")
p = doc.add_paragraph()
add_text_with_refs(p,
    "Continuous variables were summarized as median [interquartile range] and compared using "
    "the Mann-Whitney U test. Categorical variables were summarized as n (%) and compared using "
    "the chi-square test or Fisher's exact test where expected cell counts were <5."
)

p = doc.add_paragraph()
add_text_with_refs(p,
    "Multivariable logistic regression was performed to evaluate the association between "
    "twin pregnancy and IONV, adjusting for known risk factors: "
    "age, body mass index (BMI), gestational age, emergency cesarean section, prior cesarean section, "
    "hypertensive disorders of pregnancy (HDP), epidural anesthesia, surgery time, and "
    "intraoperative hypotension (systolic blood pressure <90 mmHg).{7,8} "
    "Uterine exteriorization was not included in the model because of excessive missing data "
    "(available in only approximately 20% of the cohort). "
    "Results are presented as odds ratios (OR) with 95% confidence intervals (CI). "
    "A two-sided P < 0.05 was considered statistically significant. "
    "All analyses were performed using Python 3.12 with statsmodels 0.14, scipy 1.14, "
    "and scikit-learn 1.6."
)

add_heading(doc, "Sensitivity Analyses")
p = doc.add_paragraph()
add_text_with_refs(p,
    "Six sensitivity analyses were performed using alternative IONV definitions to assess "
    "the robustness of the primary finding. Definition C restricted IONV to the post-delivery "
    "phase only (antiemetic use after delivery but not before). Definition D identified severe "
    "IONV, defined as use of two or more different antiemetic drugs. Definition E used a "
    "drug-specific approach, counting only 5-HT3 receptor antagonists (ondansetron or "
    "granisetron) as IONV markers. Definition F excluded patients whose only antiemetic was "
    "dexamethasone (frequently used prophylactically). Definition G treated the number of "
    "antiemetic drugs used as a count outcome and applied Poisson regression to estimate "
    "incidence rate ratios (IRR). Definition H stratified the primary analysis by the presence "
    "or absence of intraoperative hypotension, to assess whether the twin–IONV association "
    "differed by hypotension status."
)

p = doc.add_paragraph()
add_text_with_refs(p,
    "For definitions with sparse events (events-per-variable ratio <5), a reduced model with "
    "six covariates (twin, age, BMI, gestational age, emergency, and hypotension) was used "
    "to ensure model convergence."
)

# ========================= RESULTS ========================= #
add_heading(doc, "Results", level=1)
add_heading(doc, "Study Population")

n_total = S["n_total_before_exclusion"]
n_excl = S["n_excluded"]
n_analysis = S["n_analysis"]
n_s = S["n_single"]
n_t = S["n_twin"]
n_pre_ae = S["n_pre_anesthesia_antiemetic_excluded"]

p = doc.add_paragraph()
add_text_with_refs(p,
    f"A total of {n_total:,} women who underwent cesarean section during the study period "
    f"were identified ({S['n_single_raw']:,} singleton and {S['n_twin_raw']:,} twin pregnancies). "
    f"After applying exclusion criteria, {S['n_after_exclusion']:,} patients remained "
    f"({n_total:,} total; {S['excl_general_anesthesia']} "
    f"excluded for general anesthesia, {S['excl_intrauterine_fetal_death']} for intrauterine fetal death, "
    f"{S['excl_vanishing_twin']} for vanishing twin, {S['excl_triplet_pregnancy']} for triplet pregnancy, "
    f"{S['excl_pre-anesthesia_sbp_lt_90_mmhg']} for pre-anesthesia hypotension, and "
    f"{S['excl_no_anesthesia_data_available']} for unavailable anesthesia data). "
    f"An additional {n_pre_ae} patients who received pre-anesthesia antiemetics were excluded "
    f"from the outcome analysis, yielding a final analysis cohort of {n_analysis:,} patients "
    f"({n_s:,} singleton, {n_t:,} twin) (Fig. 1)."
)

# Table 1
add_heading(doc, "Patient Characteristics (Table 1)")
p = doc.add_paragraph()
add_text_with_refs(p,
    "Baseline patient characteristics are summarized in Table 1. "
    "Compared with the singleton group, the twin group had significantly lower age "
    "(32.0 vs 34.0 years, P < 0.001), lower gestational age (37.0 vs 38.0 weeks, P < 0.001), "
    "fewer emergency cesarean sections (35.4% vs 46.8%, P < 0.001), "
    "fewer prior cesarean sections (11.7% vs 45.0%, P < 0.001), "
    "lower hypotension incidence (44.2% vs 52.8%, P = 0.003), "
    "and higher blood loss (1195 vs 650 mL, P < 0.001). "
    "BMI, epidural use, HDP, and preoperative steroid use did not differ significantly "
    "between groups."
)

# Insert Table 1
tab1 = pd.read_csv(BASE / "tables" / "table1_characteristics.csv")
table = doc.add_table(rows=1, cols=4)
table.style = "Table Grid"
table.alignment = WD_TABLE_ALIGNMENT.CENTER
hdr = table.rows[0].cells
for i, h_text in enumerate(["Variable", "Singleton", "Twin", "P-value"]):
    hdr[i].text = h_text
    for para in hdr[i].paragraphs:
        for run in para.runs:
            run.bold = True
            run.font.size = Pt(9)

for _, row in tab1.iterrows():
    cells = table.add_row().cells
    cells[0].text = row["Variable"]
    cells[1].text = row["Singleton"]
    cells[2].text = row["Twin"]
    p_val = row["P-value"]
    cells[3].text = f"{p_val:.3f}" if p_val >= 0.001 else "< 0.001"
    for c in cells:
        for para in c.paragraphs:
            for run in para.runs:
                run.font.size = Pt(9)

doc.add_paragraph()  # spacer

# IONV outcomes
add_heading(doc, "IONV Outcomes")

p_s = S["ionv_primary_single_pct"]
p_t = S["ionv_primary_twin_pct"]
s_s = S["ionv_secondary_single_pct"]
s_t = S["ionv_secondary_twin_pct"]

p = doc.add_paragraph()
add_text_with_refs(p,
    f"The primary outcome (any intraoperative antiemetic use) occurred in "
    f"{S['ionv_primary_single_n']} of {n_s:,} singleton patients ({p_s:.1f}%) and "
    f"{S['ionv_primary_twin_n']} of {n_t:,} twin patients ({p_t:.1f}%), "
    f"with no significant difference between groups (P = 0.731) (Table 2, Fig. 1). "
    f"The secondary outcome (antiemetic use before delivery) was similarly not significantly "
    f"different: {S['ionv_secondary_single_n']} ({s_s:.1f}%) in singletons vs "
    f"{S['ionv_secondary_twin_n']} ({s_t:.1f}%) in twins (P = 0.572). "
    f"Post-delivery antiemetic use rates were also comparable (16.0% vs 15.8%, P = 0.938)."
)

# Insert Fig 1
doc.add_paragraph()
p_fig = doc.add_paragraph()
p_fig.alignment = WD_ALIGN_PARAGRAPH.CENTER
if (FIG / "fig1_ionv_rates.png").exists():
    p_fig.add_run().add_picture(str(FIG / "fig1_ionv_rates.png"), width=Inches(5.5))
p_cap = doc.add_paragraph("Figure 1. IONV rates by group.")
p_cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
p_cap.runs[0].italic = True
p_cap.runs[0].font.size = Pt(9)
p_cap.paragraph_format.space_before = Pt(12)

# Table 2 (IONV outcomes)
tab2 = pd.read_csv(BASE / "tables" / "table2_ionv_outcomes.csv")
table2 = doc.add_table(rows=1, cols=4)
table2.style = "Table Grid"
table2.alignment = WD_TABLE_ALIGNMENT.CENTER
hdr2 = table2.rows[0].cells
for i, h_text in enumerate(["Outcome", "Singleton", "Twin", "P-value"]):
    hdr2[i].text = h_text
    for para in hdr2[i].paragraphs:
        for run in para.runs:
            run.bold = True
            run.font.size = Pt(9)

for _, row in tab2.iterrows():
    cells = table2.add_row().cells
    cells[0].text = row["Outcome"]
    cells[1].text = row["Singleton"]
    cells[2].text = row["Twin"]
    p_val = row["P-value"]
    cells[3].text = f"{p_val:.3f}" if p_val >= 0.001 else "< 0.001"
    for c in cells:
        for para in c.paragraphs:
            for run in para.runs:
                run.font.size = Pt(9)

# Antiemetic drug usage
add_heading(doc, "Antiemetic Drug Usage (Table 3)")

p = doc.add_paragraph()
add_text_with_refs(p,
    "The most commonly used antiemetic was metoclopramide (12.3% in singletons, 11.1% in twins), "
    "followed by droperidol (4.0% and 4.1%, respectively) (Table 3). "
    "Ondansetron use was significantly higher in twins (2.0% vs 0.9%, P = 0.041). "
    "Other antiemetic agents did not differ significantly between groups."
)

tab3 = pd.read_csv(BASE / "tables" / "table3_antiemetic_drugs.csv")
table3 = doc.add_table(rows=1, cols=4)
table3.style = "Table Grid"
table3.alignment = WD_TABLE_ALIGNMENT.CENTER
hdr3 = table3.rows[0].cells
for i, h_text in enumerate(["Drug", "Singleton", "Twin", "P-value"]):
    hdr3[i].text = h_text
    for para in hdr3[i].paragraphs:
        for run in para.runs:
            run.bold = True
            run.font.size = Pt(9)

for _, row in tab3.iterrows():
    cells = table3.add_row().cells
    cells[0].text = row["Drug"]
    cells[1].text = row["Singleton"]
    cells[2].text = row["Twin"]
    p_val = row["P-value"]
    cells[3].text = f"{p_val:.3f}" if p_val >= 0.001 else "< 0.001"
    for c in cells:
        for para in c.paragraphs:
            for run in para.runs:
                run.font.size = Pt(9)

# Multivariable regression
add_heading(doc, "Multivariable Logistic Regression")

if "primary_or_twin" in S:
    or_val = S["primary_or_twin"]
    ci_lo = S["primary_or_twin_ci_lower"]
    ci_hi = S["primary_or_twin_ci_upper"]
    p_val = S["primary_or_twin_p"]
    model_n = int(S["primary_model_n"])
    model_events = int(S["primary_model_events"])

    p_str = f"P = {p_val:.3f}" if p_val >= 0.001 else "P < 0.001"
    p = doc.add_paragraph()
    add_text_with_refs(p,
        f"In multivariable logistic regression (n = {model_n:,}; {model_events} events), "
        f"twin pregnancy was not independently associated with IONV "
        f"(adjusted OR {or_val:.2f}, 95% CI {ci_lo:.2f}–{ci_hi:.2f}; {p_str}) (Fig. 2, Table 4). "
        f"Factors significantly associated with IONV were epidural anesthesia "
        f"(OR 1.66, 95% CI 1.36–2.01; P < 0.001), longer surgery time "
        f"(OR 1.01 per minute, 95% CI 1.01–1.02; P < 0.001), and prior cesarean section "
        f"(OR 0.77, 95% CI 0.62–0.96; P = 0.021; protective)."
    )

# Insert Fig 2 - Forest plot (primary)
if (FIG / "fig2_forest_primary.png").exists():
    doc.add_paragraph()
    p_fig2 = doc.add_paragraph()
    p_fig2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_fig2.add_run().add_picture(str(FIG / "fig2_forest_primary.png"), width=Inches(5.5))
    p_cap2 = doc.add_paragraph("Figure 2. Forest plot: multivariable logistic regression for IONV (primary outcome).")
    p_cap2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_cap2.runs[0].italic = True
    p_cap2.runs[0].font.size = Pt(9)
    p_cap2.paragraph_format.space_before = Pt(12)

# Insert Table 4
tab4 = pd.read_csv(BASE / "tables" / "table4_logistic_primary.csv")
table4 = doc.add_table(rows=1, cols=5)
table4.style = "Table Grid"
table4.alignment = WD_TABLE_ALIGNMENT.CENTER
hdr4 = table4.rows[0].cells
for i, h_text in enumerate(["Variable", "OR", "95% CI lower", "95% CI upper", "P-value"]):
    hdr4[i].text = h_text
    for para in hdr4[i].paragraphs:
        for run in para.runs:
            run.bold = True
            run.font.size = Pt(9)

for _, row in tab4.iterrows():
    cells = table4.add_row().cells
    cells[0].text = str(row["Variable"])
    cells[1].text = f"{row['OR']:.2f}"
    cells[2].text = f"{row['95% CI lower']:.2f}"
    cells[3].text = f"{row['95% CI upper']:.2f}"
    p_val = row["P-value"]
    cells[4].text = f"{p_val:.3f}" if p_val >= 0.001 else "< 0.001"
    for c in cells:
        for para in c.paragraphs:
            for run in para.runs:
                run.font.size = Pt(9)

# Secondary outcome
if "secondary_or_twin" in S:
    or_s = S["secondary_or_twin"]
    ci_s_lo = S["secondary_or_twin_ci_lower"]
    ci_s_hi = S["secondary_or_twin_ci_upper"]
    p_s_val = S["secondary_or_twin_p"]
    model_s_n = int(S["secondary_model_n"])
    model_s_events = int(S["secondary_model_events"])

    p_s_str = f"P = {p_s_val:.3f}" if p_s_val >= 0.001 else "P < 0.001"
    p = doc.add_paragraph()
    add_text_with_refs(p,
        f"For the secondary outcome (IONV before delivery; n = {model_s_n:,}, "
        f"{model_s_events} events), twin pregnancy was likewise not significantly associated "
        f"(adjusted OR {or_s:.2f}, 95% CI {ci_s_lo:.2f}–{ci_s_hi:.2f}; {p_s_str}) (Fig. 3, Table 5)."
    )

# Insert Fig 3 - Forest plot (secondary)
if (FIG / "fig3_forest_secondary.png").exists():
    doc.add_paragraph()
    p_fig3 = doc.add_paragraph()
    p_fig3.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_fig3.add_run().add_picture(str(FIG / "fig3_forest_secondary.png"), width=Inches(5.5))
    p_cap3 = doc.add_paragraph("Figure 3. Forest plot: multivariable logistic regression for IONV (secondary outcome).")
    p_cap3.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_cap3.runs[0].italic = True
    p_cap3.runs[0].font.size = Pt(9)
    p_cap3.paragraph_format.space_before = Pt(12)

# IONV timing
add_heading(doc, "IONV Timing Distribution")
p = doc.add_paragraph()
add_text_with_refs(p,
    "The majority of IONV events occurred in the post-delivery phase (delivery to operating "
    "room exit), with comparable rates between singletons (16.0%) and twins (15.8%) (Fig. 4). "
    "IONV during the anesthesia-to-delivery phase was infrequent in both groups "
    "(1.9% vs 1.5%)."
)

if (FIG / "fig4_ionv_timing.png").exists():
    doc.add_paragraph()
    p_fig4 = doc.add_paragraph()
    p_fig4.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_fig4.add_run().add_picture(str(FIG / "fig4_ionv_timing.png"), width=Inches(5.0))
    p_cap4 = doc.add_paragraph("Figure 4. IONV rates by timing phase.")
    p_cap4.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_cap4.runs[0].italic = True
    p_cap4.runs[0].font.size = Pt(9)
    p_cap4.paragraph_format.space_before = Pt(12)

# Temporal trend
add_heading(doc, "Temporal Trends")
p = doc.add_paragraph()
add_text_with_refs(p,
    "Temporal trends of IONV rates are shown in Figure 5. "
    "IONV rates remained relatively stable over the study period in both groups, "
    "without a clear upward or downward trend."
)

if (FIG / "fig5_temporal_trend.png").exists():
    doc.add_paragraph()
    p_fig5 = doc.add_paragraph()
    p_fig5.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_fig5.add_run().add_picture(str(FIG / "fig5_temporal_trend.png"), width=Inches(5.0))
    p_cap5 = doc.add_paragraph("Figure 5. Temporal trend of IONV rates by year.")
    p_cap5.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_cap5.runs[0].italic = True
    p_cap5.runs[0].font.size = Pt(9)
    p_cap5.paragraph_format.space_before = Pt(12)

# Antiemetic drug comparison figure
if (FIG / "fig6_antiemetic_drugs.png").exists():
    doc.add_paragraph()
    p_fig6 = doc.add_paragraph()
    p_fig6.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_fig6.add_run().add_picture(str(FIG / "fig6_antiemetic_drugs.png"), width=Inches(5.0))
    p_cap6 = doc.add_paragraph("Figure 6. Antiemetic drug usage comparison: Singleton vs Twin.")
    p_cap6.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_cap6.runs[0].italic = True
    p_cap6.runs[0].font.size = Pt(9)
    p_cap6.paragraph_format.space_before = Pt(12)

# ========================= SENSITIVITY ANALYSIS RESULTS ========================= #
add_heading(doc, "Sensitivity Analyses")

# Load sensitivity stats
sens_json_path = BASE / "sensitivity_stats.json"
if sens_json_path.exists():
    with open(sens_json_path) as f:
        SS = json.load(f)

    # Sensitivity rates table
    p = doc.add_paragraph()
    add_text_with_refs(p,
        "IONV rates according to alternative definitions are shown in Figure 7 and Table 6. "
        "With the primary definition (A), IONV rates were comparable between singletons "
        f"({SS['definitions']['A']['singleton_pct']:.1f}%) and twins "
        f"({SS['definitions']['A']['twin_pct']:.1f}%). "
        "Most alternative definitions yielded similar non-significant differences. "
        "However, when IONV was defined by 5-HT3 antagonist use only (Definition E), "
        f"twins had a significantly higher rate ({SS['definitions']['E']['twin_pct']:.1f}% "
        f"vs {SS['definitions']['E']['singleton_pct']:.1f}%, P = 0.020)."
    )

    # Insert Fig 7 - Sensitivity forest plot
    if (FIG / "fig7_sensitivity_forest.png").exists():
        doc.add_paragraph()
        p_fig7 = doc.add_paragraph()
        p_fig7.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p_fig7.add_run().add_picture(str(FIG / "fig7_sensitivity_forest.png"), width=Inches(5.5))
        p_cap7 = doc.add_paragraph(
            "Figure 7. Forest plot: Effect of twin pregnancy on IONV across alternative definitions. "
            "Definitions A–F: adjusted odds ratios from multivariable logistic regression. "
            "Definition G: incidence rate ratio from Poisson regression. "
            "Definition H: stratified by hypotension status.")
        p_cap7.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p_cap7.runs[0].italic = True
        p_cap7.runs[0].font.size = Pt(9)
        p_cap7.paragraph_format.space_before = Pt(12)

    # Sensitivity summary table (Table 6)
    sens_table_path = BASE / "tables" / "sensitivity_summary.csv"
    if sens_table_path.exists():
        sens_tab = pd.read_csv(sens_table_path)
        table6 = doc.add_table(rows=1, cols=6)
        table6.style = "Table Grid"
        table6.alignment = WD_TABLE_ALIGNMENT.CENTER
        hdr6 = table6.rows[0].cells
        for i, h_text in enumerate(["Definition", "Label", "n", "Events", "Twin aOR/IRR [95% CI]", "P-value"]):
            hdr6[i].text = h_text
            for para in hdr6[i].paragraphs:
                for run in para.runs:
                    run.bold = True
                    run.font.size = Pt(9)

        for _, row in sens_tab.iterrows():
            cells = table6.add_row().cells
            cells[0].text = str(row["Definition"])
            cells[1].text = str(row["Label"])
            cells[2].text = str(int(row["n"])) if not pd.isna(row["n"]) else "—"
            cells[3].text = str(int(row["Events"])) if not pd.isna(row["Events"]) else "—"
            if pd.notna(row["twin_OR"]):
                cells[4].text = f"{row['twin_OR']:.2f} [{row['twin_CI_lower']:.2f}–{row['twin_CI_upper']:.2f}]"
                p_val = row["twin_P"]
                cells[5].text = f"{p_val:.3f}" if p_val >= 0.001 else "< 0.001"
            else:
                cells[4].text = "—"
                cells[5].text = "—"
            for c in cells:
                for para in c.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(9)

        p_tab6 = doc.add_paragraph("Table 6. Sensitivity analyses: twin pregnancy and IONV across alternative definitions.")
        p_tab6.runs[0].italic = True
        p_tab6.runs[0].font.size = Pt(9)

    # Insert Fig 8 - Sensitivity rates bar chart
    if (FIG / "fig8_sensitivity_rates.png").exists():
        doc.add_paragraph()
        p_fig8 = doc.add_paragraph()
        p_fig8.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p_fig8.add_run().add_picture(str(FIG / "fig8_sensitivity_rates.png"), width=Inches(5.5))
        p_cap8 = doc.add_paragraph("Figure 8. IONV rates by definition: Singleton vs Twin.")
        p_cap8.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p_cap8.runs[0].italic = True
        p_cap8.runs[0].font.size = Pt(9)
        p_cap8.paragraph_format.space_before = Pt(12)

    # Regression results text
    p = doc.add_paragraph()
    add_text_with_refs(p,
        "In multivariable logistic regression, twin pregnancy was not significantly associated "
        "with IONV under most alternative definitions (Table 6, Fig. 7). "
    )

    # 5-HT3 finding
    e_result = [r for r in SS["regression"] if r["Definition"] == "E"]
    if e_result and e_result[0].get("twin_OR") is not None and not pd.isna(e_result[0]["twin_OR"]):
        e = e_result[0]
        p_e_str = f"P = {e['twin_P']:.3f}" if e["twin_P"] >= 0.001 else "P < 0.001"
        p2 = doc.add_paragraph()
        add_text_with_refs(p2,
            f"A notable exception was Definition E (5-HT3 antagonist use), where twin pregnancy "
            f"was associated with significantly higher 5-HT3 antagonist use "
            f"(adjusted OR {e['twin_OR']:.2f}, 95% CI {e['twin_CI_lower']:.2f}–{e['twin_CI_upper']:.2f}; "
            f"{p_e_str}; reduced model with 6 covariates due to sparse events, {e['Events']} events). "
            f"This suggests that while overall antiemetic use is similar, twins may have a different "
            f"pattern of antiemetic prescribing, with more frequent use of serotonin receptor antagonists."
        )

    # Poisson
    g_result = [r for r in SS["regression"] if r["Definition"] == "G"]
    if g_result and g_result[0].get("twin_OR") is not None and not pd.isna(g_result[0]["twin_OR"]):
        g = g_result[0]
        p_g_str = f"P = {g['twin_P']:.3f}" if g["twin_P"] >= 0.001 else "P < 0.001"
        p3 = doc.add_paragraph()
        add_text_with_refs(p3,
            f"In Poisson regression (Definition G), the antiemetic drug count did not differ "
            f"significantly between groups (twin IRR {g['twin_OR']:.2f}, 95% CI "
            f"{g['twin_CI_lower']:.2f}–{g['twin_CI_upper']:.2f}; {p_g_str})."
        )

    # Stratified analysis (Definition H)
    add_heading(doc, "Hypotension-Stratified Analysis (Definition H)")
    strata = SS.get("stratified", [])
    if strata:
        p4 = doc.add_paragraph()
        hypo_pos = [x for x in strata if "+" in x["Stratum"]]
        hypo_neg = [x for x in strata if "−" in x["Stratum"]]
        text_parts = []
        if hypo_pos and not pd.isna(hypo_pos[0]["twin_OR"]):
            h_p = hypo_pos[0]
            text_parts.append(
                f"In the hypotension-present subgroup (n = {h_p['n']}), "
                f"IONV rates were {h_p['singleton_rate']:.1f}% (singleton) vs "
                f"{h_p['twin_rate']:.1f}% (twin), with twin aOR {h_p['twin_OR']:.2f} "
                f"(95% CI {h_p['twin_CI_lower']:.2f}–{h_p['twin_CI_upper']:.2f}; "
                f"P = {h_p['twin_P']:.3f})"
            )
        if hypo_neg and not pd.isna(hypo_neg[0]["twin_OR"]):
            h_n = hypo_neg[0]
            text_parts.append(
                f"In the hypotension-absent subgroup (n = {h_n['n']}), "
                f"rates were {h_n['singleton_rate']:.1f}% vs "
                f"{h_n['twin_rate']:.1f}%, with twin aOR {h_n['twin_OR']:.2f} "
                f"(95% CI {h_n['twin_CI_lower']:.2f}–{h_n['twin_CI_upper']:.2f}; "
                f"P = {h_n['twin_P']:.3f})"
            )
        add_text_with_refs(p4,
            ". ".join(text_parts) + ". "
            "Neither stratum showed a significant association between twin pregnancy and IONV, "
            "indicating that the null finding was consistent regardless of hypotension status (Fig. 9)."
        )

    # Insert Fig 9 - Stratified hypotension
    if (FIG / "fig9_stratified_hypotension.png").exists():
        doc.add_paragraph()
        p_fig9 = doc.add_paragraph()
        p_fig9.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p_fig9.add_run().add_picture(str(FIG / "fig9_stratified_hypotension.png"), width=Inches(5.0))
        p_cap9 = doc.add_paragraph("Figure 9. IONV rates stratified by intraoperative hypotension.")
        p_cap9.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p_cap9.runs[0].italic = True
        p_cap9.runs[0].font.size = Pt(9)
        p_cap9.paragraph_format.space_before = Pt(12)

# Save
doc.save(str(BASE / "manuscript_methods_results.docx"))
print("Manuscript saved:", BASE / "manuscript_methods_results.docx")

# ============================================================
# PPTX - editable figures (one per slide)
# ============================================================
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt

prs = Presentation()
prs.slide_width = PptInches(13.333)
prs.slide_height = PptInches(7.5)

fig_files = sorted(FIG.glob("*.png"))
fig_titles = {
    "fig1_ionv_rates.png": "Figure 1: IONV Rates by Group",
    "fig2_forest_primary.png": "Figure 2: Forest Plot — Primary Outcome",
    "fig3_forest_secondary.png": "Figure 3: Forest Plot — Secondary Outcome",
    "fig4_ionv_timing.png": "Figure 4: IONV Rates by Timing Phase",
    "fig5_temporal_trend.png": "Figure 5: Temporal Trend of IONV Rates",
    "fig6_antiemetic_drugs.png": "Figure 6: Antiemetic Drug Usage Comparison",
    "fig7_sensitivity_forest.png": "Figure 7: Sensitivity — Twin aOR Across Definitions",
    "fig8_sensitivity_rates.png": "Figure 8: Sensitivity — IONV Rates by Definition",
    "fig9_stratified_hypotension.png": "Figure 9: Definition H — Hypotension-Stratified IONV",
}

for fig_path in fig_files:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    
    # Title
    from pptx.util import Emu
    from pptx.dml.color import RGBColor as PptRGB
    txBox = slide.shapes.add_textbox(PptInches(0.5), PptInches(0.2), PptInches(12), PptInches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = fig_titles.get(fig_path.name, fig_path.stem)
    p.font.size = PptPt(20)
    p.font.bold = True
    
    # Image (centered)
    img = slide.shapes.add_picture(str(fig_path), PptInches(1.5), PptInches(1.0), PptInches(10), PptInches(5.5))

prs.save(str(BASE / "figures.pptx"))
print("PPTX saved:", BASE / "figures.pptx")

print("Done!")
