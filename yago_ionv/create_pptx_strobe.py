"""
Generate PPTX with ALL figures as fully editable PowerPoint objects.
- Fig 1: Flowchart → editable shapes (rectangles, arrows, textboxes)
- Fig 2, 6: Bar charts → native PowerPoint charts (ClusteredBar)
- Fig 3, 4, 5, 7: Forest plots → editable shapes (lines, ovals, textboxes)

All text, shapes, and chart data are editable in PowerPoint.
"""
import json
from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

BASE = Path(__file__).resolve().parent

# Load all stats
with open(BASE / "def_e_stats.json") as f:
    M = json.load(f)
with open(BASE / "excl_sensitivity_stats.json") as f:
    E = json.load(f)
with open(BASE / "flowchart_counts.json") as f:
    F = json.load(f)
with open(BASE / "bootstrap_results.json") as f:
    B = json.load(f)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x4C, 0x72, 0xB0)
ORANGE = RGBColor(0xDD, 0x84, 0x52)
RED = RGBColor(0xC4, 0x4E, 0x52)
GRAY = RGBColor(0x8C, 0x8C, 0x8C)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BLUE = RGBColor(0x1A, 0x23, 0x7E)

# Color palettes for flowchart
BOX_MAIN_FILL = RGBColor(0xE8, 0xF0, 0xFE)
BOX_MAIN_EDGE = RGBColor(0x15, 0x65, 0xC0)
BOX_EXCL_FILL = RGBColor(0xFF, 0xF3, 0xE0)
BOX_EXCL_EDGE = RGBColor(0xE6, 0x51, 0x00)
BOX_FINAL_FILL = RGBColor(0xE8, 0xF5, 0xE9)
BOX_FINAL_EDGE = RGBColor(0x2E, 0x7D, 0x32)
BOX_SUB_FILL = RGBColor(0xF3, 0xE5, 0xF5)
BOX_SUB_EDGE = RGBColor(0x6A, 0x1B, 0x9A)


def add_slide_title(slide, text):
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER


def add_slide_caption(slide, text, top=Inches(6.8)):
    txBox = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.333), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER


def add_box(slide, left, top, width, height, fill_color, edge_color, texts,
            font_sizes=None, bolds=None, alignments=None):
    """Add a rounded rectangle with multiple lines of text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = edge_color
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)

    for i, text in enumerate(texts):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        fs = Pt(9)
        if font_sizes and i < len(font_sizes):
            fs = font_sizes[i]
        for run in p.runs:
            run.font.size = fs
            run.font.color.rgb = BLACK
        if bolds and i < len(bolds) and bolds[i]:
            for run in p.runs:
                run.font.bold = True
    return shape


def add_arrow_down(slide, x_center, y_start, y_end):
    """Add a downward arrow connector."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, x_center - Inches(0.1), y_start,
        Inches(0.2), y_end - y_start)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x55, 0x55, 0x55)
    shape.line.fill.background()
    return shape


def add_arrow_right(slide, x_start, y_center, x_end):
    """Add a rightward arrow connector."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, x_start, y_center - Inches(0.08),
        x_end - x_start, Inches(0.16))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xAA, 0x55, 0x00)
    shape.line.fill.background()
    return shape


# ============================================================
# SLIDE 1: STROBE Flowchart (editable shapes)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_title(slide, "Figure 1. STROBE Flow Diagram — Participant Selection")

# Layout coordinates
cx = Inches(4.0)    # center of main column
ex = Inches(9.5)    # center of exclusion column
box_w = Inches(4.5)
box_w_excl = Inches(3.5)

# Row 1: Total assessed
y1 = Inches(0.8)
h1 = Inches(0.85)
add_box(slide, cx - box_w/2, y1, box_w, h1,
        BOX_MAIN_FILL, BOX_MAIN_EDGE,
        ["Cesarean deliveries assessed for eligibility",
         f"(Apr 2014 – Oct 2024)",
         f"N = {F['total']['n']:,}  (Singleton {F['total']['n_s']:,}  /  Twin {F['total']['n_t']:,})"],
        font_sizes=[Pt(9), Pt(8), Pt(9)],
        bolds=[True, False, False])

# Arrow down
add_arrow_down(slide, cx, y1 + h1, y1 + h1 + Inches(0.3))

# Arrow right to exclusion
add_arrow_right(slide, cx + box_w/2, y1 + h1 + Inches(0.15), ex - box_w_excl/2)

# Exclusion box
excl_texts = [f"Excluded  (n = {F['total_excluded']['n']})"]
for step in F["exclusion_steps"]:
    excl_texts.append(f"{step['reason']}: n={step['n']} (S {step['n_s']}, T {step['n_t']})")
excl_h = Inches(0.3 + 0.22 * len(F["exclusion_steps"]))
add_box(slide, ex - box_w_excl/2, y1 + Inches(0.25), box_w_excl, excl_h,
        BOX_EXCL_FILL, BOX_EXCL_EDGE,
        excl_texts,
        font_sizes=[Pt(8)] + [Pt(7)] * len(F["exclusion_steps"]),
        bolds=[True] + [False] * len(F["exclusion_steps"]))

# Row 2: Eligible
y2 = y1 + h1 + Inches(0.3) + Inches(0.2)
h2 = Inches(0.65)
add_box(slide, cx - box_w/2, y2, box_w, h2,
        BOX_MAIN_FILL, BOX_MAIN_EDGE,
        ["Eligible for analysis",
         f"N = {F['eligible']['n']:,}  (Singleton {F['eligible']['n_s']:,}  /  Twin {F['eligible']['n_t']:,})"],
        font_sizes=[Pt(9), Pt(9)],
        bolds=[True, False])

# Arrow down
add_arrow_down(slide, cx, y2 + h2, y2 + h2 + Inches(0.25))

# Arrow right to preop exclusion
add_arrow_right(slide, cx + box_w/2, y2 + h2 + Inches(0.12), ex - box_w_excl/2)

# Preop antiemetic exclusion
preop_y = y2 + h2 - Inches(0.05)
add_box(slide, ex - box_w_excl/2, preop_y, box_w_excl, Inches(0.55),
        BOX_EXCL_FILL, BOX_EXCL_EDGE,
        [f"Prophylactic antiemetic before anesthesia",
         f"n = {F['preop_antiemetic']['n']}  (S {F['preop_antiemetic']['n_s']}, T {F['preop_antiemetic']['n_t']})"],
        font_sizes=[Pt(8), Pt(8)],
        bolds=[True, False])

# Row 3: Primary analysis
y3 = y2 + h2 + Inches(0.25) + Inches(0.15)
h3 = Inches(0.65)
add_box(slide, cx - box_w/2, y3, box_w, h3,
        BOX_FINAL_FILL, BOX_FINAL_EDGE,
        ["Primary analysis cohort",
         f"N = {F['primary_analysis']['n']:,}  (Singleton {F['primary_analysis']['n_s']:,}  /  Twin {F['primary_analysis']['n_t']:,})"],
        font_sizes=[Pt(9), Pt(9)],
        bolds=[True, False])

# Split arrows to Singleton / Twin
y_split1 = y3 + h3 + Inches(0.1)
s_box_w = Inches(2.2)
s_box_h = Inches(0.55)

# Singleton box
s_left = cx - Inches(2.5)
add_arrow_down(slide, cx - Inches(1.2), y3 + h3, y_split1 + Inches(0.2))
add_box(slide, s_left, y_split1 + Inches(0.2), s_box_w, s_box_h,
        BOX_FINAL_FILL, BOX_FINAL_EDGE,
        [f"Singleton", f"n = {F['primary_analysis']['n_s']:,}"],
        font_sizes=[Pt(9), Pt(9)],
        bolds=[True, False])

# Twin box
t_left = cx + Inches(0.3)
add_arrow_down(slide, cx + Inches(1.2), y3 + h3, y_split1 + Inches(0.2))
add_box(slide, t_left, y_split1 + Inches(0.2), s_box_w, s_box_h,
        BOX_FINAL_FILL, BOX_FINAL_EDGE,
        [f"Twin", f"n = {F['primary_analysis']['n_t']:,}"],
        font_sizes=[Pt(9), Pt(9)],
        bolds=[True, False])

# Row 4: Sensitivity subgroup
y4 = y_split1 + Inches(0.2) + s_box_h + Inches(0.15)
add_arrow_down(slide, cx, y_split1 + Inches(0.2) + s_box_h, y4 + Inches(0.1))

# Additional exclusion box (right)
sens_texts = ["Additional exclusion for sensitivity analysis",
              f"n = {F['exclusion_sensitivity_total']['n']:,} (with overlap)"]
for name, counts in F["exclusion_sensitivity"].items():
    sens_texts.append(f"{name}: n={counts['n']}")
sens_h = Inches(0.3 + 0.2 * (len(F["exclusion_sensitivity"]) + 1))
add_box(slide, ex - box_w_excl/2, y4, box_w_excl, sens_h,
        BOX_EXCL_FILL, BOX_EXCL_EDGE,
        sens_texts,
        font_sizes=[Pt(7)] * len(sens_texts),
        bolds=[True, True] + [False] * len(F["exclusion_sensitivity"]))

add_arrow_right(slide, cx + box_w/2 - Inches(1), y4 + Inches(0.15), ex - box_w_excl/2)

# Subgroup box
y5 = y4 + Inches(0.25) + Inches(0.1)
h5 = Inches(0.7)
add_box(slide, cx - box_w/2, y5, box_w, h5,
        BOX_SUB_FILL, BOX_SUB_EDGE,
        ["Sensitivity analysis subgroup (Elective, low-risk)",
         f"N = {F['subgroup_analysis']['n']:,}  (Singleton {F['subgroup_analysis']['n_s']:,}  /  Twin {F['subgroup_analysis']['n_t']:,})"],
        font_sizes=[Pt(9), Pt(9)],
        bolds=[True, False])

# Split to subgroup singleton/twin
y_split2 = y5 + h5 + Inches(0.05)
add_arrow_down(slide, cx - Inches(1.2), y5 + h5, y_split2 + Inches(0.12))
add_box(slide, s_left, y_split2 + Inches(0.12), s_box_w, s_box_h,
        BOX_SUB_FILL, BOX_SUB_EDGE,
        [f"Singleton", f"n = {F['subgroup_analysis']['n_s']:,}"],
        font_sizes=[Pt(9), Pt(9)],
        bolds=[True, False])

add_arrow_down(slide, cx + Inches(1.2), y5 + h5, y_split2 + Inches(0.12))
add_box(slide, t_left, y_split2 + Inches(0.12), s_box_w, s_box_h,
        BOX_SUB_FILL, BOX_SUB_EDGE,
        [f"Twin", f"n = {F['subgroup_analysis']['n_t']:,}"],
        font_sizes=[Pt(9), Pt(9)],
        bolds=[True, False])

print("Slide 1 (Flowchart) created")


# ============================================================
# SLIDE 2: IONV Rates Bar Chart (native PowerPoint chart)
# ============================================================
def create_bar_chart_slide(prs, title, caption, outcomes, keys, bar_labels):
    """Create a grouped bar chart slide with native PowerPoint chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, title)

    chart_data = CategoryChartData()
    chart_data.categories = bar_labels

    s_vals = [outcomes[k]["singleton_pct"] for k in keys]
    t_vals = [outcomes[k]["twin_pct"] for k in keys]
    chart_data.add_series("Singleton", s_vals)
    chart_data.add_series("Twin", t_vals)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1.0), Inches(0.9), Inches(11.0), Inches(5.5),
        chart_data)
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(11)

    # Style series
    series_s = chart.series[0]
    series_t = chart.series[1]
    series_s.format.fill.solid()
    series_s.format.fill.fore_color.rgb = BLUE
    series_t.format.fill.solid()
    series_t.format.fill.fore_color.rgb = ORANGE

    # Data labels
    for series in [series_s, series_t]:
        series.has_data_labels = True
        dl = series.data_labels
        dl.number_format = '0.0"%"'
        dl.font.size = Pt(9)
        dl.font.bold = True

    # Y-axis
    val_axis = chart.value_axis
    val_axis.axis_title.text_frame.paragraphs[0].text = "Rate (%)" if val_axis.has_title else ""
    try:
        val_axis.has_title = True
        val_axis.axis_title.text_frame.paragraphs[0].text = "Rate (%)"
        val_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(11)
    except Exception:
        pass

    add_slide_caption(slide, caption)
    return slide


# Fig 2: Full cohort bar chart
create_bar_chart_slide(
    prs,
    "Figure 2. IONV Rates: Broad vs Narrow Antiemetic Definition",
    "Comparison of antiemetic use between singleton and twin groups. "
    "Broad = all 7 antiemetics; Narrow = 5-HT3 antagonists (ondansetron, granisetron) only.",
    M["outcomes"],
    ["A-Primary", "A-Secondary", "E-Primary", "E-Secondary"],
    ["Broad: Primary\n(before delivery)",
     "Broad: Secondary\n(any phase)",
     "Narrow: Primary\n(5-HT3, before delivery)",
     "Narrow: Secondary\n(5-HT3, any phase)"]
)
print("Slide 2 (Bar chart — full cohort) created")


# ============================================================
# FOREST PLOT HELPER (editable shapes)
# ============================================================
def create_forest_slide(prs, title, caption, data_rows, x_label="Adjusted Odds Ratio (95% CI)",
                        ref_line=1.0, max_x=None):
    """
    Create a forest plot using editable shapes.
    data_rows: list of dicts with keys: label, or_val, ci_lo, ci_hi, p_val, highlight (bool)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, title)

    n = len(data_rows)
    plot_left = Inches(4.0)
    plot_right = Inches(11.5)
    plot_top = Inches(1.2)
    row_height = min(Inches(5.0) / max(n, 1), Inches(0.45))
    plot_height = row_height * n

    if max_x is None:
        max_x = max(r["ci_hi"] for r in data_rows) * 1.3
        max_x = max(max_x, ref_line * 3)
    min_x = 0

    def x_to_pos(val):
        frac = (val - min_x) / (max_x - min_x)
        return plot_left + int(frac * (plot_right - plot_left))

    # Reference line at OR=1
    ref_x = x_to_pos(ref_line)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, ref_x, plot_top, Pt(1.5), plot_height)
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xFF, 0x00, 0x00)
    line.line.fill.background()

    # Label "OR = 1"
    txb = slide.shapes.add_textbox(ref_x - Inches(0.3), plot_top + plot_height + Inches(0.05),
                                    Inches(0.6), Inches(0.25))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.text = "OR=1"
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    p.alignment = PP_ALIGN.CENTER

    # X-axis label
    txb = slide.shapes.add_textbox(plot_left, plot_top + plot_height + Inches(0.3),
                                    plot_right - plot_left, Inches(0.3))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.text = x_label
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER

    for i, row in enumerate(data_rows):
        y_center = plot_top + row_height * i + row_height / 2

        # Label (left side)
        txb = slide.shapes.add_textbox(Inches(0.3), y_center - Inches(0.15),
                                        Inches(3.5), Inches(0.3))
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = row["label"]
        p.font.size = Pt(8)
        p.alignment = PP_ALIGN.RIGHT

        # CI line (horizontal bar)
        ci_lo_x = x_to_pos(max(row["ci_lo"], min_x))
        ci_hi_x = x_to_pos(min(row["ci_hi"], max_x))
        color = RED if row.get("highlight", False) else BLUE
        if ci_hi_x > ci_lo_x:
            ci_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, ci_lo_x, y_center - Pt(1.5),
                ci_hi_x - ci_lo_x, Pt(3))
            ci_line.fill.solid()
            ci_line.fill.fore_color.rgb = color
            ci_line.line.fill.background()

        # Point estimate (diamond)
        or_x = x_to_pos(min(row["or_val"], max_x))
        diamond_size = Inches(0.12)
        diamond = slide.shapes.add_shape(
            MSO_SHAPE.DIAMOND, or_x - diamond_size/2, y_center - diamond_size/2,
            diamond_size, diamond_size)
        diamond.fill.solid()
        diamond.fill.fore_color.rgb = color
        diamond.line.color.rgb = color

        # Annotation (right side)
        sig = ""
        if row["p_val"] < 0.001:
            sig = "***"
        elif row["p_val"] < 0.01:
            sig = "**"
        elif row["p_val"] < 0.05:
            sig = "*"
        p_str = "P<0.001" if row["p_val"] < 0.001 else f"P={row['p_val']:.3f}"
        annotation = f"OR {row['or_val']:.2f} [{row['ci_lo']:.2f}–{row['ci_hi']:.2f}] {p_str}{sig}"

        txb = slide.shapes.add_textbox(plot_right + Inches(0.1), y_center - Inches(0.12),
                                        Inches(1.5), Inches(0.25))
        tf = txb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = annotation
        p.font.size = Pt(7)
        p.font.color.rgb = color

    add_slide_caption(slide, caption)
    return slide


# ============================================================
# SLIDE 3: Forest Plot — Narrow Definition Primary (Full cohort)
# ============================================================
or_table_e = pd.read_csv(BASE / "tables_e" / "logistic_E-Primary.csv")
label_map = {
    "twin": "Twin pregnancy", "年齢(歳)": "Age (per year)",
    "BMI": "BMI (per kg/m²)", "GA_weeks": "GA (per week)",
    "emergency": "Emergency CS", "prior_cs": "Prior CS",
    "HDP": "HDP", "epidural": "Epidural anesthesia",
    "手術時間_min": "Surgery time (per min)", "hypotension": "Hypotension (SBP<90)",
}
or_table_e["label"] = or_table_e["Variable"].map(label_map).fillna(or_table_e["Variable"])
or_table_e = or_table_e.sort_values("P-value", ascending=False)

forest_data_3 = []
for _, row in or_table_e.iterrows():
    forest_data_3.append({
        "label": row["label"],
        "or_val": row["OR"],
        "ci_lo": row["95% CI lower"],
        "ci_hi": row["95% CI upper"],
        "p_val": row["P-value"],
        "highlight": row["Variable"] == "twin",
    })

create_forest_slide(
    prs,
    "Figure 3. Forest Plot — Narrow Antiemetic Definition (5-HT3)",
    "Multivariable logistic regression for 5-HT3 antagonist use (primary outcome: before delivery). "
    "Reduced model with 6 covariates. N = 3,153.",
    forest_data_3)
print("Slide 3 (Forest — E-Primary) created")


# ============================================================
# SLIDE 4: Broad vs Narrow Comparison Forest (Full cohort)
# ============================================================
compare_data_4 = []
for key in ["E-Secondary", "E-Primary", "A-Secondary", "A-Primary"]:
    r = M["regression"][key]
    compare_data_4.append({
        "label": r["label"],
        "or_val": r["twin_OR"],
        "ci_lo": r["twin_CI_lower"],
        "ci_hi": r["twin_CI_upper"],
        "p_val": r["twin_P"],
        "highlight": r["twin_P"] < 0.05,
    })

create_forest_slide(
    prs,
    "Figure 4. Twin Effect: Broad vs Narrow Antiemetic Definition",
    "Adjusted odds ratios for twin pregnancy across broad (all 7 drugs) "
    "and narrow (5-HT3 antagonists only) IONV definitions. N = 3,153.",
    compare_data_4)
print("Slide 4 (Broad vs Narrow) created")


# ============================================================
# SLIDE 5: Covariate Sensitivity Forest Plot
# ============================================================
cov_df = pd.read_csv(BASE / "tables_e" / "covariate_sensitivity.csv")
cov_df_rev = cov_df.iloc[::-1].reset_index(drop=True)

forest_data_5 = []
for _, row in cov_df_rev.iterrows():
    forest_data_5.append({
        "label": row["Model"],
        "or_val": row["aOR"],
        "ci_lo": row["CI_lower"],
        "ci_hi": row["CI_upper"],
        "p_val": row["P"],
        "highlight": row["P"] < 0.05,
    })

create_forest_slide(
    prs,
    "Figure 5. Covariate Sensitivity Analysis",
    "Robustness of twin effect across 18 covariate models for narrow-definition "
    "primary outcome. All models show P < 0.05.",
    forest_data_5)
print("Slide 5 (Covariate sensitivity) created")


# ============================================================
# SLIDE 6: IONV Rates Bar Chart — Subgroup
# ============================================================
create_bar_chart_slide(
    prs,
    "Figure 6. IONV Rates — Elective Low-Risk Subgroup (N=663)",
    "IONV rates in sensitivity subgroup after excluding emergency CS, "
    "prior CS, HDP, preoperative steroid. Singleton n=479, Twin n=184.",
    E["outcomes"],
    ["A-Primary", "A-Secondary", "E-Primary", "E-Secondary"],
    ["Broad: Primary\n(before delivery)",
     "Broad: Secondary\n(any phase)",
     "Narrow: Primary\n(5-HT3, before delivery)",
     "Narrow: Secondary\n(5-HT3, any phase)"]
)
print("Slide 6 (Bar chart — subgroup) created")


# ============================================================
# SLIDE 7: Broad vs Narrow Forest — Subgroup
# ============================================================
compare_data_7 = []
for key in ["E-Secondary", "E-Primary", "A-Secondary", "A-Primary"]:
    r = E["regression"][key]
    compare_data_7.append({
        "label": r["label"],
        "or_val": r["twin_OR"],
        "ci_lo": r["twin_CI_lower"],
        "ci_hi": r["twin_CI_upper"],
        "p_val": r["twin_P"],
        "highlight": r["twin_P"] < 0.05,
    })

create_forest_slide(
    prs,
    "Figure 7. Twin Effect — Elective Low-Risk Subgroup",
    "Adjusted odds ratios in sensitivity subgroup (N=663). "
    "Narrow-definition aOR increased from 3.18 (full) to 13.59 (subgroup).",
    compare_data_7,
    max_x=120)
print("Slide 7 (Subgroup forest) created")


# ============================================================
# SLIDE 8: Bootstrap Validation Forest Plot (editable shapes)
# ============================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_title(slide8, "Figure 8. Wald vs Bootstrap Confidence Interval Comparison")
add_slide_caption(slide8,
    "Stratified bootstrap (10,000 resamples) validation of Wald-based CIs. "
    "BCa CI for narrow primary (full cohort) excludes 1, confirming robust twin effect.")

# Draw bootstrap comparison as a table of editable textboxes
bm = B["main_cohort"]
bs = B["subgroup"]

def fmt_bca_pptx(r):
    lo = r.get("bca_CI_lower")
    hi = r.get("bca_CI_upper")
    if lo is not None and hi is not None and hi < 1e6:
        return f"{lo:.2f}\u2013{hi:.2f}"
    return "Unstable"

boot_table_data = [
    ["Full cohort", "Broad, Primary", bm["A-Primary"]],
    ["Full cohort", "Broad, Secondary", bm["A-Secondary"]],
    ["Full cohort", "Narrow, Primary", bm["E-Primary"]],
    ["Full cohort", "Narrow, Secondary", bm["E-Secondary"]],
    ["Low-risk", "Broad, Primary", bs["A-Primary"]],
    ["Low-risk", "Broad, Secondary", bs["A-Secondary"]],
    ["Low-risk", "Narrow, Primary", bs["E-Primary"]],
    ["Low-risk", "Narrow, Secondary", bs["E-Secondary"]],
]

headers = ["Cohort", "Outcome", "aOR", "Wald 95% CI", "BCa 95% CI", "Conv."]
col_widths = [Inches(1.5), Inches(2.0), Inches(1.0), Inches(2.2), Inches(2.2), Inches(1.0)]
table_left = Inches(1.5)
table_top = Inches(1.0)
row_height = Inches(0.45)

# Header row
x = table_left
for h_idx, (header, w) in enumerate(zip(headers, col_widths)):
    box = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, table_top, w, row_height)
    box.fill.solid()
    box.fill.fore_color.rgb = DARK_BLUE
    box.line.color.rgb = WHITE
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = header
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = WHITE
    x += w

# Data rows
for r_idx, (cohort, outcome, r) in enumerate(boot_table_data):
    y = table_top + row_height * (r_idx + 1)
    row_fill = RGBColor(0xF5, 0xF5, 0xF5) if r_idx % 2 == 0 else WHITE
    # Highlight narrow primary full cohort
    is_key = (cohort == "Full cohort" and "Narrow, Primary" in outcome)
    if is_key:
        row_fill = RGBColor(0xE8, 0xF5, 0xE9)
    vals = [
        cohort if r_idx in [0, 4] else "",
        outcome,
        f"{r['point_aOR']:.2f}",
        f"{r['wald_CI_lower']:.2f}\u2013{r['wald_CI_upper']:.2f}",
        fmt_bca_pptx(r),
        f"{r['convergence_pct']}%",
    ]
    x = table_left
    for v_idx, (val, w) in enumerate(zip(vals, col_widths)):
        box = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, row_height)
        box.fill.solid()
        box.fill.fore_color.rgb = row_fill
        box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        box.line.width = Pt(0.5)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = val
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(9)
            run.font.color.rgb = BLACK
            if is_key:
                run.font.bold = True
        x += w

print("Slide 8 (Bootstrap comparison) created")


# ============================================================
# SAVE
# ============================================================
out_path = BASE / "figures_strobe.pptx"
prs.save(str(out_path))
print(f"\nEditable PPTX saved to {out_path}")
print("All 8 slides with editable shapes/charts generated.")
