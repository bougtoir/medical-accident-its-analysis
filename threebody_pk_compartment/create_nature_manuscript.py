#!/usr/bin/env python3
"""
Generate Nature Article manuscript: Three-body scattering as pharmacokinetic
compartment model.

Outputs:
  - manuscript_nature.docx (full manuscript with inline figures)
  - figures_editable.pptx (editable figures, 1 per slide)
"""

from __future__ import annotations

import os
import re

from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT

BASEDIR = os.path.dirname(os.path.abspath(__file__))
FIGDIR_NATURE = os.path.join(BASEDIR, "figures_nature")
FIGDIR_ORIG = os.path.join(BASEDIR, "figures")
# Prefer Nature-quality figures; fall back to original
FIGDIR = FIGDIR_NATURE if os.path.isdir(FIGDIR_NATURE) else FIGDIR_ORIG
OUTDIR = os.path.join(BASEDIR, "manuscript")
os.makedirs(OUTDIR, exist_ok=True)


# ---------------------------------------------------------------------------
# Helper: superscript citations using font-based superscript
# ---------------------------------------------------------------------------

def add_text_with_refs(paragraph, text):
    """
    Parse text with {N} or {N-M} citation markers and add as runs
    with font superscript.
    """
    parts = re.split(r'(\{[^}]+\})', text)
    for part in parts:
        if part.startswith('{') and part.endswith('}'):
            run = paragraph.add_run(part[1:-1])
            run.font.superscript = True
            run.font.size = Pt(8)
        else:
            run = paragraph.add_run(part)
            run.font.size = Pt(11)


def add_heading(doc, text, level=1):
    """Add a heading with appropriate formatting."""
    h = doc.add_heading(text, level=level)
    return h


def add_paragraph_with_refs(doc, text, style=None):
    """Add a paragraph with citation superscripts."""
    p = doc.add_paragraph(style=style)
    add_text_with_refs(p, text)
    return p


def add_figure(doc, fig_path, caption, fig_num):
    """Add a figure inline with caption below."""
    if os.path.exists(fig_path):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run()
        run.add_picture(fig_path, width=Inches(6.0))

    # Caption
    cap = doc.add_paragraph()
    cap.paragraph_format.space_before = Pt(6)
    cap.paragraph_format.space_after = Pt(12)
    run_bold = cap.add_run(f"Figure {fig_num} | ")
    run_bold.bold = True
    run_bold.font.size = Pt(10)
    add_text_with_refs(cap, caption)


# ---------------------------------------------------------------------------
# Manuscript content
# ---------------------------------------------------------------------------

TITLE = ("Pharmacokinetic compartment theory solves the statistical "
         "three-body problem")

AUTHORS = (
    "Tatsuki Onishi"
)

AFFILIATIONS = (
    "Independent Researcher"
)

ABSTRACT = (
    "The gravitational three-body problem remains analytically intractable "
    "after 250 years. Here we show that chaotic three-body scattering maps "
    "onto a three-compartment pharmacokinetic (PK) model, where "
    "binary configurations are compartments, configuration transitions are "
    "inter-compartmental transfers, and system dissolution is elimination. "
    "Using 15,000 N-body simulations across 64 mass configurations, we "
    "demonstrate that: (i) the linear PK model reproduces the multi-exponential "
    "lifetime distribution predicted by random-walk theory; (ii) a nonlinear "
    "(Michaelis-Menten) elimination term quantitatively captures the power-law "
    "tail from sticky chaos, reducing tail RMSE by 36-42%; and (iii) population "
    "PK (mixed-effects) modelling reveals allometric scaling of mean residence "
    "time with mass ratios (MRT proportional to mu_12 to the power 2.08 times "
    "mu_out to the power -1.68, R-squared = 0.67). The reverse mapping yields a "
    "stability analysis framework for target-mediated drug disposition derived "
    "from celestial mechanics. "
    "This cross-disciplinary bridge opens new computational strategies in both "
    "celestial mechanics and quantitative pharmacology."
)

INTRO = [
    ("The three-body problem\u2014predicting the long-term motion of three "
     "gravitationally interacting masses\u2014has resisted analytical solution "
     "since Newton.{1} Poincar\u00e9 proved its non-integrability,{2} and "
     "modern numerical studies confirm that generic three-body encounters are "
     "chaotic, with outcomes sensitive to initial conditions at the level of "
     "floating-point precision.{3} Nevertheless, the statistical properties "
     "of the outcomes are remarkably regular: the lifetime distribution follows "
     "a multi-exponential decay,{4,5} and the final-state distributions can be "
     "predicted from phase-space volume arguments.{6,7}"),

    ("Recently, Ginat and Perets showed that chaotic three-body scattering "
     "can be decomposed into a sequence of independent probabilistic "
     "excursions\u2014each a temporary binary plus a distant single body\u2014and "
     "derived an analytical, statistical solution as a random walk in "
     "binary binding energy.{5} Stone and Leigh independently demonstrated "
     "the predictive power of flux-based statistical mechanics for three-body "
     "outcomes.{6} These advances establish that the ergodic core of the "
     "three-body problem is amenable to stochastic modelling."),

    ("Here we observe that this stochastic structure is formally identical to "
     "pharmacokinetic (PK) compartmental analysis\u2014the standard framework "
     "for modelling drug absorption, distribution, and elimination in the "
     "body.{8,9} In PK theory, the body is divided into compartments "
     "(e.g. plasma, tissues), drug transfers between compartments at "
     "first-order rates, and is eliminated irreversibly. The mathematics is "
     "a continuous-time Markov chain (CTMC) whose generator is the PK rate "
     "matrix.{10}"),

    ("We propose that three-body scattering is a three-compartment PK system: "
     "each binary configuration (which body pair is bound) is a compartment; "
     "transitions between configurations are inter-compartmental transfers; "
     "escape of one body is elimination. This mapping is not merely analogical "
     "\u2014it is an exact structural equivalence at the level of the master "
     "equation. We exploit it to import the full toolkit of PK analysis "
     "(half-lives, mean residence time, bioequivalence, population modelling, "
     "nonlinear kinetics) into celestial mechanics, and conversely to bring "
     "dynamical-systems stability theory into pharmacology."),
]

RESULTS_SECTIONS = {
    "Linear PK model reproduces multi-exponential decay": [
        ("We performed 15,000 three-body scattering simulations using a "
         "symplectic leapfrog integrator (Methods), spanning three mass "
         "configurations: equal mass (1:1:1), unequal mass (1:2:0.5), and "
         "democratic initial conditions (1:1:1, symmetric start). For each "
         "simulation, we recorded the full sequence of binary configurations "
         "visited and dwell times in each (Fig. 1)."),

        ("From the transition counts and dwell times, we estimated the 9 PK "
         "rate parameters (6 inter-compartmental transfer rates k_ij and 3 "
         "elimination rates k_ei) by maximum-likelihood estimation for a "
         "continuous-time Markov chain: k_ij = N_ij / T_i, where N_ij is the "
         "number of observed i-to-j transitions and T_i is total dwell time in "
         "compartment i.{10}"),

        ("The resulting linear PK model predicts a survival function S(t) = "
         "sum of w_i times exp(-lambda_i times t), a sum of three exponentials "
         "with decay rates given by the eigenvalues of the rate matrix (Fig. 2). "
         "Kolmogorov-Smirnov tests reject the single-exponential hypothesis "
         "(p < 0.001 for all three configurations), confirming that the "
         "three-compartment model is necessary and sufficient to describe "
         "the intermediate-timescale dynamics."),

        ("The MLE-estimated rates agree quantitatively with the phase-space "
         "flux prediction from statistical mechanics.{6} For equal "
         "masses, all transfer rates converge to the same value by symmetry "
         "(k_ij approximately 0.015 per dynamical time), validating both the "
         "PK framework and the ergodic hypothesis underlying the statistical "
         "approach."),
    ],

    "Nonlinear PK captures sticky chaos": [
        ("While the linear PK model fits the bulk of the lifetime distribution, "
         "it systematically underestimates the long-lived tail\u2014a signature of "
         "'sticky chaos', where trajectories become trapped near periodic "
         "orbits for anomalously long times.{11,12} In PK language, this "
         "corresponds to saturable (capacity-limited) elimination: when the "
         "system lingers near a stable periodic orbit, the escape rate "
         "decreases rather than remaining constant."),

        ("We model this by replacing the first-order elimination k_e times P "
         "with Michaelis-Menten kinetics: V_max times P divided by (K_m + P). "
         "This hybrid model\u2014linear inter-compartmental transfer with nonlinear "
         "elimination\u2014produces a survival function whose tail decays as "
         "t to the power negative alpha times exp(-lambda_slow times t), "
         "capturing the power-law correction (Fig. 3)."),

        ("Fitting the hybrid model to all three mass configurations yields "
         "power-law exponents alpha approximately equal to 1.1-2.6, with tail "
         "RMSE improvement of 36-42% over the linear model for unequal and "
         "democratic configurations. For the equal-mass case, the linear model "
         "is already adequate (no significant tail excess), consistent with the "
         "known weaker stickiness in the symmetric case.{13}"),
    ],

    "Population PK reveals allometric scaling": [
        ("Clinical PK uses population (mixed-effects) modelling to relate "
         "individual PK parameters to patient covariates such as body weight "
         "and renal function.{14} We apply the same methodology to the "
         "three-body problem: the 'covariates' are the mass ratios of the "
         "three bodies."),

        ("Scanning 64 mass configurations (m_2, m_3 in {0.25, 0.5, 0.75, 1, "
         "1.5, 2, 3, 4} with m_1 = 1), we fit a log-linear population model: "
         "log(MRT) = beta_0 + beta_1 times log(mu_12) + beta_2 times "
         "log(mu_out) + beta_3 times log(M) + eta, where mu_12 is the reduced "
         "mass of the initial binary, mu_out is the binary-single reduced mass, "
         "and M is total mass (Fig. 4)."),

        ("The population model yields MRT proportional to mu_12 to the 2.08 "
         "times mu_out to the -1.68 times M to the -0.62, with R-squared = 0.67 "
         "and inter-individual variability omega = 0.35. The positive exponent on "
         "mu_12 indicates that heavier binaries survive longer (deeper potential "
         "wells), while the negative exponent on mu_out indicates that heavier "
         "incoming bodies disrupt the system faster\u2014both physically intuitive "
         "and quantitatively novel."),

        ("Additionally, the escape probability of the lightest body follows a "
         "logistic model in mass fraction, analogous to the dose-response "
         "relationship in pharmacology. This provides a closed-form predictor "
         "for which body will be ejected, a result of practical value for "
         "stellar dynamics simulations."),
    ],

    "Reverse mapping: TMDD stability analysis": [
        ("The structural equivalence works in both directions. Target-mediated "
         "drug disposition (TMDD)\u2014where a drug binds a receptor to form a "
         "complex that is then internalised\u2014is a three-state system with "
         "identical matrix structure to the three-body compartment model.{15,16}"),

        ("We demonstrate that the Jacobian eigenvalue analysis used to assess "
         "Lagrange point stability in celestial mechanics provides a natural "
         "framework for TMDD stability analysis (Fig. 5). The dose-response "
         "bifurcation curve\u2014identifying stable versus unstable steady states "
         "as a function of drug input rate\u2014is the pharmacological analogue "
         "of the Mardling-Aarseth stability criterion for hierarchical "
         "triples.{17}"),

        ("This reverse mapping suggests that the rich stability theory of "
         "N-body systems (KAM tori, invariant manifolds, homoclinic tangles) "
         "could be imported wholesale into nonlinear pharmacokinetics, "
         "particularly for complex biologics with multiple binding sites."),
    ],

    "3D robustness and gravitational-wave dissipation": [
        ("To test whether the PK mapping is an artefact of the 2D planar "
         "geometry or the conservative (energy-preserving) assumption, we "
         "extended the simulations to fully three-dimensional scattering with "
         "randomised orbital planes (uniform on SO(3)), thermal eccentricity "
         "(f(e) = 2e), and isotropic third-body approach directions "
         "(2,000 runs per mass configuration)."),

        ("The three-compartment PK model remains valid in all 3D configurations, "
         "with KS statistics of 0.16\u20130.31 across equal, unequal, and "
         "democratic mass setups\u2014comparable to the 2D values (Fig. 6). "
         "The median lifetimes shift modestly (3D conservative: 246\u2013438 "
         "dynamical times versus 2D: 239\u20131,043), reflecting the larger "
         "phase space available in 3D."),

        ("Adding 2.5 post-Newtonian gravitational-wave radiation reaction{19} "
         "(c = 100 in code units; v_orbital approximately 1) systematically "
         "shortens lifetimes without destroying the compartment structure "
         "(Fig. 7). The lifetime ratio tau_diss/tau_cons ranges from 0.98 "
         "(equal mass) to 0.83 (democratic), with median fractional energy "
         "loss |Delta E/E| = 0.10\u20130.86% per encounter. The democratic "
         "configuration shows the strongest dissipation effect (17% lifetime "
         "reduction), consistent with the longer interaction times providing "
         "more opportunity for gravitational-wave energy extraction."),

        ("These results indicate that the PK mapping is robust across "
         "dimensionality, eccentricity distribution, and the presence of "
         "dissipation\u2014and is therefore applicable to astrophysically "
         "realistic encounters."),
    ],
}

DISCUSSION = [
    ("We have demonstrated a quantitative structural correspondence between two "
     "seemingly unrelated fields: celestial mechanics and pharmacokinetics. "
     "The key insight is that both describe irreversible first-passage "
     "problems in multi-state continuous-time Markov chains, and that the "
     "mathematical structure\u2014not merely the analogy\u2014is shared."),

    ("The practical implications are bidirectional. For celestial mechanics: "
     "(i) the population PK framework provides a principled way to parameterise "
     "outcome statistics as a function of mass ratios, replacing ad hoc fitting "
     "functions;{5-7} (ii) the Michaelis-Menten model offers a minimal "
     "parametrisation of sticky chaos that interpolates between the ergodic "
     "(linear) regime and the trapped (power-law) regime; (iii) PK software "
     "(NONMEM, Monolix) can be directly applied to N-body simulation output."),

    ("For pharmacology: (i) the dynamical-systems stability framework "
     "(Lagrange points, invariant manifolds) provides geometric insight into "
     "TMDD steady states that goes beyond the standard eigenvalue analysis; "
     "(ii) the concept of 'sticky chaos' maps to the clinically important "
     "phenomenon of prolonged drug-receptor residence time;{18} (iii) the "
     "allometric scaling framework transfers predictive power across species "
     "and across mass configurations."),

    ("The robustness of the PK mapping to dimensionality and dissipation "
     "(Figs. 6\u20137) has astrophysical implications. The 17% "
     "lifetime reduction in democratic configurations under gravitational-wave "
     "dissipation may inform merger-rate estimates for "
     "LISA-band sources: three-body encounters in dense clusters will resolve "
     "faster than energy-conserving estimates suggest, potentially increasing the "
     "predicted rate of observable in-spirals. The population PK framework "
     "may provide closed-form event-rate formulae parameterised by cluster "
     "mass function, potentially complementing expensive Monte Carlo surveys."),

    ("Remaining limitations include the perturbative treatment of radiation "
     "reaction (valid only for c >> v_orbital) and the absence of tidal "
     "dissipation relevant for stellar encounters. The linear PK model "
     "assumes ergodicity, which breaks down for very long lifetimes; the "
     "nonlinear extension addresses this but at the cost of additional "
     "parameters."),

    ("More broadly, the correspondence inverts the conventional flow of "
     "methodology between the physical sciences and clinical medicine. "
     "The mathematical theory of pharmacokinetics was distilled from bedside "
     "observations—therapeutic drug monitoring, dose–response studies, "
     "clinical trial design—and aged over six decades of patient-facing "
     "application into a robust quantitative infrastructure. That "
     "infrastructure has been casked in regulatory-grade software (NONMEM, "
     "Monolix) and in saturation-kinetics theory refined through countless "
     "dosing studies (Michaelis–Menten, TMDD). The present work decants this "
     "accumulated clinical knowledge into an entirely new vessel—celestial "
     "mechanics—where, aerated by contact with a different physical context, "
     "it reveals latent capabilities that were invisible within medicine alone. "
     "Several concrete directions follow. "
     "(i) The three-compartment model generalises naturally to N-body "
     "hierarchies: quadruple and higher-order multiples in dense star clusters "
     "map onto N-compartment cascade models whose analytical solutions are "
     "textbook clinical pharmacokinetics.{8,9} "
     "(ii) Population PK covariate modelling—originally designed to account "
     "for patient weight, renal function, and genotype in clinical trials—can "
     "treat mass ratios, angular momenta, and cluster environments as "
     "covariates in nonlinear mixed-effects regressions, enabling systematic "
     "astrophysical parameter scans on platforms built for drug regulation. "
     "(iii) The closed-form MRT scaling (MRT proportional to "
     "mu_12 to the alpha times mu_out to the beta times M to the gamma) can be "
     "convolved with a cluster mass function to yield gravitational-wave "
     "event-rate estimates for LISA-band sources without expensive Monte Carlo "
     "surveys. "
     "(iv) Conversely, the geometric stability theory of celestial mechanics "
     "(KAM tori, invariant manifolds) may inform the rational design of drugs "
     "with controlled target residence times—a key determinant of in vivo "
     "efficacy that is central to modern drug discovery.{18} "
     "(v) Because the mapping relies only on the CTMC structure of "
     "multi-state first-passage processes, it is not limited to gravitational "
     "systems; chemical reaction networks, ecological predator–prey dynamics, "
     "and other multi-state chaotic systems may admit analogous descriptions "
     "rooted in clinical pharmacokinetic theory, suggesting a broader "
     "programme of 'pharmacokinetics of chaos.'"),

    ("In conclusion, the pharmacokinetic compartment model provides both a "
     "conceptual lens and a computational toolkit for the three-body problem "
     "that complements existing statistical-mechanical approaches. By "
     "demonstrating that clinical medicine—through its pharmacokinetic "
     "tradition—can export, not merely import, quantitative frameworks to "
     "the physical sciences, this work illustrates how disciplinary "
     "boundaries may conceal unexploited structural connections."),
]

METHODS = [
    ("N-body simulations", [
        "Three-body scattering simulations were performed using a symplectic "
        "leapfrog (Stormer-Verlet) integrator implemented in Julia 1.10. "
        "Initial conditions: binary with semi-major axis a = 1 and eccentricity "
        "e = 0; third body approaching from r = 10a with velocity at infinity "
        "v_inf = 0.1 (parabolic-like encounters). Impact parameter b sampled "
        "uniformly in b-squared from 0 to b_max = 5a. Integration continued until "
        "escape (body reaching r > 20a with positive energy) or t_max = 10,000 "
        "dynamical times. Adaptive timestep: dt = 0.01 times r_min / v_max. "
        "Energy conservation better than 10 to the -8 in all runs.",

        "For each mass configuration, 5,000 scattering experiments were "
        "performed (total: 15,000 for Phase 1; 5,000 times 64 = 320,000 for "
        "the population PK scan). Binary configuration was identified at each "
        "timestep by finding the most-bound pair (most negative pairwise "
        "energy). A configuration transition was recorded when the most-bound "
        "pair identity changed and persisted for more than 5 timesteps.",
    ]),

    ("PK model estimation", [
        "Transition rates were estimated by maximum-likelihood for a "
        "continuous-time Markov chain: k_ij = N_ij / T_i, where N_ij is the "
        "count of transitions from compartment i to j, and T_i is total "
        "observed dwell time in compartment i. The 3 times 3 rate matrix A was "
        "constructed such that A_ij = k_ji for i not equal to j, and "
        "A_ii = negative sum of (k_ij + k_ei).",

        "Survival function: S(t) = ones transpose times exp(A times t) times P_0. "
        "Eigenvalues of A give the three decay rates; half-lives t_1/2 = "
        "ln(2) / |lambda_i|. Mean residence time (MRT) = negative ones transpose "
        "times A inverse times P_0.",
    ]),

    ("Nonlinear PK model", [
        "The hybrid model uses first-order inter-compartmental transfer and "
        "Michaelis-Menten elimination: dP_i/dt = sum of k_ji times P_j minus "
        "(sum of k_ij) times P_i minus V_max,i times P_i / (K_m + P_i). "
        "The stretched-exponential survival S(t) = sum of w_i times "
        "exp(-lambda_i times t) times t to the -alpha was fitted by "
        "minimising log-space squared error against the empirical survival "
        "function (Nelder-Mead, 10,000 iterations).",
    ]),

    ("Population PK analysis", [
        "For 64 mass configurations (m_1 = 1 fixed; m_2, m_3 in "
        "{0.25, 0.5, 0.75, 1, 1.5, 2, 3, 4}), PK parameters were estimated "
        "independently. The population model: log(MRT_i) = beta_0 + beta_1 "
        "times log(mu_12,i) + beta_2 times log(mu_out,i) + beta_3 times "
        "log(M_i) + eta_i, where eta ~ N(0, omega-squared), was fitted by "
        "ordinary least squares. R-squared and omega (inter-individual "
        "variability) were computed from residuals.",
    ]),

    ("TMDD analysis", [
        "The TMDD ODE system (drug C, receptor R, complex CR) was solved "
        "with typical monoclonal antibody parameters: k_syn = 0.1, k_deg = 0.01, "
        "k_on = 0.1, k_off = 0.001, k_int = 0.05, k_el = 0.01. Steady states "
        "were found numerically (fsolve with multiple initial guesses). "
        "Stability was assessed by Jacobian eigenvalues. Dose-response "
        "bifurcation was computed over dose in [10 to the -3, 10 to the 2].",
    ]),

    ("Phase-space flux comparison", [
        "Theoretical transition rates were predicted from phase-space volume "
        "scaling: k_ij proportional to (mu_to to the 3/2 times m_single to the 3/2) "
        "divided by (mu_from to the 3/2 times m_single_from to the 3/2), following "
        "Stone and Leigh.{6} Escape rates scale as (m_single / M) to the 3/2. "
        "Rates were normalised to the dynamical timescale t_dyn = 1 / sqrt(M).",
    ]),

    ("3D extension and gravitational-wave dissipation", [
        "The 3D extension used fully randomised initial conditions: binary "
        "orbital plane sampled uniformly on SO(3), eccentricity from the "
        "thermal distribution f(e) = 2e (capped at e = 0.95), and third-body "
        "approach direction uniform on S-squared. For each of 3 mass "
        "configurations, 2,000 scattering experiments were performed.",

        "Gravitational-wave dissipation was implemented via the 2.5PN "
        "radiation reaction acceleration (leading-order Burke-Thorne term).{19} "
        "For each pair (i,j) with relative separation r and velocity v, the "
        "dissipative acceleration is: a_2.5PN = (8/5) times mu times M-squared "
        "divided by (c-to-the-5 times r-cubed) times "
        "[{3v-squared + (17/3)M/r} times v_r times n-hat minus "
        "{v-squared + 3M/r} times v]. We set c = 100 in code units "
        "(v_orbital approximately 1), giving |Delta E/E| approximately 0.1-0.9% "
        "per encounter\u2014consistent with the weak-field limit. Integration used "
        "velocity Verlet with the PN correction evaluated at the half-step "
        "velocity.",
    ]),
]

REFERENCES = [
    "Newton, I. Philosophiae Naturalis Principia Mathematica (1687).",
    "Poincar\u00e9, H. Les m\u00e9thodes nouvelles de la m\u00e9canique c\u00e9leste (Gauthier-Villars, 1892).",
    "Boekholt, T. & Portegies Zwart, S. On the reliability of N-body simulations. Comput. Astrophys. Cosmol. 2, 2 (2015).",
    "Heggie, D. C. Binary evolution in stellar dynamics. Mon. Not. R. Astron. Soc. 173, 729\u2013787 (1975).",
    "Ginat, Y. B. & Perets, H. B. Analytical, statistical approximate solution of dissipative and nondissipative binary-single stellar encounters. Nature 593, 395\u2013398 (2021).",
    "Stone, N. C. & Leigh, N. W. C. A statistical solution to the chaotic, non-hierarchical three-body problem. Nature 576, 406\u2013410 (2019).",
    "Monaghan, J. J. A statistical theory of the disruption of three-body systems\u2014II. Mon. Not. R. Astron. Soc. 177, 583\u2013594 (1976).",
    "Rowland, M. & Tozer, T. N. Clinical Pharmacokinetics and Pharmacodynamics: Concepts and Applications 4th edn (Lippincott Williams & Wilkins, 2011).",
    "Gabrielsson, J. & Weiner, D. Pharmacokinetic and Pharmacodynamic Data Analysis 5th edn (Swedish Pharmaceutical Press, 2016).",
    "Norris, J. R. Markov Chains (Cambridge Univ. Press, 1997).",
    "Mikkola, S. & Tanikawa, K. Explicit symplectic algorithms for the three-body problem. Celest. Mech. Dyn. Astron. 74, 287\u2013295 (1999).",
    "Manwadkar, V., Kol, B., Trani, A. A. & Leigh, N. W. C. Chaos and Levy flights in the three-body problem. arXiv:2011.01850 (2020).",
    "Samsing, J. & Ilan, T. Topology counts: force distributions in circular three-body systems. Nature 603, 237\u2013240 (2022).",
    "Mould, D. R. & Upton, R. N. Basic concepts in population modeling, simulation, and model-based drug development. CPT Pharmacometrics Syst. Pharmacol. 1, e6 (2012).",
    "Mager, D. E. & Jusko, W. J. General pharmacokinetic model for drugs exhibiting target-mediated drug disposition. J. Pharmacokinet. Pharmacodyn. 28, 507\u2013532 (2001).",
    "Gibiansky, L. & Gibiansky, E. Target-mediated drug disposition model: approximations, identifiability of model parameters and applications to the population pharmacokinetic-pharmacodynamic modeling of biologics. Expert Opin. Drug Metab. Toxicol. 5, 803\u2013812 (2009).",
    "Mardling, R. A. & Aarseth, S. J. Tidal interactions in star cluster simulations. Mon. Not. R. Astron. Soc. 321, 398\u2013420 (2001).",
    "Copeland, R. A. The drug\u2013target residence time model: a 10-year retrospective. Nat. Rev. Drug Discov. 15, 87\u201395 (2016).",
    "Peters, P. C. Gravitational radiation and the motion of two point masses. Phys. Rev. 136, B1224\u2013B1232 (1964).",
]


# ---------------------------------------------------------------------------
# Build manuscript docx
# ---------------------------------------------------------------------------

def build_manuscript():
    doc = Document()

    # --- Style setup ---
    style = doc.styles['Normal']
    style.font.name = 'Times New Roman'
    style.font.size = Pt(11)
    style.paragraph_format.line_spacing = 1.5

    # --- Title ---
    title_p = doc.add_paragraph()
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_p.paragraph_format.space_after = Pt(6)
    run = title_p.add_run(TITLE)
    run.bold = True
    run.font.size = Pt(16)

    # --- Authors ---
    auth_p = doc.add_paragraph()
    auth_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = auth_p.add_run(AUTHORS)
    run.font.size = Pt(12)

    # --- Affiliations ---
    aff_p = doc.add_paragraph()
    aff_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    aff_p.paragraph_format.space_after = Pt(18)
    run = aff_p.add_run(AFFILIATIONS)
    run.font.size = Pt(10)
    run.italic = True

    # --- Abstract ---
    doc.add_paragraph()
    abs_h = doc.add_paragraph()
    run = abs_h.add_run("Abstract")
    run.bold = True
    run.font.size = Pt(12)

    abs_p = doc.add_paragraph()
    abs_p.paragraph_format.space_after = Pt(12)
    run = abs_p.add_run(ABSTRACT)
    run.font.size = Pt(11)

    # --- Introduction ---
    add_heading(doc, "Introduction", level=1)
    for para_text in INTRO:
        add_paragraph_with_refs(doc, para_text)

    # --- Figure 1: Conceptual mapping (after introduction) ---
    fig1_path = os.path.join(FIGDIR, "fig1_conceptual.png")
    if not os.path.exists(fig1_path):
        fig1_path = os.path.join(FIGDIR_ORIG, "fig1_conceptual_mapping.png")
    add_figure(doc, fig1_path, (
        "Mapping three-body scattering onto pharmacokinetic compartments. "
        "(a) Three-body scattering: each binary configuration (which body pair "
        "is gravitationally bound) corresponds to a compartment; transitions "
        "between configurations are inter-compartmental transfers at rates k_ij; "
        "escape of one body is irreversible elimination at rate k_ei. "
        "(b) Equivalent three-compartment pharmacokinetic model."
    ), 1)

    # --- Results ---
    add_heading(doc, "Results", level=1)

    fig_counter = 2
    for section_title, paragraphs in RESULTS_SECTIONS.items():
        add_heading(doc, section_title, level=2)
        for para_text in paragraphs:
            add_paragraph_with_refs(doc, para_text)

        # Insert relevant figures
        if "Linear PK" in section_title:
            fig2_path = os.path.join(FIGDIR, "fig2_survival.png")
            if not os.path.exists(fig2_path):
                fig2_path = os.path.join(FIGDIR_ORIG, "fig2_lifetime_Equal_mass_(1:1:1).png")
            add_figure(doc, fig2_path, (
                "Survival curves for three mass configurations. "
                "Blue dots: empirical survival from 5,000 N-body simulations. "
                "Red line: linear 3-compartment PK model prediction (sum of "
                "three exponentials). The multi-exponential structure is "
                "statistically necessary (KS test p < 0.001 versus single "
                "exponential). (a) Equal mass (1:1:1). (b) Unequal mass "
                "(1:2:0.5). (c) Democratic initial conditions."
            ), fig_counter)
            fig_counter += 1

        elif "Nonlinear" in section_title:
            fig3_path = os.path.join(FIGDIR, "fig3_nonlinear.png")
            if not os.path.exists(fig3_path):
                fig3_path = os.path.join(FIGDIR_ORIG, "fig6_nonlinear_comparison.png")
            add_figure(doc, fig3_path, (
                "Linear versus nonlinear PK model comparison. Log-log survival "
                "plots for three mass configurations. Blue dots: N-body "
                "simulation data. Red: linear PK (multi-exponential). Green "
                "dashed: hybrid Michaelis-Menten model capturing the power-law "
                "tail from sticky chaos. (a) Equal mass: linear model adequate. "
                "(b) Unequal mass: 42% tail RMSE improvement. "
                "(c) Democratic: 36% improvement."
            ), fig_counter)
            fig_counter += 1

        elif "Population" in section_title:
            fig4_path = os.path.join(FIGDIR, "fig4_population_pk.png")
            if not os.path.exists(fig4_path):
                fig4_path = os.path.join(FIGDIR_ORIG, "fig7_population_pk.png")
            add_figure(doc, fig4_path, (
                "Population PK analysis of mass-ratio dependence. "
                "(a) MRT versus reduced mass of binary-single system. "
                "(b) Median lifetime versus incoming body mass. "
                "(c) Mean excursions versus lightest mass fraction. "
                "(d) Lightest-body escape probability with logistic fit. "
                "(e) Population model predicted versus observed MRT "
                "(R-squared = 0.67). "
                "(f) Heatmap of slowest half-life as function of m_2 and m_3."
            ), fig_counter)
            fig_counter += 1

        elif "TMDD" in section_title:
            fig5_path = os.path.join(FIGDIR, "fig5_tmdd.png")
            if not os.path.exists(fig5_path):
                fig5_path = os.path.join(FIGDIR_ORIG, "fig8_tmdd_analogy.png")
            add_figure(doc, fig5_path, (
                "TMDD reverse application. (a) Dose-response bifurcation curve "
                "showing stable (blue) and unstable (red) steady states. "
                "(b) Receptor occupancy versus dose rate. "
                "(c) Slowest Jacobian eigenvalue versus dose, identifying the "
                "stability boundary (analogue of Lagrange point stability in "
                "celestial mechanics)."
            ), fig_counter)
            fig_counter += 1

        elif "3D robustness" in section_title:
            fig6_path = os.path.join(FIGDIR, "fig6_3d_dissipative_comparison.png")
            if os.path.exists(fig6_path):
                add_figure(doc, fig6_path, (
                    "PK model validity across dimensionality and dissipation. "
                    "Rows: 2D conservative, 3D conservative, 3D dissipative "
                    "(c = 100). Columns: equal mass (1:1:1), unequal mass "
                    "(1:2:0.5), democratic. Blue dots: N-body survival data. "
                    "Red: linear 3-compartment PK fit. KS statistics annotated."
                ), fig_counter)
                fig_counter += 1

            fig7_path = os.path.join(FIGDIR, "fig7_dissipation_effect.png")
            if os.path.exists(fig7_path):
                add_figure(doc, fig7_path, (
                    "Effect of gravitational-wave dissipation on lifetime "
                    "distributions. Blue: 3D conservative. Red: 3D with 2.5PN "
                    "radiation reaction (c = 100). Dashed: PK model fits. "
                    "Dissipation systematically shortens lifetimes "
                    "(tau_diss/tau_cons = 0.83-0.98) without destroying the "
                    "PK model structure."
                ), fig_counter)
                fig_counter += 1

    # --- Discussion ---
    add_heading(doc, "Discussion", level=1)
    for idx, para_text in enumerate(DISCUSSION):
        add_paragraph_with_refs(doc, para_text)

    # --- Methods ---
    doc.add_page_break()
    add_heading(doc, "Methods", level=1)
    for method_title, method_paras in METHODS:
        add_heading(doc, method_title, level=2)
        for para_text in method_paras:
            add_paragraph_with_refs(doc, para_text)

    # --- References ---
    doc.add_page_break()
    add_heading(doc, "References", level=1)
    for i, ref in enumerate(REFERENCES, 1):
        p = doc.add_paragraph()
        run_num = p.add_run(f"{i}. ")
        run_num.font.size = Pt(10)
        run_text = p.add_run(ref)
        run_text.font.size = Pt(10)

    # --- Save ---
    out_path = os.path.join(OUTDIR, "manuscript_nature.docx")
    doc.save(out_path)
    print(f"Saved: {out_path}")
    return out_path


# ---------------------------------------------------------------------------
# Build editable PPTX with figures
# ---------------------------------------------------------------------------

def build_figures_pptx():
    from pptx import Presentation
    from pptx.util import Inches as PInches, Pt as PPt
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    # Widescreen dimensions
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)

    figures = [
        ("fig1_conceptual.png", "Figure 1",
         "Three-body scattering mapped onto PK compartments"),
        ("fig2_survival.png", "Figure 2",
         "Survival curves: linear PK model fits N-body data"),
        ("fig3_nonlinear.png", "Figure 3",
         "Nonlinear PK captures sticky chaos tail"),
        ("fig4_population_pk.png", "Figure 4",
         "Population PK: mass-ratio dependence (6 panels)"),
        ("fig5_tmdd.png", "Figure 5",
         "TMDD reverse application: dose-response bifurcation"),
    ]

    blank_layout = prs.slide_layouts[6]  # blank slide

    for fname, fig_label, caption in figures:
        fpath = os.path.join(FIGDIR_NATURE, fname)
        if not os.path.exists(fpath):
            fpath = os.path.join(FIGDIR_ORIG, fname)
        if not os.path.exists(fpath):
            continue

        slide = prs.slides.add_slide(blank_layout)

        # Title at top
        from pptx.util import Emu
        txBox = slide.shapes.add_textbox(
            PInches(0.5), PInches(0.2), PInches(12), PInches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = fig_label
        p.font.size = PPt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Image centered
        img_left = PInches(1.0)
        img_top = PInches(1.0)
        img_width = PInches(11.0)
        try:
            slide.shapes.add_picture(fpath, img_left, img_top, width=img_width)
        except Exception:
            pass

        # Caption at bottom
        cap_box = slide.shapes.add_textbox(
            PInches(0.5), PInches(6.5), PInches(12.333), PInches(0.8))
        cap_tf = cap_box.text_frame
        cap_tf.word_wrap = True
        cap_p = cap_tf.paragraphs[0]
        cap_p.text = caption
        cap_p.font.size = PPt(14)
        cap_p.alignment = PP_ALIGN.CENTER

    out_path = os.path.join(OUTDIR, "figures_editable.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")
    return out_path


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    print("=" * 60)
    print("  Generating Nature Manuscript")
    print("=" * 60)

    docx_path = build_manuscript()
    pptx_path = build_figures_pptx()

    print(f"\nOutputs:")
    print(f"  Manuscript: {docx_path}")
    print(f"  Figures PPTX: {pptx_path}")
    print("=" * 60)
