#!/usr/bin/env python3
"""
Generate the MNRAS (Monthly Notices of the Royal Astronomical Society)
submission package:

  - manuscript_mnras.docx   (author-year citations, numbered sections,
                             figures embedded inline after first mention)
  - cover_letter_mnras.docx
  - figures_mnras_editable.pptx  (one figure per slide, English, editable)
  - individual figure PNGs (copied into manuscript/figures_mnras/)
  - threebody_pk_mnras_submission.zip

MNRAS formatting choices (Instructions to Authors):
  * Author-year (Harvard) citation style, reference list alphabetical.
  * Numbered sections (1, 2, 3, ...) and subsections (3.1, ...).
  * No word limit; sectioned IMRaD + Conclusions.
  * Figures embedded inline in the manuscript (also supplied separately).

Story framing (per author): inspired by PK/PD, we re-examine the three-body
problem; the PK toolkit turns out to be a genuinely useful set of tools for
three-body statistics; the two-way traffic of methods is a productive
synergy. The reverse-mapping (TMDD) material is deliberately condensed.
"""

from __future__ import annotations

import json
import os
import re
import shutil
import zipfile
from datetime import date

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

from create_prx_submission import (
    _add_run, _render_math_segment, setup_style,
)

BASEDIR = os.path.dirname(os.path.abspath(__file__))
FIGDIR = os.path.join(BASEDIR, "figures_nature")
OUTDIR = os.path.join(BASEDIR, "manuscript")
FIGOUT = os.path.join(OUTDIR, "figures_mnras")
os.makedirs(OUTDIR, exist_ok=True)
os.makedirs(FIGOUT, exist_ok=True)

SUMMARY_PATH = os.path.join(BASEDIR, "data", "mnras_extensions_summary.json")
SUMMARY = {}
if os.path.exists(SUMMARY_PATH):
    with open(SUMMARY_PATH) as f:
        SUMMARY = json.load(f)


# ---------------------------------------------------------------------------
# References — author-year (Harvard) style
# ---------------------------------------------------------------------------
# key -> (in-text "Author Year", full reference string)
REFS = {
    1:  ("Newton 1687",
         "Newton I., 1687, Philosophiae Naturalis Principia Mathematica. "
         "Joseph Streater, London"),
    2:  ("Poincar\u00e9 1892",
         "Poincar\u00e9 H., 1892, Les M\u00e9thodes Nouvelles de la M\u00e9canique "
         "C\u00e9leste. Gauthier-Villars, Paris"),
    3:  ("Boekholt & Portegies Zwart 2015",
         "Boekholt T., Portegies Zwart S., 2015, Comput. Astrophys. Cosmol., 2, 2"),
    4:  ("Heggie 1975",
         "Heggie D. C., 1975, MNRAS, 173, 729"),
    5:  ("Ginat & Perets 2021",
         "Ginat Y. B., Perets H. B., 2021, Phys. Rev. X, 11, 031020"),
    6:  ("Stone & Leigh 2019",
         "Stone N. C., Leigh N. W. C., 2019, Nature, 576, 406"),
    7:  ("Monaghan 1976",
         "Monaghan J. J., 1976, MNRAS, 177, 583"),
    8:  ("Rowland & Tozer 2011",
         "Rowland M., Tozer T. N., 2011, Clinical Pharmacokinetics and "
         "Pharmacodynamics, 4th edn. Lippincott Williams & Wilkins, Philadelphia"),
    9:  ("Gabrielsson & Weiner 2016",
         "Gabrielsson J., Weiner D., 2016, Pharmacokinetic and Pharmacodynamic "
         "Data Analysis, 5th edn. Swedish Pharmaceutical Press, Stockholm"),
    10: ("Norris 1997",
         "Norris J. R., 1997, Markov Chains. Cambridge Univ. Press, Cambridge"),
    11: ("Mikkola & Tanikawa 1999",
         "Mikkola S., Tanikawa K., 1999, Celest. Mech. Dyn. Astron., 74, 287"),
    12: ("Manwadkar et al. 2020",
         "Manwadkar V., Kol B., Trani A. A., Leigh N. W. C., 2020, MNRAS, 497, 3694"),
    13: ("Samsing et al. 2022",
         "Samsing J., et al., 2022, Nature, 603, 237"),
    14: ("Mould & Upton 2012",
         "Mould D. R., Upton R. N., 2012, CPT: Pharmacometrics Syst. Pharmacol., 1, e6"),
    15: ("Mager & Jusko 2001",
         "Mager D. E., Jusko W. J., 2001, J. Pharmacokinet. Pharmacodyn., 28, 507"),
    16: ("Gibiansky & Gibiansky 2009",
         "Gibiansky L., Gibiansky E., 2009, Expert Opin. Drug Metab. Toxicol., 5, 803"),
    17: ("Mardling & Aarseth 2001",
         "Mardling R. A., Aarseth S. J., 2001, MNRAS, 321, 398"),
    18: ("Copeland 2016",
         "Copeland R. A., 2016, Nat. Rev. Drug Discov., 15, 87"),
    19: ("Peters 1964",
         "Peters P. C., 1964, Phys. Rev., 136, B1224"),
    # --- new references for the MNRAS extensions ---
    20: ("Hut & Bahcall 1983",
         "Hut P., Bahcall J. N., 1983, ApJ, 268, 319"),
    21: ("Valtonen & Karttunen 2006",
         "Valtonen M., Karttunen H., 2006, The Three-Body Problem. "
         "Cambridge Univ. Press, Cambridge"),
    22: ("Blanchet 2014",
         "Blanchet L., 2014, Living Rev. Relativ., 17, 2"),
    23: ("Kroupa 2001",
         "Kroupa P., 2001, MNRAS, 322, 231"),
    24: ("Rodriguez et al. 2016",
         "Rodriguez C. L., Chatterjee S., Rasio F. A., 2016, Phys. Rev. D, 93, 084029"),
    25: ("Antonini & Rasio 2016",
         "Antonini F., Rasio F. A., 2016, ApJ, 831, 187"),
    26: ("Samsing et al. 2014",
         "Samsing J., MacLeod M., Ramirez-Ruiz E., 2014, ApJ, 784, 71"),
    27: ("Amaro-Seoane et al. 2017",
         "Amaro-Seoane P., et al., 2017, preprint (arXiv:1702.00786)"),
}


def _intext(nums, paren=True):
    """Render a citation group given a list of reference numbers."""
    cites = [REFS[n][0] for n in nums]
    inner = "; ".join(cites)
    return f"({inner})" if paren else inner


def _expand_marker(marker):
    """'{5}' -> [5]; '{5,6}' -> [5,6]; '{5\u20137}' -> [5,6,7]."""
    body = marker[1:-1]
    nums = []
    for chunk in re.split(r"[,\s]+", body):
        chunk = chunk.strip()
        if not chunk:
            continue
        m = re.match(r"^(\d+)\s*[\u2013\u2014\u2212–-]\s*(\d+)$", chunk)
        if m:
            nums.extend(range(int(m.group(1)), int(m.group(2)) + 1))
        elif chunk.isdigit():
            nums.append(int(chunk))
    return nums


def add_text_with_refs(paragraph, text, font_size=Pt(11)):
    """Render text: convert {N} numeric markers to author-year citations,
    and render _{}/^{} math notation."""
    parts = re.split(r'(?<![_^])(\{[\d,\u2013\u2014\u2212 –-]+\})', text)
    for part in parts:
        if not part:
            continue
        if (part.startswith('{') and part.endswith('}')
                and re.match(r'^\{[\d,\u2013\u2014\u2212 –-]+\}$', part)):
            nums = _expand_marker(part)
            if nums:
                run = paragraph.add_run(" " + _intext(nums))
                run.font.size = font_size
        else:
            _render_math_segment(paragraph, part, font_size=font_size)


def add_heading(doc, text, level=1):
    return doc.add_heading(text, level=level)


def add_paragraph_with_refs(doc, text):
    p = doc.add_paragraph()
    add_text_with_refs(p, text)
    return p


def add_figure(doc, fig_file, caption, label):
    """Embed a figure inline with an MNRAS-style caption below it."""
    path = os.path.join(FIGDIR, fig_file)
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(14)
    p.paragraph_format.space_after = Pt(4)
    if os.path.exists(path):
        run = p.add_run()
        run.add_picture(path, width=Inches(5.6))
    else:
        run = p.add_run(f"[{label}]")
        run.bold = True
    cap = doc.add_paragraph()
    cap.paragraph_format.space_before = Pt(4)
    cap.paragraph_format.space_after = Pt(14)
    rb = cap.add_run(f"{label}. ")
    rb.bold = True
    rb.font.size = Pt(9)
    add_text_with_refs(cap, caption, font_size=Pt(9))


# ---------------------------------------------------------------------------
# Convenience accessors into the extension-analysis summary
# ---------------------------------------------------------------------------

def g(*path, default=None):
    d = SUMMARY
    for k in path:
        if isinstance(d, dict) and k in d:
            d = d[k]
        else:
            return default
    return d


def _fmt(x, spec=".2f", default="\u2014"):
    return format(x, spec) if isinstance(x, (int, float)) else default


# ---------------------------------------------------------------------------
# Front matter
# ---------------------------------------------------------------------------

TITLE = ("A pharmacokinetic toolkit for the statistical three-body problem: "
         "compartment models, allometric scaling, and astrophysical "
         "predictions")

AUTHORS = "Tatsuki Onishi"

AFFILIATIONS = ("Faculty of Data Science, Shiga University, "
                "1-1-1 Banba, Hikone, Shiga 522-8522, Japan")

CORRESPONDING = "bougtoir@gmail.com"

ABSTRACT = (
    "The statistical outcome of chaotic three-body scattering is regular "
    "enough to be described by rate equations, yet the community lacks a "
    "compact, transferable toolkit for turning those regularities into "
    "quantitative astrophysical predictions. Motivated by pharmacokinetics "
    "(PK) \u2014 the clinical discipline modelling how drugs distribute and "
    "clear \u2014 we re-examine the gravitational three-body "
    "problem through a compartmental lens: each bound binary configuration is "
    "a compartment, configuration changes are inter-compartmental transfers, "
    "and escape is elimination. Using 15,000 planar scattering experiments "
    "across 64 mass configurations, plus 27,000 new three-dimensional "
    "encounters spanning Newtonian, 1PN-precessing, and tidally dissipative "
    "dynamics, we show that the PK toolkit is a practical set of "
    "instruments for three-body statistics. A linear three-compartment model "
    "reproduces the multi-exponential lifetime distribution and yields the "
    "mean residence time (MRT) in closed form; a nonlinear "
    "(Michaelis\u2013Menten) elimination term captures the sticky-chaos tail; "
    "and population-PK (mixed-effects) modelling delivers an allometric "
    "scaling law for the interaction time. We then make astrophysical "
    "predictions: (i) Peters-inspiral merger "
    "fractions and the eccentricity boost of the gravitational-wave (GW) "
    "merger rate; (ii) a closed-form, dose\u2013response predictor for which "
    "body is ejected, convolved with a realistic black-hole mass function; "
    "(iii) eccentricity and ejection-velocity distributions of the surviving "
    "binaries; and (iv) resonance-versus-encounter timescales in globular "
    "clusters, nuclear star clusters, and AGN discs. The compartment "
    "structure survives 1PN precession and tidal drag. The mapping is bidirectional \u2014 celestial-mechanics "
    "stability theory also informs target-mediated drug disposition \u2014 "
    "illustrating that this two-way traffic of methods is a productive "
    "synergy."
)

INTRO = [
    ("The three-body problem \u2014 predicting the motion of three "
     "gravitationally interacting masses \u2014 has resisted closed-form "
     "solution since Newton.{1} Poincar\u00e9 proved its non-integrability,{2} "
     "and modern numerical work confirms that generic encounters are chaotic, "
     "with outcomes sensitive to initial conditions at the level of "
     "floating-point precision.{3} Yet the statistical properties of the "
     "outcomes are strikingly regular: lifetimes follow a multi-exponential "
     "decay,{4,7} and final-state distributions can be predicted from "
     "phase-space volume arguments.{6,20}"),

    ("This regularity has been placed on firm theoretical footing. Ginat & "
     "Perets decomposed chaotic scattering into a sequence of independent "
     "excursions and solved it as a random walk in binary binding energy,{5} "
     "while Stone & Leigh built a flux-based statistical mechanics that "
     "predicts three-body outcomes.{6} These frameworks establish that the "
     "ergodic core of the problem is amenable to stochastic modelling. What "
     "the field still lacks is a compact, portable set of estimators that "
     "turn simulated encounter statistics directly into the quantities "
     "astrophysicists want \u2014 mean interaction times, ejection "
     "probabilities, and event rates \u2014 without bespoke fitting for every "
     "new mass spectrum or environment."),

    ("Clinical pharmacokinetics (PK) is precisely such a toolkit, refined "
     "over six decades of patient-facing use.{8,9} PK describes how a drug "
     "distributes among body compartments (plasma, tissues) at first-order "
     "rates and is eliminated irreversibly; the underlying mathematics is a "
     "continuous-time Markov chain whose generator is the PK rate "
     "matrix.{10} Around this core, pharmacology has built mean-residence-"
     "time calculus, saturable (Michaelis\u2013Menten) kinetics for "
     "capacity-limited clearance, dose\u2013response models, and "
     "population (mixed-effects) regression that ties individual rate "
     "constants to covariates such as body weight and organ function.{14}"),

    ("Here, motivated by this analogy, we re-examine three-body scattering as "
     "a three-compartment PK system: each binary configuration is a "
     "compartment, configuration transitions are transfers, and escape is "
     "elimination. This is not merely analogical \u2014 it is a structural "
     "equivalence at the level of the master equation. Our aim is deliberately "
     "practical: we ask whether the PK toolkit, imported wholesale, is "
     "actually useful for three-body statistics, and whether it yields new, "
     "quantitative astrophysical predictions. We find that it does, on both "
     "counts, and that the traffic runs both ways \u2014 celestial-mechanics "
     "stability theory in turn illuminates target-mediated drug disposition. "
     "The remainder of the paper sets out the methods (Section 2), establishes "
     "the compartment mapping and its nonlinear and population extensions "
     "(Sections 3.1\u20133.3), condenses the reverse mapping (Section 3.4), and "
     "then applies the toolkit to make astrophysical predictions and "
     "robustness tests (Sections 3.5\u20133.11)."),
]


METHODS = [
    ("N-body simulations", [
        "Planar three-body scattering was integrated with a symplectic "
        "leapfrog (St\u00f6rmer\u2013Verlet) scheme in code units (G = 1). "
        "Initial conditions: a binary with semi-major axis a = 1 and a third "
        "body approaching from r = 10a; impact parameter sampled uniformly in "
        "b^2 up to b_{max} = 4a; velocity at infinity v_{\\infty} = 0.8. "
        "Integration continued until escape (a body reaching r > 100a with "
        "positive energy) or t_{max} = 10^5 dynamical times, with adaptive "
        "timestep and energy conservation better than 10^{\u22128}. For each of "
        "three canonical mass configurations \u2014 equal (1:1:1), unequal "
        "(1:2:0.5), and democratic \u2014 5,000 experiments were run "
        "(15,000 total), plus 5,000 \u00d7 64 = 320,000 for the population "
        "scan. The bound pair was identified at each step as the most-bound "
        "(most negative pairwise energy); a transition was recorded when the "
        "most-bound pair changed and persisted for more than five steps.",

        "Three-dimensional encounters were generated with a Numba-compiled "
        "velocity-Verlet integrator (validated against the planar Julia code: "
        "median lifetimes agree to within 5 per cent). Orbital planes were "
        "sampled uniformly on SO(3), eccentricities from the thermal "
        "distribution f(e) = 2e (capped at 0.95), and approach directions "
        "isotropically. Three dynamical regimes were run at 3,000 experiments "
        "per mass configuration (27,000 total): Newtonian, 1PN "
        "(pairwise post-Newtonian precession), and tidal (a phenomenological "
        "velocity-dependent drag mimicking equilibrium-tide dissipation). "
        "For every escape we recorded the final-binary semi-major axis and "
        "eccentricity, and the asymptotic ejection speed.",
    ]),

    ("PK model estimation and closed-form residence times", [
        "Transfer rates were estimated by maximum likelihood for a "
        "continuous-time Markov chain, k_{ij} = N_{ij} / T_i, with N_{ij} the "
        "count of i\u2192j transitions and T_i the total dwell time in "
        "compartment i. The 3\u00d73 generator A has off-diagonal A_{ij} = "
        "k_{ji} and diagonal A_{ii} = \u2212\u03a3(k_{ij} + k_{ei}). The "
        "survival function S(t) = 1^T exp(At) P_0 is a sum of three "
        "exponentials; half-lives are t_{1/2} = ln 2 / |\u03bb_i|, and the "
        "mean residence time is MRT = \u22121^T A^{\u22121} P_0. The MRT is "
        "the single closed-form scalar we carry through to the astrophysical "
        "applications.",
    ]),

    ("Nonlinear and population PK", [
        "Sticky chaos was modelled by replacing first-order elimination with "
        "Michaelis\u2013Menten kinetics, V_{max} P / (K_m + P), yielding a "
        "survival tail t^{\u2212\u03b1} exp(\u2212\u03bb_{slow} t) fitted in "
        "log space (Nelder\u2013Mead). The population model regressed "
        "log(MRT) on the log reduced masses: log(MRT) = \u03b2_0 + "
        "\u03b2_1 log(\u03bc_{12}) + \u03b2_2 log(\u03bc_{out}) + "
        "\u03b2_3 log(M) + \u03b7, with \u03b7 ~ N(0,\u03c9^2), by ordinary "
        "least squares over the 64-configuration grid. The escape probability "
        "of the lightest body was fitted as a logistic function of the "
        "minimum mass ratio \u2014 a dose\u2013response relationship in PK "
        "language.",
    ]),

    ("Astrophysical post-processing", [
        "Merger times were computed from the Peters (1964) formula, "
        "t_{GW} \u221d a^4 (1\u2212e^2)^{7/2} / (m_1 m_2 M),{19} using the "
        "measured final-binary a and e, scaled to a fiducial black-hole "
        "binary (10 + 10 M_\u2609). A binary was counted as a merger if "
        "t_{GW} was shorter than a Hubble time. We use the orbit-averaged "
        "circularisation-time expression, which is exact in the low-e limit "
        "and accurate to within a factor of order unity even at high "
        "eccentricity; the highest-e escapes (e > 0.9), which account for "
        "6\u201324 per cent of encounters depending on configuration, are "
        "therefore treated conservatively and dominate the merger fraction "
        "regardless of this approximation. The eccentricity boost is the "
        "ratio of merger fractions for the measured (thermal) e-distribution "
        "versus circular orbits. Which-body-ejected predictions applied the "
        "logistic escape model to a Kroupa-sampled black-hole mass function "
        "(5\u201350 M_\u2609).{23} Environmental timescales combined the "
        "PK-model MRT (in dynamical times) with a gravitational-focusing "
        "encounter rate, n \u03c3 v \u03a3, for a globular-cluster core, a "
        "nuclear star cluster, and an AGN-disc migration trap.{24,25}",

        "The 1PN acceleration follows the standard pairwise post-Newtonian "
        "expansion;{22} the tidal term is a weak velocity-dependent drag "
        "applied inside a fixed pericentre scale. Robustness was quantified "
        "by (i) the two-sample Kolmogorov\u2013Smirnov (KS) distance between "
        "the perturbed and Newtonian lifetime distributions and (ii) the KS "
        "distance between the perturbed data and its own three-compartment PK "
        "fit.",
    ]),
]


# ---------------------------------------------------------------------------
# Results (built at call time so the latest summary JSON is used)
# ---------------------------------------------------------------------------

def get_results_blocks():
    # numbers pulled from the extension-analysis summary
    a2_ks = g("A2_flux_benchmark", "ks_pk_equal")
    a2_r2 = g("A2_flux_benchmark", "mrt_vs_lifetime_r2")
    a2_n = g("A2_flux_benchmark", "n_configs", default=64)

    a3_frac = g("A3_ejection_mf", "frac_lightest_ejected_powerlaw")
    a3_mej = g("A3_ejection_mf", "mean_ejected_mass")
    a3_mmem = g("A3_ejection_mf", "mean_member_mass")

    a1_th = g("A1_gw_merger", "merger_fraction_a0_0.5AU_thermal")
    a1_ci = g("A1_gw_merger", "merger_fraction_a0_0.5AU_circular")
    a1_boost = g("A1_gw_merger", "ecc_merger_boost")
    a1_ratio = g("A1_gw_merger", "dissipative_lifetime_ratio", default={})

    a4 = g("A4_ecc_kick", default={}) or {}
    b56 = g("B56_pn_tidal", default={}) or {}
    envs = g("B7_environments", "environments", default=[]) or []

    def ratio_range(d):
        vals = [v for v in (d or {}).values() if isinstance(v, (int, float))]
        return (min(vals), max(vals)) if vals else (None, None)

    lo, hi = ratio_range(a1_ratio)

    blocks = []

    blocks.append(("Linear PK model reproduces multi-exponential decay", [
        ("We first establish the mapping. Each snapshot of a three-body "
         "encounter has one bound pair and one distant body; labelling the "
         "bound pair defines three compartments, and the reshuffling of which "
         "pair is bound defines the transfers, with permanent escape as "
         "elimination (Fig. 1). Estimating the nine rate constants by maximum "
         "likelihood gives a generator whose survival function is a sum of "
         "three exponentials."),
        ("This multi-exponential prediction matches the empirical survival "
         "curves across all three mass configurations (Fig. 2); a single "
         "exponential is rejected by the Kolmogorov\u2013Smirnov (KS) test "
         "(p < 0.001), confirming that three compartments are necessary and "
         "sufficient for the intermediate-time dynamics. For equal masses all "
         "transfer rates converge by symmetry (k_{ij} \u2248 0.015 per "
         "dynamical time), consistent with the ergodic hypothesis underlying "
         "flux-based theory.{6}"),
    ], [
        ("fig1_conceptual.png",
         "Mapping three-body scattering onto pharmacokinetic compartments. "
         "(a) Each bound pair is a compartment; configuration changes are "
         "inter-compartmental transfers at rates k_{ij}; escape is "
         "irreversible elimination at rate k_{ei}. (b) The equivalent "
         "three-compartment PK model.", "Figure 1"),
        ("fig2_survival.png",
         "Survival curves for the three mass configurations. Blue: empirical "
         "survival from 5,000 N-body experiments. Red: linear three-"
         "compartment PK prediction. The multi-exponential form is "
         "statistically required (KS test p < 0.001 versus a single "
         "exponential). (a) Equal mass, (b) unequal mass, (c) democratic.",
         "Figure 2"),
    ]))

    blocks.append(("Nonlinear PK captures sticky chaos", [
        ("The linear model underestimates the long-lived tail \u2014 the "
         "signature of sticky chaos, where trajectories linger near periodic "
         "orbits.{11,12} In PK language this is saturable, capacity-limited "
         "elimination. Replacing first-order clearance with a "
         "Michaelis\u2013Menten term yields a tail "
         "t^{\u2212\u03b1} exp(\u2212\u03bb_{slow} t) that captures the "
         "power-law correction (Fig. 3), improving the tail RMSE by "
         "36\u201342 per cent for the unequal and democratic configurations; "
         "the symmetric equal-mass case needs no correction, consistent with "
         "its known weaker stickiness."),
    ], [
        ("fig3_nonlinear.png",
         "Linear versus nonlinear PK on log\u2013log survival axes. Blue: "
         "N-body data; red: linear PK; green dashed: hybrid "
         "Michaelis\u2013Menten model capturing the sticky-chaos tail. "
         "(a) Equal mass (linear adequate), (b) unequal (42 per cent tail "
         "RMSE improvement), (c) democratic (36 per cent).", "Figure 3"),
    ]))

    pe12 = g("A_population_scaling", "exp_mu12")
    peout = g("A_population_scaling", "exp_muout")
    peM = g("A_population_scaling", "exp_M")
    pr2 = g("A_population_scaling", "r_squared")
    pomega = g("A_population_scaling", "omega")
    pvif = g("A_population_scaling", "vif", default=[float("nan")] * 3)
    plmg = g("A_population_scaling", "shapley_r2", default=[float("nan")] * 3)
    portho = g("A_population_scaling", "r2_ortho")
    pint_r2 = g("A_population_scaling", "interaction_model", "r_squared")
    pint_beta = g("A_population_scaling", "interaction_model", "beta",
                  default=[float("nan")] * 5)
    pint_cross = pint_beta[4] if len(pint_beta) > 4 else float("nan")

    blocks.append(("Population PK reveals allometric scaling", [
        (f"Clinical PK relates individual rate constants to covariates through "
         f"mixed-effects regression.{{14}} Treating the reduced masses as "
         f"covariates over the 64-configuration grid, and using the tail-robust "
         f"median lifetime as the response, we obtain "
         f"t_{{1/2}} \u221d \u03bc_{{12}}^{{{_fmt(pe12, '.2f')}}} \u00d7 "
         f"\u03bc_{{out}}^{{{_fmt(peout, '.2f')}}} \u00d7 "
         f"M^{{{_fmt(peM, '.2f')}}} (R^2 = {_fmt(pr2, '.2f')}, "
         f"inter-individual variability \u03c9 = {_fmt(pomega, '.2f')}; "
         f"Fig. 4). Heavier initial binaries and heavier binary\u2013single "
         f"pairs survive longer (deeper potential wells; positive \u03bc_{{12}} "
         f"and \u03bc_{{out}} exponents), while a heavier total mass shortens "
         f"the interaction (negative M exponent), consistent with the shorter "
         f"dynamical time of a more massive triple."),
        (f"The three reduced-mass covariates are strongly collinear "
         f"(variance-inflation factors "
         f"{_fmt(pvif[0], '.1f')}/{_fmt(pvif[1], '.1f')}/{_fmt(pvif[2], '.1f')} "
         f"for \u03bc_{{12}}/\u03bc_{{out}}/M), so the individual exponents "
         f"should be read jointly rather than as independent partial "
         f"derivatives. An LMG/Shapley decomposition, which attributes the "
         f"explained variance fairly across correlated predictors, assigns "
         f"R^2 shares of "
         f"{_fmt(plmg[0], '.02f')}/{_fmt(plmg[1], '.02f')}/{_fmt(plmg[2], '.02f')} "
         f"to \u03bc_{{12}}/\u03bc_{{out}}/M, confirming that all three carry "
         f"real, comparable information. Refitting on the orthogonal design "
         f"axes (log m_2, log m_3; correlation "
         f"{_fmt(g('A_population_scaling', 'corr_m2m3'), '.2f')}) removes the "
         f"collinearity and yields a consistent trend "
         f"(R^2 = {_fmt(portho, '.2f')}). The modest R^2 reflects the scatter "
         f"intrinsic to chaotic scattering rather than estimation noise: "
         f"re-running the scan on a wider, finer three-dimensional grid "
         f"({int(g('A_population_scaling', 'expanded_scan_check', 'n_configs', default=0))} "
         f"configurations at 2500 encounters each) leaves the law essentially "
         f"unchanged (R^2 = "
         f"{_fmt(g('A_population_scaling', 'expanded_scan_check', 'r_squared'), '.2f')}, "
         f"same exponent signs), confirming that a several-fold increase in "
         f"sampling does not tighten the relation. As a secondary model, "
         f"adding a \u03bc_{{12}}\u00d7\u03bc_{{out}} interaction term raises "
         f"the explained variance to R^2 = {_fmt(pint_r2, '.2f')} "
         f"(from {_fmt(pr2, '.2f')}), with a positive cross-coefficient "
         f"(+{_fmt(pint_cross, '.2f')}) "
         f"indicating that a deep initial binary and a massive intruder "
         f"reinforce each other's effect on the interaction time; we retain "
         f"the pure power law as the primary result for interpretability. "
         f"The value of the law is "
         f"that it replaces ad hoc outcome fits with a closed-form scaling "
         f"that, as we show below, propagates directly into event-rate "
         f"estimates."),
    ], [
        ("fig4_population_pk.png",
         "Population-PK analysis of the mass-ratio dependence. (a) MRT versus "
         "binary\u2013single reduced mass; (b) median lifetime versus intruder "
         "mass; (c) mean excursions versus lightest mass fraction; "
         "(d) lightest-body escape probability with logistic fit; "
         "(e) predicted versus observed median lifetime for the allometric "
         "fit, annotated with the R^2, the LMG/Shapley R^2 shares and the "
         "variance-inflation factors for "
         "(\u03bc_{12}, \u03bc_{out}, M); (f) slowest "
         "half-life across the (m_2, m_3) grid.", "Figure 4"),
    ]))

    blocks.append(("Reverse mapping: a brief demonstration", [
        ("The equivalence runs both ways. Target-mediated drug disposition "
         "(TMDD) \u2014 drug, receptor, and internalised complex \u2014 shares "
         "the three-state generator structure of the compartment model.{15,16} "
         "The Jacobian-eigenvalue stability analysis used for Lagrange points "
         "in celestial mechanics therefore transfers directly to locating "
         "stable and unstable TMDD steady states (Fig. 5), the "
         "pharmacological analogue of the Mardling\u2013Aarseth criterion for "
         "hierarchical triples.{17} We keep this demonstration brief; our "
         "focus here is the forward direction. Its significance is that the "
         "two-way traffic of methods is itself the point \u2014 a productive "
         "synergy rather than a one-off analogy."),
    ], [
        ("fig5_tmdd.png",
         "Reverse mapping to TMDD. (a) Dose\u2013response bifurcation with "
         "stable (blue) and unstable (red) steady states; (b) receptor "
         "occupancy versus dose rate; (c) slowest Jacobian eigenvalue versus "
         "dose, marking the stability boundary \u2014 the analogue of Lagrange-"
         "point stability.", "Figure 5"),
    ]))

    blocks.append(("Robustness across dimensionality and dissipation", [
        ("To check that the mapping is not an artefact of planar, "
         "energy-conserving dynamics, we extended to fully three-dimensional "
         "encounters with thermal eccentricities and isotropic approach "
         "directions. The three-compartment model remains valid (KS "
         "statistics 0.16\u20130.31; Fig. 6), with median lifetimes shifting "
         "modestly relative to 2D. Adding 2.5PN gravitational-wave radiation "
         "reaction (c = 100 in code units) shortens lifetimes without "
         "destroying the compartment structure (Fig. 7), the democratic "
         "configuration showing the strongest effect."),
    ], [
        ("fig6_3d_dissipative_comparison.png",
         "PK validity across dimensionality and dissipation. Rows: 2D "
         "conservative, 3D conservative, 3D dissipative (c = 100). Columns: "
         "equal, unequal, democratic. Blue: N-body survival; red: linear PK "
         "fit; KS statistics annotated.", "Figure 6"),
        ("fig7_dissipation_effect.png",
         "Effect of 2.5PN GW radiation reaction on lifetimes. Blue: 3D "
         "conservative; red: 3D dissipative; dashed: PK fits. Dissipation "
         "shortens lifetimes (\u03c4_{diss}/\u03c4_{cons} = 0.83\u20130.98) "
         "while preserving the PK structure.", "Figure 7"),
    ]))

    blocks.append(("Gravitational-wave merger rates from the surviving binaries", [
        (f"Feeding the measured final-binary semi-major axes and "
         f"eccentricities into the Peters formula{{19}} turns the PK output "
         f"into a merger prediction. For a fiducial 10 + 10 M_\u2609 binary at "
         f"a_0 = 0.5 AU, the measured thermal eccentricity distribution gives "
         f"a merger fraction of {_fmt(a1_th, '.2f')} within a Hubble time, "
         f"against {_fmt(a1_ci, '.4f')} for circular orbits \u2014 an "
         f"eccentricity boost of the GW merger rate of order "
         f"{_fmt(a1_boost, '.0f')}\u00d7 (Fig. 8). Under 2.5PN dissipation "
         f"the compartment lifetimes contract by factors "
         f"{_fmt(lo, '.2f')}\u2013{_fmt(hi, '.2f')} across configurations, so "
         f"resonant encounters resolve faster than energy-conserving "
         f"estimates assume, plausibly raising the predicted inspiral rate in "
         f"dense systems.{{13,24,26}}"),
    ], [
        ("fig8_gw_merger.png",
         "Gravitational-wave merger outcomes from the PK-selected binaries. "
         "(a) Eccentricity versus semi-major axis of the surviving binaries. "
         "(b) Peters-inspiral merger fraction (t_{GW} < t_{H}) versus the "
         "assumed physical scale of the code-unit a = 1, for the measured "
         "thermal eccentricity distribution against circular orbits, showing "
         "the eccentricity boost of the merger rate. (c) Lifetime ratio "
         "\u03c4_{diss}/\u03c4_{cons} under 2.5PN dissipation for the three "
         "mass configurations.", "Figure 8"),
    ]))

    blocks.append(("Benchmarking against phase-space flux theory", [
        (f"The PK model reproduces, rather than competes with, established "
         f"statistical theory. The multi-exponential survival it predicts "
         f"tracks the empirical curves of all three canonical configurations "
         f"\u2014 the same ergodic/flux structure derived by Stone & Leigh and "
         f"Ginat & Perets \u2014 with a mean PK\u2013data KS distance of "
         f"{_fmt(a2_ks, '.2f')} (Fig. 9a; the long-lived tail of the unequal "
         f"case is the sticky-chaos excess captured by the nonlinear term of "
         f"Section 3.2). Across the {a2_n} mass "
         f"configurations, the closed-form PK MRT recovers the empirical mean "
         f"interaction time essentially exactly (R^2 = {_fmt(a2_r2, '.3f')}, "
         f"slope = {_fmt(g('A2_flux_benchmark', 'mrt_vs_lifetime_slope'), '.2f')}; "
         f"Fig. 9b). This near-perfect agreement validates the "
         f"compartment (Markov) approximation for the mean and shows that the "
         f"single MRT scalar is an unbiased closed-form estimator of the mean "
         f"interaction time \u2014 a quantity flux theory targets \u2014 rather "
         f"than a merely correlative fit.{{5,6,21}}"),
    ], [
        ("fig9_flux_benchmark.png",
         "Benchmark against phase-space flux theory. (a) Empirical survival "
         "(dots) with the PK multi-exponential (ergodic) prediction (lines) "
         "for the three canonical mass configurations. (b) PK-model MRT "
         "versus empirical mean lifetime across 64 mass configurations, with "
         "the y = x line.", "Figure 9"),
    ]))

    blocks.append(("Eccentricity and ejection-velocity distributions", [
        ("The surviving binaries and ejected singles carry observable "
         "signatures. The final-binary eccentricities are broadly consistent "
         "with a thermal distribution f(e) = 2e, with a superthermal excess "
         "at high e that feeds the GW channel above, and the ejection speeds "
         "peak at a fraction of the critical velocity of the system, with a "
         "tail extending beyond it (Fig. 10). These "
         "distributions are read straight off the PK-classified escape "
         "events, requiring no additional modelling."),
    ], [
        ("fig10_ecc_kick.png",
         "Distributions from 3D Newtonian encounters. Top: final-binary "
         "eccentricity with the thermal 2e reference; bottom: ejection speed "
         "in units of the critical velocity. Columns: equal, unequal, "
         "democratic.", "Figure 10"),
    ]))

    blocks.append(("Which body is ejected: coupling to the mass function", [
        (f"The logistic escape model \u2014 a dose\u2013response relation in "
         f"mass fraction \u2014 predicts which body leaves. Convolving it with "
         f"a Kroupa black-hole mass function (5\u201350 M_\u2609){{23}} gives a "
         f"lightest-body ejection fraction of "
         f"{_fmt((a3_frac or 0) * 100, '.0f')} per cent, above the "
         f"mass-blind expectation of 33 per cent, and an ejected-body mass "
         f"({_fmt(a3_mej, '.1f')} M_\u2609) systematically below the mean "
         f"member mass ({_fmt(a3_mmem, '.1f')} M_\u2609) \u2014 a clean "
         f"mass-segregation signal with direct consequences for black-hole "
         f"retention in clusters (Fig. 11).{{24,25}}"),
    ], [
        ("fig11_ejection_mf.png",
         "Which body is ejected. (a) Lightest-body escape probability versus "
         "minimum mass ratio with the fitted dose\u2013response logistic "
         "(points: simulations). (b) Ejected-mass distribution from a "
         "Kroupa-sampled black-hole mass function, showing preferential "
         "ejection of the lighter members.", "Figure 11"),
    ]))

    def pn_line(mode):
        d = b56.get(mode, {}) if isinstance(b56, dict) else {}
        ks = [v.get("ks_pk_fit") for v in d.values() if isinstance(v, dict)]
        rr = [v.get("lifetime_ratio") for v in d.values() if isinstance(v, dict)]
        ks = [x for x in ks if isinstance(x, (int, float))]
        rr = [x for x in rr if isinstance(x, (int, float))]
        ks_s = f"{min(ks):.2f}\u2013{max(ks):.2f}" if ks else "\u2014"
        rr_s = f"{min(rr):.2f}\u2013{max(rr):.2f}" if rr else "\u2014"
        return ks_s, rr_s

    pn_ks, pn_rr = pn_line("pn1")
    td_ks, td_rr = pn_line("tidal")

    blocks.append(("Robustness to precession and tidal dissipation", [
        (f"Finally we test two physical effects absent from the baseline. "
         f"Adding pairwise 1PN precession leaves the compartment structure "
         f"intact (PK-fit KS = {pn_ks}) with only a mild lifetime change "
         f"(\u03c4'/\u03c4 = {pn_rr}). A phenomenological tidal drag, relevant "
         f"to stellar rather than black-hole triples, likewise preserves the "
         f"mapping (PK-fit KS = {td_ks}; \u03c4'/\u03c4 = {td_rr}) while "
         f"shortening lifetimes as energy is bled off (Fig. 12). The PK "
         f"toolkit is thus robust across the dissipation mechanisms that "
         f"matter in real environments.{{17,22}}"),
    ], [
        ("fig12_pn_tidal.png",
         "Robustness to additional physics. Survival functions for 3D "
         "encounters with 1PN precession (top) and tidal dissipation "
         "(bottom), each against the Newtonian baseline and its own PK fit "
         "(dashed). Columns: equal, unequal, democratic; KS and lifetime "
         "ratios annotated.", "Figure 12"),
    ]))

    env_txt = ""
    if envs:
        parts = []
        for e in envs:
            parts.append(
                f"{e.get('env','')} (t_enc \u2248 "
                f"{e.get('t_enc_yr', float('nan')):.0e} yr, t_res \u2248 "
                f"{e.get('t_resonance_yr', float('nan')):.0e} yr)")
        env_txt = "; ".join(parts)

    blocks.append(("Application to dense stellar environments", [
        (f"Converting the PK-model MRT (in dynamical times) into physical "
         f"resonance durations and comparing with gravitational-focusing "
         f"encounter rates places the framework in astrophysical context "
         f"(Fig. 13). {('For ' + env_txt + ', ') if env_txt else ''}"
         f"the ratio of resonance duration to encounter interval controls "
         f"whether triples process one encounter at a time or overlap. This "
         f"single MRT scalar, propagated through the population-PK scaling, "
         f"yields environment-specific hardening and merger timescales "
         f"without bespoke Monte-Carlo surveys \u2014 the practical payoff of "
         f"importing the PK toolkit.{{24,25,27}}"),
    ], [
        ("fig13_environments.png",
         "Application to dense environments. Encounter interval (blue) versus "
         "PK-model resonance duration MRT \u00d7 t_dyn (orange) for a "
         "globular-cluster core, a nuclear star cluster, and an AGN-disc "
         "migration trap.", "Figure 13"),
    ]))

    return blocks


DISCUSSION = [
    ("We asked whether the pharmacokinetic toolkit, imported wholesale, is "
     "useful for the statistical three-body problem, and it is. A linear "
     "three-compartment model reproduces the multi-exponential lifetime law "
     "and gives the mean residence time in closed form; a Michaelis\u2013Menten "
     "term parametrises sticky chaos with one or two constants; and "
     "population-PK regression reduces the mass-ratio dependence to an "
     "allometric scaling law. These are not restatements of existing theory "
     "but estimators that sit on top of it and carry over to new problems: "
     "the PK MRT matches the empirical mean lifetime across all 64 "
     "configurations, and the three-compartment survival curve follows the "
     "ergodic prediction where a single flux timescale cannot."),

    ("It helps to separate what the toolkit adds from what statistical "
     "three-body theory already supplied. Flux-based and random-walk "
     "treatments give lifetime distributions, branching ratios and ejection "
     "statistics,{5,6} and we claim no new dynamics: the compartment reading "
     "reproduces them (Section 3.7). What it supplies are working capabilities "
     "that were previously awkward or missing. First, a single closed-form "
     "scalar \u2014 the mean residence time MRT = \u22121^T A^{\u22121} P0 \u2014 "
     "that enters merger-rate and environment-timescale estimates directly, "
     "with no further Monte-Carlo integration. Second, a population regression "
     "that folds the entire mass-spectrum dependence into one re-usable "
     "scaling law and weighs each covariate through variance-inflation factors "
     "and an LMG/Shapley decomposition, in place of fitting every spectrum by "
     "hand. Third, a one- or two-parameter Michaelis\u2013Menten form that "
     "joins the ergodic and sticky-chaos regimes smoothly instead of appending "
     "a bare power-law tail. Fourth, a route in the opposite direction: "
     "celestial-mechanics stability analysis exported to pharmacology "
     "(Section 3.4), which earlier three-body work had no reason to consider. "
     "The gains are in transferability and in how quickly a new problem can be "
     "set up, not in fundamental physics."),

    ("These estimators translate into concrete astrophysics. The same "
     "closed-form outputs give Peters-inspiral merger fractions and a large "
     "eccentricity boost of the GW merger rate, a dose\u2013response predictor "
     "for which body is ejected (and hence mass segregation once convolved "
     "with a realistic mass function), the eccentricity and kick distributions "
     "of the products, and resonance-versus-encounter timescales for globular "
     "clusters, nuclear star clusters and AGN discs. Since all of this follows "
     "from a handful of fitted rate constants, the framework re-tunes to a new "
     "mass spectrum or environment at little cost, complementing expensive "
     "direct surveys.{13,24,25}"),

    ("The mapping holds up: it survives the move from two to three dimensions, "
     "thermal eccentricities, 2.5PN radiation reaction, 1PN precession and "
     "tidal drag, with the compartment structure intact throughout. That range "
     "matters, because it is the regime \u2014 three-dimensional, eccentric, "
     "weakly dissipative \u2014 in which cluster and disc encounters actually "
     "occur."),

    ("The traffic runs both ways, and this is the part we find most "
     "suggestive. Celestial-mechanics stability theory describes TMDD steady "
     "states as naturally as PK describes three-body statistics: the "
     "Jacobian-eigenvalue analysis used for Lagrange points locates the stable "
     "and unstable receptor-binding states, so the geometric stability theory "
     "of celestial mechanics may guide the design of drugs with controlled "
     "target residence times.{18} Because the bridge rests only on the "
     "continuous-time Markov structure shared by multi-state first-passage "
     "processes, it need not stop at these two fields \u2014 chemical reaction "
     "networks and ecological dynamics share the same structure \u2014 a "
     "programme we would call the \u2018pharmacokinetics of chaos\u2019."),

    ("This is more than a formal curiosity. A mature clinical science, with "
     "decades of tested estimators, diagnostics and software behind it, is "
     "shown to carry real quantitative weight in an unrelated physical domain, "
     "and the same correspondence sends methods back the other way. Borrowing "
     "an established toolkit instead of rebuilding the equivalent machinery "
     "from scratch is a practical economy, and it supports a broader claim: "
     "moving concepts deliberately between disciplines is a research strategy "
     "in its own right, not merely an aesthetic analogy."),

    ("The compartmental reading is not tied to three bodies. A compartment is "
     "just a metastable configuration and a transfer a reshuffling between "
     "configurations, so the general N-body problem admits an N-compartment "
     "description whose compartments are the distinct bound sub-structures "
     "\u2014 binaries, stable hierarchical triples, higher multiples \u2014 and "
     "whose transfers are their formation and disruption, with escape once "
     "more the elimination channel. The number of compartments grows "
     "combinatorially with N; for small-N encounters it stays tractable, and "
     "for larger systems the generator can be coarse-grained onto a few "
     "dynamically distinct states (\u2018hard binary present\u2019 versus "
     "\u2018democratic\u2019), much as physiological PK lumps many tissues into "
     "a handful of kinetically distinguishable compartments. The closed-form "
     "residence-time and population-scaling machinery used here for N = 3 then "
     "carries over with little change, which points to the compartmental "
     "toolkit as a general instrument for the statistical mechanics of few- "
     "and many-body gravitational dynamics rather than a three-body special "
     "case. Testing this for N = 4\u20135 scattering and for small sub-clusters "
     "is the obvious next step."),

    ("Some limitations remain. The radiation-reaction treatment is "
     "perturbative (valid for c \u226b v_{orbital}); the tidal term is "
     "phenomenological rather than a full equilibrium-tide model; and the "
     "linear PK model assumes ergodicity, which the nonlinear extension only "
     "partly repairs, at the cost of extra parameters. None of these changes "
     "the central result, that the compartment structure is preserved and "
     "remains predictive."),
]

CONCLUSIONS = [
    ("Motivated by pharmacokinetics, we re-examined the three-body problem "
     "and found that the PK toolkit \u2014 compartment survival analysis, "
     "mean residence time, saturable kinetics, and population scaling \u2014 "
     "is a useful, portable instrument for three-body statistics, "
     "reproducing established theory while adding transferable estimators."),
    ("Using those estimators we derived quantitative astrophysical "
     "predictions: eccentricity-boosted GW merger fractions, a closed-form "
     "which-body-ejected predictor coupled to the black-hole mass function, "
     "eccentricity and kick distributions, and environment-specific "
     "interaction timescales \u2014 all robust to precession and dissipation."),
    ("The bidirectional exchange of methods between pharmacology and "
     "celestial mechanics is a productive synergy, and we expect the same "
     "compartmental bridge to reach other multi-state chaotic systems."),
]

ACKNOWLEDGMENTS = (
    "We thank the developers of NumPy, SciPy, Numba, and Matplotlib. "
    "This work made no use of proprietary data.")

DATA_AVAILABILITY = (
    "The simulation code, generated datasets, and analysis scripts that "
    "support this article are available in the project repository and from "
    "the corresponding author on reasonable request.")


# ---------------------------------------------------------------------------
# Builders
# ---------------------------------------------------------------------------

def build_manuscript():
    doc = Document()
    setup_style(doc)

    title_p = doc.add_paragraph()
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_p.paragraph_format.space_after = Pt(12)
    r = title_p.add_run(TITLE)
    r.bold = True
    r.font.size = Pt(15)

    auth_p = doc.add_paragraph()
    auth_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = auth_p.add_run(AUTHORS)
    r.font.size = Pt(12)

    aff_p = doc.add_paragraph()
    aff_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    aff_p.paragraph_format.space_after = Pt(4)
    r = aff_p.add_run(AFFILIATIONS)
    r.font.size = Pt(10)
    r.italic = True

    corr_p = doc.add_paragraph()
    corr_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    corr_p.paragraph_format.space_after = Pt(16)
    r = corr_p.add_run("E-mail: " + CORRESPONDING)
    r.font.size = Pt(10)

    # Abstract
    h = doc.add_paragraph()
    r = h.add_run("ABSTRACT")
    r.bold = True
    r.font.size = Pt(11)
    add_paragraph_with_refs(doc, ABSTRACT)

    kw = doc.add_paragraph()
    r = kw.add_run("Key words: ")
    r.bold = True
    r.font.size = Pt(10)
    r2 = kw.add_run("celestial mechanics \u2013 gravitation \u2013 "
                    "gravitational waves \u2013 methods: statistical \u2013 "
                    "stars: kinematics and dynamics.")
    r2.font.size = Pt(10)
    kw.paragraph_format.space_after = Pt(12)

    # 1 Introduction
    add_heading(doc, "1 INTRODUCTION", level=1)
    for para in INTRO:
        add_paragraph_with_refs(doc, para)

    # 2 Methods
    add_heading(doc, "2 METHODS", level=1)
    for i, (mtitle, paras) in enumerate(METHODS, 1):
        add_heading(doc, f"2.{i} {mtitle}", level=2)
        for para in paras:
            add_paragraph_with_refs(doc, para)

    # 3 Results
    add_heading(doc, "3 RESULTS", level=1)
    for i, (stitle, paras, figs) in enumerate(get_results_blocks(), 1):
        add_heading(doc, f"3.{i} {stitle}", level=2)
        for para in paras:
            add_paragraph_with_refs(doc, para)
        for fig_file, caption, label in figs:
            add_figure(doc, fig_file, caption, label)

    # 4 Discussion
    add_heading(doc, "4 DISCUSSION", level=1)
    for para in DISCUSSION:
        add_paragraph_with_refs(doc, para)

    # 5 Conclusions
    add_heading(doc, "5 CONCLUSIONS", level=1)
    for para in CONCLUSIONS:
        add_paragraph_with_refs(doc, para)

    # Back matter
    add_heading(doc, "ACKNOWLEDGEMENTS", level=2)
    doc.add_paragraph(ACKNOWLEDGMENTS)
    add_heading(doc, "DATA AVAILABILITY", level=2)
    doc.add_paragraph(DATA_AVAILABILITY)

    # References (alphabetical, author-year)
    add_heading(doc, "REFERENCES", level=1)
    refs_sorted = sorted(REFS.values(), key=lambda t: t[1].lower())
    for _, full in refs_sorted:
        p = doc.add_paragraph()
        p.paragraph_format.left_indent = Inches(0.3)
        p.paragraph_format.first_line_indent = Inches(-0.3)
        p.paragraph_format.space_after = Pt(3)
        r = p.add_run(full)
        r.font.size = Pt(10)

    out = os.path.join(OUTDIR, "manuscript_mnras.docx")
    doc.save(out)
    return out


def build_cover_letter():
    doc = Document()
    setup_style(doc)
    today = date.today().strftime("%d %B %Y")
    for line in [today, "", "The Editors", "Monthly Notices of the Royal "
                 "Astronomical Society", ""]:
        doc.add_paragraph(line)
    doc.add_paragraph("Dear Editors,")
    body = [
        ("Please consider our manuscript, \u201c" + TITLE + "\u201d, for "
         "publication in MNRAS as a research article."),
        ("The paper imports the toolkit of clinical pharmacokinetics \u2014 "
         "compartmental survival analysis, mean-residence-time calculus, "
         "saturable kinetics, and population (mixed-effects) scaling \u2014 "
         "into the statistical three-body problem, and shows that it is a "
         "genuinely useful set of instruments rather than a loose analogy. "
         "Using 15,000 planar and 27,000 new three-dimensional scattering "
         "experiments, we demonstrate that the compartment model reproduces "
         "the multi-exponential lifetime law and the phase-space-flux "
         "predictions of Stone & Leigh (2019) and Ginat & Perets (2021), "
         "while adding closed-form estimators that established theory does "
         "not provide."),
        ("We then use those estimators to make quantitative astrophysical "
         "predictions of direct interest to MNRAS readers: eccentricity-"
         "boosted gravitational-wave merger fractions, a closed-form "
         "predictor for which body is ejected (coupled to a realistic "
         "black-hole mass function), the eccentricity and ejection-velocity "
         "distributions of the products, and environment-specific "
         "resonance-versus-encounter timescales for globular clusters, "
         "nuclear star clusters, and AGN discs. The compartment structure is "
         "shown to survive 1PN precession, tidal dissipation, and 2.5PN "
         "radiation reaction."),
        ("The correspondence is bidirectional: celestial-mechanics stability "
         "theory in turn informs target-mediated drug disposition. We believe "
         "this two-way exchange, and the concrete astrophysical predictions "
         "it enables, will interest the broad MNRAS readership working on "
         "dynamical formation of compact-object binaries."),
        ("The manuscript is original, not under consideration elsewhere, and "
         "the author declares no conflict of interest."),
    ]
    for para in body:
        add_paragraph_with_refs(doc, para)
    doc.add_paragraph("")
    doc.add_paragraph("Sincerely,")
    doc.add_paragraph(AUTHORS)
    doc.add_paragraph(AFFILIATIONS)
    doc.add_paragraph("E-mail: " + CORRESPONDING)
    out = os.path.join(OUTDIR, "cover_letter_mnras.docx")
    doc.save(out)
    return out


def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches as PInches, Pt as PPt
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]
    figs = []
    for _, _, block_figs in get_results_blocks():
        figs.extend(block_figs)
    for fig_file, caption, label in figs:
        path = os.path.join(FIGDIR, fig_file)
        if not os.path.exists(path):
            continue
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(PInches(0.4), PInches(0.15),
                                      PInches(12.5), PInches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label
        run.font.bold = True
        run.font.size = PPt(20)
        # image
        from PIL import Image
        with Image.open(path) as im:
            w, h = im.size
        avail_w, avail_h = 12.5, 5.4
        ar = w / h
        disp_w = avail_w
        disp_h = disp_w / ar
        if disp_h > avail_h:
            disp_h = avail_h
            disp_w = disp_h * ar
        left = (13.333 - disp_w) / 2
        slide.shapes.add_picture(path, PInches(left), PInches(0.9),
                                 width=PInches(disp_w))
        cb = slide.shapes.add_textbox(PInches(0.4), PInches(6.5),
                                      PInches(12.5), PInches(0.9))
        ctf = cb.text_frame
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        crun = cp.add_run()
        cap_plain = re.sub(r"[_^]\{([^}]*)\}", r"\1", caption)
        crun.text = cap_plain
        crun.font.size = PPt(11)
    out = os.path.join(OUTDIR, "figures_mnras_editable.pptx")
    prs.save(out)
    return out


def copy_figures():
    figs = []
    for _, _, block_figs in get_results_blocks():
        figs.extend(block_figs)
    copied = []
    for fig_file, _, _ in figs:
        src = os.path.join(FIGDIR, fig_file)
        if os.path.exists(src):
            dst = os.path.join(FIGOUT, fig_file)
            shutil.copy2(src, dst)
            copied.append(dst)
            pdf = src.replace(".png", ".pdf")
            if os.path.exists(pdf):
                shutil.copy2(pdf, os.path.join(FIGOUT, os.path.basename(pdf)))
    return copied


def create_zip(paths):
    out = os.path.join(OUTDIR, "threebody_pk_mnras_submission.zip")
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as z:
        for p in paths:
            if p and os.path.exists(p):
                z.write(p, os.path.basename(p))
        for f in sorted(os.listdir(FIGOUT)):
            z.write(os.path.join(FIGOUT, f), os.path.join("figures", f))
    return out


def main():
    print("Building MNRAS submission package ...")
    ms = build_manuscript()
    print("  manuscript:", ms)
    cl = build_cover_letter()
    print("  cover letter:", cl)
    figs = copy_figures()
    print(f"  copied {len(figs)} figures -> {FIGOUT}")
    try:
        pptx = build_pptx()
        print("  pptx:", pptx)
    except Exception as e:
        pptx = None
        print("  pptx FAILED:", e)
    z = create_zip([ms, cl, pptx])
    print("  zip:", z)


if __name__ == "__main__":
    main()
