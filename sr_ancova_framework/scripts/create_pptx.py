#!/usr/bin/env python3
"""Create editable PPTX with one figure per slide."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pathlib import Path

OUT_DIR = Path(__file__).resolve().parent.parent / 'pre_submission'

FIGURES = [
    {
        'file': 'fig1_schematic.png',
        'title': 'Figure 1',
        'caption': (
            'Conceptual schematic of the covariate-adjusted stochastic resonance framework. '
            '(a) Threshold-based event detector: a weak periodic signal s(t) embedded in '
            'Gaussian noise n(t) triggers events when |x(t)| exceeds threshold θ. '
            '(b) Stochastic resonance: event rate modulation is maximized at an intermediate '
            'noise level (green), suppressed at too-low noise (red), and washed out at too-high '
            'noise (blue). '
            '(c) Covariate adjustment narrows the residual noise distribution, effectively '
            'shifting the operating point on the SR curve.'
        ),
    },
    {
        'file': 'fig2_sr_curves.png',
        'title': 'Figure 2',
        'caption': (
            'Stochastic resonance curves for a threshold detector (A/θ = 0.3). '
            'Solid lines: analytical SNR from two-state theory. Points with error bars: '
            'Monte Carlo validation via cross-correlation. Covariate adjustment (ρ > 0) '
            'shifts the SR peak to higher input noise levels, meaning the detector tolerates '
            'more environmental noise when a good noise model is available.'
        ),
    },
    {
        'file': 'fig3_optimal_rho.png',
        'title': 'Figure 3',
        'caption': (
            'Optimal noise model accuracy ρ* and the corresponding SNR improvement. '
            '(a) When σ < θ (SR regime, yellow shading), ρ* = 0: any noise removal '
            'degrades performance. When σ > θ (excess noise regime, blue shading), '
            'ρ* increases to bring effective noise to the SR optimum. '
            '(b) SNR improvement at optimal ρ* grows exponentially with input noise.'
        ),
    },
    {
        'file': 'fig4_detection_probability.png',
        'title': 'Figure 4',
        'caption': (
            'Detection probability P_D and false alarm probability P_FA as functions of '
            'input noise level for different covariate model accuracies ρ. '
            'Higher ρ suppresses both P_D and P_FA by reducing effective noise; '
            'the net effect on detection performance depends on the operating regime.'
        ),
    },
    {
        'file': 'fig5_roc_comparison.png',
        'title': 'Figure 5',
        'caption': (
            'ROC curves at fixed input noise σ_n/θ = 1.5 (excess noise regime) '
            'for different covariate adjustment levels (A/θ = 0.4). In this regime, '
            'higher ρ improves the ROC curve, confirming that covariate adjustment is '
            'beneficial when the system operates above the SR optimum.'
        ),
    },
    {
        'file': 'fig6_mutual_information.png',
        'title': 'Figure 6',
        'caption': (
            'Mutual information I(S; E) between the periodic signal and the event stream '
            'as a function of input noise level. Covariate adjustment (ρ = 0.8) shifts '
            'the MI peak to higher input noise, consistent with the SR curve analysis.'
        ),
    },
    {
        'file': 'fig7_dvs_application.png',
        'title': 'Figure 7',
        'caption': (
            'Application to dynamic vision sensor (DVS) astronomical observation. '
            '(a) ROC-AUC for noise classification: the Fano filter (covariate adjustment '
            'approach) achieves AUC = 0.866, far exceeding temporal filtering and neural '
            'methods. (b) NRR vs SPR trade-off: the Fano filter preserves 93.9% of signal '
            'events while removing 71.3% of noise.'
        ),
    },
]


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    for fig_info in FIGURES:
        slide = prs.slides.add_slide(blank_layout)
        img_path = OUT_DIR / fig_info['file']
        if not img_path.exists():
            print(f"  WARNING: {img_path} not found, skipping")
            continue

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = fig_info['title']
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        # Image — centered, scaled to fit
        from PIL import Image
        with Image.open(img_path) as img:
            img_w, img_h = img.size
        aspect = img_w / img_h

        max_w = Inches(12)
        max_h = Inches(5.0)
        if aspect > (max_w / max_h):
            w = max_w
            h = int(w / aspect)
        else:
            h = max_h
            w = int(h * aspect)

        left = int((prs.slide_width - w) / 2)
        top = Inches(0.9)
        slide.shapes.add_picture(str(img_path), left, top, w, h)

        # Caption
        cap_top = top + h + Inches(0.15)
        txBox2 = slide.shapes.add_textbox(Inches(0.5), cap_top, Inches(12), Inches(1.2))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = fig_info['caption']
        p2.font.size = Pt(12)
        p2.alignment = PP_ALIGN.LEFT

    out_path = OUT_DIR / 'figures.pptx'
    prs.save(str(out_path))
    print(f"PPTX saved: {out_path}")


if __name__ == '__main__':
    main()
