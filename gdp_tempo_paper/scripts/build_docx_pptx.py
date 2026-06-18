"""Build final manuscript .docx (JA + EN) with inline figures/tables
and a single editable .pptx (EN figures) for per-slide editing.

Formatted for Journal of Macroeconomics (Elsevier):
  - Double-blind: title page with author info separated from manuscript body
  - 1.5× line spacing, Times New Roman 12pt body text
  - APA 7 author-date references (alphabetical, with DOIs)
  - Decimal section numbering (1, 1.1, 1.2, …)
  - Highlights (3-5 bullets, ≤85 chars each)
  - Figures inline in manuscript + separate files for submission

Usage:  python build_docx_pptx.py
Outputs into ../manuscript/ :
  - manuscript_en.docx / .pdf  (PDF for Editorial Manager submission)
  - manuscript_ja.docx / .pdf
  - figures_en.pptx
  - table1_correspondence.docx
  - table2_model_metrics.docx
  - table3_rpim.docx
  - table4_extended_oos.docx
  - table5_tempo_artifact.docx
  - cover_letter_en.docx / .pdf
"""
from __future__ import annotations

import os
import re
import json
from dataclasses import dataclass

import pandas as pd
from docx import Document
from docx.shared import Pt, Inches, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_ALIGN_VERTICAL
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGB

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.abspath(os.path.join(HERE, ".."))
MS = os.path.join(ROOT, "manuscript")
FIG = os.path.join(ROOT, "figures")
TAB = os.path.join(ROOT, "tables")
os.makedirs(MS, exist_ok=True)

# ---------------------------------------------------------------------------
# Figure/table metadata

FIG_LIST = [
    ("fig1", "Fig. 1", "図1",
     "In-sample 1-year GDP growth RMSE across M0-M4 for 39 countries (lower is better).",
     "39カ国における単年GDP成長率RMSEの標本内比較（M0-M4、小さいほど良い）。",
     "fig1_m_ranking_{lang}.png"),
    ("fig2", "Fig. 2", "図2",
     "Out-of-sample MAPE on 2015-19 held-out window; M2 achieves 13% relative improvement vs M0.",
     "2015-19年のホールドアウト窓における標本外MAPE。M2はM0比で相対13%改善。",
     "fig2_oos_{lang}.png"),
    ("fig3", "Fig. 3", "図3",
     "PIM stock K_tang+betaK_I versus CWON PCA (within-country demeaned log).",
     "PIM資本ストックK_tang+βK_IとCWON PCAの軌跡比較（国内平均除去対数）。",
     "fig3_trajectories_{lang}.png"),
    ("fig4", "Fig. 4", "図4",
     "gamma_price sensitivity of PIM/CWON log-ratio; Japan gap closes around gamma=+0.02.",
     "PIM/CWON対数比のγ_price感度。日本の乖離はγ=+0.02で閉じる。",
     "fig4_gamma_price_{lang}.png"),
    ("fig5", "Fig. 5", "図5",
     "Conceptual diagram of the population-capital tempo correspondence.",
     "人口・資本テンポ対応関係の概念図。",
     "fig5_concept_{lang}.png"),
    ("fig6", "Fig. 6", "図6",
     "Relational PIM diagnostics: rho_2 across 39 countries under M0 vs M4.",
     "関係型PIM診断: M0とM4における39カ国のρ̂₂。",
     "fig6_rpim_{lang}.png"),
    ("fig7", "Fig. 7", "図7",
     "Delta-mu sensitivity: mu_hat under +/-20% depreciation perturbation.",
     "δ-μ感度分析: 減価償却率±20%変動下のμ̂。",
     "fig7_delta_sensitivity_{lang}.png"),
    ("fig8", "Fig. 8", "図8",
     "Conditional OOS evaluation: interior-solution vs boundary countries.",
     "条件付き標本外評価: 内点解国 vs 境界解国。",
     "fig8_conditional_oos_{lang}.png"),
    ("fig9", "Fig. 9", "図9",
     "Cross-sectional regression of rho_2 on R&D intensity.",
     "ρ̂₂のR&D強度に対するクロスセクション回帰。",
     "fig9_rho2_regression_{lang}.png"),
    ("fig10", "Fig. 10", "図10",
     "Solow-residual decomposition: M0 vs tempo-adjusted (M2) vs joint (M4) for six representative countries.",
     "ソロー残差の分解: M0 vs テンポ調整(M2) vs 統合(M4)、代表6カ国。",
     "fig10_solow_decomp_{lang}.png"),
    ("fig11", "Fig. 11", "図11",
     "National wealth: CWON official vs intangible-adjusted produced capital (2019).",
     "国富: CWON公式値 vs 無形資本調整後の生産資本（2019年）。",
     "fig11_counterfactual_wealth_{lang}.png"),
]


# ---------------------------------------------------------------------------
# docx helpers

def set_font(run, name="Times New Roman", size=11, bold=False, italic=False):
    run.font.name = name
    run.font.size = Pt(size)
    run.bold = bold
    run.italic = italic
    # Apply East Asian font for proper JA rendering
    rPr = run._element.get_or_add_rPr()
    rFonts = rPr.find(qn("w:rFonts"))
    if rFonts is None:
        rFonts = OxmlElement("w:rFonts")
        rPr.append(rFonts)
    rFonts.set(qn("w:ascii"), name)
    rFonts.set(qn("w:hAnsi"), name)
    rFonts.set(qn("w:eastAsia"), "MS Mincho")


def add_heading(doc, text, level, lang):
    """Add heading with JEG / Springer formatting.

    Level 1 (## in md): section heading — left-aligned, 14pt bold.
    Level 2 (## special): same as level 1 but 14pt.
    Level 3 (### in md): subsection — flush left, 12pt bold.
    All levels: sentence case (no ALL CAPS), left-aligned.
    """
    h = doc.add_paragraph()
    h.paragraph_format.space_before = Pt(18)
    h.paragraph_format.space_after = Pt(6)
    h.paragraph_format.line_spacing = 1.5
    # JEG/Springer: left-aligned, no ALL CAPS
    display_text = text
    run = h.add_run(display_text)
    set_font(run, size={1: 16, 2: 14, 3: 12}[level], bold=True)
    return h


GREEK_LETTERS = set("αβγδεζηθικλμνξοπρσςτυφχψω"
                     "ΑΒΓΔΕΖΗΘΙΚΛΜΝΞΟΠΡΣΤΥΦΧΨΩ")

# Tokens that should be rendered in italic when they appear as standalone
# variable-like identifiers in math context.
MATH_VARS = re.compile(
    r"""(?<![A-Za-z_])      # not preceded by letter/underscore
    (                        # group 1: the variable token
      [A-Z]                  # single capital letter ...
      (?:_[A-Za-z0-9]+)?    # ... optionally followed by _subscript
      |                      # OR
      [αβγδεζηθικλμνξοπρσςτυφχψωΑΒΓΔΕΖΗΘΙΚΛΜΝΞΟΠΡΣΤΥΦΧΨΩ] # Greek letter
      (?:_[A-Za-z0-9]+)?    # optional subscript
    )
    (?![A-Za-z])             # not followed by letter
    """,
    re.VERBOSE,
)

# Pattern to detect subscript notations: _{...} or _X (single char)
SUB_RE = re.compile(r'_\{([^}]+)\}|_([A-Za-z0-9])')
# Pattern to detect superscript notations: ^{...} or ^X (single char) or trailing *
SUP_RE = re.compile(r'\^\{([^}]+)\}|\^([A-Za-z0-9])')

# Displayed equation: 4-space indent, optional label like (M0), (1), (2)
EQUATION_RE = re.compile(r'^    (.+?)\s{2,}\(([A-Za-z0-9]+)\)\s*$')
EQUATION_NOLABEL_RE = re.compile(r'^    (.+?)\s*$')


def _is_greek(ch):
    return ch in GREEK_LETTERS


def _add_run(paragraph, text, size=12, bold=False, italic=False,
             superscript=False, subscript=False, font_name="Times New Roman"):
    """Add a single run with specified formatting."""
    run = paragraph.add_run(text)
    set_font(run, name=font_name, size=size, bold=bold, italic=italic)
    if superscript:
        run.font.superscript = True
    if subscript:
        run.font.subscript = True
    return run


def add_math_runs(paragraph, text, size=12, base_italic=False):
    """Parse *text* and emit Word runs with proper math formatting.

    Rules applied (in order of priority):
    1. Markdown-style *italic* spans → italic runs.
    2. Subscript notation  _{...}  or  _x  → subscript runs.
    3. Superscript notation ^{...} or ^x  → superscript runs.
    4. Greek letters → always italic.
    5. Single uppercase letters that look like variables (K, I, Y, L, A, W, S,
       N, B) → italic when adjacent to math context (parens, subscripts, Greek).
    6. Everything else → roman (upright).
    """
    if not text:
        return

    # Step 1: split on markdown italic markers *...*
    # We process each segment for math notation
    parts = re.split(r'(\*[^*]+\*)', text)

    for part in parts:
        if not part:
            continue
        if part.startswith('*') and part.endswith('*') and len(part) > 2:
            # Markdown italic span — render entire span in italic, then parse
            # subscripts/superscripts within
            inner = part[1:-1]
            _emit_math_segment(paragraph, inner, size, force_italic=True)
        else:
            _emit_math_segment(paragraph, inner=part, size=size,
                               force_italic=base_italic)


def _emit_math_segment(paragraph, inner, size, force_italic=False,
                       force_bold=False):
    """Emit a text segment with subscript/superscript/Greek handling."""
    pos = 0
    while pos < len(inner):
        ch = inner[pos]

        # Check for subscript notation
        if ch == '_':
            if pos + 1 >= len(inner):
                # Trailing _ — emit as literal
                _add_run(paragraph, ch, size=size,
                         bold=force_bold, italic=force_italic)
                pos += 1
                continue
            if inner[pos + 1] == '{':
                # _{...} subscript
                end = inner.find('}', pos + 2)
                if end != -1:
                    sub_text = inner[pos + 2:end]
                    _add_run(paragraph, sub_text, size=size,
                             bold=force_bold, italic=True, subscript=True)
                    pos = end + 1
                    continue
                else:
                    # Unclosed brace — emit _{ as literal
                    _add_run(paragraph, '_{', size=size,
                             bold=force_bold, italic=force_italic)
                    pos += 2
                    continue
            else:
                # _x single-char subscript
                sub_text = inner[pos + 1]
                _add_run(paragraph, sub_text, size=size,
                         bold=force_bold, italic=True, subscript=True)
                pos += 2
                continue

        # Check for superscript notation
        if ch == '^':
            if pos + 1 >= len(inner):
                # Trailing ^ — emit as literal
                _add_run(paragraph, ch, size=size,
                         bold=force_bold, italic=force_italic)
                pos += 1
                continue
            if inner[pos + 1] == '{':
                end = inner.find('}', pos + 2)
                if end != -1:
                    sup_text = inner[pos + 2:end]
                    _add_run(paragraph, sup_text, size=size,
                             bold=force_bold, italic=True, superscript=True)
                    pos = end + 1
                    continue
                else:
                    # Unclosed brace — emit ^{ as literal
                    _add_run(paragraph, '^{', size=size,
                             bold=force_bold, italic=force_italic)
                    pos += 2
                    continue
            else:
                sup_text = inner[pos + 1]
                _add_run(paragraph, sup_text, size=size,
                         bold=force_bold, italic=True, superscript=True)
                pos += 2
                continue

        # Greek letter → always italic
        if _is_greek(ch):
            _add_run(paragraph, ch, size=size, bold=force_bold, italic=True)
            pos += 1
            continue

        # Collect a run of "normal" characters (not _ ^ or Greek)
        start = pos
        while pos < len(inner):
            if inner[pos] in ('_', '^') or _is_greek(inner[pos]):
                break
            pos += 1
        chunk = inner[start:pos]
        if chunk:
            _add_run(paragraph, chunk, size=size,
                     bold=force_bold, italic=force_italic)


def add_para(doc, text, lang, italic=False):
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(6)
    p.paragraph_format.line_spacing = 1.5
    p.paragraph_format.first_line_indent = Pt(24)
    add_math_runs(p, text, size=12, base_italic=italic)
    return p


def add_rich_para(doc, text, lang, bullet=False):
    """Add a paragraph with bold **...** spans and math-aware formatting."""
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(6)
    p.paragraph_format.line_spacing = 1.5
    if bullet:
        p.paragraph_format.left_indent = Pt(36)
        p.paragraph_format.first_line_indent = Pt(-18)
    else:
        p.paragraph_format.first_line_indent = Pt(24)

    # Unescape markdown backslash-escaped characters (e.g. \* → *)
    text = re.sub(r'\\([*_\\`])', r'\1', text)

    # Split on **bold** markers first
    parts = re.split(r'(\*\*[^*]+\*\*)', text)
    for part in parts:
        if not part:
            continue
        if part.startswith('**') and part.endswith('**') and len(part) > 4:
            inner = part[2:-2]
            # Bold span — emit with bold + math formatting
            _emit_bold_math(p, inner, size=12)
        else:
            # Normal text — emit with math formatting (handles *italic*, Greek, sub/sup)
            add_math_runs(p, part, size=12, base_italic=False)
    return p


def _emit_bold_math(paragraph, text, size=12):
    """Emit bold text with math-aware formatting (Greek italic, sub/sup)."""
    parts = re.split(r'(\*[^*]+\*)', text)
    for part in parts:
        if not part:
            continue
        if part.startswith('*') and part.endswith('*') and len(part) > 2:
            inner = part[1:-1]
            _emit_math_segment(paragraph, inner, size, force_italic=True,
                               force_bold=True)
        else:
            _emit_math_segment(paragraph, part, size, force_italic=False,
                               force_bold=True)


def add_equation_block(doc, equation_text, label=None):
    """Add a displayed equation: centered equation with right-aligned label."""
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(12)
    p.paragraph_format.space_after = Pt(12)
    p.paragraph_format.line_spacing = 1.5
    p.paragraph_format.first_line_indent = Pt(0)

    # Render equation body with math formatting
    add_math_runs(p, equation_text.strip(), size=12, base_italic=True)

    # Add equation label right-aligned using a tab stop
    if label:
        # Add tab + label in upright (non-italic)
        tab_run = p.add_run("\t")
        set_font(tab_run, size=12)
        label_run = p.add_run(f"({label})")
        set_font(label_run, size=12, italic=False)
        # Set a right-aligned tab stop at the right margin
        pPr = p._element.get_or_add_pPr()
        tabs = OxmlElement('w:tabs')
        tab = OxmlElement('w:tab')
        tab.set(qn('w:val'), 'right')
        tab.set(qn('w:pos'), '9072')  # ~16cm from left = right margin
        tab.set(qn('w:leader'), 'none')
        tabs.append(tab)
        pPr.append(tabs)


def add_figure(doc, png_path, caption_prefix, caption_text):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(18)
    p.paragraph_format.space_after = Pt(6)
    run = p.add_run()
    run.add_picture(png_path, width=Inches(6.2))
    cap = doc.add_paragraph()
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    cap.paragraph_format.space_after = Pt(12)
    r1 = cap.add_run(f"{caption_prefix}. ")
    set_font(r1, size=10, bold=True)
    r2 = cap.add_run(caption_text)
    set_font(r2, size=10)


def add_dataframe_as_table(doc, df: pd.DataFrame, col_widths=None, font_size=10):
    tbl = doc.add_table(rows=1 + len(df), cols=len(df.columns))
    tbl.style = "Table Grid"
    # JoM/Elsevier style: no vertical rules in tables
    tbl_element = tbl._tbl
    tblPr = tbl_element.tblPr
    if tblPr is None:
        tblPr = OxmlElement('w:tblPr')
        tbl_element.insert(0, tblPr)
    borders = OxmlElement('w:tblBorders')
    for border_name in ('left', 'right', 'insideV'):
        border_el = OxmlElement(f'w:{border_name}')
        border_el.set(qn('w:val'), 'none')
        border_el.set(qn('w:sz'), '0')
        border_el.set(qn('w:space'), '0')
        border_el.set(qn('w:color'), 'auto')
        borders.append(border_el)
    # Keep horizontal rules
    for border_name in ('top', 'bottom', 'insideH'):
        border_el = OxmlElement(f'w:{border_name}')
        border_el.set(qn('w:val'), 'single')
        border_el.set(qn('w:sz'), '4')
        border_el.set(qn('w:space'), '0')
        border_el.set(qn('w:color'), '000000')
        borders.append(border_el)
    tblPr.append(borders)
    # Header — math-aware formatting for Greek/subscript in column names
    for j, col in enumerate(df.columns):
        c = tbl.rows[0].cells[j]
        c.text = ""
        p = c.paragraphs[0]
        _emit_math_segment(p, str(col), size=font_size, force_bold=True)
    for i, row in enumerate(df.itertuples(index=False), start=1):
        for j, v in enumerate(row):
            c = tbl.rows[i].cells[j]
            c.text = ""
            p = c.paragraphs[0]
            _emit_math_segment(p, str(v), size=font_size)
    if col_widths:
        for row in tbl.rows:
            for idx, w in enumerate(col_widths):
                row.cells[idx].width = w


def add_table_block(doc, title, df, caption, widths=None):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(18)
    p.paragraph_format.space_after = Pt(6)
    r = p.add_run(title)
    set_font(r, size=11, bold=True)
    r2 = p.add_run(f"  {caption}")
    set_font(r2, size=11)
    add_dataframe_as_table(doc, df, col_widths=widths)


# ---------------------------------------------------------------------------
# Build docx manuscript from markdown, replacing figure/table placeholders

INSERT_FIG_RE = re.compile(r"^\*\*\[(Insert |)[Ff]ig(ure|\.)? ?(\d+)[^\]]*\]\*\*$")
INSERT_FIG_JA_RE = re.compile(r"^\*\*［図\s*(\d+)[^］]*］\*\*$")
INSERT_TAB_RE = re.compile(r"^\*\*\[(Insert |)[Tt]able ?(\d+)[^\]]*\]\*\*$")
INSERT_TAB_JA_RE = re.compile(r"^\*\*［表\s*(\d+)[^］]*］\*\*$")


def build_manuscript(lang: str):
    md_path = os.path.join(MS, f"manuscript_{lang}.md")
    with open(md_path, encoding="utf-8") as fh:
        lines = fh.readlines()

    # Load tables
    t1 = pd.read_csv(os.path.join(TAB, "table1_model_metrics.csv"))
    t2 = pd.read_csv(os.path.join(TAB, "table2_correspondence.csv"))
    t3_path = os.path.join(TAB, "table3_rpim.csv")
    t3 = pd.read_csv(t3_path) if os.path.exists(t3_path) else None
    t4_path = os.path.join(TAB, "table4_extended_oos.csv")
    t4 = pd.read_csv(t4_path) if os.path.exists(t4_path) else None
    t5_path = os.path.join(TAB, "table5_tempo_artifact.csv")
    t5 = pd.read_csv(t5_path) if os.path.exists(t5_path) else None

    # Figure caption lookup by index
    fig_cap = {}
    for key, en_prefix, ja_prefix, en_cap, ja_cap, pattern in FIG_LIST:
        idx = int(key.replace("fig", ""))
        fig_cap[idx] = {
            "en": (en_prefix, en_cap, pattern.format(lang="en")),
            "ja": (ja_prefix, ja_cap, pattern.format(lang="ja")),
        }

    t1_cap_en = "Population-capital tempo correspondence."
    t1_cap_ja = "人口・資本テンポ対応関係。"
    t2_cap_en = "M0-M4: in-sample and out-of-sample performance across 39 countries. "\
                "Medians across countries; IQR in brackets."
    t2_cap_ja = "M0-M4: 39カ国の標本内・標本外パフォーマンス（国間中央値、IQRを括弧内）。"
    t3_cap_en = "Relational PIM diagnostics: rho_2 summary under M0, M1, M2, M4."
    t3_cap_ja = "関係型PIM診断: M0, M1, M2, M4におけるρ̂₂の要約。"
    t4_cap_en = "Extended OOS metrics: direction accuracy and CWON trajectory RMSE."
    t4_cap_ja = "拡張標本外指標: 方向精度およびCWON軌跡RMSE。"
    t5_cap_en = "Tempo-artifact share of TFP-growth variance: percentage reduction in Var(d log TFP) from M0 to M2 (tempo) and M0 to M4 (joint)."
    t5_cap_ja = "テンポ・アーティファクトのTFP成長率分散シェア: M0→M2（テンポ）およびM0→M4（統合）。"

    doc = Document()
    # Use sensible page margins
    for section in doc.sections:
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    # JoM double-blind: skip title-page block in the anonymized manuscript
    # The manuscript .md contains [TITLE PAGE ...] ... [END TITLE PAGE] markers
    in_title_page = False

    i = 0
    while i < len(lines):
        line = lines[i].rstrip("\n")
        stripped = line.strip()
        if not stripped:
            i += 1
            continue

        # Skip title-page section (author info goes to separate file)
        if stripped.startswith("[TITLE PAGE") or stripped.startswith("[タイトルページ"):
            in_title_page = True
            i += 1
            continue
        if stripped.startswith("[END TITLE PAGE") or stripped.startswith("[タイトルページ終了"):
            in_title_page = False
            i += 1
            continue
        if in_title_page:
            i += 1
            continue

        # Skip manuscript-body marker line
        if stripped.startswith("[MANUSCRIPT") or stripped.startswith("[原稿"):
            i += 1
            continue

        # Headings
        if stripped.startswith("# "):
            add_heading(doc, stripped[2:], 1, lang)
        elif stripped.startswith("## "):
            add_heading(doc, stripped[3:], 2, lang)
        elif stripped.startswith("### "):
            add_heading(doc, stripped[4:], 3, lang)
        elif stripped.startswith("---"):
            pass  # ignore horizontal rules
        else:
            # Figure inline placeholder (both lang patterns)
            m_fig = INSERT_FIG_RE.match(stripped) or INSERT_FIG_JA_RE.match(stripped)
            m_tab = INSERT_TAB_RE.match(stripped) or INSERT_TAB_JA_RE.match(stripped)
            if m_fig:
                try:
                    idx = int(m_fig.group(m_fig.lastindex))
                except Exception:
                    idx = None
                if idx in fig_cap:
                    prefix, cap, fname = fig_cap[idx][lang]
                    png = os.path.join(FIG, fname)
                    if os.path.exists(png):
                        add_figure(doc, png, prefix, cap)
            elif m_tab:
                try:
                    idx = int(m_tab.group(m_tab.lastindex))
                except Exception:
                    idx = None
                if idx == 1:
                    add_table_block(
                        doc,
                        "Table 1." if lang == "en" else "表1.",
                        t2,
                        t1_cap_en if lang == "en" else t1_cap_ja,
                        widths=[Inches(1.3), Inches(2.5), Inches(2.5)],
                    )
                elif idx == 2:
                    add_table_block(
                        doc,
                        "Table 2." if lang == "en" else "表2.",
                        t1,
                        t2_cap_en if lang == "en" else t2_cap_ja,
                    )
                elif idx == 3 and t3 is not None:
                    add_table_block(
                        doc,
                        "Table 3." if lang == "en" else "表3.",
                        t3,
                        t3_cap_en if lang == "en" else t3_cap_ja,
                    )
                elif idx == 4 and t4 is not None:
                    add_table_block(
                        doc,
                        "Table 4." if lang == "en" else "表4.",
                        t4,
                        t4_cap_en if lang == "en" else t4_cap_ja,
                    )
                elif idx == 5 and t5 is not None:
                    add_table_block(
                        doc,
                        "Table 5." if lang == "en" else "表5.",
                        t5,
                        t5_cap_en if lang == "en" else t5_cap_ja,
                    )
            else:
                # Displayed equation: 4-space indent with label (M0), (1), etc.
                m_eq = EQUATION_RE.match(line)
                if m_eq:
                    add_equation_block(doc, m_eq.group(1), m_eq.group(2))
                    i += 1
                    continue

                text = stripped
                # section separator
                if text in ("Tables", "表", "References", "参考文献"):
                    add_heading(doc, text, 2, lang)
                # Bullet list items: * text... or - text...
                elif text.startswith("* "):
                    add_rich_para(doc, "•  " + text[2:], lang, bullet=True)
                elif text.startswith("- "):
                    add_rich_para(doc, "•  " + text[2:], lang, bullet=True)
                else:
                    add_rich_para(doc, text, lang)
        i += 1

    out = os.path.join(MS, f"manuscript_{lang}.docx")
    doc.save(out)
    print("wrote", out)


# ---------------------------------------------------------------------------
# Separate table .docx files

def build_standalone_tables():
    t1 = pd.read_csv(os.path.join(TAB, "table1_model_metrics.csv"))
    t2 = pd.read_csv(os.path.join(TAB, "table2_correspondence.csv"))
    t3_path = os.path.join(TAB, "table3_rpim.csv")
    t3 = pd.read_csv(t3_path) if os.path.exists(t3_path) else None
    table_list = [
        ("table1_correspondence.docx", t2,
         "Table 1. Population-capital tempo correspondence.",
         [Inches(1.3), Inches(2.5), Inches(2.5)]),
        ("table2_model_metrics.docx", t1,
         "Table 2. M0-M4: in-sample and out-of-sample performance across 39 countries.",
         None),
    ]
    if t3 is not None:
        table_list.append(
            ("table3_rpim.docx", t3,
             "Table 3. Relational PIM diagnostics: rho_2 summary under M0, M1, M2, M4.",
             None))
    t4_path = os.path.join(TAB, "table4_extended_oos.csv")
    t4 = pd.read_csv(t4_path) if os.path.exists(t4_path) else None
    if t4 is not None:
        table_list.append(
            ("table4_extended_oos.docx", t4,
             "Table 4. Extended OOS metrics: direction accuracy and CWON trajectory RMSE.",
             None))
    t5_path = os.path.join(TAB, "table5_tempo_artifact.csv")
    t5 = pd.read_csv(t5_path) if os.path.exists(t5_path) else None
    if t5 is not None:
        table_list.append(
            ("table5_tempo_artifact.docx", t5,
             "Table 5. Tempo-artifact share of TFP-growth variance.",
             None))
    for name, df, cap, widths in table_list:
        d = Document()
        add_heading(d, cap, 2, "en")
        add_dataframe_as_table(d, df, col_widths=widths)
        out = os.path.join(MS, name)
        d.save(out)
        print("wrote", out)


# ---------------------------------------------------------------------------
# Editable pptx (English figures only, one per slide)

def build_pptx():
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]

    for key, en_prefix, _ja, en_cap, _jc, pattern in FIG_LIST:
        png = os.path.join(FIG, pattern.format(lang="en"))
        if not os.path.exists(png):
            print("skip missing", png)
            continue
        slide = prs.slides.add_slide(blank)
        # Title
        tb = slide.shapes.add_textbox(
            PptInches(0.4), PptInches(0.2),
            PptInches(12.5), PptInches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{en_prefix}. {en_cap}"
        r.font.size = PptPt(20)
        r.font.bold = True
        # Image: fit to 12 x 5.5 area, centered
        slide.shapes.add_picture(
            png, PptInches(1.5), PptInches(1.1),
            width=PptInches(10.3), height=PptInches(5.6))
        # Caption bar at bottom
        cb = slide.shapes.add_textbox(
            PptInches(0.4), PptInches(6.85),
            PptInches(12.5), PptInches(0.5))
        cf = cb.text_frame
        cf.word_wrap = True
        pp = cf.paragraphs[0]
        rr = pp.add_run()
        rr.text = en_cap
        rr.font.size = PptPt(14)
        rr.font.italic = True
        rr.font.color.rgb = PptRGB(0x44, 0x44, 0x44)

    out = os.path.join(MS, "figures_en.pptx")
    prs.save(out)
    print("wrote", out)


# ---------------------------------------------------------------------------
# PDF conversion via LibreOffice (fonts embedded by default)

def convert_docx_to_pdf(docx_path: str):
    """Convert a .docx file to PDF using LibreOffice.

    LibreOffice embeds fonts by default, satisfying the Editorial Express
    requirement that all fonts be embedded in the PDF.
    """
    import subprocess
    out_dir = os.path.dirname(docx_path)
    cmd = [
        "libreoffice", "--headless", "--convert-to", "pdf",
        "--outdir", out_dir, docx_path,
    ]
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
    if result.returncode != 0:
        print(f"WARNING: PDF conversion failed for {docx_path}")
        print(result.stderr)
    else:
        pdf_path = docx_path.replace(".docx", ".pdf")
        print(f"wrote {pdf_path}")


def build_title_page():
    """Build separate title page .docx with author info (for double-blind submission)."""
    md_path = os.path.join(MS, "manuscript_en.md")
    if not os.path.exists(md_path):
        return
    with open(md_path, encoding="utf-8") as fh:
        lines = fh.readlines()

    # Extract title-page content between markers
    title_lines = []
    in_tp = False
    for raw_line in lines:
        line = raw_line.rstrip("\n").strip()
        if line.startswith("[TITLE PAGE"):
            in_tp = True
            continue
        if line.startswith("[END TITLE PAGE"):
            break
        if in_tp:
            title_lines.append(raw_line.rstrip("\n"))

    if not title_lines:
        print("no title page content found")
        return

    doc = Document()
    for section in doc.sections:
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    for raw_line in title_lines:
        stripped = raw_line.strip()
        if not stripped:
            spacer = doc.add_paragraph()
            spacer.paragraph_format.space_after = Pt(0)
            spacer.paragraph_format.space_before = Pt(0)
            spacer.paragraph_format.line_spacing = 1.5
            continue
        if stripped.startswith("# "):
            add_heading(doc, stripped[2:], 1, "en")
        else:
            add_rich_para(doc, stripped, "en")

    out = os.path.join(MS, "title_page_en.docx")
    doc.save(out)
    print(f"wrote {out}")
    convert_docx_to_pdf(out)


def build_cover_letter():
    """Build cover letter docx from markdown source."""
    CL = os.path.join(ROOT, "cover_letter")
    md_path = os.path.join(CL, "cover_letter_en.md")
    if not os.path.exists(md_path):
        print(f"cover letter not found: {md_path}")
        return
    with open(md_path, encoding="utf-8") as fh:
        lines = fh.readlines()

    doc = Document()
    for section in doc.sections:
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    for raw_line in lines:
        line = raw_line.rstrip("\n")
        stripped = line.strip()

        if not stripped:
            # Blank line → small spacer paragraph
            spacer = doc.add_paragraph()
            spacer.paragraph_format.space_after = Pt(0)
            spacer.paragraph_format.space_before = Pt(0)
            spacer.paragraph_format.line_spacing = 1.5
            continue

        # Bullet list items
        if stripped.startswith("- "):
            p = doc.add_paragraph()
            p.paragraph_format.space_after = Pt(2)
            p.paragraph_format.line_spacing = 1.5
            p.paragraph_format.left_indent = Pt(36)
            p.paragraph_format.first_line_indent = Pt(-18)
            text = "•  " + stripped[2:]
            text = re.sub(r'\\([*_\\`])', r'\1', text)
            add_math_runs(p, text, size=12, base_italic=False)
            continue

        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(2)
        p.paragraph_format.line_spacing = 1.5

        # Handle bold + italic markdown
        text = stripped
        text = re.sub(r'\\([*_\\`])', r'\1', text)

        # Split on **bold** markers
        parts = re.split(r'(\*\*[^*]+\*\*)', text)
        for part in parts:
            if not part:
                continue
            if part.startswith('**') and part.endswith('**') and len(part) > 4:
                inner = part[2:-2]
                # Bold — further split on *italic*
                sub_parts = re.split(r'(\*[^*]+\*)', inner)
                for sp in sub_parts:
                    if not sp:
                        continue
                    if sp.startswith('*') and sp.endswith('*') and len(sp) > 2:
                        _emit_math_segment(p, sp[1:-1], size=12,
                                           force_italic=True, force_bold=True)
                    else:
                        _emit_math_segment(p, sp, size=12,
                                           force_italic=False, force_bold=True)
            else:
                # Normal text — handle *italic* within
                add_math_runs(p, part, size=12, base_italic=False)

    out = os.path.join(CL, "cover_letter_en.docx")
    doc.save(out)
    print(f"wrote {out}")
    convert_docx_to_pdf(out)


def main():
    build_manuscript("en")
    build_manuscript("ja")
    build_standalone_tables()
    build_pptx()
    build_title_page()
    build_cover_letter()
    # Convert manuscripts to PDF
    for lang in ("en", "ja"):
        docx = os.path.join(MS, f"manuscript_{lang}.docx")
        if os.path.exists(docx):
            convert_docx_to_pdf(docx)


if __name__ == "__main__":
    main()
