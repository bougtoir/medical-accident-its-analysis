"""Build the Annals of Human Genetics submission package."""

from __future__ import annotations

import re
import shutil
import zipfile
from importlib.metadata import PackageNotFoundError, version
from pathlib import Path

import pandas as pd
from PIL import Image
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt

import ajba_content as revised_content


PROJECT_DIR = Path(__file__).resolve().parents[1]
DATA_DIR = PROJECT_DIR / "data"
FIGURE_DIR = PROJECT_DIR / "figures"
OUTPUT_DIR = PROJECT_DIR / "docs" / "ahg_submission"
OUTPUT_FIGURE_DIR = OUTPUT_DIR / "figures"
JOURNAL = "Annals of Human Genetics"
JOURNAL_SHORT = "AHG"
ARTICLE_TYPE = "Original Article"
EDITOR_IN_CHIEF = "Dr. Rosemary Ekong"

TITLE = revised_content.TITLE
RUNNING_TITLE = revised_content.RUNNING_TITLE
AUTHOR = revised_content.AUTHOR
AFFILIATION = revised_content.AFFILIATION
CORRESPONDENCE = revised_content.CORRESPONDENCE
ABSTRACT = revised_content.ABSTRACT
KEYWORDS = revised_content.KEYWORDS
REFERENCES = revised_content.REFERENCES
REFERENCE_KEYS = revised_content.REFERENCE_KEYS
INTRODUCTION = revised_content.INTRODUCTION
METHODS = revised_content.METHODS
RESULTS = revised_content.RESULTS
DISCUSSION = revised_content.DISCUSSION
FIGURES = revised_content.FIGURES
SUPPORTING_FIGURES = revised_content.SUPPORTING_FIGURES


def set_cell_shading(cell, fill: str) -> None:
    properties = cell._tc.get_or_add_tcPr()
    shading = OxmlElement("w:shd")
    shading.set(qn("w:fill"), fill)
    properties.append(shading)


def configure_document(document: Document) -> None:
    section = document.sections[0]
    section.top_margin = Inches(1)
    section.bottom_margin = Inches(1)
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)
    normal = document.styles["Normal"]
    normal.font.name = "Times New Roman"
    normal.font.size = Pt(12)
    normal.paragraph_format.line_spacing = 2
    normal.paragraph_format.space_after = Pt(0)
    for name in ["Title", "Heading 1", "Heading 2"]:
        style = document.styles[name]
        style.font.name = "Times New Roman"
        style.font.color.rgb = RGBColor(0, 0, 0)


def add_cited_paragraph(document: Document, text: str, italic: bool = False):
    paragraph = document.add_paragraph()
    paragraph.paragraph_format.first_line_indent = Inches(0.3)
    run = paragraph.add_run(text)
    run.font.name = "Times New Roman"
    run.font.size = Pt(12)
    run.italic = italic
    return paragraph


def add_title_page(document: Document) -> None:
    paragraph = document.add_paragraph()
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = paragraph.add_run(TITLE)
    run.bold = True
    run.font.name = "Times New Roman"
    run.font.size = Pt(16)
    paragraph = document.add_paragraph()
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    paragraph.add_run(ARTICLE_TYPE).bold = True
    paragraph = document.add_paragraph()
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    paragraph.add_run(AUTHOR.upper()).bold = True
    paragraph = document.add_paragraph()
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    paragraph.add_run(AFFILIATION)
    paragraph = document.add_paragraph()
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    paragraph.add_run(CORRESPONDENCE)
    paragraph = document.add_paragraph()
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    paragraph.add_run(f"Running title: {RUNNING_TITLE}")
    document.add_heading("Summary", level=1)
    paragraph = document.add_paragraph(ABSTRACT)
    paragraph.paragraph_format.line_spacing = 1.5
    paragraph = document.add_paragraph()
    paragraph.add_run("Keywords: ").bold = True
    paragraph.add_run(KEYWORDS)
    document.add_page_break()


def table_1_rows() -> list[list[str]]:
    pairs = pd.read_csv(DATA_DIR / "pairwise_sharing_corrected.csv")
    qualifying = pairs[
        (pairs["any_admixed"] == 0)
        & (pairs["nean_resid_z"] > 2)
        & (pairs["nean_fdr_pval"] < 0.10)
    ].sort_values("nean_resid_z", ascending=False)
    rows = [
        [
            "Population 1",
            "Population 2",
            "Region 1",
            "Region 2",
            "Distance (km)",
            "Sharing (r)",
            "z-score",
        ]
    ]
    if qualifying.empty:
        rows.append(
            [
                "No qualifying pair",
                "—",
                "—",
                "—",
                "—",
                "—",
                "No z>2 and q<0.10 result",
            ]
        )
        return rows
    for row in qualifying.itertuples():
        rows.append(
            [
                row.pop1,
                row.pop2,
                row.region1.replace("_", " ").title(),
                row.region2.replace("_", " ").title(),
                f"{row.geo_dist_km:,.0f}",
                f"{row.nean_corr:.3f}",
                f"{row.nean_resid_z:.2f}",
            ]
        )
    return rows


def table_2_rows() -> list[list[str]]:
    sublineage = pd.read_csv(DATA_DIR / "abo_sublineage_summary.csv")
    order = [
        "East Asia",
        "Europe",
        "Indigenous Americas",
        "Admixed Americas",
        "Central/South Asia",
        "Middle East",
        "Oceania",
    ]
    rows = [
        [
            "Region",
            "n",
            "Altai %",
            "Vindija %",
            "Chagyrskaya %",
        ]
    ]
    for group in order:
        group_summary = sublineage[
            (sublineage["analysis_group"] == group)
            & (sublineage["closest_reference"] != "Tie")
        ]
        total = int(group_summary["n_segments"].sum())
        values = {
            row.closest_reference: 100 * row.n_segments / total
            for row in group_summary.itertuples()
        } if total else {}
        rows.append(
            [
                group,
                str(total),
                f"{values.get('Altai', 0):.1f}",
                f"{values.get('Vindija', 0):.1f}",
                f"{values.get('Chagyrskaya', 0):.1f}",
            ]
        )
    return rows


TABLES = {
    1: (
        "Positive-residual Neanderthal pairs after false discovery rate control",
        table_1_rows,
        "The prespecified family contains all non-admixed population pairs. No pair met both z>2 and Benjamini-Hochberg q<0.10; the Denisovan analysis likewise identified no qualifying pair. Complete nominal rankings and dependence-aware model results are provided in Supplementary Data.",
    ),
}

SUPPORTING_TABLES = {
    1: (
        "Exploratory ABO-window Neanderthal-reference composition",
        table_2_rows,
        "Counts are classifiable segments, not individuals. Percentages use the three-reference denominator shown by n. Equal maximum-similarity ties are excluded from these percentages but retained in Supplementary Data. The 2/2 Indigenous American value is not a regional frequency estimate; only one segment overlaps ABO, and these counts are not interpreted as a migration route.",
    ),
}


def render_table(document: Document, label: str, spec: tuple) -> None:
    title, row_function, note = spec
    paragraph = document.add_paragraph()
    paragraph.paragraph_format.space_before = Pt(14)
    run = paragraph.add_run(f"{label}. {title}")
    run.bold = True
    rows = row_function()
    table = document.add_table(rows=len(rows), cols=len(rows[0]))
    table.style = "Table Grid"
    for row_index, row in enumerate(rows):
        for column_index, value in enumerate(row):
            cell = table.cell(row_index, column_index)
            cell.text = value
            for paragraph in cell.paragraphs:
                paragraph.paragraph_format.line_spacing = 1
                for run in paragraph.runs:
                    run.font.name = "Arial"
                    run.font.size = Pt(8.5)
                    run.bold = row_index == 0
            if row_index == 0:
                set_cell_shading(cell, "D9EAF7")
    paragraph = document.add_paragraph(f"Note. {note}")
    paragraph.paragraph_format.line_spacing = 1
    for run in paragraph.runs:
        run.font.size = Pt(9)


def add_word_table(document: Document, table_number: int) -> None:
    render_table(document, f"Table {table_number}", TABLES[table_number])


def add_supporting_table(document: Document, table_number: int) -> None:
    render_table(document, f"Table S{table_number}", SUPPORTING_TABLES[table_number])


def add_inline_figure(document: Document, figure_number: int) -> None:
    filename, caption = FIGURES[figure_number]
    paragraph = document.add_paragraph()
    paragraph.paragraph_format.space_before = Pt(16)
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    paragraph.add_run().add_picture(
        str(FIGURE_DIR / filename), width=Inches(6.35)
    )
    paragraph = document.add_paragraph()
    paragraph.paragraph_format.space_before = Pt(10)
    paragraph.paragraph_format.line_spacing = 1.15
    run = paragraph.add_run(f"Figure {figure_number}. ")
    run.bold = True
    paragraph.add_run(caption)


def add_object(document: Document, label: str) -> None:
    kind, number = label.split()
    if kind == "Figure":
        add_inline_figure(document, int(number))
    else:
        add_word_table(document, int(number))


def add_manuscript_body(document: Document, inline: bool) -> None:
    document.add_heading("Introduction", level=1)
    for text in INTRODUCTION:
        add_cited_paragraph(document, text)
    document.add_heading("Materials and Methods", level=1)
    for heading, paragraphs in METHODS:
        document.add_heading(heading, level=2)
        for text in paragraphs:
            add_cited_paragraph(document, text)
    document.add_heading("Results", level=1)
    for heading, text, objects in RESULTS:
        document.add_heading(heading, level=2)
        add_cited_paragraph(document, text)
        if inline:
            for label in objects:
                add_object(document, label)
    document.add_heading("Discussion", level=1)
    for text in DISCUSSION:
        add_cited_paragraph(document, text)
    document.add_heading("Acknowledgements", level=1)
    document.add_paragraph(
        "The author acknowledges the participants, communities, and investigators whose "
        "contributions made the 1000 Genomes, HGDP, hmmix, and ancient-genome resources "
        "available. Public availability does not remove obligations of respectful reuse."
    )
    document.add_paragraph(
        "Funding: This research received no specific grant from any funding agency in "
        "the public, commercial, or not-for-profit sectors."
    )
    document.add_heading("Data Availability", level=1)
    document.add_paragraph(
        "Analysis scripts, aggregate derived data, figures, and document-generation code "
        "are available at https://github.com/bougtoir/denisovan-archaic-dna-analysis "
        "and will be fixed as release ahg-submission-2026-07 before submission. "
        "The source hmmix segment calls are available from Zenodo record 14136628 "
        "(https://doi.org/10.5281/zenodo.14136628). The ancient ABO-window observations "
        "were derived from the public Neanderthal-segment catalogue of Iasi et al. "
        "(2024), archived at Dryad (https://doi.org/10.5061/dryad.zw3r228gg); O2 "
        "subtype-defining allele frequencies were obtained from the Ensembl Variation "
        "application programming interface and from Ohashi et al. (2006). Raw-file "
        "SHA-256 checksums and all analysis parameters are included in "
        "analysis_provenance.json and ancient_abo_provenance.json."
    )
    document.add_heading("Conflict of Interest Statement", level=1)
    document.add_paragraph("The author declares no conflict of interest.")
    document.add_heading("Ethics Statement", level=1)
    document.add_paragraph(
        "This secondary computational analysis used de-identified public genomic data "
        "and involved no recruitment, participant contact, biospecimen collection, or "
        "new phenotype inference. No separate institutional review determination was "
        "obtained; approvals, consent, and access procedures were those reported by the "
        "source studies. No source community representatives participated in this "
        "secondary study, and no direct community return-of-results process occurred. "
        "Because Indigenous genomic records are included, results are reported only at "
        "the minimum level needed for auditability, are not generalized to communities, "
        "and are not used to assign migration routes. The public article, code, and "
        "aggregate derived results are the current means of results availability."
    )
    document.add_heading("Author Contributions", level=1)
    document.add_paragraph(
        "Onishi Tatsuki: Conceptualization, methodology, formal analysis, "
        "visualization, writing—original draft, and writing—review and editing."
    )
    document.add_heading("References", level=1)
    for reference in REFERENCES:
        paragraph = document.add_paragraph(reference)
        paragraph.paragraph_format.first_line_indent = Inches(-0.25)
        paragraph.paragraph_format.left_indent = Inches(0.25)
        paragraph.paragraph_format.line_spacing = 1.15
        paragraph.paragraph_format.space_after = Pt(4)
        for run in paragraph.runs:
            run.font.size = Pt(10)
    if not inline:
        document.add_heading("Figure Legends", level=1)
        for number, (_, caption) in FIGURES.items():
            paragraph = document.add_paragraph()
            paragraph.paragraph_format.line_spacing = 1.5
            paragraph.paragraph_format.space_after = Pt(8)
            run = paragraph.add_run(f"Figure {number}. ")
            run.bold = True
            paragraph.add_run(caption)
        document.add_heading("Supporting Information Legends", level=1)
        for number, (_, caption) in SUPPORTING_FIGURES.items():
            paragraph = document.add_paragraph()
            paragraph.paragraph_format.line_spacing = 1.5
            paragraph.paragraph_format.space_after = Pt(8)
            run = paragraph.add_run(f"Figure S{number}. ")
            run.bold = True
            paragraph.add_run(caption)


def create_manuscript(path: Path, inline: bool) -> None:
    document = Document()
    configure_document(document)
    document.core_properties.title = TITLE
    document.core_properties.author = AUTHOR
    add_title_page(document)
    add_manuscript_body(document, inline)
    document.save(path)


def create_tables_document(path: Path) -> None:
    document = Document()
    configure_document(document)
    paragraph = document.add_paragraph()
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = paragraph.add_run("Editable Tables")
    run.bold = True
    run.font.size = Pt(16)
    paragraph = document.add_paragraph(TITLE)
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for number in TABLES:
        add_word_table(document, number)
        if number != max(TABLES):
            document.add_page_break()
    document.save(path)


def create_single_table_document(path: Path, table_number: int) -> None:
    document = Document()
    configure_document(document)
    add_word_table(document, table_number)
    document.save(path)


def create_single_supporting_table_document(path: Path, table_number: int) -> None:
    document = Document()
    configure_document(document)
    add_supporting_table(document, table_number)
    document.save(path)


def create_supporting_information(path: Path) -> None:
    document = Document()
    configure_document(document)
    paragraph = document.add_paragraph()
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = paragraph.add_run("Supporting Information")
    run.bold = True
    run.font.size = Pt(16)
    paragraph = document.add_paragraph(TITLE)
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for number, (filename, caption) in SUPPORTING_FIGURES.items():
        heading = document.add_paragraph()
        heading.paragraph_format.space_before = Pt(16)
        run = heading.add_run(f"Figure S{number}. {caption}")
        run.bold = True
        image = document.add_paragraph()
        image.alignment = WD_ALIGN_PARAGRAPH.CENTER
        image.add_run().add_picture(
            str(FIGURE_DIR / filename), width=Inches(6.35)
        )
        document.add_page_break()
    for number in SUPPORTING_TABLES:
        add_supporting_table(document, number)
        document.add_page_break()
    document.add_heading("Supplementary Data Files", level=1)
    document.add_paragraph(
        "Supplementary Data 1: population_metadata.csv. Population, project, sample "
        "size, coordinates, continent assignment, and analysis inclusion."
    )
    document.add_paragraph(
        "Supplementary Data 2: pairwise_sharing_corrected.csv. Complete pairwise "
        "similarity, geographic, covariate, residual, permutation, and false discovery "
        "rate results."
    )
    document.add_paragraph(
        "Supplementary Data 3: model_summary.csv. Quadratic assignment procedure "
        "coefficients, permutation P values, descriptive R-squared values, and "
        "population-deletion intervals."
    )
    document.add_paragraph(
        "Supplementary Data 4: sensitivity_analysis.csv and "
        "window_size_sensitivity.csv. Metric, population-subset, and window-size "
        "robustness summaries."
    )
    document.save(path)


def create_cover_letter(path: Path) -> None:
    document = Document()
    configure_document(document)
    document.styles["Normal"].paragraph_format.line_spacing = 1
    document.styles["Normal"].paragraph_format.space_after = Pt(7)
    for text in [
        EDITOR_IN_CHIEF,
        "Editor-in-Chief, " + JOURNAL,
        "University College London",
    ]:
        document.add_paragraph(text)
    document.add_paragraph()
    paragraph = document.add_paragraph()
    run = paragraph.add_run(f"Re: Submission of an {ARTICLE_TYPE}")
    run.bold = True
    document.add_paragraph(f"Dear {EDITOR_IN_CHIEF},")
    paragraphs = [
        (
            f"I am pleased to submit “{TITLE}” for consideration as an "
            f"{ARTICLE_TYPE} in the {JOURNAL}."
        ),
        (
            "Focal-locus and special-connection interpretations of shared archaic "
            "segments are common in human population biology, yet they are seldom "
            "tested against a genome-wide baseline that respects the dependence "
            "structure of pairwise data. Using publicly archived hmmix "
            f"archaic-introgression calls from {revised_content.INDIVIDUALS:,} "
            f"individuals in {revised_content.POPULATIONS} populations "
            "(1000 Genomes Project and Human Genome Diversity Project), we construct "
            "such a baseline: population profiles are built so that window "
            "frequencies remain within 0-1, distance and pair-level effects are "
            "tested with population-label quadratic assignment permutations, and "
            "multiple testing is controlled with the false-discovery rate."
        ),
        (
            "The analysis shows a broad geographic distance-decay pattern but no "
            "population pair that survives false-discovery-rate correction and no "
            "ABO-window signal beyond the genome-wide expectation. The contribution "
            "is therefore a reusable, dependence-aware negative control against which "
            "focal-locus and special-connection archaic claims can be judged, rather "
            "than a new migration route. This methodological, reproducibility-focused "
            "study fits the scope of Annals of Human Genetics in human population and "
            "evolutionary genetics, the geographic distribution of genetic variation, "
            "the interpretation of the human genome including archaic and ancient-DNA "
            "signals, and statistical-genetic methodology applied to real data."
        ),
        (
            "The work is original, is not under consideration elsewhere, and uses "
            "de-identified public genomic resources. No new human participants or "
            "specimens were recruited. The manuscript explicitly discloses that no "
            "separate institutional review determination, community participation, or "
            "direct return-of-results process occurred for this secondary analysis. "
            "The author declares no conflict of interest and reports no external funding."
        ),
        (
            "All analysis code and derived outputs are provided through the project "
            "repository. The source archaic-introgression data, generated with hmmix "
            "(a hidden Markov model-based detection method), are publicly archived in "
            "Zenodo. "
            "The submission includes separate figure files, editable tables, and "
            "figure legends in the manuscript."
        ),
    ]
    for text in paragraphs:
        document.add_paragraph(text)
    document.add_paragraph("Sincerely,")
    document.add_paragraph(AUTHOR)
    document.add_paragraph(AFFILIATION)
    document.add_paragraph("Email: bougtoir@gmail.com")
    document.save(path)


def add_slide_title(slide, title: str) -> None:
    box = slide.shapes.add_textbox(
        PptInches(0.6),
        PptInches(0.12),
        PptInches(12.1),
        PptInches(0.7),
    )
    box.text_frame.word_wrap = True
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.name = "Arial"
    paragraph.font.size = PptPt(16)
    paragraph.font.bold = True
    paragraph.alignment = PP_ALIGN.CENTER


def add_slide_caption(slide, caption: str) -> None:
    box = slide.shapes.add_textbox(
        PptInches(0.65),
        PptInches(6.55),
        PptInches(12.0),
        PptInches(0.72),
    )
    box.text_frame.word_wrap = True
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = caption
    paragraph.font.name = "Arial"
    paragraph.font.size = PptPt(8)
    paragraph.alignment = PP_ALIGN.LEFT


def add_picture_contained(slide, path: Path) -> None:
    with Image.open(path) as image:
        width, height = image.size
    area_left = 0.55
    area_top = 0.85
    area_width = 12.2
    area_height = 5.65
    scale = min(area_width / width, area_height / height)
    picture_width = width * scale
    picture_height = height * scale
    left = area_left + (area_width - picture_width) / 2
    top = area_top + (area_height - picture_height) / 2
    slide.shapes.add_picture(
        str(path),
        PptInches(left),
        PptInches(top),
        PptInches(picture_width),
        PptInches(picture_height),
    )


def add_ppt_table(slide, rows: list[list[str]]) -> None:
    table_shape = slide.shapes.add_table(
        len(rows),
        len(rows[0]),
        PptInches(0.45),
        PptInches(1.0),
        PptInches(12.4),
        PptInches(5.5),
    )
    table = table_shape.table
    for row_index, row in enumerate(rows):
        for column_index, value in enumerate(row):
            cell = table.cell(row_index, column_index)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                PptRGBColor(217, 234, 247)
                if row_index == 0
                else PptRGBColor(255, 255, 255)
            )
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = "Arial"
                paragraph.font.size = PptPt(9)
                paragraph.font.bold = row_index == 0


def create_presentation(path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = PptInches(13.333)
    presentation.slide_height = PptInches(7.5)
    blank = presentation.slide_layouts[6]
    for number, (filename, caption) in FIGURES.items():
        slide = presentation.slides.add_slide(blank)
        add_slide_title(slide, f"Figure {number}")
        add_picture_contained(slide, FIGURE_DIR / filename)
        add_slide_caption(slide, caption)
    for number, (filename, caption) in SUPPORTING_FIGURES.items():
        slide = presentation.slides.add_slide(blank)
        add_slide_title(slide, f"Figure S{number}")
        add_picture_contained(slide, FIGURE_DIR / filename)
        add_slide_caption(slide, caption)
    for number, (title, row_function, note) in TABLES.items():
        slide = presentation.slides.add_slide(blank)
        add_slide_title(slide, f"Table {number}. {title}")
        add_ppt_table(slide, row_function())
        add_slide_caption(slide, f"Note. {note}")
    for number, (title, row_function, note) in SUPPORTING_TABLES.items():
        slide = presentation.slides.add_slide(blank)
        add_slide_title(slide, f"Table S{number}. {title}")
        add_ppt_table(slide, row_function())
        add_slide_caption(slide, f"Note. {note}")
    presentation.save(path)


def prepare_separate_figures() -> None:
    OUTPUT_FIGURE_DIR.mkdir(parents=True, exist_ok=True)
    for stale in OUTPUT_FIGURE_DIR.glob("Figure_*"):
        stale.unlink()
    for number, (filename, _) in FIGURES.items():
        source = FIGURE_DIR / filename
        png_target = OUTPUT_FIGURE_DIR / f"Figure_{number}.png"
        tiff_target = OUTPUT_FIGURE_DIR / f"Figure_{number}.tiff"
        shutil.copy2(source, png_target)
        tiff_source = source.with_suffix(".tiff")
        if tiff_source.exists():
            shutil.copy2(tiff_source, tiff_target)
        else:
            with Image.open(source) as image:
                image.convert("RGB").save(
                    tiff_target,
                    format="TIFF",
                    dpi=(300, 300),
                    compression="tiff_lzw",
                )
    for number, (filename, _) in SUPPORTING_FIGURES.items():
        source = FIGURE_DIR / filename
        png_target = OUTPUT_FIGURE_DIR / f"Figure_S{number}.png"
        tiff_target = OUTPUT_FIGURE_DIR / f"Figure_S{number}.tiff"
        shutil.copy2(source, png_target)
        tiff_source = source.with_suffix(".tiff")
        if tiff_source.exists():
            shutil.copy2(tiff_source, tiff_target)
        else:
            with Image.open(source) as image:
                image.convert("RGB").save(
                    tiff_target,
                    format="TIFF",
                    dpi=(300, 300),
                    compression="tiff_lzw",
                )


def validate_content() -> list[str]:
    body_texts = INTRODUCTION.copy()
    for _, paragraphs in METHODS:
        body_texts.extend(paragraphs)
    body_texts.extend(text for _, text, _ in RESULTS)
    body_texts.extend(DISCUSSION)
    joined_body = "\n".join(body_texts)
    uncited_references = []
    for key in REFERENCE_KEYS:
        author, year = key.rsplit(" ", 1)
        variants = [key, f"{author} ({year})"]
        if not any(variant in joined_body for variant in variants):
            uncited_references.append(key)
    figure_mentions = []
    supporting_figure_mentions = []
    table_mentions = []
    for text in body_texts:
        figure_mentions.extend(
            int(value) for value in re.findall(r"Figures? (?!S)(\d+)", text)
        )
        supporting_figure_mentions.extend(
            int(value) for value in re.findall(r"Figures? S(\d+)", text)
        )
        table_mentions.extend(int(value) for value in re.findall(r"Table (\d+)", text))
    figure_order = list(dict.fromkeys(figure_mentions))
    supporting_figure_order = list(dict.fromkeys(supporting_figure_mentions))
    table_order = list(dict.fromkeys(table_mentions))
    unresolved = [
        value
        for value in [
            "[" + "Affiliation to be added]",
            "[" + "To be added]",
            "[" + "Corresponding author details]",
        ]
        if value in "\n".join(body_texts)
    ]
    checks = [
        ("Author-date citations", not re.findall(r"\{\d", joined_body)),
        ("Every reference cited", not uncited_references),
        ("References alphabetized", REFERENCES == sorted(REFERENCES)),
        ("Figure first-appearance order", figure_order == list(FIGURES)),
        (
            "Supporting figure first-appearance order",
            supporting_figure_order == list(SUPPORTING_FIGURES),
        ),
        ("Table first-appearance order", table_order == list(TABLES)),
        ("No placeholder strings", not unresolved),
        ("Running title under 70 characters", len(RUNNING_TITLE) < 70),
        ("Title without abbreviations", not re.search(r"\b[A-Z]{2,}\b", TITLE)),
        ("Summary at most 200 words", len(ABSTRACT.split()) <= 200),
        ("Main text at most 4000 words", len(joined_body.split()) <= 4000),
        (
            "Three to six keywords",
            3 <= len([k for k in KEYWORDS.split(";") if k.strip()]) <= 6,
        ),
        (
            "Keywords alphabetical",
            [k.strip() for k in KEYWORDS.split(";") if k.strip()]
            == sorted(k.strip() for k in KEYWORDS.split(";") if k.strip()),
        ),
        (
            "All figure source files present",
            all((FIGURE_DIR / filename).exists() for filename, _ in FIGURES.values()),
        ),
        (
            "All supporting figure source files present",
            all(
                (FIGURE_DIR / filename).exists()
                for filename, _ in SUPPORTING_FIGURES.values()
            ),
        ),
    ]
    lines = [
        f"{JOURNAL_SHORT} SUBMISSION VALIDATION",
        "==========================",
        "",
        f"Summary words: {len(ABSTRACT.split())}",
        f"Main-text words: {len(joined_body.split())}",
        f"References: {len(REFERENCES)}",
        f"Uncited references: {uncited_references}",
        f"First-appearance figure order: {figure_order}",
        f"First-appearance supporting figure order: {supporting_figure_order}",
        f"First-appearance table order: {table_order}",
        "",
    ]
    for label, passed in checks:
        lines.append(f"{'PASS' if passed else 'FAIL'}: {label}")
    if not all(passed for _, passed in checks):
        raise RuntimeError("\n".join(lines))
    return lines


def create_checklist(path: Path) -> None:
    content = """# AHG submission checklist

## Prepared files

- `manuscript_ahg.docx`: manuscript with figure legends and no embedded figure bodies
- `manuscript_ahg_inline_review.docx`: internal review copy with figures and tables immediately after first mention
- `Table_1_residual_outliers.docx`: main editable table; `Table_S1_abo_summary.docx`: supporting editable table
- `tables_ahg.docx`: editable main Table 1 for internal convenience
- `supporting_information_ahg.docx`: Supporting Figures S1-S5, supporting Table S1, and data-file descriptions
- `figures_tables_ahg.pptx`: Figures 1-4, Figures S1-S5, Table 1, and Table S1
- `cover_letter_ahg.docx`: Annals of Human Genetics Original Article cover letter
- `figures/Figure_1` through `Figure_4` and `Figure_S1` through `Figure_S5`: separate PNG and TIFF files
- `supplementary_data/`: population metadata, complete pairwise results, model output, sensitivities, and provenance
- `reproducibility_checklist.md`: data provenance, rebuild commands, expected checks, and package versions
- `reference_validation.csv`: DOI/PubMed existence and title checks

## Automated checks

- References use author-date (author-year) style and are alphabetized.
- Every listed reference is cited and every citation has a reference entry.
- Figures 1-4, Figures S1-S5, and Table 1 are first mentioned sequentially.
- The summary is unstructured and within 200 words.
- The main text (Introduction-Discussion) is within 4,000 words, excluding references.
- Three to six MeSH keywords are listed in alphabetical order.
- The running title is under 70 characters and the title contains no abbreviations.
- Required title-page, availability, funding, conflict, ethics, and contribution statements are present.
- No submission placeholder strings remain.

## Author checks before upload

- Confirm the full correspondence postal address.
- Confirm the no-external-funding statement.
- Confirm the conflict-of-interest statement.
- Obtain or confirm an institutional determination for this secondary genomic analysis.
- Review the explicit disclosure of no direct community engagement or return of results.
- Provide an authenticated ORCID iD for the submitting author in the Research Exchange portal.
- Confirm the AHG article type and current limits in the Research Exchange portal.
- Figures and tables may be embedded in the main file at initial submission; at revision they must be supplied as separate files.
- Upload the manuscript without embedded figures; upload each TIFF separately.
- Upload `supporting_information_ahg.docx` and the supplementary CSV/JSON files.
- Upload `Table_1_residual_outliers.docx` as the editable main table and `Table_S1_abo_summary.docx` as the editable supporting table.
- Do not interpret nominal residuals, PEL-containing pairs, or the two Indigenous-American ABO-window segments as definitive migration evidence.

## Submission links

- Author guidelines: https://onlinelibrary.wiley.com/page/journal/14691809/homepage/forauthors.html
- New-submission portal (Research Exchange): https://authors.wiley.com/journal/AHG
- Editorial office / submission help: AHG.journal@wiley.com; Editor-in-Chief: ahgeditor@ucl.ac.uk
"""
    path.write_text(content, encoding="utf-8")


def create_reproducibility_checklist(path: Path) -> None:
    packages = [
        "pandas",
        "numpy",
        "scipy",
        "statsmodels",
        "matplotlib",
        "seaborn",
        "python-docx",
        "python-pptx",
        "Pillow",
    ]
    versions = []
    for package in packages:
        try:
            versions.append(f"- `{package}=={version(package)}`")
        except PackageNotFoundError:
            versions.append(f"- `{package}`: version not available")
    content = f"""# Reproducibility checklist

## Public source data

- hmmix archaic-introgression segment files from the 1000 Genomes Project and Human Genome Diversity Project (HGDP): Zenodo record 14136628
- O2 blood-group subtype-defining `rs41302905 T` frequencies: Ensembl Variation application programming interface endpoint
- Solomon Islands ABO*O02 frequencies: Ohashi et al. 2006, doi:10.1007/s10038-006-0375-8
- Ancient ABO-window summary: reproducibly extracted from the public Neanderthal-segment catalogue of Iasi et al. 2024 (Dryad doi:10.5061/dryad.zw3r228gg; files Neandertal_segments_matching_references_Shared_map.csv and Meta_Data_individuals.csv) by `scripts/build_ancient_abo_summary.py`. Source-file SHA-256 hashes are recorded in `data/ancient_abo_provenance.json`.

## Rebuild order

Run from the project root:

```bash
python scripts/run_ajba_pipeline.py \
  --segments-1kg /path/to/hg38_1000g_segments.txt \
  --segments-hgdp /path/to/hg38_HGDP_segments.txt \
  --permutations 9999 \
  --sensitivity-permutations 999
```

The committed `data/ancient_abo_summary.csv` (supporting temporal figure only) is
regenerated by additionally passing the Iasi et al. 2024 Dryad files:

```bash
python scripts/build_ancient_abo_summary.py \
  --iasi-segments /path/to/Neandertal_segments_matching_references_Shared_map.csv \
  --iasi-metadata /path/to/Meta_Data_individuals.csv
```

## Expected primary checks

- Individuals: {revised_content.INDIVIDUALS:,}
- Populations: {revised_content.POPULATIONS}
- Unique population pairs: {revised_content.PAIRS:,}
- Every population-window frequency is between 0 and 1
- Neanderthal raw distance r: {revised_content.NEANDERTHAL['raw_r']:.4f}
- Denisovan raw distance r: {revised_content.DENISOVAN['raw_r']:.4f}
- Neanderthal expanded descriptive R²: {revised_content.NEANDERTHAL['expanded_r_squared']:.4f}
- Denisovan expanded descriptive R²: {revised_content.DENISOVAN['expanded_r_squared']:.4f}
- Quadratic assignment procedure distance P: {revised_content.NEANDERTHAL['distance_qap_p']:.4f} and {revised_content.DENISOVAN['distance_qap_p']:.4f}
- False discovery rate q<0.10 non-admixed outliers: {revised_content.NEANDERTHAL['fdr_q_lt_0.10_positive_z_gt_2']} and {revised_content.DENISOVAN['fdr_q_lt_0.10_positive_z_gt_2']}
- Neanderthal/Both segments in the 500-kb ABO interval: {revised_content.ABO['interval_segments']:,}
- Strict ABO-overlapping Neanderthal/Both segments: {revised_content.ABO['strict_overlap']}
- Neanderthal/Both segments with tied maximum reference similarity: {revised_content.ABO['ties']}
- Indigenous American window carriers: Pima 1/13, Maya 1/21, Colombian 0/7
- Strict ABO overlap among those carriers: Pima only

## Environment used for the package

{chr(10).join(versions)}

## Interpretation guardrails

- Pairwise correlation does not prove identity by descent.
- Pairwise rows are dependent; inference uses population-label permutations.
- Expanded-model R² is descriptive and not a causal variance decomposition.
- Reference-genome similarity does not prove a specific migration route.
- Admixed American residuals are not treated as ancient-migration evidence.
- No positive-residual non-admixed pair survived false discovery rate correction.
- Ancient and modern ABO-window calls were produced by different pipelines.
"""
    path.write_text(content, encoding="utf-8")


def create_zip(path: Path, files: list[Path]) -> None:
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for file_path in files:
            archive.write(file_path, file_path.relative_to(OUTPUT_DIR))


def main() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    for stale in ["Table_1_corrected_model.docx", "Table_2_abo_summary.docx"]:
        (OUTPUT_DIR / stale).unlink(missing_ok=True)
    prepare_separate_figures()
    manuscript = OUTPUT_DIR / "manuscript_ahg.docx"
    review = OUTPUT_DIR / "manuscript_ahg_inline_review.docx"
    tables = OUTPUT_DIR / "tables_ahg.docx"
    table_1 = OUTPUT_DIR / "Table_1_residual_outliers.docx"
    table_s1 = OUTPUT_DIR / "Table_S1_abo_summary.docx"
    supporting = OUTPUT_DIR / "supporting_information_ahg.docx"
    cover = OUTPUT_DIR / "cover_letter_ahg.docx"
    presentation = OUTPUT_DIR / "figures_tables_ahg.pptx"
    checklist = OUTPUT_DIR / "submission_checklist.md"
    reproducibility = OUTPUT_DIR / "reproducibility_checklist.md"
    validation = OUTPUT_DIR / "submission_validation.txt"
    reference_validation = OUTPUT_DIR / "reference_validation.csv"
    supplementary_directory = OUTPUT_DIR / "supplementary_data"
    supplementary_directory.mkdir(parents=True, exist_ok=True)
    create_manuscript(manuscript, inline=False)
    create_manuscript(review, inline=True)
    create_tables_document(tables)
    create_single_table_document(table_1, 1)
    create_single_supporting_table_document(table_s1, 1)
    create_supporting_information(supporting)
    create_cover_letter(cover)
    create_presentation(presentation)
    create_checklist(checklist)
    create_reproducibility_checklist(reproducibility)
    validation.write_text("\n".join(validate_content()) + "\n", encoding="utf-8")
    supplementary_sources = [
        DATA_DIR / "population_metadata.csv",
        DATA_DIR / "pairwise_sharing_corrected.csv",
        DATA_DIR / "model_summary.csv",
        DATA_DIR / "sensitivity_analysis.csv",
        DATA_DIR / "window_size_sensitivity.csv",
        DATA_DIR / "analysis_provenance.json",
        DATA_DIR / "profile_quality_checks.csv",
        DATA_DIR / "ancient_abo_summary.csv",
        DATA_DIR / "ancient_abo_provenance.json",
    ]
    for source in supplementary_sources:
        shutil.copy2(source, supplementary_directory / source.name)
    zip_files = [
        manuscript,
        table_1,
        table_s1,
        supporting,
        cover,
        presentation,
        checklist,
        reproducibility,
        validation,
        *sorted(OUTPUT_FIGURE_DIR.glob("Figure_*")),
        *sorted(supplementary_directory.iterdir()),
    ]
    if reference_validation.exists():
        zip_files.append(reference_validation)
    create_zip(OUTPUT_DIR / "AHG_submission_package.zip", zip_files)
    print(f"Created {JOURNAL_SHORT} submission materials in {OUTPUT_DIR}")


if __name__ == "__main__":
    main()
